from __future__ import annotations

import json
import os
import sys
import textwrap
import importlib.util
import itertools
import uuid  # FIXED: stable UUID keys for fund lines (Bug 5)
from datetime import date
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

if importlib.util.find_spec("matplotlib") is not None:
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
    MATPLOTLIB_ERROR = ""
else:
    plt = None
    MATPLOTLIB_AVAILABLE = False
    MATPLOTLIB_ERROR = "matplotlib non installé"
import numpy as np
import pandas as pd
import requests
import streamlit as st
import altair as alt
if importlib.util.find_spec("pypfopt") is not None:
    from pypfopt import EfficientFrontier, risk_models, expected_returns
    PYPFOPT_AVAILABLE = True
    PYPFOPT_ERROR = ""
else:
    EfficientFrontier = None
    risk_models = None
    expected_returns = None
    PYPFOPT_AVAILABLE = False
    PYPFOPT_ERROR = "pyportfolioopt non installé"
if importlib.util.find_spec("reportlab") is not None:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak, HRFlowable
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
    REPORTLAB_ERROR = ""
else:
    SimpleDocTemplate = Paragraph = Spacer = Image = Table = TableStyle = PageBreak = HRFlowable = None
    A4 = None
    getSampleStyleSheet = None
    colors = None
    ParagraphStyle = None
    canvas = None
    REPORTLAB_AVAILABLE = False
    REPORTLAB_ERROR = "reportlab non installé"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import pptx.util as pptx_util
    PPTX_AVAILABLE = True
    PPTX_ERROR = ""
except ImportError as e:
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    PPTX_AVAILABLE = False
    PPTX_ERROR = str(e)


# ------------------------------------------------------------
# Constantes & univers de fonds recommandés
# ------------------------------------------------------------
TODAY = pd.Timestamp.today().normalize()
APP_TITLE = "Analyse Patrimoniale"
APP_SUBTITLE = "Outil de conseil en gestion de patrimoine"
ANNUAL_FEE_EURO_PCT = 0.9
ANNUAL_FEE_UC_PCT = 1.2
RISK_FREE_RATE_FALLBACK = 0.026  # 2.6 % — valeur de repli si API Bund indisponible

# Registre des contrats disponibles
# Clé = label affiché à l'utilisateur
# Valeur = dict avec chemin des fichiers
CONTRACTS_REGISTRY: Dict[str, Any] = {
    "Linxea Avenir 2": {
        "path": "linxea/avenir2",
        "funds_filename": "Liste des fonds LINXEA AVENIR 2.csv",
        "euro_funds": {
            "Suravenir Opportunités 2": "suravenir_opportunites2_historique.csv",
            "Suravenir Rendement 2":    "suravenir_rendement2_historique.csv",
        },
        "assureur": "Suravenir",
        "entry_fee_max_pct": 0.0,
    },
    "Linxea Spirit 2": {
        "path": "linxea/spirit2",
        "funds_filename": "SPIRIT2_UC_frais_avec_categories.csv",
        "euro_funds": {
            "Euro Nouvelle Génération": "spirit2_euro_nouvelle_generation.csv",
            "Euro Objectif Climat":     "spirit2_euro_objectif_climat.csv",
        },
        "assureur": "Spirica",
        "entry_fee_max_pct": 0.0,
    },
}

RECO_FUNDS_CORE = [
    ("R-co Valor C EUR", "FR0011253624"),
    ("Vivalor International", "FR0014001LS1"),
    ("CARMIGNAC Investissement A EUR", "FR0010148981"),
    ("FIDELITY FUNDS - WORLD FUND", "LU0069449576"),
    ("CLARTAN VALEURS", "LU1100076550"),
    ("CARMIGNAC PATRIMOINE", "FR0010135103"),
]

RECO_FUNDS_DEF = [
    ("Fonds en euros (EUROFUND)", "EUROFUND"),
    ("SYCOYIELD 2030 RC", "FR001400MCQ6"),
    ("R-Co Target 2029 HY", "FR0014002XJ3"),
    ("Euro Bond 1-3 Years", "LU0321462953"),
]

FUND_NAME_MAP = {isin: name for name, isin in RECO_FUNDS_CORE + RECO_FUNDS_DEF}

# Univers UC explicites (hors fonds en euros)
EQUITY_FUNDS = [isin for _, isin in RECO_FUNDS_CORE]
BOND_FUNDS = [isin for _, isin in RECO_FUNDS_DEF if isin != "EUROFUND"]

# Libellés FR -> codes internes pour l'affectation des versements
ALLOC_LABELS = {
    "Répartition égale": "equal",
    "Personnalisé": "custom",
    "Tout sur une ligne": "single",
}


# ------------------------------------------------------------
# Utils format
# ------------------------------------------------------------

def params_hash(values: Tuple[Any, ...]) -> str:
    try:
        payload = json.dumps(values, default=str, sort_keys=True)
    except Exception:
        payload = repr(values)
    return str(hash(payload))

def to_eur(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return s + " €"


def fmt_date(x: Any) -> str:
    try:
        return pd.Timestamp(x).strftime("%d/%m/%Y")
    except Exception:
        return "—"


def fmt_eur_fr(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return f"{s} €"


def fmt_pct_fr(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return f"{s} %"


# ------------------------------------------------------------
# XIRR
# ------------------------------------------------------------

def _npv(rate: float, cfs: List[Tuple[pd.Timestamp, float]]) -> float:
    t0 = cfs[0][0]
    return sum(cf / ((1 + rate) ** ((t - t0).days / 365.25)) for t, cf in cfs)


def xirr(cash_flows: List[Tuple[pd.Timestamp, float]], guess: float = 0.1) -> Optional[float]:
    if not cash_flows:
        return None
    cfs = sorted(cash_flows, key=lambda x: x[0])
    try:
        r = guess
        for _ in range(100):
            f = _npv(r, cfs)
            h = 1e-6
            f1 = _npv(r + h, cfs)
            d = (f1 - f) / h
            if abs(d) < 1e-12:
                break
            r2 = r - f / d
            if abs(r2 - r) < 1e-9:
                r = r2
                break
            r = r2
        if abs(_npv(r, cfs)) > 1.0:  # NPV non nulle = pas convergé
            return None
        return r
    except Exception:
        return None


# ------------------------------------------------------------
# API EODHD
# ------------------------------------------------------------

def _get_api_key() -> str:
    return st.secrets.get("EODHD_API_KEY", "")


@st.cache_data(show_spinner=False, ttl=3600)
def eodhd_get(path: str, params: Dict[str, Any] | None = None) -> Any:
    # IMPORTANT : ne jamais appeler st.* (spinner, warning, etc.)
    # dans une fonction @st.cache_data — Streamlit intercepte ces appels
    # lors du replay du cache au premier rendu et lève une exception
    # silencieuse qui vide le résultat.
    base = "https://eodhd.com/api"
    token = _get_api_key()
    p = {"api_token": token, "fmt": "json"}
    if params:
        p.update(params)
    r = requests.get(f"{base}{path}", params=p, timeout=20)
    r.raise_for_status()
    try:
        return r.json()
    except Exception:
        return None


@st.cache_data(show_spinner=False, ttl=3600)
def eodhd_search(q: str) -> List[Dict[str, Any]]:
    try:
        js = eodhd_get(f"/search/{q}")
        if isinstance(js, list):
            return js
    except Exception:
        pass
    return []


@st.cache_data(show_spinner=False, ttl=86400)
def _fetch_bund_rate_from_api() -> Optional[float]:
    """
    Tente de récupérer le rendement du Bund allemand 10 ans via EODHD.
    Retourne le taux en décimal (ex: 0.026) ou None si indisponible.
    Essaie plusieurs tickers dans l'ordre en cas d'échec.
    TTL : 24 h (le taux change peu en intra-journalier).
    """
    tickers = ["DE10Y.GBOND", "BUND10Y.INDX", "IRLTLT01DEM156N.FRED"]
    for ticker in tickers:
        try:
            data = eodhd_get(
                f"/eod/{ticker}",
                {"fmt": "json", "order": "d", "limit": 1},
            )
            if isinstance(data, list) and len(data) > 0:
                rate_pct = float(data[0].get("close", 0.0))
                if 0.0 < rate_pct < 15.0:
                    return rate_pct / 100.0
        except Exception:
            continue
    return None


@st.cache_data(show_spinner=False, ttl=86400)
def _fetch_us_treasury_rate_from_api() -> Optional[float]:
    """Récupère le rendement du US Treasury 10 ans via EODHD."""
    tickers = ["US10Y.GBOND", "DGS10.FRED", "TNX.INDX"]
    for ticker in tickers:
        try:
            data = eodhd_get(
                f"/eod/{ticker}",
                {"fmt": "json", "order": "d", "limit": 1},
            )
            if isinstance(data, list) and len(data) > 0:
                rate_pct = float(data[0].get("close", 0.0))
                if 0.0 < rate_pct < 15.0:
                    return rate_pct / 100.0
        except Exception:
            continue
    return None


def _is_leveraged_product(name: str) -> bool:
    """Retourne True si le fonds est un produit à levier/inverse."""
    _kw = ("leveraged", "2x", "3x", "daily (2x", "daily (3x", "short",
           "bear", "ultra", "inverse", "levier")
    return any(k in str(name).lower() for k in _kw)


def get_risk_free_rate() -> float:
    """
    Retourne le taux sans risque actif en décimal.
    Priorité :
    1. API EODHD (Bund ou US Treasury selon choix utilisateur)
    2. Saisie manuelle dans la sidebar (RISK_FREE_RATE_MANUAL, stocké en décimal)
    3. Fallback RISK_FREE_RATE_FALLBACK
    """
    _ref = st.session_state.get("RFR_BOND_REF", "Bund allemand 10 ans")
    if _ref == "US Treasury 10 ans":
        api_rate = _fetch_us_treasury_rate_from_api()
    else:
        api_rate = _fetch_bund_rate_from_api()
    if api_rate is not None:
        return api_rate
    return float(
        st.session_state.get("RISK_FREE_RATE_MANUAL", RISK_FREE_RATE_FALLBACK)
    )


@st.cache_data(show_spinner=False, ttl=86400)
def eodhd_prices_daily(symbol: str) -> pd.DataFrame:
    try:
        js = eodhd_get(f"/eod/{symbol}", params={"period": "d"})
        if not isinstance(js, list) or len(js) == 0:
            return pd.DataFrame()
        df = pd.DataFrame(js)
        df["date"] = pd.to_datetime(df["date"])
        df.set_index("date", inplace=True)
        if "adjusted_close" in df.columns and pd.notnull(df["adjusted_close"]).any():
            df["Close"] = df["adjusted_close"].astype(float)
        elif "close" in df.columns:
            df["Close"] = df["close"].astype(float)
        else:
            return pd.DataFrame()
        return df[["Close"]].sort_index()
    except Exception:
        return pd.DataFrame()


def _symbol_candidates(isin_or_name: str) -> List[str]:
    val = str(isin_or_name).strip()
    if not val:
        return []
    if val.upper() == "EUROFUND":
        return ["EUROFUND"]
    candidates = [f"{val}.EUFUND", f"{val}.FUND", val]
    try:
        res = eodhd_search(val)
        for it in res:
            code = it.get("Code")
            exch = it.get("Exchange")
            if code and exch:
                candidates.append(f"{code}.{exch}")
    except Exception:
        pass
    seen = set()
    uniq = []
    for c in candidates:
        if c not in seen:
            seen.add(c)
            uniq.append(c)
    return uniq


def _get_close_on(df: pd.DataFrame, d: pd.Timestamp) -> float:
    if df.empty:
        return np.nan
    if d in df.index:
        return float(df.loc[d, "Close"])
    after = df.loc[df.index >= d]
    if not after.empty:
        return float(after.iloc[0]["Close"])
    return float(df.iloc[-1]["Close"])


def apply_annual_fee(
    df: pd.DataFrame,
    annual_fee_pct: float,
    buy_date: Optional[pd.Timestamp] = None,  # FIXED: anchor fee deduction to buy date, not fund inception (Bug 1)
) -> pd.DataFrame:
    # ⚠️ USAGE CORRECT de cette fonction :
    # - Toujours fournir buy_date (date d'achat du client), jamais None
    # - Réservée aux calculs de simulation et de performance client
    #   (XIRR, val_now = qty × last_px_net)
    # - NE PAS utiliser pour afficher une VL de marché absolue
    # - Pour la VL de marché brute : utiliser get_current_nav(isin)
    #   ou get_price_series(isin, None, 0.0)
    # - Pour les calculs de corrélation/vol/Sharpe : utiliser
    #   get_series_for_line(..., apply_fees=False)
    if df.empty or annual_fee_pct == 0:
        return df
    df = df.copy()
    fee_rate = float(annual_fee_pct) / 100.0
    # FIXED: use buy_date as base to avoid applying 20+ years of fictitious fees (Bug 1)
    base_date = pd.Timestamp(buy_date).normalize() if buy_date is not None else df.index[0]
    # FIXED: clamp to zero so pre-buy-date rows (if any) are never penalised (Bug 1)
    day_offsets = np.maximum(0.0, (df.index - base_date).days.astype(float))
    fee_factors = (1.0 - fee_rate) ** (day_offsets / 365.25)
    df["Close"] = df["Close"].astype(float).to_numpy() * fee_factors
    return df


# ------------------------------------------------------------
# Loaders données contrat
# ------------------------------------------------------------

# Chemin racine des données contrats
CONTRACT_DATA_ROOT = "data"


@st.cache_data(show_spinner=False, ttl=3600)
def load_contract_funds(contract_path: str, funds_filename: str) -> pd.DataFrame:
    """
    Charge le fichier Liste des fonds d'un contrat.
    Nettoie les frais (virgule→point, supprime %) et les convertit en float.
    Retourne un DataFrame avec colonnes normalisées :
      isin, name, manager, category, sri, fee_uc_pct, fee_contract_pct, fee_total_pct
    """
    path = os.path.join(CONTRACT_DATA_ROOT, contract_path, funds_filename)
    try:
        df = pd.read_csv(path, sep=";", encoding="utf-8-sig", dtype=str)
    except FileNotFoundError:
        st.error(
            f"Fichier référentiel introuvable : {path}\n"
            f"Vérifiez que le fichier existe bien dans le dépôt GitHub "
            f"avec ce nom exact (attention aux espaces et majuscules)."
        )
        return pd.DataFrame()
    except Exception as e:
        st.error(f"Erreur chargement référentiel : {e}")
        return pd.DataFrame()
    df.columns = df.columns.str.strip()

    def _parse_pct(val: str) -> float:
        try:
            return float(str(val).replace("%", "").replace(",", ".").strip())
        except Exception:
            return 0.0

    result = pd.DataFrame({
        "isin":             df["Code ISIN"].str.strip(),
        "name":             df["Libellé du fonds"].str.strip(),
        "manager":          df["Société de gestion"].str.strip(),
        "category":         df["Catégorie Morningstar"].str.strip(),
        "sri":              pd.to_numeric(df["Risque (SRI)"], errors="coerce").fillna(0).astype(int),
        "fee_uc_pct":       df["Frais UC (B) %"].apply(_parse_pct),
        "fee_contract_pct": df["Frais contrat (C) %"].apply(_parse_pct),
        "fee_total_pct":    df["Frais totaux (B+C) %"].apply(_parse_pct),
    })
    return result.dropna(subset=["isin"]).reset_index(drop=True)


@st.cache_data(show_spinner=False, ttl=3600)
def load_euro_fund_history(contract_path: str, fund_filename: str) -> pd.DataFrame:
    """
    Charge l'historique d'un fonds euros.
    Retourne un DataFrame avec colonnes : annee (int), taux_net_publie_pct (float).
    La colonne utilisée pour la simulation est taux_net_publie_% (net de frais,
    avant prélèvements sociaux — c'est le rendement servi sur le contrat).
    """
    path = os.path.join(CONTRACT_DATA_ROOT, contract_path, fund_filename)
    df = pd.read_csv(path, sep=";", encoding="utf-8-sig")
    df.columns = df.columns.str.strip()
    df["annee"] = pd.to_numeric(df["annee"], errors="coerce").astype("Int64")
    df["taux_net_publie_pct"] = pd.to_numeric(
        df["taux_net_publie_%"], errors="coerce"
    )
    return df[["annee", "taux_net_publie_pct"]].dropna().reset_index(drop=True)


def get_euro_fund_avg_rate(history_df: pd.DataFrame, years: int = 5) -> float:
    """
    Calcule la moyenne du taux net publié sur les N dernières années disponibles.
    Utilisé comme taux par défaut du fonds euros dans la simulation.
    """
    if history_df.empty:
        return 2.0
    recent = history_df.nlargest(years, "annee")
    return round(float(recent["taux_net_publie_pct"].mean()), 2)


def _compute_auto_euro_rate(
    history_df: pd.DataFrame,
    start: pd.Timestamp,
    end: pd.Timestamp,
) -> Optional[float]:
    """
    Calcule le taux moyen du fonds en euros sur les années couvertes par la fenêtre
    [start, end]. Retourne None si l'historique est insuffisant.
    """
    if history_df.empty:
        return None
    start_year = int(start.year)
    end_year = int(end.year)
    mask = (history_df["annee"] >= start_year) & (history_df["annee"] <= end_year)
    filtered = history_df[mask]
    if filtered.empty:
        return None
    return round(float(filtered["taux_net_publie_pct"].mean()), 2)


@st.cache_data(show_spinner=False, ttl=86400)
def _mstarpy_nav_series(isin: str) -> pd.DataFrame:
    """Fallback mstarpy désactivé — retourne toujours un DataFrame vide."""
    return pd.DataFrame()


def get_price_series(
    isin_or_name: str, start: Optional[pd.Timestamp], euro_rate: float
) -> Tuple[pd.DataFrame, str, str]:
    """
    Retourne la série de prix BRUTE (sans frais appliqués).
    • EUROFUND : start et euro_rate comptent dans la clé de cache (ils définissent la série).
    • Instruments EODHD : toujours appeler avec start=None pour maximiser le taux de cache ;
      le filtrage et l'application des frais sont délégués à get_price_series_with_fees().
    """
    # FIXED: fees and start-filtering for EODHD removed from cached fn (Correction 1 — cache explosion)
    debug = {"cands": []}
    val = str(isin_or_name).strip()
    if not val:
        return pd.DataFrame(), "", json.dumps(debug)

    # ✅ Fonds en euros — capitalisation annualisée (jours calendaires)
    if val.upper() == "EUROFUND":
        start_dt = (
            pd.Timestamp(start).normalize()
            if start is not None
            else pd.Timestamp("2000-01-03")
        )
        start_dt = max(start_dt, pd.Timestamp("2000-01-03"))

        idx = pd.bdate_range(start=start_dt, end=TODAY, freq="B")
        if len(idx) == 0:
            return pd.DataFrame(), "", "{}"

        # FIXED: vectorise with NumPy (Bug 8) — fees applied in get_price_series_with_fees()
        days = (idx - idx[0]).days.values.astype(float)
        r = float(euro_rate) / 100.0
        close_values = (1.0 + r) ** (days / 365.25)
        df = pd.DataFrame({"Close": close_values}, index=idx)
        return df, "EUROFUND", "{}"

    # Source prioritaire : VL importée manuellement (prend le dessus sur EODHD et mstarpy)
    # get_price_series() n'est pas décorée @st.cache_data — accès direct à st.session_state autorisé
    _manual_store = st.session_state.get("MANUAL_NAV_STORE", {})
    if val in _manual_store:
        return _manual_store[val], val, json.dumps({"source": "manual_upload"})

    # ✅ Instruments EODHD — historique complet, sans filtrage ni frais
    # Passer start=None depuis les appelants pour qu'un seul enregistrement de cache
    # serve toutes les lignes du même fonds quelle que soit leur date d'achat.
    cands = _symbol_candidates(val)
    debug["cands"] = cands

    for sym in cands:
        df = eodhd_prices_daily(sym)
        if not df.empty:
            return df, sym, json.dumps(debug)

    # Fallback mstarpy si EODHD n'a rien retourné
    # Uniquement pour les ISIN valides (12 caractères, 2 premières lettres alpha)
    _is_isin = len(val) == 12 and val[:2].isalpha()
    if _is_isin:
        _mstar_df = _mstarpy_nav_series(val)
        if not _mstar_df.empty:
            return _mstar_df, val, json.dumps({"source": "mstarpy_nav"})

    return pd.DataFrame(), "", json.dumps(debug)


def get_price_series_with_fees(
    isin_or_name: str,
    start: Optional[pd.Timestamp],
    euro_rate: float,
    annual_uc_fee_pct: float = ANNUAL_FEE_UC_PCT,  # frais UC de la ligne (contract-aware)
) -> Tuple[pd.DataFrame, str, str]:
    """
    Wrapper non-caché : appelle get_price_series() (start=None pour les instruments EODHD
    afin de maximiser le taux de cache) puis applique les frais annuels ancrés sur la
    date d'achat.  La fonction cachée est ainsi partagée entre toutes les lignes d'un
    même fonds indépendamment de leur date d'achat individuelle.

    ⚠️ USAGE CORRECT de cette fonction :
    - Toujours fournir start = date d'achat du client (jamais None)
    - Réservée aux calculs de simulation et de performance client
    - NE PAS utiliser pour afficher une VL de marché absolue
    - Pour la VL de marché brute : utiliser get_current_nav(isin)
    - Pour les calculs de corrélation/vol/Sharpe : utiliser
      get_series_for_line(..., apply_fees=False)
    """
    # FIXED: fee application and start-trimming moved outside the cache (Correction 1)
    val = str(isin_or_name).strip()

    if val.upper() == "EUROFUND":
        # EUROFUND : euro_rate est déjà net de tous frais de gestion assureur
        # → pas d'apply_annual_fee pour ne pas doubler-compter les frais
        df, sym, dbg = get_price_series(isin_or_name, start, euro_rate)
        return df, sym, dbg

    # Instrument EODHD : start=None → un seul enregistrement de cache par ISIN
    df, sym, dbg = get_price_series(isin_or_name, None, euro_rate)
    if df.empty:
        return df, sym, dbg
    if start is not None:
        df = df.loc[df.index >= pd.Timestamp(start)]
    if not df.empty:
        df = apply_annual_fee(df, annual_uc_fee_pct, buy_date=start)
    return df, sym, dbg


def get_current_nav(
    isin_or_name: str,
) -> Tuple[float, Optional[pd.Timestamp]]:
    """
    Retourne (dernière_VL_brute, date_de_cette_VL).
    Retourne (np.nan, None) si indisponible.
    Ne jamais appeler apply_annual_fee() ici.
    """
    val = str(isin_or_name).strip().upper()
    if val in ("EUROFUND", "STRUCTURED", ""):
        return np.nan, None
    df, _, _ = get_price_series(isin_or_name, None, 0.0)
    if df.empty:
        return np.nan, None
    last_date = df.index[-1]
    return float(df["Close"].iloc[-1]), pd.Timestamp(last_date)


@st.cache_data(show_spinner=False, ttl=3600)
def structured_series(
    start: pd.Timestamp,
    end: pd.Timestamp,
    annual_rate_pct: float,
    redemption_years: int,
) -> pd.DataFrame:
    """
    Série synthétique autocall (simplifiée) :
    - Prix d'achat = 1.0
    - Plat jusqu'à la date de remboursement estimée
    - Saut à 1 + (rate * years) le jour de remboursement, puis figé
    """
    start_dt = pd.Timestamp(start).normalize()
    end_dt = pd.Timestamp(end).normalize()
    idx = pd.bdate_range(start=start_dt, end=end_dt, freq="B")
    if len(idx) == 0:
        return pd.DataFrame()

    df = pd.DataFrame(index=idx, columns=["Close"], dtype=float)
    df.iloc[0, 0] = 1.0

    r = float(annual_rate_pct) / 100.0
    yrs = int(redemption_years)

    redemption_dt = start_dt + pd.DateOffset(years=yrs)

    # série plate + saut à partir du 1er jour ouvré >= redemption_dt
    redeemed = False
    for i in range(1, len(df)):
        d = df.index[i]
        df.iloc[i, 0] = df.iloc[i - 1, 0]

        if (not redeemed) and (d >= redemption_dt):
            df.iloc[i, 0] = 1.0 + r * yrs
            df.iloc[i:, 0] = df.iloc[i, 0]
            redeemed = True
            break

    # sécurité : propagation si besoin
    for i in range(1, len(df)):
        if pd.isna(df.iloc[i, 0]):
            df.iloc[i, 0] = df.iloc[i - 1, 0]

    return df


def _warn_once(key: str, msg: str, silent: bool = False) -> None:
    seen = st.session_state.get("WARN_ONCE")
    if not isinstance(seen, set):
        if isinstance(seen, list):
            seen = set(seen)
        else:
            seen = set()
    if key in seen:
        return
    seen.add(key)
    st.session_state["WARN_ONCE"] = seen
    if silent:
        _bw = st.session_state.setdefault("PP_BUILD_WARNINGS", [])
        _bw.append(msg)
    else:
        st.warning(msg)


def _safe_struct_params(line: Dict[str, Any]) -> Tuple[float, int]:
    try:
        years = int(line.get("struct_years", 6))
    except Exception:
        years = 6
    years = max(1, years)
    try:
        rate = float(line.get("struct_rate", 8.0))
    except Exception:
        rate = 0.0
    rate = max(0.0, min(rate, 25.0))
    return rate, years


def get_series_for_line(
    line: Dict[str, Any],
    start: Optional[pd.Timestamp],
    euro_rate: float,
    apply_fees: bool = True,  # FIXED: route to fee-aware wrapper or raw cache fn (Correction 1)
) -> Tuple[pd.DataFrame, str]:
    isin_or_name = str(line.get("isin") or line.get("name") or "").strip()
    sym_upper = isin_or_name.upper()

    if sym_upper == "EUROFUND":
        if apply_fees:
            df, _, _ = get_price_series_with_fees("EUROFUND", start, euro_rate)
        else:
            df, _, _ = get_price_series("EUROFUND", start, euro_rate)
        if df.empty:
            _warn_once("series_empty:EUROFUND", "Série indisponible pour le fonds en euros (EUROFUND).")
        return df, "EUROFUND"

    if sym_upper == "STRUCTURED":
        buy_ts = pd.Timestamp(line.get("buy_date") or start or TODAY).normalize()
        rate, years = _safe_struct_params(line)
        df = structured_series(
            start=buy_ts,
            end=TODAY,
            annual_rate_pct=rate,
            redemption_years=years,
        )
        if start is not None:
            df = df.loc[df.index >= pd.Timestamp(start)]
        if df.empty:
            label = line.get("name") or "Produit structuré"
            _warn_once(
                f"series_empty:STRUCTURED:{label}",
                f"Série indisponible pour {label} (STRUCTURED) : ligne ignorée.",
            )
        return df, "STRUCTURED"

    # Appliquer uniquement fee_contract_pct (C) — la VL EODHD est déjà
    # nette de fee_uc_pct (B/TER). Utiliser fee_total_pct uniquement en
    # fallback si fee_contract_pct est absent (ligne sans référentiel contrat).
    _raw_fee_gsl = line.get("fee_contract_pct")
    if _raw_fee_gsl is None or _raw_fee_gsl == "":
        _raw_fee_gsl = line.get("fee_total_pct")
        if _raw_fee_gsl is None or _raw_fee_gsl == "":
            fee_uc = ANNUAL_FEE_UC_PCT
        else:
            try:
                fee_uc = float(_raw_fee_gsl)
            except (TypeError, ValueError):
                fee_uc = ANNUAL_FEE_UC_PCT
    else:
        try:
            fee_uc = float(_raw_fee_gsl)
        except (TypeError, ValueError):
            fee_uc = ANNUAL_FEE_UC_PCT
    # fee_uc == 0.0 explicite → pas de frais appliqués (ETF sans frais contrat)

    if apply_fees:
        df, sym, _ = get_price_series_with_fees(isin_or_name, start, euro_rate,
                                                 annual_uc_fee_pct=fee_uc)
    else:
        df, sym, _ = get_price_series(isin_or_name, None, euro_rate)
        if start is not None and not df.empty:
            df = df.loc[df.index >= pd.Timestamp(start)]
    if df.empty:
        label = line.get("name") or isin_or_name or "Ligne"
        _warn_once(
            f"series_empty:{isin_or_name}",
            f"Série indisponible pour {label} ({isin_or_name}) : ligne ignorée.",
        )
    return df, sym or isin_or_name

# ------------------------------------------------------------
# Alternatives si date < 1ère VL
# ------------------------------------------------------------

def suggest_alternative_funds(buy_date: pd.Timestamp, euro_rate: float) -> List[Tuple[str, str, pd.Timestamp]]:
    """
    Propose des fonds recommandés (core + défensifs) dont la première VL
    est antérieure ou égale à la date d'achat donnée.
    Retourne (nom, isin, date_inception).
    """
    alternatives: List[Tuple[str, str, pd.Timestamp]] = []
    universe = RECO_FUNDS_CORE + RECO_FUNDS_DEF

    for name, isin in universe:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty:
            continue
        inception = df.index.min()
        if inception <= buy_date:
            alternatives.append((name, isin, inception))

    return alternatives


# ------------------------------------------------------------
# Calendrier versements & poids
# ------------------------------------------------------------

def _month_schedule(d0: pd.Timestamp, d1: pd.Timestamp) -> List[pd.Timestamp]:
    if d0 > d1:
        return []
    out = []
    cur = pd.Timestamp(d0.year, d0.month, 1)
    stop = pd.Timestamp(d1.year, d1.month, 1)
    while cur <= stop:
        bdays = pd.bdate_range(start=cur, end=cur + pd.offsets.MonthEnd(0))
        if len(bdays) > 0:
            out.append(bdays[0])
        cur = cur + pd.offsets.MonthBegin(1)
    return out


def _weights_for(
    lines: List[Dict[str, Any]],
    alloc_mode: str,
    custom_weights: Dict[Any, float],
    single_target: Optional[Any],
) -> Dict[Any, float]:
    # FIXED: use stable UUID key (ln.get("id")) with id(ln) fallback for old lines (Bug 5)
    keys = [ln.get("id") or id(ln) for ln in lines]
    if len(keys) == 0:
        return {}
    if alloc_mode == "equal":
        w = 1.0 / len(keys)
        return {k: w for k in keys}
    if alloc_mode == "custom":
        tot = sum(max(0.0, float(custom_weights.get(ln.get("id") or id(ln), 0.0))) for ln in lines)
        if tot <= 0:
            w = 1.0 / len(keys)
            return {k: w for k in keys}
        return {ln.get("id") or id(ln): max(0.0, float(custom_weights.get(ln.get("id") or id(ln), 0.0))) / tot for ln in lines}
    if alloc_mode == "single":
        target = single_target
        return {ln.get("id") or id(ln): (1.0 if (ln.get("id") or id(ln)) == target else 0.0) for ln in lines}
    w = 1.0 / len(keys)
    return {k: w for k in keys}


# ------------------------------------------------------------
# Calcul par ligne (avec frais)
# ------------------------------------------------------------

def compute_line_metrics(
    line: Dict[str, Any], fee_pct: float, euro_rate: float
) -> Tuple[float, float, float]:
    amt_brut = float(line.get("amount_gross", 0.0))
    net_amt = amt_brut * (1.0 - fee_pct / 100.0)
    buy_ts = pd.Timestamp(line.get("buy_date"))
    px_manual = line.get("buy_px", None)

    dfp, sym_used = get_series_for_line(line, buy_ts, euro_rate)
    if dfp.empty:
        return float(net_amt), np.nan, 0.0

    sym_upper = str(sym_used or line.get("isin") or "").upper()
    if sym_upper == "EUROFUND":
        px = _get_close_on(dfp, buy_ts)
    else:
        if px_manual not in (None, "", 0, "0"):
            try:
                px = float(px_manual)
            except Exception:
                px = _get_close_on(dfp, buy_ts)
        else:
            px = _get_close_on(dfp, buy_ts)

    qty = (net_amt / px) if px and px > 0 else 0.0
    return float(net_amt), float(px), float(qty)


# ------------------------------------------------------------
# Simulation d'un portefeuille (avec contrôle 1ère VL)
# + distinction poids mensuels / ponctuels
# ------------------------------------------------------------

def simulate_portfolio(
    lines: List[Dict[str, Any]],
    monthly_amt_gross: float,
    one_amt_gross: float,
    one_date: date,
    alloc_mode: str,
    custom_weights_monthly: Optional[Dict[int, float]],
    custom_weights_oneoff: Optional[Dict[int, float]],
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    portfolio_label: str = "",
) -> Tuple[pd.DataFrame, float, float, float, Optional[float], pd.Timestamp, pd.Timestamp]:
    if not lines:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    price_map: Dict[Any, pd.Series] = {}
    eff_buy_date: Dict[Any, pd.Timestamp] = {}
    buy_price_used: Dict[Any, float] = {}

    invalid_found = False
    date_warnings = st.session_state.setdefault("DATE_WARNINGS", [])

    for ln in lines:
        key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key, avoids rerun ID changes (Bug 5)
        # FIXED (Bug B): lire la série BRUTE (sans frais) pour déterminer l'inception
        # et vérifier la date d'achat ; les frais seront appliqués ci-dessous
        # ancrés sur d_buy et non sur l'inception du fonds.
        df_full, sym = get_series_for_line(ln, None, euro_rate, apply_fees=False)

        # Sécurité
        if df_full.empty:
            continue

        inception = df_full.index.min()
        d_buy = pd.Timestamp(ln["buy_date"])

        if d_buy < inception:
            invalid_found = True
            ln["invalid_date"] = True
            ln["inception_date"] = inception

            alts = suggest_alternative_funds(d_buy, euro_rate)
            if alts:
                alt_lines = [
                    f"- {name} ({isin}), historique depuis le {fmt_date(incep)}"
                    for name, isin, incep in alts
                ]
                alt_msg = "\n".join(alt_lines)
            else:
                alt_msg = "Aucun fonds recommandé ne dispose d'un historique suffisant pour cette date."

            date_warnings.append(
                f"[{portfolio_label}] {ln.get('name','(sans nom)')} "
                f"({ln.get('isin','—')}) :\n"
                f"- Date d'achat saisie : {fmt_date(d_buy)}\n"
                f"- 1ère VL disponible : {fmt_date(inception)}\n\n"
                f"Impossible de simuler ce fonds sur toute la période demandée.\n"
                f"Propositions d'alternatives pour l'analyse historique :\n{alt_msg}"
            )
            continue

        # FIXED (Bug B): appliquer les frais ancrés sur d_buy (date client réelle)
        # FIXED (Bug E): distinguer fee_total_pct absent (fallback) de 0.0 explicite
        # Appliquer uniquement fee_contract_pct (C) — VL EODHD déjà nette du TER.
        # Fallback fee_total_pct si fee_contract_pct absent, puis ANNUAL_FEE_UC_PCT.
        _isin_sim = str(ln.get("isin", "")).upper()
        _raw_fee_sim = ln.get("fee_contract_pct")
        if _raw_fee_sim is None or _raw_fee_sim == "":
            _raw_fee_sim = ln.get("fee_total_pct")
            if _raw_fee_sim is None or _raw_fee_sim == "":
                _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
            else:
                try:
                    _fee_sim = float(_raw_fee_sim)
                except (TypeError, ValueError):
                    _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
        else:
            try:
                _fee_sim = float(_raw_fee_sim)
            except (TypeError, ValueError):
                _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
        # EUROFUND et STRUCTURED restent à 0.0 (inchangé)
        if _fee_sim > 0 and _isin_sim not in ("EUROFUND", "STRUCTURED"):
            df = apply_annual_fee(df_full.copy(), _fee_sim, buy_date=d_buy)
        else:
            df = df_full

        ln["sym_used"] = sym

        if d_buy in df.index:
            px_buy = float(df.loc[d_buy, "Close"])
            eff_dt = d_buy
        else:
            after = df.loc[df.index >= d_buy]
            if after.empty:
                px_buy = float(df.iloc[-1]["Close"])
                eff_dt = df.index[-1]
            else:
                px_buy = float(after.iloc[0]["Close"])
                eff_dt = after.index[0]

        px_manual = ln.get("buy_px", None)
        px_for_qty = float(px_manual) if (px_manual not in (None, "", 0, "0")) else px_buy

        price_map[key_id] = df["Close"].astype(float)
        eff_buy_date[key_id] = eff_dt
        buy_price_used[key_id] = px_for_qty

    if invalid_found and not price_map:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY
    if not price_map:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    start_min = min(eff_buy_date.values())
    start_full = max(eff_buy_date.values())

    bidx = pd.bdate_range(start=start_min, end=TODAY, freq="B")
    prices = pd.DataFrame(index=bidx)
    for key_id, s in price_map.items():
        prices[key_id] = s.reindex(bidx).ffill()

    qty_events = pd.DataFrame(0.0, index=bidx, columns=prices.columns)
    total_brut = 0.0
    total_net = 0.0
    cash_flows: List[Tuple[pd.Timestamp, float]] = []

    # Achats initiaux
    for ln in lines:
        key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
        if key_id not in prices.columns:
            continue
        brut = float(ln.get("amount_gross", 0.0))
        net = brut * (1.0 - fee_pct / 100.0)
        px = float(buy_price_used[key_id])
        dt = eff_buy_date[key_id]
        if brut > 0 and px > 0:
            q = net / px
            tgt = dt if dt in qty_events.index else qty_events.index[qty_events.index >= dt][0]
            qty_events.loc[tgt, key_id] += q
            total_brut += brut
            total_net += net
            cash_flows.append((tgt, -brut))

    # Poids pour versements mensuels / ponctuels
    weights_monthly = _weights_for(
        lines,
        alloc_mode,
        custom_weights_monthly or {},
        single_target,
    )
    weights_oneoff = _weights_for(
        lines,
        alloc_mode,
        custom_weights_oneoff or {},
        single_target,
    )

    # Versement ponctuel
    if one_amt_gross > 0:
        dt = pd.Timestamp(one_date)
        if dt not in qty_events.index:
            after = qty_events.index[qty_events.index >= dt]
            if len(after) > 0:
                dt = after[0]
            else:
                dt = None
        if dt is not None:
            net_amt = one_amt_gross * (1.0 - fee_pct / 100.0)
            for ln in lines:
                key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
                w = weights_oneoff.get(key_id, 0.0)
                if w <= 0 or key_id not in prices.columns:
                    continue
                px = float(prices.loc[dt, key_id])
                if px > 0:
                    qty_events.loc[dt, key_id] += (net_amt * w) / px
            total_brut += float(one_amt_gross)
            total_net += float(net_amt)
            cash_flows.append((dt, -float(one_amt_gross)))

    # Mensuels
    if monthly_amt_gross > 0:
        sched = _month_schedule(start_min, TODAY)
        for dt in sched:
            if dt not in qty_events.index:
                after = qty_events.index[qty_events.index >= dt]
                if len(after) == 0:
                    continue
                dt = after[0]
            net_m = monthly_amt_gross * (1.0 - fee_pct / 100.0)
            for ln in lines:
                key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
                w = weights_monthly.get(key_id, 0.0)
                if w <= 0 or key_id not in prices.columns:
                    continue
                px = float(prices.loc[dt, key_id])
                if px > 0:
                    qty_events.loc[dt, key_id] += (net_m * w) / px
            total_brut += float(monthly_amt_gross)
            total_net += float(net_m)
            cash_flows.append((dt, -float(monthly_amt_gross)))

    qty_cum = qty_events.cumsum()
    values = (qty_cum * prices).sum(axis=1)
    df_val = pd.DataFrame({"Valeur": values})
    final_val = float(df_val["Valeur"].iloc[-1]) if not df_val.empty else 0.0

    cash_flows.append((TODAY, final_val))
    irr = xirr(cash_flows)

    return df_val, total_brut, total_net, final_val, (irr * 100.0 if irr is not None else None), start_min, start_full


@st.cache_data(show_spinner=False, ttl=3600)
def _simulate_portfolio_cached(
    lines_json: str,
    monthly_amt_gross: float,
    one_amt_gross: float,
    one_date_str: str,
    alloc_mode: str,
    custom_weights_monthly_json: str,
    custom_weights_oneoff_json: str,
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    portfolio_label: str = "",
) -> Tuple[pd.DataFrame, float, float, float, Optional[float], pd.Timestamp, pd.Timestamp]:
    """Wrapper caché de simulate_portfolio. Tous les paramètres sont hashables.
    La clé de cache est le hash de la sérialisation JSON des lignes + scalaires.
    Si le portefeuille n'a pas changé entre deux reruns, résultat instantané."""
    lines = json.loads(lines_json)
    custom_weights_monthly = json.loads(custom_weights_monthly_json)
    custom_weights_oneoff = json.loads(custom_weights_oneoff_json)
    one_date = pd.Timestamp(one_date_str).date()
    return simulate_portfolio(
        lines=lines,
        monthly_amt_gross=monthly_amt_gross,
        one_amt_gross=one_amt_gross,
        one_date=one_date,
        alloc_mode=alloc_mode,
        custom_weights_monthly=custom_weights_monthly or None,
        custom_weights_oneoff=custom_weights_oneoff or None,
        single_target=single_target,
        euro_rate=euro_rate,
        fee_pct=fee_pct,
        portfolio_label=portfolio_label,
    )


def _make_sim_args(
    lines: List[Dict[str, Any]],
    monthly_amt: float,
    one_amt: float,
    one_date,
    alloc_mode: str,
    custom_weights_monthly: Optional[Dict],
    custom_weights_oneoff: Optional[Dict],
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    label: str,
) -> dict:
    """Sérialise les arguments pour _simulate_portfolio_cached.
    Convertit les clés en str pour la sérialisation JSON."""
    def _str_keys(d: Optional[Dict]) -> Dict[str, float]:
        if not d:
            return {}
        return {str(k): v for k, v in d.items()}
    return dict(
        lines_json=json.dumps(lines, default=str, sort_keys=True),
        monthly_amt_gross=float(monthly_amt),
        one_amt_gross=float(one_amt),
        one_date_str=str(pd.Timestamp(one_date).date()),
        alloc_mode=str(alloc_mode),
        custom_weights_monthly_json=json.dumps(_str_keys(custom_weights_monthly)),
        custom_weights_oneoff_json=json.dumps(_str_keys(custom_weights_oneoff)),
        single_target=single_target,
        euro_rate=float(euro_rate),
        fee_pct=float(fee_pct),
        portfolio_label=str(label),
    )


# ------------------------------------------------------------
# Cartes lignes (édition / suppression)
# ------------------------------------------------------------

def _line_card(line: Dict[str, Any], idx: int, port_key: str):
    # FIXED: use stable UUID-based key so widget state survives st.rerun() (Bug 5)
    _card_id = line.get("id", str(idx))
    state_key = f"edit_mode_{port_key}_{_card_id}"
    if state_key not in st.session_state:
        st.session_state[state_key] = False

    fee_pct = st.session_state.get("FEE_A", 0.0) if port_key == "A_lines" else st.session_state.get("FEE_B", 0.0)
    # FIXED: use per-portfolio euro rate instead of shared EURO_RATE_PREVIEW (Résidu Bug 4)
    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )
    net_amt, buy_px, qty_disp = compute_line_metrics(line, fee_pct, euro_rate)

    with st.container(border=True):
        cols = st.columns([3, 2, 2, 2, 1])
        with cols[0]:
            st.markdown(f"**{line.get('name','—')}**")
            st.caption(f"ISIN / Code : `{line.get('isin','—')}`")
            st.caption(f"Symbole EODHD : `{line.get('sym_used','—')}`")
            if line.get("invalid_date"):
                st.markdown(
                    f"⚠️ Date d'achat antérieure à la 1ère VL ({fmt_date(line.get('inception_date'))}).",
                )
        with cols[1]:
            st.markdown(f"Investi (brut)\n\n**{to_eur(line.get('amount_gross', 0.0))}**")
            st.caption(f"Net après frais {fee_pct:.1f}% : **{to_eur(net_amt)}**")
            st.caption(f"Date d'achat : {fmt_date(line.get('buy_date'))}")
            if line.get("date_overridden"):
                st.caption("📌 Date individuelle (non liée à la date globale)")
            _buy_ts_check = pd.Timestamp(line.get("buy_date", TODAY))
            if _buy_ts_check > pd.Timestamp(TODAY):
                st.caption("⚠️ Date d'achat dans le futur — simulation non disponible.")
        with cols[2]:
            st.markdown(f"VL d'achat\n\n**{to_eur(buy_px)}**")
            st.caption(f"Quantité : {qty_disp:.6f}")
            if line.get("note"):
                st.caption(line["note"])
        with cols[3]:
            try:
                # FIXED (Bug A + P2): VL brute sans frais + date de cotation
                last, nav_date = get_current_nav(
                    line.get("isin") or line.get("name") or ""
                )
                if last == last:  # not nan
                    date_str = fmt_date(nav_date) if nav_date is not None else "—"
                    st.markdown(f"Dernière VL connue ({date_str}) : **{to_eur(last)}**")
                    st.caption(
                        "VL publiée par EODHD (clôture J-1 à J-3 selon le fonds)"
                    )
                else:
                    st.markdown("Dernière VL connue : —")
            except Exception:
                st.markdown("Dernière VL connue : —")
        with cols[4]:
            if not st.session_state[state_key]:
                # FIXED: use stable UUID-based widget key (Bug 5)
                if st.button("✏️", key=f"edit_{port_key}_{_card_id}", help="Modifier"):
                    st.session_state[state_key] = True
                    st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)
            # FIXED: use stable UUID-based widget key (Bug 5)
            if st.button("🗑️", key=f"del_{port_key}_{_card_id}", help="Supprimer"):
                st.session_state[port_key].pop(idx)
                st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)

        if st.session_state[state_key]:
            # FIXED: use stable UUID-based form key (Bug 5)
            with st.form(key=f"form_edit_{port_key}_{_card_id}", clear_on_submit=False):
                c1, c2, c3, c4 = st.columns([2, 2, 2, 1])
                with c1:
                    new_amount = st.text_input("Montant investi (brut) €", value=str(line.get("amount_gross", "")))
                with c2:
                    new_date = st.date_input("Date d'achat", value=pd.Timestamp(line.get("buy_date")).date())
                with c3:
                    new_px = st.text_input("Prix d'achat (optionnel)", value=str(line.get("buy_px", "")))
                with c4:
                    st.caption(" ")
                    submitted = st.form_submit_button("💾 Enregistrer")
                if submitted:
                    try:
                        amt_gross = float(str(new_amount).replace(" ", "").replace(",", "."))
                        assert amt_gross > 0
                    except Exception:
                        st.warning("Montant brut invalide.")
                        st.stop()
                    buy_ts = pd.Timestamp(new_date)
                    line["amount_gross"] = float(amt_gross)
                    line["buy_date"] = buy_ts
                    line["date_overridden"] = True
                    if new_px.strip():
                        try:
                            line["buy_px"] = float(str(new_px).replace(",", "."))
                        except Exception:
                            line["buy_px"] = ""
                    else:
                        line["buy_px"] = ""
                    line.pop("invalid_date", None)
                    line.pop("inception_date", None)
                    st.session_state[state_key] = False
                    st.success("Ligne mise à jour.")
                    st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)

    # Import VL manuelle si la série est vide (EODHD + mstarpy ont échoué)
    _isin_card = str(line.get("isin") or "").strip()
    _has_data = not get_series_for_line(
        line, pd.Timestamp(line.get("buy_date", TODAY)), euro_rate, apply_fees=False
    )[0].empty
    _has_manual = _isin_card in st.session_state.get("MANUAL_NAV_STORE", {})
    if not _has_data and not _has_manual and _isin_card not in ("EUROFUND", "STRUCTURED", ""):
        with st.expander(
            f"📥 Importer l'historique VL — {line.get('name', _isin_card)}", expanded=False
        ):
            st.caption(
                "Ce fonds n'est pas disponible via EODHD ni Morningstar. "
                "Importez son historique de VL au format CSV (colonnes : date ; vl). "
                "Sources recommandées : OPCVM360, Fundkis, Morningstar Direct."
            )
            _uploaded_nav = st.file_uploader(
                "CSV historique VL (séparateur ; décimale ,)",
                type=["csv"],
                key=f"nav_upload_{_isin_card}_{_card_id}",
            )
            if _uploaded_nav is not None:
                try:
                    _nav_raw = pd.read_csv(
                        _uploaded_nav, sep=";", decimal=",",
                        encoding="utf-8-sig", dtype=str,
                    )
                    _nav_raw.columns = _nav_raw.columns.str.strip().str.lower()
                    _vl_col = next(
                        (c for c in _nav_raw.columns
                         if c in ("vl", "nav", "close", "valeur liquidative", "price")),
                        None,
                    )
                    if _vl_col and "date" in _nav_raw.columns:
                        _nav_raw["date"] = pd.to_datetime(
                            _nav_raw["date"], dayfirst=True, errors="coerce"
                        )
                        _nav_raw[_vl_col] = pd.to_numeric(
                            _nav_raw[_vl_col].str.replace(",", "."), errors="coerce"
                        )
                        _nav_raw = _nav_raw.dropna(subset=["date", _vl_col])
                        _nav_df = (
                            _nav_raw.set_index("date")
                            .rename(columns={_vl_col: "Close"})[["Close"]]
                            .sort_index()
                        )
                        if not _nav_df.empty:
                            _store = st.session_state.get("MANUAL_NAV_STORE", {})
                            _store[_isin_card] = _nav_df
                            st.session_state["MANUAL_NAV_STORE"] = _store
                            st.success(
                                f"✅ {len(_nav_df)} VL importées pour "
                                f"{line.get('name', _isin_card)}. Relancez le calcul."
                            )
                        else:
                            st.error("Fichier valide mais aucune donnée exploitable.")
                    else:
                        st.error(
                            "Colonnes attendues : 'date' et 'vl' (ou 'nav' / 'close'). "
                            f"Colonnes détectées : {list(_nav_raw.columns)}"
                        )
                except Exception as _e:
                    st.error(f"Erreur lecture CSV : {_e}")


def build_positions_dataframe(port_key: str) -> pd.DataFrame:
    """
    Construit un DataFrame par ligne :
    Nom, ISIN, Date d'achat, Net investi, Valeur actuelle, Perf € et Perf %.
    """
    fee_pct = (
        st.session_state.get("FEE_A", 0.0)
        if port_key == "A_lines"
        else st.session_state.get("FEE_B", 0.0)
    )

    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )

    lines = st.session_state.get(port_key, [])
    rows: List[Dict[str, Any]] = []

    for ln in lines:
        buy_ts = pd.Timestamp(ln.get("buy_date"))
        net_amt, buy_px, qty = compute_line_metrics(ln, fee_pct, euro_rate)
        dfl, _ = get_series_for_line(ln, buy_ts, euro_rate)
        if dfl.empty:
            continue

        # FIXED: removed dead `if not dfl.empty` block — unreachable after `continue` above (Bug 9)
        last_px = float(dfl["Close"].iloc[-1])

        val_now = qty * last_px if last_px == last_px else 0.0
        perf_abs = val_now - net_amt
        perf_pct = (val_now / net_amt - 1.0) * 100.0 if net_amt > 0 else np.nan

        # FIXED (Bug D + P2): VL marché brute séparée de val_now ; date ignorée ici
        nav_brute, _ = get_current_nav(ln.get("isin") or "")

        rows.append(
            {
                "Nom": ln.get("name", ""),
                "ISIN / Code": ln.get("isin", ""),
                "Date d'achat": fmt_date(ln.get("buy_date")),
                "Date individuelle": "📌" if ln.get("date_overridden") else "",
                "Net investi €": net_amt,
                "Valeur actuelle €": val_now,
                "VL marché": round(nav_brute, 4) if nav_brute == nav_brute else None,
                "Perf €": perf_abs,
                "Perf %": perf_pct,
            }
        )

    return pd.DataFrame(rows)

# ------------------------------------------------------------
# Tableau synthétique par ligne (un seul tableau par portefeuille)
# ------------------------------------------------------------

def positions_table(title: str, port_key: str):
    """
    Affiche un seul tableau synthétique par portefeuille :
    Nom, ISIN, Date d'achat, Net investi, Valeur actuelle, Perf € et Perf %.
    """
    fee_pct = (
        st.session_state.get("FEE_A", 0.0)
        if port_key == "A_lines"
        else st.session_state.get("FEE_B", 0.0)
    )

    # ✅ Taux fonds euros par portefeuille (au lieu de EURO_RATE_PREVIEW)
    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )

    lines = st.session_state.get(port_key, [])
    rows: List[Dict[str, Any]] = []

    for ln in lines:
        buy_ts = pd.Timestamp(ln.get("buy_date"))

        # Montant net investi, VL d'achat et quantité
        net_amt, buy_px, qty = compute_line_metrics(ln, fee_pct, euro_rate)

        # ✅ IMPORTANT : on récupère la série "depuis buy_ts" pour éviter le mismatch EUROFUND
        dfl, _ = get_series_for_line(ln, buy_ts, euro_rate)
        if dfl.empty:
            continue

        # FIXED: removed dead `if not dfl.empty` block — unreachable after `continue` above (Bug 9)
        last_px = float(dfl["Close"].iloc[-1])

        # Valeur actuelle et performance
        val_now = qty * last_px if last_px == last_px else 0.0
        perf_abs = val_now - net_amt
        perf_pct = (val_now / net_amt - 1.0) * 100.0 if net_amt > 0 else np.nan

        # FIXED (Bug D + P2): VL marché brute séparée de val_now ; date ignorée ici
        nav_brute, _ = get_current_nav(ln.get("isin") or "")

        rows.append(
            {
                "Nom": ln.get("name", ""),
                "ISIN / Code": ln.get("isin", ""),
                "Date d'achat": fmt_date(ln.get("buy_date")),
                "Date individuelle": "📌" if ln.get("date_overridden") else "",
                "Net investi €": net_amt,
                "Valeur actuelle €": val_now,
                "VL marché": round(nav_brute, 4) if nav_brute == nav_brute else None,
                "Perf €": perf_abs,
                "Perf %": perf_pct,
            }
        )

    st.markdown(f"### {title}")
    df = pd.DataFrame(rows)
    if df.empty:
        st.info("Aucune ligne.")
    else:
        st.dataframe(
            df.style.format(
                {
                    "Net investi €": to_eur,
                    "Valeur actuelle €": to_eur,
                    "VL marché": "{:.4f}".format,
                    "Perf €": to_eur,
                    "Perf %": "{:,.2f}%".format,
                }
            ),
            hide_index=True,
            use_container_width=True,
        )
        st.caption("VL marché : dernière cotation brute EODHD (hors frais de gestion)")


def _prepare_pie_df(df_positions: pd.DataFrame, max_items: int = 8, min_pct: float = 0.03) -> pd.DataFrame:
    if df_positions.empty:
        return df_positions
    df = df_positions.copy()
    df = df[df["Valeur actuelle €"] > 0]
    if df.empty:
        return df
    total = df["Valeur actuelle €"].sum()
    df["Part %"] = df["Valeur actuelle €"] / total
    df = df.sort_values("Valeur actuelle €", ascending=False)
    if len(df) > max_items:
        df_main = df.iloc[:max_items].copy()
        df_other = df.iloc[max_items:]
        df_main = pd.concat(
            [
                df_main,
                pd.DataFrame(
                    {
                        "Nom": ["Autres"],
                        "Valeur actuelle €": [df_other["Valeur actuelle €"].sum()],
                        "Part %": [df_other["Valeur actuelle €"].sum() / total],
                    }
                ),
            ],
            ignore_index=True,
        )
        df = df_main
    else:
        small = df[df["Part %"] < min_pct]
        if not small.empty and len(df) > 1:
            df_main = df[df["Part %"] >= min_pct]
            df_other = pd.DataFrame(
                {
                    "Nom": ["Autres"],
                    "Valeur actuelle €": [small["Valeur actuelle €"].sum()],
                    "Part %": [small["Valeur actuelle €"].sum() / total],
                }
            )
            df = pd.concat([df_main, df_other], ignore_index=True)
    df["Part %"] = df["Part %"] * 100.0
    return df


# ------------------------------------------------------------
# Analytics internes : retours, corrélation, volatilité
# ------------------------------------------------------------


def _build_returns_df(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    start_date: Optional[pd.Timestamp] = None,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Construit un DataFrame de rendements journaliers (pct_change)
    pour toutes les lignes du portefeuille avec un historique suffisant.
    Index = dates, colonnes = "Nom (ISIN)".
    """
    cutoff = pd.Timestamp(start_date) if start_date is not None else TODAY - pd.Timedelta(days=365 * 3)
    series_map: Dict[str, pd.Series] = {}

    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()

        # Exclure le fonds euros — capital garanti, corrélation non pertinente
        if isin.upper() == "EUROFUND":
            continue
        # Exclure les structurés — série synthétique non pertinente
        if isin.upper() == "STRUCTURED":
            continue

        key = f"{label} ({isin})" if isin else label

        # FIXED (Bug F): apply_fees=False — les frais biaisent les rendements si
        # ancrés sur l'inception du fonds ; pct_change() ne nécessite pas de frais nets
        df, _ = get_series_for_line(ln, None, euro_rate, apply_fees=False)
        if df.empty:
            continue

        s = df["Close"].astype(float)
        # Vérifier que le fonds a des données AVANT le cutoff (historique complet)
        if cutoff is not None and not s.empty and s.index.min() > cutoff + pd.Timedelta(days=30):
            continue  # Fonds trop récent — exclu pour ne pas fausser les calculs
        s = s[s.index >= cutoff] if cutoff is not None else s
        if s.size < min_points:
            continue

        series_map[key] = s

    if not series_map:
        return pd.DataFrame()

    df_prices = pd.DataFrame(series_map).dropna(how="any")
    if df_prices.shape[0] < min_points:
        return pd.DataFrame()

    returns = df_prices.pct_change().dropna(how="any")
    return returns



def correlation_matrix_from_lines(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    start_date: Optional[pd.Timestamp] = None,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Matrice de corrélation entre les lignes du portefeuille,
    basée sur les rendements journaliers.
    """
    rets = _build_returns_df(lines, euro_rate, start_date=start_date, min_points=min_points)
    if rets.empty:
        return pd.DataFrame()
    return rets.corr()


def volatility_table_from_lines(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    start_date: Optional[pd.Timestamp] = None,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Volatilité annuelle par ligne (et écart-type quotidien).
    """
    rets = _build_returns_df(lines, euro_rate, start_date=start_date, min_points=min_points)
    if rets.empty:
        return pd.DataFrame()

    rows = []
    for col in rets.columns:
        r = rets[col]
        daily_std = float(r.std())
        ann_std = daily_std * np.sqrt(252.0)
        rows.append(
            {
                "Ligne": col,
                "Écart-type quotidien %": daily_std * 100.0,
                "Volatilité annuelle %": ann_std * 100.0,
                "Nombre de points": int(r.count()),
            }
        )
    df = pd.DataFrame(rows)
    return df.sort_values("Volatilité annuelle %", ascending=False)


def portfolio_risk_stats(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    start_date: Optional[pd.Timestamp] = None,
    min_points: int = 60,
    fee_pct: float = 0.0,  # FIXED: explicit fee rate eliminates max(net_A, net_B) proxy (Bug 3)
) -> Optional[Dict[str, float]]:
    """
    Calcule quelques stats globales de risque pour le portefeuille :
    - volatilité annuelle
    - max drawdown sur la période.
    Pondération par montant net investi.
    """
    rets = _build_returns_df(lines, euro_rate, start_date=start_date, min_points=min_points)
    if rets.empty:
        return None

    # FIXED: use caller-supplied fee_pct instead of guessing A vs B from session state (Bug 3)
    net_by_col: Dict[str, float] = {}
    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()
        key = f"{label} ({isin})" if isin else label

        net, _, _ = compute_line_metrics(ln, fee_pct, euro_rate)
        if net > 0:
            net_by_col[key] = net

    # normalisation des poids
    weights = {}
    tot = sum(net_by_col.get(c, 0.0) for c in rets.columns)
    if tot <= 0:
        return None
    for c in rets.columns:
        w = net_by_col.get(c, 0.0) / tot
        weights[c] = w

    # série de rendement du portefeuille
    w_vec = np.array([weights[c] for c in rets.columns])
    rp = rets.to_numpy().dot(w_vec)
    rp = pd.Series(rp, index=rets.index)

    daily_std = float(rp.std())
    vol_ann = daily_std * np.sqrt(252.0)

    # max drawdown
    cum = (1.0 + rp).cumprod()
    running_max = cum.cummax()
    dd = cum / running_max - 1.0
    max_dd = float(dd.min())  # négatif

    return {
        "vol_ann_pct": vol_ann * 100.0,
        "max_dd_pct": max_dd * 100.0,
    }


def compute_diversification_score(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    start_date: Optional[pd.Timestamp] = None,
) -> Optional[Dict[str, Any]]:
    """
    Score de diversification 0–100 basé sur la corrélation moyenne hors-diagonale.
    0 = tout corrélé (fausse diversification), 100 = parfaitement décorrélé.
    Retourne aussi les paires fortement corrélées (doublons) et en vigilance.
    """
    corr = correlation_matrix_from_lines(lines, euro_rate, start_date=start_date)
    if corr.empty or corr.shape[0] < 2:
        n_uc = corr.shape[0] if not corr.empty else 0
        if n_uc == 1:
            return {
                "score": None,
                "avg_corr": 0.0,
                "doublons": [],
                "vigilance": [],
                "n_lines": 1,
                "n_effective": 1,
                "message": "Un seul support UC — pas de diversification à analyser",
            }
        return None

    avg_corr = _avg_offdiag_corr(corr)
    # Score non-linéaire calibré sur les réalités financières
    # Exposant 0.45 : corrélations "normales" (0.50-0.70) donnent des scores corrects
    _raw_score = 100.0 * ((1.0 - min(1.0, max(0.0, avg_corr))) ** 0.45)

    doublons: List[tuple] = []   # > 0.90 : quasi-identiques
    vigilance: List[tuple] = []  # > 0.80 : diversification limitée
    cols = list(corr.columns)
    for i in range(len(cols)):
        for j in range(i + 1, len(cols)):
            c = float(corr.iloc[i, j])
            name_i = cols[i].split(" (")[0][:30]
            name_j = cols[j].split(" (")[0][:30]
            if c > 0.90:
                doublons.append((name_i, name_j, c))
            elif c > 0.80:
                vigilance.append((name_i, name_j, c))

    # Nombre de fonds réellement utiles (diversification effective)
    # Un fonds est "utile" s'il n'est pas corrélé >0.80 avec un fonds déjà retenu
    effective_funds: List[str] = []
    for col_i in cols:
        is_redundant = False
        for col_e in effective_funds:
            if float(corr.loc[col_i, col_e]) > 0.80:
                is_redundant = True
                break
        if not is_redundant:
            effective_funds.append(col_i)
    n_effective = len(effective_funds)

    # Pénalité de redondance adoucie par racine carrée
    # 1 effectif sur 2 → facteur 0.71 (au lieu de 0.50 brut)
    _effectiveness_ratio = (n_effective / corr.shape[0]) ** 0.5 if corr.shape[0] > 0 else 1.0
    score = max(0.0, min(100.0, _raw_score * _effectiveness_ratio))

    return {
        "score": round(score, 0),
        "avg_corr": avg_corr,
        "doublons": doublons,
        "vigilance": vigilance,
        "n_lines": corr.shape[0],
        "n_effective": n_effective,
    }


def _portfolio_weighted_returns(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    min_points: int = 60,
    start_date: Optional[pd.Timestamp] = None,
) -> Optional[pd.Series]:
    """Série de rendements quotidiens pondérés du portefeuille. Réutilisé par Sharpe/Sortino/Beta."""
    rets = _build_returns_df(lines, euro_rate, start_date=start_date, min_points=min_points)
    if rets.empty:
        return None
    net_by_col: Dict[str, float] = {}
    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()
        key = f"{label} ({isin})" if isin else label
        net, _, _ = compute_line_metrics(ln, fee_pct, euro_rate)
        if net > 0:
            net_by_col[key] = net
    tot = sum(net_by_col.get(c, 0.0) for c in rets.columns)
    if tot <= 0:
        return None
    w_vec = np.array([net_by_col.get(c, 0.0) / tot for c in rets.columns])
    rp = pd.Series(rets.to_numpy().dot(w_vec), index=rets.index)
    return rp if len(rp) >= min_points else None


def compute_sharpe_ratio(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
    start_date: Optional[pd.Timestamp] = None,
) -> Optional[float]:
    """
    Ratio de Sharpe annualisé.
    Taux sans risque = euro_rate / 100 si non fourni (alternative naturelle en AV).
    Retourne None si données insuffisantes (< 60 jours) ou volatilité nulle.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points, start_date=start_date)
    if rp is None:
        return None
    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann = rp.mean() * 252.0
    vol_ann = rp.std() * np.sqrt(252.0)
    if vol_ann == 0:
        return None
    return float((ret_ann - rf) / vol_ann)


def compute_sortino_ratio(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
    start_date: Optional[pd.Timestamp] = None,
) -> Optional[float]:
    """
    Ratio de Sortino annualisé — pénalise uniquement la volatilité à la baisse.
    Retourne None si données insuffisantes ou aucun rendement négatif.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points, start_date=start_date)
    if rp is None:
        return None
    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann = rp.mean() * 252.0
    downside = rp[rp < 0]
    if len(downside) == 0:
        return None
    downside_dev = downside.std() * np.sqrt(252.0)
    if downside_dev == 0:
        return None
    return float((ret_ann - rf) / downside_dev)


def compute_beta_alpha(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    benchmark_symbol: str = "CW8.PA",
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
    start_date: Optional[pd.Timestamp] = None,
) -> Optional[Dict[str, Any]]:
    """
    Bêta et Alpha de Jensen par rapport à l'indice de référence.
    Retourne {"beta": float, "alpha_pct": float, "benchmark_name": str} ou None.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points, start_date=start_date)
    if rp is None:
        return None

    # Récupérer la série benchmark avec fallbacks
    df_bench: pd.DataFrame = pd.DataFrame()
    bench_name = benchmark_symbol
    for sym in [benchmark_symbol, "CW8.PA", "IWDA.AS"]:
        df_b, _, _ = get_price_series(sym, None, 0.0)
        if not df_b.empty:
            df_bench = df_b
            bench_name = sym
            break
    if df_bench.empty:
        return None

    rb = df_bench["Close"].astype(float).pct_change().dropna()
    rb.index = pd.DatetimeIndex([pd.Timestamp(x).normalize() for x in rb.index])

    # Aligner les deux séries sur les dates communes
    common = rp.index.intersection(rb.index)
    if len(common) < min_points:
        return None
    rp_c = rp.loc[common]
    rb_c = rb.loc[common]

    var_b = float(rb_c.var())
    if var_b == 0:
        return None
    beta = float(np.cov(rp_c.values, rb_c.values)[0, 1] / var_b)

    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann_p = rp_c.mean() * 252.0
    ret_ann_b = rb_c.mean() * 252.0
    alpha_pct = float((ret_ann_p - (rf + beta * (ret_ann_b - rf))) * 100.0)

    return {"beta": beta, "alpha_pct": alpha_pct, "benchmark_name": bench_name}


def _corr_heatmap_chart(corr: pd.DataFrame, title: str) -> Optional[alt.Chart]:
    """
    Heatmap de corrélation — triangle inférieur uniquement,
    palette divergente RdYlGn inversée (rouge=corrélé, vert=décorrélé).
    """
    if corr.empty or corr.shape[0] < 2:
        return None
    labels = list(corr.index)
    rows = []
    for i, r in enumerate(labels):
        for j, c in enumerate(labels):
            if j <= i:  # triangle inférieur + diagonale
                rows.append({
                    "Ligne1": r,
                    "Ligne2": c,
                    "corr": float(corr.loc[r, c]),
                    "sort_i": i,
                    "sort_j": j,
                })
    df_melt = pd.DataFrame(rows)
    base = alt.Chart(df_melt).encode(
        x=alt.X("Ligne1:O", sort=labels, title="", axis=alt.Axis(labelAngle=-35)),
        y=alt.Y("Ligne2:O", sort=labels[::-1], title=""),
    )
    heat = base.mark_rect().encode(
        color=alt.Color(
            "corr:Q",
            scale=alt.Scale(
                domain=[-1, 0, 1],
                range=["#1a9641", "#ffffbf", "#d7191c"],
                type="linear",
            ),
            legend=alt.Legend(title="Corrélation", format=".2f"),
        ),
        tooltip=[
            alt.Tooltip("Ligne1:N", title="Ligne 1"),
            alt.Tooltip("Ligne2:N", title="Ligne 2"),
            alt.Tooltip("corr:Q", title="Corrélation", format=".2f"),
        ],
    )
    text = base.mark_text(
        baseline="middle",
        fontSize=11,
        fontWeight="bold",
    ).encode(
        text=alt.Text("corr:Q", format=".2f"),
        color=alt.condition(
            "datum.corr > 0.5 || datum.corr < -0.5",
            alt.value("white"),
            alt.value("#333333"),
        ),
    )
    return (heat + text).properties(title=title, height=max(220, len(labels) * 48))

# ------------------------------------------------------------
# Blocs de saisie : soit fonds recommandés, soit saisie libre
# FIXED: removed 2 duplicate comment blocks (Bug 9)
# ------------------------------------------------------------

def _add_from_reco_block(port_key: str, label: str):
    st.subheader(label)

    cat = st.selectbox(
        "Catégorie",
        ["Core (référence)", "Défensif", "Produits structurés"],
        key=f"reco_cat_{port_key}",
    )

    # ✅ Date d'achat centralisée (versement initial uniquement)
    buy_date = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if port_key == "A_lines"
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    # ============================
    # CAS 1 — PRODUIT STRUCTURÉ
    # ============================
    if cat == "Produits structurés":
        st.markdown("### Produit structuré (Autocall)")

        c1, c2 = st.columns(2)
        with c1:
            amount = st.text_input(
                "Montant investi (brut) €",
                value="",
                key=f"struct_amt_{port_key}",
            )
        with c2:
            struct_years = st.number_input(
                "Durée estimée avant remboursement (années)",
                min_value=1,
                max_value=12,
                value=6,
                step=1,
                key=f"struct_years_{port_key}",
            )

        struct_rate = st.number_input(
            "Rendement annuel estimé (%)",
            min_value=0.0,
            max_value=25.0,
            value=8.0,
            step=0.10,
            key=f"struct_rate_{port_key}",
        )

        st.caption(
            f"Date d'investissement initiale : {pd.Timestamp(buy_date).strftime('%d/%m/%Y')}"
        )

        if st.button("➕ Ajouter le produit structuré", key=f"struct_add_{port_key}"):
            try:
                amt = float(str(amount).replace(" ", "").replace(",", "."))
                assert amt > 0
            except Exception:
                st.warning("Montant invalide.")
                return

            ln = {
                "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
                "name": f"Produit structuré ({struct_rate:.2f}% / {int(struct_years)} ans)",
                "isin": "STRUCTURED",
                "amount_gross": float(amt),
                "buy_date": pd.Timestamp(buy_date),
                "buy_px": 1.0,
                "struct_rate": float(struct_rate),
                "struct_years": int(struct_years),
                "note": "",
                "sym_used": "STRUCTURED",
            }
            st.session_state[port_key].append(ln)
            st.toast("✅ Produit structuré ajouté.", icon="✅")
            st.rerun()
        return  # ✅ IMPORTANT : on sort de la fonction pour ne pas afficher la partie fonds

    # ============================
    # CAS 2 — FONDS CLASSIQUES
    # ============================
    if cat == "Core (référence)":
        fonds_list = RECO_FUNDS_CORE
    else:
        fonds_list = RECO_FUNDS_DEF

    options = [f"{nm} ({isin})" for nm, isin in fonds_list]
    choice = st.selectbox("Fonds recommandé", options, key=f"reco_choice_{port_key}")
    idx = options.index(choice) if choice in options else 0
    name, isin = fonds_list[idx]

    c1, c2 = st.columns([2, 2])
    with c1:
        amount = st.text_input("Montant investi (brut) €", value="", key=f"reco_amt_{port_key}")
    with c2:
        st.caption(f"Date d'achat (versement initial) : {pd.Timestamp(buy_date).strftime('%d/%m/%Y')}")

    px = st.text_input("Prix d'achat (optionnel)", value="", key=f"reco_px_{port_key}")

    if st.button("➕ Ajouter ce fonds recommandé", key=f"reco_add_{port_key}"):
        try:
            amt = float(str(amount).replace(" ", "").replace(",", "."))
            assert amt > 0
        except Exception:
            st.warning("Montant invalide.")
            return

        ln = {
            "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
            "name": name,
            "isin": isin,
            "amount_gross": float(amt),
            "buy_date": pd.Timestamp(buy_date),
            "buy_px": float(str(px).replace(",", ".")) if px.strip() else "",
            "note": "",
            "sym_used": "",
        }
        st.session_state[port_key].append(ln)
        st.toast("✅ Fonds recommandé ajouté.", icon="✅")
        st.rerun()


def _add_line_form_free(port_key: str, label: str):
    st.subheader(label)

    # ✅ Date d'achat centralisée (versement initial)
    buy_date_central = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if port_key == "A_lines"
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    with st.form(key=f"form_add_free_{port_key}", clear_on_submit=False):
        c1, c2 = st.columns([3, 2])

        with c1:
            name = st.text_input("Nom du fonds (libre)", value="")
            isin = st.text_input("ISIN ou code (peut être 'EUROFUND')", value="")

        with c2:
            amount = st.text_input("Montant investi (brut) €", value="")
            st.caption(
                f"Date d'achat (versement initial) : "
                f"{pd.Timestamp(buy_date_central).strftime('%d/%m/%Y')}"
            )

        px = st.text_input("Prix d'achat (optionnel)", value="")
        note = st.text_input("Note (optionnel)", value="")
        add_btn = st.form_submit_button("➕ Ajouter cette ligne")

    if not add_btn:
        return

    isin_final = isin.strip()
    name_final = name.strip()

    # Si nom vide mais ISIN renseigné : tentative de récupération du nom
    if not name_final and isin_final:
        res = eodhd_search(isin_final)
        match = None
        for it in res:
            if it.get("ISIN") == isin_final:
                match = it
                break
        if match is None and res:
            match = res[0]
        if match:
            name_final = match.get("Name", isin_final)

    if not name_final and isin_final.upper() == "EUROFUND":
        name_final = "Fonds en euros (EUROFUND)"

    if not name_final:
        name_final = isin_final or "—"

    try:
        amt = float(str(amount).replace(" ", "").replace(",", "."))
        assert amt > 0
    except Exception:
        st.warning("Montant invalide.")
        return

    ln = {
        "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
        "name": name_final,
        "isin": isin_final or name_final,
        "amount_gross": float(amt),
        "buy_date": pd.Timestamp(buy_date_central),  # ✅ applique la date centrale
        "buy_px": float(str(px).replace(",", ".")) if px.strip() else "",
        "note": note.strip(),
        "sym_used": "",
    }

    st.session_state[port_key].append(ln)
    st.success("Ligne ajoutée.")


@st.cache_data(show_spinner=False, ttl=3600)
def _build_returns_matrix(
    isins: Tuple[str, ...],
    euro_rate: float,
    start: pd.Timestamp,
    end: pd.Timestamp,
    min_points: int = 60,
) -> Tuple[pd.DataFrame, List[str], Dict[str, str]]:
    series_map: Dict[str, pd.Series] = {}
    warnings: List[str] = []
    status: Dict[str, str] = {}

    for isin in isins:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty or df["Close"].dropna().shape[0] < min_points:
            warnings.append(f"{isin} : historique insuffisant")
            status[isin] = "insufficient"
            continue
        s = df["Close"].astype(float)
        s = s[(s.index >= start) & (s.index <= end)]
        if s.dropna().shape[0] < min_points:
            warnings.append(f"{isin} : historique insuffisant sur la fenêtre")
            status[isin] = "insufficient"
            continue
        series_map[isin] = s
        status[isin] = "ok"

    if not series_map:
        return pd.DataFrame(), warnings, status

    prices = pd.DataFrame(series_map).ffill().dropna(how="any")
    if prices.shape[0] < min_points:
        return pd.DataFrame(), warnings, status
    returns = prices.pct_change().dropna(how="any")
    return returns, warnings, status


def _suggest_weights(
    returns: pd.DataFrame,
    max_weight: float,
    min_funds: int,
) -> Dict[str, float]:
    if returns.empty:
        return {}
    corr = returns.corr()
    vols = returns.std() * np.sqrt(252.0)
    avg_corr = corr.mean()
    ranked = avg_corr.sort_values().index.tolist()
    min_funds = max(1, min(min_funds, len(ranked)))
    selected = ranked[:min_funds]
    inv_vol = 1 / vols[selected]
    weights = (inv_vol / inv_vol.sum()).to_dict()

    # cap weights if needed
    if max_weight > 0 and max_weight < 1:
        for _ in range(10):
            over = {k: v for k, v in weights.items() if v > max_weight}
            if not over:
                break
            excess = sum(v - max_weight for v in over.values())
            for k in over:
                weights[k] = max_weight
            remaining = {k: v for k, v in weights.items() if k not in over}
            if not remaining:
                break
            total_remaining = sum(remaining.values())
            for k in remaining:
                weights[k] += excess * (remaining[k] / total_remaining)

    total = sum(weights.values())
    if total > 0:
        weights = {k: v / total for k, v in weights.items()}
    return weights


def _round_allocations(amounts: Dict[str, float]) -> Dict[str, int]:
    floors = {k: int(np.floor(v)) for k, v in amounts.items()}
    remainder = int(round(sum(amounts.values()) - sum(floors.values())))
    if remainder <= 0:
        return floors
    fractions = sorted(
        ((k, amounts[k] - floors[k]) for k in amounts),
        key=lambda x: x[1],
        reverse=True,
    )
    for i in range(min(remainder, len(fractions))):
        k, _ = fractions[i]
        floors[k] += 1
    return floors


def _apply_weight_caps(weights: Dict[str, float], max_weight: float) -> Dict[str, float]:
    if not weights or max_weight <= 0:
        return weights
    weights = weights.copy()
    for _ in range(10):
        over = {k: v for k, v in weights.items() if v > max_weight}
        if not over:
            break
        excess = sum(v - max_weight for v in over.values())
        for k in over:
            weights[k] = max_weight
        remaining = {k: v for k, v in weights.items() if k not in over}
        if not remaining or excess <= 0:
            break
        total_remaining = sum(remaining.values())
        for k in remaining:
            weights[k] += excess * (remaining[k] / total_remaining)
    total = sum(weights.values())
    if total > 0:
        weights = {k: v / total for k, v in weights.items()}
    return weights


def _apply_min_floor_preserve_count(
    weights: Dict[str, float],
    floor: float,
) -> Dict[str, float]:
    if not weights:
        return {}
    keys = list(weights.keys())
    n = len(keys)
    if n == 0:
        return {}

    cleaned: Dict[str, float] = {}
    for k, v in weights.items():
        try:
            w = float(v)
        except Exception:
            w = 0.0
        if not np.isfinite(w) or w < 0:
            w = 0.0
        cleaned[k] = w

    total = sum(cleaned.values())
    if total <= 0:
        base = {k: 1.0 / n for k in keys}
    else:
        base = {k: cleaned[k] / total for k in keys}

    max_floor = (1.0 / n) - 1e-9
    eff_floor = min(max(floor, 0.0), max_floor) if max_floor > 0 else 0.0
    residual = 1.0 - n * eff_floor
    if residual < 0:
        residual = 0.0

    out = {k: eff_floor + residual * base[k] for k in keys}
    return out


def _round_weights_to_step_preserve_count(
    weights: Dict[str, float],
    step: float,
    min_w: float,
    max_w: Optional[float] = None,
) -> Dict[str, float]:
    if not weights:
        return {}
    if step <= 0:
        return weights

    keys = list(weights.keys())
    n = len(keys)
    if n == 0:
        return {}

    def _round_half_up(v: float) -> int:
        return int(np.floor((v / step) + 0.5 + 1e-12))

    eff_min = min_w
    if n > 0:
        max_floor = (1.0 / n) - 1e-9
        if max_floor > 0:
            eff_min = min(min_w, max_floor)
        else:
            eff_min = 0.0
    min_units = int(np.ceil(eff_min / step)) if eff_min > 0 else 0
    max_units = int(np.floor(max_w / step)) if (max_w is not None and max_w > 0) else None
    if max_units is not None and max_units < min_units:
        max_units = min_units
    target_units = int(round(1.0 / step))

    units: Dict[str, int] = {}
    for k in keys:
        u = _round_half_up(weights.get(k, 0.0))
        if u < min_units:
            u = min_units
        if max_units is not None and u > max_units:
            u = max_units
        units[k] = u

    delta = target_units - sum(units.values())
    if delta > 0:
        while delta > 0:
            candidates = [k for k in keys if max_units is None or units[k] < max_units]
            if not candidates:
                break
            k = max(candidates, key=lambda x: units[x])
            units[k] += 1
            delta -= 1
    elif delta < 0:
        while delta < 0:
            candidates = [k for k in keys if units[k] > min_units]
            if not candidates:
                break
            k = max(candidates, key=lambda x: units[x])
            units[k] -= 1
            delta += 1

    if delta != 0:
        k = max(keys, key=lambda x: units[x])
        units[k] = max(units[k] + delta, min_units)

    result = {k: units[k] * step for k in keys}
    total = sum(result.values())
    if total > 0 and abs(total - 1.0) > 1e-9:
        result = {k: v / total for k, v in result.items()}
    return result


def _apply_practical_constraints(
    weights: Dict[str, float],
    min_w: float = 0.10,
    step: float = 0.05,
    max_w: Optional[float] = None,
) -> Dict[str, float]:
    if not weights:
        return {}

    floored = _apply_min_floor_preserve_count(weights, min_w)

    if max_w is not None and max_w > 0:
        floored = _apply_weight_caps(floored, max_w)
        total = sum(floored.values())
        if total > 0:
            floored = {k: v / total for k, v in floored.items()}

    rounded = _round_weights_to_step_preserve_count(floored, step, min_w, max_w=max_w)
    return rounded


def _returns_for_isins(
    isins: List[str],
    start: pd.Timestamp,
    end: pd.Timestamp,
    euro_rate: float,
    min_points: int = 60,
) -> Tuple[pd.DataFrame, Dict[str, str]]:
    series_map: Dict[str, pd.Series] = {}
    status: Dict[str, str] = {}
    for isin in isins:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty:
            status[isin] = "insufficient"
            continue
        s = df["Close"].astype(float)
        s = s[(s.index >= start) & (s.index <= end)]
        if s.dropna().shape[0] < min_points:
            status[isin] = "insufficient"
            continue
        series_map[isin] = s
        status[isin] = "ok"
    if not series_map:
        return pd.DataFrame(), status
    prices = pd.DataFrame(series_map).ffill().dropna(how="any")
    if prices.empty or prices.shape[0] < min_points:
        return pd.DataFrame(), status
    returns = prices.pct_change().dropna(how="any")
    return returns, status


def _annualized_stats(returns: pd.DataFrame) -> Tuple[pd.Series, pd.Series]:
    if returns.empty:
        return pd.Series(dtype=float), pd.Series(dtype=float)
    # FIXED: arithmetic compound formula — was mean()*252 (linear, not compounded) (Résidu Bug 2)
    ann_return = (1 + returns.mean()) ** 252 - 1
    ann_vol = returns.std() * np.sqrt(252.0)
    return ann_return, ann_vol


def _avg_correlation(corr: pd.DataFrame) -> float:
    if corr.empty or corr.shape[0] < 2:
        return np.nan
    vals = corr.values[np.triu_indices_from(corr, 1)]
    return float(np.nanmean(vals)) if vals.size else np.nan


def _avg_offdiag_corr(corr: pd.DataFrame) -> float:
    if corr.empty or corr.shape[0] < 2:
        return np.nan
    vals = corr.values[np.triu_indices_from(corr, 1)]
    return float(np.nanmean(vals)) if vals.size else np.nan


def _select_min_corr_subset(
    candidates: List[str],
    returns: pd.DataFrame,
    k: int,
    anchor: Optional[str] = None,
) -> List[str]:
    if k <= 0 or not candidates:
        return []
    if anchor and anchor not in candidates:
        return []
    if k >= len(candidates):
        return candidates
    corr = returns.corr()
    if corr.empty:
        return candidates[:k]
    if anchor:
        pool = [c for c in candidates if c != anchor]
        best_combo = [anchor]
        best_score = None
        max_checks = 2000
        combos = itertools.combinations(pool, k - 1)
        for idx, combo in enumerate(combos):
            if idx >= max_checks:
                break
            subset = [anchor, *combo]
            score = _avg_offdiag_corr(corr.loc[subset, subset])
            if best_score is None or score < best_score:
                best_score = score
                best_combo = subset
        if best_score is None:
            best_combo = [anchor] + pool[: k - 1]
        return best_combo
    max_checks = 2000
    best_combo: List[str] = []
    best_score = None
    for idx, combo in enumerate(itertools.combinations(candidates, k)):
        if idx >= max_checks:
            break
        subset = list(combo)
        score = _avg_offdiag_corr(corr.loc[subset, subset])
        if best_score is None or score < best_score:
            best_score = score
            best_combo = subset
    if not best_combo:
        best_combo = candidates[:k]
    return best_combo


def _greedy_select(
    candidates: List[str],
    returns: pd.DataFrame,
    target_count: int,
    forced: Optional[str] = None,
    corr_penalty: float = 0.35,
) -> List[str]:
    selected: List[str] = []
    if forced and forced in candidates:
        selected.append(forced)
    remaining = [c for c in candidates if c not in selected]
    if returns.empty or target_count <= 0:
        return selected
    ann_return, ann_vol = _annualized_stats(returns[remaining + selected] if remaining else returns)
    rfr = get_risk_free_rate()  # taux sans risque dynamique (Bund 10 ans ou saisie manuelle)
    sharpe = (ann_return - rfr) / ann_vol.replace(0, np.nan)
    corr = returns.corr() if not returns.empty else pd.DataFrame()

    while len(selected) < target_count and remaining:
        best = None
        best_score = None
        for cand in remaining:
            base_sharpe = float(sharpe.get(cand, np.nan))
            if np.isnan(base_sharpe):
                base_sharpe = -1e9
            if selected and not corr.empty:
                avg_corr = float(corr.loc[cand, selected].mean())
            else:
                avg_corr = 0.0
            score = base_sharpe - corr_penalty * avg_corr
            if best_score is None or score > best_score:
                best_score = score
                best = cand
        if best is None:
            break
        selected.append(best)
        remaining = [c for c in remaining if c != best]
    return selected


def _optimize_uc_weights(
    returns: pd.DataFrame,
    objective: str,
    min_weight: float,
    max_weight: float,
    target_vol: Optional[float],
    target_return: Optional[float],
) -> Dict[str, float]:
    if returns.empty:
        return {}
    n_assets = returns.shape[1]
    if n_assets == 0:
        return {}
    bounds = (min_weight, max_weight)
    try:
        if not PYPFOPT_AVAILABLE:
            raise RuntimeError(PYPFOPT_ERROR)
        mu = expected_returns.mean_historical_return(returns, frequency=252)
        cov = risk_models.sample_cov(returns, frequency=252)
        ef = EfficientFrontier(mu, cov, weight_bounds=bounds)
        if objective == "max_sharpe":
            weights = ef.max_sharpe(risk_free_rate=get_risk_free_rate())
        elif objective == "target_vol" and target_vol is not None:
            weights = ef.efficient_risk(target_vol)
        elif objective == "target_return" and target_return is not None:
            weights = ef.efficient_return(target_return)
        else:
            weights = ef.max_sharpe(risk_free_rate=get_risk_free_rate())
        cleaned = ef.clean_weights()
        total = sum(cleaned.values())
        if total > 0:
            return {k: v / total for k, v in cleaned.items()}
    except Exception:
        pass
    equal_weight = 1.0 / n_assets
    return {col: equal_weight for col in returns.columns}


def _select_min_corr_combo(
    returns: pd.DataFrame,
    k: int,
    anchor: Optional[str] = None,
) -> List[str]:
    if returns.empty:
        return []
    cols = list(returns.columns)
    if anchor and anchor not in cols:
        return []
    if k <= 0:
        return []
    if k > len(cols):
        return cols
    corr = returns.corr()
    best_combo: List[str] = []
    best_score = None
    if anchor:
        others = [c for c in cols if c != anchor]
        for combo in itertools.combinations(others, k - 1):
            candidate = [anchor, *combo]
            sub = corr.loc[candidate, candidate].to_numpy()
            triu = sub[np.triu_indices_from(sub, 1)]
            score = float(triu.mean()) if triu.size else 1.0
            if best_score is None or score < best_score:
                best_score = score
                best_combo = candidate
    else:
        for combo in itertools.combinations(cols, k):
            sub = corr.loc[list(combo), list(combo)].to_numpy()
            triu = sub[np.triu_indices_from(sub, 1)]
            score = float(triu.mean()) if triu.size else 1.0
            if best_score is None or score < best_score:
                best_score = score
                best_combo = list(combo)
    return best_combo


def _compute_drawdown(returns: pd.Series) -> Optional[float]:
    if returns.empty:
        return None
    cum = (1.0 + returns).cumprod()
    running_max = cum.cummax()
    dd = cum / running_max - 1.0
    return float(dd.min())


@st.cache_data(show_spinner="📊 Mise à jour du classement des fonds (une fois par semaine)...", ttl=604800)
def compute_category_rankings(
    contract_path: str,
    funds_filename: str,
    euro_rate: float,
    _cache_key: str = "",
) -> Dict[str, List[Dict[str, Any]]]:
    """
    Calcule rendement annualisé, Sharpe et volatilité de chaque fonds
    du référentiel sur 3 ans. Cache pendant 7 jours.
    Retourne {catégorie: [{isin, name, ann_ret_pct, sharpe, vol_pct, ter, sri}, ...]}
    trié par Sharpe décroissant dans chaque catégorie.
    """
    funds_df = load_contract_funds(contract_path, funds_filename)
    if funds_df.empty:
        return {}

    rfr = get_risk_free_rate()
    cutoff = pd.Timestamp.today().normalize() - pd.DateOffset(years=3)
    rankings: Dict[str, List[Dict[str, Any]]] = {}

    for _, row in funds_df.iterrows():
        isin = row["isin"]
        cat = row.get("category", "Autre")
        if not cat or str(cat).strip() == "":
            cat = "Autre"

        try:
            df, _, _ = get_price_series(isin, None, euro_rate)
            if df.empty:
                continue

            s = df["Close"].astype(float)
            s = s[s.index >= cutoff]
            if len(s) < 60:
                continue

            returns = s.pct_change().dropna()
            if len(returns) < 60:
                continue

            ann_ret = float((1 + returns.mean()) ** 252 - 1)
            ann_vol = float(returns.std() * np.sqrt(252))
            sharpe = float((ann_ret - rfr) / ann_vol) if ann_vol > 0.001 else 0.0

            rankings.setdefault(cat, []).append({
                "isin": isin,
                "name": str(row["name"]),
                "ann_ret_pct": round(ann_ret * 100, 2),
                "sharpe": round(sharpe, 2),
                "vol_pct": round(ann_vol * 100, 2),
                "ter": float(row.get("fee_uc_pct", 0)),
                "sri": int(row.get("sri", 0)),
                "is_etf": any(kw in str(row["name"]).lower() for kw in ("etf", "tracker", "ishares", "lyxor", "amundi index", "xtrackers", "vanguard", "invesco")) or float(row.get("fee_uc_pct", 1.0)) < 0.50,
            })
        except Exception:
            continue

    for cat in rankings:
        rankings[cat].sort(key=lambda x: x["sharpe"], reverse=True)

    return rankings


@st.cache_data(show_spinner=False, ttl=3600)
def _compute_rankings_for_period(
    contract_path: str,
    funds_filename: str,
    euro_rate: float,
    start_date_str: str,
    end_date_str: str,
    _cache_key: str = "",
) -> Dict[str, List[Dict[str, Any]]]:
    """
    Calcule rendement annualisé et Sharpe de chaque fonds du référentiel
    sur la période [start_date, end_date]. Cache 1h (les séries de prix
    sous-jacentes sont déjà cachées 24h).
    """
    _funds = load_contract_funds(contract_path, funds_filename)
    if _funds.empty:
        return {}

    _rfr = get_risk_free_rate()
    _start = pd.Timestamp(start_date_str)
    _end = pd.Timestamp(end_date_str)
    _rankings: Dict[str, List[Dict[str, Any]]] = {}

    for _, _row in _funds.iterrows():
        _isin = _row["isin"]
        _cat = str(_row.get("category", "Autre")).strip()
        if not _cat:
            _cat = "Autre"
        try:
            _df, _, _ = get_price_series(_isin, None, euro_rate)
            if _df.empty:
                continue
            _s_full = _df["Close"].astype(float)
            _first_vl = _s_full.index.min()
            _s = _s_full[(_s_full.index >= _start) & (_s_full.index <= _end)]
            if len(_s) < 60:
                continue
            # Exclure si la 1ère VL est plus de 30 jours après le début de la période
            if _first_vl > _start + pd.Timedelta(days=30):
                continue
            # Exclure les produits à levier des classements
            if _is_leveraged_product(str(_row["name"])):
                continue
            _rets = _s.pct_change().dropna()
            if len(_rets) < 60:
                continue
            _ann_ret = float((1 + _rets.mean()) ** 252 - 1)
            _ann_vol = float(_rets.std() * np.sqrt(252))
            _sharpe = float((_ann_ret - _rfr) / _ann_vol) if _ann_vol > 0.001 else 0.0
            _is_etf = (
                any(kw in str(_row["name"]).lower() for kw in
                    ("etf", "tracker", "ishares", "lyxor", "amundi index", "xtrackers", "vanguard", "invesco"))
                or float(_row.get("fee_uc_pct", 1.0)) < 0.50
            )
            _rankings.setdefault(_cat, []).append({
                "isin": _isin,
                "name": str(_row["name"]),
                "ann_ret_pct": round(_ann_ret * 100, 2),
                "sharpe": round(_sharpe, 2),
                "vol_pct": round(_ann_vol * 100, 2),
                "ter": float(_row.get("fee_uc_pct", 0)),
                "sri": int(_row.get("sri", 0)),
                "is_etf": _is_etf,
                "first_vl": _first_vl.strftime("%d/%m/%Y"),
            })
        except Exception:
            continue
    for _cat in _rankings:
        _rankings[_cat].sort(key=lambda x: x["sharpe"], reverse=True)
    return _rankings


def _fund_name(isin: str) -> str:
    return FUND_NAME_MAP.get(isin, isin)


def _safe_fund_label(name: str, isin: str) -> str:
    cleaned = str(name or "").strip()
    if cleaned:
        return f"{cleaned} ({isin})"
    if isin:
        res = eodhd_search(isin)
        if res:
            maybe = res[0].get("Name") or res[0].get("name")
            if maybe:
                return f"{maybe} ({isin})"
    return f"{isin} ({isin})"


def _build_gauge_chart(score: float, label: str) -> Optional[alt.Chart]:
    """Donut Altair avec le score au centre."""
    try:
        _data = pd.DataFrame({
            "category": ["Score", "Restant"],
            "value": [score, max(0.0, 100.0 - score)],
        })
        _gc = "#2E7D32" if score >= 70 else ("#FF9800" if score >= 40 else "#E53935")
        _base = alt.Chart(_data).encode(
            theta=alt.Theta("value:Q", stack=True),
            color=alt.Color("category:N",
                scale=alt.Scale(domain=["Score", "Restant"], range=[_gc, "#E8EEF7"]),
                legend=None,
            ),
        )
        _donut = _base.mark_arc(innerRadius=50, outerRadius=80)
        _text = alt.Chart(pd.DataFrame({"text": [f"{score:.0f}"]})).mark_text(
            size=28, fontWeight="bold", color="#1B2A4A"
        ).encode(text="text:N")
        return (_donut + _text).properties(height=180, title=label)
    except Exception:
        return None


def render_portfolio_builder():
    st.title("Construction de portefeuille optimisé")

    # ── Sélecteur de mode ──────────────────────────────────────
    _is_improve = st.session_state.get("PP_IMPROVE_MODE", False)
    _mode_options = ["✨ Créer un portefeuille", "🔧 Améliorer un portefeuille existant"]
    _mode_default = 1 if _is_improve else 0
    _pp_mode = st.radio(
        "Mode",
        _mode_options,
        index=_mode_default,
        horizontal=True,
        key="PP_BUILD_MODE",
        label_visibility="collapsed",
    )
    _is_improve_mode = "Améliorer" in _pp_mode

    st.session_state.setdefault("PP_RUN", False)
    st.session_state.setdefault("PP_PARAMS_HASH", "")
    st.session_state.setdefault("PP_OPT_START_DATE", (TODAY - pd.DateOffset(years=3)).date())
    st.session_state.setdefault("PP_OPT_END_DATE", TODAY.date())
    st.session_state.setdefault("PP_CONTRACT_LABEL", list(CONTRACTS_REGISTRY.keys())[0])
    st.session_state.setdefault("PP_TOTAL_UC_COUNT", 4)

    # Charger le référentiel du contrat sélectionné (indépendant de render_app)
    pp_contract_label = st.session_state.get("PP_CONTRACT_LABEL", list(CONTRACTS_REGISTRY.keys())[0])
    pp_contract_cfg = CONTRACTS_REGISTRY.get(pp_contract_label, list(CONTRACTS_REGISTRY.values())[0])
    funds_df = load_contract_funds(
        pp_contract_cfg["path"],
        pp_contract_cfg["funds_filename"],
    )
    CONTRACT_FUND_NAMES = dict(zip(funds_df["isin"], funds_df["name"])) if not funds_df.empty else FUND_NAME_MAP.copy()
    CONTRACT_FUND_FEES  = dict(zip(funds_df["isin"], funds_df["fee_total_pct"])) if not funds_df.empty else {}

    profile_map = {
        "Prudent": 50,
        "Équilibré": 30,
        "Dynamique": 15,
        "Offensif": 5,
    }

    # ── Pré-remplissage des paramètres en mode amélioration ──────
    # Flag _PP_IMPROVE_APPLIED : appliqué une seule fois à l'arrivée dans le builder.
    # Utilise = (pas setdefault) car les keys des widgets existent déjà après le 1er rendu.
    if st.session_state.get("PP_IMPROVE_MODE") and not st.session_state.get("_PP_IMPROVE_APPLIED"):
        _source_improve_pf = st.session_state.get("PP_IMPROVE_SOURCE", "client")
        _is_client_imp_pf = _source_improve_pf == "client"

        # Pré-remplir le contrat (key = PP_CONTRACT_LABEL)
        _contract_imp_pf = st.session_state.get(
            "CONTRACT_LABEL_A" if _is_client_imp_pf else "CONTRACT_LABEL_B", ""
        )
        if _contract_imp_pf:
            st.session_state["PP_CONTRACT_LABEL"] = _contract_imp_pf

        # Pré-remplir le budget (key = PP_BUDGET, consumed by pop → value= du widget PP_TOTAL_BUDGET)
        _val_imp_pf = float(st.session_state.get(
            "_LAST_VAL_A" if _is_client_imp_pf else "_LAST_VAL_B", 0
        ) or 0)
        if _val_imp_pf > 0:
            st.session_state["PP_BUDGET"] = _val_imp_pf

        # Pré-remplir le % fonds euros (key = PP_EURO_PCT)
        _lines_imp_pf = st.session_state.get("PP_IMPROVE_LINES", [])
        if _lines_imp_pf:
            _total_imp_pf = sum(float(ln.get("amount_gross", 0) or 0) for ln in _lines_imp_pf)
            _euro_imp_pf = sum(
                float(ln.get("amount_gross", 0) or 0)
                for ln in _lines_imp_pf
                if str(ln.get("isin", "")).upper() == "EUROFUND"
            )
            _euro_pct_imp_pf = int(_euro_imp_pf / _total_imp_pf * 100) if _total_imp_pf > 0 else 30
            st.session_state["PP_EURO_PCT"] = _euro_pct_imp_pf

        # Pré-remplir le taux fonds euros (key = PP_EURO_RATE)
        _euro_rate_imp_pf = float(st.session_state.get(
            "EURO_RATE_A" if _is_client_imp_pf else "EURO_RATE_B", 2.0
        ))
        st.session_state["PP_EURO_RATE"] = _euro_rate_imp_pf

        # Pré-remplir la fenêtre d'analyse (keys = PP_WINDOW_MODE, PP_OPT_START_DATE, PP_OPT_END_DATE)
        _imp_start_dt = st.session_state.get("PP_IMPROVE_START_DATE")
        _imp_end_dt = st.session_state.get("PP_IMPROVE_END_DATE")
        if _imp_start_dt and _imp_end_dt:
            st.session_state["PP_WINDOW_MODE"] = "Dates personnalisées"
            st.session_state["PP_OPT_START_DATE"] = pd.Timestamp(_imp_start_dt).date()
            st.session_state["PP_OPT_END_DATE"] = pd.Timestamp(_imp_end_dt).date()
        else:
            _dates_imp_pf = [ln.get("buy_date") for ln in _lines_imp_pf if ln.get("buy_date")]
            if _dates_imp_pf:
                _min_date_imp_pf = min(pd.Timestamp(d) for d in _dates_imp_pf)
                _years_diff_pf = (pd.Timestamp.today() - _min_date_imp_pf).days / 365.25
                if _years_diff_pf >= 10:
                    st.session_state["PP_WINDOW_MODE"] = "10 ans"
                elif _years_diff_pf >= 5:
                    st.session_state["PP_WINDOW_MODE"] = "5 ans"
                else:
                    st.session_state["PP_WINDOW_MODE"] = "3 ans"

        st.session_state["_PP_IMPROVE_APPLIED"] = True
        st.rerun()

    with st.sidebar:
        st.header("Paramètres clés")

        # ── Contrat & fonds en euros ───────────────────────────
        pp_contract_label = st.selectbox(
            "Contrat",
            list(CONTRACTS_REGISTRY.keys()),
            key="PP_CONTRACT_LABEL",
        )
        pp_contract_cfg = CONTRACTS_REGISTRY[pp_contract_label]
        pp_euro_options = list(pp_contract_cfg.get("euro_funds", {}).keys())
        pp_euro_fund_label = st.selectbox(
            "Fonds en euros",
            pp_euro_options if pp_euro_options else ["—"],
            key="PP_EURO_FUND_LABEL",
        )
        st.markdown("---")

        # ── Profil & budget ────────────────────────────────────
        profile = st.selectbox(
            "Profil de risque client",
            list(profile_map.keys()),
            index=1,
            key="PP_PROFILE",
        )
        _recommended_euro_pct = profile_map[profile]
        st.info(
            f"📋 Allocation fonds en euros recommandée "
            f"pour un profil **{profile}** : **{_recommended_euro_pct}%**"
        )
        euro_pct = int(st.number_input(
            "Part fonds en euros (%)",
            min_value=0,
            max_value=100,
            value=int(st.session_state.get("PP_EURO_PCT", _recommended_euro_pct)),
            step=5,
            key="PP_EURO_PCT",
            help="Valeur recommandée selon le profil. Modifiable selon la situation du client.",
        ))

        # Transférer PP_BUDGET → PP_TOTAL_BUDGET si présent (arrivée depuis le comparateur)
        _pp_budget_from_comparator = st.session_state.pop("PP_BUDGET", None)
        if _pp_budget_from_comparator is not None and "PP_TOTAL_BUDGET" not in st.session_state:
            st.session_state["PP_TOTAL_BUDGET"] = int(_pp_budget_from_comparator)
        total_budget = st.number_input(
            "Montant investi (EUR)",
            min_value=0,
            max_value=10_000_000,
            value=100_000,
            step=10,
            key="PP_TOTAL_BUDGET",
        )

        # ── Fenêtre d'analyse (détermine opt_start / opt_end) ──
        if st.session_state.get("PP_WINDOW_MODE") == "1 an":
            st.session_state["PP_WINDOW_MODE"] = "3 ans"
        opt_window_mode = st.radio(
            "Fenêtre d'analyse",
            ["3 ans", "5 ans", "10 ans", "Dates personnalisées"],
            horizontal=False,
            key="PP_WINDOW_MODE",
        )

        if opt_window_mode == "Dates personnalisées":
            # Initialiser les defaults AVANT le widget si absents
            if "PP_OPT_START_DATE" not in st.session_state:
                st.session_state["PP_OPT_START_DATE"] = (TODAY - pd.DateOffset(years=3)).date()
            if "PP_OPT_END_DATE" not in st.session_state:
                st.session_state["PP_OPT_END_DATE"] = TODAY.date()
            st.date_input("Date de début", key="PP_OPT_START_DATE")
            st.date_input("Date de fin", key="PP_OPT_END_DATE")
            opt_start = pd.Timestamp(st.session_state["PP_OPT_START_DATE"])
            opt_end = pd.Timestamp(st.session_state["PP_OPT_END_DATE"])
        else:
            years_map = {"3 ans": 3, "5 ans": 5, "10 ans": 10}
            opt_years = years_map[opt_window_mode]
            opt_start = TODAY - pd.DateOffset(years=opt_years)
            opt_end = TODAY

        # ── Taux fonds en euros (auto depuis historique) ───────
        st.markdown("---")
        st.markdown("**Rendement fonds en euros**")
        try:
            _pp_ef_filename = pp_contract_cfg["euro_funds"].get(pp_euro_fund_label, "")
            _pp_ef_history = (
                load_euro_fund_history(pp_contract_cfg["path"], _pp_ef_filename)
                if _pp_ef_filename else pd.DataFrame()
            )
        except Exception:
            _pp_ef_history = pd.DataFrame()

        _auto_euro_rate = _compute_auto_euro_rate(_pp_ef_history, opt_start, opt_end)
        if _auto_euro_rate is not None:
            st.success(
                f"Taux net moyen **{pp_euro_fund_label}** "
                f"sur la fenêtre : **{_auto_euro_rate:.2f}%**"
            )
        else:
            st.info("Historique insuffisant sur la fenêtre — saisissez le taux manuellement.")

        _override_euro = st.checkbox(
            "Saisir un taux personnalisé",
            value=False,
            key="PP_OVERRIDE_EURO_RATE",
        )
        if _override_euro:
            euro_rate = float(st.number_input(
                "Rendement annuel du fonds en euros (%)",
                min_value=0.0, max_value=10.0,
                value=float(_auto_euro_rate if _auto_euro_rate is not None else 2.0),
                step=0.10,
                key="PP_EURO_RATE",
            ))
        else:
            euro_rate = float(_auto_euro_rate if _auto_euro_rate is not None else 2.0)
            st.session_state["PP_EURO_RATE"] = euro_rate
            st.caption(f"Taux utilisé : **{euro_rate:.2f}%**")

        # ── Taux sans risque (Sharpe) ──────────────────────────
        st.markdown("---")
        st.markdown("**Taux sans risque (Sharpe)**")
        _rfr_ref_choice = st.selectbox(
            "Obligation de référence",
            ["Bund allemand 10 ans", "US Treasury 10 ans"],
            key="PP_RFR_BOND_REF",
            help="Le taux sans risque est utilisé pour le calcul du ratio de Sharpe.",
        )
        st.session_state["RFR_BOND_REF"] = _rfr_ref_choice

        if _rfr_ref_choice == "US Treasury 10 ans":
            rfr_api = _fetch_us_treasury_rate_from_api()
            _rfr_label = "US Treasury 10 ans"
        else:
            rfr_api = _fetch_bund_rate_from_api()
            _rfr_label = "Bund 10 ans"

        if rfr_api is not None:
            st.success(
                f"📡 {_rfr_label} : **{rfr_api * 100:.2f}%** — récupéré automatiquement"
            )
            st.session_state["RISK_FREE_RATE_SOURCE"] = "api"
            st.session_state["RISK_FREE_RATE_VALUE"] = rfr_api
        else:
            st.warning(f"⚠️ {_rfr_label} indisponible via API")
        default_pct = (
            rfr_api * 100.0
            if rfr_api is not None
            else float(st.session_state.get("RISK_FREE_RATE_MANUAL_PCT",
                                            RISK_FREE_RATE_FALLBACK * 100.0))
        )
        manual_rfr_pct = st.number_input(
            "Taux sans risque (%)",
            min_value=0.0,
            max_value=15.0,
            value=default_pct,
            step=0.05,
            format="%.2f",
            help=(
                f"Taux de référence : rendement de l'obligation {_rfr_label}."
            ),
            key="RFR_MANUAL_INPUT",
        )
        st.session_state["RISK_FREE_RATE_MANUAL_PCT"] = manual_rfr_pct
        st.session_state["RISK_FREE_RATE_MANUAL"] = manual_rfr_pct / 100.0
        if rfr_api is None:
            st.session_state["RISK_FREE_RATE_SOURCE"] = "manual"

    # Lookup global isin -> name (nécessaire pour le diagnostic et les sélections)
    _global_isin_name = {}
    if not funds_df.empty:
        for _, _r in funds_df.iterrows():
            _global_isin_name[_r["isin"]] = _r["name"]

    # ── Diagnostic du portefeuille (mode améliorer) ───────────────
    _rlookup: Dict[str, Dict] = {}
    _rby_cat: Dict[str, List[Dict]] = {}

    if _is_improve_mode and not funds_df.empty:
        st.markdown("### 🔧 Diagnostic de votre portefeuille")
        st.caption(f"Période d'analyse : du {fmt_date(opt_start)} au {fmt_date(opt_end)}")
        _diag_sort = st.selectbox(
            "Trier les suggestions par",
            ["Meilleur rapport rendement/risque", "Meilleur rendement", "Volatilité la plus basse"],
            key="PP_DIAG_SORT_CRITERIA",
        )

        _selected_all_diag = list(st.session_state.get("PP_SELECTED_ACTIONS", [])) + list(st.session_state.get("PP_SELECTED_BONDS", []))

        if not _selected_all_diag:
            st.info("Aucun fonds chargé. Chargez un portefeuille depuis le comparateur ou ajoutez des fonds ci-dessous.")
        else:
            _rankings_diag = _compute_rankings_for_period(
                pp_contract_cfg["path"],
                pp_contract_cfg["funds_filename"],
                float(euro_rate),
                str(opt_start.date()),
                str(opt_end.date()),
                _cache_key=f"{pp_contract_label}_{opt_start.date()}_{opt_end.date()}",
            )

            for _cat_rk, _funds_rk in _rankings_diag.items():
                _rby_cat[_cat_rk] = _funds_rk
                for _f_rk in _funds_rk:
                    _rlookup[_f_rk["isin"]] = {**_f_rk, "category": _cat_rk}

            _euro_rate_diag = float(euro_rate)
            _uc_isins_diag = [i for i in _selected_all_diag if i.upper() not in ("EUROFUND", "STRUCTURED")]
            _diag_div = None
            if len(_uc_isins_diag) >= 2:
                try:
                    _diag_div = compute_diversification_score(
                        [{"isin": i, "name": _global_isin_name.get(i, i), "amount_gross": 1000} for i in _uc_isins_diag],
                        _euro_rate_diag,
                        start_date=opt_start,
                    )
                except Exception:
                    pass

            _doublons_pairs: List[Tuple[str, str, float]] = []
            if _diag_div and _diag_div.get("doublons"):
                _doublons_pairs = _diag_div["doublons"]

            # ── Carte par fonds ──────────────────────────────────
            for _isin_d in _selected_all_diag:
                if _isin_d.upper() in ("EUROFUND", "STRUCTURED"):
                    continue

                _info_d = _rlookup.get(_isin_d)
                _fname_d = _global_isin_name.get(_isin_d, _info_d["name"] if _info_d else _isin_d)

                if not _info_d:
                    with st.container(border=True):
                        st.markdown(f"**📊 {_fname_d}** ({_isin_d})")
                        st.caption("Données de performance non disponibles pour ce fonds.")
                        if st.button("❌ Retirer", key=f"PP_DIAG_REMOVE_{_isin_d}"):
                            for _lk in ("PP_SELECTED_ACTIONS", "PP_SELECTED_BONDS"):
                                _ll = st.session_state.get(_lk, [])
                                _to_rm = [v for v in _ll if v.upper() == _isin_d.upper()]
                                for _tr in _to_rm:
                                    _ll.remove(_tr)
                                st.session_state[_lk] = _ll
                            st.rerun()
                    continue

                _cat_d = _info_d["category"]
                _cat_funds_d = _rby_cat.get(_cat_d, [])
                _rank_d = next((i + 1 for i, f in enumerate(_cat_funds_d) if f["isin"] == _isin_d), None)
                _total_cat_d = len(_cat_funds_d)

                _doublon_with = None
                for _dn1, _dn2, _dc in _doublons_pairs:
                    if _fname_d[:30] in _dn1 or _fname_d[:30] in _dn2:
                        _doublon_with = (_dn2 if _fname_d[:30] in _dn1 else _dn1, _dc)
                        break

                _already_up = set(i.upper() for i in _selected_all_diag)
                _is_bond_d = _is_bond_category(_cat_d)

                # Trier selon le critère choisi
                if _diag_sort == "Meilleur rendement":
                    _cat_funds_sorted = sorted(_cat_funds_d, key=lambda x: x["ann_ret_pct"], reverse=True)
                    _best_d = next(
                        (f for f in _cat_funds_sorted
                         if f["isin"] != _isin_d and f["isin"].upper() not in _already_up
                         and f["ann_ret_pct"] > _info_d["ann_ret_pct"]),
                        None,
                    )
                elif _diag_sort == "Volatilité la plus basse":
                    _cat_funds_sorted = sorted(_cat_funds_d, key=lambda x: x["vol_pct"], reverse=False)
                    _best_d = next(
                        (f for f in _cat_funds_sorted
                         if f["isin"] != _isin_d and f["isin"].upper() not in _already_up
                         and f["vol_pct"] < _info_d["vol_pct"]),
                        None,
                    )
                else:  # Meilleur rapport rendement/risque
                    if _is_bond_d:
                        _cat_funds_sorted = sorted(_cat_funds_d, key=lambda x: x["ann_ret_pct"] / max(x["vol_pct"], 0.1), reverse=True)
                        _curr_ratio = _info_d["ann_ret_pct"] / max(_info_d["vol_pct"], 0.1)
                        _best_d = next(
                            (f for f in _cat_funds_sorted
                             if f["isin"] != _isin_d and f["isin"].upper() not in _already_up
                             and f["ann_ret_pct"] / max(f["vol_pct"], 0.1) > _curr_ratio),
                            None,
                        )
                    else:
                        _cat_funds_sorted = sorted(_cat_funds_d, key=lambda x: x["sharpe"], reverse=True)
                        _best_d = next(
                            (f for f in _cat_funds_sorted
                             if f["isin"] != _isin_d and f["isin"].upper() not in _already_up
                             and f["sharpe"] > _info_d["sharpe"]),
                            None,
                        )

                _best_etf_d = next(
                    (f for f in _cat_funds_d if f.get("is_etf") and f["isin"] != _isin_d
                     and f["isin"].upper() not in _already_up
                     and f["ann_ret_pct"] >= _info_d["ann_ret_pct"] - 2.0),
                    None,
                ) if not _info_d.get("is_etf") else None

                with st.container(border=True):
                    _badges = []
                    if _rank_d and _rank_d == 1:
                        _badges.append("🏆 1er")
                    elif _rank_d:
                        _badges.append(f"rang {_rank_d}/{_total_cat_d}")
                    if _doublon_with:
                        _badges.append("⚠️ doublon")
                    _badge_txt = " · ".join(_badges)

                    st.markdown(f"**📊 {_fname_d}** — {_cat_d}" + (f" ({_badge_txt})" if _badge_txt else ""))
                    if _is_bond_d:
                        st.caption(
                            f"Votre fonds : {_info_d['ann_ret_pct']:+.1f}%/an  │  "
                            f"Vol {_info_d['vol_pct']:.1f}%  │  "
                            f"TER {_info_d['ter']:.2f}%"
                        )
                    else:
                        st.caption(
                            f"Votre fonds : {_info_d['ann_ret_pct']:+.1f}%/an  │  "
                            f"Sharpe {_info_d['sharpe']:.2f}  │  "
                            f"TER {_info_d['ter']:.2f}%"
                        )
                    if _info_d.get("first_vl"):
                        st.caption(f"📅 Première VL disponible : {_info_d['first_vl']}")

                    if _doublon_with:
                        st.caption(f"⚠️ Corrélation **{_doublon_with[1]:.0%}** avec {_doublon_with[0]} — comportement quasi-identique")

                    if _best_d and _rank_d and _rank_d > 1:
                        _delta_rdt = _best_d["ann_ret_pct"] - _info_d["ann_ret_pct"]
                        _etf_tag = " 🏷️ ETF" if _best_d.get("is_etf") else ""
                        st.markdown(
                            f"🏆 **Meilleur de la catégorie :** {_best_d['name'][:40]}{_etf_tag}  —  "
                            f"{_best_d['ann_ret_pct']:+.1f}%/an  │  Sharpe {_best_d['sharpe']:.2f}"
                        )
                        if _delta_rdt > 0.5:
                            st.caption(f"Soit **+{_delta_rdt:.1f}%/an** de rendement supplémentaire")

                    if _best_etf_d:
                        st.caption(
                            f"💡 **Alternative ETF :** {_best_etf_d['name'][:40]}  —  "
                            f"{_best_etf_d['ann_ret_pct']:+.1f}%/an pour un TER de {_best_etf_d['ter']:.2f}% "
                            f"(vs {_info_d['ter']:.2f}%)"
                        )

                    _btn_cols = st.columns([2, 2, 2, 1])
                    with _btn_cols[0]:
                        if not _best_d:
                            st.caption("✅ **Meilleur fonds de sa catégorie sur cette période** — aucun remplacement à suggérer")
                        elif _best_d and _best_d["sharpe"] > _info_d["sharpe"]:
                            if st.button(f"🔄 Remplacer par {_best_d['name'][:25]}", key=f"PP_DREPLACE_{_isin_d}"):
                                for _list_key in ("PP_SELECTED_ACTIONS", "PP_SELECTED_BONDS"):
                                    _lst = st.session_state.get(_list_key, [])
                                    _idx_found = next((_i for _i, _v in enumerate(_lst) if _v.upper() == _isin_d.upper()), None)
                                    if _idx_found is not None:
                                        _lst[_idx_found] = _best_d["isin"]
                                        st.session_state[_list_key] = _lst
                                        break
                                st.rerun()
                    with _btn_cols[1]:
                        if st.button("🔍 Alternatives", key=f"PP_DALTS_{_isin_d}"):
                            st.session_state[f"_PP_SHOW_ALTS_{_isin_d}"] = not st.session_state.get(f"_PP_SHOW_ALTS_{_isin_d}", False)
                            st.rerun()
                    with _btn_cols[2]:
                        if _best_d:
                            if st.button("🔄 Simuler", key=f"PP_DSIM_{_isin_d}"):
                                st.session_state[f"_PP_SHOW_SIM_{_isin_d}"] = not st.session_state.get(f"_PP_SHOW_SIM_{_isin_d}", False)
                                st.rerun()
                    with _btn_cols[3]:
                        if st.button("❌", key=f"PP_DREMOVE_{_isin_d}", help="Retirer"):
                            for _list_key in ("PP_SELECTED_ACTIONS", "PP_SELECTED_BONDS"):
                                _lst = st.session_state.get(_list_key, [])
                                _to_rm = [v for v in _lst if v.upper() == _isin_d.upper()]
                                for _tr in _to_rm:
                                    _lst.remove(_tr)
                                st.session_state[_list_key] = _lst
                            st.rerun()

                    if st.session_state.get(f"_PP_SHOW_ALTS_{_isin_d}"):
                        st.markdown("**Top 5 de cette catégorie :**")
                        _alts_top = [f for f in _cat_funds_d if f["isin"] != _isin_d and f["isin"].upper() not in _already_up][:5]
                        if _alts_top:
                            for _alt in _alts_top:
                                _at1, _at2, _at3 = st.columns([4, 4, 1])
                                _etf_b = " 🏷️ ETF" if _alt.get("is_etf") else ""
                                with _at1:
                                    st.caption(f"{_alt['name'][:40]}{_etf_b}")
                                with _at2:
                                    st.caption(f"{_alt['ann_ret_pct']:+.1f}%/an │ Sharpe {_alt['sharpe']:.2f} │ TER {_alt['ter']:.2f}%")
                                with _at3:
                                    if st.button("➕", key=f"PP_ALTADD_{_isin_d}_{_alt['isin']}", help="Ajouter en plus"):
                                        _alt_isin = _alt["isin"]
                                        _alt_cat = _rlookup.get(_alt_isin, {}).get("category", "")
                                        if _is_bond_category(_alt_cat):
                                            _tgt_list = "PP_SELECTED_BONDS"
                                            _tgt_limit = 10
                                        else:
                                            _tgt_list = "PP_SELECTED_ACTIONS"
                                            _tgt_limit = 40
                                        _tgt_current = st.session_state.get(_tgt_list, [])
                                        if _alt_isin not in _tgt_current and len(_tgt_current) < _tgt_limit:
                                            _tgt_current.append(_alt_isin)
                                            st.session_state[_tgt_list] = _tgt_current
                                            st.rerun()
                        else:
                            st.caption("Aucune alternative disponible dans cette catégorie.")

                    if st.session_state.get(f"_PP_SHOW_SIM_{_isin_d}") and _best_d:
                        try:
                            _curr_lines_s = [{"isin": i, "name": _global_isin_name.get(i, i), "amount_gross": 1000} for i in _uc_isins_diag]
                            _new_isins_s = [i if i.upper() != _isin_d.upper() else _best_d["isin"] for i in _uc_isins_diag]
                            _new_lines_s = [{"isin": i, "name": _global_isin_name.get(i, _best_d["name"] if i == _best_d["isin"] else i), "amount_gross": 1000} for i in _new_isins_s]

                            with st.spinner("Calcul de l'impact..."):
                                _cs = compute_sharpe_ratio(_curr_lines_s, _euro_rate_diag, 0.0, start_date=opt_start)
                                _ns = compute_sharpe_ratio(_new_lines_s, _euro_rate_diag, 0.0, start_date=opt_start)
                                _cd = compute_diversification_score(_curr_lines_s, _euro_rate_diag, start_date=opt_start)
                                _nd = compute_diversification_score(_new_lines_s, _euro_rate_diag, start_date=opt_start)

                            _sc1, _sc2 = st.columns(2)
                            with _sc1:
                                st.metric("Sharpe portefeuille", f"{(_ns or 0):.2f}", delta=f"{((_ns or 0) - (_cs or 0)):+.2f}")
                            with _sc2:
                                _oc = _cd["avg_corr"] if _cd else 0
                                _nc = _nd["avg_corr"] if _nd else 0
                                st.metric("Corrélation moyenne", f"{_nc:.0%}", delta=f"{(_nc - _oc):.0%}", delta_color="inverse")
                        except Exception as _e_sim:
                            st.warning(f"Simulation impossible : données insuffisantes. ({_e_sim})")

            # ── Suggestions ─────────────────────────────────────
            _has_bonds_diag = any(
                _rlookup.get(i, {}).get("category", "") and
                any(kw in _rlookup.get(i, {}).get("category", "").upper() for kw in
                    ("OBLIGATION", "BOND", "FIXED INCOME", "TAUX", "REVENUS FIXES"))
                for i in _selected_all_diag
            )
            if not _has_bonds_diag and _selected_all_diag:
                with st.expander("🛡️ Ajouter un fonds obligataire ?", expanded=False):
                    st.markdown(
                        "Votre portefeuille ne contient aucun fonds obligataire. "
                        "Les obligations ont un comportement généralement décorrélé des actions : "
                        "elles protègent le capital en cas de baisse des marchés."
                    )

                    _bond_subcats = {}
                    for _bc_cat, _bc_funds in _rby_cat.items():
                        _bc_up = _bc_cat.upper()
                        if not any(kw in _bc_up for kw in ("OBLIGATION", "BOND", "FIXED INCOME", "TAUX", "REVENUS FIXES", "CONVERTIBLE")):
                            continue
                        _bc_filtered = [f for f in _bc_funds if f["isin"].upper() not in _already_up]
                        if _bc_filtered:
                            _bond_subcats[_bc_cat] = _bc_filtered

                    for _bsc_cat, _bsc_funds in _bond_subcats.items():
                        st.caption(f"**{_bsc_cat}** :")
                        for _bf in _bsc_funds[:3]:
                            _bf1, _bf2, _bf3 = st.columns([4, 4, 1])
                            _etf_bf = " 🏷️ ETF" if _bf.get("is_etf") else ""
                            with _bf1:
                                st.caption(f"{_bf['name'][:40]}{_etf_bf}")
                            with _bf2:
                                st.caption(f"{_bf['ann_ret_pct']:+.1f}%/an │ Sharpe {_bf['sharpe']:.2f} │ TER {_bf['ter']:.2f}%")
                            with _bf3:
                                if st.button("➕", key=f"PP_SUGBOND_{_bf['isin']}", help="Ajouter"):
                                    _bonds_lst = st.session_state.get("PP_SELECTED_BONDS", [])
                                    if _bf["isin"] not in _bonds_lst and len(_bonds_lst) < 10:
                                        _bonds_lst.append(_bf["isin"])
                                        st.session_state["PP_SELECTED_BONDS"] = _bonds_lst
                                        st.rerun()


                with st.expander("ℹ️ Comprendre les types d'obligations", expanded=False):
                    st.markdown("""
**Obligations d'État (souveraines)** — Faible risque, rendement modéré. Très décorrélées des actions. Idéal pour sécuriser le capital en cas de crise.

**Obligations d'entreprise (corporate)** — Risque intermédiaire, rendement plus élevé. Bon compromis rendement/protection.

**High Yield (haut rendement)** — Risque plus élevé, corrélation plus forte avec les actions. Ajoute du rendement mais protège moins en cas de baisse.

**Obligations datées (target date, ex: 2029)** — Capital remboursé à l'échéance si pas de défaut. Visibilité sur le rendement final.

**Convertibles** — Mi-chemin entre action et obligation. Participe à la hausse des actions avec un coussin en baisse.
                    """)

        # ── Liste compacte "Fonds dans l'allocation" ────────────
        st.markdown("### 📋 Fonds dans l'allocation")
        _all_in_alloc = list(st.session_state.get("PP_SELECTED_ACTIONS", [])) + list(st.session_state.get("PP_SELECTED_BONDS", []))
        if _all_in_alloc:
            for _i_alloc, _isin_alloc in enumerate(_all_in_alloc):
                _nm_alloc = _global_isin_name.get(_isin_alloc, _isin_alloc)
                _info_alloc = _rlookup.get(_isin_alloc, {})
                _cat_alloc = _info_alloc.get("category", "—") if _info_alloc else "—"
                _ca1, _ca2, _ca3 = st.columns([4, 4, 1])
                with _ca1:
                    st.caption(f"{_nm_alloc}")
                with _ca2:
                    st.caption(f"{_cat_alloc}")
                with _ca3:
                    if st.button("❌", key=f"PP_ALLOC_RM_{_i_alloc}_{_isin_alloc}", help="Retirer"):
                        for _lk in ("PP_SELECTED_ACTIONS", "PP_SELECTED_BONDS"):
                            _ll = st.session_state.get(_lk, [])
                            if _isin_alloc in _ll:
                                _ll.remove(_isin_alloc)
                                st.session_state[_lk] = _ll
                                break
                        st.rerun()
        else:
            st.caption("Aucun fonds dans l'allocation.")

        st.markdown("---")

        st.caption(
            "Ces analyses sont basées sur les performances historiques nettes des frais de gestion "
            "internes (TER intégré dans les valeurs liquidatives publiées). Un fonds avec des frais "
            "élevés peut être le meilleur choix s'il surperforme sa catégorie — le rendement net prime "
            "sur le niveau de frais. Les corrélations passées ne garantissent pas les corrélations futures."
        )

    # ── Sélection des fonds UC ─────────────────────────────────
    def _is_etf(row) -> bool:
        _etf_kw = ("etf", "tracker", "ishares", "lyxor", "amundi index", "xtrackers", "vanguard")
        return any(kw in str(row.get("name", "")).lower() for kw in _etf_kw) or float(row.get("fee_uc_pct", 1.0)) < 0.50

    if "PP_SELECTED_ACTIONS" not in st.session_state:
        st.session_state["PP_SELECTED_ACTIONS"] = []
    if "PP_SELECTED_BONDS" not in st.session_state:
        st.session_state["PP_SELECTED_BONDS"] = []

    if not funds_df.empty:
        actions_df = funds_df[~funds_df["category"].apply(_is_bond_category)].copy()
        bonds_df = funds_df[funds_df["category"].apply(_is_bond_category)].copy()
    else:
        actions_df = pd.DataFrame()
        bonds_df = pd.DataFrame()

    def _apply_filters(base_df, search, filter_cats, filter_type, filter_sri, filter_ter, filter_managers):
        df = base_df.copy()
        if df.empty:
            return df
        if search:
            _s = search.lower()
            df = df[
                df["name"].str.lower().str.contains(_s, na=False)
                | df["isin"].str.lower().str.contains(_s, na=False)
            ]
        if filter_cats:
            df = df[df["category"].isin(filter_cats)]
        df = df[df["sri"].fillna(7) <= filter_sri]
        df = df[df["fee_uc_pct"].fillna(3.0) <= filter_ter]
        if filter_managers:
            df = df[df["manager"].isin(filter_managers)]
        if filter_type == "ETF uniquement":
            df = df[df.apply(_is_etf, axis=1)]
        elif filter_type == "OPCVM uniquement":
            df = df[~df.apply(_is_etf, axis=1)]
        return df

    # En mode Améliorer : section de sélection dans un expander fermé
    # En mode Créer : affiché normalement (pas d'expander)
    if _is_improve_mode:
        _sel_container = st.expander("➕ Ajouter un nouveau fonds au portefeuille", expanded=False)
    else:
        _sel_container = st.container()

    with _sel_container:
        st.markdown("### 📋 Sélection des fonds UC")

        # Pré-calcul des ISINs valides (historique suffisant) pour le filtrage en mode Créer
        if not _is_improve_mode:
            _rankings_for_filter = _compute_rankings_for_period(
                pp_contract_cfg["path"],
                pp_contract_cfg["funds_filename"],
                float(euro_rate),
                str(opt_start.date()),
                str(opt_end.date()),
                _cache_key=f"{pp_contract_label}_{opt_start.date()}_{opt_end.date()}",
            )
            _valid_isins_for_period = set()
            for _cat_vf, _funds_vf in _rankings_for_filter.items():
                for _f_vf in _funds_vf:
                    _valid_isins_for_period.add(_f_vf["isin"])
        else:
            _valid_isins_for_period = None

        st.markdown("#### 📈 Fonds Actions & Diversifiés")

        # ── SECTION ACTIONS ────────────────────────────────────────
        _act_managers = sorted(actions_df["manager"].dropna().unique().tolist()) if not actions_df.empty else []
        _act_cats = sorted(actions_df["category"].dropna().unique().tolist()) if not actions_df.empty else []

        # Appliquer le filtre "Similaires" si demandé (clé tampon → default du multiselect)
        _similar_cat_act = st.session_state.pop("_PP_SIMILAR_CAT_ACT", None)
        if _similar_cat_act and _similar_cat_act in _act_cats:
            st.session_state["PP_FILTER_CAT_ACT"] = [_similar_cat_act]

        _ac1, _ac2, _ac3 = st.columns(3)
        with _ac1:
            _act_search = st.text_input("🔍 Rechercher", key="PP_SEARCH_ACT", placeholder="Nom ou ISIN…")
            _act_filter_cats = st.multiselect("Catégorie", _act_cats, key="PP_FILTER_CAT_ACT")
        with _ac2:
            _act_filter_type = st.selectbox("Type", ["Tous", "ETF uniquement", "OPCVM uniquement"], key="PP_FILTER_TYPE_ACT")
            _act_filter_sri = st.slider("SRI max", 1, 7, 7, key="PP_FILTER_SRI_ACT")
        with _ac3:
            _act_filter_ter = st.slider("Frais de gestion max (TER %/an)", 0.0, 3.0, 3.0, 0.1, key="PP_FILTER_TER_ACT")
            _act_filter_mgr = st.multiselect("Société de gestion", _act_managers, key="PP_FILTER_MANAGER_ACT")

        _act_filtered = _apply_filters(actions_df, _act_search, _act_filter_cats, _act_filter_type, _act_filter_sri, _act_filter_ter, _act_filter_mgr)
        st.caption(f"{len(_act_filtered)} fonds correspondent aux filtres")

        # Filtrer les fonds avec historique insuffisant (mode Créer)
        if _valid_isins_for_period is not None and not _act_filtered.empty:
            _act_valid = _act_filtered[_act_filtered["isin"].isin(_valid_isins_for_period)]
            _act_invalid = _act_filtered[~_act_filtered["isin"].isin(_valid_isins_for_period)]
        else:
            _act_valid = _act_filtered
            _act_invalid = pd.DataFrame()

        if not _act_valid.empty:
            _act_disp = _act_valid[["name", "isin", "category", "sri", "fee_uc_pct", "fee_contract_pct", "fee_total_pct"]].copy()
            _act_disp.columns = ["Nom", "ISIN", "Catégorie", "SRI", "TER (%/an)", "Frais contrat (%/an)", "Total (%/an)"]
            _act_disp["Type"] = _act_valid.apply(_is_etf, axis=1).map({True: "ETF", False: "OPCVM"})
            try:
                st.dataframe(_act_disp, use_container_width=True, hide_index=True, height=300)
            except Exception:
                pass
            st.caption(f"{len(_act_valid)} fonds avec historique suffisant sur la période")

            if not _act_invalid.empty:
                with st.expander(f"⚠️ {len(_act_invalid)} fonds exclus (1ère VL postérieure au {fmt_date(opt_start)})", expanded=False):
                    _inv_disp = _act_invalid[["name", "isin", "category"]].copy()
                    _inv_disp.columns = ["Nom", "ISIN", "Catégorie"]
                    try:
                        st.dataframe(_inv_disp, use_container_width=True, hide_index=True)
                    except Exception:
                        pass
                    st.caption("Ces fonds n'ont pas un historique suffisant pour la période d'analyse choisie.")

        _already_selected_upper = set(
            [i.upper() for i in st.session_state.get("PP_SELECTED_ACTIONS", [])] +
            [i.upper() for i in st.session_state.get("PP_SELECTED_BONDS", [])]
        )

        # ── Top 5 contextuel (uniquement si catégorie filtrée) ────────
        if _act_filter_cats:
            _rankings_data = _compute_rankings_for_period(
                pp_contract_cfg["path"],
                pp_contract_cfg["funds_filename"],
                float(euro_rate),
                str(opt_start.date()),
                str(opt_end.date()),
                _cache_key=f"{pp_contract_label}_{opt_start.date()}_{opt_end.date()}",
            )

            _top5_candidates_act = []
            for _cat_sel in _act_filter_cats:
                for _f_r in _rankings_data.get(_cat_sel, []):
                    if _f_r["isin"].upper() not in _already_selected_upper:
                        _top5_candidates_act.append(_f_r)

            if _top5_candidates_act:
                _is_bond_top5_act = all(_is_bond_category(c) for c in _act_filter_cats)
                _tc1, _tc2 = st.columns([3, 1])
                with _tc1:
                    st.markdown("**🏆 Top 5 de cette catégorie**")
                with _tc2:
                    _sort_opts_act = ["Rendement/an", "Volatilité la plus basse", "TER le plus bas"] if _is_bond_top5_act else ["Sharpe", "Rendement/an", "TER le plus bas"]
                    _sort_crit_act = st.selectbox(
                        "Trier par", _sort_opts_act,
                        key="PP_TOP5_SORT_ACT", label_visibility="collapsed",
                    )

                _sort_map_act = {
                    "Sharpe": ("sharpe", True), "Rendement/an": ("ann_ret_pct", True),
                    "TER le plus bas": ("ter", False), "Volatilité la plus basse": ("vol_pct", False),
                }
                _sort_key_act, _sort_reverse_act = _sort_map_act[_sort_crit_act]
                _top5_sorted_act = sorted(_top5_candidates_act, key=lambda x: x[_sort_key_act], reverse=_sort_reverse_act)[:5]

                if _top5_sorted_act:
                    _max_val_act = max(abs(f[_sort_key_act]) for f in _top5_sorted_act) if _top5_sorted_act else 1
                    if _max_val_act <= 0:
                        _max_val_act = 1

                    for _f_t in _top5_sorted_act:
                        _bar_pct = min(100, abs(_f_t[_sort_key_act]) / _max_val_act * 100)
                        _bar_color = "#1A7A4A" if _f_t["ann_ret_pct"] > 0 else "#CC2200"
                        _etf_badge = " 🏷️ ETF" if _f_t.get("is_etf") else ""

                        _bt1, _bt2, _bt3 = st.columns([4, 4, 1])
                        with _bt1:
                            st.markdown(
                                f'<div style="background:linear-gradient(90deg, {_bar_color}40 {_bar_pct:.0f}%, transparent {_bar_pct:.0f}%); '
                                f'padding:4px 8px; border-radius:4px; font-size:13px;">'
                                f'{_f_t["name"][:40]}{_etf_badge}</div>',
                                unsafe_allow_html=True,
                            )
                        with _bt2:
                            if _is_bond_top5_act:
                                st.caption(
                                    f'{_f_t["ann_ret_pct"]:+.1f}%/an  │  '
                                    f'Vol {_f_t["vol_pct"]:.1f}%  │  '
                                    f'TER {_f_t["ter"]:.2f}%  │  '
                                    f'1ère VL : {_f_t.get("first_vl", "—")}'
                                )
                            else:
                                st.caption(
                                    f'{_f_t["ann_ret_pct"]:+.1f}%/an  │  '
                                    f'Sharpe {_f_t["sharpe"]:.2f}  │  '
                                    f'TER {_f_t["ter"]:.2f}%  │  '
                                    f'1ère VL : {_f_t.get("first_vl", "—")}'
                                )
                        with _bt3:
                            if st.button("➕", key=f"PP_TOP5_ADD_ACT_{_f_t['isin']}", help="Ajouter"):
                                if _f_t["isin"] not in st.session_state["PP_SELECTED_ACTIONS"]:
                                    if len(st.session_state["PP_SELECTED_ACTIONS"]) < 40:
                                        st.session_state["PP_SELECTED_ACTIONS"].append(_f_t["isin"])
                                        st.rerun()

                    st.caption("Performance nette des frais de gestion (TER intégré dans les VL). Données 3 ans, mises à jour chaque semaine.")

        _act_options = ["— sélectionner —"] + ([f"{r['name']} ({r['isin']})" for _, r in _act_valid.iterrows() if r["isin"].upper() not in _already_selected_upper] if not _act_valid.empty else [])
        _act_lookup = {f"{r['name']} ({r['isin']})": r["isin"] for _, r in _act_valid.iterrows() if r["isin"].upper() not in _already_selected_upper} if not _act_valid.empty else {}

        _col_sel_a, _col_btn_a1, _col_btn_a2 = st.columns([4, 1, 2])
        with _col_sel_a:
            _act_choice = st.selectbox("Sélectionner un fonds", _act_options, key="PP_PICK_ACT")
        with _col_btn_a1:
            st.write("")
            st.write("")
            if st.button("➕ Ajouter", key="PP_ADD_ACT"):
                _isin_to_add = _act_lookup.get(_act_choice, "")
                if _isin_to_add and _isin_to_add not in st.session_state["PP_SELECTED_ACTIONS"] and _isin_to_add not in st.session_state["PP_SELECTED_BONDS"]:
                    if len(st.session_state["PP_SELECTED_ACTIONS"]) >= 40:
                        st.warning("Maximum 40 fonds actions atteint.")
                    else:
                        st.session_state["PP_SELECTED_ACTIONS"].append(_isin_to_add)
                        st.rerun()
        with _col_btn_a2:
            st.write("")
            st.write("")
            if st.button("📋 Ajouter tout", key="PP_ADD_ALL_ACT"):
                _added_a = 0
                for _, _row_a in _act_valid.iterrows():
                    if len(st.session_state["PP_SELECTED_ACTIONS"]) >= 40:
                        break
                    _isin_a = _row_a["isin"]
                    if _isin_a not in st.session_state["PP_SELECTED_ACTIONS"] and _isin_a not in st.session_state["PP_SELECTED_BONDS"]:
                        st.session_state["PP_SELECTED_ACTIONS"].append(_isin_a)
                        _added_a += 1
                if _added_a:
                    st.rerun()

        st.markdown("**Fonds actions sélectionnés :**")
        if len(st.session_state.get("PP_SELECTED_ACTIONS", [])) > 20:
            st.warning("⚠️ Plus de 20 fonds sélectionnés — le calcul d'optimisation peut prendre plusieurs minutes.")
        for _i_a, _isin_a in enumerate(list(st.session_state["PP_SELECTED_ACTIONS"])):
            _fname_a = _global_isin_name.get(_isin_a, _isin_a)
            _c1_a, _c2_a, _c3_a, _c4_a = st.columns([4, 2, 1, 1])
            with _c1_a:
                st.caption(f"{_fname_a} ({_isin_a})")
            with _c2_a:
                st.checkbox("📌 Conserver", key=f"PP_ANCHOR_{_isin_a}", value=False)
            with _c3_a:
                if st.button("🔍", key=f"PP_SIMILAR_ACT_{_i_a}", help="Voir les fonds de la même catégorie"):
                    _match_sim_a = funds_df[funds_df["isin"].str.upper() == _isin_a.upper()]
                    if not _match_sim_a.empty:
                        _cat_sim_a = str(_match_sim_a.iloc[0]["category"])
                        if _is_bond_category(_cat_sim_a):
                            st.session_state["_PP_SIMILAR_CAT_OBL"] = _cat_sim_a
                        else:
                            st.session_state["_PP_SIMILAR_CAT_ACT"] = _cat_sim_a
                    st.rerun()
            with _c4_a:
                if st.button("❌", key=f"PP_REMOVE_ACT_{_i_a}"):
                    st.session_state["PP_SELECTED_ACTIONS"].remove(_isin_a)
                    st.rerun()

        # ── SECTION OBLIGATIONS ────────────────────────────────────
        st.markdown("---")
        st.markdown("#### 🛡️ Fonds Obligataires")
        _obl_managers = sorted(bonds_df["manager"].dropna().unique().tolist()) if not bonds_df.empty else []
        _obl_cats = sorted(bonds_df["category"].dropna().unique().tolist()) if not bonds_df.empty else []

        # Appliquer le filtre "Similaires" si demandé (clé tampon → default du multiselect)
        _similar_cat_obl = st.session_state.pop("_PP_SIMILAR_CAT_OBL", None)
        if _similar_cat_obl and _similar_cat_obl in _obl_cats:
            st.session_state["PP_FILTER_CAT_OBL"] = [_similar_cat_obl]

        _oc1, _oc2, _oc3 = st.columns(3)
        with _oc1:
            _obl_search = st.text_input("🔍 Rechercher", key="PP_SEARCH_OBL", placeholder="Nom ou ISIN…")
            _obl_filter_cats = st.multiselect("Catégorie", _obl_cats, key="PP_FILTER_CAT_OBL")
        with _oc2:
            _obl_filter_type = st.selectbox("Type", ["Tous", "ETF uniquement", "OPCVM uniquement"], key="PP_FILTER_TYPE_OBL")
            _obl_filter_sri = st.slider("SRI max", 1, 7, 7, key="PP_FILTER_SRI_OBL")
        with _oc3:
            _obl_filter_ter = st.slider("Frais de gestion max (TER %/an)", 0.0, 3.0, 3.0, 0.1, key="PP_FILTER_TER_OBL")
            _obl_filter_mgr = st.multiselect("Société de gestion", _obl_managers, key="PP_FILTER_MANAGER_OBL")

        _obl_filtered = _apply_filters(bonds_df, _obl_search, _obl_filter_cats, _obl_filter_type, _obl_filter_sri, _obl_filter_ter, _obl_filter_mgr)
        st.caption(f"{len(_obl_filtered)} fonds correspondent aux filtres")

        # Filtrer les fonds obligataires avec historique insuffisant (mode Créer)
        if _valid_isins_for_period is not None and not _obl_filtered.empty:
            _obl_valid = _obl_filtered[_obl_filtered["isin"].isin(_valid_isins_for_period)]
            _obl_invalid = _obl_filtered[~_obl_filtered["isin"].isin(_valid_isins_for_period)]
        else:
            _obl_valid = _obl_filtered
            _obl_invalid = pd.DataFrame()

        if not _obl_valid.empty:
            _obl_disp = _obl_valid[["name", "isin", "category", "sri", "fee_uc_pct", "fee_contract_pct", "fee_total_pct"]].copy()
            _obl_disp.columns = ["Nom", "ISIN", "Catégorie", "SRI", "TER (%/an)", "Frais contrat (%/an)", "Total (%/an)"]
            _obl_disp["Type"] = _obl_valid.apply(_is_etf, axis=1).map({True: "ETF", False: "OPCVM"})
            try:
                st.dataframe(_obl_disp, use_container_width=True, hide_index=True, height=300)
            except Exception:
                pass
            st.caption(f"{len(_obl_valid)} fonds avec historique suffisant sur la période")

            if not _obl_invalid.empty:
                with st.expander(f"⚠️ {len(_obl_invalid)} fonds exclus (1ère VL postérieure au {fmt_date(opt_start)})", expanded=False):
                    _inv_obl_disp = _obl_invalid[["name", "isin", "category"]].copy()
                    _inv_obl_disp.columns = ["Nom", "ISIN", "Catégorie"]
                    try:
                        st.dataframe(_inv_obl_disp, use_container_width=True, hide_index=True)
                    except Exception:
                        pass
                    st.caption("Ces fonds n'ont pas un historique suffisant pour la période d'analyse choisie.")

        # ── Top 5 contextuel (uniquement si catégorie filtrée) ────────
        if _obl_filter_cats:
            _rankings_data_obl = _compute_rankings_for_period(
                pp_contract_cfg["path"],
                pp_contract_cfg["funds_filename"],
                float(euro_rate),
                str(opt_start.date()),
                str(opt_end.date()),
                _cache_key=f"{pp_contract_label}_{opt_start.date()}_{opt_end.date()}",
            )

            _top5_candidates_obl = []
            for _cat_sel_o in _obl_filter_cats:
                for _f_r in _rankings_data_obl.get(_cat_sel_o, []):
                    if _f_r["isin"].upper() not in _already_selected_upper:
                        _top5_candidates_obl.append(_f_r)

            if _top5_candidates_obl:
                _is_bond_top5_obl = all(_is_bond_category(c) for c in _obl_filter_cats)
                _tc1o, _tc2o = st.columns([3, 1])
                with _tc1o:
                    st.markdown("**🏆 Top 5 de cette catégorie**")
                with _tc2o:
                    _sort_opts_obl = ["Rendement/an", "Volatilité la plus basse", "TER le plus bas"] if _is_bond_top5_obl else ["Sharpe", "Rendement/an", "TER le plus bas"]
                    _sort_crit_obl = st.selectbox(
                        "Trier par", _sort_opts_obl,
                        key="PP_TOP5_SORT_OBL", label_visibility="collapsed",
                    )

                _sort_map_obl = {
                    "Sharpe": ("sharpe", True), "Rendement/an": ("ann_ret_pct", True),
                    "TER le plus bas": ("ter", False), "Volatilité la plus basse": ("vol_pct", False),
                }
                _sort_key_obl, _sort_reverse_obl = _sort_map_obl[_sort_crit_obl]
                _top5_sorted_obl = sorted(_top5_candidates_obl, key=lambda x: x[_sort_key_obl], reverse=_sort_reverse_obl)[:5]

                if _top5_sorted_obl:
                    _max_val_obl = max(abs(f[_sort_key_obl]) for f in _top5_sorted_obl) if _top5_sorted_obl else 1
                    if _max_val_obl <= 0:
                        _max_val_obl = 1

                    for _f_t in _top5_sorted_obl:
                        _bar_pct = min(100, abs(_f_t[_sort_key_obl]) / _max_val_obl * 100)
                        _bar_color = "#1A7A4A" if _f_t["ann_ret_pct"] > 0 else "#CC2200"
                        _etf_badge = " 🏷️ ETF" if _f_t.get("is_etf") else ""

                        _bt1o, _bt2o, _bt3o = st.columns([4, 4, 1])
                        with _bt1o:
                            st.markdown(
                                f'<div style="background:linear-gradient(90deg, {_bar_color}40 {_bar_pct:.0f}%, transparent {_bar_pct:.0f}%); '
                                f'padding:4px 8px; border-radius:4px; font-size:13px;">'
                                f'{_f_t["name"][:40]}{_etf_badge}</div>',
                                unsafe_allow_html=True,
                            )
                        with _bt2o:
                            if _is_bond_top5_obl:
                                st.caption(
                                    f'{_f_t["ann_ret_pct"]:+.1f}%/an  │  '
                                    f'Vol {_f_t["vol_pct"]:.1f}%  │  '
                                    f'TER {_f_t["ter"]:.2f}%  │  '
                                    f'1ère VL : {_f_t.get("first_vl", "—")}'
                                )
                            else:
                                st.caption(
                                    f'{_f_t["ann_ret_pct"]:+.1f}%/an  │  '
                                    f'Sharpe {_f_t["sharpe"]:.2f}  │  '
                                    f'TER {_f_t["ter"]:.2f}%  │  '
                                    f'1ère VL : {_f_t.get("first_vl", "—")}'
                                )
                        with _bt3o:
                            if st.button("➕", key=f"PP_TOP5_ADD_OBL_{_f_t['isin']}", help="Ajouter"):
                                if _f_t["isin"] not in st.session_state["PP_SELECTED_BONDS"]:
                                    if len(st.session_state["PP_SELECTED_BONDS"]) < 10:
                                        st.session_state["PP_SELECTED_BONDS"].append(_f_t["isin"])
                                        st.rerun()

                    st.caption("Performance nette des frais de gestion (TER intégré dans les VL). Données 3 ans, mises à jour chaque semaine.")

        _obl_options = ["— sélectionner —"] + ([f"{r['name']} ({r['isin']})" for _, r in _obl_valid.iterrows() if r["isin"].upper() not in _already_selected_upper] if not _obl_valid.empty else [])
        _obl_lookup = {f"{r['name']} ({r['isin']})": r["isin"] for _, r in _obl_valid.iterrows() if r["isin"].upper() not in _already_selected_upper} if not _obl_valid.empty else {}

        _col_sel_o, _col_btn_o1, _col_btn_o2 = st.columns([4, 1, 2])
        with _col_sel_o:
            _obl_choice = st.selectbox("Sélectionner un fonds", _obl_options, key="PP_PICK_OBL")
        with _col_btn_o1:
            st.write("")
            st.write("")
            if st.button("➕ Ajouter", key="PP_ADD_OBL"):
                _isin_to_add = _obl_lookup.get(_obl_choice, "")
                if _isin_to_add and _isin_to_add not in st.session_state["PP_SELECTED_ACTIONS"] and _isin_to_add not in st.session_state["PP_SELECTED_BONDS"]:
                    if len(st.session_state["PP_SELECTED_BONDS"]) >= 10:
                        st.warning("Maximum 10 fonds obligataires atteint.")
                    else:
                        st.session_state["PP_SELECTED_BONDS"].append(_isin_to_add)
                        st.rerun()
        with _col_btn_o2:
            st.write("")
            st.write("")
            if st.button("📋 Ajouter tout", key="PP_ADD_ALL_OBL"):
                _added_o = 0
                for _, _row_o in _obl_valid.iterrows():
                    if len(st.session_state["PP_SELECTED_BONDS"]) >= 10:
                        break
                    _isin_o = _row_o["isin"]
                    if _isin_o not in st.session_state["PP_SELECTED_ACTIONS"] and _isin_o not in st.session_state["PP_SELECTED_BONDS"]:
                        st.session_state["PP_SELECTED_BONDS"].append(_isin_o)
                        _added_o += 1
                if _added_o:
                    st.rerun()

        st.markdown("**Fonds obligataires sélectionnés :**")
        for _i_o, _isin_o in enumerate(list(st.session_state["PP_SELECTED_BONDS"])):
            _fname_o = _global_isin_name.get(_isin_o, _isin_o)
            _c1_o, _c2_o, _c3_o, _c4_o = st.columns([4, 2, 1, 1])
            with _c1_o:
                st.caption(f"{_fname_o} ({_isin_o})")
            with _c2_o:
                st.checkbox("📌 Conserver", key=f"PP_ANCHOR_{_isin_o}", value=False)
            with _c3_o:
                if st.button("🔍", key=f"PP_SIMILAR_OBL_{_i_o}", help="Voir les fonds de la même catégorie"):
                    _match_sim_o = funds_df[funds_df["isin"].str.upper() == _isin_o.upper()]
                    if not _match_sim_o.empty:
                        _cat_sim_o = str(_match_sim_o.iloc[0]["category"])
                        if _is_bond_category(_cat_sim_o):
                            st.session_state["_PP_SIMILAR_CAT_OBL"] = _cat_sim_o
                        else:
                            st.session_state["_PP_SIMILAR_CAT_ACT"] = _cat_sim_o
                    st.rerun()
            with _c4_o:
                if st.button("❌", key=f"PP_REMOVE_OBL_{_i_o}"):
                    st.session_state["PP_SELECTED_BONDS"].remove(_isin_o)
                    st.rerun()

    # ── Récap global ───────────────────────────────────────────
    selected_isins = st.session_state.get("PP_SELECTED_ACTIONS", []) + st.session_state.get("PP_SELECTED_BONDS", [])
    forced_isin = next((isin for isin in selected_isins if st.session_state.get(f"PP_ANCHOR_{isin}", False)), None)
    _nb = len(selected_isins)
    st.caption(f"Univers d'investissement : {_nb} fonds sélectionnés ({len(st.session_state.get('PP_SELECTED_ACTIONS', []))} actions/diversifiés + {len(st.session_state.get('PP_SELECTED_BONDS', []))} obligataires)")
    if 0 < _nb < 2:
        st.warning("Sélectionnez au moins 2 fonds pour lancer l'optimisation.")

    # ── Répartition de l'allocation ──────────────────────────────
    st.markdown("### 🎯 Répartition de l'allocation")
    _alloc_mode_choice = st.radio(
        "Mode de répartition",
        ["Par catégorie", "Par nombre total"],
        horizontal=True,
        key="PP_ALLOC_MODE",
        help="Par catégorie : vous choisissez combien de fonds par catégorie. Par nombre total : vous choisissez le nombre de fonds UC et obligataires, l'algorithme sélectionne la meilleure combinaison.",
    )

    _cats_in_selection: Dict[str, List[str]] = {}
    for _isin in selected_isins:
        _row_c = funds_df.loc[funds_df["isin"].str.upper() == _isin.upper()] if not funds_df.empty else pd.DataFrame()
        _cat_c = _row_c["category"].values[0] if not _row_c.empty else "Autre"
        _cats_in_selection.setdefault(_cat_c, []).append(_isin)

    _cat_counts: Dict[str, int] = {}

    if _alloc_mode_choice == "Par catégorie":
        st.caption(
            "Définissez combien de fonds de chaque catégorie dans l'allocation finale. "
            "L'algorithme sélectionnera les meilleurs candidats."
        )
        for _cat, _isins_c in sorted(_cats_in_selection.items()):
            _cc1, _cc2 = st.columns([4, 1])
            with _cc1:
                st.caption(f"{_cat} — {len(_isins_c)} fonds disponibles")
            with _cc2:
                _n = st.number_input(
                    "Nb", min_value=0, max_value=len(_isins_c),
                    value=min(1, len(_isins_c)), step=1,
                    key=f"PP_CAT_COUNT_{_cat}", label_visibility="collapsed"
                )
                _cat_counts[_cat] = int(_n)

    else:  # Par nombre total
        st.caption(
            "Choisissez le nombre total de fonds UC et obligataires. "
            "L'algorithme sélectionnera la meilleure combinaison parmi vos fonds sélectionnés."
        )
        _nb_act_avail = max(1, len(st.session_state.get("PP_SELECTED_ACTIONS", [])))
        _nb_obl_avail = len(st.session_state.get("PP_SELECTED_BONDS", []))
        _tc1, _tc2 = st.columns(2)
        with _tc1:
            _nb_uc_total = st.number_input(
                "Nombre de fonds UC (actions/diversifiés)",
                min_value=1,
                max_value=min(20, _nb_act_avail),
                value=min(4, _nb_act_avail),
                step=1,
                key="PP_NB_UC_TOTAL",
            )
        with _tc2:
            _nb_obl_total = st.number_input(
                "Nombre de fonds obligataires",
                min_value=0,
                max_value=min(10, _nb_obl_avail) if _nb_obl_avail > 0 else 0,
                value=min(1, _nb_obl_avail),
                step=1,
                key="PP_NB_OBL_TOTAL",
            )
        _cat_counts = {"_UC_TOTAL": _nb_uc_total, "_OBL_TOTAL": _nb_obl_total}

    _total_requested = sum(_cat_counts.values())
    st.caption(f"Total : {_total_requested} fonds dans l'allocation finale")

    # ── Architecture 3 poches ──────────────────────────────────
    st.markdown("---")
    st.markdown("**Répartition du budget par poche**")
    _obl_default_map = {"Prudent": 30, "Équilibré": 20, "Dynamique": 10, "Offensif": 5}
    _has_bonds_selected = len(st.session_state.get("PP_SELECTED_BONDS", [])) > 0
    _obl_max = max(0, 100 - int(euro_pct))
    if not _has_bonds_selected:
        obl_pct = 0
        st.session_state["PP_OBL_PCT"] = 0
        st.caption("Aucun fonds obligataire sélectionné — poche obligations désactivée.")
    else:
        obl_pct = int(st.number_input(
            "Part obligations (%)",
            min_value=0,
            max_value=_obl_max,
            value=min(_obl_default_map.get(profile, 20), _obl_max),
            step=5,
            key="PP_OBL_PCT",
            help="Part du budget total allouée aux fonds obligataires.",
        ))
    act_pct = max(0, 100 - int(euro_pct) - obl_pct)
    st.caption(
        f"Budget total : {to_eur(total_budget)}  │  "
        f"Fonds euros ({int(euro_pct)}%) : {to_eur(total_budget * int(euro_pct) / 100)}  │  "
        f"UC Actions ({act_pct}%) : {to_eur(total_budget * act_pct / 100)}  │  "
        f"UC Obligations ({obl_pct}%) : {to_eur(total_budget * obl_pct / 100)}"
    )

    st.markdown("---")
    _obj_c1, _obj_c2 = st.columns(2)
    with _obj_c1:
        objective_choice = st.selectbox(
            "Objectif global (poche actions)",
            ["Maximiser le Sharpe", "Minimiser la volatilité"],
            key="PP_OBJ_GLOBAL",
            help="Objectif d'optimisation pour la poche actions/diversifiés.",
        )
    with _obj_c2:
        if obl_pct > 0:
            _obj_oblig = st.selectbox(
                "Objectif poche obligataire",
                ["Le plus stable possible", "Meilleur rendement-volatilité", "Meilleur rendement"],
                key="PP_OBJ_OBLIG",
                help="Critère de sélection et de répartition des fonds obligataires.",
            )
        else:
            _obj_oblig = "Le plus stable possible"
            st.caption("Poche obligataire désactivée")

    if st.session_state.get("PP_PRACTICAL_MODE", True):
        _uc_frac_check = act_pct / 100.0
        _max_n_check = max(1, int(_uc_frac_check / 0.10)) if _uc_frac_check > 0 else 1
        if _total_requested > _max_n_check:
            st.warning(
                f"⚠️ Mode terrain (min 10%) : max {_max_n_check} fonds actions avec {act_pct}% actions. "
                f"Vous en demandez {_total_requested}."
            )

    practical_mode = st.checkbox(
        "Optimisation terrain (min 10% + arrondi 5%)",
        value=True,
        key="PP_PRACTICAL_MODE",
    )

    params = (
        profile,
        int(euro_pct),
        obl_pct,
        float(euro_rate),
        int(total_budget),
        opt_window_mode,
        str(opt_start.date()),
        str(opt_end.date()),
        tuple(sorted(selected_isins)),
        objective_choice,
        _obj_oblig,
        forced_isin,
        _alloc_mode_choice,
        tuple(sorted(_cat_counts.items())),
    )
    current_hash = params_hash(params)

    if st.session_state.get("PP_RUN") and st.session_state.get("PP_PARAMS_HASH") != current_hash:
        st.session_state["PP_RUN"] = False
        st.info("Paramètres modifiés. Cliquez de nouveau sur '✅ Lancer l'optimisation'.")

    run_clicked = st.button("✅ Lancer l'optimisation", type="primary")
    if run_clicked:
        st.session_state["PP_RUN"] = True
        st.session_state["PP_PARAMS_HASH"] = current_hash

    if not st.session_state.get("PP_RUN"):
        st.info("Renseignez tous les paramètres puis cliquez sur '✅ Lancer l'optimisation'.")
        return

    if opt_start > opt_end:
        st.warning("La date de debut doit etre anterieure a la date de fin.")
        return

    try:
        if funds_df.empty:
            st.warning(
                "Référentiel de fonds non chargé. "
                "Sélectionnez un contrat dans la barre latérale."
            )
            return

        # Univers de candidats depuis la sélection UI
        all_candidates = list(selected_isins)
        if not all_candidates:
            st.info("Sélectionnez au moins 2 fonds UC pour lancer l'optimisation.")
            return

        returns_all, status_all = _returns_for_isins(all_candidates, opt_start, opt_end, euro_rate=float(euro_rate))

        insufficient = [isin for isin, status in status_all.items() if status != "ok"]
        valid_all = [isin for isin in all_candidates if status_all.get(isin) == "ok"]

        if insufficient:
            st.warning("Certains fonds ont été exclus (historique insuffisant sur la fenêtre).")

        if forced_isin and forced_isin not in valid_all:
            st.warning("Le fonds conservé est indisponible sur la période et a été ignoré.")
            forced_isin = None

        if returns_all.empty:
            st.info("Historique insuffisant pour calculer les rendements sur la fenêtre choisie.")
            return

        def _zscore(s: pd.Series) -> pd.Series:
            if s.empty:
                return s
            std = float(s.std())
            if std == 0 or np.isnan(std):
                return pd.Series(0.0, index=s.index)
            return (s - s.mean()) / std

        def _stats_for(candidates: List[str], ret_df: pd.DataFrame) -> Tuple[pd.Series, pd.Series, pd.Series, pd.DataFrame, pd.Series]:
            if not candidates or ret_df.empty:
                return pd.Series(dtype=float), pd.Series(dtype=float), pd.Series(dtype=float), pd.DataFrame(), pd.Series(dtype=float)
            # FIXED: arithmetic annualisation formula — was incorrectly treating pct_change() returns
            # as log-returns then applying exp(), double-compounding the result (Bug 2)
            rfr = get_risk_free_rate()  # taux sans risque dynamique (Bund 10 ans ou saisie manuelle)
            ann_return = (1 + ret_df[candidates].mean()) ** 252 - 1
            ann_vol = ret_df[candidates].std() * np.sqrt(252.0)
            sharpe = ((ann_return - rfr) / ann_vol.replace(0, np.nan)).replace([np.inf, -np.inf], np.nan)
            corr = ret_df[candidates].corr()
            avg_corr = corr.apply(lambda row: row.drop(labels=row.name, errors="ignore").mean(), axis=1).fillna(0.0) if not corr.empty else pd.Series(0.0, index=candidates)
            return ann_return, ann_vol, sharpe, corr, avg_corr

        def _greedy_min_corr(candidates: List[str], corr: pd.DataFrame, k: int, seed: Optional[List[str]] = None) -> List[str]:
            if k <= 0 or not candidates:
                return []
            selected = [x for x in (seed or []) if x in candidates]
            remaining = [c for c in candidates if c not in selected]
            while len(selected) < k and remaining:
                best, best_score = None, None
                for cand in remaining:
                    score = float(corr.loc[cand, selected].mean()) if selected and not corr.empty else 0.0
                    if best_score is None or score < best_score:
                        best, best_score = cand, score
                if best is None:
                    break
                selected.append(best)
                remaining = [c for c in remaining if c != best]
            return selected

        def _select_by_objective(candidates: List[str], ret_df: pd.DataFrame, k: int, objective: str, forced: Optional[str] = None) -> List[str]:
            if k <= 0 or not candidates:
                return []
            if ret_df.empty:
                base = candidates[:k]
                if forced and forced in candidates and forced not in base:
                    base = [forced] + [x for x in base if x != forced]
                    base = base[:k]
                return base

            ann_ret, ann_vol, sharpe, corr, avg_corr = _stats_for(candidates, ret_df)
            sharpe_rank = sharpe.fillna(-1e9).sort_values(ascending=False)

            if objective == "Rendement prioritaire (diversifié)":
                selected = ann_ret.sort_values(ascending=False).index.tolist()[:k]
            elif objective == "Sécuriser le capital (diversifié)":
                selected = ann_vol.sort_values(ascending=True).index.tolist()[:k]
            elif objective == "Diversification maximale":
                seed = sharpe_rank.index.tolist()[:1]
                selected = _greedy_min_corr(candidates, corr, k, seed=seed)
            else:
                # Maximiser Sharpe (diversifié) + fallback
                selected = sharpe_rank.index.tolist()[:k]

            if forced and forced in candidates:
                if forced not in selected:
                    if objective == "Diversification maximale":
                        selected = _greedy_min_corr(candidates, corr, k, seed=[forced])
                    else:
                        rest = [x for x in selected if x != forced]
                        selected = [forced] + rest
                        selected = selected[:k]
                else:
                    selected = [forced] + [x for x in selected if x != forced]
            return selected[:k]

        # Sélection directe depuis l'UI — filtrage par validité des données
        selected_isins = [isin for isin in all_candidates if isin in valid_all]
        _isin_slot_label: Dict[str, str] = {}
        for _isin in selected_isins:
            _row = funds_df.loc[funds_df["isin"] == _isin]
            _isin_slot_label[_isin] = _row["category"].values[0] if not _row.empty else "UC"
        if not selected_isins:
            st.info("Aucun fonds selectionné après filtrage (historique insuffisant sur la période).")
            return

        # ── Pré-sélection ────────────────────────────────────
        _preselected_isins: List[str] = []

        if _alloc_mode_choice == "Par nombre total":
            # Sélection directe par l'objectif choisi, sans contrainte de catégorie
            _valid_actions = [isin for isin in st.session_state.get("PP_SELECTED_ACTIONS", []) if isin in valid_all]
            _valid_bonds = [isin for isin in st.session_state.get("PP_SELECTED_BONDS", []) if isin in valid_all]

            # Sélectionner les meilleurs UC
            if _valid_actions and returns_all is not None and not returns_all.empty:
                _uc_candidates = [c for c in _valid_actions if c in returns_all.columns]
                _preselected_uc = _select_by_objective(
                    _uc_candidates,
                    returns_all[_uc_candidates] if _uc_candidates else pd.DataFrame(),
                    _nb_uc_total,
                    objective_choice,
                    forced=forced_isin if forced_isin in _uc_candidates else None,
                )
            else:
                _preselected_uc = _valid_actions[:_nb_uc_total]

            # Sélectionner les meilleurs obligataires
            if _valid_bonds and _nb_obl_total > 0 and returns_all is not None and not returns_all.empty:
                _obl_candidates = [c for c in _valid_bonds if c in returns_all.columns]
                _preselected_obl = _select_by_objective(
                    _obl_candidates,
                    returns_all[_obl_candidates] if _obl_candidates else pd.DataFrame(),
                    _nb_obl_total,
                    objective_choice,
                    forced=None,
                )
            else:
                _preselected_obl = _valid_bonds[:_nb_obl_total]

            _preselected_isins = _preselected_uc + _preselected_obl

            # Afficher les fonds non retenus
            _removed_in_presel = [isin for isin in valid_all if isin not in _preselected_isins]
            if _removed_in_presel:
                st.info(f"ℹ️ {len(_removed_in_presel)} fonds non retenus par l'algorithme")
                with st.expander("Voir le détail des fonds non retenus", expanded=False):
                    for _isin_nr in _removed_in_presel:
                        _name_nr = _global_isin_name.get(_isin_nr, _isin_nr)
                        if status_all.get(_isin_nr) != "ok":
                            _reason = "Historique insuffisant sur la période d'analyse"
                        else:
                            _reason = "Non sélectionné par l'objectif d'optimisation choisi"
                        st.caption(f"• {_name_nr} — {_reason}")

            all_candidates = _preselected_isins
            selected_isins = _preselected_isins
            returns_all = returns_all[[c for c in returns_all.columns if c in all_candidates]]

        else:  # Par catégorie
            for _cat, _n_wanted in _cat_counts.items():
                if _n_wanted <= 0:
                    continue
                _cat_isins = [isin for isin in _cats_in_selection.get(_cat, []) if isin in valid_all]
                if len(_cat_isins) <= _n_wanted:
                    _preselected_isins.extend(_cat_isins)
                else:
                    _cat_returns = returns_all[_cat_isins] if not returns_all.empty else pd.DataFrame()
                    if not _cat_returns.empty:
                        rfr = get_risk_free_rate()
                        _sharpes: Dict[str, float] = {}
                        for isin in _cat_isins:
                            _ret = float((1 + _cat_returns[isin].mean()) ** 252 - 1)
                            _vol = float(_cat_returns[isin].std() * np.sqrt(252))
                            _sharpes[isin] = (_ret - rfr) / _vol if _vol > 0 else 0.0
                        _sorted_by_sharpe = sorted(_sharpes.items(), key=lambda x: x[1], reverse=True)
                        _preselected_isins.extend([isin for isin, _ in _sorted_by_sharpe[:_n_wanted]])
                    else:
                        _preselected_isins.extend(_cat_isins[:_n_wanted])
            # S'assurer que le fonds conservé est inclus
            if forced_isin and forced_isin not in _preselected_isins and forced_isin in valid_all:
                _preselected_isins.append(forced_isin)
            if _preselected_isins:
                _removed_in_presel = [isin for isin in valid_all if isin not in _preselected_isins]
                if _removed_in_presel:
                    st.info(f"ℹ️ {len(_removed_in_presel)} fonds non retenus par l'algorithme")
                    with st.expander("Voir le détail des fonds non retenus", expanded=False):
                        for _isin_nr in _removed_in_presel:
                            _name_nr = _global_isin_name.get(_isin_nr, _isin_nr)
                            if status_all.get(_isin_nr) != "ok":
                                _reason = "Historique insuffisant sur la période d'analyse"
                            else:
                                _reason = "Non sélectionné par l'objectif d'optimisation choisi"
                            st.caption(f"• {_name_nr} — {_reason}")
                all_candidates = _preselected_isins
                selected_isins = _preselected_isins
                returns_all = returns_all[[c for c in returns_all.columns if c in all_candidates]]

        returns_selected = returns_all[selected_isins].dropna(how="any")
        if returns_selected.empty:
            st.info("Historique insuffisant apres selection des UC.")
            return

        # ── Architecture 3 poches : euros / actions / obligations ──
        _act_frac = act_pct / 100.0
        _obl_frac = obl_pct / 100.0
        _euro_frac = float(euro_pct) / 100.0

        if _act_frac <= 0.0 and _obl_frac <= 0.0:
            st.info("Part UC nulle avec le profil choisi.")
            return

        # Séparer les fonds sélectionnés en actions et obligations
        _sel_actions = [i for i in selected_isins if i in st.session_state.get("PP_SELECTED_ACTIONS", [])]
        _sel_bonds = [i for i in selected_isins if i in st.session_state.get("PP_SELECTED_BONDS", [])]

        # ── POCHE ACTIONS ─────────────────────────────────
        weights_act_raw: Dict[str, float] = {}
        if _sel_actions and _act_frac > 0:
            _act_in_ret = [c for c in _sel_actions if c in returns_selected.columns]
            if _act_in_ret:
                _ret_act = returns_selected[_act_in_ret]
                cap_act = 0.20 / _act_frac if _act_frac > 0 else 1.0
                cap_act = min(cap_act, 1.0)

                if PYPFOPT_AVAILABLE:
                    try:
                        mu = expected_returns.mean_historical_return(_ret_act, returns_data=True, frequency=252)
                        cov = risk_models.sample_cov(_ret_act, returns_data=True, frequency=252)
                        ef = EfficientFrontier(mu, cov, weight_bounds=(0.0, cap_act))
                        if objective_choice == "Minimiser la volatilité":
                            ef.min_volatility()
                        else:
                            ef.max_sharpe(risk_free_rate=get_risk_free_rate())
                        weights_act_raw = ef.clean_weights()
                    except Exception:
                        weights_act_raw = {}

                if not weights_act_raw:
                    rfr = get_risk_free_rate()
                    if objective_choice == "Minimiser la volatilité":
                        vol = _ret_act.std() * np.sqrt(252.0)
                        score = (1.0 / vol.replace(0, np.nan)).fillna(0.0)
                    else:
                        ann_ret = (1 + _ret_act.mean()) ** 252 - 1
                        vol = _ret_act.std() * np.sqrt(252.0)
                        score = ((ann_ret - rfr) / vol.replace(0, np.nan)).fillna(0.0).clip(lower=0.0)
                    weights_act_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}

                if not weights_act_raw:
                    weights_act_raw = {i: 1.0 / len(_act_in_ret) for i in _act_in_ret}

                weights_act_raw = {k: float(v) for k, v in weights_act_raw.items() if k in _act_in_ret}
                _total_wa = sum(weights_act_raw.values())
                if _total_wa > 0:
                    weights_act_raw = {k: v / _total_wa for k, v in weights_act_raw.items()}

        # ── POCHE OBLIGATIONS ─────────────────────────────
        weights_obl_raw: Dict[str, float] = {}
        if _sel_bonds and _obl_frac > 0:
            _obl_in_ret = [c for c in _sel_bonds if c in returns_selected.columns]
            if _obl_in_ret:
                _ret_obl = returns_selected[_obl_in_ret]
                _obl_ann_ret = ((1 + _ret_obl.mean()) ** 252 - 1)
                _obl_ann_vol = _ret_obl.std() * np.sqrt(252.0)

                if _obj_oblig == "Meilleur rendement":
                    # Tri par rendement décroissant, répartition équipondérée
                    _ranked = _obl_ann_ret.sort_values(ascending=False).index.tolist()
                    _n_obl = len(_ranked)
                    weights_obl_raw = {i: 1.0 / _n_obl for i in _ranked}

                elif _obj_oblig == "Meilleur rendement-volatilité":
                    # Ratio rendement / vol (sans soustraire taux sans risque)
                    _ratio_rv = (_obl_ann_ret / _obl_ann_vol.replace(0, np.nan)).fillna(0.0)
                    # Répartition par inverse de vol
                    _inv_vol = (1.0 / _obl_ann_vol.replace(0, np.nan)).fillna(0.0)
                    weights_obl_raw = (_inv_vol / _inv_vol.sum()).to_dict() if float(_inv_vol.sum()) > 0 else {}

                else:  # "Le plus stable possible"
                    # Exclure rendement négatif si possible (garder au moins tous)
                    _pos_ret = [i for i in _obl_in_ret if float(_obl_ann_ret.get(i, 0)) >= 0]
                    _pool = _pos_ret if len(_pos_ret) >= 1 else _obl_in_ret
                    _inv_vol = pd.Series({i: 1.0 / max(float(_obl_ann_vol.get(i, 0.01)), 0.001) for i in _pool})
                    weights_obl_raw = (_inv_vol / _inv_vol.sum()).to_dict() if float(_inv_vol.sum()) > 0 else {}

                if not weights_obl_raw:
                    weights_obl_raw = {i: 1.0 / len(_obl_in_ret) for i in _obl_in_ret}

                weights_obl_raw = {k: float(v) for k, v in weights_obl_raw.items()}
                _total_wo = sum(weights_obl_raw.values())
                if _total_wo > 0:
                    weights_obl_raw = {k: v / _total_wo for k, v in weights_obl_raw.items()}

        # ── Fusionner dans weights_uc_raw (pour compatibilité filtres doublons en aval) ──
        weights_uc_raw: Dict[str, float] = {}
        _total_uc_frac = _act_frac + _obl_frac
        for isin, w in weights_act_raw.items():
            weights_uc_raw[isin] = w * (_act_frac / _total_uc_frac) if _total_uc_frac > 0 else w
        for isin, w in weights_obl_raw.items():
            weights_uc_raw[isin] = w * (_obl_frac / _total_uc_frac) if _total_uc_frac > 0 else w

        if not weights_uc_raw:
            weights_uc_raw = {isin: 1.0 / max(len(selected_isins), 1) for isin in selected_isins}

        uc_total = _total_uc_frac
        cap_uc_final = 0.20
        uc_max_bound = min(cap_uc_final / max(uc_total, 0.01), 1.0)

        weights_uc_raw = _apply_weight_caps(weights_uc_raw, uc_max_bound)
        total_uc_raw = float(sum(weights_uc_raw.values()))
        if total_uc_raw > 0:
            weights_uc_raw = {k: v / total_uc_raw for k, v in weights_uc_raw.items()}

        # ── Filtre anti-doublon ─────────────────────────────────
        def _normalize_fund_name(name: str) -> str:
            import re
            n = name.lower().strip()
            # Supprimer les variantes de devise / classe de parts / version
            n = re.sub(
                r'\b(hedged?|acc|accumulation|dist|distribution|inc|income|'
                r'eur|usd|gbp|chf|jpy|sek|nok|dkk|'
                r'a|b|c|d|e|f|i|r|x|y|z|'
                r'retail|institutional|instit|'
                r'h1|h2|c1|c2|i1|i2|r1|r2|'
                r'class|share|part|tranche|'
                r'cap|capi|capitalisation|capitalization|'
                r'\d{4})\b',
                '',
                n,
            )
            n = re.sub(r'[^a-z0-9]+', ' ', n).strip()
            return n

        # FILTRE 1 — Doublons de nom
        if len(selected_isins) >= 2:
            _name_remove: set = set()
            for _i_n, _isin_i in enumerate(selected_isins):
                if _isin_i in _name_remove:
                    continue
                _raw_i = CONTRACT_FUND_NAMES.get(_isin_i, FUND_NAME_MAP.get(_isin_i, _isin_i))
                _norm_i = _normalize_fund_name(_raw_i)
                for _j_n, _isin_j in enumerate(selected_isins):
                    if _j_n <= _i_n or _isin_j in _name_remove:
                        continue
                    _raw_j = CONTRACT_FUND_NAMES.get(_isin_j, FUND_NAME_MAP.get(_isin_j, _isin_j))
                    _norm_j = _normalize_fund_name(_raw_j)
                    _shorter, _longer = (
                        (_norm_i, _norm_j) if len(_norm_i) <= len(_norm_j) else (_norm_j, _norm_i)
                    )
                    _is_doublon = (
                        _norm_i == _norm_j
                        or (
                            _shorter in _longer
                            and len(_shorter) >= 0.70 * len(_longer)
                        )
                    )
                    if _is_doublon:
                        _w_i = weights_uc_raw.get(_isin_i, 0)
                        _w_j = weights_uc_raw.get(_isin_j, 0)
                        if forced_isin == _isin_i:
                            _name_remove.add(_isin_j)
                        elif forced_isin == _isin_j:
                            _name_remove.add(_isin_i)
                        elif _w_i >= _w_j:
                            _name_remove.add(_isin_j)
                        else:
                            _name_remove.add(_isin_i)
            if _name_remove:
                _name_removed_labels = [CONTRACT_FUND_NAMES.get(i, FUND_NAME_MAP.get(i, i)) for i in _name_remove]
                st.warning(
                    f"⚠️ {len(_name_remove)} doublon(s) de nom retirés : "
                    + ", ".join(_name_removed_labels)
                )
                for _isin_rm in _name_remove:
                    weights_uc_raw.pop(_isin_rm, None)
                _total_w = sum(weights_uc_raw.values())
                if _total_w > 0:
                    weights_uc_raw = {k: v / _total_w for k, v in weights_uc_raw.items()}
                selected_isins = [isin for isin in selected_isins if isin not in _name_remove]

        # FILTRE 2 — Doublons de corrélation (>0.85)
        if len(selected_isins) >= 2 and not returns_all.empty:
            _corr_matrix = returns_all[[c for c in selected_isins if c in returns_all.columns]].corr()
            _to_remove: set = set()
            _checked: set = set()
            for _i_cd, _isin_i in enumerate(selected_isins):
                if _isin_i in _to_remove:
                    continue
                for _j_cd, _isin_j in enumerate(selected_isins):
                    if _j_cd <= _i_cd or _isin_j in _to_remove:
                        continue
                    _pair = tuple(sorted([_isin_i, _isin_j]))
                    if _pair in _checked:
                        continue
                    _checked.add(_pair)
                    try:
                        _c = float(_corr_matrix.loc[_isin_i, _isin_j])
                    except (KeyError, ValueError):
                        continue
                    if _c > 0.85:
                        _w_i = weights_uc_raw.get(_isin_i, 0)
                        _w_j = weights_uc_raw.get(_isin_j, 0)
                        if forced_isin == _isin_i:
                            _to_remove.add(_isin_j)
                        elif forced_isin == _isin_j:
                            _to_remove.add(_isin_i)
                        elif _w_i >= _w_j:
                            _to_remove.add(_isin_j)
                        else:
                            _to_remove.add(_isin_i)
            if _to_remove:
                _removed_names = [CONTRACT_FUND_NAMES.get(i, FUND_NAME_MAP.get(i, i)) for i in _to_remove]
                st.warning(
                    f"⚠️ {len(_to_remove)} fonds retirés car trop corrélés (>85%) : "
                    + ", ".join(_removed_names)
                )
                for _isin_rm in _to_remove:
                    weights_uc_raw.pop(_isin_rm, None)
                _total_w = sum(weights_uc_raw.values())
                if _total_w > 0:
                    weights_uc_raw = {k: v / _total_w for k, v in weights_uc_raw.items()}
                selected_isins = [isin for isin in selected_isins if isin not in _to_remove]

        # ── Limitation mode terrain : Top-N fonds ─────────────
        if practical_mode:
            _uc_frac = 1.0 - float(euro_pct) / 100.0
            _max_n = max(1, int(_uc_frac / 0.10))
            if len(weights_uc_raw) > _max_n:
                _sorted_by_weight = sorted(weights_uc_raw.items(), key=lambda x: x[1], reverse=True)
                _kept = dict(_sorted_by_weight[:_max_n])
                if forced_isin and forced_isin not in _kept:
                    _last_key = list(_kept.keys())[-1]
                    _kept.pop(_last_key)
                    _kept[forced_isin] = weights_uc_raw.get(forced_isin, 0.1)
                _kept_total = sum(_kept.values())
                weights_uc_raw = {k: v / _kept_total for k, v in _kept.items()} if _kept_total > 0 else {k: 1.0 / len(_kept) for k in _kept}
                _n_removed = len(selected_isins) - len(_kept)
                selected_isins = list(weights_uc_raw.keys())
                st.info(f"ℹ️ Mode terrain : {_n_removed} fonds retirés (poids insuffisants). {len(selected_isins)} fonds conservés.")

        if practical_mode:
            weights_uc_raw = _apply_practical_constraints(
                weights_uc_raw,
                min_w=0.10,
                step=0.05,
                max_w=uc_max_bound,
            )

        # ── Assembler l'allocation par poche (budgets séparés) ──
        # Recalculer les poids par poche après filtres doublons
        _act_remaining = [i for i in selected_isins if i in weights_act_raw]
        _obl_remaining = [i for i in selected_isins if i in weights_obl_raw]

        # Renormaliser chaque poche séparément
        _act_w_sum = sum(weights_act_raw.get(i, 0) for i in _act_remaining)
        _obl_w_sum = sum(weights_obl_raw.get(i, 0) for i in _obl_remaining)

        # Appliquer les contraintes terrain par poche
        if practical_mode and _act_remaining:
            _act_w_normed = {i: weights_act_raw[i] / _act_w_sum for i in _act_remaining} if _act_w_sum > 0 else {i: 1.0 / len(_act_remaining) for i in _act_remaining}
            _act_w_normed = _apply_practical_constraints(_act_w_normed, min_w=0.10, step=0.05, max_w=min(0.25 / max(_act_frac, 0.01), 1.0))
        elif _act_remaining:
            _act_w_normed = {i: weights_act_raw[i] / _act_w_sum for i in _act_remaining} if _act_w_sum > 0 else {i: 1.0 / len(_act_remaining) for i in _act_remaining}
        else:
            _act_w_normed = {}

        if practical_mode and _obl_remaining:
            _obl_w_normed = {i: weights_obl_raw[i] / _obl_w_sum for i in _obl_remaining} if _obl_w_sum > 0 else {i: 1.0 / len(_obl_remaining) for i in _obl_remaining}
            _obl_w_normed = _apply_practical_constraints(_obl_w_normed, min_w=0.10, step=0.05, max_w=min(0.25 / max(_obl_frac, 0.01), 1.0) if _obl_frac > 0 else 1.0)
        elif _obl_remaining:
            _obl_w_normed = {i: weights_obl_raw[i] / _obl_w_sum for i in _obl_remaining} if _obl_w_sum > 0 else {i: 1.0 / len(_obl_remaining) for i in _obl_remaining}
        else:
            _obl_w_normed = {}

        # Budgets par poche
        act_budget = float(total_budget) * _act_frac
        obl_budget = float(total_budget) * _obl_frac
        euro_budget = float(total_budget) * _euro_frac

        # Si aucun fonds obligataire retenu, redistribuer le budget obl vers actions
        if not _obl_remaining and obl_budget > 0 and _act_remaining:
            act_budget += obl_budget
            obl_budget = 0.0

        # Construire weights_final pour la compatibilité avec le code aval
        weights_final: Dict[str, float] = {}
        _total_b = float(total_budget)
        if _total_b > 0:
            if euro_budget > 0:
                weights_final["EUROFUND"] = euro_budget / _total_b
            for i, w in _act_w_normed.items():
                weights_final[i] = (act_budget * w) / _total_b
            for i, w in _obl_w_normed.items():
                weights_final[i] = (obl_budget * w) / _total_b
        else:
            weights_final = {"EUROFUND": 1.0}

        # Security clamp + exact renormalization
        weights_final = {k: max(0.0, float(v)) for k, v in weights_final.items()}
        total_weight = float(sum(weights_final.values()))
        if total_weight > 0:
            weights_final = {k: v / total_weight for k, v in weights_final.items()}
        # Update selected_isins to match
        selected_isins = [i for i in selected_isins if i in weights_final and i != "EUROFUND"]

        def _round_allocations_to_step(amounts: Dict[str, float], step: int, total: int) -> Dict[str, int]:
            if step <= 0:
                step = 1
            rounded = {k: int(np.floor(max(0.0, v) / step)) * step for k, v in amounts.items()}
            remaining = int(total - sum(rounded.values()))
            if remaining > 0 and rounded:
                frac_rank = sorted(
                    ((k, (max(0.0, amounts[k]) / step) - np.floor(max(0.0, amounts[k]) / step)) for k in rounded.keys()),
                    key=lambda x: x[1],
                    reverse=True,
                )
                idx = 0
                while remaining >= step and frac_rank:
                    k = frac_rank[idx % len(frac_rank)][0]
                    rounded[k] += step
                    remaining -= step
                    idx += 1
                if remaining != 0:
                    last_key = list(rounded.keys())[-1]
                    rounded[last_key] = max(0, rounded[last_key] + remaining)
            elif remaining < 0 and rounded:
                last_key = list(rounded.keys())[-1]
                rounded[last_key] = max(0, rounded[last_key] + remaining)

            delta = int(total - sum(rounded.values()))
            if rounded and delta != 0:
                last_key = list(rounded.keys())[-1]
                rounded[last_key] = max(0, rounded[last_key] + delta)
            return rounded

        amounts_raw = {k: float(total_budget) * float(w) for k, w in weights_final.items()}
        amounts = _round_allocations_to_step(amounts_raw, step=10, total=int(total_budget))

        rows = [
            {
                "Support": "Fonds en euros (EUROFUND)",
                "Categorie": "EUROFUND",
                "%": weights_final.get("EUROFUND", 0.0) * 100.0,
                "Montant EUR": amounts.get("EUROFUND", 0),
            }
        ]
        for isin in selected_isins:
            rows.append(
                {
                    "Support": CONTRACT_FUND_NAMES.get(isin, FUND_NAME_MAP.get(isin, isin)),
                    "Categorie": _isin_slot_label.get(isin, "UC"),
                    "%": weights_final.get(isin, 0.0) * 100.0,
                    "Montant EUR": amounts.get(isin, 0),
                }
            )
        df_alloc = pd.DataFrame(rows)

        st.markdown("**Allocation finale**")
        st.dataframe(
            df_alloc.style.format({"%": "{:,.2f}%".format, "Montant EUR": to_eur}),
            use_container_width=True,
            hide_index=True,
        )

        if MATPLOTLIB_AVAILABLE:
            fig, ax = plt.subplots(figsize=(5.0, 3.2))
            ax.pie(df_alloc["Montant EUR"], labels=df_alloc["Support"], autopct="%1.1f%%")
            ax.set_title("Repartition finale")
            st.pyplot(fig)
            plt.close(fig)
        else:
            st.warning(f"Graphique indisponible ({MATPLOTLIB_ERROR}).")

        # ── Performance du portefeuille global (3 poches) ──
        _all_opt_lines = []
        for _isin_m, _amt_m in amounts.items():
            if _amt_m <= 0:
                continue
            _all_opt_lines.append({
                "isin": _isin_m,
                "name": _global_isin_name.get(_isin_m, CONTRACT_FUND_NAMES.get(_isin_m, _isin_m)),
                "amount_gross": float(_amt_m),
                "buy_date": opt_start,
            })

        _global_sharpe = compute_sharpe_ratio(_all_opt_lines, float(euro_rate), 0.0, start_date=opt_start)
        _global_risk = portfolio_risk_stats(_all_opt_lines, float(euro_rate), start_date=opt_start, fee_pct=0.0)
        _global_rp = _portfolio_weighted_returns(_all_opt_lines, float(euro_rate), 0.0, start_date=opt_start)

        st.markdown("#### Performance du portefeuille global")
        st.caption(f"Période : du {fmt_date(opt_start)} au {fmt_date(opt_end)}")

        if _global_rp is not None:
            _global_ann_ret = float(_global_rp.mean() * 252 * 100)
            _mc1, _mc2, _mc3 = st.columns(3)
            with _mc1:
                st.metric("Rendement annualisé global", f"{_global_ann_ret:+.1f}%")
            with _mc2:
                st.metric("Volatilité annualisée", f"{_global_risk['vol_ann_pct']:.1f}%" if _global_risk else "—")
            with _mc3:
                st.metric("Sharpe du portefeuille", f"{_global_sharpe:.2f}" if _global_sharpe else "—")
        else:
            st.caption("Données insuffisantes pour calculer les métriques globales.")

        # Détail par poche
        with st.expander("Détail par poche", expanded=False):
            _pp_selected_actions = list(st.session_state.get("PP_SELECTED_ACTIONS", []))
            _pp_selected_bonds = list(st.session_state.get("PP_SELECTED_BONDS", []))
            _act_lines_m = [l for l in _all_opt_lines if l["isin"] in _pp_selected_actions]
            _obl_lines_m = [l for l in _all_opt_lines if l["isin"] in _pp_selected_bonds]
            _act_rp_m = _portfolio_weighted_returns(_act_lines_m, float(euro_rate), 0.0, start_date=opt_start) if _act_lines_m else None
            _obl_rp_m = _portfolio_weighted_returns(_obl_lines_m, float(euro_rate), 0.0, start_date=opt_start) if _obl_lines_m else None
            _pc1, _pc2, _pc3 = st.columns(3)
            with _pc1:
                _act_ret_ann = float(_act_rp_m.mean() * 252 * 100) if _act_rp_m is not None else 0
                st.metric(f"Poche Actions ({act_pct:.0f}%)", f"{_act_ret_ann:+.1f}%/an")
            with _pc2:
                _obl_ret_ann = float(_obl_rp_m.mean() * 252 * 100) if _obl_rp_m is not None else 0
                st.metric(f"Poche Obligations ({obl_pct:.0f}%)", f"{_obl_ret_ann:+.1f}%/an")
            with _pc3:
                st.metric(f"Fonds euros ({int(euro_pct)}%)", f"+{float(euro_rate):.1f}%/an")

        rfr_display = get_risk_free_rate()
        rfr_src = st.session_state.get("RISK_FREE_RATE_SOURCE", "manual")
        rfr_label = (
            f"📡 {st.session_state.get('RFR_BOND_REF', 'Bund 10 ans')} ({rfr_display * 100:.2f}%)"
            if rfr_src == "api"
            else f"✏️ Taux saisi manuellement ({rfr_display * 100:.2f}%)"
        )
        st.caption(
            f"Ratio de Sharpe calculé avec un taux sans risque de "
            f"**{rfr_display * 100:.2f}%** — {rfr_label}"
        )

        st.markdown(
            """
- Allocation calibree sur la fenetre d'analyse et l'objectif choisi.
- EUROFUND est traite uniquement via son taux annuel parametre.
- Les UC respectent un cap strict de 25% du portefeuille final par fonds.
- Si un fonds est marqué 📌 Conserver, il est gardé dans la poche actions.
- En cas de donnees insuffisantes, la selection est reduite automatiquement.
"""
        )

        if practical_mode:
            st.caption(
                "Les ponderations UC sont ajustees pour rester coherentes en pratique "
                "(min 10 % et arrondi par paliers de 5 %)."
            )

        if insufficient:
            st.caption("Fonds exclus : " + ", ".join(insufficient))
        st.caption(f"Fenetre utilisee : {fmt_date(opt_start)} a {fmt_date(opt_end)}")

        # ── Tableau détaillé par fonds ─────────────────────────────
        st.markdown("#### Détail par fonds")
        _detail_rows = []
        _all_isins_detail = list(amounts.keys()) if amounts else []
        for _isin_det in _all_isins_detail:
            if _isin_det.upper() in ("EUROFUND", "STRUCTURED"):
                continue
            _amt_det = amounts.get(_isin_det, 0)
            if _amt_det <= 0:
                continue
            _name_det = CONTRACT_FUND_NAMES.get(_isin_det, _isin_det)
            _cat_det_row = funds_df.loc[funds_df["isin"].str.upper() == _isin_det.upper()] if not funds_df.empty else pd.DataFrame()
            _cat_det = str(_cat_det_row["category"].values[0]) if not _cat_det_row.empty else ""
            _is_bond_det = _is_bond_category(_cat_det)
            try:
                _df_det, _, _ = get_price_series(_isin_det, None, float(euro_rate))
                if _df_det.empty:
                    continue
                _s_det = _df_det["Close"].astype(float)
                _first_vl_det = _s_det.index.min()
                _s_det = _s_det[(_s_det.index >= opt_start) & (_s_det.index <= opt_end)]
                if len(_s_det) < 60:
                    continue
                _rets_det = _s_det.pct_change().dropna()
                if len(_rets_det) < 60:
                    continue
                _ann_ret_det = float((1 + _rets_det.mean()) ** 252 - 1)
                _ann_vol_det = float(_rets_det.std() * np.sqrt(252))
                _rfr_det = get_risk_free_rate()
                _sharpe_det = float((_ann_ret_det - _rfr_det) / _ann_vol_det) if _ann_vol_det > 0.001 else 0.0
                _detail_rows.append({
                    "Fonds": _name_det[:40],
                    "Rdt/an": f"{_ann_ret_det*100:+.1f}%",
                    "Sharpe": f"{_sharpe_det:.2f}" if not _is_bond_det else "—",
                    "Vol/an": f"{_ann_vol_det*100:.1f}%",
                    "1ère VL": _first_vl_det.strftime("%d/%m/%Y"),
                })
            except Exception:
                continue
        if _detail_rows:
            try:
                st.dataframe(pd.DataFrame(_detail_rows), hide_index=True, use_container_width=True)
            except Exception:
                pass
        st.caption(f"Ratios calculés sur la période du {fmt_date(opt_start)} au {fmt_date(opt_end)}")

        # ── Rendu enrichi des résultats (mode Créer uniquement) ──
        _is_improve_active = _is_improve_mode and st.session_state.get("PP_IMPROVE_LINES")
        if not _is_improve_active:
            _opt_lines_for_div = [
                {"isin": _isin_opt, "name": CONTRACT_FUND_NAMES.get(_isin_opt, _isin_opt), "amount_gross": float(_amt_opt)}
                for _isin_opt, _amt_opt in amounts.items()
                if _amt_opt > 0 and _isin_opt.upper() not in ("EUROFUND", "STRUCTURED")
            ]
            _opt_div = compute_diversification_score(_opt_lines_for_div, float(euro_rate), start_date=opt_start) if len(_opt_lines_for_div) >= 2 else None

            # Réutiliser les métriques globales déjà calculées
            _create_sharpe = _global_sharpe
            _create_sortino = compute_sortino_ratio(_all_opt_lines, float(euro_rate), 0.0, start_date=opt_start)
            _create_risk = _global_risk
            _create_rp = _global_rp

            st.markdown("---")

            # 0) Courbe d'évolution en EUR (mode Créer)
            _montant_investi = float(total_budget)
            if _create_rp is not None and MATPLOTLIB_AVAILABLE:
                try:
                    st.markdown("#### Évolution de la valeur du portefeuille")
                    st.caption(f"Période : du {fmt_date(opt_start)} au {fmt_date(opt_end)}")
                    _cum_create = (1 + _create_rp).cumprod() * _montant_investi
                    fig, ax = plt.subplots(figsize=(10, 4))
                    ax.plot(_cum_create.index, _cum_create.values, color="#1A7A4A", linewidth=1.5, label="Portefeuille")
                    ax.axhline(y=_montant_investi, color="#999", linestyle="--", alpha=0.5, linewidth=0.8)
                    ax.set_ylabel("Valeur du portefeuille (€)")
                    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:,.0f}".replace(",", " ")))
                    ax.legend(loc="upper left")
                    ax.grid(True, alpha=0.3)
                    ax.set_title("Évolution de la valeur du portefeuille")
                    st.pyplot(fig)
                    plt.close(fig)

                    _montant_final = float(_cum_create.iloc[-1])
                    _gain_cumule = _montant_final - _montant_investi
                    _gain_pct = (_montant_final / _montant_investi - 1) * 100 if _montant_investi > 0 else 0
                    _years_create = len(_create_rp) / 252.0
                    _gain_ann = ((_montant_final / _montant_investi) ** (1 / _years_create) - 1) * 100 if _years_create > 0 else 0
                    _mc1, _mc2, _mc3, _mc4 = st.columns(4)
                    with _mc1:
                        st.metric("Montant investi", to_eur(_montant_investi))
                    with _mc2:
                        st.metric("Montant final", to_eur(_montant_final))
                    with _mc3:
                        st.metric("Gain cumulé", f"{to_eur(_gain_cumule)} ({_gain_pct:+.1f}%)")
                    with _mc4:
                        st.metric("Gain annualisé", f"{_gain_ann:+.1f}%/an")
                except Exception as _e_create_graph:
                    st.caption(f"Graphique non disponible : {_e_create_graph}")
            elif _create_rp is None:
                st.info("Graphique non disponible : données insuffisantes pour les fonds sélectionnés.")

            # A) Donut diversification
            _opt_div_score = float(_opt_div["score"]) if _opt_div and _opt_div.get("score") is not None else 0
            st.markdown("#### 📊 Diversification du portefeuille")
            try:
                _gauge_create = _build_gauge_chart(_opt_div_score, "Portefeuille")
                if _gauge_create:
                    st.altair_chart(_gauge_create, use_container_width=True)
                else:
                    st.metric("Score diversification", f"{_opt_div_score:.0f}/100")
            except Exception:
                st.metric("Score diversification", f"{_opt_div_score:.0f}/100")

            if _opt_div_score >= 70:
                st.success(f"Bonne diversification ({_opt_div_score:.0f}/100)")
            elif _opt_div_score >= 40:
                st.info(f"Diversification correcte ({_opt_div_score:.0f}/100) — envisagez d'ajouter des fonds décorrélés")
            else:
                st.warning(f"Diversification insuffisante ({_opt_div_score:.0f}/100) — les fonds sont trop corrélés entre eux")

            # B) Métriques avec phrases contextuelles
            st.markdown("#### Indicateurs de performance et de risque")
            st.caption(f"Période d'analyse : du {fmt_date(opt_start)} au {fmt_date(opt_end)}")

            # Rendement annualisé
            if _create_rp is not None:
                _create_ann_ret = float(_create_rp.mean() * 252 * 100)
                st.metric("Rendement annualisé", f"{_create_ann_ret:+.1f}%")
                if _create_ann_ret > 8:
                    st.caption("Rendement élevé — à mettre en perspective avec le risque pris.")
                elif _create_ann_ret > 4:
                    st.caption("Rendement solide sur la période analysée.")
                elif _create_ann_ret > 0:
                    st.caption("Rendement positif mais modéré.")

            # Sharpe
            if _create_sharpe:
                st.metric("Rendement / Risque (Sharpe)", f"{_create_sharpe:.2f}")
                if _create_sharpe >= 0.8:
                    st.caption("Excellent rapport rendement/risque.")
                elif _create_sharpe >= 0.4:
                    st.caption("Bon rapport rendement/risque.")
                else:
                    st.caption("Rapport rendement/risque à améliorer — envisagez de diversifier davantage.")

            # Sortino
            if _create_sortino:
                st.metric("Protection baissière (Sortino)", f"{_create_sortino:.2f}")
                if _create_sortino >= 1.0:
                    st.caption("Très bonne protection en cas de baisse des marchés.")
                elif _create_sortino >= 0.5:
                    st.caption("Protection correcte en cas de baisse.")

            # Volatilité
            if _create_risk:
                _create_vol = _create_risk["vol_ann_pct"]
                _create_dd = _create_risk["max_dd_pct"]
                st.metric("Volatilité annuelle", f"{_create_vol:.1f}%")
                if _create_vol < 8:
                    st.caption("Faible volatilité — portefeuille défensif.")
                elif _create_vol < 15:
                    st.caption("Volatilité modérée — adapté à un profil équilibré.")
                else:
                    st.caption("Volatilité élevée — portefeuille dynamique.")

                st.metric("Perte maximale historique", f"{_create_dd:.1f}%")
                st.caption(f"En cas de crise comparable sur la période, la perte maximale estimée serait de {abs(_create_dd):.1f}%.")

            # Corrélation
            if _opt_div and _opt_div.get("avg_corr") is not None:
                _create_corr = _opt_div["avg_corr"]
                st.metric("Corrélation moyenne", f"{_create_corr:.0%}")
                if _create_corr < 0.4:
                    st.caption("Les fonds ont des comportements très indépendants — excellente diversification.")
                elif _create_corr < 0.6:
                    st.caption("Bonne indépendance entre les fonds.")
                else:
                    st.caption("Les fonds tendent à bouger ensemble — diversification limitée.")

            st.markdown("---")

            # C) Camembert composition
            st.markdown("#### Composition du portefeuille")
            _new_total_create = sum(float(a) for a in amounts.values())
            _create_pie_data = []
            for _isin_cp, _amt_cp in amounts.items():
                if _amt_cp <= 0:
                    continue
                _nm_cp = CONTRACT_FUND_NAMES.get(_isin_cp, _isin_cp)[:30]
                _pct_cp = _amt_cp / _new_total_create * 100 if _new_total_create > 0 else 0
                _create_pie_data.append({"Fonds": _nm_cp, "Part": _pct_cp})
            if _create_pie_data:
                try:
                    _create_pie_df = pd.DataFrame(_create_pie_data)
                    _create_pie = alt.Chart(_create_pie_df).mark_arc(innerRadius=45, outerRadius=85).encode(
                        theta=alt.Theta("Part:Q", stack=True),
                        color=alt.Color("Fonds:N",
                            scale=alt.Scale(scheme="tableau10"),
                            legend=alt.Legend(title=None, orient="right", labelLimit=250),
                        ),
                        tooltip=[
                            alt.Tooltip("Fonds:N", title="Fonds"),
                            alt.Tooltip("Part:Q", title="%", format=".1f"),
                        ],
                    ).properties(height=280)
                    _create_pie_text = alt.Chart(_create_pie_df).mark_text(radius=70, size=11, fontWeight="bold").encode(
                        theta=alt.Theta("Part:Q", stack=True),
                        text=alt.Text("Part:Q", format=".0f"),
                        color=alt.value("#333"),
                    )
                    st.altair_chart((_create_pie + _create_pie_text), use_container_width=True)
                except Exception:
                    for _cpd in _create_pie_data:
                        st.caption(f"{_cpd['Fonds']} — {_cpd['Part']:.1f}%")

        st.markdown("---")
        if st.button("📤 Envoyer vers Portefeuille Cabinet", type="primary", key="PP_SEND_TO_COMPARATOR"):
            new_lines = []
            for isin, amount in amounts.items():
                if amount <= 0:
                    continue
                fund_name = CONTRACT_FUND_NAMES.get(isin, FUND_NAME_MAP.get(isin, isin))
                fund_row = funds_df[funds_df["isin"] == isin]
                fee_uc = float(fund_row["fee_uc_pct"].iloc[0]) if not fund_row.empty else 0.0
                fee_contract = float(fund_row["fee_contract_pct"].iloc[0]) if not fund_row.empty else 0.0
                new_lines.append({
                    "id": str(uuid.uuid4()),
                    "name": fund_name,
                    "isin": isin,
                    "amount_gross": float(amount),
                    "buy_date": pd.Timestamp.today().normalize(),
                    "buy_px": "",
                    "note": f"Construction optimisée — {objective_choice}",
                    "sym_used": "",
                    "fee_uc_pct": fee_uc,
                    "fee_contract_pct": fee_contract,
                    "fee_total_pct": fee_uc + fee_contract,
                })
            st.session_state["B_lines"] = new_lines
            st.session_state["MODE_ANALYSE"] = "compare"
            st.session_state["CONTRACT_LABEL_B"] = pp_contract_label
            st.session_state["CONTRACT_FUNDS_DF_B"] = funds_df
            st.success(f"✅ {len(new_lines)} lignes envoyées vers le Portefeuille Cabinet.")
            st.caption("Basculez vers '⚖️ Comparateur' dans le menu pour voir la comparaison.")

        # ── Comparaison avant / après (mode amélioration) ──────────
        if _is_improve_mode and st.session_state.get("PP_IMPROVE_LINES"):
            st.markdown("---")
            st.markdown("### 📊 Comparaison avant / après")
            st.caption(f"Période d'analyse : du {fmt_date(opt_start)} au {fmt_date(opt_end)}")

            _old_lines_cmp = st.session_state.get("PP_IMPROVE_LINES", [])
            _euro_rate_cmp = float(euro_rate)
            _fee_cmp = 0.0

            # Date de début commune = plus ancienne date d'achat de l'ancien portefeuille
            _oldest_buy = min(
                (pd.Timestamp(l.get("buy_date", pd.Timestamp.today())) for l in _old_lines_cmp),
                default=pd.Timestamp.today(),
            ).normalize()

            # Construire les nouvelles lignes
            _new_lines_cmp = []
            for _isin_cmp, _amt_cmp in amounts.items():
                if _amt_cmp <= 0:
                    continue
                _fund_row_cmp = funds_df[funds_df["isin"] == _isin_cmp]
                _new_lines_cmp.append({
                    "isin": _isin_cmp,
                    "name": CONTRACT_FUND_NAMES.get(_isin_cmp, _isin_cmp),
                    "amount_gross": float(_amt_cmp),
                    "buy_date": _oldest_buy,
                    "fee_uc_pct": float(_fund_row_cmp["fee_uc_pct"].iloc[0]) if not _fund_row_cmp.empty else 0.0,
                    "fee_contract_pct": float(_fund_row_cmp["fee_contract_pct"].iloc[0]) if not _fund_row_cmp.empty else 0.0,
                })

            # Calculer toutes les métriques
            try:
                _old_sharpe_c = compute_sharpe_ratio(_old_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
                _new_sharpe_c = compute_sharpe_ratio(_new_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
                _old_sortino_c = compute_sortino_ratio(_old_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
                _new_sortino_c = compute_sortino_ratio(_new_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
                _old_risk_c = portfolio_risk_stats(_old_lines_cmp, _euro_rate_cmp, start_date=opt_start, fee_pct=_fee_cmp)
                _new_risk_c = portfolio_risk_stats(_new_lines_cmp, _euro_rate_cmp, start_date=opt_start, fee_pct=_fee_cmp)
                _old_div_c = compute_diversification_score(_old_lines_cmp, _euro_rate_cmp, start_date=opt_start)
                _new_div_c = compute_diversification_score(_new_lines_cmp, _euro_rate_cmp, start_date=opt_start)
            except Exception:
                _old_sharpe_c = _new_sharpe_c = _old_sortino_c = _new_sortino_c = None
                _old_risk_c = _new_risk_c = _old_div_c = _new_div_c = None

            # Rendements pondérés pour le graphique et les indicateurs
            try:
                _rp_old = _portfolio_weighted_returns(_old_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
                _rp_new = _portfolio_weighted_returns(_new_lines_cmp, _euro_rate_cmp, _fee_cmp, start_date=opt_start)
            except Exception:
                _rp_old = _rp_new = None

            # ── A) Courbe d'évolution en EUR ──────────────────────
            _montant_compare = float(total_budget)
            if MATPLOTLIB_AVAILABLE:
                try:
                    if _rp_old is not None and _rp_new is not None:
                        _common_idx = _rp_old.index.intersection(_rp_new.index)
                        if len(_common_idx) > 60:
                            _rp_old_c = _rp_old.loc[_common_idx]
                            _rp_new_c = _rp_new.loc[_common_idx]

                            _cum_old = (1 + _rp_old_c).cumprod() * _montant_compare
                            _cum_new = (1 + _rp_new_c).cumprod() * _montant_compare

                            fig, ax = plt.subplots(figsize=(10, 4))
                            ax.plot(_cum_old.index, _cum_old.values,
                                    color="#CC2200", label="Ancien portefeuille", linewidth=1.5)
                            ax.plot(_cum_new.index, _cum_new.values,
                                    color="#1A7A4A", label="Portefeuille amélioré", linewidth=1.5)
                            ax.fill_between(_cum_new.index, _cum_old.values, _cum_new.values,
                                            where=_cum_new.values >= _cum_old.values,
                                            alpha=0.08, color="#1A7A4A")
                            ax.fill_between(_cum_new.index, _cum_old.values, _cum_new.values,
                                            where=_cum_new.values < _cum_old.values,
                                            alpha=0.08, color="#CC2200")
                            ax.set_ylabel("Valeur du portefeuille (€)")
                            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:,.0f}".replace(",", " ")))
                            ax.legend(loc="upper left")
                            ax.grid(True, alpha=0.3)
                            ax.set_title("Évolution comparée de la valeur du portefeuille")
                            st.pyplot(fig)
                            plt.close(fig)

                            _final_old = float(_cum_old.iloc[-1])
                            _final_new = float(_cum_new.iloc[-1])
                            _years_period = len(_common_idx) / 252.0

                            _perf_c1, _perf_c2 = st.columns(2)
                            with _perf_c1:
                                st.metric("Ancien portefeuille", to_eur(_final_old),
                                          delta=f"{((_final_old / _montant_compare - 1) * 100):+.1f}%")
                            with _perf_c2:
                                st.metric("Portefeuille amélioré", to_eur(_final_new),
                                          delta=f"{to_eur(_final_new - _final_old)} vs ancien")
                        else:
                            st.caption("Historique commun insuffisant pour le graphique d'évolution.")
                except Exception as _e_graph:
                    st.caption(f"Graphique non disponible : {_e_graph}")

            # ── B) Donuts diversification côte à côte ────────────────
            st.markdown("#### Diversification")
            _old_score_c = float(_old_div_c["score"]) if _old_div_c and _old_div_c.get("score") is not None else 0
            _new_score_c = float(_new_div_c["score"]) if _new_div_c and _new_div_c.get("score") is not None else 0

            _dv1, _dv2 = st.columns(2)
            with _dv1:
                st.markdown("**Ancien portefeuille**")
                try:
                    _gauge_old = _build_gauge_chart(_old_score_c, "Avant")
                    if _gauge_old:
                        st.altair_chart(_gauge_old, use_container_width=True)
                    else:
                        st.metric("Score", f"{_old_score_c:.0f}/100")
                except Exception:
                    st.metric("Score", f"{_old_score_c:.0f}/100")
            with _dv2:
                st.markdown("**Portefeuille amélioré**")
                try:
                    _gauge_new = _build_gauge_chart(_new_score_c, "Après")
                    if _gauge_new:
                        st.altair_chart(_gauge_new, use_container_width=True)
                    else:
                        st.metric("Score", f"{_new_score_c:.0f}/100")
                except Exception:
                    st.metric("Score", f"{_new_score_c:.0f}/100")

            _delta_div_c = _new_score_c - _old_score_c
            if _delta_div_c > 5:
                st.success(f"✅ Diversification améliorée : +{_delta_div_c:.0f} points")
            elif _delta_div_c < -5:
                st.warning(f"⚠️ Diversification en baisse : {_delta_div_c:.0f} points")

            # ── C) Métriques en colonnes avec phrases contextuelles ──
            st.markdown("#### Indicateurs de performance et de risque")
            _mc1, _mc2 = st.columns(2)

            with _mc1:
                st.markdown("**Ancien portefeuille**")
            with _mc2:
                st.markdown("**Portefeuille amélioré**")

            # Rendement annualisé
            _ret1, _ret2 = st.columns(2)
            with _ret1:
                _old_ann_ret_display = float(_rp_old.mean() * 252 * 100) if _rp_old is not None else 0
                st.metric("Rendement annualisé", f"{_old_ann_ret_display:+.1f}%")
            with _ret2:
                _new_ann_ret_display = float(_rp_new.mean() * 252 * 100) if _rp_new is not None else 0
                _ret_delta = _new_ann_ret_display - _old_ann_ret_display
                st.metric("Rendement annualisé", f"{_new_ann_ret_display:+.1f}%",
                          delta=f"{_ret_delta:+.1f}%" if abs(_ret_delta) > 0.1 else None)
            if _new_ann_ret_display > _old_ann_ret_display + 0.5:
                st.caption(f"💬 Le nouveau portefeuille génère **+{(_new_ann_ret_display - _old_ann_ret_display):.1f}%/an de rendement supplémentaire** sur la période analysée.")

            # Sharpe
            _sh1, _sh2 = st.columns(2)
            with _sh1:
                st.metric("Rendement / Risque (Sharpe)", f"{(_old_sharpe_c or 0):.2f}")
            with _sh2:
                _sh_delta = ((_new_sharpe_c or 0) - (_old_sharpe_c or 0))
                st.metric("Rendement / Risque (Sharpe)", f"{(_new_sharpe_c or 0):.2f}",
                          delta=f"{_sh_delta:+.2f}" if _old_sharpe_c else None)
            if _old_sharpe_c and _new_sharpe_c and _old_sharpe_c > 0:
                _ratio_sh = _new_sharpe_c / _old_sharpe_c
                if _ratio_sh > 1.1:
                    st.caption(f"💬 Le nouveau portefeuille génère **{_ratio_sh:.1f}× plus de rendement** par unité de risque prise.")
                elif _ratio_sh < 0.9:
                    st.caption(f"💬 Le nouveau portefeuille génère {_ratio_sh:.1f}× de rendement par unité de risque (en baisse).")

            # Sortino
            _so1, _so2 = st.columns(2)
            with _so1:
                st.metric("Protection baissière (Sortino)", f"{(_old_sortino_c or 0):.2f}")
            with _so2:
                _so_delta = ((_new_sortino_c or 0) - (_old_sortino_c or 0))
                st.metric("Protection baissière (Sortino)", f"{(_new_sortino_c or 0):.2f}",
                          delta=f"{_so_delta:+.2f}" if _old_sortino_c else None)
            if _old_sortino_c and _new_sortino_c and _old_sortino_c > 0:
                _ratio_so = _new_sortino_c / _old_sortino_c
                if _ratio_so > 1.1:
                    st.caption(f"💬 En cas de baisse des marchés, le nouveau portefeuille offre **{_ratio_so:.1f}× plus de protection**.")

            # Volatilité
            _vol_old = _old_risk_c["vol_ann_pct"] if _old_risk_c else 0
            _vol_new = _new_risk_c["vol_ann_pct"] if _new_risk_c else 0
            _vl1, _vl2 = st.columns(2)
            with _vl1:
                st.metric("Volatilité annuelle", f"{_vol_old:.1f}%")
            with _vl2:
                st.metric("Volatilité annuelle", f"{_vol_new:.1f}%",
                          delta=f"{(_vol_new - _vol_old):.1f}%", delta_color="inverse")
            if _vol_old > 0 and _vol_new > 0:
                _red_vol = (1 - _vol_new / _vol_old) * 100
                if _red_vol > 5:
                    st.caption(f"💬 Le nouveau portefeuille est **{_red_vol:.0f}% moins volatil** — les fluctuations de valeur sont réduites.")
                elif _red_vol < -5:
                    st.caption(f"💬 Le nouveau portefeuille est {abs(_red_vol):.0f}% plus volatil — les fluctuations augmentent en échange d'un meilleur rendement potentiel.")

            # Max Drawdown
            _dd_old = _old_risk_c["max_dd_pct"] if _old_risk_c else 0
            _dd_new = _new_risk_c["max_dd_pct"] if _new_risk_c else 0
            _dd1, _dd2 = st.columns(2)
            with _dd1:
                st.metric("Perte maximale historique", f"{_dd_old:.1f}%")
            with _dd2:
                _dd_delta_display = abs(_dd_old) - abs(_dd_new)
                st.metric("Perte maximale historique", f"{_dd_new:.1f}%",
                          delta=f"{-_dd_delta_display:+.1f} pts", delta_color="inverse")
            _delta_dd = abs(_dd_new) - abs(_dd_old)
            if _delta_dd < -1:
                st.caption(f"💬 En cas de crise, la perte maximale passe de **{abs(_dd_old):.1f}%** à **{abs(_dd_new):.1f}%** — soit **{abs(_delta_dd):.1f} points de moins**.")
            elif _delta_dd > 1:
                st.caption(f"💬 La perte maximale augmente de {abs(_delta_dd):.1f} points — à mettre en perspective avec le gain de rendement.")

            # Corrélation
            _corr_old = _old_div_c["avg_corr"] if _old_div_c else 0
            _corr_new = _new_div_c["avg_corr"] if _new_div_c else 0
            _cr1, _cr2 = st.columns(2)
            with _cr1:
                st.metric("Corrélation moyenne", f"{_corr_old:.0%}")
            with _cr2:
                st.metric("Corrélation moyenne", f"{_corr_new:.0%}",
                          delta=f"{(_corr_new - _corr_old):.0%}", delta_color="inverse")
            if _corr_new < _corr_old - 0.05:
                st.caption("💬 Les fonds du nouveau portefeuille ont des **comportements beaucoup plus indépendants** — meilleure protection en cas de choc sectoriel.")
            elif _corr_new > _corr_old + 0.05:
                st.caption("💬 La corrélation augmente — les fonds ont tendance à bouger davantage ensemble.")

            st.markdown("---")

            # ── D) Composition en camemberts côte à côte ─────────────
            st.markdown("#### Composition des portefeuilles")
            _cp1, _cp2 = st.columns(2)

            with _cp1:
                st.markdown("**Ancien portefeuille**")
                _old_total_c = sum(float(l.get("amount_gross", 0) or 0) for l in _old_lines_cmp)
                _old_pie_data = []
                for _lc in _old_lines_cmp:
                    _nm_lc = str(_lc.get("name", _lc.get("isin", "—")))[:30]
                    _amt_lc = float(_lc.get("amount_gross", 0) or 0)
                    _pct_lc = _amt_lc / _old_total_c * 100 if _old_total_c > 0 else 0
                    _old_pie_data.append({"Fonds": _nm_lc, "Part": _pct_lc})
                if _old_pie_data:
                    try:
                        _old_pie_df = pd.DataFrame(_old_pie_data)
                        _old_pie = alt.Chart(_old_pie_df).mark_arc(innerRadius=45, outerRadius=85).encode(
                            theta=alt.Theta("Part:Q", stack=True),
                            color=alt.Color("Fonds:N",
                                scale=alt.Scale(scheme="tableau10"),
                                legend=alt.Legend(title=None, orient="right", labelLimit=250),
                            ),
                            tooltip=[
                                alt.Tooltip("Fonds:N", title="Fonds"),
                                alt.Tooltip("Part:Q", title="%", format=".1f"),
                            ],
                        ).properties(height=280)
                        _old_pie_text = alt.Chart(_old_pie_df).mark_text(radius=70, size=11, fontWeight="bold").encode(
                            theta=alt.Theta("Part:Q", stack=True),
                            text=alt.Text("Part:Q", format=".0f"),
                            color=alt.value("#333"),
                        )
                        st.altair_chart((_old_pie + _old_pie_text), use_container_width=True)
                    except Exception:
                        for _op in _old_pie_data:
                            st.caption(f"{_op['Fonds']} — {_op['Part']:.1f}%")

            with _cp2:
                st.markdown("**Portefeuille amélioré**")
                _new_total_c = sum(float(a) for a in amounts.values())
                _new_pie_data = []
                for _isin_nc, _amt_nc in amounts.items():
                    if _amt_nc <= 0:
                        continue
                    _nm_nc = CONTRACT_FUND_NAMES.get(_isin_nc, _isin_nc)[:30]
                    _pct_nc = _amt_nc / _new_total_c * 100 if _new_total_c > 0 else 0
                    _new_pie_data.append({"Fonds": _nm_nc, "Part": _pct_nc})
                if _new_pie_data:
                    try:
                        _new_pie_df = pd.DataFrame(_new_pie_data)
                        _new_pie = alt.Chart(_new_pie_df).mark_arc(innerRadius=45, outerRadius=85).encode(
                            theta=alt.Theta("Part:Q", stack=True),
                            color=alt.Color("Fonds:N",
                                scale=alt.Scale(scheme="tableau10"),
                                legend=alt.Legend(title=None, orient="right", labelLimit=250),
                            ),
                            tooltip=[
                                alt.Tooltip("Fonds:N", title="Fonds"),
                                alt.Tooltip("Part:Q", title="%", format=".1f"),
                            ],
                        ).properties(height=280)
                        _new_pie_text = alt.Chart(_new_pie_df).mark_text(radius=70, size=11, fontWeight="bold").encode(
                            theta=alt.Theta("Part:Q", stack=True),
                            text=alt.Text("Part:Q", format=".0f"),
                            color=alt.value("#333"),
                        )
                        st.altair_chart((_new_pie + _new_pie_text), use_container_width=True)
                    except Exception:
                        for _np in _new_pie_data:
                            st.caption(f"{_np['Fonds']} — {_np['Part']:.1f}%")

            # ── Tableau détaillé par fonds — Ancien portefeuille ──────
            st.markdown("#### Détail par fonds — Ancien portefeuille")
            _detail_old_rows = []
            for _ln_det_old in _old_lines_cmp:
                _isin_do = str(_ln_det_old.get("isin", "")).strip()
                if _isin_do.upper() in ("EUROFUND", "STRUCTURED", ""):
                    continue
                _name_do = str(_ln_det_old.get("name", _isin_do))[:40]
                _cat_do_row = funds_df.loc[funds_df["isin"].str.upper() == _isin_do.upper()] if not funds_df.empty else pd.DataFrame()
                _cat_do = str(_cat_do_row["category"].values[0]) if not _cat_do_row.empty else ""
                _is_bond_do = _is_bond_category(_cat_do)
                try:
                    _df_do, _, _ = get_price_series(_isin_do, None, float(_euro_rate_cmp))
                    if _df_do.empty:
                        continue
                    _s_do = _df_do["Close"].astype(float)
                    _first_vl_do = _s_do.index.min()
                    _s_do = _s_do[(_s_do.index >= opt_start) & (_s_do.index <= opt_end)]
                    if len(_s_do) < 60:
                        continue
                    _rets_do = _s_do.pct_change().dropna()
                    if len(_rets_do) < 60:
                        continue
                    _ann_ret_do = float((1 + _rets_do.mean()) ** 252 - 1)
                    _ann_vol_do = float(_rets_do.std() * np.sqrt(252))
                    _rfr_do = get_risk_free_rate()
                    _sharpe_do = float((_ann_ret_do - _rfr_do) / _ann_vol_do) if _ann_vol_do > 0.001 else 0.0
                    _detail_old_rows.append({
                        "Fonds": _name_do,
                        "Rdt/an": f"{_ann_ret_do*100:+.1f}%",
                        "Sharpe": f"{_sharpe_do:.2f}" if not _is_bond_do else "—",
                        "Vol/an": f"{_ann_vol_do*100:.1f}%",
                        "1ère VL": _first_vl_do.strftime("%d/%m/%Y"),
                    })
                except Exception:
                    continue
            if _detail_old_rows:
                try:
                    st.dataframe(pd.DataFrame(_detail_old_rows), hide_index=True, use_container_width=True)
                except Exception:
                    pass

            st.markdown("#### Détail par fonds — Portefeuille amélioré")
            _detail_new_rows = []
            for _isin_dn, _amt_dn in amounts.items():
                if _isin_dn.upper() in ("EUROFUND", "STRUCTURED"):
                    continue
                if _amt_dn <= 0:
                    continue
                _name_dn = CONTRACT_FUND_NAMES.get(_isin_dn, _isin_dn)[:40]
                _cat_dn_row = funds_df.loc[funds_df["isin"].str.upper() == _isin_dn.upper()] if not funds_df.empty else pd.DataFrame()
                _cat_dn = str(_cat_dn_row["category"].values[0]) if not _cat_dn_row.empty else ""
                _is_bond_dn = _is_bond_category(_cat_dn)
                try:
                    _df_dn, _, _ = get_price_series(_isin_dn, None, float(_euro_rate_cmp))
                    if _df_dn.empty:
                        continue
                    _s_dn = _df_dn["Close"].astype(float)
                    _first_vl_dn = _s_dn.index.min()
                    _s_dn = _s_dn[(_s_dn.index >= opt_start) & (_s_dn.index <= opt_end)]
                    if len(_s_dn) < 60:
                        continue
                    _rets_dn = _s_dn.pct_change().dropna()
                    if len(_rets_dn) < 60:
                        continue
                    _ann_ret_dn = float((1 + _rets_dn.mean()) ** 252 - 1)
                    _ann_vol_dn = float(_rets_dn.std() * np.sqrt(252))
                    _rfr_dn = get_risk_free_rate()
                    _sharpe_dn = float((_ann_ret_dn - _rfr_dn) / _ann_vol_dn) if _ann_vol_dn > 0.001 else 0.0
                    _detail_new_rows.append({
                        "Fonds": _name_dn,
                        "Rdt/an": f"{_ann_ret_dn*100:+.1f}%",
                        "Sharpe": f"{_sharpe_dn:.2f}" if not _is_bond_dn else "—",
                        "Vol/an": f"{_ann_vol_dn*100:.1f}%",
                        "1ère VL": _first_vl_dn.strftime("%d/%m/%Y"),
                    })
                except Exception:
                    continue
            if _detail_new_rows:
                try:
                    st.dataframe(pd.DataFrame(_detail_new_rows), hide_index=True, use_container_width=True)
                except Exception:
                    pass

            st.caption(f"Ratios calculés sur la période du {fmt_date(opt_start)} au {fmt_date(opt_end)}")

            # Message de synthèse
            _improvements_c = []
            if _old_div_c and _new_div_c and _new_div_c.get("score") and _old_div_c.get("score"):
                _d_sc = float(_new_div_c["score"]) - float(_old_div_c["score"])
                if _d_sc > 5:
                    _improvements_c.append(f"diversification +{_d_sc:.0f} pts")
            if _old_sharpe_c and _new_sharpe_c and _new_sharpe_c > _old_sharpe_c + 0.1:
                _improvements_c.append(f"Sharpe {_old_sharpe_c:.2f} → {_new_sharpe_c:.2f}")
            if _old_div_c and _new_div_c and _new_div_c.get("avg_corr", 1) < _old_div_c.get("avg_corr", 0) - 0.05:
                _improvements_c.append(f"corrélation réduite ({_old_div_c['avg_corr']:.0%} → {_new_div_c['avg_corr']:.0%})")

            if _improvements_c:
                st.success(f"✅ Améliorations : {', '.join(_improvements_c)}.")
            elif _new_sharpe_c and _old_sharpe_c and _new_sharpe_c >= _old_sharpe_c:
                st.info("L'allocation optimisée maintient les mêmes niveaux de performance.")
            else:
                st.warning("⚠️ L'allocation optimisée n'améliore pas tous les indicateurs. Ajustez la sélection des fonds.")

    except Exception as e:
        st.error("Une erreur est survenue dans le builder.")
        st.exception(e)

    # ── Avertissements collectés pendant le build ──────────────
    _bw = st.session_state.get("PP_BUILD_WARNINGS", [])
    if _bw:
        with st.expander(f"⚠️ {len(_bw)} avertissements", expanded=False):
            for _w in _bw:
                st.caption(f"• {_w}")
        st.session_state["PP_BUILD_WARNINGS"] = []

    # ── Méthodologie et avertissements ──────────────────────────
    with st.expander("ℹ️ Méthodologie et avertissements", expanded=False):
        st.markdown("""
**Données de performance**
Les rendements affichés sont calculés à partir des valeurs liquidatives publiées (EODHD), qui intègrent les frais de gestion internes du fonds (TER). Il s'agit donc de la performance **nette des frais de gestion du fonds**, telle que perçue par l'investisseur.

**Frais du contrat d'assurance-vie**
Les frais de gestion prélevés par l'assureur (frais contrat, ex: 0.50%/an) s'appliquent en supplément de manière uniforme sur tous les supports UC. Ils n'affectent pas le classement relatif des fonds entre eux.

**Ratio de Sharpe**
Mesure le rendement par unité de risque. Plus le Sharpe est élevé, meilleur est le rapport rendement/risque. Calculé sur 3 ans avec le taux sans risque du Bund 10 ans.

**Score de diversification**
Basé sur la corrélation entre les supports UC (hors fonds euros). Un score élevé signifie que les fonds ont des comportements différents — en cas de baisse d'un marché, les autres supports ne suivent pas la même trajectoire.

**Avertissement**
Les performances passées ne préjugent pas des performances futures. Un fonds avec des frais élevés peut être le meilleur choix s'il surperforme sa catégorie. Les corrélations historiques ne garantissent pas les corrélations futures. Ces simulations sont fournies à titre indicatif et ne constituent pas un conseil en investissement.
        """)


# ------------------------------------------------------------
# Sélection fonds depuis le référentiel contrat
# ------------------------------------------------------------

_BOND_CATEGORY_PREFIXES = ("OBLIGATIONS",)
_BOND_CATEGORY_KEYWORDS = (
    "CONVERTIBLES",
    "SUBORDINATED BOND",
    "REVENUS FIXES",
)


def _is_bond_category(cat: str) -> bool:
    """
    Retourne True si la catégorie Morningstar est de nature obligataire.
    Couvre :
    - Toutes catégories commençant par "OBLIGATIONS" (ex : OBLIGATIONS EUR...)
    - Convertibles (CONVERTIBLES EUR, CONVERTIBLES EUROPE...)
    - Dettes subordonnées (EUR SUBORDINATED BOND)
    - Revenus fixes (REVENUS FIXES EUR — fonds à échéance)
    - Produits leveragés/inversés sur obligations
      (TRADING - LEVERAGED/INVERSE OBLIGATIONS)
    """
    c = str(cat).upper().strip()
    return (
        c.startswith(_BOND_CATEGORY_PREFIXES)
        or any(k in c for k in _BOND_CATEGORY_KEYWORDS)
        or "OBLIGATIONS" in c
    )


def _get_assureur(contract_label: str) -> str:
    cfg = CONTRACTS_REGISTRY.get(contract_label, {})
    return cfg.get("assureur", "—")


def _add_fund_from_contract(port_key: str, label: str):
    """
    Interface d'ajout de fonds structurée en deux blocs :
      Bloc A — Actifs défensifs (fonds euros + fonds obligataires)
      Bloc B — Autres UC (actions, diversifiés, etc.)
    Utilise le contrat spécifique au portefeuille (A ou B).
    """
    st.subheader(label)

    is_a = port_key == "A_lines"
    contract_label = st.session_state.get(
        "CONTRACT_LABEL_A" if is_a else "CONTRACT_LABEL_B", ""
    )
    contract_cfg = CONTRACTS_REGISTRY.get(contract_label, {})
    funds_df = st.session_state.get(
        "CONTRACT_FUNDS_DF_A" if is_a else "CONTRACT_FUNDS_DF_B", pd.DataFrame()
    )
    buy_date_central = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if is_a
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    if not contract_label or not contract_cfg:
        st.warning("Contrat non sélectionné — configurez-le dans les paramètres.")
        return

    assureur = contract_cfg.get("assureur", "—")
    st.caption(
        f"Contrat : **{contract_label}** ({assureur}) "
        f"— {len(funds_df)} UC disponibles "
        f"| Date d'achat : {pd.Timestamp(buy_date_central).strftime('%d/%m/%Y')}"
    )

    # ═══════════════════════════════════════════════════════════
    # BLOC A — Actifs défensifs
    # ═══════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 🛡️ Actifs défensifs")

    # ── Sous-section 1 : Fonds en euros ──────────────────────
    st.markdown("#### Fonds en euros — capital garanti à 98%")

    ef_options = list(contract_cfg.get("euro_funds", {}).keys())
    if ef_options:
        ef_selected = st.selectbox("Fonds en euros", ef_options, key=f"ef_sel_{port_key}")
        ef_filename = contract_cfg["euro_funds"][ef_selected]
        try:
            ef_history = load_euro_fund_history(contract_cfg["path"], ef_filename)
        except Exception:
            ef_history = pd.DataFrame()
        ef_avg = get_euro_fund_avg_rate(ef_history, years=5)
        ef_last = (
            float(ef_history.iloc[-1]["taux_net_publie_pct"])
            if not ef_history.empty else ef_avg
        )
        with st.container(border=True):
            hc1, hc2 = st.columns([3, 1])
            with hc1:
                st.markdown(f"**{ef_selected}**")
                st.caption(
                    f"Assureur : {assureur} | SRI 1 | Capital garanti à 98%"
                    f" | Taux net {int(ef_history['annee'].max()) if not ef_history.empty else '—'}"
                    f" : **{ef_last:.2f}%**"
                )
            with hc2:
                st.metric("Moy. 5 ans", f"{ef_avg:.2f}%")
        euro_amt = st.number_input(
            "Montant (€)",
            min_value=0.0, max_value=10_000_000.0,
            value=0.0, step=1000.0,
            key=f"euro_amt_{port_key}",
        )
        if st.button("➕ Ajouter fonds en euros", key=f"add_euro_{port_key}"):
            if euro_amt <= 0:
                st.warning("Montant invalide.")
            else:
                st.session_state[port_key].append({
                    "id":            str(uuid.uuid4()),
                    "name":          ef_selected,
                    "isin":          "EUROFUND",
                    "amount_gross":  float(euro_amt),
                    "buy_date":      pd.Timestamp(buy_date_central),
                    "buy_px":        "",
                    "note":          f"Contrat : {contract_label}",
                    "sym_used":      "EUROFUND",
                    "fee_total_pct": 0.0,
                })
                st.toast(f"✅ {ef_selected} ajouté.", icon="✅")
                st.rerun()

    # ── Sous-section 2 : Fonds obligataires ──────────────────
    st.markdown("#### Fonds obligataires")

    if funds_df.empty:
        st.info("Référentiel de fonds non chargé.")
    else:
        bond_df = funds_df[funds_df["category"].apply(_is_bond_category)].copy()

        if bond_df.empty:
            st.info("Aucun fonds obligataire dans le référentiel.")
        else:
            bf1, bf2, bf3 = st.columns([3, 2, 1])
            with bf1:
                bond_search = st.text_input(
                    "Rechercher (nom, société de gestion)",
                    value="", key=f"bond_search_{port_key}",
                    placeholder="Ex : AXA, Amundi...",
                )
            with bf2:
                bond_cats = ["Toutes"] + sorted(bond_df["category"].dropna().unique().tolist())
                bond_cat = st.selectbox("Catégorie", bond_cats, key=f"bond_cat_{port_key}")
            with bf3:
                bond_sri = st.selectbox("SRI", ["Tous"] + [str(i) for i in range(1, 8)], key=f"bond_sri_{port_key}")

            if bond_search.strip():
                q = bond_search.strip().upper()
                bond_df = bond_df[
                    bond_df["name"].str.upper().str.contains(q, na=False) |
                    bond_df["manager"].str.upper().str.contains(q, na=False)
                ]
            if bond_cat != "Toutes":
                bond_df = bond_df[bond_df["category"] == bond_cat]
            if bond_sri != "Tous":
                bond_df = bond_df[bond_df["sri"] == int(bond_sri)]

            if bond_df.empty:
                st.info("Aucun résultat.")
            else:
                # Téléchargement liste filtrée
                _bond_xlsx = BytesIO()
                bond_df[["isin", "name", "manager", "category", "sri", "fee_total_pct"]].rename(
                    columns={
                        "isin": "ISIN", "name": "Libellé", "manager": "Société",
                        "category": "Catégorie", "sri": "SRI", "fee_total_pct": "Frais totaux %",
                    }
                ).to_excel(_bond_xlsx, index=False, engine="openpyxl")
                st.download_button(
                    f"⬇️ Télécharger la liste filtrée ({len(bond_df)} fonds)",
                    data=_bond_xlsx.getvalue(),
                    file_name=f"obligataires_{contract_label.replace(' ', '_')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"bond_dl_{port_key}",
                )
                # Sélection directe par selectbox
                bond_options = [
                    f"{row['name']}  |  {row['manager']}  |  SRI {row['sri']}  |  {row['fee_total_pct']:.2f}%"
                    for _, row in bond_df.iterrows()
                ]
                bond_sel_idx = st.selectbox(
                    f"{len(bond_df)} fonds obligataires disponibles — sélectionner",
                    options=range(len(bond_options)),
                    format_func=lambda i: bond_options[i],
                    key=f"bond_sel_{port_key}",
                )
                bond_row = bond_df.iloc[bond_sel_idx]
                with st.container(border=True):
                    bc1, bc2 = st.columns([3, 1])
                    with bc1:
                        st.markdown(f"**{bond_row['name']}**")
                        st.caption(
                            f"{bond_row['manager']}  ·  {bond_row['category']}"
                            f"  ·  ISIN : `{bond_row['isin']}`"
                        )
                    with bc2:
                        st.metric("Frais totaux", f"{bond_row['fee_total_pct']:.2f}%")
                bb1, bb2, bb3 = st.columns([2, 2, 1])
                with bb1:
                    bond_amt = st.text_input(
                        "Montant (€)", value="", key=f"bond_amt_{port_key}",
                        placeholder="Ex : 10000",
                    )
                with bb2:
                    bond_px = st.text_input(
                        "Prix d'achat (optionnel)", value="", key=f"bond_px_{port_key}",
                    )
                with bb3:
                    if st.button("➕ Ajouter", key=f"bond_add_{port_key}"):
                        try:
                            amt = float(str(bond_amt).replace(" ", "").replace(",", "."))
                            assert amt > 0
                        except Exception:
                            st.warning("Montant invalide.")
                        else:
                            st.session_state[port_key].append({
                                "id":               str(uuid.uuid4()),
                                "name":             bond_row["name"],
                                "isin":             bond_row["isin"],
                                "amount_gross":     float(amt),
                                "buy_date":         pd.Timestamp(buy_date_central),
                                "buy_px":           float(str(bond_px).replace(",", ".")) if bond_px.strip() else "",
                                "note":             f"Contrat : {contract_label} | SRI {bond_row['sri']}",
                                "sym_used":         "",
                                "fee_uc_pct":       bond_row["fee_uc_pct"],
                                "fee_contract_pct": bond_row["fee_contract_pct"],
                                "fee_total_pct":    bond_row["fee_total_pct"],
                            })
                            st.toast(f"✅ {bond_row['name']} ajouté.", icon="✅")
                            st.rerun()

    # ═══════════════════════════════════════════════════════════
    # BLOC B — Autres UC
    # ═══════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 📈 Autres UC — Actions, diversifiés, etc.")

    if funds_df.empty:
        st.info("Référentiel de fonds non chargé.")
        return

    other_df = funds_df[~funds_df["category"].apply(_is_bond_category)].copy()

    if other_df.empty:
        st.info("Aucun fonds UC dans le référentiel.")
        return

    oc1, oc2, oc3 = st.columns([3, 2, 1])
    with oc1:
        other_search = st.text_input(
            "Rechercher (nom, ISIN)",
            value="", key=f"other_search_{port_key}",
            placeholder="Ex : Carmignac, FR0010135103...",
        )
    with oc2:
        other_cats = ["Toutes"] + sorted(other_df["category"].dropna().unique().tolist())
        other_cat = st.selectbox("Catégorie", other_cats, key=f"other_cat_{port_key}")
    with oc3:
        other_sri = st.selectbox("SRI", ["Tous"] + [str(i) for i in range(1, 8)], key=f"other_sri_{port_key}")

    if other_search.strip():
        q = other_search.strip().upper()
        other_df = other_df[
            other_df["isin"].str.upper().str.contains(q, na=False) |
            other_df["name"].str.upper().str.contains(q, na=False)
        ]
    if other_cat != "Toutes":
        other_df = other_df[other_df["category"] == other_cat]
    if other_sri != "Tous":
        other_df = other_df[other_df["sri"] == int(other_sri)]

    if other_df.empty:
        st.info("Aucun fonds ne correspond aux critères.")
        return

    # Téléchargement liste filtrée
    _other_xlsx = BytesIO()
    other_df[["isin", "name", "manager", "category", "sri", "fee_total_pct"]].rename(
        columns={
            "isin": "ISIN", "name": "Libellé", "manager": "Société",
            "category": "Catégorie", "sri": "SRI", "fee_total_pct": "Frais totaux %",
        }
    ).to_excel(_other_xlsx, index=False, engine="openpyxl")
    st.download_button(
        f"⬇️ Télécharger la liste filtrée ({len(other_df)} UC)",
        data=_other_xlsx.getvalue(),
        file_name=f"uc_{contract_label.replace(' ', '_')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key=f"other_dl_{port_key}",
    )
    # Sélection directe par selectbox
    other_options = [
        f"{row['name']}  |  {row['manager']}  |  SRI {row['sri']}  |  {row['fee_total_pct']:.2f}%"
        for _, row in other_df.iterrows()
    ]
    other_sel_idx = st.selectbox(
        f"{len(other_df)} UC disponibles — sélectionner",
        options=range(len(other_options)),
        format_func=lambda i: other_options[i],
        key=f"other_sel_{port_key}",
    )
    other_row = other_df.iloc[other_sel_idx]
    with st.container(border=True):
        oc1, oc2 = st.columns([3, 1])
        with oc1:
            st.markdown(f"**{other_row['name']}**")
            st.caption(
                f"{other_row['manager']}  ·  {other_row['category']}"
                f"  ·  ISIN : `{other_row['isin']}`"
            )
        with oc2:
            st.metric("Frais totaux", f"{other_row['fee_total_pct']:.2f}%")
    ob1, ob2, ob3 = st.columns([2, 2, 1])
    with ob1:
        other_amt = st.text_input(
            "Montant (€)", value="", key=f"other_amt_{port_key}",
            placeholder="Ex : 10000",
        )
    with ob2:
        other_px = st.text_input(
            "Prix d'achat (optionnel)", value="", key=f"other_px_{port_key}",
        )
    with ob3:
        if st.button("➕ Ajouter", key=f"other_add_{port_key}"):
            try:
                amt = float(str(other_amt).replace(" ", "").replace(",", "."))
                assert amt > 0
            except Exception:
                st.warning("Montant invalide.")
            else:
                st.session_state[port_key].append({
                    "id":               str(uuid.uuid4()),
                    "name":             other_row["name"],
                    "isin":             other_row["isin"],
                    "amount_gross":     float(amt),
                    "buy_date":         pd.Timestamp(buy_date_central),
                    "buy_px":           float(str(other_px).replace(",", ".")) if other_px.strip() else "",
                    "note":             f"Contrat : {contract_label} | SRI {other_row['sri']}",
                    "sym_used":         "",
                    "fee_uc_pct":       other_row["fee_uc_pct"],
                    "fee_contract_pct": other_row["fee_contract_pct"],
                    "fee_total_pct":    other_row["fee_total_pct"],
                })
                st.toast(f"✅ {other_row['name']} ajouté.", icon="✅")
                st.rerun()


def _build_onepager_pdf(report: Dict[str, Any]) -> bytes:
    """
    Génère un PDF one-pager 'Proposition d'Arbitrage' prêt à présenter.
    Inclut : état actuel vs cible, gain estimé, économie frais/fiscalité.
    """
    if not REPORTLAB_AVAILABLE:
        return b""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor, white, black
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    )
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from io import BytesIO
    NAVY   = HexColor("#1B2A4A")
    GOLD   = HexColor("#C9A84C")
    STEEL  = HexColor("#4A6FA5")
    LIGHT  = HexColor("#F5F7FA")
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=1.8*cm, rightMargin=1.8*cm,
        topMargin=1.5*cm, bottomMargin=1.5*cm,
    )
    styles = getSampleStyleSheet()
    story = []
    def _h(txt, color=NAVY, size=16):
        return Paragraph(
            f'<font color="#{color.hexval()[2:]}" size="{size}"><b>{txt}</b></font>',
            styles["Normal"]
        )
    def _p(txt, size=9):
        return Paragraph(f'<font size="{size}">{txt}</font>', styles["Normal"])
    def _kpi(label, value, delta=None):
        val_str = f'<font size="14" color="#1B2A4A"><b>{value}</b></font>'
        lbl_str = f'<font size="8" color="#666666">{label}</font>'
        d_str = f'<br/><font size="9" color="#2E86AB">{delta}</font>' if delta else ""
        return Paragraph(f'{lbl_str}<br/>{val_str}{d_str}', styles["Normal"])
    # ── En-tête ──────────────────────────────────────────────────────
    nom_client = report.get("nom_client", "Client") or "Client"
    nom_cabinet = report.get("nom_cabinet", "Cabinet") or "Cabinet"
    as_of = report.get("as_of", "")
    header_data = [[
        Paragraph(f'<font color="white" size="18"><b>PROPOSITION D\'ARBITRAGE</b></font>'
                  f'<br/><font color="white" size="9">{nom_cabinet}</font>',
                  styles["Normal"]),
        Paragraph(f'<font color="white" size="10"><b>{nom_client}</b></font>'
                  f'<br/><font color="white" size="8">Au {as_of}</font>',
                  styles["Normal"]),
    ]]
    header_tbl = Table(header_data, colWidths=[11*cm, 6*cm])
    header_tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,-1), NAVY),
        ("TEXTCOLOR",   (0,0), (-1,-1), white),
        ("ALIGN",       (1,0), (1,0),   "RIGHT"),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",  (0,0), (-1,-1), 12),
        ("BOTTOMPADDING",(0,0),(-1,-1), 12),
        ("LEFTPADDING", (0,0), (-1,-1), 12),
        ("ROUNDEDCORNERS", [4]),
    ]))
    story.append(header_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 1 : État actuel vs Cible ────────────────────────────────
    story.append(_h("① État actuel vs Portefeuille cible", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    synth_a = report.get("client_summary", {})
    synth_b = report.get("valority_summary", {})
    comp    = report.get("comparison", {})
    def _eur(v):
        try:
            return f"{float(v):,.0f} €".replace(",", " ")
        except Exception:
            return "—"
    def _pct(v):
        try:
            return f"{float(v):.2f} %"
        except Exception:
            return "—"
    compare_data = [
        ["", "Portefeuille actuel (Client)", "Portefeuille cible (Cabinet)"],
        ["Valeur actuelle",
         _eur(synth_a.get("val", 0)), _eur(synth_b.get("val", 0))],
        ["Versements nets",
         _eur(synth_a.get("net", 0)), _eur(synth_b.get("net", 0))],
        ["Gain net",
         _eur(synth_a.get("val", 0) - synth_a.get("net", 0)),
         _eur(synth_b.get("val", 0) - synth_b.get("net", 0))],
        ["XIRR (annualisé)",
         _pct(synth_a.get("irr_pct", 0)), _pct(synth_b.get("irr_pct", 0))],
        ["Performance totale",
         _pct(synth_a.get("perf_tot_pct", 0)), _pct(synth_b.get("perf_tot_pct", 0))],
    ]
    cmp_tbl = Table(compare_data, colWidths=[5*cm, 5.5*cm, 5.5*cm])
    cmp_tbl.setStyle(TableStyle([
        ("BACKGROUND",   (0,0), (-1,0),  NAVY),
        ("TEXTCOLOR",    (0,0), (-1,0),  white),
        ("BACKGROUND",   (0,1), (-1,-1), LIGHT),
        ("BACKGROUND",   (2,1), (2,-1),  HexColor("#EAF4FB")),
        ("FONTSIZE",     (0,0), (-1,-1), 8),
        ("FONTNAME",     (0,0), (-1,0),  "Helvetica-Bold"),
        ("ALIGN",        (1,0), (-1,-1), "CENTER"),
        ("ROWBACKGROUNDS",(0,1),(-1,-1), [LIGHT, white]),
        ("GRID",         (0,0), (-1,-1), 0.3, HexColor("#CCCCCC")),
        ("TOPPADDING",   (0,0), (-1,-1), 5),
        ("BOTTOMPADDING",(0,0),(-1,-1),  5),
    ]))
    story.append(cmp_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 2 : Gain de performance estimé ──────────────────────────
    story.append(_h("② Gain de performance net estimé", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    delta_val  = comp.get("delta_val", 0.0)
    delta_xirr = comp.get("delta_perf_pct", 0.0)
    kpi_data = [[
        _kpi("Gain de valeur estimé", _eur(delta_val),
             "Valeur Cabinet − Valeur Client"),
        _kpi("Gain XIRR annualisé", _pct(delta_xirr),
             "Performance supplémentaire annualisée"),
        _kpi("Horizon recommandé", "Long terme (≥8 ans)",
             "Pour bénéficier de l'abattement AV"),
    ]]
    kpi_tbl = Table(kpi_data, colWidths=[5.5*cm, 5.5*cm, 5.5*cm])
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0), (-1,-1), LIGHT),
        ("ALIGN",         (0,0), (-1,-1), "CENTER"),
        ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0), (-1,-1), 10),
        ("BOTTOMPADDING", (0,0), (-1,-1), 10),
        ("BOX",           (0,0), (-1,-1), 1, GOLD),
        ("INNERGRID",     (0,0), (-1,-1), 0.5, HexColor("#CCCCCC")),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 3 : Économie frais / fiscalité ──────────────────────────
    story.append(_h("③ Économie de frais & levier fiscal", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    fees = report.get("fees_analysis", {})
    val_a = float(synth_a.get("val", 0) or 0)
    val_b = float(synth_b.get("val", 0) or 0)
    fee_a_pct = float(report.get("fee_a_pct", 0.6))
    fee_b_pct = float(report.get("fee_b_pct", 0.6))
    econo_frais_an = val_a * max(0.0, fee_a_pct - fee_b_pct) / 100.0
    eco_data = [
        ["Levier", "Impact estimé", "Note"],
        ["Économie frais contrat/an",
         _eur(econo_frais_an),
         f"Si passage de {fee_a_pct:.2f}% → {fee_b_pct:.2f}%/an"],
        ["Abattement AV (célibataire, ≥8 ans)",
         "4 600 €/an d'IR évité",
         "Sur la quote-part de gains dans le rachat"],
        ["Abattement AV (couple, ≥8 ans)",
         "9 200 €/an d'IR évité",
         "Sur la quote-part de gains dans le rachat"],
        ["Transmission (Art. 990I)",
         "152 500 € exonérés/bénéficiaire",
         "Pour versements avant 70 ans"],
    ]
    eco_tbl = Table(eco_data, colWidths=[5.5*cm, 4.5*cm, 7*cm])
    eco_tbl.setStyle(TableStyle([
        ("BACKGROUND",   (0,0), (-1,0),  STEEL),
        ("TEXTCOLOR",    (0,0), (-1,0),  white),
        ("BACKGROUND",   (0,1), (-1,-1), LIGHT),
        ("ROWBACKGROUNDS",(0,1),(-1,-1), [LIGHT, white]),
        ("FONTSIZE",     (0,0), (-1,-1), 8),
        ("FONTNAME",     (0,0), (-1,0),  "Helvetica-Bold"),
        ("GRID",         (0,0), (-1,-1), 0.3, HexColor("#CCCCCC")),
        ("TOPPADDING",   (0,0), (-1,-1), 5),
        ("BOTTOMPADDING",(0,0),(-1,-1),  5),
    ]))
    story.append(eco_tbl)
    story.append(Spacer(1, 14))
    # ── Pied de page légal ───────────────────────────────────────────
    story.append(HRFlowable(width="100%", thickness=0.5, color=HexColor("#CCCCCC")))
    story.append(Spacer(1, 4))
    story.append(_p(
        "Document à titre indicatif — ne constitue pas un conseil en investissement. "
        "Performances passées ne préjugent pas des performances futures. "
        f"Généré par {nom_cabinet} • {as_of}" if nom_cabinet else f"Document confidentiel • {as_of}",
        size=7,
    ))
    doc.build(story)
    return buf.getvalue()


# ── Helpers d'interprétation des ratios (module level) ──────────────────────

def _interpret_sharpe(val: Optional[float]) -> Tuple[str, str, str]:
    """Retourne (emoji_couleur, verdict_court, phrase_client)."""
    if val is None:
        return "⚪", "Données insuffisantes", "Pas assez d'historique pour calculer ce ratio."
    if val >= 1.0:
        return "🟢", "Excellent", f"Pour chaque unité de risque prise, le portefeuille génère {val:.2f} unité de rendement excédentaire. C'est un très bon ratio."
    if val >= 0.5:
        return "🟢", "Bon", f"Le portefeuille offre un bon compromis rendement/risque (Sharpe {val:.2f})."
    if val >= 0.0:
        return "🟠", "Moyen", f"Le rendement compense à peine le risque pris (Sharpe {val:.2f}). Une optimisation est possible."
    return "🔴", "Négatif", f"Le portefeuille perd de l'argent par rapport au fonds euros après ajustement du risque (Sharpe {val:.2f})."


def _interpret_sharpe_short(val: Optional[float]) -> str:
    """Verdict court pour st.metric delta."""
    if val is None:
        return "—"
    if val >= 1.0:
        return "🟢 Excellent"
    if val >= 0.5:
        return "🟢 Bon"
    if val >= 0.0:
        return "🟠 Moyen"
    return "🔴 Négatif"


def _interpret_sortino_short(val: Optional[float]) -> str:
    if val is None:
        return "—"
    if val >= 1.5:
        return "🟢 Excellente protection"
    if val >= 0.8:
        return "🟢 Bonne protection"
    if val >= 0.0:
        return "🟠 Protection limitée"
    return "🔴 Vulnérable"


def _interpret_beta_short(val: Optional[float]) -> str:
    if val is None:
        return "—"
    if val < 0.5:
        return "🟢 Très défensif"
    if val < 0.8:
        return "🟢 Défensif"
    if val < 1.2:
        return "🟠 Neutre (marché)"
    return "🔴 Agressif"


def _interpret_sortino(val: Optional[float]) -> Tuple[str, str, str]:
    if val is None:
        return "⚪", "—", ""
    if val >= 1.5:
        return "🟢", "Excellente protection", f"Le Sortino de {val:.2f} montre que le portefeuille limite très bien les baisses."
    if val >= 0.8:
        return "🟢", "Bonne protection", f"Sortino de {val:.2f} : les pertes sont contenues par rapport aux gains."
    if val >= 0.0:
        return "🟠", "Protection limitée", f"Sortino de {val:.2f} : le portefeuille est vulnérable en cas de baisse."
    return "🔴", "Vulnérable", f"Sortino négatif ({val:.2f}) : le portefeuille souffre particulièrement dans les phases baissières."


def _interpret_beta(val: Optional[float]) -> Tuple[str, str, str]:
    if val is None:
        return "⚪", "—", ""
    if val < 0.5:
        return "🟢", "Très défensif", f"Bêta de {val:.2f} : le portefeuille ne bouge que de {val*100:.0f}% quand le marché bouge de 100%. Très protecteur."
    if val < 0.8:
        return "🟢", "Défensif", f"Bêta de {val:.2f} : le portefeuille amortit les mouvements du marché."
    if val < 1.2:
        return "🟠", "Neutre", f"Bêta de {val:.2f} : le portefeuille suit globalement le marché."
    return "🔴", "Agressif", f"Bêta de {val:.2f} : le portefeuille amplifie les mouvements du marché. Plus de potentiel mais plus de risque."


def _render_ratios_card(
    sharpe: Optional[float],
    sortino: Optional[float],
    beta_alpha: Optional[Dict[str, Any]],
    euro_rate: float,
) -> None:
    """Affiche les 3 ratios dans un container bordé avec interprétation colorée."""
    with st.container(border=True):
        # ── Sharpe ──
        _sh_emoji, _sh_verdict, _sh_phrase = _interpret_sharpe(sharpe)
        st.metric(
            f"{_sh_emoji} Ratio de Sharpe",
            f"{sharpe:.2f}" if sharpe is not None else "—",
            delta=_sh_verdict,
        )
        st.caption(_sh_phrase)

        st.markdown("---")

        # ── Sortino ──
        _so_emoji, _so_verdict, _so_phrase = _interpret_sortino(sortino)
        st.metric(
            f"{_so_emoji} Ratio de Sortino",
            f"{sortino:.2f}" if sortino is not None else "—",
            delta=_so_verdict,
        )
        if _so_phrase:
            st.caption(_so_phrase)

        st.markdown("---")

        # ── Bêta ──
        if beta_alpha is not None:
            _b = beta_alpha["beta"]
            _a = beta_alpha["alpha_pct"]
            _bn = beta_alpha["benchmark_name"]
            _be_emoji, _be_verdict, _be_phrase = _interpret_beta(_b)
            st.metric(
                f"{_be_emoji} Bêta (vs {_bn})",
                f"{_b:.2f}",
                delta=_be_verdict,
            )
            st.caption(_be_phrase)
            if _a is not None:
                _alpha_color = "🟢" if _a > 0 else "🔴"
                st.caption(
                    f"{_alpha_color} Alpha de Jensen : {_a:+.2f}%/an — "
                    + ("le portefeuille crée de la valeur au-delà du marché."
                       if _a > 0
                       else "le portefeuille sous-performe par rapport à ce que son niveau de risque devrait générer.")
                )
        else:
            st.metric("⚪ Bêta", "—")
            st.caption("Indice de référence indisponible ou historique insuffisant.")

        # ── Note méthodologique ──
        st.caption(
            f"_Taux sans risque utilisé : {euro_rate:.2f}% (taux fonds euros du contrat). "
            f"Calculs sur l'historique complet des VL disponibles._"
        )




# ─────────────────────────────────────────────────────────────────────────
# Fonctions de génération de rapports (globales pour accès depuis tous les modules)
# ─────────────────────────────────────────────────────────────────────────

def _chart_to_pptx_image(
    df: pd.DataFrame,
    col_y: str,
    label: str,
    color: str = "#1B2A4A",
) -> Optional[bytes]:
    """
    Génère un graphique matplotlib en mémoire (PNG bytes)
    pour insertion dans le PPTX.
    df doit avoir un DatetimeIndex ou colonne Date.
    """
    if not MATPLOTLIB_AVAILABLE or df is None or df.empty:
        return None
    try:
        import io as _io
        _fig, _ax = plt.subplots(figsize=(7, 2.8))
        _ax.fill_between(
            df.index if hasattr(df.index, "dtype") and "datetime" in str(df.index.dtype)
            else range(len(df)),
            df[col_y].ffill(),
            alpha=0.18,
            color=color,
        )
        _ax.plot(
            df.index if hasattr(df.index, "dtype") and "datetime" in str(df.index.dtype)
            else range(len(df)),
            df[col_y].ffill(),
            color=color,
            linewidth=1.8,
        )
        _ax.set_ylabel("€", fontsize=8)
        _ax.set_title(label, fontsize=9, color="#1B2A4A", fontweight="bold")
        _ax.spines["top"].set_visible(False)
        _ax.spines["right"].set_visible(False)
        _ax.yaxis.set_major_formatter(
            plt.FuncFormatter(lambda x, _: f"{x:,.0f}")
        )
        _fig.tight_layout()
        _buf = _io.BytesIO()
        _fig.savefig(_buf, format="png", dpi=130, bbox_inches="tight")
        plt.close(_fig)
        _buf.seek(0)
        return _buf.read()
    except Exception:
        return None


def _add_table_to_story(
    story: List[Any],
    df: pd.DataFrame,
    col_widths: Optional[List[float]] = None,
    font_size: int = 9,
):
    if df.empty:
        story.append(Paragraph("Données indisponibles.", getSampleStyleSheet()["Normal"]))
        return
    headers = list(df.columns)
    fmt_rows = []
    for _, row in df.iterrows():
        formatted = []
        for col, val in row.items():
            if "€" in col:
                formatted.append(val if isinstance(val, str) else fmt_eur_fr(val))
            elif "%" in col:
                formatted.append(val if isinstance(val, str) else fmt_pct_fr(val))
            else:
                formatted.append(str(val))
        fmt_rows.append(formatted)
    data = [headers] + fmt_rows
    table = Table(data, repeatRows=1, colWidths=col_widths)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), font_size),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
    ]
    for i in range(1, len(data)):
        if i % 2 == 0:
            style.append(("BACKGROUND", (0, i), (-1, i), colors.whitesmoke))
    table.setStyle(TableStyle(style))
    story.append(table)


def _fig_to_rl_image(fig: plt.Figure, width: float = 480, height: float = 270) -> Image:
    if not MATPLOTLIB_AVAILABLE:
        return None
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return Image(buf, width=width, height=height)


def _build_value_chart(df_map: Dict[str, pd.DataFrame]) -> Optional[Image]:
    if not MATPLOTLIB_AVAILABLE:
        return None
    if not df_map:
        return None
    plt.rcParams["font.family"] = "DejaVu Sans"
    _CHART_COLORS = ["#1F3B6D", "#C8963E"]
    fig, ax = plt.subplots(figsize=(6, 3))
    fig.patch.set_facecolor("white")
    has_data = False
    for idx, (label, df) in enumerate(df_map.items()):
        if df is None or df.empty or "Valeur" not in df.columns:
            continue
        ax.plot(
            df.index, df["Valeur"],
            label=label,
            color=_CHART_COLORS[idx % len(_CHART_COLORS)],
            linewidth=1.8,
        )
        has_data = True
    if not has_data:
        plt.close(fig)
        return None
    from matplotlib.ticker import FuncFormatter
    ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: f"{x:,.0f}".replace(",", " ")))
    ax.set_title("Évolution de la valeur du portefeuille")
    ax.set_xlabel("Date")
    ax.set_ylabel("Valeur (€)")
    ax.legend(loc="best")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", color="#E2E8F0", linewidth=0.5)
    fig.autofmt_xdate()
    return _fig_to_rl_image(fig)


def _wrap_label(label: str, width: int = 28) -> str:
    if not label:
        return "—"
    return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)


def _allocation_from_positions(df_positions: pd.DataFrame) -> Tuple[pd.DataFrame, str]:
    df = df_positions.copy()
    df = df[df["Valeur actuelle €"] >= 0]
    df["Nom"] = df["Nom"].fillna("—")
    df["ISIN / Code"] = df["ISIN / Code"].fillna("—")

    total_value = df["Valeur actuelle €"].sum()
    if total_value > 0:
        df["Poids"] = df["Valeur actuelle €"] / total_value
        basis_label = "Valeur actuelle"
    else:
        total_net = df["Net investi €"].sum()
        if total_net > 0:
            df["Poids"] = df["Net investi €"] / total_net
        else:
            df["Poids"] = 0.0
            if len(df) > 0:
                df.loc[df.index[0], "Poids"] = 1.0
        basis_label = "Net investi"

    df = df.sort_values("Poids", ascending=False)
    if len(df) > 8:
        df_main = df.iloc[:8].copy()
        df_other = df.iloc[8:]
        other_row = pd.DataFrame(
            {
                "Nom": ["Autres"],
                "ISIN / Code": ["—"],
                "Net investi €": [df_other["Net investi €"].sum()],
                "Valeur actuelle €": [df_other["Valeur actuelle €"].sum()],
                "Poids": [df_other["Poids"].sum()],
            }
        )
        df = pd.concat([df_main, other_row], ignore_index=True)

    df["Part %"] = df["Poids"] * 100.0
    return df, basis_label


def _build_allocation_donut(
    df_alloc: pd.DataFrame,
    title: str,
    figsize: Tuple[float, float] = (6.0, 3.4),
) -> Optional[Image]:
    if not MATPLOTLIB_AVAILABLE or df_alloc.empty:
        return None
    plt.rcParams["font.family"] = "DejaVu Sans"
    _DONUT_COLORS = ["#1F3B6D", "#C8963E", "#1A7A4A", "#F1F4F9", "#2E5FA3",
                     "#4B5563", "#6B7280", "#9CA3AF"]
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("white")
    wedges, _ = ax.pie(
        df_alloc["Poids"],
        startangle=90,
        labels=None,
        colors=_DONUT_COLORS[:len(df_alloc)],
        wedgeprops=dict(width=0.35, edgecolor="white"),
    )
    labels = [
        f"{_wrap_label(nm)} ({pct:.1f}%)"
        for nm, pct in zip(df_alloc["Nom"], df_alloc["Part %"])
    ]
    ax.legend(
        wedges,
        labels,
        loc="center left",
        bbox_to_anchor=(1.02, 0.5),
        frameon=False,
        fontsize=8,
    )
    ax.set_title(title)
    ax.set_aspect("equal")
    fig.tight_layout(rect=[0, 0, 0.78, 1])
    return _fig_to_rl_image(fig, width=460, height=280)


def _build_envelope_breakdown(
    lines: List[Dict[str, Any]],
    title: str,
) -> Tuple[Optional[Image], Optional[str]]:
    if not lines:
        return None, "Répartition par enveloppe : —"
    categories = {"Fonds euros": 0.0, "UC": 0.0, "Structurés": 0.0}
    for ln in lines:
        isin = str(ln.get("isin", "")).upper()
        val = float(ln.get("value", 0.0))
        if isin == "EUROFUND":
            categories["Fonds euros"] += val
        elif isin == "STRUCTURED":
            categories["Structurés"] += val
        else:
            categories["UC"] += val
    total = sum(categories.values())
    if total <= 0:
        return None, "Répartition par enveloppe : —"

    shares = {k: v / total for k, v in categories.items() if v > 0}
    major = max(shares.items(), key=lambda x: x[1])
    if sum(1 for v in shares.values() if v >= 0.01) < 2:
        return None, f"Répartition par enveloppe : {major[1] * 100:.1f}% {major[0]}"

    if not MATPLOTLIB_AVAILABLE:
        return None, None
    labels = list(shares.keys())
    values = [shares[k] * 100 for k in labels]
    fig, ax = plt.subplots(figsize=(6.0, 1.8))
    ax.barh(labels, values, color="#4C78A8")
    ax.set_xlim(0, 100)
    ax.set_xlabel("%")
    ax.set_title(title)
    for i, v in enumerate(values):
        ax.text(min(v + 1, 98), i, f"{v:.1f}%", va="center", fontsize=8)
    fig.tight_layout()
    return _fig_to_rl_image(fig, width=460, height=140), None


def _build_contribution_bar(df_positions: pd.DataFrame) -> Optional[Image]:
    if not MATPLOTLIB_AVAILABLE or df_positions.empty:
        return None
    df = df_positions.copy()
    if not {"Nom", "Valeur actuelle €", "Net investi €"}.issubset(df.columns):
        return None
    plt.rcParams["font.family"] = "DejaVu Sans"
    df["Contribution €"] = df["Valeur actuelle €"] - df["Net investi €"]
    df = df.sort_values("Contribution €", ascending=False)
    bar_colors = ["#1A7A4A" if v >= 0 else "#CC2200" for v in df["Contribution €"]]
    fig_height = max(2.0, min(4.2, 0.35 * len(df) + 1.2))
    fig, ax = plt.subplots(figsize=(6.2, fig_height))
    fig.patch.set_facecolor("white")
    ax.barh(df["Nom"], df["Contribution €"], color=bar_colors)
    ax.invert_yaxis()
    ax.set_title("Contribution à la performance (€)")
    ax.axvline(0, color="black", linewidth=0.5)
    ax.tick_params(axis="y", labelsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i, v in enumerate(df["Contribution €"]):
        offset = 0.01 * abs(v) if v != 0 else 0.5
        x_pos = v + offset if v >= 0 else v - offset
        ax.text(x_pos, i, fmt_eur_fr(v), va="center", fontsize=8)
    fig.tight_layout()
    return _fig_to_rl_image(fig, width=460, height=200)


def _build_single_line_bar(label: str, value: float, title: str) -> Optional[Image]:
    if not MATPLOTLIB_AVAILABLE:
        return None
    fig, ax = plt.subplots(figsize=(6.0, 1.3))
    ax.barh([label], [100], color="#4C78A8")
    ax.set_xlim(0, 100)
    ax.set_title(title)
    ax.set_xticks([])
    ax.tick_params(axis="y", labelsize=8)
    fig.tight_layout()
    return _fig_to_rl_image(fig, width=460, height=90)


def _years_between(d0: pd.Timestamp, d1: pd.Timestamp) -> float:
    return max(0.0, (d1 - d0).days / 365.25)


def generate_pdf_report(report: Dict[str, Any]) -> bytes:
    if not REPORTLAB_AVAILABLE:
        raise RuntimeError(f"PDF indisponible: {REPORTLAB_ERROR}")
    if not MATPLOTLIB_AVAILABLE:
        raise RuntimeError(f"PDF indisponible: {MATPLOTLIB_ERROR}")

    class NumberedCanvas(canvas.Canvas):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, **kwargs)
            self._saved_page_states = []

        def showPage(self):
            self._saved_page_states.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            page_count = len(self._saved_page_states)
            for state in self._saved_page_states:
                self.__dict__.update(state)
                self._draw_header_footer(page_count)
                super().showPage()
            super().save()

        def _draw_header_footer(self, page_count: int):
            width, height = A4
            self.setFillColor(colors.HexColor("#1F3B6D"))
            self.setFont("Helvetica-Bold", 10)
            header_text = report.get("nom_cabinet", "") or "Rapport de portefeuille"
            self.drawString(36, height - 30, header_text)
            self.setFillColor(colors.grey)
            self.setFont("Helvetica", 8)
            as_of_str = report.get("as_of", "")
            nom_cli_str = report.get("nom_client", "").strip()
            right_str = f"{nom_cli_str}  ·  {as_of_str}" if nom_cli_str else as_of_str
            self.drawRightString(width - 36, height - 30, right_str)
            self.setFillColor(colors.grey)
            self.setFont("Helvetica", 7)
            self.drawString(
                36,
                24,
                "Document informatif, ne constitue pas un conseil en investissement.",
            )
            self.drawRightString(width - 36, 24, f"Page {self.getPageNumber()} / {page_count}")

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=36,
        rightMargin=36,
        topMargin=54,
        bottomMargin=48,
    )
    base_styles = getSampleStyleSheet()
    _NAVY = "#1F3B6D"
    _GREEN = "#1A7A4A"
    _RED = "#CC2200"
    _GREY_TXT = "#647480"
    styles = {
        "title": ParagraphStyle(
            "Title",
            parent=base_styles["Title"],
            fontName="Helvetica-Bold",
            textColor=colors.HexColor(_NAVY),
            fontSize=20,
            spaceAfter=12,
        ),
        "h1": ParagraphStyle(
            "H1",
            parent=base_styles["Heading1"],
            fontName="Helvetica-Bold",
            textColor=colors.HexColor(_NAVY),
            fontSize=16,
            spaceAfter=10,
        ),
        "h2": ParagraphStyle(
            "H2",
            parent=base_styles["Heading2"],
            fontName="Helvetica-Bold",
            textColor=colors.HexColor(_NAVY),
            fontSize=12,
            spaceAfter=6,
        ),
        "small": ParagraphStyle(
            "Small",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=7,
            textColor=colors.HexColor(_GREY_TXT),
        ),
        "kpi": ParagraphStyle(
            "KPI",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=9,
            textColor=colors.HexColor("#111827"),
        ),
        "verdict_pos": ParagraphStyle(
            "VerdictPos",
            parent=base_styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=10,
            textColor=colors.HexColor(_GREEN),
            spaceBefore=6,
            spaceAfter=6,
        ),
        "verdict_neg": ParagraphStyle(
            "VerdictNeg",
            parent=base_styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=10,
            textColor=colors.HexColor(_RED),
            spaceBefore=6,
            spaceAfter=6,
        ),
        "interp_pos": ParagraphStyle(
            "InterpPos",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=8,
            textColor=colors.HexColor(_GREEN),
        ),
        "interp_neg": ParagraphStyle(
            "InterpNeg",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=8,
            textColor=colors.HexColor(_RED),
        ),
        "interp_neutral": ParagraphStyle(
            "InterpNeutral",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=8,
            textColor=colors.HexColor(_GREY_TXT),
        ),
        "fee_high": ParagraphStyle(
            "FeeHigh",
            parent=base_styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=9,
            textColor=colors.HexColor(_RED),
        ),
        "av_green": ParagraphStyle(
            "AVGreen",
            parent=base_styles["Normal"],
            fontName="Helvetica",
            fontSize=9,
            textColor=colors.HexColor(_GREEN),
        ),
    }
    story: List[Any] = []

    def _fe_pdf(x) -> str:
        try:
            return f"{float(x):,.0f} €".replace(",", "\u202f")
        except Exception:
            return "—"

    def _tbl_style(hdr_color: str = "#1F3B6D") -> TableStyle:
        return TableStyle([
            ("BACKGROUND",   (0, 0), (-1, 0), colors.HexColor(hdr_color)),
            ("TEXTCOLOR",    (0, 0), (-1, 0), colors.white),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.white]),
            ("BOX",          (0, 0), (-1, -1), 0.5, colors.lightgrey),
            ("INNERGRID",    (0, 0), (-1, -1), 0.25, colors.lightgrey),
            ("VALIGN",       (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING",  (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING",   (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 4),
        ])

    mode    = report.get("mode", "compare")
    synthA  = report.get("client_summary", {})
    synthB  = report.get("valority_summary", {})

    # ══════════════════════════════════════════════════════════
    # PAGE 1 — Synthèse exécutive
    # ══════════════════════════════════════════════════════════
    _nom_cli = report.get("nom_client", "").strip()
    _nom_cab = report.get("nom_cabinet", "").strip()
    _as_of   = report.get("as_of", "")
    _title_p1 = f"Rapport de gestion — {_nom_cli}" if _nom_cli else "Rapport de gestion"
    story.append(Paragraph(_title_p1, styles["title"]))
    _sub_p1 = "  ·  ".join(filter(None, [
        f"Cabinet : {_nom_cab}" if _nom_cab else "",
        f"Date : {_as_of}" if _as_of else "",
    ]))
    if _sub_p1:
        story.append(Paragraph(_sub_p1, styles["small"]))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#1F3B6D"), spaceAfter=8))

    _kpi_labels_p1 = [
        "Valeur actuelle", "Net investi", "Brut versé",
        "Performance totale", "XIRR annualisé",
    ]

    def _synth_vals_p1(s):
        return [
            fmt_eur_fr(s.get("val", 0)),
            fmt_eur_fr(s.get("net", 0)),
            fmt_eur_fr(s.get("brut", 0)),
            fmt_pct_fr(s.get("perf_tot_pct", 0)),
            fmt_pct_fr(s.get("irr_pct", 0)),
        ]

    if mode == "compare":
        _hdr_p1 = [
            Paragraph("<b>Indicateur</b>",       styles["h2"]),
            Paragraph("<b>Votre portefeuille</b>", styles["h2"]),
            Paragraph("<b>Notre proposition</b>", styles["h2"]),
        ]
        _kpi_data_p1 = [_hdr_p1]
        for _lbl, _vA, _vB in zip(_kpi_labels_p1, _synth_vals_p1(synthA), _synth_vals_p1(synthB)):
            _kpi_data_p1.append([
                Paragraph(_lbl, styles["small"]),
                Paragraph(_vA, styles["kpi"]),
                Paragraph(_vB, styles["kpi"]),
            ])
        _kpi_tbl_p1 = Table(_kpi_data_p1, colWidths=[160, 150, 150])
        _kpi_tbl_p1.setStyle(_tbl_style())
        story.append(_kpi_tbl_p1)
        story.append(Spacer(1, 10))
        # Verdict delta
        _val_A_p1   = float(synthA.get("val") or 0)
        _val_B_p1   = float(synthB.get("val") or 0)
        _perf_A_p1  = float(synthA.get("perf_tot_pct") or 0)
        _perf_B_p1  = float(synthB.get("perf_tot_pct") or 0)
        _delta_p1   = _val_B_p1 - _val_A_p1
        _dperf_p1   = _perf_B_p1 - _perf_A_p1
        if _delta_p1 > 0:
            _verdict = (
                f"Notre proposition génère <b>{_fe_pdf(_delta_p1)} de plus-value supplémentaire</b> "
                f"par rapport au portefeuille actuel ({_dperf_p1:+.1f} pt de performance)."
            )
            story.append(Paragraph(_verdict, styles["verdict_pos"]))
        elif _delta_p1 < 0:
            _verdict = (
                f"Le portefeuille actuel affiche une valeur supérieure de {_fe_pdf(-_delta_p1)} "
                f"à notre proposition à ce jour."
            )
            story.append(Paragraph(_verdict, styles["verdict_neg"]))
        else:
            _verdict = "Les deux allocations présentent une valeur équivalente à ce jour."
            story.append(Paragraph(_verdict, styles["kpi"]))
    else:
        _synth_main = synthB if mode == "valority" else synthA
        _lbl_main   = "Notre proposition (Cabinet)" if mode == "valority" else "Votre portefeuille (Client)"
        _hdr_p1_s   = [
            Paragraph("<b>Indicateur</b>",       styles["h2"]),
            Paragraph(f"<b>{_lbl_main}</b>",     styles["h2"]),
        ]
        _kpi_data_p1_s = [_hdr_p1_s]
        for _lbl, _val in zip(_kpi_labels_p1, _synth_vals_p1(_synth_main)):
            _kpi_data_p1_s.append([Paragraph(_lbl, styles["small"]), Paragraph(_val, styles["kpi"])])
        _kpi_tbl_p1_s = Table(_kpi_data_p1_s, colWidths=[200, 260])
        _kpi_tbl_p1_s.setStyle(_tbl_style())
        story.append(_kpi_tbl_p1_s)

    story.append(Spacer(1, 15))
    story.append(Paragraph(
        "Performance nette des frais de gestion. "
        "Les frais internes des fonds (TER) sont reflétés dans les valeurs liquidatives publiées.",
        styles["small"],
    ))

    # ══════════════════════════════════════════════════════════
    # PAGE 2 — Évolution de la valeur
    # ══════════════════════════════════════════════════════════
    story.append(PageBreak())
    story.append(Paragraph("Évolution de la valeur", styles["h1"]))
    _vc = _build_value_chart(report.get("df_map", {}))
    if _vc is not None:
        story.append(_vc)
    else:
        story.append(Paragraph(
            "Données de valorisation insuffisantes pour tracer le graphique.", styles["small"]
        ))

    # ══════════════════════════════════════════════════════════
    # PAGE 3 — Diversification & risque
    # ══════════════════════════════════════════════════════════
    story.append(PageBreak())
    story.append(Paragraph("Diversification & risque", styles["h1"]))

    div_A    = report.get("diversification_A") or {}
    div_B    = report.get("diversification_B") or {}
    risk_A   = report.get("risk_A") or {}
    risk_B   = report.get("risk_B") or {}
    ratios_A = report.get("ratios_A") or {}
    ratios_B = report.get("ratios_B") or {}

    def _div_lbl(score):
        if score >= 70:   return "Bonne"
        if score >= 40:   return "Moyenne"
        return "Fausse diversification"

    def _build_div_rows_p3(div, risk):
        rows = []
        if div:
            rows.append(("Score",                f"{div.get('score', 0):.0f}/100 — {_div_lbl(div.get('score', 0))}"))
            rows.append(("Corrélation moyenne",  f"{div.get('avg_corr', 0):.1%}"))
            rows.append(("Lignes analysées",     str(div.get("n_lines", "—"))))
            rows.append(("Positions indép.",     str(div.get("n_effective", "—"))))
            _dbl = div.get("doublons", [])
            rows.append(("Recoupements",         str(len(_dbl)) if _dbl else "Aucun"))
        if risk:
            rows.append(("Volatilité annualisée", f"{risk.get('vol_ann_pct', 0):.1f}%"))
            rows.append(("Drawdown max.",         f"{risk.get('max_dd_pct', 0):.1f}%"))
        return rows

    if div_A or div_B or risk_A or risk_B:
        story.append(Paragraph("Analyse de la diversification et du risque", styles["h2"]))
        if mode == "compare":
            _rows_A = _build_div_rows_p3(div_A, risk_A)
            _rows_B = _build_div_rows_p3(div_B, risk_B)
            _n = max(len(_rows_A), len(_rows_B))
            _rows_A += [("", "")] * (_n - len(_rows_A))
            _rows_B += [("", "")] * (_n - len(_rows_B))
            _div_data = [[
                Paragraph("<b>Indicateur</b>",        styles["h2"]),
                Paragraph("<b>Votre portefeuille</b>", styles["h2"]),
                Paragraph("<b>Notre proposition</b>",  styles["h2"]),
            ]]
            for (_la, _va), (_lb, _vb) in zip(_rows_A, _rows_B):
                _div_data.append([
                    Paragraph(_la or _lb, styles["small"]),
                    Paragraph(_va, styles["kpi"]),
                    Paragraph(_vb, styles["kpi"]),
                ])
            _div_tbl = Table(_div_data, colWidths=[160, 150, 150])
            _div_tbl.setStyle(_tbl_style("#4B5563"))
            story.append(_div_tbl)
        else:
            _d_s = div_B if mode == "valority" else div_A
            _r_s = risk_B if mode == "valority" else risk_A
            _lbl_s = "Notre proposition" if mode == "valority" else "Votre portefeuille"
            _rows_s = _build_div_rows_p3(_d_s, _r_s)
            if _rows_s:
                _div_data_s = [[
                    Paragraph("<b>Indicateur</b>",    styles["h2"]),
                    Paragraph(f"<b>{_lbl_s}</b>",     styles["h2"]),
                ]]
                for _lbl, _val in _rows_s:
                    _div_data_s.append([Paragraph(_lbl, styles["small"]), Paragraph(_val, styles["kpi"])])
                _div_tbl_s = Table(_div_data_s, colWidths=[200, 260])
                _div_tbl_s.setStyle(_tbl_style("#4B5563"))
                story.append(_div_tbl_s)
        story.append(Spacer(1, 15))

    # Ratios avec colonne interprétation
    _has_ratios_p3 = bool(ratios_A or ratios_B)
    if _has_ratios_p3:
        story.append(Paragraph("Indicateurs techniques de risque", styles["h2"]))
        _ba_A = ratios_A.get("beta_alpha") or {}
        _ba_B = ratios_B.get("beta_alpha") or {}

        def _fmt_r(v):
            try:
                return f"{float(v):.2f}" if v not in (None, "", 0.0) else "—"
            except Exception:
                return "—"

        def _interp(name, va, vb):
            """Returns (text, sentiment) where sentiment is 'pos', 'neg', or 'neutral'."""
            try:
                fa, fb = float(va or 0), float(vb or 0)
                if name == "Sharpe":
                    if fb > fa > 0: return f"↑ +{fb-fa:.2f} — meilleur rendement/risque", "pos"
                    if fa > fb > 0: return f"↓ {fb-fa:.2f} — moins favorable", "neg"
                elif name == "Sortino":
                    if fb > fa > 0: return "↑ meilleure résistance aux baisses", "pos"
                    if fa > fb > 0: return "↓ davantage exposé aux baisses", "neg"
                elif name == "Bêta":
                    if 0 < fb < fa: return f"↓ moins sensible aux marchés ({fb:.2f})", "pos"
                    if fb > fa:     return f"↑ plus amplifié ({fb:.2f})", "neg"
                elif name == "Alpha (%/an)":
                    if fb > 0: return f"Surperformance {fb:.2f}%/an vs indice", "pos"
                    if fb < 0: return f"Sous-performance {fb:.2f}%/an vs indice", "neg"
            except Exception:
                pass
            return "", "neutral"

        def _interp_para(name, va, vb):
            txt, sentiment = _interp(name, va, vb)
            _sty = styles["interp_pos"] if sentiment == "pos" else (
                styles["interp_neg"] if sentiment == "neg" else styles["interp_neutral"]
            )
            return Paragraph(txt, _sty)

        _ratio_defs = [
            ("Sharpe",      ratios_A.get("sharpe"),      ratios_B.get("sharpe")),
            ("Sortino",     ratios_A.get("sortino"),     ratios_B.get("sortino")),
            ("Bêta",        _ba_A.get("beta"),           _ba_B.get("beta")),
            ("Alpha (%/an)", _ba_A.get("alpha_pct"),    _ba_B.get("alpha_pct")),
        ]
        if mode == "compare":
            _r_data = [[
                Paragraph("<b>Ratio</b>",         styles["h2"]),
                Paragraph("<b>Votre portef.</b>", styles["h2"]),
                Paragraph("<b>Notre propos.</b>", styles["h2"]),
                Paragraph("<b>Lecture</b>",       styles["h2"]),
            ]]
            for _rn, _rva, _rvb in _ratio_defs:
                _r_data.append([
                    Paragraph(_rn, styles["small"]),
                    Paragraph(_fmt_r(_rva), styles["kpi"]),
                    Paragraph(_fmt_r(_rvb), styles["kpi"]),
                    _interp_para(_rn, _rva, _rvb),
                ])
            _r_tbl = Table(_r_data, colWidths=[80, 90, 90, 200])
            _r_tbl.setStyle(_tbl_style("#4B5563"))
            story.append(_r_tbl)
        else:
            _rat_s = ratios_B if mode == "valority" else ratios_A
            _ba_s  = (_rat_s.get("beta_alpha") or {})
            _ratio_defs_s = [
                ("Sharpe",       _rat_s.get("sharpe")),
                ("Sortino",      _rat_s.get("sortino")),
                ("Bêta",         _ba_s.get("beta")),
                ("Alpha (%/an)", _ba_s.get("alpha_pct")),
            ]
            _r_data_s = [[
                Paragraph("<b>Ratio</b>", styles["h2"]),
                Paragraph("<b>Valeur</b>", styles["h2"]),
            ]]
            for _rn, _rv in _ratio_defs_s:
                _r_data_s.append([Paragraph(_rn, styles["small"]), Paragraph(_fmt_r(_rv), styles["kpi"])])
            _r_tbl_s = Table(_r_data_s, colWidths=[200, 260])
            _r_tbl_s.setStyle(_tbl_style("#4B5563"))
            story.append(_r_tbl_s)
        story.append(Spacer(1, 12))
        story.append(Paragraph(
            "Sharpe = rendement annualisé / volatilité (plus élevé = meilleur). "
            "Sortino = idem, volatilité baissière uniquement. "
            "Bêta = sensibilité aux marchés (1 = marché). "
            "Alpha = surperformance annualisée vs indice.",
            styles["small"],
        ))

    # ══════════════════════════════════════════════════════════
    # PAGE 4 — Transparence des frais
    # ══════════════════════════════════════════════════════════
    story.append(PageBreak())
    story.append(Paragraph("Transparence des frais", styles["h1"]))

    def _build_fee_table_p4(lines_list):
        if not lines_list:
            return None, 0.0
        _hdr = [
            Paragraph("<b>Fonds</b>",          styles["h2"]),
            Paragraph("<b>TER (UC)</b>",        styles["h2"]),
            Paragraph("<b>Frais contrat</b>",   styles["h2"]),
            Paragraph("<b>Total</b>",           styles["h2"]),
        ]
        _data = [_hdr]
        _tot_val = _w_ter = _w_ctr = _w_tot = 0.0
        for ln in lines_list:
            _isin = str(ln.get("isin", "")).upper()
            _name = ln.get("name") or _isin or "—"
            _val  = float(ln.get("value", 0) or 0)
            _tot_val += _val
            if _isin == "EUROFUND":
                _data.append([
                    Paragraph(_name[:42], styles["kpi"]),
                    Paragraph("inclus", styles["kpi"]),
                    Paragraph("inclus", styles["kpi"]),
                    Paragraph("inclus", styles["kpi"]),
                ])
                continue
            _ter_f = float(ln.get("fee_uc_pct") or 0)
            _ctr_f = float(ln.get("fee_contract_pct") or 0)
            _tot_f = float(ln.get("fee_total_pct") or 0)
            if _tot_f == 0 and (_ter_f > 0 or _ctr_f > 0):
                _tot_f = _ter_f + _ctr_f
            _s_ter = f"{_ter_f:.2f}%" if _ter_f > 0 else "—"
            _s_ctr = f"{_ctr_f:.2f}%" if _ctr_f > 0 else "—"
            _s_tot = f"{_tot_f:.2f}%" if _tot_f > 0 else "—"
            _w_ter += _ter_f * _val
            _w_ctr += _ctr_f * _val
            _w_tot += _tot_f * _val
            _tot_style = styles["fee_high"] if _tot_f > 2.0 else styles["kpi"]
            _data.append([
                Paragraph(_name[:42], styles["kpi"]),
                Paragraph(_s_ter,     styles["kpi"]),
                Paragraph(_s_ctr,     styles["kpi"]),
                Paragraph(_s_tot,     _tot_style),
            ])
        if _tot_val > 0:
            _data.append([
                Paragraph("<b>Moyenne pondérée</b>",              styles["h2"]),
                Paragraph(f"<b>{_w_ter/_tot_val:.2f}%</b>",      styles["h2"]),
                Paragraph(f"<b>{_w_ctr/_tot_val:.2f}%</b>",      styles["h2"]),
                Paragraph(f"<b>{_w_tot/_tot_val:.2f}%</b>",      styles["h2"]),
            ])
        _tbl = Table(_data, colWidths=[195, 85, 95, 85])
        _ts = _tbl_style()
        if _tot_val > 0:
            _ts.add("BACKGROUND", (0, len(_data) - 1), (-1, len(_data) - 1), colors.HexColor("#EEF2F7"))
        _tbl.setStyle(_ts)
        return _tbl, _w_tot / _tot_val if _tot_val > 0 else 0.0

    _lines_A_p4 = report.get("lines_client", [])
    _lines_B_p4 = report.get("lines_valority", [])
    _fee_A_p4   = float(report.get("fee_gestion_A") or 0)
    _fee_B_p4   = float(report.get("fee_gestion_B") or 0)

    if mode == "compare":
        _fee_tbl_A, _ = _build_fee_table_p4(_lines_A_p4)
        _fee_tbl_B, _ = _build_fee_table_p4(_lines_B_p4)
        if _fee_tbl_A:
            story.append(Paragraph("Votre portefeuille", styles["h2"]))
            story.append(_fee_tbl_A)
            story.append(Spacer(1, 8))
        if _fee_tbl_B:
            story.append(Paragraph("Notre proposition", styles["h2"]))
            story.append(_fee_tbl_B)
            story.append(Spacer(1, 8))
        if _fee_A_p4 > 0 and _fee_B_p4 > 0 and _fee_B_p4 < _fee_A_p4:
            _eco_p4 = (
                100_000 * ((1.05)**15 - (1.05 - _fee_A_p4/100)**15) -
                100_000 * ((1.05)**15 - (1.05 - _fee_B_p4/100)**15)
            )
            _eco_ann = float(report.get("client_summary", {}).get("val", 100_000)) * (
                _fee_A_p4 - _fee_B_p4) / 100.0
            story.append(Paragraph(
                f"Frais actuels : <b>{_fee_A_p4:.2f}%/an</b>  →  "
                f"Notre proposition : <b>{_fee_B_p4:.2f}%/an</b>.",
                styles["kpi"],
            ))
            story.append(Paragraph(
                f"Économie annuelle estimée : <b>{_fe_pdf(_eco_ann)}/an</b> en faveur du cabinet. "
                f"Sur 15 ans : <b>{_fe_pdf(_eco_p4)}</b> "
                f"(base 100 000 €, rendement brut 5%/an).",
                styles["verdict_pos"],
            ))
        elif _fee_A_p4 > 0 or _fee_B_p4 > 0:
            story.append(Paragraph(
                f"Frais de gestion — Votre portefeuille : {_fee_A_p4:.2f}%/an  ·  "
                f"Notre proposition : {_fee_B_p4:.2f}%/an.",
                styles["small"],
            ))
    else:
        _lines_p4 = _lines_B_p4 if mode == "valority" else _lines_A_p4
        _fee_tbl_s, _ = _build_fee_table_p4(_lines_p4)
        if _fee_tbl_s:
            story.append(_fee_tbl_s)
            story.append(Spacer(1, 8))

    story.append(Spacer(1, 15))
    story.append(Paragraph(
        "TER (Total Expense Ratio / frais UC) : frais internes du fonds, déjà intégrés dans la valeur liquidative. "
        "Frais contrat : frais annuels prélevés par l'assureur. "
        "Total = TER + frais contrat.",
        styles["small"],
    ))

    # ══════════════════════════════════════════════════════════
    # PAGES 5-6 — Détail des portefeuilles
    # ══════════════════════════════════════════════════════════
    def add_portfolio_details_section(
        label: str,
        positions_df: pd.DataFrame,
        lines_values: List[Dict[str, Any]],
    ):
        story.append(PageBreak())
        story.append(Paragraph(f"Détail du portefeuille — {label}", styles["h1"]))

        if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
            df_alloc, basis_label = _allocation_from_positions(positions_df)

            if len(df_alloc) >= 2:
                story.append(Paragraph(f"Allocation par ligne ({basis_label})", styles["h2"]))
                donut = _build_allocation_donut(df_alloc, "Allocation par ligne")
                if donut is not None:
                    story.append(donut)
                    story.append(Spacer(1, 6))
            else:
                line = df_alloc.iloc[0] if not df_alloc.empty else None
                name = line["Nom"] if line is not None else "—"
                story.append(
                    Paragraph(
                        f"Portefeuille concentré : 100% sur <b>{name}</b>.",
                        styles["kpi"],
                    )
                )
                bar = _build_single_line_bar(_wrap_label(name), 100.0, "Répartition 100%")
                if bar is not None:
                    story.append(bar)
                    story.append(Spacer(1, 6))

            alloc_table = df_alloc[
                ["Nom", "ISIN / Code", "Part %", "Net investi €", "Valeur actuelle €"]
            ].copy()
            _add_table_to_story(
                story,
                alloc_table,
                col_widths=[170, 80, 55, 95, 95],
                font_size=9,
            )
            story.append(Spacer(1, 6))

            envelope_chart, envelope_text = _build_envelope_breakdown(
                lines_values,
                "Répartition par enveloppe",
            )
            if envelope_chart is not None:
                story.append(envelope_chart)
            elif envelope_text:
                story.append(Paragraph(envelope_text, styles["small"]))
        else:
            story.append(Paragraph("Données indisponibles pour la composition.", styles["small"]))

        story.append(PageBreak())
        story.append(Paragraph(f"Contribution & positions — {label}", styles["h1"]))

        if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
            if len(positions_df) == 1:
                ln = positions_df.iloc[0]
                story.append(
                    Paragraph(
                        f"Contribution : <b>{ln['Nom']}</b> = {fmt_eur_fr(ln['Valeur actuelle €'] - ln['Net investi €'])}",
                        styles["kpi"],
                    )
                )
                bar = _build_single_line_bar(_wrap_label(ln["Nom"]), 100.0, "Contribution (ligne unique)")
                if bar is not None:
                    story.append(bar)
                    story.append(Spacer(1, 6))
            else:
                contrib_chart = _build_contribution_bar(positions_df)
                if contrib_chart is not None:
                    story.append(contrib_chart)
                    story.append(Spacer(1, 6))
        else:
            story.append(Paragraph("Données indisponibles pour la contribution.", styles["small"]))

        story.append(Paragraph("Positions", styles["h2"]))
        if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
            positions_table = positions_df[
                [
                    "Nom",
                    "ISIN / Code",
                    "Date d'achat",
                    "Net investi €",
                    "Valeur actuelle €",
                    "Perf €",
                    "Perf %",
                ]
            ].copy()
            _add_table_to_story(
                story,
                positions_table,
                col_widths=[150, 70, 65, 80, 80, 50, 45],
                font_size=9,
            )
        else:
            story.append(Paragraph("Données indisponibles.", styles["small"]))

    if mode == "compare":
        add_portfolio_details_section(
            "Client",
            report.get("positions_df_client", pd.DataFrame()),
            report.get("lines_client", []),
        )
        add_portfolio_details_section(
            "Cabinet",
            report.get("positions_df_valority", pd.DataFrame()),
            report.get("lines_valority", []),
        )
    else:
        add_portfolio_details_section(
            "Cabinet" if mode == "valority" else "Client",
            report.get("positions_df", pd.DataFrame()),
            report.get("lines", []),
        )

    # ══════════════════════════════════════════════════════════
    # PAGE 7 — Fiscalité (si rachat > 0)
    # ══════════════════════════════════════════════════════════
    _rachat_A_p7 = float(report.get("rachat_optimal_A") or 0)
    _rachat_B_p7 = float(report.get("rachat_optimal_B") or 0)
    if _rachat_A_p7 > 0 or _rachat_B_p7 > 0:
        story.append(PageBreak())
        story.append(Paragraph("Optimisation fiscale", styles["h1"]))
        _sit_p7  = report.get("situation_familiale", "")
        _abat_p7 = 9_200.0 if "Couple" in str(_sit_p7) else 4_600.0
        # Tableau rachat
        _rach_data = [[
            Paragraph("<b>Paramètre</b>",      styles["h2"]),
            Paragraph("<b>Votre situation</b>", styles["h2"]),
        ]]
        _rach_data.append([Paragraph("Situation familiale",    styles["small"]),
                           Paragraph(str(_sit_p7) or "—",     styles["kpi"])])
        _rach_data.append([Paragraph("Abattement annuel AV",   styles["small"]),
                           Paragraph(_fe_pdf(_abat_p7) + "/an", styles["kpi"])])
        if _rachat_A_p7 > 0:
            _rach_data.append([
                Paragraph("Rachat sans IR — portefeuille actuel", styles["small"]),
                Paragraph(f"≈ {_fe_pdf(_rachat_A_p7)}/an",        styles["kpi"]),
            ])
        if _rachat_B_p7 > 0:
            _rach_data.append([
                Paragraph("Rachat sans IR — notre proposition",    styles["small"]),
                Paragraph(f"≈ {_fe_pdf(_rachat_B_p7)}/an",        styles["kpi"]),
            ])
        _rach_tbl = Table(_rach_data, colWidths=[260, 200])
        _rach_tbl.setStyle(_tbl_style())
        story.append(_rach_tbl)
        story.append(Spacer(1, 12))
        # Comparatif AV / CTO / PEA
        story.append(Paragraph("Comparatif fiscal AV / CTO / PEA (2026)", styles["h2"]))
        _cmp_data = [[
            Paragraph("<b>Critère</b>",              styles["h2"]),
            Paragraph("<b>AV (≥ 8 ans)</b>",         styles["h2"]),
            Paragraph("<b>CTO</b>",                  styles["h2"]),
            Paragraph("<b>PEA</b>",                  styles["h2"]),
        ]]
        _cmp_rows_p7 = [
            ("IR à la sortie",       "7,5% (après abatt.)",          "12,8% flat tax",        "0% (après 5 ans)", False),
            ("Prélèvements sociaux", "17,2%",                         "18,6%",                 "18,6%",            True),
            ("Abattement annuel",    f"{_fe_pdf(_abat_p7)}/pers.",    "Aucun",                 "Aucun",            True),
            ("Transmission",         "Hors succession (art. 990I)", "Succession classique",  "Succession classique", True),
            ("Plafond versements",   "Illimité",                      "Illimité",              f"{_fe_pdf(150_000)}", False),
        ]
        for _c0, _c1, _c2, _c3, _av_highlight in _cmp_rows_p7:
            _av_style = styles["av_green"] if _av_highlight else styles["kpi"]
            _cmp_data.append([
                Paragraph(_c0, styles["small"]),
                Paragraph(_c1, _av_style),
                Paragraph(_c2, styles["kpi"]),
                Paragraph(_c3, styles["kpi"]),
            ])
        _cmp_tbl = Table(_cmp_data, colWidths=[145, 120, 120, 115])
        _cmp_tbl.setStyle(_tbl_style("#1F3B6D"))
        story.append(_cmp_tbl)
        story.append(Spacer(1, 15))
        # Transmission 990I
        story.append(Paragraph("<b>Transmission — art. 990I CGI</b>", styles["h2"]))
        story.append(Paragraph(
            "Les capitaux versés avant 70 ans bénéficient d'un abattement de 152 500 € par bénéficiaire "
            "(art. 990I CGI), puis d'un prélèvement forfaitaire de 20% jusqu'à 852 500 € et 31,25% au-delà. "
            "Ces sommes sont transmises hors succession, sans droits de mutation classiques. "
            "Une désignation bénéficiaire précise et régulièrement mise à jour est indispensable "
            "pour optimiser ce dispositif.",
            styles["kpi"],
        ))

    doc.build(story, canvasmaker=NumberedCanvas)
    buffer.seek(0)
    return buffer.read()




def generate_pptx_report(report: Dict[str, Any]) -> bytes:
    """Génère une présentation 16:9 de 5 slides pour le client."""
    if not PPTX_AVAILABLE:
        raise RuntimeError(f"PPTX indisponible: {PPTX_ERROR}")

    # ── Palette ──────────────────────────────────────────────────────────
    NAVY   = RGBColor(0x1F, 0x3B, 0x6D)
    GOLD   = RGBColor(0xC8, 0x96, 0x3E)
    ICE    = RGBColor(0xE8, 0xEE, 0xF7)
    GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
    GREY   = RGBColor(0x64, 0x74, 0x8B)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGREY  = RGBColor(0xF1, 0xF4, 0xF9)
    BORDER = RGBColor(0xD1, 0xD9, 0xE6)
    BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
    BLACK  = RGBColor(0x11, 0x18, 0x27)

    # ── Données du rapport ────────────────────────────────────────────────
    nom_cab  = report.get("nom_cabinet", "") or "Votre cabinet"
    nom_cli  = report.get("nom_client", "") or ""
    as_of    = report.get("as_of", "")
    contrat  = report.get("contrat_label", "Assurance-vie") or "Assurance-vie"
    date_ouv = report.get("date_ouverture", "")
    mode     = report.get("mode", "compare")
    _include_fiscal = report.get("include_fiscal_slides", True)

    synthA  = report.get("client_summary", {})
    synthB  = report.get("valority_summary", {})
    comp    = report.get("comparison", {})

    # Synthèse principale selon le mode
    if mode == "valority":
        synth_main = synthB
        _raw_df = report.get("positions_df_valority")
        positions_df_main = _raw_df if isinstance(_raw_df, pd.DataFrame) else pd.DataFrame()
    else:
        synth_main = synthA
        _raw_df = report.get("positions_df_client")
        positions_df_main = _raw_df if isinstance(_raw_df, pd.DataFrame) else pd.DataFrame()

    val  = synth_main.get("val", 0.0)
    net  = synth_main.get("net", 0.0)
    brut = synth_main.get("brut", 0.0)
    perf = synth_main.get("perf_tot_pct", 0.0)
    xirr_val = synth_main.get("irr_pct", 0.0)
    # Guard : si val = 0, utiliser valB ou valA selon disponibilité
    if val == 0.0 and mode == "compare":
        val  = report.get("valority_summary", {}).get("val", 0.0) or \
               report.get("client_summary", {}).get("val", 0.0)
        net  = report.get("valority_summary", {}).get("net", 0.0) or \
               report.get("client_summary", {}).get("net", 0.0)
        brut = report.get("valority_summary", {}).get("brut", 0.0) or \
               report.get("client_summary", {}).get("brut", 0.0)

    def _fe(x: Any) -> str:
        try:
            return f"{float(x):,.0f} €".replace(",", " ")
        except Exception:
            return "— €"

    def _fp(x: Any) -> str:
        try:
            return f"{float(x):+.2f}%"
        except Exception:
            return "—%"

    # ── Helpers ──────────────────────────────────────────────────────────
    def _add_rect(slide, l: float, t: float, w: float, h: float,
                  fill_rgb: RGBColor, line_rgb: Optional[RGBColor] = None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE RECTANGLE
            Inches(l), Inches(t), Inches(w), Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if line_rgb is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line_rgb
        return shape

    def _add_text(slide, text: str, l: float, t: float, w: float, h: float,
                  size: float, color: RGBColor, bold: bool = False,
                  align: str = "left", italic: bool = False) -> None:
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.RIGHT if align == "right" else PP_ALIGN.LEFT

    def _footer(slide):
        parts = [nom_cab]
        if nom_cli:
            parts.append(f"Document pour {nom_cli}")
        if as_of:
            parts.append(as_of)
        _add_text(slide, "  ·  ".join(parts),
                  0.28, 5.38, 9.44, 0.20, 7, GREY, align="center")

    def _left_stripe(slide, color: RGBColor = NAVY):
        _add_rect(slide, 0, 0, 0.18, 5.625, color)

    # ── Présentation ──────────────────────────────────────────────────────
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _In, Pt as _Pt, Emu as _Emu

    prs = _Prs()
    prs.slide_width  = _In(10)
    prs.slide_height = _In(5.625)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Page de garde
    # ══════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    _add_rect(s1, 0, 0, 10, 5.625, NAVY)
    _add_rect(s1, 0, 0, 0.18, 5.625, GOLD)
    _add_rect(s1, 0.55, 0.70, 8.80, 2.00, LGREY)

    _add_text(s1, nom_cab.upper(),
              0.55, 0.80, 8.80, 0.32, 9, GOLD, bold=False, align="center")
    if mode == "client":
        _titre_s1 = "Bilan Patrimonial"
        _sous_titre_s1 = f"Situation actuelle de {nom_cli}" if nom_cli else "Situation actuelle"
    elif mode == "valority":
        _titre_s1 = "Proposition d'Investissement"
        _sous_titre_s1 = f"Une stratégie élaborée par {nom_cab}" if nom_cab else "Stratégie conseiller"
    else:
        _titre_s1 = "Bilan & Recommandation"
        _sous_titre_s1 = f"Audit & Recommandation {nom_cab}" if nom_cab else "Comparatif complet"
    _add_text(s1, _titre_s1,
              0.55, 1.00, 8.80, 0.55, 28, NAVY, bold=True, align="center")
    _add_text(s1, _sous_titre_s1,
              0.55, 1.58, 8.80, 0.32, 13, GREY, align="center")

    _add_rect(s1, 3.20, 2.80, 3.60, 0.04, GOLD)
    _add_text(s1, nom_cli if nom_cli else "—",
              0.55, 2.95, 8.80, 0.42, 20, WHITE, bold=True, align="center")
    _add_text(s1, contrat,
              0.55, 3.50, 8.80, 0.28, 9, ICE, align="center")
    _dates_str = ""
    if date_ouv:
        _dates_str = f"Ouverture : {date_ouv}"
    if as_of:
        _dates_str += (f"  ·  Arrêté au : {as_of}" if _dates_str else f"Arrêté au : {as_of}")
    _add_text(s1, _dates_str,
              0.55, 3.80, 8.80, 0.28, 8, GREY, align="center")
    _add_text(s1,
              "Document à usage interne — Simulation historique. "
              "Les performances passées ne préjugent pas des performances futures.",
              0.55, 5.20, 8.80, 0.28, 6, GREY, align="center")

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Votre portefeuille aujourd'hui
    # ══════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    _add_rect(s2, 0, 0, 10, 5.625, LGREY)
    _left_stripe(s2)
    _footer(s2)

    _add_text(s2, "Votre portefeuille aujourd'hui",
              0.30, 0.12, 9.30, 0.42, 20, NAVY, bold=True)

    # Message humain
    _vc_s2 = float(val or 0) - float(net or 0)
    if _vc_s2 >= 0:
        _msg_s2 = (f"Vous avez investi {_fe(net)} sur ce contrat. "
                   f"Il vaut aujourd'hui {_fe(val)}, soit un gain de {_fe(_vc_s2)}.")
    else:
        _msg_s2 = (f"Vous avez investi {_fe(net)} sur ce contrat. "
                   f"Il vaut aujourd'hui {_fe(val)}.")
    _add_text(s2, _msg_s2, 0.30, 0.54, 9.30, 0.30, 9, GREY)

    # 4 KPIs
    _xirr_s2_str = f"{float(xirr_val):+.2f}%" if xirr_val not in (None, 0.0, "") else "—"
    _kpi_data = [
        ("Valeur actuelle",     _fe(val),  f"Versé : {_fe(brut)}", NAVY),
        ("Capital net investi", _fe(net),  "Après frais d'entrée",  NAVY),
        ("Performance totale",  f"{float(perf or 0):+.2f}%", "Depuis l'ouverture", GREEN),
        ("Rendement annualisé", _xirr_s2_str, "XIRR", BLUE),
    ]
    _kpi_x = [0.28, 2.58, 4.88, 7.18]
    for idx, (lbl, val_str, sub_lbl, val_color) in enumerate(_kpi_data):
        _xk = _kpi_x[idx]
        _add_rect(s2, _xk, 1.02, 2.20, 1.00, WHITE, BORDER)
        _add_text(s2, lbl,     _xk + 0.08, 1.07, 2.04, 0.22, 8, GREY)
        _add_text(s2, val_str, _xk + 0.08, 1.29, 2.04, 0.30, 17, val_color, bold=True)
        _add_text(s2, sub_lbl, _xk + 0.08, 1.61, 2.04, 0.22, 7, GREY)

    # Composition du portefeuille
    _add_text(s2, "Composition du portefeuille",
              0.28, 2.18, 9.10, 0.28, 11, NAVY, bold=True)

    if not positions_df_main.empty:
        _df_tbl = positions_df_main.copy()
        _row_y = 2.46
        _col_w = [3.5, 0.85, 1.5, 1.65, 1.0]
        _col_x_s2 = [0.28]
        for _cw in _col_w[:-1]:
            _col_x_s2.append(_col_x_s2[-1] + _cw)
        _headers = ["Fonds", "Part %", "Net investi", "Valeur actuelle", "Perf %"]
        for ci, (_cx, _ch) in enumerate(zip(_col_x_s2, _headers)):
            _add_rect(s2, _cx, _row_y, _col_w[ci], 0.25, NAVY)
            _add_text(s2, _ch, _cx + 0.04, _row_y + 0.02, _col_w[ci] - 0.08, 0.22, 9, WHITE, bold=True)

        def _get_row_val(row_series, *names):
            for n in names:
                v = row_series.get(n)
                if v is not None and v == v:
                    return v
            return 0

        _val_col = "Valeur actuelle €" if "Valeur actuelle €" in _df_tbl.columns else (
            "Valeur" if "Valeur" in _df_tbl.columns else None)
        _total_val_tbl = _df_tbl[_val_col].sum() if _val_col else 0.0

        for ri, (_, row) in enumerate(_df_tbl.head(7).iterrows()):
            _ry = _row_y + 0.25 + ri * 0.26
            _row_bg = WHITE if ri % 2 == 0 else LGREY
            _add_rect(s2, 0.28, _ry, 9.10, 0.26, _row_bg)
            _nom_tbl = str(row.get("Nom", "—"))[:38]
            _net_tbl = float(_get_row_val(row, "Net investi €", "Net investi", "Net") or 0)
            _val_tbl = float(_get_row_val(row, "Valeur actuelle €", "Valeur actuelle", "Valeur") or 0)
            _perf_tbl = ((_val_tbl / _net_tbl - 1) * 100) if _net_tbl > 0 else 0.0
            _part_tbl = (_val_tbl / _total_val_tbl * 100) if _total_val_tbl > 0 else 0.0
            _perf_color = GREEN if _perf_tbl >= 0 else RGBColor(0xCC, 0x22, 0x00)
            _row_vals = [_nom_tbl, f"{_part_tbl:.1f}%", _fe(_net_tbl), _fe(_val_tbl), f"{_perf_tbl:+.2f}%"]
            for ci2, (_cx2, _rv) in enumerate(zip(_col_x_s2, _row_vals)):
                _tc = _perf_color if ci2 == 4 else BLACK
                _add_text(s2, _rv, _cx2 + 0.04, _ry + 0.03, _col_w[ci2] - 0.08, 0.20, 8, _tc)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Ce que vous pouvez faire aujourd'hui
    # ══════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    _add_rect(s3, 0, 0, 10, 5.625, LGREY)
    _left_stripe(s3)
    _footer(s3)

    _add_text(s3, "Ce que vous pouvez faire aujourd'hui",
              0.50, 0.15, 9.00, 0.42, 20, NAVY, bold=True)
    _add_text(s3, f"Situation au {as_of} — {contrat}",
              0.50, 0.57, 9.00, 0.25, 9, GREY)

    # Bloc Retrait
    _rachat_opt = float(report.get("rachat_optimal_A") or 0)
    _rachat_ps  = float(report.get("rachat_ps_A") or 0)
    _add_rect(s3, 0.50, 1.00, 4.20, 1.90, WHITE, BORDER)
    _add_text(s3, "Retrait sans impôt sur le revenu",
              0.70, 1.10, 3.80, 0.25, 12, NAVY, bold=True)
    if _rachat_opt > 0:
        _add_text(s3, _fe(_rachat_opt),
                  0.70, 1.46, 3.80, 0.40, 26, GREEN, bold=True)
        _add_text(s3, f"Net en main  ·  PS uniquement : {_fe(_rachat_ps)}",
                  0.70, 1.94, 3.80, 0.22, 9, GREY)
        _add_text(s3, "Sans payer d'IR — abattement contrat ≥ 8 ans",
                  0.70, 2.16, 3.80, 0.22, 8, GREY)
    else:
        _add_text(s3, "Contrat < 8 ans",
                  0.70, 1.46, 3.80, 0.30, 14, GREY)
        _add_text(s3, "Flat tax 30% applicable — maturité fiscale non atteinte",
                  0.70, 1.86, 3.80, 0.22, 9, GREY)

    # Bloc Transmission
    _nb_b_s3   = int(report.get("nb_beneficiaires") or 2)
    _age_s3    = int(report.get("age_souscripteur") or 55)
    _val_A_s3  = float(synthA.get("val") or val or 0)
    _add_rect(s3, 5.30, 1.00, 4.20, 1.90, WHITE, BORDER)
    _add_text(s3, "Transmission à vos proches",
              5.50, 1.10, 3.80, 0.25, 12, NAVY, bold=True)
    if _age_s3 < 70:
        _add_text(s3, f"{_fe(_val_A_s3)} exonérés",
                  5.50, 1.46, 3.80, 0.40, 22, GREEN, bold=True)
        _add_text(s3, f"Art. 990I — {_fe(152_500)}/bénéficiaire × {_nb_b_s3}",
                  5.50, 1.94, 3.80, 0.22, 9, GREY)
        _add_text(s3, "Hors succession · Versements avant 70 ans",
                  5.50, 2.16, 3.80, 0.22, 8, GREY)
    else:
        _add_text(s3, "Art. 757B — Abattement 30 500 €",
                  5.50, 1.46, 3.80, 0.30, 14, GREY)
        _add_text(s3, "Intérêts et plus-values exonérés",
                  5.50, 1.86, 3.80, 0.22, 9, GREY)

    # Bloc Projection 10/15/20 ans
    _xirr_s3_pct = float(synthA.get("irr_pct") or 0) or 3.0
    _net_s3 = float(synthA.get("net") or net or 0)
    _sit_s3 = report.get("situation_familiale", "")
    _abat_s3 = 9_200.0 if "Couple" in str(_sit_s3) else 4_600.0
    _xirr_s3 = _xirr_s3_pct / 100.0
    _add_rect(s3, 0.50, 3.10, 9.00, 2.05, WHITE, BORDER)
    _add_text(s3, f"Projection si le rendement se maintient ({_xirr_s3_pct:+.1f}%/an)",
              0.70, 3.20, 8.60, 0.25, 11, NAVY, bold=True)
    for _i, _horizon in enumerate([10, 15, 20]):
        _val_proj_s3 = _val_A_s3 * (1 + _xirr_s3) ** _horizon if _val_A_s3 > 0 else 0
        _col_x_s3 = 0.70 + _i * 2.90
        _add_rect(s3, _col_x_s3, 3.52, 2.60, 1.47, LGREY)
        _add_text(s3, f"Dans {_horizon} ans",
                  _col_x_s3 + 0.08, 3.57, 2.45, 0.22, 9, GREY, bold=True)
        _add_text(s3, f"~{_fe(_val_proj_s3)}",
                  _col_x_s3 + 0.08, 3.79, 2.45, 0.28, 13, NAVY, bold=True)
        _taux_pv_proj = max(0.0, 1.0 - _net_s3 / max(_val_proj_s3, 1))
        _rachat_proj = (min(_val_proj_s3, _abat_s3 / _taux_pv_proj)
                        if _taux_pv_proj > 0 else _val_proj_s3)
        _add_text(s3, f"Retrait IR-free : ~{_fe(_rachat_proj)}/an",
                  _col_x_s3 + 0.08, 4.09, 2.45, 0.20, 8, GREEN)
        _add_text(s3, f"Transmission : ~{_fe(_val_proj_s3)}",
                  _col_x_s3 + 0.08, 4.32, 2.45, 0.20, 8, NAVY)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Ce que nous aurions pu faire ensemble  (mode compare uniquement)
    # ══════════════════════════════════════════════════════════════════════
    if mode == "compare":
        s4 = prs.slides.add_slide(blank_layout)
        _add_rect(s4, 0, 0, 10, 5.625, LGREY)
        _left_stripe(s4)
        _footer(s4)

        _val_A_s4   = float(synthA.get("val") or 0)
        _val_B_s4   = float(synthB.get("val") or 0)
        _xirr_A_s4  = float(synthA.get("irr_pct") or 0)
        _xirr_B_s4  = float(synthB.get("irr_pct") or 0)
        _perf_A_s4  = float(synthA.get("perf_tot_pct") or 0)
        _perf_B_s4  = float(synthB.get("perf_tot_pct") or 0)
        _delta_s4   = _val_B_s4 - _val_A_s4
        _delta_sign_s4 = "+" if _delta_s4 >= 0 else ""

        _add_text(s4, "Ce que nous aurions pu faire ensemble",
                  0.30, 0.12, 9.30, 0.42, 20, NAVY, bold=True)
        _add_text(s4, "Simulation historique sur la même période avec l'allocation conseillée",
                  0.30, 0.54, 9.30, 0.25, 9, GREY)

        # Bloc choc central
        _add_rect(s4, 0.30, 0.92, 9.40, 1.10, NAVY)
        _add_text(s4, "La différence entre les deux allocations :",
                  0.50, 0.96, 9.00, 0.28, 10, ICE)
        _add_text(s4, f"{_delta_sign_s4}{_fe(_delta_s4)}",
                  0.30, 1.08, 9.40, 0.70, 38, GOLD, bold=True, align="center")

        # Graphique superposé : client + cabinet
        _dfA_val_s4 = report.get("dfA_val")
        _dfB_val_s4 = report.get("dfB_val")
        # Fallback sur df_map si les séries directes sont vides
        if (not isinstance(_dfA_val_s4, pd.DataFrame) or _dfA_val_s4.empty):
            _df_map_s4 = report.get("df_map", {})
            _dfA_val_s4 = _df_map_s4.get("A") if isinstance(_df_map_s4, dict) else None
        if (not isinstance(_dfB_val_s4, pd.DataFrame) or _dfB_val_s4.empty):
            _df_map_s4 = report.get("df_map", {})
            _dfB_val_s4 = _df_map_s4.get("B") if isinstance(_df_map_s4, dict) else None
        _has_A_s4 = isinstance(_dfA_val_s4, pd.DataFrame) and not _dfA_val_s4.empty and "Valeur" in _dfA_val_s4.columns
        _has_B_s4 = isinstance(_dfB_val_s4, pd.DataFrame) and not _dfB_val_s4.empty and "Valeur" in _dfB_val_s4.columns
        if _has_A_s4 or _has_B_s4:
            _fig_s4, _ax_s4 = plt.subplots(figsize=(8.5, 3.2), dpi=150)
            _fig_s4.patch.set_facecolor("#F5F5F5")
            _ax_s4.set_facecolor("#F5F5F5")
            _yA_s4 = _yB_s4 = None
            if _has_A_s4:
                _dfA_plt = _dfA_val_s4.copy()
                if "Date" in _dfA_plt.columns:
                    _dfA_plt = _dfA_plt.set_index("Date")
                _yA_s4 = _dfA_plt["Valeur"]
                _ax_s4.plot(_yA_s4.index, _yA_s4.values,
                            color="#CC2200", linewidth=2, label="Votre portefeuille")
            if _has_B_s4:
                _dfB_plt = _dfB_val_s4.copy()
                if "Date" in _dfB_plt.columns:
                    _dfB_plt = _dfB_plt.set_index("Date")
                _yB_s4 = _dfB_plt["Valeur"]
                _ax_s4.plot(_yB_s4.index, _yB_s4.values,
                            color="#1A7A4A", linewidth=2, label="Notre proposition")
            if _yA_s4 is not None and _yB_s4 is not None:
                _idx_common = _yA_s4.index.intersection(_yB_s4.index)
                if len(_idx_common):
                    _ax_s4.fill_between(_idx_common,
                                        _yA_s4.reindex(_idx_common),
                                        _yB_s4.reindex(_idx_common),
                                        color="#1A7A4A", alpha=0.10)
            _ax_s4.yaxis.set_major_formatter(
                plt.FuncFormatter(lambda v, _: f"{int(v):,}".replace(",", " "))
            )
            _ax_s4.set_title("Évolution comparée de la valeur", fontsize=10, color="#1F3B6D", pad=6)
            _ax_s4.tick_params(labelsize=7)
            _ax_s4.spines[["top", "right"]].set_visible(False)
            _ax_s4.legend(loc="lower center", bbox_to_anchor=(0.5, -0.22),
                          ncol=2, fontsize=8, frameon=False)
            _buf_s4 = BytesIO()
            _fig_s4.tight_layout()
            _fig_s4.savefig(_buf_s4, format="png", bbox_inches="tight")
            plt.close(_fig_s4)
            _buf_s4.seek(0)
            s4.shapes.add_picture(
                _buf_s4,
                left=Inches(0.50), top=Inches(2.10),
                width=Inches(9.00), height=Inches(2.10),
            )

        # 3 KPIs horizontaux
        _xirr_A_s4_str = f"{_xirr_A_s4:+.2f}%/an" if _xirr_A_s4 not in (None, 0.0) else "—"
        _xirr_B_s4_str = f"{_xirr_B_s4:+.2f}%/an" if _xirr_B_s4 not in (None, 0.0) else "—"
        _kpi_defs4 = [
            ("Votre portefeuille",  f"{_perf_A_s4:+.1f}% ({_xirr_A_s4_str})",   NAVY),
            ("Notre proposition",   f"{_perf_B_s4:+.1f}% ({_xirr_B_s4_str})",   NAVY),
            ("Différence",          f"{_delta_sign_s4}{_fe(_delta_s4)}",          GREEN),
        ]
        _kpi_x4 = [0.30, 3.55, 6.80]
        for _ki4, (_klbl4, _kval4, _kcol4) in enumerate(_kpi_defs4):
            _kx4 = _kpi_x4[_ki4]
            _add_rect(s4, _kx4, 4.28, 3.10, 0.92, LGREY)
            _add_text(s4, _klbl4, _kx4 + 0.12, 4.34, 2.86, 0.22, 9, GREY)
            _add_text(s4, _kval4, _kx4 + 0.12, 4.57, 2.86, 0.38, 18, _kcol4, bold=True)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 5bis — Pourquoi cette différence ?  (mode compare uniquement)
    # ══════════════════════════════════════════════════════════════════════
    if mode == "compare":
        _div_A  = report.get("diversification_A") or {}
        _div_B  = report.get("diversification_B") or {}
        _sharpe_A  = float(report.get("sharpe_A")  or 0)
        _sharpe_B  = float(report.get("sharpe_B")  or 0)
        _sortino_A = float(report.get("sortino_A") or 0)
        _sortino_B = float(report.get("sortino_B") or 0)
        _fee_A  = float(report.get("fee_gestion_A") or 0)
        _fee_B  = float(report.get("fee_gestion_B") or 0)

        # — Calculs bloc Diversification
        _show_div = bool(_div_A or _div_B)
        _n_red_A  = int(_div_A.get("n_redundant", 0)) if _div_A else 0
        _n_eff_A  = int(_div_A.get("n_effective", 0)) if _div_A else 0
        _n_lin_A  = int(_div_A.get("n_lines",     0)) if _div_A else 0
        _n_eff_B  = int(_div_B.get("n_effective", 0)) if _div_B else 0
        _n_lin_B  = int(_div_B.get("n_lines",     0)) if _div_B else 0

        # — Calculs bloc Efficacité
        _show_eff = (_sharpe_A > 0 or _sharpe_B > 0 or _sortino_A > 0 or _sortino_B > 0)
        _ratio_sharpe      = (_sharpe_B / _sharpe_A) if _sharpe_A > 0 else 0
        _reduction_sortino = (
            (1 - (1 / _sortino_B) / (1 / _sortino_A)) * 100
            if _sortino_A > 0 and _sortino_B > 0 else 0
        )

        # — Calculs bloc Frais
        _show_frais = (_fee_A > 0 or _fee_B > 0)
        _drag_A_15  = 100_000 * ((1.05) ** 15 - (1.05 - _fee_A / 100) ** 15)
        _drag_B_15  = 100_000 * ((1.05) ** 15 - (1.05 - _fee_B / 100) ** 15)
        _eco_15     = _drag_A_15 - _drag_B_15

        # Ne créer la slide que si au moins un bloc a des données
        if _show_div or _show_eff or _show_frais:
            s5b = prs.slides.add_slide(blank_layout)
            _add_rect(s5b, 0, 0, 10, 5.625, WHITE)
            _left_stripe(s5b)
            _footer(s5b)

            _add_text(s5b, "Pourquoi cette différence ?",
                      0.40, 0.12, 9.20, 0.42, 20, NAVY, bold=True)

            # Bloc 1 — Diversification
            if _show_div:
                _add_rect(s5b, 0.40, 0.85, 9.20, 1.30, LGREY)
                _add_text(s5b, "Une meilleure répartition de vos placements",
                          0.55, 0.92, 8.90, 0.28, 13, NAVY, bold=True)
                if _n_red_A > 0:
                    _add_text(s5b,
                              f"Votre portefeuille actuel contient {_n_red_A} fonds qui se recoupent : "
                              f"vous pensez diversifier, mais une partie de votre argent est placée "
                              f"deux fois sur les mêmes marchés.",
                              0.55, 1.26, 8.90, 0.60, 12, GREY)
                else:
                    _add_text(s5b,
                              f"Votre portefeuille compte {_n_lin_A} lignes pour {_n_eff_A} expositions "
                              f"réellement distinctes. Notre allocation optimise cette répartition : "
                              f"{_n_lin_B} lignes, {_n_eff_B} zones de marché indépendantes.",
                              0.55, 1.26, 8.90, 0.60, 12, GREY)

            # Bloc 2 — Efficacité
            if _show_eff:
                _add_rect(s5b, 0.40, 2.30, 9.20, 1.30, LGREY)
                _add_text(s5b, "Un meilleur équilibre entre rendement et risque",
                          0.55, 2.37, 8.90, 0.28, 13, NAVY, bold=True)
                if _ratio_sharpe > 1:
                    _eff_body = (
                        f"Pour chaque euro de risque pris, notre allocation génère "
                        f"{_ratio_sharpe:.1f}× plus de rendement que votre portefeuille actuel. "
                        f"Concrètement, vous visez le même objectif en vous exposant moins aux à-coups des marchés."
                    )
                elif _reduction_sortino > 0:
                    _eff_body = (
                        f"Notre allocation réduit les pertes en cas de forte baisse de marché "
                        f"d'environ {abs(_reduction_sortino):.0f}%, tout en conservant un potentiel de gain comparable."
                    )
                else:
                    _eff_body = (
                        "Le profil rendement/risque de notre allocation est optimisé pour votre horizon "
                        "d'investissement, avec une volatilité maîtrisée."
                    )
                _add_text(s5b, _eff_body, 0.55, 2.71, 8.90, 0.72, 12, GREY)

            # Bloc 3 — Frais
            if _show_frais:
                _add_rect(s5b, 0.40, 3.75, 9.20, 1.10, LGREY)
                _add_text(s5b, "Des frais réduits qui s'accumulent sur la durée",
                          0.55, 3.82, 8.90, 0.28, 13, NAVY, bold=True)
                if _eco_15 > 0:
                    _frais_body = (
                        f"Vos frais annuels actuels : {_fee_A:.2f}%  —  Notre proposition : {_fee_B:.2f}%. "
                        f"Sur 15 ans, cette différence représente une économie estimée à {_fe(_eco_15)} "
                        f"(base 100 000 €, rendement brut 5%/an)."
                    )
                else:
                    _frais_body = (
                        f"Les niveaux de frais sont comparables ({_fee_A:.2f}% vs {_fee_B:.2f}%/an). "
                        f"L'avantage de notre allocation porte principalement sur la performance et la diversification."
                    )
                _add_text(s5b, _frais_body, 0.55, 4.12, 8.90, 0.60, 12, GREY)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Ce que ça change pour vous concrètement  (mode compare uniquement)
    # ══════════════════════════════════════════════════════════════════════
    if mode == "compare":
        s5 = prs.slides.add_slide(blank_layout)
        _add_rect(s5, 0, 0, 10, 5.625, LGREY)
        _left_stripe(s5)
        _footer(s5)

        _add_text(s5, "Ce que ça change pour vous concrètement",
                  0.30, 0.12, 9.30, 0.42, 20, NAVY, bold=True)
        _add_text(s5, "Projections comparées — simulation historique extrapolée",
                  0.30, 0.54, 9.30, 0.25, 9, GREY)

        _val_A_s5    = float(synthA.get("val") or 0)
        _val_B_s5    = float(synthB.get("val") or 0)
        _xirr_A_s5   = float(synthA.get("irr_pct") or 0) or 3.0
        _xirr_B_s5   = float(synthB.get("irr_pct") or 0) or 3.0
        _rachat_B_s5 = float(report.get("rachat_optimal_B") or 0)
        _net_B_s5    = float(synthB.get("net") or 0)
        _val_A_15s5  = _val_A_s5 * (1 + _xirr_A_s5 / 100) ** 15 if _val_A_s5 > 0 else 0
        _val_B_15s5  = _val_B_s5 * (1 + _xirr_B_s5 / 100) ** 15 if _val_B_s5 > 0 else 0
        _val_A_20s5  = _val_A_s5 * (1 + _xirr_A_s5 / 100) ** 20 if _val_A_s5 > 0 else 0
        _val_B_20s5  = _val_B_s5 * (1 + _xirr_B_s5 / 100) ** 20 if _val_B_s5 > 0 else 0
        _sit_s5      = report.get("situation_familiale", "")
        _abat_s5     = 9_200.0 if "Couple" in str(_sit_s5) else 4_600.0
        _taux_pv_B_s5 = max(0.0, 1.0 - (_net_B_s5 / max(_val_B_s5, 1)))
        _rachat_B_wir = (min(_val_B_s5, _abat_s5 / _taux_pv_B_s5)
                         if _taux_pv_B_s5 > 0 else _rachat_B_s5)

        # Graphique barres groupées : 3 horizons × 2 barres
        _groups5  = ["Aujourd'hui", "Dans 15 ans", "Dans 20 ans"]
        _vals_A5  = [_val_A_s5,   _val_A_15s5, _val_A_20s5]
        _vals_B5  = [_val_B_s5,   _val_B_15s5, _val_B_20s5]
        _fig5, _ax5 = plt.subplots(figsize=(9.0, 3.6), dpi=150)
        _fig5.patch.set_facecolor("#F5F7FA")
        _ax5.set_facecolor("#F5F7FA")
        _x5 = range(len(_groups5))
        _w5 = 0.35
        _bars_A5 = _ax5.bar([xi - _w5 / 2 for xi in _x5], _vals_A5, _w5,
                            color="#CC2200", label="Votre portefeuille")
        _bars_B5 = _ax5.bar([xi + _w5 / 2 for xi in _x5], _vals_B5, _w5,
                            color="#1A7A4A", label="Notre proposition")
        # Montants au-dessus des barres
        for _bar5 in list(_bars_A5) + list(_bars_B5):
            _h5 = _bar5.get_height()
            if _h5 > 0:
                _ax5.text(
                    _bar5.get_x() + _bar5.get_width() / 2, _h5 * 1.02,
                    _fe(_h5), ha="center", va="bottom", fontsize=6.5,
                    color="#333333", fontweight="bold"
                )
        # Delta en vert entre les paires
        for _gi5 in _x5:
            _dlt5 = _vals_B5[_gi5] - _vals_A5[_gi5]
            if _dlt5 != 0:
                _mid5 = (_vals_A5[_gi5] + _vals_B5[_gi5]) / 2
                _sign5 = "+" if _dlt5 >= 0 else ""
                _ax5.text(
                    _gi5, _mid5,
                    f"{_sign5}{_fe(_dlt5)}", ha="center", va="center",
                    fontsize=7, color="#1A7A4A", fontweight="bold",
                    bbox=dict(boxstyle="round,pad=0.15", fc="white", ec="#1A7A4A", lw=0.6)
                )
        _ax5.set_xticks(list(_x5))
        _ax5.set_xticklabels(_groups5, fontsize=9)
        _ax5.yaxis.set_major_formatter(
            plt.FuncFormatter(lambda v, _: f"{v/1e3:.0f} k€" if v >= 1000 else f"{v:.0f} €")
        )
        _ax5.tick_params(axis="y", labelsize=8)
        _ax5.spines["top"].set_visible(False)
        _ax5.spines["right"].set_visible(False)
        _ax5.legend(fontsize=8, loc="upper left", frameon=False)
        _ax5.set_ylabel("")
        _fig5.tight_layout(pad=0.4)
        _buf5 = BytesIO()
        _fig5.savefig(_buf5, format="png", bbox_inches="tight")
        plt.close(_fig5)
        _buf5.seek(0)
        s5.shapes.add_picture(_buf5, left=Inches(0.40), top=Inches(0.82),
                               width=Inches(9.20), height=Inches(3.60))
        _add_text(s5,
                  "Les performances passées ne préjugent pas des performances futures — simulation à titre indicatif.",
                  0.30, 4.62, 9.40, 0.22, 7, GREY, italic=True)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Pourquoi l'assurance-vie ?
    # ══════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank_layout)
    _add_rect(s6, 0, 0, 10, 5.625, LGREY)
    _left_stripe(s6, GOLD)
    _footer(s6)

    _add_text(s6, "Pourquoi l'assurance-vie ?",
              0.30, 0.12, 9.30, 0.42, 20, NAVY, bold=True)
    _add_text(s6, "Comparatif fiscal 2026 — AV (≥8 ans) vs CTO vs PEA",
              0.30, 0.54, 9.30, 0.25, 9, GREY)

    # Tableau comparatif AV / CTO / PEA
    _tbl6_headers = ["Critère", "AV (≥ 8 ans)", "CTO", "PEA"]
    _tbl6_rows = [
        ("IR à la sortie",       "7,5% (après abatt.)",        "12,8% flat tax",     "0% (après 5 ans)"),
        ("Prélèvements sociaux", "17,2%",                       "18,6%",              "18,6%"),
        ("Abattement annuel",    f"{_fe(4_600)}/pers.",         "Aucun",              "Aucun"),
        ("Transmission",         "Hors succession (art. 990I)", "Succession classique", "Succession classique"),
        ("Plafond versements",   "Illimité",                    "Illimité",           f"{_fe(150_000)}"),
    ]
    _tbl6_col_w = [2.60, 2.20, 2.20, 2.20]
    _tbl6_col_x = [0.40, 3.00, 5.20, 7.40]
    _tbl6_y = 0.92
    for ci6, (_cx6, _ch6) in enumerate(zip(_tbl6_col_x, _tbl6_headers)):
        _bg6h = NAVY if ci6 == 0 else (GOLD if ci6 == 1 else RGBColor(0x44, 0x55, 0x77))
        _add_rect(s6, _cx6, _tbl6_y, _tbl6_col_w[ci6], 0.30, _bg6h)
        _add_text(s6, _ch6, _cx6 + 0.06, _tbl6_y + 0.05, _tbl6_col_w[ci6] - 0.12, 0.22, 9, WHITE, bold=True)
    for ri6, _row_vals6 in enumerate(_tbl6_rows):
        _ry6 = _tbl6_y + 0.30 + ri6 * 0.38
        _row_bg6 = WHITE if ri6 % 2 == 0 else LGREY
        _add_rect(s6, 0.40, _ry6, 9.20, 0.38, _row_bg6)
        for ci6b, (_cx6b, _rv6) in enumerate(zip(_tbl6_col_x, _row_vals6)):
            _tc6 = NAVY if ci6b == 0 else (GREEN if ci6b == 1 else BLACK)
            _add_text(s6, _rv6, _cx6b + 0.06, _ry6 + 0.08, _tbl6_col_w[ci6b] - 0.12, 0.26, 8, _tc6,
                      bold=(ci6b == 0))

    _add_text(s6,
              "L'assurance-vie reste l'enveloppe offrant le meilleur équilibre entre rendement, "
              "fiscalité à la sortie et transmission patrimoniale.",
              0.28, 3.76, 9.44, 0.40, 9, GREY)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Votre conseiller à vos côtés
    # ══════════════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(blank_layout)
    _add_rect(s7, 0, 0, 10, 5.625, NAVY)
    _add_rect(s7, 0, 0, 0.18, 5.625, GOLD)

    _add_text(s7, "Votre conseiller\nà vos côtés",
              0.75, 0.45, 5.50, 1.10, 28, WHITE, bold=True)
    _add_rect(s7, 0.75, 1.70, 2.00, 0.04, GOLD)

    _engagements7 = [
        ("Transparence",
         "Simulations basées sur les frais réels de votre contrat."),
        ("Indépendance",
         "Sélection de fonds fondée sur la performance."),
        ("Suivi",
         "Point annuel sur l'évolution de votre allocation."),
    ]
    for ei7, (etitle7, etxt7) in enumerate(_engagements7):
        _ey7 = 1.90 + ei7 * 0.95
        _add_rect(s7, 0.75, _ey7, 0.05, 0.60, GOLD)
        _add_text(s7, etitle7, 0.90, _ey7 + 0.04, 4.50, 0.28, 12, WHITE, bold=True)
        _add_text(s7, etxt7,   0.90, _ey7 + 0.32, 4.50, 0.28, 9, ICE)

    _add_text(s7, "  ·  ".join(filter(None, [nom_cab, "Document confidentiel", nom_cli, as_of])),
              0.28, 5.38, 9.44, 0.20, 7, GREY, align="center")

    # ── Sérialisation ─────────────────────────────────────────────────────
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()




def render_app(run_page_config: bool = True):
    # ------------------------------------------------------------
    # Layout principal
    # ------------------------------------------------------------
    if run_page_config:
        st.set_page_config(page_title=APP_TITLE, layout="wide")
    _nom_cab_display = st.session_state.get("NOM_CABINET", "").strip()
    st.title(_nom_cab_display if _nom_cab_display else APP_TITLE)
    st.caption(APP_SUBTITLE)
    _nom = st.session_state.get("NOM_CLIENT", "").strip()
    if _nom:
        st.markdown(f"**Dossier client :** {_nom}")
    # Init state
    st.session_state.setdefault("NOM_CLIENT", "")
    st.session_state.setdefault("NOM_CABINET", "")
    st.session_state.setdefault("A_lines", [])
    st.session_state.setdefault("B_lines", [])
    st.session_state.setdefault("ALLOC_MODE", "equal")
    st.session_state.setdefault("MANUAL_NAV_STORE", {})
    st.session_state.setdefault("DATE_WARNINGS", [])
    # ── PERSIST : restaure les clés widget supprimées lors de la navigation ──
    _COMP_WIDGET_DEFAULTS = {
        "FEE_A": 3.0, "FEE_B": 2.0,
        "EURO_RATE_A": 2.0, "EURO_RATE_B": 2.5,
        "INIT_A_DATE": pd.Timestamp("2024-01-02").date(),
        "INIT_B_DATE": pd.Timestamp("2024-01-02").date(),
        "M_A": 0.0, "M_B": 0.0,
        "ONE_A": 0.0, "ONE_B": 0.0,
        "ONE_A_DATE": pd.Timestamp("2024-07-01").date(),
        "ONE_B_DATE": pd.Timestamp("2024-07-01").date(),
    }
    for _k, _default in _COMP_WIDGET_DEFAULTS.items():
        if _k not in st.session_state:
            st.session_state[_k] = st.session_state.get(f"{_k}_PERSIST", _default)
        else:
            st.session_state.setdefault(f"{_k}_PERSIST", st.session_state[_k])

    # -------------------------------------------------------------------
    # Sidebar : paramètres globaux
    # -------------------------------------------------------------------
    with st.sidebar:
        st.divider()

        # ── Section 1 — Contrats ────────────────────────────────────
        with st.expander("① Contrats", expanded=True):
            st.markdown("**Portefeuille Client**")
            contract_label_A = st.selectbox(
                "Contrat client",
                list(CONTRACTS_REGISTRY.keys()),
                key="CONTRACT_LABEL_A",
            )
            contract_cfg_A = CONTRACTS_REGISTRY[contract_label_A]
            try:
                funds_df_A = load_contract_funds(
                    contract_cfg_A["path"], contract_cfg_A["funds_filename"]
                )
            except Exception:
                funds_df_A = pd.DataFrame()
            st.session_state["CONTRACT_FUNDS_DF_A"] = funds_df_A

            euro_fund_label_A = st.selectbox(
                "Fonds en euros client",
                list(contract_cfg_A["euro_funds"].keys()),
                key="EURO_FUND_LABEL_A",
            )
            try:
                euro_history_df_A = load_euro_fund_history(
                    contract_cfg_A["path"], contract_cfg_A["euro_funds"][euro_fund_label_A]
                )
            except Exception:
                euro_history_df_A = pd.DataFrame()
            avg_rate_A = get_euro_fund_avg_rate(euro_history_df_A, years=5)
            st.caption(
                f"Fonds euros client — moyenne 5 ans : **{avg_rate_A:.2f}%** "
                f"({euro_fund_label_A})"
            )

            st.markdown("---")
            st.markdown("**Portefeuille Cabinet**")
            contract_label_B = st.selectbox(
                "Contrat cabinet",
                list(CONTRACTS_REGISTRY.keys()),
                key="CONTRACT_LABEL_B",
            )
            contract_cfg_B = CONTRACTS_REGISTRY[contract_label_B]
            try:
                funds_df_B = load_contract_funds(
                    contract_cfg_B["path"], contract_cfg_B["funds_filename"]
                )
            except Exception:
                funds_df_B = pd.DataFrame()
            st.session_state["CONTRACT_FUNDS_DF_B"] = funds_df_B

            euro_fund_label_B = st.selectbox(
                "Fonds en euros cabinet",
                list(contract_cfg_B["euro_funds"].keys()),
                key="EURO_FUND_LABEL_B",
            )
            try:
                euro_history_df_B = load_euro_fund_history(
                    contract_cfg_B["path"], contract_cfg_B["euro_funds"][euro_fund_label_B]
                )
            except Exception:
                euro_history_df_B = pd.DataFrame()
            avg_rate_B = get_euro_fund_avg_rate(euro_history_df_B, years=5)
            st.caption(
                f"Fonds euros cabinet — moyenne 5 ans : **{avg_rate_B:.2f}%** "
                f"({euro_fund_label_B})"
            )

            # Rétrocompatibilité : render_portfolio_builder() lit CONTRACT_LABEL / CONTRACT_FUNDS_DF
            st.session_state["CONTRACT_LABEL"]    = contract_label_A
            st.session_state["CONTRACT_CFG"]      = contract_cfg_A
            st.session_state["CONTRACT_FUNDS_DF"] = funds_df_A
            st.session_state["EURO_FUND_LABEL"]   = euro_fund_label_A
            st.session_state["EURO_FUND_HISTORY"] = euro_history_df_A
            st.session_state["EURO_FUND_AVG_RATE"] = avg_rate_A

        # ── Section 2 — Paramètres de simulation ───────────────────
        with st.expander("② Paramètres de simulation", expanded=True):
            st.markdown("**Taux fonds en euros**")
            st.caption(
                f"Client : **{avg_rate_A:.2f}%** ({euro_fund_label_A}) | "
                f"Cabinet : **{avg_rate_B:.2f}%** ({euro_fund_label_B})"
            )
            override_euro = st.checkbox(
                "Utiliser un taux personnalisé (fonds euros boosté, etc.)",
                value=False,
                key="OVERRIDE_EURO_RATE",
            )
            if override_euro:
                euro_rate_value_A = st.number_input(
                    "Portefeuille Client — taux annuel (%)",
                    min_value=0.0, max_value=10.0,
                    value=float(st.session_state.get("EURO_RATE_A", avg_rate_A)),
                    step=0.10, key="EURO_RATE_A",
                )
                euro_rate_value_B = st.number_input(
                    "Portefeuille Cabinet — taux annuel (%)",
                    min_value=0.0, max_value=10.0,
                    value=float(st.session_state.get("EURO_RATE_B", avg_rate_B)),
                    step=0.10, key="EURO_RATE_B",
                )
            else:
                euro_rate_value_A = avg_rate_A
                euro_rate_value_B = avg_rate_B
                st.session_state["EURO_RATE_A"] = avg_rate_A
                st.session_state["EURO_RATE_B"] = avg_rate_B
                if avg_rate_A == avg_rate_B:
                    st.caption(f"Taux appliqué aux deux portefeuilles : **{avg_rate_A:.2f}%**")
                else:
                    st.caption(
                        f"Taux client : **{avg_rate_A:.2f}%** | Taux cabinet : **{avg_rate_B:.2f}%**"
                    )
            st.caption(
                "Frais de gestion contrat fonds euros inclus dans le taux net publié. "
                "Frais UC contrat : lus depuis le référentiel fonds."
            )

            st.markdown("---")
            st.markdown("**Frais d'entrée (%)**")
            FEE_A = st.number_input(
                "Frais d'entrée — Portefeuille 1 (Client)",
                0.0,
                10.0,
                st.session_state.get("FEE_A", 3.0),
                0.10,
                key="FEE_A",
            )
            FEE_B = st.number_input(
                "Frais d'entrée — Portefeuille 2 (Cabinet)",
                0.0,
                10.0,
                st.session_state.get("FEE_B", 2.0),
                0.10,
                key="FEE_B",
            )
            st.caption("Les frais s'appliquent sur chaque investissement (initial, mensuel, ponctuel).")

            st.markdown("---")
            st.markdown("**Date du versement initial**")
            st.date_input(
                "Portefeuille 1 (Client) — date d'investissement initiale",
                value=st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date()),
                key="INIT_A_DATE",
            )
            st.date_input(
                "Portefeuille 2 (Cabinet) — date d'investissement initiale",
                value=st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date()),
                key="INIT_B_DATE",
            )

        # ── Section 3 — Versements ──────────────────────────────────
        with st.expander("③ Versements", expanded=False):
            st.markdown("**Portefeuille 1 — Client**")
            M_A = st.number_input(
                "Mensuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("M_A", 0.0),
                100.0,
                key="M_A",
            )
            ONE_A = st.number_input(
                "Ponctuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("ONE_A", 0.0),
                100.0,
                key="ONE_A",
            )
            ONE_A_DATE = st.date_input(
                "Date du ponctuel",
                value=st.session_state.get("ONE_A_DATE", pd.Timestamp("2024-07-01").date()),
                key="ONE_A_DATE",
            )
            st.markdown("---")
            st.markdown("**Portefeuille 2 — Cabinet**")
            M_B = st.number_input(
                "Mensuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("M_B", 0.0),
                100.0,
                key="M_B",
            )
            ONE_B = st.number_input(
                "Ponctuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("ONE_B", 0.0),
                100.0,
                key="ONE_B",
            )
            ONE_B_DATE = st.date_input(
                "Date du ponctuel",
                value=st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date()),
                key="ONE_B_DATE",
            )

        # ── Section 4 — Options d'analyse ──────────────────────────
        with st.expander("④ Options d'analyse", expanded=False):
            st.markdown("**Règle d'affectation des versements**")
            current_code = st.session_state.get("ALLOC_MODE", "equal")
            inv_labels = {v: k for k, v in ALLOC_LABELS.items()}
            current_label = inv_labels.get(current_code, "Répartition égale")

            mode_label = st.selectbox(
                "Mode",
                list(ALLOC_LABELS.keys()),
                index=list(ALLOC_LABELS.keys()).index(current_label),
                help="Répartition des versements entre les lignes.",
            )
            st.session_state["ALLOC_MODE"] = ALLOC_LABELS[mode_label]

            st.markdown("---")
            st.markdown("**Mode d'analyse**")
            mode_ui = st.radio(
                "Choix",
                ["Comparer Client vs Cabinet", "Analyser uniquement Cabinet", "Analyser uniquement Client"],
                index=0,
                key="MODE_ANALYSE_UI",
            )
            if "Comparer" in mode_ui:
                st.session_state["MODE_ANALYSE"] = "compare"
            elif "Cabinet" in mode_ui:
                st.session_state["MODE_ANALYSE"] = "valority"
            else:
                st.session_state["MODE_ANALYSE"] = "client"

            st.markdown("---")
            st.markdown("**Indice de référence (Bêta)**")
            _bench_options = {
                "MSCI World (CW8.PA)": "CW8.PA",
                "CAC 40 (^FCHI)": "^FCHI",
                "Euro Stoxx 50 (^STOXX50E)": "^STOXX50E",
                "S&P 500 (^GSPC)": "^GSPC",
            }
            _bench_label = st.selectbox(
                "Indice",
                list(_bench_options.keys()),
                index=0,
                key="BENCHMARK_LABEL",
            )
            st.session_state["BENCHMARK_SYMBOL"] = _bench_options[_bench_label]

        # ── Debug (masqué en production) ────────────────────────────
        _debug_enabled = st.secrets.get("debug_mode", False)
        if _debug_enabled:
            with st.expander("🔧 Debug", expanded=False):
                st.caption("Versions & état")
                st.code(
                    f"Python: {sys.version.split()[0]}\n"
                    f"Streamlit: {st.__version__}\n"
                    f"Pandas: {pd.__version__}"
                )
                st.caption("Modules")
                st.code(
                    f"Matplotlib: {MATPLOTLIB_AVAILABLE} ({MATPLOTLIB_ERROR})\n"
                    f"Reportlab: {REPORTLAB_AVAILABLE} ({REPORTLAB_ERROR})"
                )
                st.caption("Session state (clés)")
                st.write(sorted(list(st.session_state.keys())))
                st.caption("Dernière exception")
                st.write(st.session_state.get("LAST_EXCEPTION", "—"))
                st.caption("Test rapide EODHD")
                if _get_api_key():
                    try:
                        res = eodhd_get("/status")
                        st.write(res if res is not None else "Réponse vide")
                    except Exception as e:
                        st.write(f"Erreur EODHD: {e}")
                else:
                    st.write("Token EODHD absent")


    # ── PERSIST save : sauvegarde les clés widget pour survie à la navigation ──
    for _k in _COMP_WIDGET_DEFAULTS:
        if _k in st.session_state:
            st.session_state[f"{_k}_PERSIST"] = st.session_state[_k]

    mode = st.session_state.get("MODE_ANALYSE", "compare")
    show_client = mode in ("compare", "client")
    show_valority = mode in ("compare", "valority")

    # ── Synchronisation date globale → lignes (pour cohérence carte + tableau + simulation) ──
    _global_date_A = pd.Timestamp(st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date()))
    _global_date_B = pd.Timestamp(st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date()))
    for _ln in st.session_state.get("A_lines", []):
        if not _ln.get("date_overridden"):
            _ln["buy_date"] = _global_date_A
    for _ln in st.session_state.get("B_lines", []):
        if not _ln.get("date_overridden"):
            _ln["buy_date"] = _global_date_B

    # Onglets principaux : Client / Cabinet (conditionnés au mode)
    tab_labels = []
    if show_client:
        tab_labels.append("Portefeuille Client")
    if show_valority:
        tab_labels.append("Portefeuille Cabinet")
    tabs = st.tabs(tab_labels) if tab_labels else []

    idx = 0
    if show_client:
        with tabs[idx]:
            _add_fund_from_contract("A_lines", "Ajouter un fonds — Portefeuille Client")
            st.markdown("#### Lignes actuelles — Portefeuille Client")
            for i, ln in enumerate(st.session_state.get("A_lines", [])):
                _line_card(ln, i, "A_lines")
        idx += 1

    if show_valority:
        with tabs[idx]:
            _add_fund_from_contract("B_lines", "Ajouter un fonds — Portefeuille Cabinet")
            st.markdown("#### Lignes actuelles — Portefeuille Cabinet")
            for i, ln in enumerate(st.session_state.get("B_lines", [])):
                _line_card(ln, i, "B_lines")

    # ------------------------------------------------------------
    # Simulation (selon mode)
    # ------------------------------------------------------------
    # FIXED: use stable UUID key instead of id() which changes on st.rerun() (Bug 5)
    _ln_a0 = st.session_state["A_lines"][0] if (show_client and st.session_state["A_lines"]) else None
    single_target_A = (_ln_a0.get("id") or id(_ln_a0)) if _ln_a0 is not None else None
    _ln_b0 = st.session_state["B_lines"][0] if (show_valority and st.session_state["B_lines"]) else None
    single_target_B = (_ln_b0.get("id") or id(_ln_b0)) if _ln_b0 is not None else None

    alloc_mode_code = st.session_state.get("ALLOC_MODE", "equal")

    custom_month_weights_A: Optional[Dict[int, float]] = None
    custom_oneoff_weights_A: Optional[Dict[int, float]] = None
    custom_month_weights_B: Optional[Dict[int, float]] = None
    custom_oneoff_weights_B: Optional[Dict[int, float]] = None

    if alloc_mode_code == "custom":
        if show_client:
            cmA = st.session_state.get("CUSTOM_M_A", {}) or {}
            coA = st.session_state.get("CUSTOM_O_A", {}) or {}
            tot_mA = sum(v for v in cmA.values() if v > 0)
            tot_oA = sum(v for v in coA.values() if v > 0)
            if tot_mA > 0:
                custom_month_weights_A = {k: v / tot_mA for k, v in cmA.items() if v > 0}
            if tot_oA > 0:
                custom_oneoff_weights_A = {k: v / tot_oA for k, v in coA.items() if v > 0}

        if show_valority:
            cmB = st.session_state.get("CUSTOM_M_B", {}) or {}
            coB = st.session_state.get("CUSTOM_O_B", {}) or {}
            tot_mB = sum(v for v in cmB.values() if v > 0)
            tot_oB = sum(v for v in coB.values() if v > 0)
            if tot_mB > 0:
                custom_month_weights_B = {k: v / tot_mB for k, v in cmB.items() if v > 0}
            if tot_oB > 0:
                custom_oneoff_weights_B = {k: v / tot_oB for k, v in coB.items() if v > 0}

    # Reset warnings avant chaque run
    st.session_state["DATE_WARNINGS"] = []

    # Valeurs par défaut (si on ne simule pas un des portefeuilles)
    dfA, brutA, netA, valA, xirrA, startA_min, fullA = pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY
    dfB, brutB, netB, valB, xirrB, startB_min, fullB = pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    if show_client:
        _args_A = _make_sim_args(
            lines=st.session_state.get("A_lines", []),
            monthly_amt=st.session_state.get("M_A", 0.0),
            one_amt=st.session_state.get("ONE_A", 0.0),
            one_date=st.session_state.get("ONE_A_DATE", pd.Timestamp("2024-07-01").date()),
            alloc_mode=alloc_mode_code,
            custom_weights_monthly=custom_month_weights_A,
            custom_weights_oneoff=custom_oneoff_weights_A,
            single_target=single_target_A,
            euro_rate=st.session_state.get("EURO_RATE_A", 2.0),
            fee_pct=st.session_state.get("FEE_A", 0.0),
            label="Client",
        )
        dfA, brutA, netA, valA, xirrA, startA_min, fullA = _simulate_portfolio_cached(**_args_A)
        st.session_state["_LAST_VAL_A"] = valA
        st.session_state["_LAST_NET_A"] = netA
        st.session_state["_LAST_XIRR_A"] = xirrA

    if show_valority:
        _args_B = _make_sim_args(
            lines=st.session_state.get("B_lines", []),
            monthly_amt=st.session_state.get("M_B", 0.0),
            one_amt=st.session_state.get("ONE_B", 0.0),
            one_date=st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date()),
            alloc_mode=alloc_mode_code,
            custom_weights_monthly=custom_month_weights_B,
            custom_weights_oneoff=custom_oneoff_weights_B,
            single_target=single_target_B,
            euro_rate=st.session_state.get("EURO_RATE_B", 2.5),
            fee_pct=st.session_state.get("FEE_B", 0.0),
            label="Cabinet",
        )
        dfB, brutB, netB, valB, xirrB, startB_min, fullB = _simulate_portfolio_cached(**_args_B)
        st.session_state["_LAST_VAL_B"] = valB
        st.session_state["_LAST_NET_B"] = netB
        st.session_state["_LAST_XIRR_B"] = xirrB

    # ------------------------------------------------------------
    # Avertissements sur les dates / 1ère VL
    # ------------------------------------------------------------
    if st.session_state.get("DATE_WARNINGS"):
        with st.expander("⚠️ Problèmes d'historique / dates de VL"):
            for msg in st.session_state["DATE_WARNINGS"]:
                st.warning(msg)

    # ------------------------------------------------------------
    # Graphique (évolution des portefeuilles)
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📈 Comment s'est comporté votre portefeuille")

    # Déterminer le start_plot uniquement sur les portefeuilles affichés
    full_dates: List[pd.Timestamp] = []
    if show_client and isinstance(fullA, pd.Timestamp):
        full_dates.append(fullA)
    if show_valority and isinstance(fullB, pd.Timestamp):
        full_dates.append(fullB)

    start_plot = max(full_dates) if full_dates else TODAY

    idx = pd.bdate_range(start=start_plot, end=TODAY, freq="B")
    chart_df = pd.DataFrame(index=idx)

    if show_client and not dfA.empty:
        chart_df["Client"] = dfA.reindex(idx)["Valeur"].ffill()

    if show_valority and not dfB.empty:
        chart_df["Cabinet"] = dfB.reindex(idx)["Valeur"].ffill()

    # Passage en format long pour Altair
    chart_long = chart_df.reset_index().rename(columns={"index": "Date"})
    chart_long = chart_long.melt("Date", var_name="variable", value_name="Valeur (€)")

    if chart_long.dropna().empty:
        st.info("Ajoutez des lignes et/ou vérifiez vos paramètres pour afficher le graphique.")
    else:
        # Calcul du domaine Y dynamique avec marges
        valid_vals = chart_long["Valeur (€)"].dropna()
        if not valid_vals.empty:
            y_min = float(valid_vals.min())
            y_max = float(valid_vals.max())
            padding = (y_max - y_min) * 0.05 if y_max > y_min else y_max * 0.05
            y_domain = [max(0.0, y_min - padding), y_max + padding]
        else:
            y_domain = [0, 1]

        # ── Bandeau KPIs adapté au mode ──────────────────────────────
        with st.container(border=True):
            if mode == "compare":
                _has_A = show_client and not dfA.empty and netA > 0
                _has_B = show_valority and not dfB.empty and netB > 0
                if _has_A and _has_B:
                    _kpi_cols = st.columns(4)
                    with _kpi_cols[0]:
                        st.metric(
                            "Rendement client (%/an)",
                            f"{xirrA:.2f} %" if xirrA is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _kpi_cols[1]:
                        st.metric(
                            "Rendement simulation (%/an)",
                            f"{xirrB:.2f} %" if xirrB is not None else "—",
                            delta=f"{(xirrB - xirrA):.2f} pts"
                            if (xirrA is not None and xirrB is not None) else None,
                        )
                    with _kpi_cols[2]:
                        _gain_client = (valA - netA) if netA > 0 else 0.0
                        st.metric("Gain net Client", to_eur(_gain_client))
                    with _kpi_cols[3]:
                        _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                        st.metric(
                            "Gain net Cabinet", to_eur(_gain_cabinet),
                            delta=to_eur(_gain_cabinet - _gain_client) if netA > 0 else None,
                        )
                elif _has_A:
                    _pad_l, _k0, _k1, _k2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                    with _k0:
                        st.metric(
                            "Rendement client (%/an)",
                            f"{xirrA:.2f} %" if xirrA is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _k1:
                        _gain_client = (valA - netA) if netA > 0 else 0.0
                        st.metric("Gain net", to_eur(_gain_client))
                    with _k2:
                        _perf_cli = (valA / netA - 1.0) * 100.0 if netA > 0 else 0.0
                        st.metric("Performance totale", f"{_perf_cli:+.2f}%")
                elif _has_B:
                    _pad_l, _k0, _k1, _k2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                    with _k0:
                        st.metric(
                            "Rendement simulation (%/an)",
                            f"{xirrB:.2f} %" if xirrB is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _k1:
                        _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                        st.metric("Gain net", to_eur(_gain_cabinet))
                    with _k2:
                        _perf_val = (valB / netB - 1.0) * 100.0 if netB > 0 else 0.0
                        st.metric("Performance totale", f"{_perf_val:+.2f}%")
            elif mode == "client":
                _pad_l, _kpi_cols0, _kpi_cols1, _kpi_cols2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                _kpi_cols = [_kpi_cols0, _kpi_cols1, _kpi_cols2]
                with _kpi_cols[0]:
                    st.metric(
                        "Rendement client (%/an)",
                        f"{xirrA:.2f} %" if xirrA is not None else "—",
                        help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                    )
                with _kpi_cols[1]:
                    _gain_client = (valA - netA) if netA > 0 else 0.0
                    st.metric("Gain net", to_eur(_gain_client))
                with _kpi_cols[2]:
                    _perf_cli = (valA / netA - 1.0) * 100.0 if netA > 0 else 0.0
                    st.metric("Performance totale", f"{_perf_cli:+.2f}%")
            else:  # valority
                _pad_l, _kpi_cols0, _kpi_cols1, _kpi_cols2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                _kpi_cols = [_kpi_cols0, _kpi_cols1, _kpi_cols2]
                with _kpi_cols[0]:
                    st.metric(
                        "Rendement simulation (%/an)",
                        f"{xirrB:.2f} %" if xirrB is not None else "—",
                        help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                    )
                with _kpi_cols[1]:
                    _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                    st.metric("Gain net", to_eur(_gain_cabinet))
                with _kpi_cols[2]:
                    _perf_val = (valB / netB - 1.0) * 100.0 if netB > 0 else 0.0
                    st.metric("Performance totale", f"{_perf_val:+.2f}%")

        # ── Graphique line chart (rendu direct, sans placeholder) ────
        base = (
            alt.Chart(chart_long)
            .mark_line()
            .encode(
                x=alt.X("Date:T", title="Date"),
                y=alt.Y(
                    "Valeur (€):Q",
                    title="Valeur (€)",
                    scale=alt.Scale(domain=y_domain),
                    axis=alt.Axis(format=",.0f"),
                ),
                color=alt.Color("variable:N", title="Portefeuille"),
                tooltip=[
                    alt.Tooltip("Date:T", title="Date"),
                    alt.Tooltip("variable:N", title="Portefeuille"),
                    alt.Tooltip("Valeur (€):Q", title="Valeur", format=",.2f"),
                ],
            )
            .properties(height=360, width="container")
        )
        st.altair_chart(base, use_container_width=True)

        # ── Bloc storytelling percutant ──────────────────────────────
        if mode == "compare" and netA > 0 and netB > 0:
            _perf_a = (valA / netA - 1.0) * 100.0
            _perf_b = (valB / netB - 1.0) * 100.0
            _gain_a = valA - netA
            _gain_b = valB - netB
            _manque = _gain_b - _gain_a
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric(
                        "Votre rendement sur la période",
                        f"{_perf_a:+.2f}%",
                        help="Performance totale nette des frais contrat",
                    )
                    st.metric(
                        "Vous avez gagné",
                        to_eur(_gain_a),
                        help="Valeur actuelle − versements nets investis",
                    )
                with _st_c2:
                    st.metric(
                        "Rendement de la proposition cabinet",
                        f"{_perf_b:+.2f}%",
                        delta=f"{(_perf_b - _perf_a):+.2f} pts vs votre situation actuelle",
                    )
                    st.metric(
                        "Avec le cabinet, vous auriez gagné",
                        to_eur(_gain_b),
                        delta=to_eur(_manque),
                        help="Manque à gagner par rapport à la proposition",
                    )
                if _manque > 0:
                    st.info(
                        f"💡 La proposition du cabinet aurait généré "
                        f"**{to_eur(_manque)} de plus** sur la même période "
                        f"et le même capital investi."
                    )
        elif mode == "client" and netA > 0:
            _perf_a = (valA / netA - 1.0) * 100.0
            _gain_a = valA - netA
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric("Votre rendement sur la période",
                               f"{_perf_a:+.2f}%")
                with _st_c2:
                    st.metric("Vous avez gagné", to_eur(_gain_a))
        elif mode == "valority" and netB > 0:
            _perf_b = (valB / netB - 1.0) * 100.0
            _gain_b = valB - netB
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric("Rendement de la simulation",
                               f"{_perf_b:+.2f}%")
                with _st_c2:
                    st.metric("Capital créé", to_eur(_gain_b))

    # ------------------------------------------------------------
    # Synthèse chiffrée : cartes Client / Cabinet
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("💡 Synthèse")

    mode = st.session_state.get("MODE_ANALYSE", "compare")

    # Période analysée (uniquement sur ce qui est affiché)
    period_dates: List[pd.Timestamp] = []
    if mode in ("compare", "client") and isinstance(startA_min, pd.Timestamp):
        period_dates.append(startA_min)
    if mode in ("compare", "valority") and isinstance(startB_min, pd.Timestamp):
        period_dates.append(startB_min)

    if period_dates:
        start_global = min(period_dates)
        st.caption(f"Période analysée : du **{fmt_date(start_global)}** au **{fmt_date(TODAY)}**")

    perf_tot_client = (valA / netA - 1.0) * 100.0 if (show_client and netA > 0) else None
    perf_tot_valority = (valB / netB - 1.0) * 100.0 if (show_valority and netB > 0) else None

    # ✅ 2 colonnes si compare, sinon 1 colonne (container)
    if mode == "compare":
        col_client, col_valority = st.columns(2)
    else:
        col_client = st.container()
        col_valority = st.container()

    # ----- Carte Client -----
    if mode in ("compare", "client"):
        with col_client:
            with st.container(border=True):
                st.markdown("#### 🧍 Situation actuelle — Client")
                st.metric("Valeur actuelle", to_eur(valA))
                if _nom := st.session_state.get("NOM_CLIENT", "").strip():
                    st.caption(f"Client : {_nom}")
                _s_c1, _s_c2 = st.columns(2)
                with _s_c1:
                    st.metric("Net investi", to_eur(netA),
                              help="Versements bruts − frais d'entrée")
                with _s_c2:
                    st.metric("Performance totale",
                              f"{perf_tot_client:+.2f}%" if perf_tot_client is not None else "—")
                _fee_detail_A = []
                for _ln in st.session_state.get("A_lines", []):
                    _isin = str(_ln.get("isin", "")).upper()
                    if _isin in ("EUROFUND", "STRUCTURED"):
                        continue
                    _name = _ln.get("name") or _isin
                    _fc = _ln.get("fee_contract_pct")
                    _fb = _ln.get("fee_uc_pct")
                    if _fc is not None and _fb is not None:
                        try:
                            _fee_detail_A.append(
                                f"{_name[:30]} : frais contrat {float(_fc):.2f}%/an "
                                f"(+ TER {float(_fb):.2f}%/an intégré dans la VL)"
                            )
                        except Exception:
                            pass
                _fee_note_A = (
                    "ℹ️ **Performance nette des frais de gestion du contrat assureur** "
                    f"({st.session_state.get('CONTRACT_LABEL_A', 'contrat')}). "
                    "Les frais de gestion internes des fonds (TER) sont déjà intégrés "
                    "dans les valeurs liquidatives publiées et ne sont pas déduits en plus. "
                    "À titre de comparaison, Quantalys et Morningstar affichent les "
                    "performances brutes de frais contrat."
                )
                with st.expander("ℹ️ Détail du calcul de performance", expanded=False):
                    st.markdown(_fee_note_A)
                    if _fee_detail_A:
                        st.markdown("**Frais appliqués par fonds :**")
                        for _d in _fee_detail_A:
                            st.markdown(f"- {_d}")
                    st.markdown(
                        "**Formule :** XIRR calculé sur les flux réels "
                        "(versements initiaux, mensuels, ponctuels) et la valeur "
                        "liquidative nette à ce jour."
                    )


    # ----- Carte Cabinet -----
    if mode in ("compare", "valority"):
        with col_valority:
            with st.container(border=True):
                st.markdown("#### 🏢 Simulation — Allocation Cabinet")
                st.metric("Valeur actuelle simulée", to_eur(valB))
                _s_c1b, _s_c2b = st.columns(2)
                with _s_c1b:
                    st.metric("Net investi", to_eur(netB),
                              help="Versements bruts − frais d'entrée")
                with _s_c2b:
                    st.metric("Performance totale",
                              f"{perf_tot_valority:+.2f}%" if perf_tot_valority is not None else "—")
                _fee_detail_B = []
                for _ln in st.session_state.get("B_lines", []):
                    _isin = str(_ln.get("isin", "")).upper()
                    if _isin in ("EUROFUND", "STRUCTURED"):
                        continue
                    _name = _ln.get("name") or _isin
                    _fc = _ln.get("fee_contract_pct")
                    _fb = _ln.get("fee_uc_pct")
                    if _fc is not None and _fb is not None:
                        try:
                            _fee_detail_B.append(
                                f"{_name[:30]} : frais contrat {float(_fc):.2f}%/an "
                                f"(+ TER {float(_fb):.2f}%/an intégré dans la VL)"
                            )
                        except Exception:
                            pass
                _fee_note_B = (
                    "ℹ️ **Performance nette des frais de gestion du contrat assureur** "
                    f"({st.session_state.get('CONTRACT_LABEL_B', 'contrat')}). "
                    "Les frais de gestion internes des fonds (TER) sont déjà intégrés "
                    "dans les valeurs liquidatives publiées et ne sont pas déduits en plus. "
                    "À titre de comparaison, Quantalys et Morningstar affichent les "
                    "performances brutes de frais contrat."
                )
                with st.expander("ℹ️ Détail du calcul de performance", expanded=False):
                    st.markdown(_fee_note_B)
                    if _fee_detail_B:
                        st.markdown("**Frais appliqués par fonds :**")
                        for _d in _fee_detail_B:
                            st.markdown(f"- {_d}")
                    st.markdown(
                        "**Formule :** XIRR calculé sur les flux réels "
                        "(versements initiaux, mensuels, ponctuels) et la valeur "
                        "liquidative nette à ce jour."
                    )


    def build_html_report(report: Dict[str, Any]) -> str:
        """
        Construit un rapport HTML exportable pour le client.
        Le contenu repose sur 'report', préparé plus bas dans le code.
        """
        as_of = report.get("as_of", "")
        mode = report.get("mode", "compare")
        synthA = report.get("client_summary", {})
        synthB = report.get("valority_summary", {})
        comp = report.get("comparison", {})

        dfA_lines = report.get("df_client_lines")
        dfB_lines = report.get("df_valority_lines")
        dfA_val = report.get("dfA_val")
        dfB_val = report.get("dfB_val")

        def _fmt_eur(x):
            try:
                return f"{x:,.2f} €".replace(",", " ").replace(".", ",")
            except Exception:
                return str(x)

        # Tables en HTML
        html_client_lines = dfA_lines.to_html(index=False, border=0, justify="left") if dfA_lines is not None else ""
        html_valority_lines = dfB_lines.to_html(index=False, border=0, justify="left") if dfB_lines is not None else ""

        if dfA_val is not None:
            html_A_val = dfA_val.to_html(index=False, border=0, justify="left")
        else:
            html_A_val = ""

        if dfB_val is not None:
            html_B_val = dfB_val.to_html(index=False, border=0, justify="left")
        else:
            html_B_val = ""

        if mode == "compare":
            synth_html = f"""
<div class="block">
  <h3>Situation actuelle — Client</h3>
  <ul>
    <li>Valeur actuelle : <b>{_fmt_eur(synthA.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthA.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthA.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthA.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthA.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>

<div class="block">
  <h3>Simulation — Allocation Cabinet</h3>
  <ul>
    <li>Valeur actuelle simulée : <b>{_fmt_eur(synthB.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthB.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthB.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthB.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthB.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>

<div class="block">
  <h3>Comparaison Client vs Cabinet</h3>
  <ul>
    <li>Différence de valeur finale : <b>{_fmt_eur(comp.get("delta_val", 0))}</b></li>
    <li>Écart de performance totale (Cabinet – Client) :
        <b>{comp.get("delta_perf_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Client</h3>
{html_client_lines}

<h3>Portefeuille Cabinet</h3>
{html_valority_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur des portefeuilles</h2>

<h3>Client – Valeur du portefeuille par date</h3>
{html_A_val}

<h3>Cabinet – Valeur du portefeuille par date</h3>
{html_B_val}
"""
        elif mode == "valority":
            synth_html = f"""
<div class="block">
  <h3>Simulation — Allocation Cabinet</h3>
  <ul>
    <li>Valeur actuelle simulée : <b>{_fmt_eur(synthB.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthB.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthB.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthB.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthB.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Cabinet</h3>
{html_valority_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur du portefeuille</h2>

<h3>Cabinet – Valeur du portefeuille par date</h3>
{html_B_val}
"""
        else:
            synth_html = f"""
<div class="block">
  <h3>Situation actuelle — Client</h3>
  <ul>
    <li>Valeur actuelle : <b>{_fmt_eur(synthA.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthA.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthA.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthA.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthA.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Client</h3>
{html_client_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur du portefeuille</h2>

<h3>Client – Valeur du portefeuille par date</h3>
{html_A_val}
"""

        html = f"""
<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="utf-8" />
<title>Rapport de portefeuille</title>
<style>
body {{
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
  margin: 24px;
  color: #222;
}}
h1, h2, h3 {{
  margin-top: 24px;
}}
table {{
  border-collapse: collapse;
  width: 100%;
  margin: 8px 0 16px 0;
  font-size: 14px;
}}
th, td {{
  border: 1px solid #ddd;
  padding: 6px 8px;
}}
th {{
  background-color: #f5f5f5;
  text-align: left;
}}
.small {{
  font-size: 12px;
  color: #666;
}}
.block {{
  border: 1px solid #eee;
  border-radius: 8px;
  padding: 12px 16px;
  margin-bottom: 16px;
  background-color: #fafafa;
}}
</style>
</head>
<body>

<h1>Rapport de portefeuille</h1>
<p class="small">Date de génération : {as_of}</p>

<h2>1. Synthèse chiffrée</h2>

{synth_html}
{detail_html}
{hist_html}

<p class="small">
Ce document est fourni à titre informatif uniquement et ne constitue pas un conseil en investissement
personnalisé.
</p>

</body>
    </html>
"""
        return html


    report_data = {
        "as_of": fmt_date(TODAY),
        "mode": st.session_state.get("MODE_ANALYSE", "compare"),
    }
    report_data["nom_client"] = st.session_state.get("NOM_CLIENT", "").strip()
    report_data["nom_cabinet"] = st.session_state.get("NOM_CABINET", "").strip()
    def _avg_fee_gestion(lines_key):
        lines = st.session_state.get(lines_key, [])
        fees = [float(l.get("fee_contract_pct") or 0) for l in lines if l.get("fee_contract_pct")]
        return round(sum(fees) / len(fees), 4) if fees else ANNUAL_FEE_UC_PCT
    report_data["fee_contract_pct"] = _avg_fee_gestion("A_lines")
    report_data["fee_gestion_A"] = _avg_fee_gestion("A_lines")
    report_data["fee_gestion_B"] = _avg_fee_gestion("B_lines")
    report_data["situation_familiale"] = st.session_state.get(
        "tax_situation_familiale", "Célibataire / veuf / divorcé"
    )
    report_data["contrat_label"] = st.session_state.get("CONTRACT_LABEL_A", "")
    report_data["date_ouverture"] = ""
    _all_lines = st.session_state.get("A_lines", [])
    if _all_lines:
        _dates = [ln.get("buy_date") for ln in _all_lines if ln.get("buy_date")]
        if _dates:
            _min_date = min(pd.Timestamp(d) for d in _dates)
            report_data["date_ouverture"] = _min_date.strftime("%d/%m/%Y")

    mode_report = report_data["mode"]
    df_client_lines = build_positions_dataframe("A_lines") if show_client else pd.DataFrame()
    df_valority_lines = build_positions_dataframe("B_lines") if show_valority else pd.DataFrame()

    report_data["df_client_lines"] = df_client_lines
    report_data["df_valority_lines"] = df_valority_lines
    report_data["positions_df_client"] = df_client_lines
    report_data["positions_df_valority"] = df_valority_lines
    report_data["dfA_val"] = (
        dfA.reset_index().rename(columns={"index": "Date"}) if (show_client and not dfA.empty) else pd.DataFrame()
    )
    report_data["dfB_val"] = (
        dfB.reset_index().rename(columns={"index": "Date"}) if (show_valority and not dfB.empty) else pd.DataFrame()
    )
    report_data["client_summary"] = {
        "val": valA,
        "net": netA,
        "brut": brutA,
        "perf_tot_pct": perf_tot_client or 0.0,
        "irr_pct": xirrA or 0.0,
    }
    report_data["valority_summary"] = {
        "val": valB,
        "net": netB,
        "brut": brutB,
        "perf_tot_pct": perf_tot_valority or 0.0,
        "irr_pct": xirrB or 0.0,
    }
    report_data["comparison"] = {
        "delta_val": (valB - valA) if (valA is not None and valB is not None) else 0.0,
        "delta_perf_pct": (perf_tot_valority - perf_tot_client)
        if (perf_tot_client is not None and perf_tot_valority is not None)
        else 0.0,
    }

    df_map: Dict[str, pd.DataFrame] = {}
    if mode_report in ("compare", "client") and not dfA.empty:
        df_map["Client"] = dfA
    if mode_report in ("compare", "valority") and not dfB.empty:
        df_map["Cabinet"] = dfB
    report_data["df_map"] = df_map

    if mode_report == "compare":
        positions_df = df_valority_lines if not df_valority_lines.empty else df_client_lines
        lines = st.session_state.get("B_lines", []) or st.session_state.get("A_lines", [])
    else:
        if mode_report == "valority":
            positions_df = df_valority_lines
            lines = st.session_state.get("B_lines", [])
            start_min = startB_min
            brut = brutB
            net = netB
            val = valB
        else:
            positions_df = df_client_lines
            lines = st.session_state.get("A_lines", [])
            start_min = startA_min
            brut = brutA
            net = netA
            val = valA

        years = _years_between(start_min, TODAY) if isinstance(start_min, pd.Timestamp) else 0.0
        fees_paid = max(0.0, brut - net) if brut is not None and net is not None else 0.0
        value_created = (val - net) if val is not None and net is not None else 0.0
        value_per_year = (value_created / years) if years > 0 else 0.0
        report_data["fees_analysis"] = {
            "fees_paid": fees_paid,
            "value_created": value_created,
            "value_per_year": value_per_year,
        }

    def _build_lines_with_values(
        lines_src: List[Dict[str, Any]],
        positions_src: pd.DataFrame,
    ) -> List[Dict[str, Any]]:
        items: List[Dict[str, Any]] = []
        if isinstance(positions_src, pd.DataFrame) and not positions_src.empty:
            for ln in lines_src:
                isin = ln.get("isin", "")
                name = ln.get("name", "")
                match = positions_src[
                    (positions_src["Nom"] == name) & (positions_src["ISIN / Code"] == isin)
                ]
                val = float(match["Valeur actuelle €"].iloc[0]) if not match.empty else 0.0
                items.append({
                    "isin": isin,
                    "name": name,
                    "value": val,
                    "fee_uc_pct": ln.get("fee_uc_pct"),
                    "fee_contract_pct": ln.get("fee_contract_pct"),
                    "fee_total_pct": ln.get("fee_total_pct"),
                })
        return items

    report_data["positions_df"] = positions_df
    report_data["lines"] = _build_lines_with_values(lines, positions_df)
    report_data["lines_client"] = _build_lines_with_values(
        st.session_state.get("A_lines", []),
        df_client_lines,
    )
    report_data["lines_valority"] = _build_lines_with_values(
        st.session_state.get("B_lines", []),
        df_valority_lines,
    )

    st.session_state["REPORT_DATA"] = report_data
    # Pré-générer le PDF une seule fois (évite double-rendu au 1er chargement)
    try:
        _pdf_bytes_cache = generate_pdf_report(report_data)
    except Exception:
        _pdf_bytes_cache = b""
    st.session_state["PDF_BYTES_CACHE"] = _pdf_bytes_cache

    # ------------------------------------------------------------
    # Bloc final : Comparaison OU "Frais & valeur créée"
    # ------------------------------------------------------------
    mode = st.session_state.get("MODE_ANALYSE", "compare")

    # ============================
    # CAS 1 — MODE COMPARAISON
    # ============================
    if mode == "compare":
        st.markdown("---")
        st.subheader("⚖️ Comparaison")

        gain_vs_client = (valB - valA) if (valA is not None and valB is not None) else 0.0
        delta_xirr = (xirrB - xirrA) if (xirrA is not None and xirrB is not None) else None
        perf_diff_tot = (
            (perf_tot_valority - perf_tot_client)
            if (perf_tot_client is not None and perf_tot_valority is not None)
            else None
        )

        with st.container(border=True):
            c1, c2, c3 = st.columns(3)
            with c1:
                st.metric("Gain en valeur", to_eur(gain_vs_client))
            with c2:
                st.metric(
                    "Surperformance totale",
                    f"{perf_diff_tot:+.2f}%" if perf_diff_tot is not None else "—",
                )
            with c3:
                st.metric(
                    "Surperformance annualisée (Î” XIRR)",
                    f"{delta_xirr:+.2f}%" if delta_xirr is not None else "—",
                )



    # ============================
    # CAS 2 — MODE ANALYSE SIMPLE
    # ============================
    else:
        # Sélection des variables selon le mode
        if mode == "valority":
            brut = brutB
            net = netB
            val = valB
            start_min = startB_min
            irr = xirrB
            fee_pct = st.session_state.get("FEE_B", 0.0)
            title = "🏢 Allocation Cabinet — Frais & valeur créée"
        else:  # mode == "client"
            brut = brutA
            net = netA
            val = valA
            start_min = startA_min
            irr = xirrA
            fee_pct = st.session_state.get("FEE_A", 0.0)
            title = "🧍 Portefeuille — Frais & valeur créée"

        st.markdown("---")
        st.subheader("📌 Analyse : frais & valeur créée")

        if brut > 0 and net >= 0 and val >= 0 and isinstance(start_min, pd.Timestamp):
            fees_paid = max(0.0, brut - net)     # frais d'entrée réellement payés
            value_created = val - net            # valeur créée vs capital réellement investi
            years = _years_between(start_min, TODAY)
            value_per_year = (value_created / years) if years > 0 else None

            with st.container(border=True):
                st.markdown(f"#### {title}")
                st.caption(
                    f"Période : **{fmt_date(start_min)} → {fmt_date(TODAY)}** "
                    f"• Frais d'entrée : **{fee_pct:.2f}%**"
                )

                _pad_l, c1, c2, c3, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                with c1:
                    st.metric("Frais d'entrée payés", to_eur(fees_paid))
                with c2:
                    st.metric("Valeur créée (net)", to_eur(value_created))
                with c3:
                    st.metric(
                        "Valeur créée / an (moyenne)",
                        to_eur(value_per_year) if value_per_year is not None else "—",
                    )

                st.markdown(
                    f"""
- Montants versés (brut) : **{to_eur(brut)}**
- Montants réellement investis (après frais) : **{to_eur(net)}**
- Valeur actuelle : **{to_eur(val)}**
"""
                )

                if irr is not None:
                    st.markdown(f"- Rendement annualisé (XIRR) : **{irr:.2f}%**")
                else:
                    st.markdown("- Rendement annualisé (XIRR) : **—**")

                # Message factuel — différencié selon que la valeur créée est positive ou négative
                if fees_paid > 0 and value_created > 0:
                    ratio = value_created / fees_paid
                    st.markdown(
                        f"**Lecture :** {to_eur(fees_paid)} de frais d'entrée ont généré "
                        f"**{to_eur(value_created)}** de valeur nette créée à date "
                        f"(**×{ratio:.1f}**)."
                    )
                elif fees_paid > 0 and value_created <= 0:
                    st.markdown(
                        f"**Lecture :** {to_eur(fees_paid)} de frais d'entrée payés. "
                        f"Le portefeuille affiche une moins-value nette de **{to_eur(abs(value_created))}** à date."
                    )
        else:
            st.info("Ajoutez des lignes (et/ou des versements) pour afficher l'analyse frais & valeur créée.")

    # ------------------------------------------------------------
    # Transparence des frais
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("💰 Transparence des frais")

    _lns_fee_A = st.session_state.get("A_lines", []) if show_client else []
    _lns_fee_B = st.session_state.get("B_lines", []) if show_valority else []

    def _calc_fee_cost(lines: list, positions_df: pd.DataFrame) -> tuple:
        """Retourne (cout_annuel_eur, fee_avg_pct_weighted). EUROFUND exclu."""
        if not lines:
            return 0.0, 0.0
        _val_map_fc: Dict[str, float] = {}
        if not positions_df.empty and "ISIN / Code" in positions_df.columns and "Valeur actuelle €" in positions_df.columns:
            for _, _rfc in positions_df.iterrows():
                _k_fc = str(_rfc.get("ISIN / Code", "") or "").upper()
                if _k_fc:
                    _val_map_fc[_k_fc] = float(_rfc.get("Valeur actuelle €", 0) or 0)
        _total_val_fc = 0.0
        _weighted_fee_fc = 0.0
        _cout_fc = 0.0
        for _ln_fc in lines:
            _isin_fc = str(_ln_fc.get("isin", "") or "").upper()
            if _isin_fc == "EUROFUND":
                continue
            _val_fc = _val_map_fc.get(_isin_fc, float(_ln_fc.get("amount_gross", 0) or 0))
            if _isin_fc == "STRUCTURED":
                _fee_fc = float(_ln_fc.get("fee_contract_pct") or 0)
            else:
                _ter_fc = float(_ln_fc.get("fee_uc_pct") or 0)
                _ctr_fc = float(_ln_fc.get("fee_contract_pct") or 0)
                _fee_fc = _ter_fc + _ctr_fc
            _cout_fc += _val_fc * _fee_fc / 100.0
            _weighted_fee_fc += _fee_fc * _val_fc
            _total_val_fc += _val_fc
        _avg_fc = _weighted_fee_fc / _total_val_fc if _total_val_fc > 0 else 0.0
        return _cout_fc, _avg_fc

    _cout_A, _fee_avg_A = _calc_fee_cost(_lns_fee_A, df_client_lines) if _lns_fee_A else (0.0, 0.0)
    _cout_B, _fee_avg_B = _calc_fee_cost(_lns_fee_B, df_valority_lines) if _lns_fee_B else (0.0, 0.0)

    _has_fee_data_tr = (_cout_A > 0) or (_cout_B > 0)
    _RDT_HYPO_TR = 5.0 / 100.0  # hypothèse rendement 5%/an
    _capital_ref = 100_000.0  # base standardisée pour comparaison équitable

    st.caption("Frais des fonds euros exclus du calcul (inclus dans le taux net publié).")

    if _has_fee_data_tr:
        _HORIZONS_TR = [5, 10, 15, 20]

        # ── ① Bar chart vertical — Frais annuels (sur capital réel)
        with st.container(border=True):
            st.markdown("**💸 Coût annuel des frais récurrents**")
            if mode == "compare" and _cout_A > 0 and _cout_B > 0:
                # Affichage sur base standardisée pour comparaison équitable des taux
                _cout_display_A = _capital_ref * _fee_avg_A / 100.0
                _cout_display_B = _capital_ref * _fee_avg_B / 100.0
                _bar_ann_df = pd.DataFrame([
                    {"Portefeuille": "Client", "Frais annuels (€)": _cout_display_A},
                    {"Portefeuille": "Cabinet", "Frais annuels (€)": _cout_display_B},
                ])
                _bar_ann_base = alt.Chart(_bar_ann_df)
                _bar_ann = (
                    _bar_ann_base
                    .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                    .encode(
                        x=alt.X("Portefeuille:N", title="", axis=alt.Axis(labelAngle=0)),
                        y=alt.Y("Frais annuels (€):Q", axis=alt.Axis(format=",.0f"), title="€/an"),
                        color=alt.Color(
                            "Portefeuille:N",
                            scale=alt.Scale(domain=["Client", "Cabinet"], range=["#E53935", "#2E7D32"]),
                            legend=None,
                        ),
                        tooltip=[
                            alt.Tooltip("Portefeuille:N"),
                            alt.Tooltip("Frais annuels (€):Q", format=",.0f", title="Coût/an (base 100k€)"),
                        ],
                    )
                )
                _bar_ann_lbl = (
                    _bar_ann_base
                    .mark_text(dy=-8, fontSize=11, fontWeight="bold")
                    .encode(
                        x=alt.X("Portefeuille:N"),
                        y=alt.Y("Frais annuels (€):Q"),
                        text=alt.Text("Frais annuels (€):Q", format=",.0f"),
                        color=alt.value("#333333"),
                    )
                )
                st.altair_chart((_bar_ann + _bar_ann_lbl).properties(height=250), use_container_width=True)
                st.caption(
                    f"Taux moyen pondéré : 🧍 Client **{_fee_avg_A:.2f}%/an** — "
                    f"🏢 Cabinet **{_fee_avg_B:.2f}%/an** — Base : 100 000 €"
                )
                _eco_ann = _cout_display_A - _cout_display_B
                _delta_ann = "par an en faveur du cabinet" if _eco_ann > 0 else "par an en faveur du client"
                st.metric("Économie annuelle", to_eur(abs(_eco_ann)), delta=_delta_ann)
            elif _cout_A > 0:
                _bar_s_df = pd.DataFrame([{"Portefeuille": "Client", "Frais annuels (€)": _cout_A}])
                _bar_s_base = alt.Chart(_bar_s_df)
                _bar_s = (
                    _bar_s_base
                    .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3, color="#E53935")
                    .encode(
                        x=alt.X("Portefeuille:N", title="", axis=alt.Axis(labelAngle=0)),
                        y=alt.Y("Frais annuels (€):Q", axis=alt.Axis(format=",.0f"), title="€/an"),
                        tooltip=[alt.Tooltip("Frais annuels (€):Q", format=",.0f", title="Coût/an")],
                    )
                )
                _bar_s_lbl = (
                    _bar_s_base
                    .mark_text(dy=-8, fontSize=11, fontWeight="bold")
                    .encode(
                        x=alt.X("Portefeuille:N"),
                        y=alt.Y("Frais annuels (€):Q"),
                        text=alt.Text("Frais annuels (€):Q", format=",.0f"),
                        color=alt.value("#333333"),
                    )
                )
                st.altair_chart((_bar_s + _bar_s_lbl).properties(height=200), use_container_width=True)
                st.caption(f"Taux moyen pondéré : **{_fee_avg_A:.2f}%/an**")
                st.metric("Coût annuel estimé", to_eur(_cout_A))
            elif _cout_B > 0:
                _bar_s_df = pd.DataFrame([{"Portefeuille": "Cabinet", "Frais annuels (€)": _cout_B}])
                _bar_s_base = alt.Chart(_bar_s_df)
                _bar_s = (
                    _bar_s_base
                    .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3, color="#2E7D32")
                    .encode(
                        x=alt.X("Portefeuille:N", title="", axis=alt.Axis(labelAngle=0)),
                        y=alt.Y("Frais annuels (€):Q", axis=alt.Axis(format=",.0f"), title="€/an"),
                        tooltip=[alt.Tooltip("Frais annuels (€):Q", format=",.0f", title="Coût/an")],
                    )
                )
                _bar_s_lbl = (
                    _bar_s_base
                    .mark_text(dy=-8, fontSize=11, fontWeight="bold")
                    .encode(
                        x=alt.X("Portefeuille:N"),
                        y=alt.Y("Frais annuels (€):Q"),
                        text=alt.Text("Frais annuels (€):Q", format=",.0f"),
                        color=alt.value("#333333"),
                    )
                )
                st.altair_chart((_bar_s + _bar_s_lbl).properties(height=200), use_container_width=True)
                st.caption(f"Taux moyen pondéré : **{_fee_avg_B:.2f}%/an**")
                st.metric("Coût annuel estimé", to_eur(_cout_B))

        # ── ② Bar chart groupé — Frais cumulés (base standardisée 100 000 €)
        with st.container(border=True):
            st.markdown("**📈 Impact cumulé des frais sur la durée**")
            _drag_horizon = st.slider("Horizon de projection", 5, 30, 15, 5, key="drag_horizon")
            st.caption(
                "⚠️ Simulation sur une base de **100 000 €**, hypothèse de rendement : **5%/an** "
                "(hors frais). Les rendements futurs ne sont pas garantis."
            )
            _drag_rows_tr: List[Dict] = []
            for _h_tr in _HORIZONS_TR:
                if _cout_A > 0:
                    _drag_A_h = max(0.0, _capital_ref * (
                        (1 + _RDT_HYPO_TR) ** _h_tr
                        - (1 + _RDT_HYPO_TR - _fee_avg_A / 100.0) ** _h_tr
                    ))
                    _drag_rows_tr.append({"Horizon": f"{_h_tr} ans", "Frais cumulés (€)": _drag_A_h, "Portefeuille": "Client"})
                if _cout_B > 0:
                    _drag_B_h = max(0.0, _capital_ref * (
                        (1 + _RDT_HYPO_TR) ** _h_tr
                        - (1 + _RDT_HYPO_TR - _fee_avg_B / 100.0) ** _h_tr
                    ))
                    _drag_rows_tr.append({"Horizon": f"{_h_tr} ans", "Frais cumulés (€)": _drag_B_h, "Portefeuille": "Cabinet"})
            if _drag_rows_tr:
                _drag_df_tr = pd.DataFrame(_drag_rows_tr)
                _sort_tr = [f"{h} ans" for h in _HORIZONS_TR]
                _domain_tr = [p for p in ["Client", "Cabinet"] if p in _drag_df_tr["Portefeuille"].unique()]
                _range_tr = [{"Client": "#E53935", "Cabinet": "#2E7D32"}[p] for p in _domain_tr]
                _drag_chart_tr = (
                    alt.Chart(_drag_df_tr)
                    .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                    .encode(
                        x=alt.X("Horizon:N", sort=_sort_tr, title="Horizon"),
                        y=alt.Y("Frais cumulés (€):Q", axis=alt.Axis(format=",.0f"), title="€"),
                        color=alt.Color(
                            "Portefeuille:N",
                            scale=alt.Scale(domain=_domain_tr, range=_range_tr),
                            legend=alt.Legend(title=""),
                        ),
                        xOffset=alt.XOffset("Portefeuille:N"),
                        tooltip=[
                            alt.Tooltip("Horizon:N"),
                            alt.Tooltip("Portefeuille:N"),
                            alt.Tooltip("Frais cumulés (€):Q", format=",.0f", title="Frais cumulés"),
                        ],
                    )
                    .properties(height=300)
                )
                st.altair_chart(_drag_chart_tr, use_container_width=True)

        # ── ③ Message de synthèse
        _drag_sel_A = 0.0
        _drag_sel_B = 0.0
        if _cout_A > 0:
            _drag_sel_A = max(0.0, _capital_ref * (
                (1 + _RDT_HYPO_TR) ** _drag_horizon
                - (1 + _RDT_HYPO_TR - _fee_avg_A / 100.0) ** _drag_horizon
            ))
        if _cout_B > 0:
            _drag_sel_B = max(0.0, _capital_ref * (
                (1 + _RDT_HYPO_TR) ** _drag_horizon
                - (1 + _RDT_HYPO_TR - _fee_avg_B / 100.0) ** _drag_horizon
            ))
        with st.container(border=True):
            if mode == "compare" and _drag_sel_A > 0 and _drag_sel_B > 0:
                _diff_drag_tr = _drag_sel_A - _drag_sel_B
                if _diff_drag_tr > 50:
                    st.success(
                        f"💬 Sur **{_drag_horizon} ans**, pour 100 000 € investis, les frais représentent "
                        f"**{to_eur(_drag_sel_A)}** (Client) vs **{to_eur(_drag_sel_B)}** (Cabinet). "
                        f"Économie estimée : **{to_eur(abs(_diff_drag_tr))}**."
                    )
                elif _diff_drag_tr < -50:
                    st.info(
                        f"💬 Les frais du cabinet ({_fee_avg_B:.2f}%/an) sont légèrement supérieurs. "
                        f"L'avantage de la proposition repose sur la performance et la diversification."
                    )
                else:
                    st.info(
                        f"💬 Sur **{_drag_horizon} ans**, les frais des deux portefeuilles sont comparables "
                        f"({to_eur(_drag_sel_A)} Client vs {to_eur(_drag_sel_B)} Cabinet, base 100 000 €)."
                    )
            elif _drag_sel_A > 0:
                st.caption(
                    f"💬 Sur {_drag_horizon} ans, pour 100 000 € investis, les frais récurrents "
                    f"représentent **{to_eur(_drag_sel_A)}** (hypothèse 5%/an)."
                )
            elif _drag_sel_B > 0:
                st.caption(
                    f"💬 Sur {_drag_horizon} ans, pour 100 000 € investis, les frais récurrents "
                    f"représentent **{to_eur(_drag_sel_B)}** (hypothèse 5%/an)."
                )

        # ── ④ Expander : Détail des frais par support
        with st.expander("📐 Détail des frais par support", expanded=False):
            def _build_fee_detail_df(lines: list, positions_df: pd.DataFrame) -> pd.DataFrame:
                _val_map_fd: Dict[str, float] = {}
                if not positions_df.empty and "ISIN / Code" in positions_df.columns and "Valeur actuelle €" in positions_df.columns:
                    for _, _rfd in positions_df.iterrows():
                        _k_fd = str(_rfd.get("ISIN / Code", "") or "").upper()
                        if _k_fd:
                            _val_map_fd[_k_fd] = float(_rfd.get("Valeur actuelle €", 0) or 0)
                _rows_fd: List[Dict] = []
                _tot_val_fd = 0.0
                _wtd_fee_fd = 0.0
                _tot_cout_fd = 0.0
                for _ln_fd in lines:
                    _isin_fd = str(_ln_fd.get("isin", "") or "").upper()
                    _name_fd = str(_ln_fd.get("name", "") or _isin_fd)[:35]
                    if _isin_fd == "EUROFUND":
                        _rows_fd.append({
                            "Ligne": _name_fd,
                            "ISIN": _isin_fd,
                            "TER fonds (%/an)": "inclus dans taux net",
                            "Frais contrat (%/an)": "inclus dans taux net",
                            "Total (%/an)": "0.00%",
                            "Valeur actuelle €": _val_map_fd.get(_isin_fd, float(_ln_fd.get("amount_gross", 0) or 0)),
                            "Coût annuel €": 0.0,
                        })
                        continue
                    _val_fd = _val_map_fd.get(_isin_fd, float(_ln_fd.get("amount_gross", 0) or 0))
                    if _isin_fd == "STRUCTURED":
                        _ter_fd = 0.0
                        _ctr_fd = float(_ln_fd.get("fee_contract_pct") or 0)
                    else:
                        _ter_fd = float(_ln_fd.get("fee_uc_pct") or 0)
                        _ctr_fd = float(_ln_fd.get("fee_contract_pct") or 0)
                    _total_fd = _ter_fd + _ctr_fd
                    _cout_fd = _val_fd * _total_fd / 100.0
                    _tot_val_fd += _val_fd
                    _wtd_fee_fd += _total_fd * _val_fd
                    _tot_cout_fd += _cout_fd
                    _rows_fd.append({
                        "Ligne": _name_fd,
                        "ISIN": _isin_fd,
                        "TER fonds (%/an)": f"{_ter_fd:.2f}%",
                        "Frais contrat (%/an)": f"{_ctr_fd:.2f}%",
                        "Total (%/an)": f"{_total_fd:.2f}%",
                        "Valeur actuelle €": _val_fd,
                        "Coût annuel €": _cout_fd,
                    })
                if not _rows_fd:
                    return pd.DataFrame()
                _avg_fd = _wtd_fee_fd / _tot_val_fd if _tot_val_fd > 0 else 0.0
                _rows_fd.append({
                    "Ligne": "📊 Total / Moyenne pondérée",
                    "ISIN": "",
                    "TER fonds (%/an)": "",
                    "Frais contrat (%/an)": "",
                    "Total (%/an)": f"{_avg_fd:.2f}%",
                    "Valeur actuelle €": _tot_val_fd,
                    "Coût annuel €": _tot_cout_fd,
                })
                return pd.DataFrame(_rows_fd)

            _col_cfg_fd = {
                "Valeur actuelle €": st.column_config.NumberColumn("Valeur actuelle €", format="%.0f €"),
                "Coût annuel €": st.column_config.NumberColumn("Coût annuel €", format="%.0f €"),
            }
            if mode == "compare":
                _fd_c1, _fd_c2 = st.columns(2)
                with _fd_c1:
                    st.caption("🧍 **Client**")
                    if _lns_fee_A:
                        _df_fd_A = _build_fee_detail_df(_lns_fee_A, df_client_lines)
                        if not _df_fd_A.empty:
                            st.dataframe(_df_fd_A, hide_index=True, use_container_width=True, column_config=_col_cfg_fd)
                with _fd_c2:
                    st.caption("🏢 **Cabinet**")
                    if _lns_fee_B:
                        _df_fd_B = _build_fee_detail_df(_lns_fee_B, df_valority_lines)
                        if not _df_fd_B.empty:
                            st.dataframe(_df_fd_B, hide_index=True, use_container_width=True, column_config=_col_cfg_fd)
            elif _lns_fee_A:
                _df_fd_s = _build_fee_detail_df(_lns_fee_A, df_client_lines)
                if not _df_fd_s.empty:
                    st.dataframe(_df_fd_s, hide_index=True, use_container_width=True, column_config=_col_cfg_fd)
            elif _lns_fee_B:
                _df_fd_s = _build_fee_detail_df(_lns_fee_B, df_valority_lines)
                if not _df_fd_s.empty:
                    st.dataframe(_df_fd_s, hide_index=True, use_container_width=True, column_config=_col_cfg_fd)

    # ------------------------------------------------------------
    # Indicateurs de risque & performance (2 niveaux)
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📊 Indicateurs de risque & performance")

    _period_label_parts = []
    if show_client and isinstance(startA_min, pd.Timestamp):
        _period_label_parts.append(f"Client : depuis le {fmt_date(startA_min)}")
    if show_valority and isinstance(startB_min, pd.Timestamp):
        _period_label_parts.append(f"Cabinet : depuis le {fmt_date(startB_min)}")
    if _period_label_parts:
        st.caption(" | ".join(_period_label_parts))

    _bench_sym = st.session_state.get("BENCHMARK_SYMBOL", "CW8.PA")
    _euro_rate_A = st.session_state.get("EURO_RATE_A", 2.0)
    _euro_rate_B = st.session_state.get("EURO_RATE_B", 2.5)
    _linesA = st.session_state.get("A_lines", [])
    _linesB = st.session_state.get("B_lines", [])
    _fee_A = float(st.session_state.get("FEE_A", 0.0))
    _fee_B = float(st.session_state.get("FEE_B", 0.0))

    # ── Calculs diversification ────────────────────────────────────────────
    # ── Dates réelles des portefeuilles pour les calculs de risque ──────────
    _start_A = startA_min if isinstance(startA_min, pd.Timestamp) else None
    _start_B = startB_min if isinstance(startB_min, pd.Timestamp) else None

    _div_A = compute_diversification_score(_linesA, _euro_rate_A, start_date=_start_A) if (show_client and _linesA) else None
    _div_B = compute_diversification_score(_linesB, _euro_rate_B, start_date=_start_B) if (show_valority and _linesB) else None
    _risk_A = portfolio_risk_stats(_linesA, _euro_rate_A, start_date=_start_A, fee_pct=_fee_A) if (show_client and _linesA) else None
    _risk_B = portfolio_risk_stats(_linesB, _euro_rate_B, start_date=_start_B, fee_pct=_fee_B) if (show_valority and _linesB) else None

    # ── Calculs ratios techniques ──────────────────────────────────────────
    _sharpe_A = compute_sharpe_ratio(_linesA, _euro_rate_A, _fee_A, start_date=_start_A) if (show_client and _linesA) else None
    _sharpe_B = compute_sharpe_ratio(_linesB, _euro_rate_B, _fee_B, start_date=_start_B) if (show_valority and _linesB) else None
    _sortino_A = compute_sortino_ratio(_linesA, _euro_rate_A, _fee_A, start_date=_start_A) if (show_client and _linesA) else None
    _sortino_B = compute_sortino_ratio(_linesB, _euro_rate_B, _fee_B, start_date=_start_B) if (show_valority and _linesB) else None
    _ba_A = compute_beta_alpha(_linesA, _euro_rate_A, _fee_A, _bench_sym, start_date=_start_A) if (show_client and _linesA) else None
    _ba_B = compute_beta_alpha(_linesB, _euro_rate_B, _fee_B, _bench_sym, start_date=_start_B) if (show_valority and _linesB) else None

    # Stocker le bêta pour le stress-test
    st.session_state["BETA_AUTO_A"] = _ba_A["beta"] if _ba_A else None
    st.session_state["BETA_AUTO_B"] = _ba_B["beta"] if _ba_B else None

    # Indicateurs de risque pour le PDF — inséré ICI après calcul des variables
    report_data["diversification_A"] = _div_A
    report_data["diversification_B"] = _div_B
    report_data["risk_A"] = _risk_A
    report_data["risk_B"] = _risk_B
    report_data["sharpe_A"] = _sharpe_A
    report_data["sharpe_B"] = _sharpe_B
    report_data["sortino_A"] = _sortino_A
    report_data["sortino_B"] = _sortino_B
    report_data["beta_alpha_A"] = _ba_A
    report_data["beta_alpha_B"] = _ba_B
    report_data["ratios_A"] = {
        "sharpe": _sharpe_A,
        "sortino": _sortino_A,
        "beta_alpha": _ba_A,
    }
    report_data["ratios_B"] = {
        "sharpe": _sharpe_B,
        "sortino": _sortino_B,
        "beta_alpha": _ba_B,
    }

    # ── Données enrichies pour la PPTX client ─────────────────────────────
    report_data["euro_rate_A"] = st.session_state.get("EURO_RATE_A", 2.0)
    report_data["euro_rate_B"] = st.session_state.get("EURO_RATE_B", 2.5)
    report_data["fee_A"] = float(st.session_state.get("FEE_A", 0.0))
    report_data["fee_B"] = float(st.session_state.get("FEE_B", 0.0))
    report_data["situation_familiale"] = st.session_state.get(
        "tax_situation_familiale", "Célibataire / veuf / divorcé"
    )
    report_data["age_souscripteur"] = st.session_state.get("tax_age_souscripteur", 55)
    report_data["nb_beneficiaires"] = int(st.session_state.get("tax_nb_beneficiaires", 2))
    report_data["nb_enfants"] = int(st.session_state.get("tax_nb_enfants", 2))

    # Rachat optimal Client
    if show_client and valA and valA > 0 and netA and netA > 0:
        _tax_date_r = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
        _anc_r = (date.today() - _tax_date_r).days / 365.25
        _sit_r = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")
        _net_total_r = float(st.session_state.get("tax_versements_nets_total", netA))
        if _anc_r >= 8:
            _opt_r_A = calc_optimisation_abattement(valA, netA, _anc_r, _sit_r, _net_total_r)
            report_data["rachat_optimal_A"] = _opt_r_A["rachat_optimal"]
            report_data["rachat_ps_A"] = _opt_r_A["ps_du"]
        else:
            report_data["rachat_optimal_A"] = 0.0
            report_data["rachat_ps_A"] = 0.0
    else:
        report_data["rachat_optimal_A"] = 0.0
        report_data["rachat_ps_A"] = 0.0

    # Rachat optimal Cabinet
    if show_valority and valB and valB > 0 and netB and netB > 0:
        _tax_date_r = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
        _anc_r = (date.today() - _tax_date_r).days / 365.25
        _sit_r = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")
        _net_total_r = float(st.session_state.get("tax_versements_nets_total", netB))
        if _anc_r >= 8:
            _opt_r_B = calc_optimisation_abattement(valB, netB, _anc_r, _sit_r, _net_total_r)
            report_data["rachat_optimal_B"] = _opt_r_B["rachat_optimal"]
            report_data["rachat_ps_B"] = _opt_r_B["ps_du"]
        else:
            report_data["rachat_optimal_B"] = 0.0
            report_data["rachat_ps_B"] = 0.0
    else:
        report_data["rachat_optimal_B"] = 0.0
        report_data["rachat_ps_B"] = 0.0

    # Données transmission
    _nb_benef_r = int(st.session_state.get("tax_nb_beneficiaires", 2))
    report_data["plafond_990i"] = 152_500 * _nb_benef_r

    # Stocker le report_data complet (avec toutes les données) pour la section Rapports
    st.session_state["_LAST_REPORT_DATA"] = report_data

    _mode_risk = st.session_state.get("MODE_ANALYSE", "compare")

    # ══════════════════════════════════════════════════════════════════════
    # ① DIVERSIFICATION
    # ══════════════════════════════════════════════════════════════════════
    st.markdown("##### 🧩 Analyse du portefeuille")

    # ── Allocation 3 poches + profil factuel ─────────────────────────────
    def _calc_allocation_3poches(lines: List, funds_df_key: str) -> Tuple[float, float, float, str]:
        _fdf = st.session_state.get(funds_df_key, pd.DataFrame())
        _total = 0.0
        _securise = 0.0
        _defensif = 0.0
        _dynamique = 0.0
        _bond_kw = ("OBLIGATION", "BOND", "FIXED INCOME", "MONÉTAIRE",
                     "MONETAIRE", "MONEY MARKET", "TAUX", "RATE",
                     "GOUVERNEMENT", "GOVERNMENT", "SOVEREIGN",
                     "CORPORATE", "HIGH YIELD", "AGGREGATE",
                     "REVENUS FIXES", "CONVERTIBLES", "SUBORDINATED")
        for _ln in lines:
            _amt = float(_ln.get("amount_gross", 0) or 0)
            _isin = str(_ln.get("isin", "")).upper()
            _total += _amt
            if _isin == "EUROFUND":
                _securise += _amt
            else:
                _cat = ""
                if isinstance(_fdf, pd.DataFrame) and not _fdf.empty:
                    _match = _fdf[_fdf["isin"].str.upper() == _isin]
                    if not _match.empty:
                        _cat = str(_match.iloc[0].get("category", "")).upper()
                if any(kw in _cat for kw in _bond_kw):
                    _defensif += _amt
                else:
                    _dynamique += _amt
        if _total <= 0:
            return 0.0, 0.0, 0.0, "—"
        _pct_s = _securise / _total * 100
        _pct_d = _defensif / _total * 100
        _pct_a = _dynamique / _total * 100
        _pct_expose = _pct_d + _pct_a
        if _pct_expose <= 30:
            _profil = "Sécuritaire"
        elif _pct_expose <= 50:
            _profil = "Prudent"
        elif _pct_expose <= 70:
            _profil = "Équilibré"
        elif _pct_expose <= 90:
            _profil = "Dynamique"
        else:
            _profil = "Offensif"
        return _pct_s, _pct_d, _pct_a, _profil

    def _isin_from_name_pfx(name_pfx: str, lines: List) -> str:
        for _ln in lines:
            _lbl = (_ln.get("name") or _ln.get("isin") or "").strip()
            if _lbl[:30] == name_pfx:
                return str(_ln.get("isin", "")).strip()
        return ""

    # ── Jauge score diversification ──────────────────────────────────────
    def _render_gauge_div(div_res: Optional[Dict], label: str) -> None:
        if div_res is None:
            st.info(f"{label} : ajoutez au moins 2 lignes pour calculer la diversification.")
            return
        if div_res.get("score") is None:
            st.info(div_res.get("message", "Pas de diversification à analyser"))
            return
        _sc = float(div_res["score"])
        _gc = "#2E7D32" if _sc >= 70 else ("#FF9800" if _sc >= 40 else "#E53935")
        _verdict = (
            "Bonne diversification des supports" if _sc >= 70
            else ("Diversification des supports à améliorer" if _sc >= 40
                  else "Diversification des supports insuffisante")
        )
        _gdf = pd.DataFrame([
            {"seg": "score", "val": _sc, "ord": 1},
            {"seg": "reste", "val": max(0.0, 100.0 - _sc), "ord": 2},
        ])
        _arc = (
            alt.Chart(_gdf)
            .mark_arc(innerRadius=55, outerRadius=85)
            .encode(
                theta=alt.Theta("val:Q"),
                color=alt.Color(
                    "seg:N",
                    scale=alt.Scale(domain=["score", "reste"], range=[_gc, "#E0E0E0"]),
                    legend=None,
                ),
                order=alt.Order("ord:Q"),
                tooltip=alt.value(None),
            )
        )
        _txt_c = (
            alt.Chart(pd.DataFrame([{"t": f"{_sc:.0f}"}]))
            .mark_text(size=22, fontWeight="bold", dy=5)
            .encode(text="t:N")
        )
        st.altair_chart(
            alt.layer(_arc, _txt_c).properties(height=200),
            use_container_width=True,
        )
        st.caption(f"**{label}** — {_verdict}")

    # ── ⑥ Sparklines + détails doublons (réutilisé dans les 2 expanders) ──
    def _render_doublon_detail(div_res: Dict, lines: List, euro_rate: float) -> None:
        for _ni, _nj, _c in div_res["doublons"]:
            st.warning(f"⚠️ **{_ni}** et **{_nj}** — corrélation {_c:.0%}")
            _isin_i = _isin_from_name_pfx(_ni, lines)
            _isin_j = _isin_from_name_pfx(_nj, lines)
            if _isin_i and _isin_j:
                try:
                    _s1, _, _ = get_price_series(_isin_i, None, euro_rate)
                    _s2, _, _ = get_price_series(_isin_j, None, euro_rate)
                    if not _s1.empty and not _s2.empty:
                        _cutoff_sp = pd.Timestamp.today() - pd.Timedelta(days=365)
                        _p1 = _s1["Close"].astype(float)
                        _p1 = _p1[_p1.index >= _cutoff_sp]
                        _p2 = _s2["Close"].astype(float)
                        _p2 = _p2[_p2.index >= _cutoff_sp]
                        if len(_p1) > 10 and len(_p2) > 10:
                            _cs = max(_p1.index[0], _p2.index[0])
                            _p1 = _p1[_p1.index >= _cs]
                            _p2 = _p2[_p2.index >= _cs]
                            _b1 = (_p1 / _p1.iloc[0] * 100).reset_index()
                            _b1.columns = ["Date", "Val"]
                            _b1["Fonds"] = _ni[:30]
                            _b2 = (_p2 / _p2.iloc[0] * 100).reset_index()
                            _b2.columns = ["Date", "Val"]
                            _b2["Fonds"] = _nj[:30]
                            _sp_df = pd.concat([_b1, _b2], ignore_index=True)
                            _sp_chart = (
                                alt.Chart(_sp_df)
                                .mark_line()
                                .encode(
                                    x=alt.X("Date:T", title=""),
                                    y=alt.Y("Val:Q", title="Base 100", scale=alt.Scale(zero=False)),
                                    color=alt.Color("Fonds:N", legend=alt.Legend(orient="bottom", title="")),
                                    tooltip=[
                                        alt.Tooltip("Date:T"),
                                        alt.Tooltip("Fonds:N"),
                                        alt.Tooltip("Val:Q", format=".1f", title="Base 100"),
                                    ],
                                )
                                .properties(
                                    height=200,
                                    title=f"{_ni[:25]} vs {_nj[:25]} (corr. {_c:.0%})",
                                )
                            )
                            st.altair_chart(_sp_chart, use_container_width=True)
                except Exception:
                    pass
        for _ni_v, _nj_v, _c_v in div_res["vigilance"]:
            st.caption(f"🟠 {_ni_v} et {_nj_v} — corrélation {_c_v:.0%}")
        st.caption(
            f"{div_res['n_effective']} sources de diversification réelles "
            f"sur {div_res['n_lines']} lignes"
        )
        st.caption(f"Corrélation moyenne hors-diagonale : {div_res['avg_corr']:.0%}")

    # ── Rendu selon le mode ─────────────────────────────────────────────────

    # Bandeau allocation factuelle + profil
    if _mode_risk == "compare":
        _bp1, _bp2 = st.columns(2)
        with _bp1:
            if _linesA:
                _ps_A, _pd_A, _pa_A, _prof_A = _calc_allocation_3poches(_linesA, "CONTRACT_FUNDS_DF_A")
                with st.container(border=True):
                    st.markdown(f"### Client — {_prof_A}")
                    st.caption(
                        f"🛡️ {_ps_A:.0f}% fonds euros · 🔵 {_pd_A:.0f}% obligations "
                        f"· 🔶 {_pa_A:.0f}% actions & diversifiés"
                    )
        with _bp2:
            if _linesB:
                _ps_B, _pd_B, _pa_B, _prof_B = _calc_allocation_3poches(_linesB, "CONTRACT_FUNDS_DF_B")
                with st.container(border=True):
                    st.markdown(f"### Cabinet — {_prof_B}")
                    st.caption(
                        f"🛡️ {_ps_B:.0f}% fonds euros · 🔵 {_pd_B:.0f}% obligations "
                        f"· 🔶 {_pa_B:.0f}% actions & diversifiés"
                    )
    elif _mode_risk == "client":
        if _linesA:
            _ps_A, _pd_A, _pa_A, _prof_A = _calc_allocation_3poches(_linesA, "CONTRACT_FUNDS_DF_A")
            with st.container(border=True):
                st.markdown(f"### Client — {_prof_A}")
                st.caption(
                    f"🛡️ {_ps_A:.0f}% fonds euros · 🔵 {_pd_A:.0f}% obligations "
                    f"· 🔶 {_pa_A:.0f}% actions & diversifiés"
                )
    else:
        if _linesB:
            _ps_B, _pd_B, _pa_B, _prof_B = _calc_allocation_3poches(_linesB, "CONTRACT_FUNDS_DF_B")
            with st.container(border=True):
                st.markdown(f"### Cabinet — {_prof_B}")
                st.caption(
                    f"🛡️ {_ps_B:.0f}% fonds euros · 🔵 {_pd_B:.0f}% obligations "
                    f"· 🔶 {_pa_B:.0f}% actions & diversifiés"
                )

    # Expander profils de risque
    with st.expander("ℹ️ Comprendre les profils de risque", expanded=False):
        _profil_data = pd.DataFrame({
            "Profil": ["Sécuritaire", "Prudent", "Équilibré", "Dynamique", "Offensif"],
            "Fonds euros": ["70–100%", "50–70%", "30–50%", "10–30%", "0–10%"],
            "Obligations": ["0–20%", "15–30%", "10–20%", "5–15%", "0–10%"],
            "Actions & diversifiés": ["0–10%", "10–30%", "30–60%", "60–80%", "80–100%"],
        })
        st.dataframe(_profil_data, hide_index=True, use_container_width=True)
        st.caption(
            "Ces profils sont indicatifs et correspondent aux grilles couramment utilisées "
            "par les assureurs et les CGP. Le choix du profil dépend de votre horizon "
            "d'investissement, de votre sensibilité au risque et de votre situation personnelle."
        )

    # Jauges de diversification
    if _mode_risk == "compare":
        _gp1, _gp2 = st.columns(2)
        with _gp1:
            _render_gauge_div(_div_A, "Client")
        with _gp2:
            _render_gauge_div(_div_B, "Cabinet")
    elif _mode_risk == "client":
        _render_gauge_div(_div_A, "Client")
    else:
        _render_gauge_div(_div_B, "Cabinet")

    st.caption(
        "Ce score mesure la corrélation entre vos supports UC (actions et obligations), "
        "hors fonds euros. Un score élevé signifie qu'en cas de baisse d'un marché, "
        "vos autres supports ne suivent pas la même trajectoire — votre capital est mieux protégé."
    )

    # Mini-visuel corrélations + alertes doublons nominatives
    def _render_corr_mini(lines: List, euro_rate: float) -> None:
        """Affiche les paires de corrélation avec dots colorés (max 6, triées desc)."""
        if not lines or len(lines) < 2:
            return
        try:
            _cm = correlation_matrix_from_lines(lines, euro_rate)
            if _cm is None or _cm.empty:
                return
            _pairs_cm = []
            _cols_cm = list(_cm.columns)
            for _i_cm in range(len(_cols_cm)):
                for _j_cm in range(_i_cm + 1, len(_cols_cm)):
                    _cv = float(_cm.iloc[_i_cm, _j_cm])
                    if not pd.isna(_cv):
                        _pairs_cm.append((_cols_cm[_i_cm], _cols_cm[_j_cm], _cv))
            _pairs_cm.sort(key=lambda x: -abs(x[2]))
            _pairs_cm = _pairs_cm[:6]
            if not _pairs_cm:
                return
            st.caption("**Corrélations entre fonds**")
            for _n1_cm, _n2_cm, _cv_cm in _pairs_cm:
                _dot = "🟢" if abs(_cv_cm) < 0.40 else ("🔴" if abs(_cv_cm) > 0.80 else "🟡")
                st.caption(f"{_dot} **{str(_n1_cm)[:20]}** / **{str(_n2_cm)[:20]}** — {_cv_cm:.0%}")
        except Exception:
            pass

    if _mode_risk == "compare":
        _cc1, _cc2 = st.columns(2)
        with _cc1:
            if _linesA:
                _render_corr_mini(_linesA, _euro_rate_A)
            if _div_A and _div_A["doublons"]:
                for _ni_al, _nj_al, _c_al in _div_A["doublons"]:
                    st.warning(
                        f"⚠️ {_ni_al} et {_nj_al} ont un comportement identique "
                        f"(corrélation {_c_al:.0%}). Ces deux lignes génèrent des frais de gestion "
                        f"distincts pour une exposition identique — c'est un doublon à optimiser."
                    )
                _n_red_A = _div_A["n_lines"] - _div_A["n_effective"]
                st.caption(
                    f"💰 Ces {_n_red_A} ligne(s) redondante(s) représentent des frais de gestion "
                    f"supportés sans bénéfice de diversification. Consolider ces positions "
                    f"permettrait de réduire les coûts tout en conservant la même exposition."
                )
            elif _div_A:
                st.caption(
                    "✅ Le portefeuille Client ne présente aucun doublon — "
                    "chaque ligne apporte une exposition distincte."
                )
        with _cc2:
            if _linesB:
                _render_corr_mini(_linesB, _euro_rate_B)
            if _div_B and _div_B["doublons"]:
                for _ni_al, _nj_al, _c_al in _div_B["doublons"]:
                    st.warning(
                        f"⚠️ {_ni_al} et {_nj_al} ont un comportement identique "
                        f"(corrélation {_c_al:.0%}). Ces deux lignes génèrent des frais de gestion "
                        f"distincts pour une exposition identique — c'est un doublon à optimiser."
                    )
                _n_red_B = _div_B["n_lines"] - _div_B["n_effective"]
                st.caption(
                    f"💰 Ces {_n_red_B} ligne(s) redondante(s) représentent des frais de gestion "
                    f"supportés sans bénéfice de diversification. Consolider ces positions "
                    f"permettrait de réduire les coûts tout en conservant la même exposition."
                )
            elif _div_B:
                st.caption(
                    "✅ Le portefeuille Cabinet ne présente aucun doublon — "
                    "chaque ligne apporte une exposition distincte."
                )
    else:
        _lines_s2 = _linesA if _mode_risk == "client" else _linesB
        _div_s1 = _div_A if _mode_risk == "client" else _div_B
        _eur_s2 = _euro_rate_A if _mode_risk == "client" else _euro_rate_B
        if _lines_s2:
            _render_corr_mini(_lines_s2, _eur_s2)
        if _div_s1 and _div_s1["doublons"]:
            for _ni_al, _nj_al, _c_al in _div_s1["doublons"]:
                st.warning(
                    f"⚠️ {_ni_al} et {_nj_al} ont un comportement identique "
                    f"(corrélation {_c_al:.0%}). Ces deux lignes génèrent des frais de gestion "
                    f"distincts pour une exposition identique — c'est un doublon à optimiser."
                )
            _n_red = _div_s1["n_lines"] - _div_s1["n_effective"]
            st.caption(
                f"💰 Ces {_n_red} ligne(s) redondante(s) représentent des frais de gestion "
                f"supportés sans bénéfice de diversification. Consolider ces positions "
                f"permettrait de réduire les coûts tout en conservant la même exposition."
            )

    # Narratif (1 phrase — compare uniquement)
    if _mode_risk == "compare" and _div_A and _div_B:
        _eff_A = _div_A["n_effective"]
        _eff_B = _div_B["n_effective"]
        if _eff_B > _eff_A and _eff_A > 0:
            _ratio_d = _eff_B / _eff_A
            _ratio_txt = (
                f"{_ratio_d:.0f}× plus de sources" if _ratio_d >= 2.0
                else f"{(_ratio_d - 1) * 100:.0f}% de sources en plus"
            )
            st.success(
                f"💬 La proposition offre {_ratio_txt} de diversification réelles "
                f"({_eff_B} contre {_eff_A})."
            )
        elif _eff_A > _eff_B and _eff_B > 0:
            st.warning(
                f"💬 Portefeuille actuel mieux diversifié "
                f"({_eff_A} sources réelles contre {_eff_B} pour la proposition)."
            )
        else:
            st.info(
                f"💬 Diversification comparable entre les deux portefeuilles "
                f"({_eff_A} sources réelles)."
            )

    # ── Ratios rendement / risque ────────────────────────────────────────────
    _VS_DIV = "<div style='text-align:center; padding-top:1.4rem; font-size:1.1rem; color:#aaa;'>vs</div>"

    def _render_single_ratios_narrative(sharpe, sortino, ba, euro_rate, val, risk, label):
        """Affiche les ratios en langage client pour un portefeuille seul."""
        if sharpe is not None:
            with st.container(border=True):
                _sh_eur = sharpe * 1000
                if _sh_eur >= 1000:
                    _sh_verdict = "✅ Excellent"
                elif _sh_eur >= 500:
                    _sh_verdict = "✅ Bon"
                else:
                    _sh_verdict = "⚠️ À améliorer"
                st.metric(
                    "📈 Rendement / Risque",
                    f"{_sh_eur:,.0f} €",
                    help="Pour 1 000 € de risque (volatilité), rendement net du taux sans risque."
                )
                st.caption(
                    f"{_sh_verdict} — Pour 1 000 € de risque, ce portefeuille génère "
                    f"**{_sh_eur:,.0f} €** de rendement supplémentaire."
                )
        if sortino and sortino > 0:
            with st.container(border=True):
                _so_eur = 1.0 / sortino
                if _so_eur < 0.7:
                    _so_verdict = "✅ Bonne protection baissière"
                elif _so_eur < 1.0:
                    _so_verdict = "Protection correcte"
                else:
                    _so_verdict = "⚠️ Vulnérable aux baisses"
                st.metric(
                    "🛡️ Protection baissière",
                    f"{_so_eur:,.2f} €",
                    help="Perte € en phase de baisse pour 1 € de rendement (plus bas = mieux)."
                )
                st.caption(
                    f"{_so_verdict} — Pour 1 € de rendement, perte de "
                    f"**{_so_eur:,.2f} €** en phase de baisse."
                )
        _al = ba["alpha_pct"] if ba else None
        if _al is not None and abs(_al) > 0.5:
            with st.container(border=True):
                _al_sign = "+" if _al > 0 else ""
                _al_verdict = "✅ Valeur ajoutée" if _al > 0.5 else "⚠️ Sous-performance vs indice"
                st.metric(
                    "⭐ Valeur ajoutée (Alpha)",
                    f"{_al_sign}{_al:.2f}% / an",
                    help="Surperformance annualisée vs l'indice de référence à risque égal."
                )
                st.caption(
                    f"{_al_verdict} — **{_al_sign}{_al:.2f}%** par an vs l'indice à risque équivalent."
                )

    st.markdown("---")
    st.markdown("##### 💡 Rendement, risque & valeur ajoutée")

    if _mode_risk == "compare" and (_sharpe_A is not None or _sharpe_B is not None):
        with st.container(border=True):
            st.markdown("**📈 Le rendement justifie-t-il le risque pris ?**")
            _col_a1, _col_vs1, _col_b1 = st.columns([5, 1, 5])
            with _col_a1:
                if _sharpe_A is not None:
                    st.metric("🧍 Client", f"{_sharpe_A * 1000:,.0f} €", delta="pour 1 000 € de risque")
                else:
                    st.caption("🧍 Client — données manquantes")
            with _col_vs1:
                st.markdown(_VS_DIV, unsafe_allow_html=True)
            with _col_b1:
                if _sharpe_B is not None:
                    st.metric("🏢 Cabinet", f"{_sharpe_B * 1000:,.0f} €", delta="pour 1 000 € de risque")
                else:
                    st.caption("🏢 Cabinet — données manquantes")

        with st.container(border=True):
            st.markdown("**🛡️ Comment le portefeuille résiste-t-il aux baisses ?**")
            _col_a2, _col_vs2, _col_b2 = st.columns([5, 1, 5])
            with _col_a2:
                if _sortino_A and _sortino_A > 0:
                    st.metric("🧍 Client", f"{1.0 / _sortino_A:,.2f} €", delta="pour 1 € de rendement")
                else:
                    st.caption("🧍 Client — données manquantes")
            with _col_vs2:
                st.markdown(_VS_DIV, unsafe_allow_html=True)
            with _col_b2:
                if _sortino_B and _sortino_B > 0:
                    st.metric("🏢 Cabinet", f"{1.0 / _sortino_B:,.2f} €", delta="pour 1 € de rendement")
                else:
                    st.caption("🏢 Cabinet — données manquantes")

        with st.container(border=True):
            st.markdown("**⭐ L'allocation crée-t-elle de la valeur ?**")
            _col_a3, _col_vs3, _col_b3 = st.columns([5, 1, 5])
            _al_A_v = _ba_A["alpha_pct"] if _ba_A else None
            _al_B_v = _ba_B["alpha_pct"] if _ba_B else None
            with _col_a3:
                if _al_A_v is not None:
                    st.metric("🧍 Client", f"{_al_A_v:+.2f}% / an")
                else:
                    st.caption("🧍 Client — données manquantes")
            with _col_vs3:
                st.markdown(_VS_DIV, unsafe_allow_html=True)
            with _col_b3:
                if _al_B_v is not None:
                    st.metric("🏢 Cabinet", f"{_al_B_v:+.2f}% / an")
                else:
                    st.caption("🏢 Cabinet — données manquantes")

        # Verdict global
        _wins_v = []
        _losses_v = []
        if _sharpe_A is not None and _sharpe_B is not None:
            if _sharpe_B > _sharpe_A + 0.05:
                _sh_gain = ((_sharpe_B / _sharpe_A) - 1) * 100 if _sharpe_A > 0 else 0
                _wins_v.append(f"rendement/risque {abs(_sh_gain):.0f}% plus efficace")
            elif _sharpe_A > _sharpe_B + 0.05:
                _losses_v.append("rendement/risque")
        if _sortino_A and _sortino_B and _sortino_A > 0 and _sortino_B > 0:
            if _sortino_B > _sortino_A + 0.05:
                _so_gain = (1 - (1 / _sortino_B) / (1 / _sortino_A)) * 100
                _wins_v.append(f"{abs(_so_gain):.0f}% de pertes en moins dans les baisses")
            elif _sortino_A > _sortino_B + 0.05:
                _losses_v.append("protection baissière")
        _al_A_vd = _ba_A["alpha_pct"] if _ba_A else None
        _al_B_vd = _ba_B["alpha_pct"] if _ba_B else None
        if _al_B_vd is not None and _al_A_vd is not None and _al_B_vd > _al_A_vd + 0.5:
            _wins_v.append(f"alpha supérieur de {_al_B_vd - _al_A_vd:.1f} pts/an")
        if len(_wins_v) >= 2:
            st.success(f"✅ **La proposition cabinet est plus efficace** : {', '.join(_wins_v)}.")
        elif len(_wins_v) == 1:
            st.info(
                f"Amélioration côté Cabinet : {_wins_v[0]}."
                + (f" Point d'attention : {', '.join(_losses_v)}." if _losses_v else "")
            )
        elif _losses_v:
            st.warning(f"⚠️ Le portefeuille actuel fait mieux sur : {', '.join(_losses_v)}.")
        else:
            st.info("Profils de risque comparables.")

    elif _mode_risk == "client" and _sharpe_A is not None:
        _render_single_ratios_narrative(
            _sharpe_A, _sortino_A, _ba_A, _euro_rate_A,
            st.session_state.get("VALEUR_PORTF_A", 0), _risk_A, "🧍 Client"
        )
    elif _mode_risk == "valority" and _sharpe_B is not None:
        _render_single_ratios_narrative(
            _sharpe_B, _sortino_B, _ba_B, _euro_rate_B,
            st.session_state.get("VALEUR_PORTF_B", 0), _risk_B, "🏢 Cabinet"
        )

    # ── Expanders ─────────────────────────────────────────────────────────────

    # ── Expander 1 : Détail doublons ──────────────────────────────────────
    if _mode_risk == "compare":
        _has_doublons = bool((_div_A and _div_A["doublons"]) or (_div_B and _div_B["doublons"]))
    elif _mode_risk == "client":
        _has_doublons = bool(_div_A and _div_A["doublons"])
    else:
        _has_doublons = bool(_div_B and _div_B["doublons"])

    if _has_doublons:
        with st.expander("🔬 Détail de l'analyse", expanded=False):
            if _mode_risk == "compare":
                for _lns_d, _dres_d, _eur_d, _lbl_d in [
                    (_linesA, _div_A, _euro_rate_A, "🧍 Client"),
                    (_linesB, _div_B, _euro_rate_B, "🏢 Cabinet"),
                ]:
                    if _dres_d is None or not _dres_d["doublons"]:
                        continue
                    st.markdown(f"**{_lbl_d}**")
                    _render_doublon_detail(_dres_d, _lns_d, _eur_d)
            else:
                _lns_d3 = _linesA if _mode_risk == "client" else _linesB
                _dres_d3 = _div_A if _mode_risk == "client" else _div_B
                _eur_d3 = _euro_rate_A if _mode_risk == "client" else _euro_rate_B
                if _dres_d3:
                    _render_doublon_detail(_dres_d3, _lns_d3, _eur_d3)

    # ── Expander 2 : Ratios techniques ────────────────────────────────────
    with st.expander("📐 Détails techniques — Ratios", expanded=False):
        st.caption("Données brutes à destination du conseiller.")
        _rows_tech = []
        if _sharpe_A is not None or _sharpe_B is not None:
            _rows_tech.append({
                "Indicateur": "Sharpe",
                "Client": f"{_sharpe_A:.2f}" if _sharpe_A is not None else "—",
                "Cabinet": f"{_sharpe_B:.2f}" if _sharpe_B is not None else "—",
            })
        if _sortino_A is not None or _sortino_B is not None:
            _rows_tech.append({
                "Indicateur": "Sortino",
                "Client": f"{_sortino_A:.2f}" if _sortino_A is not None else "—",
                "Cabinet": f"{_sortino_B:.2f}" if _sortino_B is not None else "—",
            })
        if _ba_A or _ba_B:
            _bench_lbl_tech = st.session_state.get("BENCHMARK_LABEL", "MSCI World")
            _rows_tech.append({
                "Indicateur": f"Bêta (vs {_bench_lbl_tech})",
                "Client": f"{_ba_A['beta']:.2f}" if _ba_A else "—",
                "Cabinet": f"{_ba_B['beta']:.2f}" if _ba_B else "—",
            })
            _rows_tech.append({
                "Indicateur": "Alpha",
                "Client": f"{_ba_A['alpha_pct']:+.2f}%" if _ba_A else "—",
                "Cabinet": f"{_ba_B['alpha_pct']:+.2f}%" if _ba_B else "—",
            })
        if _risk_A or _risk_B:
            _rows_tech.append({
                "Indicateur": "Volatilité ann.",
                "Client": f"{_risk_A['vol_ann_pct']:.2f}%" if _risk_A else "—",
                "Cabinet": f"{_risk_B['vol_ann_pct']:.2f}%" if _risk_B else "—",
            })
            _rows_tech.append({
                "Indicateur": "Max drawdown",
                "Client": f"{_risk_A['max_dd_pct']:.2f}%" if _risk_A else "—",
                "Cabinet": f"{_risk_B['max_dd_pct']:.2f}%" if _risk_B else "—",
            })
        if _rows_tech:
            st.dataframe(pd.DataFrame(_rows_tech), hide_index=True, use_container_width=True)
        st.caption(f"_Taux sans risque : fonds euros du contrat. Indice : {st.session_state.get('BENCHMARK_LABEL', 'MSCI World')}._")

    with st.expander("ℹ️ Comprendre ces indicateurs", expanded=False):
        st.markdown("""
**Diversification effective** — Nombre de fonds qui apportent réellement une exposition nouvelle. Un fonds corrélé à plus de 80% avec un autre ne réduit pas le risque global — il double les frais sans améliorer la protection.

**Rendement par unité de risque** — Mesure combien de rendement supplémentaire votre portefeuille génère au-delà du fonds euros (sans risque), pour chaque unité de volatilité subie. Plus c'est haut, mieux c'est — 1 000 € pour 1 000 € de risque est un bon seuil.

**Protection en phase de baisse** — Compare les gains en période haussière aux pertes en période baissière. Un portefeuille bien construit perd moins qu'il ne gagne.

**Valeur ajoutée de l'allocation (Alpha)** — Mesure si l'allocation fait mieux qu'un simple ETF indiciel avec le même niveau de risque. Un alpha positif signifie que le choix des fonds crée de la valeur au-delà du marché.

**Ratio de Sharpe / Sortino / Bêta** — Détails dans l'expander "Détails techniques" ci-dessus.

_Taux sans risque utilisé : le taux du fonds euros du contrat, car c'est l'alternative naturelle du client en assurance-vie._
""")

    # ------------------------------------------------------------
    # Améliorer cette allocation — passerelle vers le builder
    # ------------------------------------------------------------
    st.markdown("---")
    with st.container(border=True):
        st.markdown("##### 🔧 Améliorer cette allocation")
        st.caption("Ouvrez le constructeur de portefeuille avec vos fonds pré-chargés et des suggestions d'amélioration.")
        _improve_col1, _improve_col2 = st.columns([3, 1])
        with _improve_col1:
            _improve_target = st.radio(
                "Portefeuille à améliorer",
                ["🧍 Client", "🏢 Cabinet"],
                horizontal=True,
                key="improve_target",
            )
        with _improve_col2:
            if st.button("🔧 Ouvrir le constructeur", type="primary", key="btn_improve_alloc"):
                _is_client_improve = "Client" in _improve_target
                _lines_to_improve = st.session_state.get("A_lines" if _is_client_improve else "B_lines", [])
                _contract = st.session_state.get("CONTRACT_LABEL_A" if _is_client_improve else "CONTRACT_LABEL_B", "")

                # Pré-charger les fonds dans le builder
                _actions_imp = []
                _bonds_imp = []
                _funds_df_improve = st.session_state.get(
                    "CONTRACT_FUNDS_DF_A" if _is_client_improve else "CONTRACT_FUNDS_DF_B",
                    pd.DataFrame()
                )
                for _ln_imp in _lines_to_improve:
                    _isin_imp = str(_ln_imp.get("isin", "")).upper()
                    if _isin_imp in ("EUROFUND", "STRUCTURED"):
                        continue
                    _cat_imp = ""
                    if not _funds_df_improve.empty:
                        _match_imp = _funds_df_improve[_funds_df_improve["isin"].str.upper() == _isin_imp]
                        if not _match_imp.empty:
                            _cat_imp = str(_match_imp.iloc[0].get("category", "")).upper()
                    if _is_bond_category(_cat_imp):
                        _bonds_imp.append(_isin_imp)
                    else:
                        _actions_imp.append(_isin_imp)

                st.session_state["PP_SELECTED_ACTIONS"] = _actions_imp
                st.session_state["PP_SELECTED_BONDS"] = _bonds_imp
                st.session_state["PP_CONTRACT_LABEL"] = _contract
                st.session_state["PP_IMPROVE_MODE"] = True
                st.session_state["PP_IMPROVE_SOURCE"] = "client" if _is_client_improve else "cabinet"
                st.session_state["PP_IMPROVE_LINES"] = _lines_to_improve
                # Passer le budget = valeur actuelle du portefeuille
                _val_improve = float(st.session_state.get(
                    "_LAST_VAL_A" if _is_client_improve else "_LAST_VAL_B", 0
                ) or 0)
                if _val_improve > 0:
                    st.session_state["PP_BUDGET"] = _val_improve
                # Passer les dates d'analyse
                _improve_dates = [ln.get("buy_date") for ln in _lines_to_improve if ln.get("buy_date")]
                if _improve_dates:
                    st.session_state["PP_IMPROVE_START_DATE"] = min(pd.Timestamp(d) for d in _improve_dates)
                    st.session_state["PP_IMPROVE_END_DATE"] = pd.Timestamp.today().normalize()
                st.success("✅ Portefeuille chargé dans le constructeur. Cliquez sur **🏗️ Construction optimisée** dans le menu à gauche pour ajuster l'allocation.")

    # ------------------------------------------------------------
    # Tables positions
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📋 Positions & composition")
    if show_client:
        positions_table("Portefeuille 1 — Client", "A_lines")
    if show_valority:
        positions_table("Portefeuille 2 — Cabinet", "B_lines")


    def _render_portfolio_pie(port_key: str, title: str):
        df_positions = build_positions_dataframe(port_key)
        if df_positions.empty:
            st.info(f"{title} : Données indisponibles.")
            return
        df_pie = _prepare_pie_df(df_positions)
        if df_pie.empty:
            st.info(f"{title} : Données indisponibles.")
            return
        # Palette premium Bleu Nuit / Acier / Or / complémentaires
        _WM_PALETTE = [
            "#1B2A4A", "#4A6FA5", "#C9A84C", "#2E86AB",
            "#5C4033", "#6B8F71", "#8B5E3C", "#A0A0A0",
        ]
        df_pie["color"] = [
            _WM_PALETTE[i % len(_WM_PALETTE)] for i in range(len(df_pie))
        ]
        df_pie["Part_frac"] = df_pie["Part %"] / 100.0
        donut = (
            alt.Chart(df_pie)
            .mark_arc(innerRadius=60, outerRadius=115)
            .encode(
                theta=alt.Theta("Valeur actuelle €:Q", stack=True),
                color=alt.Color(
                    "Nom:N",
                    scale=alt.Scale(
                        domain=df_pie["Nom"].tolist(),
                        range=_WM_PALETTE[: len(df_pie)],
                    ),
                    legend=alt.Legend(title="Support", orient="right"),
                ),
                tooltip=[
                    alt.Tooltip("Nom:N", title="Support"),
                    alt.Tooltip("Valeur actuelle €:Q", title="Valeur", format=",.0f"),
                    alt.Tooltip("Part %:Q", title="Part %", format=".1f"),
                ],
            )
            .properties(title="", height=260)
        )
        st.altair_chart(donut, use_container_width=True)
        st.dataframe(
            df_pie[["Nom", "Valeur actuelle €", "Part %"]].style.format(
                {
                    "Valeur actuelle €": to_eur,
                    "Part %": "{:,.2f}%".format,
                }
            ),
            hide_index=True,
            use_container_width=True,
        )

    if show_client and show_valority:
        col_a, col_b = st.columns(2)
        with col_a:
            st.markdown("#### Portefeuille Client")
            _render_portfolio_pie("A_lines", "Portefeuille Client")
        with col_b:
            st.markdown("#### Portefeuille Cabinet")
            _render_portfolio_pie("B_lines", "Portefeuille Cabinet")
    elif show_client:
        st.markdown("#### Portefeuille Client")
        _render_portfolio_pie("A_lines", "Portefeuille Client")
    elif show_valority:
        st.markdown("#### Portefeuille Cabinet")
        _render_portfolio_pie("B_lines", "Portefeuille Cabinet")

    # ── Transparence frais & Clean Shares ────────────────────────
    st.markdown("---")
    with st.expander("🔍 Transparence des frais & Clean Shares", expanded=False):
        st.caption(
            "Décomposition des frais par couche : frais du support (TER, intégré dans la VL) "
            "et frais de l'enveloppe contrat (déduits de la performance)."
        )
        # Identifier les fonds UC du portefeuille actif
        _cs_port_key = "A_lines" if show_client else "B_lines"
        _cs_contract = st.session_state.get(
            "CONTRACT_LABEL_A" if show_client else "CONTRACT_LABEL_B", ""
        )
        _cs_assureur = CONTRACTS_REGISTRY.get(_cs_contract, {}).get("assureur", "—")
        # Clean Shares par assureur (référentiel simplifié)
        _CLEAN_SHARE_ASSUREURS = {"Spirica"}  # Spirica permet les parts I sur certains fonds
        _cs_lines = [
            ln for ln in st.session_state.get(_cs_port_key, [])
            if str(ln.get("isin", "")).upper() not in ("EUROFUND", "STRUCTURED")
        ]
        if not _cs_lines:
            st.info("Aucun fonds UC dans le portefeuille sélectionné.")
        else:
            _frais_rows = []
            for _ln in _cs_lines:
                _name = (_ln.get("name") or _ln.get("isin") or "—")[:35]
                _ter = float(_ln.get("fee_uc_pct") or 0.0)
                _contrat_fee = float(_ln.get("fee_contract_pct") or
                                     st.session_state.get("FEE_A" if show_client else "FEE_B", 0.6))
                _total = _ter + _contrat_fee
                _frais_rows.append({
                    "Support": _name,
                    "TER fonds (%/an)": round(_ter, 2),
                    "Frais contrat (%/an)": round(_contrat_fee, 2),
                    "Coût total (%/an)": round(_total, 2),
                })
            _frais_df = pd.DataFrame(_frais_rows)
            st.dataframe(
                _frais_df.style.format({
                    "TER fonds (%/an)": "{:.2f}%",
                    "Frais contrat (%/an)": "{:.2f}%",
                    "Coût total (%/an)": "{:.2f}%",
                }).background_gradient(subset=["Coût total (%/an)"], cmap="YlOrRd"),
                hide_index=True, use_container_width=True,
            )
            st.caption(
                "Le TER est intégré dans la valeur liquidative publiée. "
                "Les frais contrat sont déduits en sus par l'assureur."
            )
            # Alerte Clean Shares si assureur éligible
            if _cs_assureur in _CLEAN_SHARE_ASSUREURS:
                st.info(
                    f"✨ **Optimisation possible — {_cs_assureur}** : Ce contrat donne "
                    "accès aux parts institutionnelles (parts I / Clean Shares) sur certains "
                    "fonds éligibles, potentiellement **−0.5% à −1%/an de frais de gestion**. "
                    "Vérifiez l'éligibilité de chaque fonds directement auprès de l'assureur."
                )

    # APP – Composition
    def _wrap_label_app(label: str, width: int = 28) -> str:
        if not label:
            return "—"
        return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)

    if show_valority and not dfB.empty:
        st.markdown("---")
        st.subheader("📈 Évolution du portefeuille Cabinet")
        st.caption(
            "Comparaison année par année entre le capital net investi "
            "et la valeur du portefeuille."
        )
        # ── Construire les données année par année ────────────────
        _blines_b = st.session_state.get("B_lines", [])
        _fee_b = float(st.session_state.get("FEE_B", 0.0) or 0.0)
        _M_B_evol = float(st.session_state.get("M_B", 0.0) or 0.0)
        _ONE_B = float(st.session_state.get("ONE_B", 0.0) or 0.0)
        _ONE_B_DATE = st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date())
        _start_b = startB_min if isinstance(startB_min, pd.Timestamp) else TODAY

        _years_b = sorted(set(dfB.index.year))
        _bar_rows: List[Dict[str, Any]] = []

        for _y in _years_b:
            _mask_y = dfB.index.year == _y
            if not _mask_y.any():
                continue
            _val_end = float(dfB.loc[_mask_y, "Valeur"].iloc[-1])
            _year_end = pd.Timestamp(f"{_y}-12-31")

            # 1) Versements initiaux nets
            _net_init = sum(
                float(ln.get("amount_gross", 0)) * (1.0 - _fee_b / 100.0)
                for ln in _blines_b
                if pd.Timestamp(ln.get("buy_date") or _start_b) <= _year_end
            )

            # 2) Versement ponctuel net
            _net_one = 0.0
            if _ONE_B > 0 and pd.Timestamp(_ONE_B_DATE) <= _year_end:
                _net_one = _ONE_B * (1.0 - _fee_b / 100.0)

            # 3) Versements mensuels nets cumulés
            _net_monthly_evol = 0.0
            if _M_B_evol > 0:
                _to_date = min(_year_end, TODAY)
                _n_months = max(
                    0,
                    (_to_date.year - _start_b.year) * 12
                    + (_to_date.month - _start_b.month),
                )
                _net_monthly_evol = _n_months * _M_B_evol * (1.0 - _fee_b / 100.0)

            _net_cumul = _net_init + _net_one + _net_monthly_evol

            _bar_rows.append({
                "Année": str(_y),
                "Capital net investi": round(_net_cumul, 2),
                "Valeur du portefeuille": round(_val_end, 2),
            })

        _bar_df_evol = pd.DataFrame(_bar_rows)

        if not _bar_df_evol.empty:
            _bar_long_evol = _bar_df_evol.melt(
                "Année", var_name="Type", value_name="Montant (€)"
            )
            _evol_chart = (
                alt.Chart(_bar_long_evol)
                .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                .encode(
                    x=alt.X("Année:N", title="Année", axis=alt.Axis(labelAngle=0)),
                    y=alt.Y(
                        "Montant (€):Q",
                        title="Montant (€)",
                        axis=alt.Axis(format=",.0f"),
                        scale=alt.Scale(zero=True),
                    ),
                    color=alt.Color(
                        "Type:N",
                        scale=alt.Scale(
                            domain=["Capital net investi", "Valeur du portefeuille"],
                            range=["#1B2A4A", "#C9A84C"],
                        ),
                        legend=alt.Legend(title=None, orient="bottom"),
                    ),
                    xOffset="Type:N",
                    tooltip=[
                        alt.Tooltip("Année:N"),
                        alt.Tooltip("Type:N"),
                        alt.Tooltip("Montant (€):Q", format=",.2f"),
                    ],
                )
                .properties(height=400)
            )
            st.altair_chart(_evol_chart, use_container_width=True)

            # Récapitulatif chiffré de la dernière année disponible
            _last_row = _bar_df_evol.iloc[-1]
            _pv_total = _last_row["Valeur du portefeuille"] - _last_row["Capital net investi"]
            _pv_pct = (
                _pv_total / _last_row["Capital net investi"] * 100.0
            ) if _last_row["Capital net investi"] > 0 else 0.0
            _cols_legend = st.columns(3)
            with _cols_legend[0]:
                st.metric("Capital net investi", to_eur(_last_row["Capital net investi"]))
            with _cols_legend[1]:
                st.metric("Valeur actuelle", to_eur(_last_row["Valeur du portefeuille"]))
            with _cols_legend[2]:
                st.metric("Plus-value", to_eur(_pv_total), delta=f"{_pv_pct:+.1f}%")
        else:
            st.info("Données insuffisantes pour afficher l'évolution du capital.")

    st.markdown("---")
    st.info("📥 Pour télécharger les rapports (PDF, Présentation, One-pager), rendez-vous dans **Rapports & Présentations** dans le menu de navigation.")

    with st.expander("Aide rapide"):
        st.markdown(
            """
- Dans chaque portefeuille, vous pouvez **soit** ajouter des *fonds recommandés* (onglet dédié),
  **soit** utiliser la *saisie libre* avec ISIN / code.
- Pour le **fonds en euros**, utilisez le symbole **EUROFUND** (taux paramétrable dans la barre de gauche).
- Les frais d'entrée s'appliquent à chaque investissement.
- Le **rendement total** est la performance globale depuis l'origine (valeur actuelle / net investi).
- Le **rendement annualisé** utilise le XIRR (prise en compte des dates et montants).
- En mode **Personnalisé**, vous pouvez affecter précisément les versements mensuels et ponctuels à chaque ligne,
  avec un contrôle automatique de cohérence par rapport aux montants bruts saisis.
            """
        )

    # ------------------------------------------------------------
    # Analyse interne — Corrélation & volatilité (réservé conseiller)
    # ------------------------------------------------------------
    st.markdown("---")
    with st.expander("🔒 Analyse interne — Corrélation, volatilité et profil de risque", expanded=False):
        st.caption(
            "Section réservée au conseiller : analyse technique basée sur les valeurs liquidatives "
            "(corrélations, volatilités, drawdown)."
        )

        # FIXED: use per-portfolio euro rates instead of shared EURO_RATE_PREVIEW (Résidu Bug 4)
        euro_rate_A = st.session_state.get("EURO_RATE_A", 2.0)
        euro_rate_B = st.session_state.get("EURO_RATE_B", 2.5)
        linesA = st.session_state.get("A_lines", [])
        linesB = st.session_state.get("B_lines", [])

        # Portefeuille Client
        if show_client:
            st.markdown("### Portefeuille 1 — Client")
            corrA = correlation_matrix_from_lines(linesA, euro_rate_A, start_date=_start_A)
            volA = volatility_table_from_lines(linesA, euro_rate_A, start_date=_start_A)
            riskA = portfolio_risk_stats(linesA, euro_rate_A, start_date=_start_A, fee_pct=st.session_state.get("FEE_A", 0.0))

            if corrA.empty and volA.empty:
                st.info("Pas assez d'historique ou de lignes pour analyser ce portefeuille.")
            else:
                if riskA is not None:
                    c1, c2 = st.columns(2)
                    with c1:
                        st.metric(
                            "Volatilité annuelle estimée",
                            f"{riskA['vol_ann_pct']:.2f} %",
                        )
                    with c2:
                        st.metric(
                            "Max drawdown (historique sur la période)",
                            f"{riskA['max_dd_pct']:.2f} %",
                        )

                if not volA.empty:
                    st.markdown("**Volatilité par ligne**")
                    st.dataframe(
                        volA.style.format(
                            {
                                "Écart-type quotidien %": "{:,.2f}%".format,
                                "Volatilité annuelle %": "{:,.2f}%".format,
                            }
                        ),
                        use_container_width=True,
                    )

                if not corrA.empty:
                    chartA = _corr_heatmap_chart(corrA, "Corrélation des lignes — Portefeuille Client")
                    if chartA is not None:
                        st.altair_chart(chartA, use_container_width=True)
                    # Alerte si corrélation excessive
                    _corrA_vals = corrA.where(
                        ~pd.DataFrame(
                            [[i == j for j in range(len(corrA.columns))]
                             for i in range(len(corrA.index))],
                            index=corrA.index, columns=corrA.columns
                        )
                    )
                    if not _corrA_vals.isna().all().all() and float(_corrA_vals.max().max()) > 0.90:
                        st.caption(
                            "⚠️ Information : Certains fonds sont très fortement corrélés (>0.90), "
                            "limitant l'impact réel de la diversification."
                        )

        if show_client and show_valority:
            st.markdown("---")

        if show_valority:
            st.markdown("### Portefeuille 2 — Cabinet")
            corrB = correlation_matrix_from_lines(linesB, euro_rate_B, start_date=_start_B)
            volB = volatility_table_from_lines(linesB, euro_rate_B, start_date=_start_B)
            riskB = portfolio_risk_stats(linesB, euro_rate_B, start_date=_start_B, fee_pct=st.session_state.get("FEE_B", 0.0))

            if corrB.empty and volB.empty:
                st.info("Pas assez d'historique ou de lignes pour analyser ce portefeuille.")
            else:
                if riskB is not None:
                    c1, c2 = st.columns(2)
                    with c1:
                        st.metric(
                            "Volatilité annuelle estimée",
                            f"{riskB['vol_ann_pct']:.2f} %",
                        )
                    with c2:
                        st.metric(
                            "Max drawdown (historique sur la période)",
                            f"{riskB['max_dd_pct']:.2f} %",
                        )

                if not volB.empty:
                    st.markdown("**Volatilité par ligne**")
                    st.dataframe(
                        volB.style.format(
                            {
                                "Écart-type quotidien %": "{:,.2f}%".format,
                                "Volatilité annuelle %": "{:,.2f}%".format,
                            }
                        ),
                        use_container_width=True,
                    )

                if not corrB.empty:
                    chartB = _corr_heatmap_chart(corrB, "Corrélation des lignes — Portefeuille Cabinet")
                    if chartB is not None:
                        st.altair_chart(chartB, use_container_width=True)
                    # Alerte si corrélation excessive
                    _corrB_vals = corrB.where(
                        ~pd.DataFrame(
                            [[i == j for j in range(len(corrB.columns))]
                             for i in range(len(corrB.index))],
                            index=corrB.index, columns=corrB.columns
                        )
                    )
                    if not _corrB_vals.isna().all().all() and float(_corrB_vals.max().max()) > 0.90:
                        st.caption(
                            "⚠️ Information : Certains fonds sont très fortement corrélés (>0.90), "
                            "limitant l'impact réel de la diversification."
                        )

def simulate_market_crash(
    portfolio_value: float,
    uc_pct: float,
    euro_pct: float,
    beta: float = 0.85,
    crash_magnitude: float = 0.20,
    duration_euro: float = 3.0,
    rate_shock: float = 0.01,
) -> Dict[str, float]:
    """
    Simule l'impact d'un krach combiné (marché + taux) sur un portefeuille.
    Retourne les pertes estimées par composante et la valeur post-choc.
    """
    perte_uc     = portfolio_value * uc_pct   * beta   * (-crash_magnitude)
    perte_euro   = portfolio_value * euro_pct * (-duration_euro * rate_shock)
    perte_totale = perte_uc + perte_euro
    return {
        "perte_uc":          perte_uc,
        "perte_euro":        perte_euro,
        "perte_totale":      perte_totale,
        "valeur_post_choc":  portfolio_value + perte_totale,
        "pct_perte_total":   (perte_totale / portfolio_value * 100) if portfolio_value > 0 else 0.0,
    }


def run_comparator():
    render_app(run_page_config=False)


# ============================================================
# MODULE FISCALITÉ ASSURANCE-VIE — fonctions de calcul pures
# ============================================================

def calc_quote_part_gains(
    valeur_contrat: float,
    versements_nets: float,
    montant_rachat: float,
) -> float:
    """Quote-part de gains dans le rachat."""
    if valeur_contrat <= 0 or versements_nets >= valeur_contrat:
        return 0.0
    return max(0.0, montant_rachat * (1.0 - versements_nets / valeur_contrat))


def calc_imposition_rachat(
    gains: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
    montant_rachat: float,
    option_ir: bool = False,
    tmi: float = 0.30,
) -> Dict[str, Any]:
    PS_RATE = 0.172
    PFU_RATE = 0.128
    PFL_LOW_RATE = 0.075
    PFL_HIGH_RATE = 0.128
    ABATTEMENT_SEUL = 4_600.0
    ABATTEMENT_COUPLE = 9_200.0
    SEUIL_150K = 150_000.0

    abattement = 0.0
    base_ir = gains
    taux_ir = 0.0
    regime = ""
    tmi_applicable = ""

    _tmi_pct = int(round(tmi * 100))

    if option_ir:
        # Option barème progressif — l'abattement s'applique aussi pour ≥8 ans
        if anciennete_annees >= 8:
            abattement = ABATTEMENT_COUPLE if "Couple" in situation_familiale else ABATTEMENT_SEUL
        base_ir = max(0.0, gains - abattement)
        taux_ir = tmi
        if anciennete_annees < 8:
            regime = f"Option barème IR {_tmi_pct}% (contrat < 8 ans)"
            tmi_applicable = f"Barème progressif {_tmi_pct}% — case 2CH"
        else:
            regime = f"Option barème IR {_tmi_pct}% (contrat ≥8 ans)"
            tmi_applicable = f"Barème progressif {_tmi_pct}% + abattement — case 2CH"
    else:
        if anciennete_annees < 8:
            taux_ir = PFU_RATE
            regime = "PFU 30% (contrat < 8 ans)"
            tmi_applicable = "Flat Tax 12,8% IR + 17,2% PS"
        else:
            abattement = ABATTEMENT_COUPLE if "Couple" in situation_familiale else ABATTEMENT_SEUL
            base_ir = max(0.0, gains - abattement)
            if versements_nets_total <= SEUIL_150K:
                taux_ir = PFL_LOW_RATE
                regime = "PFL 7,5% + PS (contrat ≥8 ans, versements ≤150k€)"
                tmi_applicable = "Prélèvement forfaitaire libératoire 7,5%"
            else:
                taux_ir = PFL_HIGH_RATE
                regime = "PFU 12,8% + PS (contrat ≥8 ans, versements >150k€)"
                tmi_applicable = "Flat Tax 12,8% (versements post-27/09/2017 >150k€)"

    montant_ir = base_ir * taux_ir
    montant_ps = gains * PS_RATE

    return {
        "gains": gains,
        "regime": regime,
        "base_ir": base_ir,
        "abattement_applique": abattement,
        "taux_ir": taux_ir,
        "montant_ir": montant_ir,
        "montant_ps": montant_ps,
        "total_impots": montant_ir + montant_ps,
        "net_percu": montant_rachat - (montant_ir + montant_ps),
        "tmi_applicable": tmi_applicable,
    }


def calc_rachat_depuis_net(
    montant_net_souhaite: float,
    valeur_contrat: float,
    versements_nets: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
) -> Tuple[float, Dict[str, Any]]:
    """Bisection : trouve le brut tel que brut - impôts == net_souhaité."""
    if montant_net_souhaite <= 0:
        return 0.0, calc_imposition_rachat(0.0, anciennete_annees, situation_familiale, versements_nets_total, 0.0)

    lo, hi = montant_net_souhaite, min(valeur_contrat, montant_net_souhaite * 2.5)
    # s'assurer que hi couvre bien le cas
    for _ in range(50):
        g = calc_quote_part_gains(valeur_contrat, versements_nets, hi)
        d = calc_imposition_rachat(g, anciennete_annees, situation_familiale, versements_nets_total, hi)
        if hi - d["total_impots"] >= montant_net_souhaite:
            break
        hi = min(valeur_contrat, hi * 1.5)

    for _ in range(1000):
        mid = (lo + hi) / 2.0
        g = calc_quote_part_gains(valeur_contrat, versements_nets, mid)
        d = calc_imposition_rachat(g, anciennete_annees, situation_familiale, versements_nets_total, mid)
        net = mid - d["total_impots"]
        if abs(net - montant_net_souhaite) < 0.01:
            return mid, d
        if net < montant_net_souhaite:
            lo = mid
        else:
            hi = mid
    return mid, d


def calc_optimisation_abattement(
    valeur_contrat: float,
    versements_nets: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
    abattement_deja_utilise: float = 0.0,
) -> Dict[str, Any]:
    PS_RATE = 0.172
    ABATTEMENT_SEUL = 4_600.0
    ABATTEMENT_COUPLE = 9_200.0

    abattement = ABATTEMENT_COUPLE if "Couple" in situation_familiale else ABATTEMENT_SEUL
    abattement_restant = max(0.0, abattement - abattement_deja_utilise)

    if valeur_contrat <= 0 or versements_nets >= valeur_contrat:
        return {
            "rachat_optimal": 0.0, "gains_realises": 0.0,
            "ir_du": 0.0, "ps_du": 0.0, "economies_ir": 0.0,
            "abattement_utilise": 0.0, "abattement_restant_apres": abattement_restant,
        }

    taux_pv = 1.0 - versements_nets / valeur_contrat
    if taux_pv <= 0:
        return {
            "rachat_optimal": 0.0, "gains_realises": 0.0,
            "ir_du": 0.0, "ps_du": 0.0, "economies_ir": 0.0,
            "abattement_utilise": 0.0, "abattement_restant_apres": abattement_restant,
        }

    rachat_optimal = abattement_restant / taux_pv
    rachat_optimal = min(rachat_optimal, valeur_contrat)
    gains_realises = calc_quote_part_gains(valeur_contrat, versements_nets, rachat_optimal)

    # IR si on avait dépassé l'abattement (scénario sans optimisation)
    SEUIL_150K = 150_000.0
    taux_ir = 0.075 if versements_nets_total <= SEUIL_150K else 0.128
    economies_ir = gains_realises * taux_ir

    return {
        "rachat_optimal": rachat_optimal,
        "gains_realises": gains_realises,
        "ir_du": 0.0,
        "ps_du": gains_realises * PS_RATE,
        "economies_ir": economies_ir,
        "abattement_utilise": gains_realises,
        "abattement_restant_apres": 0.0,
    }


def calc_transmission_990I(
    capital_deces: float,
    nb_beneficiaires: int,
    parts_pct: List[float],
    types_beneficiaires: List[str],
) -> List[Dict[str, Any]]:
    ABATTEMENT_990I = 152_500.0
    TAUX_990I_T1 = 0.20
    TAUX_990I_T2 = 0.3125
    SEUIL_990I_T2 = 700_000.0

    results = []
    for i in range(nb_beneficiaires):
        t = types_beneficiaires[i] if i < len(types_beneficiaires) else "Enfant"
        p = parts_pct[i] if i < len(parts_pct) else 0.0
        part_eur = capital_deces * p / 100.0
        note = ""

        if "Conjoint" in t or "PACS" in t:
            taxe = 0.0
            taxable = 0.0
            abat = part_eur
        else:
            if "Démembré" in t:
                note = (
                    "Le démembrement de la clause bénéficiaire nécessite une évaluation "
                    "actuarielle des droits d'usufruit et nue-propriété selon le barème "
                    "fiscal de l'article 669 CGI (fonction de l'âge de l'usufruitier). "
                    "Cette simulation affiche la taxation globale sur la part démembrée "
                    "comme si elle revenait à un bénéficiaire standard, à affiner avec "
                    "un notaire."
                )
            abat = min(part_eur, ABATTEMENT_990I)
            taxable = max(0.0, part_eur - ABATTEMENT_990I)
            if taxable <= SEUIL_990I_T2:
                taxe = taxable * TAUX_990I_T1
            else:
                taxe = SEUIL_990I_T2 * TAUX_990I_T1 + (taxable - SEUIL_990I_T2) * TAUX_990I_T2

        results.append({
            "index": i + 1,
            "type": t,
            "part_pct": p,
            "part_eur": part_eur,
            "abattement": abat,
            "taxable": taxable if "Conjoint" not in t and "PACS" not in t else 0.0,
            "taxe": taxe,
            "net_recu": part_eur - taxe,
            "note": note,
        })
    return results


def calc_transmission_757B(
    primes_versees_apres_70: float,
    nb_beneficiaires: int,
    parts_pct: List[float],
    types_beneficiaires: List[str],
    capital_deces: float,
    liens_succession: List[str],
) -> List[Dict[str, Any]]:
    ABATTEMENT_757B = 30_500.0
    TAUX_LIEN = {
        "Enfant": 0.20,
        "Frère/Sœur": 0.40,
        "Neveu/Nièce": 0.55,
        "Tiers": 0.60,
        "Conjoint/PACS": 0.0,
    }

    # Parts des non-conjoints pour répartir l'abattement global
    non_conjoint_total_pct = sum(
        parts_pct[i] for i in range(nb_beneficiaires)
        if i < len(types_beneficiaires)
        and "Conjoint" not in types_beneficiaires[i]
        and "PACS" not in types_beneficiaires[i]
    )

    results = []
    for i in range(nb_beneficiaires):
        t = types_beneficiaires[i] if i < len(types_beneficiaires) else "Enfant"
        p = parts_pct[i] if i < len(parts_pct) else 0.0
        lien = liens_succession[i] if i < len(liens_succession) else "Enfant"
        part_eur = capital_deces * p / 100.0

        if "Conjoint" in t or "PACS" in t or lien == "Conjoint/PACS":
            results.append({
                "index": i + 1, "type": t, "part_pct": p, "part_eur": part_eur,
                "primes_part": 0.0, "abattement_part": 0.0,
                "taxable": 0.0, "taxe": 0.0, "net_recu": part_eur,
                "taux_applique": 0.0,
            })
            continue

        # Répartition des primes post-70 ans selon la quote-part
        primes_part = primes_versees_apres_70 * p / max(non_conjoint_total_pct, 1.0)
        # Abattement global 30 500 € réparti proportionnellement
        abat_part = ABATTEMENT_757B * p / max(non_conjoint_total_pct, 1.0)
        taxable = max(0.0, primes_part - abat_part)
        taux = TAUX_LIEN.get(lien, 0.60)
        taxe = taxable * taux

        results.append({
            "index": i + 1, "type": t, "part_pct": p, "part_eur": part_eur,
            "primes_part": primes_part, "abattement_part": abat_part,
            "taxable": taxable, "taxe": taxe, "net_recu": part_eur - taxe,
            "taux_applique": taux,
        })
    return results


# ============================================================
# MODULE FISCALITÉ ASSURANCE-VIE — onglets UI
# ============================================================

def _fmt_eur(v: float) -> str:
    """Formatage monétaire uniforme pour le module fiscal."""
    return f"{v:,.0f} €".replace(",", "\u202f")


def _tab_rachat() -> None:
    # ── Données du contrat — lues depuis le Tableau de bord ──────────
    _date_ouv_r = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
    anciennete_annees = (date.today() - _date_ouv_r).days / 365.25
    valeur_contrat = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
    versements_nets = float(st.session_state.get("tax_versements_nets", 80_000.0))
    versements_nets_total = float(st.session_state.get("tax_versements_nets_total", 80_000.0))
    situation_familiale = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

    with st.container(border=True):
        _ri1, _ri2, _ri3 = st.columns(3)
        with _ri1:
            st.metric("Valeur du contrat", to_eur(valeur_contrat))
        with _ri2:
            st.metric("Versements nets", to_eur(versements_nets))
        with _ri3:
            _anc_disp = anciennete_annees
            st.metric("Ancienneté", f"{_anc_disp:.1f} ans")
        st.caption(f"📅 Ouverture : **{_date_ouv_r.strftime('%d/%m/%Y')}** | Situation : **{situation_familiale}** — _modifiez dans le 📋 Tableau de bord_")

    # ── PS déjà prélevés — calcul automatique part fonds euros ──
    _is_client_tax = st.session_state.get("tax_is_client", True)
    _port_key_tax = "A_lines" if _is_client_tax else "B_lines"
    _tax_lines = st.session_state.get(_port_key_tax, [])
    _total_euro = 0.0
    _total_uc = 0.0
    for _ln in _tax_lines:
        _amt = float(_ln.get("amount_gross", 0.0))
        if str(_ln.get("isin", "")).upper() == "EUROFUND":
            _total_euro += _amt
        else:
            _total_uc += _amt
    _total_investi = _total_euro + _total_uc
    _pct_euro = _total_euro / _total_investi if _total_investi > 0 else 0.0

    st.markdown("#### Montant du rachat")

    _mode_rachat = st.radio(
        "Je connais :",
        ["Le montant brut à racheter", "Le montant net que je veux recevoir"],
        horizontal=True,
        key="tax_rachat_mode",
    )

    if "brut" in _mode_rachat.lower():
        montant_brut = st.number_input(
            "Montant brut du rachat (€)",
            min_value=0.0,
            max_value=float(max(valeur_contrat, 0.01)),
            value=float(st.session_state.get("tax_montant_brut", 10_000.0)),
            step=500.0,
            key="_tax_brut_input",
        )
        st.session_state["tax_montant_brut"] = montant_brut
    else:
        _montant_net_souhaite = st.number_input(
            "Montant net souhaité (€)",
            min_value=0.0,
            max_value=float(max(valeur_contrat, 0.01)),
            value=float(st.session_state.get("tax_montant_net", 9_000.0)),
            step=500.0,
            key="_tax_net_input",
        )
        st.session_state["tax_montant_net"] = _montant_net_souhaite
        montant_brut, _ = calc_rachat_depuis_net(
            _montant_net_souhaite, valeur_contrat, versements_nets,
            anciennete_annees, situation_familiale, versements_nets_total,
        )
        st.caption(f"→ Rachat brut nécessaire : **{to_eur(montant_brut)}**")

    # Calcul et affichage
    if valeur_contrat > 0 and versements_nets > 0 and montant_brut > 0:
        gains = calc_quote_part_gains(valeur_contrat, versements_nets, montant_brut)
        _ps_deja_preleves = gains * _pct_euro * 0.172
        result = calc_imposition_rachat(
            gains, anciennete_annees, situation_familiale,
            versements_nets_total, montant_brut,
        )
        st.session_state["tax_rachat_result"] = result

        capital_pur = montant_brut - gains
        ps_net = max(0.0, result["montant_ps"] - _ps_deja_preleves)
        total_impots = result["montant_ir"] + ps_net
        net_percu = montant_brut - total_impots

        # ── Bandeau régime ──
        if anciennete_annees < 8:
            st.error(f"⏱️ Contrat de {anciennete_annees:.1f} an(s) — PFU 30%")
        else:
            st.success(f"✅ Contrat de {anciennete_annees:.1f} ans — Régime favorable ≥8 ans · {result['tmi_applicable']}")

        # ── 3 métriques principales ──
        _m1, _m2, _m3 = st.columns(3)
        with _m1:
            st.metric("💰 Rachat brut", to_eur(montant_brut))
        with _m2:
            st.metric("🏛️ Impôts & PS", to_eur(total_impots),
                      delta=f"-{total_impots/montant_brut*100:.1f}%" if montant_brut > 0 else None)
        with _m3:
            st.metric("✅ Net perçu", to_eur(net_percu))

        if _ps_deja_preleves > 0:
            st.caption(
                f"ℹ️ PS déjà prélevés sur fonds euros : **{to_eur(_ps_deja_preleves)}** déduits "
                f"(part euros : {_pct_euro:.0%} du portefeuille)"
            )

        # ── Conseil stratégie ≥8 ans ──
        if anciennete_annees >= 8 and gains > 0:
            _abat_total = result["abattement_applique"]
            if _abat_total > 0:
                _brut_sans_ir = montant_brut * (_abat_total / gains) if gains > 0 else montant_brut
                _brut_sans_ir = min(_brut_sans_ir, valeur_contrat)
                _g_lim = calc_quote_part_gains(valeur_contrat, versements_nets, _brut_sans_ir)
                _ps_lim = _g_lim * 0.172
                st.info(
                    f"💡 **Stratégie** : un rachat de **{to_eur(_brut_sans_ir)}** ne génèrerait "
                    f"aucun IR, uniquement **{to_eur(_ps_lim)}** de PS."
                )

        # ── Donut chart (pleine largeur) ──
        st.markdown("##### Décomposition du rachat")
        _abat_gains = min(gains, result["abattement_applique"])
        _donut_parts = []
        if capital_pur > 0:
            _donut_parts.append({"Composante": "Capital restitué (non taxé)", "Montant": capital_pur})
        if _abat_gains > 0:
            _donut_parts.append({"Composante": "Gains exonérés (abattement)", "Montant": _abat_gains})
        if result["montant_ir"] > 0:
            _donut_parts.append({"Composante": "Impôt sur le revenu", "Montant": result["montant_ir"]})
        if ps_net > 0:
            _donut_parts.append({"Composante": "Prélèvements sociaux", "Montant": ps_net})
        if _donut_parts:
            _donut_df = pd.DataFrame(_donut_parts)
            _donut_chart = (
                alt.Chart(_donut_df)
                .mark_arc(innerRadius=60, outerRadius=120)
                .encode(
                    theta=alt.Theta("Montant:Q"),
                    color=alt.Color(
                        "Composante:N",
                        scale=alt.Scale(
                            domain=[p["Composante"] for p in _donut_parts],
                            range=["#2E7D32", "#81C784", "#E53935", "#EF9A9A"][:len(_donut_parts)],
                        ),
                        legend=alt.Legend(title="", orient="bottom"),
                    ),
                    tooltip=[
                        alt.Tooltip("Composante:N"),
                        alt.Tooltip("Montant:Q", format=",.0f", title="€"),
                    ],
                )
                .properties(height=300)
            )
            st.altair_chart(_donut_chart, use_container_width=True)

        # ── Waterfall : Brut → Net (pleine largeur) ──
        st.markdown("##### Brut → Net")
        _wf_data = pd.DataFrame([
            {"Étape": "Rachat brut", "Montant": montant_brut, "Type": "Base"},
            {"Étape": "− IR", "Montant": -result["montant_ir"], "Type": "Impôt"},
            {"Étape": "− PS", "Montant": -ps_net, "Type": "Impôt"},
            {"Étape": "Net perçu", "Montant": net_percu, "Type": "Net"},
        ])
        _wf_chart = (
            alt.Chart(_wf_data)
            .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
            .encode(
                x=alt.X("Étape:N", sort=None, axis=alt.Axis(labelAngle=0), title=""),
                y=alt.Y("Montant:Q", axis=alt.Axis(format=",.0f"), title="€"),
                color=alt.Color(
                    "Type:N",
                    scale=alt.Scale(
                        domain=["Base", "Impôt", "Net"],
                        range=["#1B2A4A", "#E53935", "#2E7D32"],
                    ),
                    legend=None,
                ),
                tooltip=[
                    alt.Tooltip("Étape:N"),
                    alt.Tooltip("Montant:Q", format=",.0f", title="€"),
                ],
            )
            .properties(height=280)
        )
        st.altair_chart(_wf_chart, use_container_width=True)

        # ── Comparatif PFU vs Option barème IR ──
        _tmi_val = float(st.session_state.get("tax_tmi_PERSIST", 0.30))
        _res_ir_opt = calc_imposition_rachat(
            gains, anciennete_annees, situation_familiale,
            versements_nets_total, montant_brut,
            option_ir=True, tmi=_tmi_val,
        )
        _ps_net_ir_opt = max(0.0, _res_ir_opt["montant_ps"] - _ps_deja_preleves)
        _total_ir_opt = _res_ir_opt["montant_ir"] + _ps_net_ir_opt
        _net_ir_opt = montant_brut - _total_ir_opt
        # PFU values already computed above
        _pfu_better = net_percu >= _net_ir_opt
        _tmi_pct_lbl = f"{int(round(_tmi_val * 100))}%"

        st.markdown("#### 🔀 PFU ou Option barème IR ?")
        _col_pfu, _col_ir_c = st.columns(2)
        with _col_pfu:
            _msg_pfu = (
                f"**Flat Tax / PFL**\n\n"
                f"Régime : {result['regime']}\n\n"
                f"IR : **{_fmt_eur(result['montant_ir'])}**\n\n"
                f"PS nets : **{_fmt_eur(ps_net)}**\n\n"
                f"Total prélevé : **{_fmt_eur(total_impots)}**\n\n"
                f"Net perçu : **{_fmt_eur(net_percu)}**"
            )
            if _pfu_better:
                st.success(_msg_pfu)
            else:
                st.info(_msg_pfu)
        with _col_ir_c:
            _msg_ir = (
                f"**Option barème IR ({_tmi_pct_lbl})**\n\n"
                f"Régime : {_res_ir_opt['regime']}\n\n"
                f"IR : **{_fmt_eur(_res_ir_opt['montant_ir'])}**\n\n"
                f"PS nets : **{_fmt_eur(_ps_net_ir_opt)}**\n\n"
                f"Total prélevé : **{_fmt_eur(_total_ir_opt)}**\n\n"
                f"Net perçu : **{_fmt_eur(_net_ir_opt)}**"
            )
            if not _pfu_better:
                st.success(_msg_ir)
            else:
                st.info(_msg_ir)

        _economie = abs(net_percu - _net_ir_opt)
        _meilleure = "Flat Tax / PFL" if _pfu_better else f"Option barème IR ({_tmi_pct_lbl})"
        if _economie > 0.5:
            st.success(f"✅ L'option **{_meilleure}** vous fait économiser **{_fmt_eur(_economie)}** sur ce rachat.")
        else:
            st.info("Les deux options sont quasi-équivalentes sur ce rachat.")

        # ── Détail du calcul (expander) ──
        with st.expander("📐 Détail du calcul", expanded=False):
            _rows_det = [
                ("Quote-part gains dans le rachat", _fmt_eur(gains)),
                ("Quote-part capital (non taxé)", _fmt_eur(capital_pur)),
                ("Abattement IR appliqué", _fmt_eur(result["abattement_applique"])),
                ("Base imposable IR", _fmt_eur(result["base_ir"])),
                (f"IR dû ({result['taux_ir']*100:.1f}%)", _fmt_eur(result["montant_ir"])),
                ("PS bruts (17,2%)", _fmt_eur(result["montant_ps"])),
                ("PS déjà prélevés (fonds €)", f"− {_fmt_eur(_ps_deja_preleves)}"),
                ("PS nets dus", _fmt_eur(ps_net)),
                ("TOTAL prélèvements", _fmt_eur(total_impots)),
                ("NET PERÇU", _fmt_eur(net_percu)),
            ]
            st.dataframe(
                pd.DataFrame(_rows_det, columns=["", "Montant"]),
                hide_index=True, use_container_width=True,
            )


def _tab_optimisation_abattement() -> None:
    st.markdown("#### Optimisation de l'abattement annuel")
    st.info(
        "L'abattement de 4 600 € (célibataire) ou 9 200 € (couple) est global "
        "tous contrats d'assurance-vie. Cette simulation calcule le rachat optimal "
        "pour 'purger' les plus-values progressivement sans payer d'IR."
    )

    # ── Données lues depuis le Tableau de bord ──────────────────
    _date_ouv_o = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
    anc2 = (date.today() - _date_ouv_o).days / 365.25
    valeur2 = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
    versements2 = float(st.session_state.get("tax_versements_nets", 80_000.0))
    versements_total2 = float(st.session_state.get("tax_versements_nets_total", 80_000.0))
    sit2 = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

    with st.container(border=True):
        _oi1, _oi2, _oi3 = st.columns(3)
        _oi1.metric("Valeur du contrat", to_eur(valeur2))
        _oi2.metric("Versements nets", to_eur(versements2))
        _oi3.metric("Ancienneté", f"{anc2:.1f} ans")
        st.caption(f"Situation : **{sit2}** — _modifiez dans le 📋 Tableau de bord_")

    abat_deja = st.number_input(
        "Abattement déjà utilisé cette année sur d'autres contrats (€)",
        min_value=0.0, max_value=9_200.0,
        value=float(st.session_state.get("tax_abattement_deja_utilise", 0.0)),
        step=100.0, key="tax_abattement_deja_utilise",
        help="L'abattement de 4 600€/9 200€ est global tous contrats AV.",
    )

    if anc2 < 8:
        st.warning(
            f"Ce contrat n'a pas encore atteint 8 ans ({anc2:.1f} ans). "
            "L'optimisation par abattement ne s'applique pas."
        )
        return

    if valeur2 <= 0 or versements2 >= valeur2:
        st.info("Pas de plus-value latente — aucun calcul nécessaire.")
        return

    res = calc_optimisation_abattement(
        valeur2, versements2, anc2, sit2, versements_total2, abat_deja,
    )

    abat_total = 9_200.0 if "Couple" in sit2 else 4_600.0
    abat_restant = max(0.0, abat_total - abat_deja)

    mc1, mc2, mc3 = st.columns(3)
    mc1.metric("Abattement disponible restant", _fmt_eur(abat_restant))
    mc2.metric("Rachat recommandé", _fmt_eur(res["rachat_optimal"]))
    mc3.metric("IR dû", _fmt_eur(res["ir_du"]))

    mc4, mc5 = st.columns(2)
    mc4.metric("PS dus (inévitables)", _fmt_eur(res["ps_du"]))
    mc5.metric("Économie IR vs rachat non optimisé", _fmt_eur(res["economies_ir"]))

    # Graphique
    chart_data = pd.DataFrame({
        "Catégorie": ["Capital récupéré", "PS dus", "IR économisé"],
        "Montant (€)": [
            res["rachat_optimal"] - res["ps_du"],
            res["ps_du"],
            res["economies_ir"],
        ],
    })
    if MATPLOTLIB_AVAILABLE:
        fig, ax = plt.subplots(figsize=(5, 2.8))
        colors_bar = ["#2ecc71", "#e67e22", "#3498db"]
        bars = ax.bar(chart_data["Catégorie"], chart_data["Montant (€)"], color=colors_bar)
        ax.bar_label(bars, fmt="%.0f €", padding=3, fontsize=9)
        ax.set_ylabel("€")
        ax.set_title("Optimisation abattement — répartition")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        st.pyplot(fig)
        plt.close(fig)
    else:
        st.bar_chart(chart_data.set_index("Catégorie"))

    # Projection pluriannuelle
    nb_annees = st.slider(
        "Nombre d'années de la stratégie",
        min_value=1, max_value=30,
        value=int(st.session_state.get("tax_nb_annees", 10)),
        key="tax_nb_annees",
    )
    economie_cumulee = res["economies_ir"] * nb_annees
    st.info(
        f"📅 Cette optimisation se renouvelle chaque année civile. "
        f"Sur **{nb_annees} années**, l'économie d'IR cumulée estimée serait de "
        f"**{_fmt_eur(economie_cumulee)}** (hypothèse : situation stable)."
    )


def _tab_transmission() -> None:
    st.markdown("#### Fiscalité au décès — Clause bénéficiaire")

    regime_versements = st.radio(
        "Les versements ont été effectués :",
        [
            "Principalement avant 70 ans (Art. 990I)",
            "Principalement après 70 ans (Art. 757B)",
            "Mix avant et après 70 ans",
        ],
        key="tax_regime_versements",
    )

    c1, c2 = st.columns(2)
    with c1:
        capital_deces = st.number_input(
            "Valeur du contrat au décès (€)",
            min_value=0.0, max_value=10_000_000.0,
            value=float(st.session_state.get("tax_capital_deces", 200_000.0)),
            step=5_000.0, key="tax_capital_deces",
        )
    with c2:
        primes_apres_70 = 0.0
        if "757B" in regime_versements or "Mix" in regime_versements:
            primes_apres_70 = st.number_input(
                "Primes versées après 70 ans (€)",
                min_value=0.0, max_value=capital_deces,
                value=0.0, step=5_000.0,
                help="Seules les primes (pas les gains) au-delà de 30 500€ sont soumises aux droits de succession.",
            )

    st.markdown("---")

    # ── Lire les bénéficiaires depuis le Tableau de bord ──────────
    _benef_data_t = st.session_state.get("tax_dash_beneficiaires", [])
    if not _benef_data_t:
        st.warning("⚠️ Renseignez les bénéficiaires dans l'onglet 📋 Tableau de bord.")
        return

    nb_benef = len(_benef_data_t)
    noms = [b.get("nom", "") for b in _benef_data_t]
    types_benef = [b.get("type", "Enfant") for b in _benef_data_t]
    parts_pct = [float(b.get("part", 100.0 / nb_benef)) for b in _benef_data_t]
    liens_succ = [b.get("type", "Enfant") for b in _benef_data_t]
    # Mapper les types AV vers liens succession
    _lien_map = {
        "Conjoint/PACS": "Conjoint/PACS",
        "Enfant": "Enfant",
        "Frère/Sœur": "Frère/Sœur",
        "Neveu/Nièce": "Neveu/Nièce",
        "Tiers": "Tiers",
    }
    liens_succ = [_lien_map.get(t, "Enfant") for t in types_benef]

    total_parts = sum(parts_pct)
    if abs(total_parts - 100.0) > 0.5:
        st.error(f"⚠️ La somme des quotes-parts est {total_parts:.1f}% — elle doit être égale à 100%.")
        return

    st.markdown("---")

    if capital_deces <= 0:
        st.info("Saisissez la valeur du contrat au décès.")
        return

    # Calcul selon régime
    show_990i = "990I" in regime_versements or "Mix" in regime_versements
    show_757b = "757B" in regime_versements or "Mix" in regime_versements

    if show_990i:
        st.markdown("##### Art. 990I — Versements avant 70 ans")
        types_990i = [
            "Conjoint/PACS" if ("Conjoint" in t or "PACS" in t) else t
            for t in types_benef
        ]
        results_990i = calc_transmission_990I(capital_deces, nb_benef, parts_pct, types_990i)
        cap_990i = capital_deces
        if "Mix" in regime_versements:
            cap_990i_pct = 1.0 - (primes_apres_70 / capital_deces) if capital_deces > 0 else 1.0
            results_990i = calc_transmission_990I(
                capital_deces * cap_990i_pct, nb_benef, parts_pct, types_990i
            )
        _render_table_transmission(results_990i, noms, "990I")
        # Note démembrement
        for r in results_990i:
            if r.get("note"):
                st.info(r["note"])

    if show_757b:
        st.markdown("##### Art. 757B — Versements après 70 ans")
        st.info(
            "Les **gains** générés par les primes versées après 70 ans sont exonérés. "
            "Seules les **primes** (hors gains) au-delà de 30 500 € (abattement global) "
            "sont soumises aux droits de succession selon le lien de parenté."
        )
        st.caption(
            "Les droits de succession réels dépendent du patrimoine global, des donations "
            "antérieures et des abattements de droit commun (ex. 100 000€ par enfant). "
            "Cette simulation utilise des taux simplifiés."
        )
        results_757b = calc_transmission_757B(
            primes_apres_70, nb_benef, parts_pct, types_benef, capital_deces, liens_succ,
        )
        _render_table_transmission(results_757b, noms, "757B")


def _render_table_transmission(results: List[Dict[str, Any]], noms: List[str], regime: str) -> None:
    """Affiche le tableau de résultats par bénéficiaire."""
    rows = []
    for r in results:
        nom = noms[r["index"] - 1] if r["index"] - 1 < len(noms) and noms[r["index"] - 1] else f"Bénéficiaire {r['index']}"
        if regime == "990I":
            rows.append({
                "Bénéficiaire": nom,
                "Lien": r["type"],
                "Part (€)": _fmt_eur(r["part_eur"]),
                "Abattement": _fmt_eur(r.get("abattement", 0.0)),
                "Taxable": _fmt_eur(r["taxable"]),
                "Taxe due": _fmt_eur(r["taxe"]),
                "Net reçu": _fmt_eur(r["net_recu"]),
            })
        else:
            rows.append({
                "Bénéficiaire": nom,
                "Lien": r["type"],
                "Part (€)": _fmt_eur(r["part_eur"]),
                "Primes part": _fmt_eur(r.get("primes_part", 0.0)),
                "Abattement": _fmt_eur(r.get("abattement_part", 0.0)),
                "Taxable": _fmt_eur(r["taxable"]),
                "Taxe due": _fmt_eur(r["taxe"]),
                "Net reçu": _fmt_eur(r["net_recu"]),
            })
    total_taxe = sum(r["taxe"] for r in results)
    total_net = sum(r["net_recu"] for r in results)
    total_part = sum(r["part_eur"] for r in results)

    if rows:
        st.dataframe(pd.DataFrame(rows), hide_index=True, use_container_width=True)

    tc1, tc2, tc3 = st.columns(3)
    tc1.metric("Capital total transmis", _fmt_eur(total_part))
    tc2.metric("Total taxes", _fmt_eur(total_taxe))
    taux_eff = (total_taxe / total_part * 100) if total_part > 0 else 0.0
    tc3.metric("Taux effectif global", f"{taux_eff:.1f}%")


def _tab_exoneration() -> None:
    st.info(
        "Dans certaines situations, les gains issus d'un rachat sont totalement exonérés "
        "d'IR **ET** de prélèvements sociaux. Ces exonérations s'appliquent sur justificatif."
    )

    st.markdown("#### Cochez votre situation")
    cas_licenciement = st.checkbox(
        "Licenciement (inscription à Pôle Emploi / France Travail)"
    )
    cas_retraite_anticipee = st.checkbox(
        "Mise à la retraite anticipée par l'employeur"
    )
    cas_invalidite = st.checkbox(
        "Invalidité 2e ou 3e catégorie (assuré, conjoint ou enfant à charge)"
    )
    cas_liquidation = st.checkbox(
        "Liquidation judiciaire (non-salariés : cessation d'activité suite à jugement de liquidation)"
    )
    cas_cessation_ns = st.checkbox(
        "Cessation d'activité non salariée suite à jugement de liquidation"
    )

    cas_coches = [
        (cas_licenciement, "Licenciement",
         "Justificatif : lettre de licenciement + attestation Pôle Emploi / France Travail"),
        (cas_retraite_anticipee, "Retraite anticipée",
         "Justificatif : notification de mise à la retraite par l'employeur"),
        (cas_invalidite, "Invalidité",
         "Justificatif : notification CPAM de classement en invalidité 2e ou 3e catégorie"),
        (cas_liquidation, "Liquidation judiciaire",
         "Justificatif : jugement du tribunal de commerce"),
        (cas_cessation_ns, "Cessation NS",
         "Justificatif : jugement du tribunal de commerce"),
    ]

    au_moins_un = any(c[0] for c in cas_coches)

    if au_moins_un:
        st.success(
            "✅ **Exonération totale applicable** : ni IR ni prélèvements sociaux "
            "ne sont dus sur les gains du rachat."
        )
        # Tableau comparatif
        rachat_result: Optional[Dict[str, Any]] = st.session_state.get("tax_rachat_result")
        if rachat_result:
            ir_normal = rachat_result.get("montant_ir", 0.0)
            ps_normal = rachat_result.get("montant_ps", 0.0)
            total_normal = rachat_result.get("total_impots", 0.0)
            df_comp = pd.DataFrame([
                {"": "IR dû",
                 "Sans exonération": _fmt_eur(ir_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(ir_normal)},
                {"": "PS dûs",
                 "Sans exonération": _fmt_eur(ps_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(ps_normal)},
                {"": "TOTAL",
                 "Sans exonération": _fmt_eur(total_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(total_normal)},
            ])
            st.dataframe(df_comp, hide_index=True, use_container_width=True)
        else:
            st.caption("Renseignez d'abord l'onglet **Rachat** pour voir l'économie chiffrée.")

        st.markdown("#### Justificatifs requis")
        for coched, label, justif in cas_coches:
            if coched:
                st.caption(f"**{label}** — {justif}")

    st.markdown("---")
    st.warning(
        "⚠️ **Procédure :** l'exonération n'est pas appliquée automatiquement par l'assureur. "
        "Le bénéficiaire doit fournir les justificatifs à l'assureur au moment du rachat "
        "**ET** mentionner l'exonération dans sa déclaration de revenus (case spécifique)."
    )


def _tab_rente_viagere() -> None:
    """Simulation de transformation du capital en rente viagère."""
    st.markdown("#### Simulation — Sortie en rente viagère")
    st.info(
        "La rente viagère issue d'un contrat d'assurance-vie bénéficie d'une "
        "fiscalité allégée. Seule une fraction de la rente est imposable à l'IR, "
        "selon l'âge au moment de la mise en rente (Art. 158-6 CGI)."
    )
    # Fractions imposables selon l'âge (Art. 158-6 CGI)
    _FRACTIONS_IMPOSABLES = {
        "Moins de 50 ans":     0.70,
        "50 à 59 ans":         0.50,
        "60 à 69 ans":         0.40,
        "70 ans et plus":      0.30,
    }
    _rc1, _rc2 = st.columns(2)
    with _rc1:
        _capital_rente = st.number_input(
            "Capital à convertir en rente (€)",
            min_value=0.0, max_value=5_000_000.0,
            value=200_000.0, step=5_000.0, key="rente_capital",
        )
        _taux_rente_annuel = st.number_input(
            "Taux de conversion annuel estimé (%)",
            min_value=1.0, max_value=10.0,
            value=4.5, step=0.1, key="rente_taux",
            help="Dépend de l'âge, du sexe et de la table de mortalité de l'assureur. "
                 "Typiquement 4-6% pour un rentier de 65 ans.",
        )
    with _rc2:
        _tranche_age = st.selectbox(
            "Tranche d'âge au moment de la mise en rente",
            list(_FRACTIONS_IMPOSABLES.keys()),
            index=2,  # 60-69 ans par défaut
            key="rente_age",
        )
        _tmi = st.selectbox(
            "Taux marginal d'imposition (TMI)",
            [0.11, 0.30, 0.41, 0.45],
            index=1,
            format_func=lambda x: f"{x*100:.0f}%",
            key="rente_tmi",
        )
    _rente_brute_annuelle = _capital_rente * (_taux_rente_annuel / 100.0)
    _rente_brute_mensuelle = _rente_brute_annuelle / 12.0
    _fraction_imposable = _FRACTIONS_IMPOSABLES[_tranche_age]
    _base_ir = _rente_brute_annuelle * _fraction_imposable
    _ir_annuel = _base_ir * _tmi
    _ps_annuel = _base_ir * 0.172  # PS à 17,2%
    _rente_nette_annuelle = _rente_brute_annuelle - _ir_annuel - _ps_annuel
    _rente_nette_mensuelle = _rente_nette_annuelle / 12.0
    _taux_prelevement_effectif = (_ir_annuel + _ps_annuel) / _rente_brute_annuelle if _rente_brute_annuelle > 0 else 0.0
    st.markdown("---")
    _rm1, _rm2, _rm3, _rm4 = st.columns(4)
    _rm1.metric("Rente brute mensuelle", f"{_rente_brute_mensuelle:,.0f} €")
    _rm2.metric("Rente nette mensuelle", f"{_rente_nette_mensuelle:,.0f} €")
    _rm3.metric("Fraction imposable", f"{_fraction_imposable*100:.0f}%")
    _rm4.metric("Taux réel de prélèvement", f"{_taux_prelevement_effectif*100:.1f}%")
    st.markdown("---")
    _detail_rows = [
        ("Rente brute annuelle", f"{_rente_brute_annuelle:,.0f} €"),
        (f"Fraction imposable ({_fraction_imposable*100:.0f}%)", f"{_base_ir:,.0f} €"),
        (f"IR ({_tmi*100:.0f}% × fraction imposable)", f"−{_ir_annuel:,.0f} €"),
        ("Prélèvements sociaux (17,2% × fraction imposable)", f"−{_ps_annuel:,.0f} €"),
        ("Rente nette annuelle", f"{_rente_nette_annuelle:,.0f} €"),
        ("Rente nette mensuelle", f"{_rente_nette_mensuelle:,.0f} €"),
    ]
    st.dataframe(
        pd.DataFrame(_detail_rows, columns=["", "Montant"]),
        hide_index=True, use_container_width=True,
    )
    st.caption(
        "⚠️ Simulation indicative. Le taux de conversion réel dépend de la table de mortalité "
        "TGH05/TGF05, du type de rente (simple, réversible, avec annuités garanties) et de "
        "l'assureur. Consultez un conseiller pour un chiffrage personnalisé."
    )


def render_tax_module() -> None:
    """Module Fiscalité & Avantages AV — entièrement autonome."""
    # Constantes fiscales (locales, non globales)
    PS_RATE            = 0.172
    PFU_RATE           = 0.128
    PFL_LOW_RATE       = 0.075
    PFL_HIGH_RATE      = 0.128
    ABATTEMENT_SEUL    = 4_600.0
    ABATTEMENT_COUPLE  = 9_200.0
    SEUIL_150K         = 150_000.0
    ABATTEMENT_990I    = 152_500.0
    TAUX_990I_TRANCHE1 = 0.20
    TAUX_990I_TRANCHE2 = 0.3125
    SEUIL_990I_T2      = 700_000.0
    ABATTEMENT_757B    = 30_500.0
    AGE_SEUIL          = 70

    # Initialisation session_state
    st.session_state.setdefault("tax_date_ouverture", date(2016, 1, 2))
    st.session_state.setdefault("tax_valeur_contrat", 100_000.0)
    st.session_state.setdefault("tax_versements_nets", 80_000.0)
    st.session_state.setdefault("tax_versements_nets_total", 80_000.0)
    st.session_state.setdefault("tax_situation_familiale", "Célibataire / veuf / divorcé")
    st.session_state.setdefault("tax_montant_brut", 10_000.0)
    st.session_state.setdefault("tax_montant_net", 9_000.0)
    st.session_state.setdefault("tax_last_edited", "brut")
    st.session_state.setdefault("tax_rachat_result", None)
    st.session_state.setdefault("tax_nb_beneficiaires", 2)
    st.session_state.setdefault("tax_regime_versements", "Principalement avant 70 ans (Art. 990I)")
    st.session_state.setdefault("tax_capital_deces", 200_000.0)
    st.session_state.setdefault("tax_abattement_deja_utilise", 0.0)
    st.session_state.setdefault("tax_nb_annees", 10)
    # ── PERSIST keys : survivent à la navigation entre modules ────────
    st.session_state.setdefault("tax_date_ouverture_PERSIST", date(2016, 1, 2))
    st.session_state.setdefault("tax_situation_familiale_PERSIST", "Célibataire / veuf / divorcé")
    st.session_state.setdefault("tax_age_PERSIST", 55)
    st.session_state.setdefault("tax_tmi_PERSIST", 0.30)
    st.session_state.setdefault("tax_nb_enfants_PERSIST", 2)
    st.session_state.setdefault("tax_nb_benef_PERSIST", 2)
    st.session_state.setdefault("tax_beneficiaires_PERSIST", [])
    st.session_state.setdefault("tax_regime_matrimonial_PERSIST", "Séparation de biens")
    # Réinjecter les clés lues par les fonctions de calcul depuis PERSIST
    st.session_state["tax_date_ouverture"] = st.session_state["tax_date_ouverture_PERSIST"]
    st.session_state["tax_regime_matrimonial"] = st.session_state["tax_regime_matrimonial_PERSIST"]
    st.session_state["tax_situation_familiale"] = st.session_state["tax_situation_familiale_PERSIST"]
    st.session_state["tax_age_souscripteur"] = st.session_state["tax_age_PERSIST"]
    st.session_state["tax_nb_enfants"] = st.session_state["tax_nb_enfants_PERSIST"]
    st.session_state["tax_nb_beneficiaires"] = st.session_state["tax_nb_benef_PERSIST"]

    st.title("Fiscalité & Avantages de l'assurance-vie")
    st.warning(
        "⚠️ Simulation indicative à titre pédagogique. Ne constitue pas un conseil "
        "fiscal ou juridique. Consultez un professionnel pour toute décision."
    )

    # ── Sélecteur de portefeuille ──────────────────────────────────
    _has_A = bool(st.session_state.get("A_lines", []))
    _has_B = bool(st.session_state.get("B_lines", []))
    _options_pf = []
    if _has_A:
        _options_pf.append("🧍 Portefeuille Client")
    if _has_B:
        _options_pf.append("🏢 Portefeuille Cabinet")
    if not _options_pf:
        _options_pf = ["🧍 Portefeuille Client"]

    _selected_pf = st.radio(
        "Portefeuille analysé",
        _options_pf,
        horizontal=True,
        key="tax_portfolio_selector",
    )
    _is_client = "Client" in _selected_pf
    st.session_state["tax_is_client"] = _is_client

    # ── Synchronisation depuis le comparateur ──────────────────────
    if _is_client:
        _sync_val = float(st.session_state.get("_LAST_VAL_A", 0.0) or 0.0)
        _sync_net = float(st.session_state.get("_LAST_NET_A", 0.0) or 0.0)
        _sync_xirr = st.session_state.get("_LAST_XIRR_A")
        _sync_contract = st.session_state.get("CONTRACT_LABEL_A", "")
        _sync_euro_rate = float(st.session_state.get("EURO_RATE_A", 2.0))
    else:
        _sync_val = float(st.session_state.get("_LAST_VAL_B", 0.0) or 0.0)
        _sync_net = float(st.session_state.get("_LAST_NET_B", 0.0) or 0.0)
        _sync_xirr = st.session_state.get("_LAST_XIRR_B")
        _sync_contract = st.session_state.get("CONTRACT_LABEL_B", "")
        _sync_euro_rate = float(st.session_state.get("EURO_RATE_B", 2.5))

    if _sync_val > 0:
        st.session_state["tax_valeur_contrat"] = _sync_val
        st.session_state["tax_capital_deces"] = _sync_val
    if _sync_net > 0:
        st.session_state["tax_versements_nets"] = _sync_net
        st.session_state["tax_versements_nets_total"] = _sync_net

    st.session_state["_tax_sync_contract"] = _sync_contract
    st.session_state["_tax_sync_xirr"] = _sync_xirr
    st.session_state["_tax_sync_euro_rate"] = _sync_euro_rate

    tab_dashboard, tab_retrait, tab_transmission_v2, tab_avantages = st.tabs([
        "📋 Tableau de bord",
        "💶 Je veux retirer de l'argent",
        "🏠 Je veux transmettre",
        "⚖️ Pourquoi l'assurance-vie ?",
    ])

    with tab_dashboard:
        _tab_dashboard_fiscal()
    with tab_retrait:
        _tab_retrait_wrapper()
    with tab_transmission_v2:
        _tab_transmission_v2()
    with tab_avantages:
        _tab_avantages_av()

    st.markdown("---")
    with st.expander("📖 Cas d'exonération & Sortie en rente", expanded=False):
        _ex1, _ex2 = st.tabs(["Cas d'exonération", "Sortie en rente"])
        with _ex1:
            _tab_exoneration()
        with _ex2:
            _tab_rente_viagere()


def _tab_dashboard_fiscal():
    """Hub central : infos contrat + situation perso + bénéficiaires + frise."""
    st.markdown("#### 📋 Votre contrat en un coup d'œil")

    _c1, _c2 = st.columns(2)
    with _c1:
        _contract = st.session_state.get("_tax_sync_contract", "") or st.session_state.get("CONTRACT_LABEL_A", "")
        if _contract:
            st.markdown(f"**Contrat** : {_contract}")
        _date_ouv = st.date_input(
            "Date d'ouverture du contrat",
            value=st.session_state.get("tax_date_ouverture_PERSIST", date(2016, 1, 2)),
            key="tax_dash_date_ouverture",
            help="Date réelle d'ouverture du contrat.",
        )
        st.session_state["tax_date_ouverture"] = _date_ouv
        st.session_state["tax_date_ouverture_PERSIST"] = _date_ouv

        _sit_fam_options = ["Célibataire / veuf / divorcé", "Couple (imposition commune)"]
        _sit_fam_persist = st.session_state.get("tax_situation_familiale_PERSIST", "Célibataire / veuf / divorcé")
        _sit_fam_idx = _sit_fam_options.index(_sit_fam_persist) if _sit_fam_persist in _sit_fam_options else 0
        _sit_fam = st.selectbox(
            "Situation familiale",
            _sit_fam_options,
            index=_sit_fam_idx,
            key="tax_dash_situation_familiale",
        )
        st.session_state["tax_situation_familiale"] = _sit_fam
        st.session_state["tax_situation_familiale_PERSIST"] = _sit_fam

        if "Couple" in _sit_fam:
            _regime_options = [
                "Séparation de biens",
                "Communauté réduite aux acquêts",
                "Communauté universelle",
            ]
            _regime_persist = st.session_state.get("tax_regime_matrimonial_PERSIST", "Séparation de biens")
            _regime_idx = _regime_options.index(_regime_persist) if _regime_persist in _regime_options else 0
            _regime = st.selectbox(
                "Régime matrimonial",
                _regime_options,
                index=_regime_idx,
                key="tax_regime_matrimonial_widget",
                help="Impacte les droits de succession entre époux.",
            )
            st.session_state["tax_regime_matrimonial"] = _regime
            st.session_state["tax_regime_matrimonial_PERSIST"] = _regime
        else:
            st.session_state["tax_regime_matrimonial"] = ""
            st.session_state.setdefault("tax_regime_matrimonial_PERSIST", "Séparation de biens")

        _age = st.number_input(
            "Âge du souscripteur",
            min_value=18, max_value=100,
            value=int(st.session_state.get("tax_age_PERSIST", 55)),
            step=1,
            key="tax_age_widget",
        )
        st.session_state["tax_age_souscripteur"] = _age
        st.session_state["tax_age_PERSIST"] = _age

        _tmi_options = [0.0, 0.11, 0.30, 0.41, 0.45]
        _tmi_persist = st.session_state.get("tax_tmi_PERSIST", 0.30)
        _tmi_idx = _tmi_options.index(_tmi_persist) if _tmi_persist in _tmi_options else 2
        _tmi = st.selectbox(
            "Taux marginal d'imposition (TMI)",
            _tmi_options,
            index=_tmi_idx,
            format_func=lambda x: f"{int(x*100)}%",
            key="tax_tmi_widget",
            help="Taux marginal de l'IR. Utilisé pour estimer les économies fiscales.",
        )
        st.session_state["tax_tmi_PERSIST"] = _tmi

    with _c2:
        _val = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
        _net = float(st.session_state.get("tax_versements_nets", 80_000.0))
        _pv = _val - _net
        _pv_pct = (_pv / _net * 100) if _net > 0 else 0
        _xirr = st.session_state.get("_tax_sync_xirr")

        st.metric("Valeur actuelle du contrat", to_eur(_val))
        st.metric("Versements nets", to_eur(_net))
        st.metric("Plus-values latentes", to_eur(_pv), delta=f"{_pv_pct:+.1f}%")
        if _xirr is not None:
            st.metric("Rendement annualisé (XIRR)", f"{float(_xirr):+.2f}%")

    # ── Frise chronologique maturité fiscale ──
    st.markdown("---")
    _anciennete = (_date_ouv.toordinal() - date.today().toordinal()) * -1 / 365.25
    _anciennete = max(0.0, _anciennete)
    _ans = int(_anciennete)
    _mois = int((_anciennete - _ans) * 12)

    # Frise simple via Altair
    _ouv_ts = pd.Timestamp(_date_ouv)
    _today_ts = pd.Timestamp(date.today())
    _4ans_ts = _ouv_ts + pd.DateOffset(years=4)
    _8ans_ts = _ouv_ts + pd.DateOffset(years=8)
    _12ans_ts = _ouv_ts + pd.DateOffset(years=12)

    _frise_pts = pd.DataFrame([
        {"label": "Ouverture", "date": _ouv_ts, "type": "jalon"},
        {"label": "4 ans", "date": _4ans_ts, "type": "jalon"},
        {"label": "8 ans (maturité)", "date": _8ans_ts, "type": "jalon_majeur"},
    ])
    _aujourd_hui_df = pd.DataFrame([
        {"label": "Aujourd'hui", "date": _today_ts, "type": "position"},
    ])

    _pos_color = "#2E7D32" if _anciennete >= 8 else ("#FF8F00" if _anciennete >= 4 else "#C62828")

    _rule_df = pd.DataFrame({"x": [_ouv_ts], "x2": [_12ans_ts], "y": [0.5]})
    _rule = (
        alt.Chart(_rule_df)
        .mark_rule(strokeWidth=3, color="#CCCCCC")
        .encode(x=alt.X("x:T"), x2="x2:T", y=alt.value(40))
    )

    _pts = (
        alt.Chart(_frise_pts)
        .mark_point(size=120, filled=True, color="#1B2A4A")
        .encode(
            x=alt.X("date:T", axis=alt.Axis(format="%Y", title="")),
            y=alt.value(40),
            tooltip=["label:N", alt.Tooltip("date:T", format="%d/%m/%Y")],
        )
    )
    _lbl = (
        alt.Chart(_frise_pts)
        .mark_text(dy=-18, fontSize=11, fontWeight="bold", color="#1B2A4A")
        .encode(x="date:T", y=alt.value(40), text="label:N")
    )
    _pos_pt = (
        alt.Chart(_aujourd_hui_df)
        .mark_point(size=200, filled=True, shape="triangle-down")
        .encode(
            x="date:T",
            y=alt.value(40),
            color=alt.value(_pos_color),
            tooltip=["label:N"],
        )
    )
    _pos_lbl = (
        alt.Chart(_aujourd_hui_df)
        .mark_text(dy=20, fontSize=11, fontWeight="bold")
        .encode(x="date:T", y=alt.value(40), text="label:N", color=alt.value(_pos_color))
    )

    _frise_chart = (_rule + _pts + _lbl + _pos_pt + _pos_lbl).properties(height=80, width="container").configure_view(strokeOpacity=0)
    st.altair_chart(_frise_chart, use_container_width=True)

    # ── Statut de maturité ──
    if _anciennete >= 8:
        _emoji_mat, _label_mat = "🟢", "Maturité fiscale atteinte"
        _abat = "9 200 €" if ("Couple" in _sit_fam) else "4 600 €"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Régime fiscal optimal. "
            f"Abattement annuel de **{_abat}** sur les gains en cas de rachat. "
            f"Taux réduit de 7,5% (PFL) pour les versements ≤ 150 000 €."
        )
    elif _anciennete >= 4:
        _restant = 8 - _anciennete
        _r_ans, _r_mois = int(_restant), int((_restant - int(_restant)) * 12) + 1
        _emoji_mat, _label_mat = "🟠", "Maturité en cours"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Maturité fiscale dans "
            f"**{_r_ans} an(s) {_r_mois} mois**. Rachats soumis au PFU 30% en attendant."
        )
    else:
        _restant = 8 - _anciennete
        _r_ans, _r_mois = int(_restant), int((_restant - int(_restant)) * 12) + 1
        _emoji_mat, _label_mat = "🔴", "Contrat récent"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Flat tax 30% sur les rachats. "
            f"Maturité fiscale dans **{_r_ans} ans {_r_mois} mois**."
        )
    st.markdown(f"### {_emoji_mat} {_label_mat}")
    st.markdown(_detail)

    # ── Situation familiale & Clause bénéficiaire ──
    st.markdown("---")
    st.markdown("#### 👨‍👩‍👧‍👦 Situation familiale & Clause bénéficiaire")

    _fc1, _fc2 = st.columns(2)
    with _fc1:
        _nb_enfants = st.number_input(
            "Nombre d'enfants",
            min_value=0, max_value=10,
            value=int(st.session_state.get("tax_nb_enfants_PERSIST", 2)),
            step=1, key="tax_nb_enfants_widget",
            help="Pour le calcul des droits de succession classique (abattement 100 000€ par enfant).",
        )
        st.session_state["tax_nb_enfants"] = _nb_enfants
        st.session_state["tax_nb_enfants_PERSIST"] = _nb_enfants
    with _fc2:
        _nb_benef = st.number_input(
            "Nombre de bénéficiaires (clause AV)",
            min_value=1, max_value=5,
            value=int(st.session_state.get("tax_nb_benef_PERSIST", 2)),
            step=1, key="tax_nb_benef_widget",
        )
        st.session_state["tax_nb_beneficiaires"] = _nb_benef
        st.session_state["tax_nb_benef_PERSIST"] = _nb_benef

    _TYPES_BENEF = ["Conjoint/PACS", "Enfant", "Frère/Sœur", "Neveu/Nièce", "Tiers"]
    _persisted_benef = st.session_state.get("tax_beneficiaires_PERSIST", [])
    _benef_data = []
    for i in range(int(_nb_benef)):
        _prev = _persisted_benef[i] if i < len(_persisted_benef) else {}
        with st.expander(f"Bénéficiaire {i + 1}", expanded=(i == 0)):
            _bc1, _bc2, _bc3 = st.columns([2, 2, 1])
            with _bc1:
                _nom = st.text_input("Nom (facultatif)", key=f"tax_dash_benef_nom_{i}", value=_prev.get("nom", ""))
            with _bc2:
                _type_prev = _prev.get("type", "Enfant")
                _type_idx = _TYPES_BENEF.index(_type_prev) if _type_prev in _TYPES_BENEF else 1
                _type = st.selectbox("Lien", _TYPES_BENEF, index=_type_idx, key=f"tax_dash_benef_type_{i}")
            with _bc3:
                _part = st.number_input(
                    "Quote-part %", min_value=0.0, max_value=100.0,
                    value=_prev.get("part", round(100.0 / int(_nb_benef), 0)), step=5.0,
                    key=f"tax_dash_benef_part_{i}",
                )
            _benef_data.append({"nom": _nom, "type": _type, "part": float(_part)})

    st.session_state["tax_dash_beneficiaires"] = _benef_data
    st.session_state["tax_beneficiaires_PERSIST"] = _benef_data

    _total_parts = sum(b["part"] for b in _benef_data)
    if abs(_total_parts - 100.0) > 0.5:
        st.error(f"⚠️ La somme des quotes-parts est {_total_parts:.1f}% — elle doit être égale à 100%.")

    # ── Transmission en bref ──
    st.markdown("---")
    st.markdown("#### 🏠 Transmission — En bref")
    if _age < 70:
        _abat_trans = 152_500 * int(_nb_benef)
        st.success(
            f"✅ Vous avez moins de 70 ans — **Art. 990I** : **{to_eur(152_500)}** "
            f"exonérés par bénéficiaire. Avec **{int(_nb_benef)} bénéficiaire(s)**, "
            f"jusqu'à **{to_eur(_abat_trans)}** transmis hors succession."
        )
        _marge = _abat_trans - _val
        if _marge > 0:
            st.info(f"💡 Marge de versement avant plafond d'exonération : **{to_eur(_marge)}**")
    else:
        st.warning(
            "⚠️ Après 70 ans — **Art. 757B** : abattement global de **30 500 €**. "
            "Les **intérêts et plus-values** restent exonérés de droits de succession."
        )


def _tab_retrait_wrapper():
    """Regroupe : combien je paie / combien je peux retirer sans impôt / dois-je attendre."""
    st.markdown("#### 💶 Retirer de l'argent — Simulateur de rachat")

    _q1, _q2, _q3 = st.tabs([
        "Combien je paie si je retire X€ ?",
        "Combien puis-je retirer sans impôt ?",
        "Ai-je intérêt à attendre ?",
    ])

    with _q1:
        _tab_rachat()

    with _q2:
        st.markdown("#### Combien puis-je retirer sans payer d'IR ?")
        _date_ouv_q2 = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
        _anc_q2 = (date.today() - _date_ouv_q2).days / 365.25
        _val_q2 = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
        _net_q2 = float(st.session_state.get("tax_versements_nets", 80_000.0))
        _net_total_q2 = float(st.session_state.get("tax_versements_nets_total", _net_q2))
        _sit_q2 = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

        if _anc_q2 < 8:
            st.warning(
                f"⚠️ Votre contrat a {_anc_q2:.1f} ans — pas encore éligible à l'abattement. "
                f"Tout rachat est soumis au PFU 30%. L'abattement s'applique après 8 ans."
            )
        elif _val_q2 > 0 and _net_q2 > 0:
            _opt_q2 = calc_optimisation_abattement(
                _val_q2, _net_q2, _anc_q2, _sit_q2, _net_total_q2, 0.0
            )
            _rachat_optimal = _opt_q2["rachat_optimal"]
            _ps_du = _opt_q2["ps_du"]
            _net_main = _rachat_optimal - _ps_du
            # Cas où les PV totales du contrat sont inférieures à l'abattement
            _gains_total = max(0.0, _val_q2 - _net_q2)
            _abat_label = "9 200 €" if "Couple" in _sit_q2 else "4 600 €"
            _abat_val = 9_200.0 if "Couple" in _sit_q2 else 4_600.0
            if _gains_total <= _abat_val:
                st.success(
                    f"✅ **Tout rachat est exonéré d'IR** : les plus-values totales du contrat "
                    f"(**{to_eur(_gains_total)}**) sont inférieures à l'abattement annuel "
                    f"(**{_abat_label}**). Vous pouvez racheter l'intégralité du contrat sans payer "
                    f"d'impôt sur le revenu. Seuls les prélèvements sociaux (17,2%) resteront dus "
                    f"sur les gains."
                )
                _mc1, _mc2, _mc3 = st.columns([1, 2, 1])
                with _mc2:
                    st.metric("Rachat sans IR possible", to_eur(_val_q2),
                              delta="Intégralité du contrat")
                    st.metric("Net en main", to_eur(_val_q2 - _gains_total * 0.172),
                              delta=f"PS uniquement : {to_eur(_gains_total * 0.172)}")
            else:
                _mc1, _mc2, _mc3 = st.columns([1, 2, 1])
                with _mc2:
                    st.metric("✅ Rachat maximum sans IR", to_eur(_rachat_optimal),
                              delta="0 € d'impôt sur le revenu")
                    st.metric("Net en main", to_eur(_net_main),
                              delta=f"PS uniquement : {to_eur(_ps_du)}")
                st.markdown("---")
                st.caption(
                    f"💡 Ce montant correspond au rachat qui génère exactement **{_abat_label}** "
                    f"de gains imposables — intégralement couverts par l'abattement annuel. "
                    f"Seuls les prélèvements sociaux (17,2%) restent dus sur les gains."
                )
            with st.expander("📐 Détail du calcul", expanded=False):
                _tab_optimisation_abattement()
        else:
            st.info("Renseignez les données du contrat dans le Tableau de bord.")

    with _q3:
        _date_ouv_att = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
        _anciennete_att = (date.today() - _date_ouv_att).days / 365.25

        if _anciennete_att >= 8:
            st.success(
                "✅ Votre contrat a déjà atteint sa maturité fiscale (>8 ans). "
                "Aucun intérêt fiscal à attendre davantage pour effectuer un rachat."
            )
        else:
            _restant_att = 8 - _anciennete_att
            _r_ans_att = int(_restant_att)
            _r_mois_att = int((_restant_att - _r_ans_att) * 12) + 1
            st.markdown(f"**Maturité fiscale dans : {_r_ans_att} an(s) et {_r_mois_att} mois**")

            _montant_test = st.number_input(
                "Montant du rachat envisagé (€)",
                min_value=0.0, max_value=5_000_000.0,
                value=30_000.0, step=1_000.0, key="attente_montant",
            )

            _val_att = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
            _net_att = float(st.session_state.get("tax_versements_nets", 80_000.0))
            _net_total_att = float(st.session_state.get("tax_versements_nets_total", _net_att))
            _sit_att = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

            if _montant_test > 0 and _val_att > 0:
                _gains_now = calc_quote_part_gains(_val_att, _net_att, _montant_test)
                _gains_after = calc_quote_part_gains(_val_att, _net_att, _montant_test)
                _res_now = calc_imposition_rachat(
                    _gains_now, _anciennete_att, _sit_att, _net_total_att, _montant_test
                )
                _res_after = calc_imposition_rachat(
                    _gains_after, 8.1, _sit_att, _net_total_att, _montant_test
                )
                _eco = _res_now["total_impots"] - _res_after["total_impots"]

                _ca1, _ca2 = st.columns(2)
                with _ca1:
                    st.markdown("**Rachat maintenant**")
                    st.metric("Impôt + PS", to_eur(_res_now["total_impots"]))
                    st.metric("Net perçu", to_eur(_res_now["net_percu"]))
                with _ca2:
                    st.markdown(f"**Rachat après maturité** ({_r_ans_att}a {_r_mois_att}m)")
                    st.metric("Impôt + PS", to_eur(_res_after["total_impots"]))
                    st.metric("Net perçu", to_eur(_res_after["net_percu"]))

                if _eco > 0:
                    st.success(
                        f"💰 **En attendant {_r_ans_att} an(s) et {_r_mois_att} mois**, vous économisez "
                        f"**{to_eur(_eco)}** d'impôt sur ce rachat de {to_eur(_montant_test)}."
                    )
                    # Bar chart comparatif
                    _bar_att_df = pd.DataFrame([
                        {"Scénario": "Rachat\nmaintenant", "Impôt (€)": _res_now["total_impots"], "Type": "Maintenant"},
                        {"Scénario": f"Rachat après\nmaturité ({_r_ans_att}a)", "Impôt (€)": _res_after["total_impots"], "Type": "Après"},
                    ])
                    _bar_att = (
                        alt.Chart(_bar_att_df)
                        .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                        .encode(
                            x=alt.X("Scénario:N", sort=None, axis=alt.Axis(labelAngle=0), title=""),
                            y=alt.Y("Impôt (€):Q", axis=alt.Axis(format=",.0f"), title="Impôt total (€)"),
                            color=alt.Color("Type:N", scale=alt.Scale(
                                domain=["Maintenant", "Après"], range=["#E53935", "#2E7D32"]
                            ), legend=None),
                            tooltip=[alt.Tooltip("Scénario:N"), alt.Tooltip("Impôt (€):Q", format=",.0f")],
                        )
                        .properties(height=250, title=f"Économie : {to_eur(_eco)}")
                    )
                    st.altair_chart(_bar_att, use_container_width=True)
                else:
                    st.info("L'écart fiscal est négligeable dans votre situation.")


def _tab_transmission_v2():
    """Transmission enrichie : calcul existant + comparaison AV vs succession + projection."""

    # ── Avertissement régime matrimonial ──
    _regime_mat = st.session_state.get("tax_regime_matrimonial", "")
    if "Communauté universelle" in _regime_mat:
        st.warning(
            "⚠️ **Communauté universelle** : en présence d'une clause d'attribution intégrale, "
            "le conjoint survivant hérite de l'intégralité des biens **sans droits de succession**. "
            "L'assurance-vie reste utile pour les autres bénéficiaires (enfants, tiers). "
            "Les calculs ci-dessous concernent ces bénéficiaires hors conjoint."
        )
    elif "Communauté réduite aux acquêts" in _regime_mat:
        st.info(
            "ℹ️ **Communauté réduite aux acquêts** : les biens propres et la part de communauté "
            "du défunt entrent dans la succession. L'AV reste hors succession pour les bénéficiaires désignés."
        )

    _q1, _q2 = st.tabs([
        "Combien mes proches recevront-ils ?",
        "Ai-je intérêt à verser avant 70 ans ?",
    ])

    with _q1:
        st.markdown("#### Que recevront vos proches ?")
        _val_t1 = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
        _age_t1 = st.session_state.get("tax_age_souscripteur", 55)
        _xirr_t1 = st.session_state.get("_tax_sync_xirr")
        _default_rdt_t1 = float(_xirr_t1) if (_xirr_t1 and _xirr_t1 > 0) else 5.0

        st.markdown("---")
        _horizon_trans = st.slider(
            "Horizon de projection (années)",
            min_value=0, max_value=30, value=0, step=1,
            key="tax_trans_horizon",
            help="0 = aujourd'hui. Déplacez pour projeter dans le futur avec le rendement du portefeuille.",
        )

        if _horizon_trans > 0:
            _rdt_trans = st.number_input(
                "Rendement annuel estimé (%)",
                min_value=-5.0, max_value=20.0,
                value=round(_default_rdt_t1, 1), step=0.5,
                key="tax_trans_rendement",
                help="Pré-rempli avec le rendement XIRR du portefeuille si disponible.",
            )
            _val_projetee = _val_t1 * (1 + _rdt_trans / 100) ** _horizon_trans
            _age_futur = _age_t1 + _horizon_trans
        else:
            _rdt_trans = 0.0
            _val_projetee = _val_t1
            _age_futur = _age_t1

        if _horizon_trans > 0:
            _col_ajd, _col_futur = st.columns(2)
            with _col_ajd:
                st.markdown("**Aujourd'hui**")
                st.metric("Valeur du contrat", to_eur(_val_t1))
                st.metric("Âge du souscripteur", f"{_age_t1} ans")
            with _col_futur:
                st.markdown(f"**Dans {_horizon_trans} ans**")
                st.metric("Valeur projetée", to_eur(_val_projetee),
                          delta=f"+{to_eur(_val_projetee - _val_t1)}")
                _age_delta = "⚠️ >70 ans" if _age_futur >= 70 and _age_t1 < 70 else None
                st.metric("Âge du souscripteur", f"{_age_futur} ans", delta=_age_delta)
                if _age_t1 < 70 and _age_futur >= 70:
                    st.warning(
                        f"⚠️ Dans {70 - _age_t1} ans vous aurez 70 ans. "
                        f"Les versements effectués après cet âge relèveront de l'article 757B "
                        f"(abattement global 30 500€) au lieu du 990I (152 500€/bénéficiaire)."
                    )

        # Mettre à jour le capital décès pour les calculs
        st.session_state["tax_capital_deces"] = _val_projetee
        _tab_transmission()

        # ── Graphique projection net transmis ──
        if _val_t1 > 0:
            st.markdown("---")
            st.markdown("##### Projection : net transmis au fil du temps")
            _benef_data_t1 = st.session_state.get("tax_dash_beneficiaires", [])
            _nb_b_t1 = len(_benef_data_t1) if _benef_data_t1 else int(st.session_state.get("tax_nb_beneficiaires", 2))
            _nb_b_t1 = max(1, _nb_b_t1)
            _plafond_990i_t1 = 152_500 * _nb_b_t1

            _proj_rows_t1 = []
            for _y in range(0, 31):
                _v = _val_t1 * (1 + _default_rdt_t1 / 100) ** _y
                _a = _age_t1 + _y
                if _a < 70:
                    _part = _v / _nb_b_t1
                    _taxable = max(0.0, _part - 152_500)
                    _droits_p = _taxable * 0.20 if _taxable <= 700_000 else (
                        700_000 * 0.20 + (_taxable - 700_000) * 0.3125
                    )
                    _net_trans = _v - _droits_p * _nb_b_t1
                else:
                    _net_trans = _v * 0.80
                _proj_rows_t1.append({
                    "Année": _y,
                    "Net transmis (€)": _net_trans,
                    "Valeur brute (€)": _v,
                })

            _proj_df_t1 = pd.DataFrame(_proj_rows_t1)
            _area_brute_t1 = (
                alt.Chart(_proj_df_t1)
                .mark_area(opacity=0.15, color="#1B2A4A")
                .encode(x=alt.X("Année:Q", title="Années"), y=alt.Y("Valeur brute (€):Q", axis=alt.Axis(format=",.0f"), title="€"))
            )
            _line_net_t1 = (
                alt.Chart(_proj_df_t1)
                .mark_line(strokeWidth=2.5, color="#2E7D32")
                .encode(
                    x="Année:Q",
                    y=alt.Y("Net transmis (€):Q"),
                    tooltip=[
                        alt.Tooltip("Année:Q"),
                        alt.Tooltip("Net transmis (€):Q", format=",.0f"),
                        alt.Tooltip("Valeur brute (€):Q", format=",.0f"),
                    ],
                )
            )
            # Annotations on same DataFrame to avoid Vega-Lite "Unrecognized data set" error
            _seuil_990i_t1 = (
                alt.Chart(_proj_df_t1)
                .transform_calculate(Seuil990I=str(int(_plafond_990i_t1)))
                .mark_line(strokeDash=[5, 5], strokeWidth=2, color="#E53935")
                .encode(x="Année:Q", y="Seuil990I:Q")
            )
            _label_seuil_t1 = (
                alt.Chart(_proj_df_t1)
                .transform_filter(alt.datum["Année"] == 20)
                .transform_calculate(Seuil990I=str(int(_plafond_990i_t1)))
                .mark_text(dy=-12, fontSize=11, color="#E53935", fontWeight="bold")
                .encode(x="Année:Q", y="Seuil990I:Q",
                        text=alt.value(f"Seuil exonération 990I : {to_eur(_plafond_990i_t1)}"))
            )
            _layers_t1 = [_area_brute_t1, _line_net_t1, _seuil_990i_t1, _label_seuil_t1]
            if _age_t1 < 70:
                _y_70 = 70 - _age_t1
                _rule_70 = (
                    alt.Chart(_proj_df_t1)
                    .transform_filter(alt.datum["Année"] == _y_70)
                    .mark_rule(strokeDash=[3, 3], color="#FF8F00", strokeWidth=1.5)
                    .encode(x="Année:Q")
                )
                _label_70 = (
                    alt.Chart(_proj_df_t1)
                    .transform_filter(alt.datum["Année"] == _y_70)
                    .mark_text(dy=-10, fontSize=11, color="#FF8F00", fontWeight="bold")
                    .encode(x="Année:Q", y="Valeur brute (€):Q", text=alt.value("70 ans"))
                )
                _layers_t1 += [_rule_70, _label_70]
            st.altair_chart(alt.layer(*_layers_t1).properties(height=400), use_container_width=True)
            st.caption(
                "Ligne verte = net transmis aux bénéficiaires. "
                "Zone grise = valeur brute du contrat. "
                "Ligne rouge pointillée = seuil exonération art. 990I. "
                + ("Ligne orange = passage des 70 ans (changement de régime fiscal)." if _age_t1 < 70 else "")
            )

    with _q2:
        st.markdown("#### Ai-je intérêt à verser avant 70 ans ?")

        _age = st.session_state.get("tax_age_souscripteur", 55)
        _nb = int(st.session_state.get("tax_nb_beneficiaires", 2))
        _val_actuelle = float(st.session_state.get("tax_valeur_contrat", 100_000.0))

        if _age >= 70:
            st.warning(
                "Vous avez déjà dépassé 70 ans. Les nouveaux versements relèvent de "
                "l'article 757B (abattement global de 30 500€). "
                "Les **intérêts et plus-values** sur ces versements restent toutefois exonérés."
            )
            return

        _ans_avant_70 = 70 - _age
        _plafond_exo = 152_500 * _nb
        _marge = max(0, _plafond_exo - _val_actuelle)

        st.markdown(f"**Vous avez {_age} ans — encore {_ans_avant_70} ans pour verser avant 70 ans.**")
        st.metric("Plafond d'exonération total (990I)", to_eur(_plafond_exo),
                  delta=f"{_nb} bénéficiaire(s) × {to_eur(152_500)}")
        st.metric("Valeur actuelle du contrat", to_eur(_val_actuelle))
        st.metric("Marge de versement avant plafond", to_eur(_marge))

        st.markdown("---")
        st.markdown("**Projection : que devient votre versement ?**")

        _xirr = st.session_state.get("_tax_sync_xirr") or st.session_state.get("_LAST_XIRR_A") or st.session_state.get("_LAST_XIRR_B")
        _default_rdt = float(_xirr) if (_xirr and _xirr > 0) else 5.0

        _p1, _p2 = st.columns(2)
        with _p1:
            _versement_proj = st.number_input(
                "Versement envisagé (€)",
                min_value=0.0, max_value=5_000_000.0,
                value=min(_marge, 100_000.0) if _marge > 0 else 50_000.0,
                step=5_000.0, key="proj_versement",
            )
        with _p2:
            _rdt_proj = st.number_input(
                "Rendement annuel estimé (%)",
                min_value=-5.0, max_value=20.0,
                value=round(_default_rdt, 1), step=0.5,
                key="proj_rendement",
                help="Pré-rempli avec le rendement réel du portefeuille (XIRR) si disponible.",
            )

        _horizon = _ans_avant_70
        if _horizon > 0 and _versement_proj > 0:
            _capital_futur = _versement_proj * (1 + _rdt_proj / 100) ** _horizon
            _gains = _capital_futur - _versement_proj

            st.markdown("---")
            _pc1, _pc2, _pc3 = st.columns(3)
            with _pc1:
                st.metric("Versement aujourd'hui", to_eur(_versement_proj))
            with _pc2:
                st.metric(f"Valeur dans {_horizon} ans", to_eur(_capital_futur))
            with _pc3:
                st.metric("Plus-values générées", to_eur(_gains),
                          delta=f"+{(_capital_futur / _versement_proj - 1) * 100:.0f}%")

            if _capital_futur <= _plafond_exo:
                st.success(
                    f"✅ Si vous versez **{to_eur(_versement_proj)}** aujourd'hui et que le portefeuille "
                    f"fait **{_rdt_proj:.1f}%/an**, dans **{_horizon} ans** ce sera "
                    f"**{to_eur(_capital_futur)}** transmis à vos {_nb} bénéficiaire(s) — "
                    f"**entièrement exonéré de droits de succession** (art. 990I)."
                )
            else:
                st.info(
                    f"Le capital projeté ({to_eur(_capital_futur)}) dépasse le plafond d'exonération "
                    f"({to_eur(_plafond_exo)}). La part au-delà sera taxée à 20% puis 31.25%."
                )

            # Area chart projection + seuil exonération
            _proj_years = list(range(_horizon + 1))
            _proj_vals = [_versement_proj * (1 + _rdt_proj / 100) ** y for y in _proj_years]
            _proj_df = pd.DataFrame({"Année": _proj_years, "Valeur (€)": _proj_vals})

            _area_proj = (
                alt.Chart(_proj_df)
                .mark_area(opacity=0.25, color="#1B2A4A", interpolate="monotone")
                .encode(
                    x=alt.X("Année:Q", title="Années après versement"),
                    y=alt.Y("Valeur (€):Q", axis=alt.Axis(format=",.0f"), title="Valeur (€)"),
                )
            )
            _line_proj = (
                alt.Chart(_proj_df)
                .mark_line(strokeWidth=2.5, color="#1B2A4A", interpolate="monotone")
                .encode(
                    x="Année:Q",
                    y="Valeur (€):Q",
                    tooltip=[alt.Tooltip("Année:Q"), alt.Tooltip("Valeur (€):Q", format=",.0f")],
                )
            )
            # Annotations on same DataFrame to avoid Vega-Lite "Unrecognized data set" error
            _seuil_line = (
                alt.Chart(_proj_df)
                .transform_calculate(SeuilExo=str(int(_plafond_exo)))
                .mark_line(strokeDash=[5, 5], strokeWidth=2, color="#E53935")
                .encode(x="Année:Q", y=alt.Y("SeuilExo:Q"))
            )
            _mid_year = int(_horizon * 0.5)
            _seuil_label = (
                alt.Chart(_proj_df)
                .transform_filter(alt.datum["Année"] == _mid_year)
                .transform_calculate(SeuilExo=str(int(_plafond_exo)))
                .mark_text(dy=-12, fontSize=11, color="#E53935", fontWeight="bold")
                .encode(x="Année:Q", y="SeuilExo:Q",
                        text=alt.value(f"Seuil exonération : {to_eur(_plafond_exo)}"))
            )
            st.altair_chart(
                (_area_proj + _line_proj + _seuil_line + _seuil_label).properties(height=320),
                use_container_width=True,
            )


def _tab_avantages_av():
    """Comparaison AV vs CTO / PEA / PER — taux LFSS 2026 codés en dur."""

    # ── Taux 2026 (LFSS 2026) ──
    PS_AV = 0.172            # AV exclue de la hausse CSG → maintenu à 17.2%
    PS_CTO_PEA = 0.186       # CTO et PEA → 18.6% (CSG 10.6%)
    IR_PFU = 0.128           # IR forfaitaire (flat tax)
    IR_AV_8ANS = 0.075       # IR AV ≥8 ans, versements ≤ 150k€
    FLAT_TAX_CTO = IR_PFU + PS_CTO_PEA   # 31.4%
    PEA_PS = PS_CTO_PEA      # 18.6% — 0% IR après 5 ans
    PLAFOND_PEA = 150_000.0
    SEUIL_AV_150K = 150_000.0

    # ── Helper succession classique ──
    def _calc_succession_classique(capital_par_benef: float) -> float:
        abat = 100_000.0
        taxable = max(0.0, capital_par_benef - abat)
        if taxable <= 8_072:
            tax = taxable * 0.05
        elif taxable <= 12_109:
            tax = 8_072 * 0.05 + (taxable - 8_072) * 0.10
        elif taxable <= 15_932:
            tax = 8_072 * 0.05 + (12_109 - 8_072) * 0.10 + (taxable - 12_109) * 0.15
        elif taxable <= 552_324:
            tax = (8_072 * 0.05 + (12_109 - 8_072) * 0.10 + (15_932 - 12_109) * 0.15
                   + (taxable - 15_932) * 0.20)
        else:
            tax = (8_072 * 0.05 + (12_109 - 8_072) * 0.10 + (15_932 - 12_109) * 0.15
                   + (552_324 - 15_932) * 0.20 + (taxable - 552_324) * 0.30)
        return tax

    # ── Lecture session_state ──
    _is_client = st.session_state.get("tax_is_client", True)
    _xirr = st.session_state.get("_LAST_XIRR_A" if _is_client else "_LAST_XIRR_B")
    _val = float(st.session_state.get("_LAST_VAL_A" if _is_client else "_LAST_VAL_B", 0) or 0)
    _default_rdt = float(_xirr) if (_xirr and _xirr > 0) else 6.0
    _situation = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")
    _abat_av = 9_200.0 if "Couple" in _situation else 4_600.0
    _tmi = float(st.session_state.get("tax_tmi_PERSIST", 0.30))
    _nb_benef = int(st.session_state.get("tax_nb_beneficiaires", 2) or 2)

    # ── Inputs ──
    st.markdown("#### Comparaison des enveloppes fiscales")
    _ic1, _ic2 = st.columns(2)
    with _ic1:
        _capital_init = st.number_input(
            "Capital initial (€)", min_value=1_000.0, max_value=5_000_000.0,
            value=max(1_000.0, _val) if _val > 0 else 100_000.0,
            step=5_000.0, key="avcto_capital",
        )
        _perf_brute = st.number_input(
            "Performance brute annuelle (%)", min_value=0.0, max_value=20.0,
            value=round(_default_rdt, 1), step=0.5, key="avcto_perf",
            help="Pré-rempli avec le rendement réel du portefeuille si disponible.",
        )
    with _ic2:
        _horizon_av = st.slider(
            "Horizon (années)", min_value=5, max_value=40, value=20, key="avcto_horizon",
        )
        _freq_arb = st.selectbox(
            "Fréquence d'arbitrage (CTO)",
            ["Annuel", "Tous les 2 ans", "Tous les 5 ans"],
            index=0, key="avcto_freq",
            help="En CTO, chaque arbitrage déclenche la flat tax sur les PV réalisées.",
        )
    _freq_years = {"Annuel": 1, "Tous les 2 ans": 2, "Tous les 5 ans": 5}[_freq_arb]

    # ── Calculs communs ──
    # AV
    _val_av_brut = _capital_init * (1 + _perf_brute / 100) ** _horizon_av
    _gains_av = _val_av_brut - _capital_init
    _gains_taxables_av = max(0.0, _gains_av - _abat_av)
    if _capital_init <= SEUIL_AV_150K:
        _tax_av = _gains_taxables_av * (IR_AV_8ANS + PS_AV)
    else:
        _tax_av = _gains_taxables_av * (IR_PFU + PS_AV)
    _net_av = _val_av_brut - _tax_av

    # CTO — frottement à chaque arbitrage
    _series_cto = [_capital_init]
    _cto_val = _capital_init
    _cto_base = _capital_init
    for _y in range(1, _horizon_av + 1):
        _cto_val *= (1 + _perf_brute / 100)
        if _y % _freq_years == 0 and _y < _horizon_av:
            _pv = _cto_val - _cto_base
            if _pv > 0:
                _cto_val -= _pv * FLAT_TAX_CTO
                _cto_base = _cto_val
        _series_cto.append(_cto_val)
    _gains_cto_final = max(0.0, _cto_val - _capital_init)
    _tax_cto_final = _gains_cto_final * FLAT_TAX_CTO
    _net_cto = _cto_val - _tax_cto_final

    # PEA — plafonné, PS uniquement, 0% IR après 5 ans
    _capital_pea = min(_capital_init, PLAFOND_PEA)
    _val_pea_base = _capital_pea * (1 + _perf_brute / 100) ** _horizon_av
    _gains_pea = _val_pea_base - _capital_pea
    _tax_pea = _gains_pea * PEA_PS
    _net_pea = _val_pea_base - _tax_pea
    _note_pea_surplus = False
    if _capital_init > PLAFOND_PEA:
        _surplus = _capital_init - PLAFOND_PEA
        _surplus_val = _surplus * (1 + _perf_brute / 100) ** _horizon_av
        _surplus_tax = max(0.0, _surplus_val - _surplus) * FLAT_TAX_CTO
        _net_pea += _surplus_val - _surplus_tax
        _note_pea_surplus = True

    # PER
    _eco_entree_per = _capital_init * _tmi
    _val_per = _capital_init * (1 + _perf_brute / 100) ** _horizon_av
    _tax_per_sortie = _val_per * (_tmi + PS_CTO_PEA)
    _net_per = _val_per - _tax_per_sortie + _eco_entree_per

    # ── Rachat annuel sans IR (C3) ──
    _ratio_gains = _gains_av / _val_av_brut if _val_av_brut > 0 else 0.0
    _rachat_annuel_sans_ir = _abat_av / _ratio_gains if _ratio_gains > 0 else _val_av_brut

    # ── Transmission — calculé avant les onglets (C6) ──
    _nb_benef = max(1, _nb_benef)
    _parts_pct = [100.0 / _nb_benef] * _nb_benef
    _types_benef = ["Enfant"] * _nb_benef
    try:
        _av_transm_results = calc_transmission_990I(
            _val_av_brut, _nb_benef, _parts_pct, _types_benef
        )
        _total_net_av_transmission = sum(r["net_recu"] for r in _av_transm_results)  # C1
        _av_tax_total = sum(r.get("taxe", 0) for r in _av_transm_results)
    except Exception:
        _av_tax_total = max(0.0, _val_av_brut - 152_500.0 * _nb_benef) * 0.20
        _total_net_av_transmission = _val_av_brut - _av_tax_total
    _cto_par_benef = _net_cto / _nb_benef  # C1 — après flat tax
    _taxe_cto_total = sum(_calc_succession_classique(_cto_par_benef) for _ in range(_nb_benef))
    _total_net_cto_transmission = _net_cto - _taxe_cto_total
    _delta_transm = _total_net_av_transmission - _total_net_cto_transmission

    # ── Onglets ──
    _tab_synth, _tab_cto, _tab_pea, _tab_per = st.tabs([
        "📋 Synthèse", "AV (≥8 ans) vs CTO", "AV (≥8 ans) vs PEA", "AV vs PER",
    ])

    # ────────────────────────────────────────────
    # ONGLET 1 — Synthèse
    # ────────────────────────────────────────────
    with _tab_synth:
        _ir_av_lbl = "7.5%" if _capital_init <= SEUIL_AV_150K else "12.8%"
        _ir_av_cell = (
            f"{_ir_av_lbl} après 8 ans (≤150k€ versés)"
            if _capital_init <= SEUIL_AV_150K
            else "12.8% après 8 ans"
        )
        _synth_df = pd.DataFrame({
            "Critère": [
                "Fiscalité pendant la détention",
                "PS applicables (2026)",
                "IR à la sortie",
                "Abattement",
                "Rachat annuel sans IR",
                "Plafond versements",
                "Univers d'investissement",
                "Disponibilité",
                "Transmission",
                "Déduction à l'entrée",
                f"Net après impôt ({_horizon_av} ans)",
            ],
            "Assurance-vie": [
                "Aucune",
                "17.2% ✅",
                _ir_av_cell,
                f"{to_eur(_abat_av)}/an après 8 ans",
                f"≈ {to_eur(_rachat_annuel_sans_ir)}/an (contrat ≥8 ans)",  # C3
                "Aucun",
                "UC + Fonds euros + Structurés",
                "✅ Totale",
                "✅ 152 500€/bénéf. hors succession",
                "Non",
                to_eur(_net_av),
            ],
            "CTO": [
                f"Flat tax {FLAT_TAX_CTO*100:.1f}% à chaque arbitrage",
                "18.6%",
                "12.8%",
                "Aucun",
                "Aucun",
                "Aucun",
                "Tout",
                "✅ Totale",
                "❌ Succession classique",
                "Non",
                to_eur(_net_cto),
            ],
            "PEA": [
                "Aucune",
                "18.6%",
                "0% (≥5 ans)",
                "Aucun",
                "Aucun",
                "150 000 €",
                "Actions européennes + ETF synth.",
                "✅ Après 5 ans",
                "❌ Succession classique",
                "Non",
                to_eur(_net_pea),
            ],
            f"PER (TMI {_tmi*100:.0f}%)": [
                "Aucune",
                "18.6%",
                f"TMI {_tmi*100:.0f}%",
                "10% sur rente",
                "Aucun",
                "~10% revenus N-1",
                "UC + Fonds euros",
                "❌ Bloqué retraite",
                "❌ Succession classique",
                f"✅ TMI {_tmi*100:.0f}%",
                to_eur(_net_per),
            ],
        })
        st.dataframe(_synth_df, hide_index=True, use_container_width=True)

        # C4 — message nuancé : AV toujours présentée comme meilleur compromis global
        if _net_pea > _net_av:
            st.success(
                f"✅ Sur {_horizon_av} ans, le PEA offre un léger avantage fiscal "
                f"({to_eur(_net_pea)} vs {to_eur(_net_av)}). "
                f"Cependant, **l'AV (contrat ≥8 ans) reste le meilleur compromis global** : "
                f"disponibilité totale, PS maintenus à 17.2% (vs 18.6% ailleurs), "
                f"transmission hors succession (152 500€/bénéf.) et pas de plafond. "
                f"Rachat annuel sans IR jusqu'à {to_eur(_rachat_annuel_sans_ir)}/an "
                f"(abattement {to_eur(_abat_av)})."
            )
        else:
            st.success(
                f"✅ Sur {_horizon_av} ans, **l'assurance-vie (contrat ≥8 ans) est l'enveloppe "
                f"la plus avantageuse** avec un net après impôt de {to_eur(_net_av)}. "
                f"Elle combine disponibilité totale, PS maintenus à 17.2%, "
                f"transmission hors succession et abattement annuel de {to_eur(_abat_av)}. "
                f"Rachat annuel sans IR jusqu'à {to_eur(_rachat_annuel_sans_ir)}/an."
            )
        if _note_pea_surplus:
            st.caption(
                "⚠️ Capital > 150k€ : le PEA est plafonné, le surplus est simulé en CTO à 31.4%."
            )
        st.caption(
            "⚠️ PEA : rendement optimiste — l'univers est limité aux actions européennes "
            "et ETF synthétiques. La diversification est plus contrainte qu'en AV."
        )

    # ────────────────────────────────────────────
    # ONGLET 2 — AV vs CTO
    # ────────────────────────────────────────────
    with _tab_cto:
        st.markdown("#### L'effet capitalisation : AV (≥8 ans) vs CTO")

        _years = list(range(_horizon_av + 1))
        _series_av_list = [_capital_init * (1 + _perf_brute / 100) ** y for y in _years]

        _lbl_av_cto = "AV (≥8 ans) — PS 17.2%"
        _lbl_cto = f"CTO — flat tax {FLAT_TAX_CTO*100:.1f}%"
        _chart_df = pd.DataFrame({
            "Année": _years * 2,
            "Valeur (€)": _series_av_list + _series_cto,
            "Enveloppe": (
                [_lbl_av_cto] * len(_years)
                + [_lbl_cto] * len(_years)
            ),
            "AV": _series_av_list * 2,
            "CTO": _series_cto * 2,
        })

        _area_cto = (
            alt.Chart(_chart_df)
            .transform_filter(alt.datum["Enveloppe"] == _lbl_av_cto)
            .mark_area(opacity=0.12, color="#2E7D32")
            .encode(
                x=alt.X("Année:Q"),
                y=alt.Y("CTO:Q"),
                y2="AV:Q",
            )
        )
        _line_cto = (
            alt.Chart(_chart_df)
            .mark_line(strokeWidth=2.5)
            .encode(
                x=alt.X("Année:Q", title="Années"),
                y=alt.Y(
                    "Valeur (€):Q",
                    axis=alt.Axis(format=",.0f"),
                    title="Valeur (€)",
                    scale=alt.Scale(zero=False),  # C5
                ),
                color=alt.Color(
                    "Enveloppe:N",
                    scale=alt.Scale(
                        domain=[_lbl_av_cto, _lbl_cto],
                        range=["#1B2A4A", "#CC2200"],
                    ),
                    legend=alt.Legend(title="", orient="bottom"),
                ),
                tooltip=[
                    alt.Tooltip("Année:Q"),
                    alt.Tooltip("Enveloppe:N"),
                    alt.Tooltip("Valeur (€):Q", format=",.0f"),
                ],
            )
            .properties(height=380)
        )
        st.altair_chart((_area_cto + _line_cto).properties(height=380), use_container_width=True)

        _ecart_cto = _net_av - _net_cto
        _ecart_cto_pct = (_ecart_cto / _net_cto * 100) if _net_cto > 0 else 0.0
        _mc1, _mc2, _mc3 = st.columns(3)
        with _mc1:
            st.metric("Net après impôt — AV", to_eur(_net_av))
        with _mc2:
            st.metric("Net après impôt — CTO", to_eur(_net_cto))
        with _mc3:
            st.metric(
                "Avantage AV", to_eur(_ecart_cto),
                delta=f"+{_ecart_cto_pct:.0f}%" if _ecart_cto > 0 else f"{_ecart_cto_pct:.0f}%",
            )
        st.caption(
            f"💬 Sur {_horizon_av} ans, l'AV (contrat de plus de 8 ans) génère "
            f"{to_eur(_ecart_cto)} de plus grâce à la capitalisation sans frottement. "
            f"PS AV : 17.2% vs CTO : 18.6% (LFSS 2026)."
        )

    # ────────────────────────────────────────────
    # ONGLET 3 — AV vs PEA
    # ────────────────────────────────────────────
    with _tab_pea:
        st.markdown("#### AV (≥8 ans) vs PEA")

        _series_pea_list = [
            min(_capital_init, PLAFOND_PEA) * (1 + _perf_brute / 100) ** y
            for y in _years
        ]

        _lbl_av_pea = "AV (≥8 ans) — PS 17.2%"
        _lbl_pea = "PEA — PS 18.6%, 0% IR"
        _chart_pea_df = pd.DataFrame({
            "Année": _years * 2,
            "Valeur (€)": _series_av_list + _series_pea_list,
            "Enveloppe": (
                [_lbl_av_pea] * len(_years)
                + [_lbl_pea] * len(_years)
            ),
        })

        _line_pea = (
            alt.Chart(_chart_pea_df)
            .mark_line(strokeWidth=2.5)
            .encode(
                x=alt.X("Année:Q", title="Années"),
                y=alt.Y(
                    "Valeur (€):Q",
                    axis=alt.Axis(format=",.0f"),
                    title="Valeur (€)",
                    scale=alt.Scale(zero=False),  # C5
                ),
                color=alt.Color(
                    "Enveloppe:N",
                    scale=alt.Scale(
                        domain=[_lbl_av_pea, _lbl_pea],
                        range=["#1B2A4A", "#2E7D32"],
                    ),
                    legend=alt.Legend(title="", orient="bottom"),
                ),
                tooltip=[
                    alt.Tooltip("Année:Q"),
                    alt.Tooltip("Enveloppe:N"),
                    alt.Tooltip("Valeur (€):Q", format=",.0f"),
                ],
            )
            .properties(height=380)
        )
        st.altair_chart(_line_pea, use_container_width=True)

        _ecart_pea = _net_av - _net_pea
        _mp1, _mp2, _mp3 = st.columns(3)
        with _mp1:
            st.metric("Net après impôt — AV", to_eur(_net_av))
        with _mp2:
            st.metric("Net après impôt — PEA", to_eur(_net_pea))
        with _mp3:
            st.metric(
                "Delta AV vs PEA", to_eur(_ecart_pea),
                delta=f"+{_ecart_pea/_net_pea*100:.0f}%" if (_net_pea > 0 and _ecart_pea > 0) else None,
            )

        if _capital_init <= PLAFOND_PEA and _net_pea > _net_av:
            st.info(
                f"💬 Le PEA offre un meilleur net ({to_eur(_net_pea)}) grâce à l'exonération d'IR. "
                f"Mais il est plafonné à 150k€, limité aux actions européennes/ETF synthétiques, "
                f"et n'offre aucun avantage transmission."
            )
        elif _capital_init > PLAFOND_PEA:
            st.info(
                f"💬 Le PEA est plafonné à 150 000 €. Au-delà, le surplus est investi en CTO "
                f"(flat tax {FLAT_TAX_CTO*100:.1f}%). L'AV n'a aucun plafond et offre une "
                f"transmission avantageuse."
            )
        else:
            st.success(
                f"✅ L'AV génère {to_eur(abs(_ecart_pea))} de plus que le PEA sur {_horizon_av} ans, "
                f"tout en offrant disponibilité totale et avantage de transmission."
            )
        st.caption(
            "⚠️ Rendement PEA optimiste : univers limité aux actions européennes et ETF synthétiques."
        )

    # ────────────────────────────────────────────
    # ONGLET 4 — AV vs PER
    # ────────────────────────────────────────────
    with _tab_per:
        st.markdown("#### AV vs PER")

        _comp_per_df = pd.DataFrame({
            "Critère": [
                "Disponibilité du capital",
                "Fiscalité à l'entrée",
                "Fiscalité à la sortie (capital)",
                "Fiscalité à la sortie (rente)",
                "Transmission",
                "Plafond de versement",
                "Cas de déblocage anticipé",
            ],
            "Assurance-vie": [
                "✅ Totale, à tout moment",
                "Aucune déduction",
                f"IR {IR_AV_8ANS*100:.1f}% + PS {PS_AV*100:.1f}% après 8 ans (avec abattement)",
                "Fraction imposable selon l'âge (30–70%)",
                "✅ Hors succession (152 500€/bénéficiaire)",
                "Aucun plafond",
                "Non applicable (capital toujours disponible)",
            ],
            "PER": [
                "❌ Bloqué jusqu'à la retraite",
                f"✅ Déductible — économie {to_eur(_eco_entree_per)} à votre TMI {_tmi*100:.0f}%",
                f"IR {_tmi*100:.0f}% + PS {PS_CTO_PEA*100:.1f}% sur la totalité",
                f"IR {_tmi*100:.0f}% + PS {PS_CTO_PEA*100:.1f}% sur la totalité",
                "❌ Succession classique (sauf exceptions)",
                "Plafonné (10% revenus N-1)",
                "Achat résidence principale, invalidité, décès conjoint...",
            ],
        })
        st.dataframe(_comp_per_df, hide_index=True, use_container_width=True)

        _ecart_per = _net_av - _net_per
        _mr1, _mr2, _mr3 = st.columns(3)
        with _mr1:
            st.metric("Net après impôt — AV", to_eur(_net_av))
        with _mr2:
            st.metric("Net après impôt — PER", to_eur(_net_per))
        with _mr3:
            st.metric("Delta AV vs PER", to_eur(_ecart_per))

        if _tmi >= 0.41:
            st.info(
                f"Le PER est avantageux pour votre TMI ({_tmi*100:.0f}%) grâce à la déduction "
                f"immédiate de {to_eur(_eco_entree_per)}. Mais le capital est bloqué jusqu'à la retraite."
            )
        else:
            st.success(
                f"✅ L'AV est plus avantageuse à votre TMI ({_tmi*100:.0f}%). "
                f"Elle offre disponibilité + transmission + PS à 17.2%."
            )

    # ────────────────────────────────────────────
    # BLOC TRANSMISSION (toujours visible, une seule fois — C6)
    # ────────────────────────────────────────────
    st.markdown("---")
    st.markdown("##### 🏠 Avantage transmission")

    _mt1, _mt2, _mt3 = st.columns(3)
    with _mt1:
        st.metric("Net transmis — AV", to_eur(_total_net_av_transmission))
    with _mt2:
        st.metric("Net transmis — CTO", to_eur(_total_net_cto_transmission))
    with _mt3:
        st.metric("Avantage AV en transmission", to_eur(_delta_transm))

    if _delta_transm > 0:
        st.success(
            f"✅ L'AV (contrat de plus de 8 ans) permet de transmettre {to_eur(_delta_transm)} "
            f"de plus grâce à l'abattement de 152 500 € par bénéficiaire (art. 990 I CGI). "
            f"Rachat annuel sans IR possible jusqu'à {to_eur(_rachat_annuel_sans_ir)}/an "
            f"(abattement {to_eur(_abat_av)}/an)."
        )


def run_perfect_portfolio():
    render_portfolio_builder()


# ──────────────────────────────────────────────────────────────────────────────
# Section Rapports & Présentations
# ──────────────────────────────────────────────────────────────────────────────

def render_reports_section() -> None:
    st.title("📥 Rapports & Présentations")
    _nom = st.session_state.get("NOM_CLIENT", "").strip()
    if _nom:
        st.markdown(f"**Client :** {_nom}")

    report_data = st.session_state.get("_LAST_REPORT_DATA")
    if not report_data:
        st.warning(
            "⚠️ Aucune analyse disponible. Configurez un portefeuille dans le "
            "**Comparateur** et lancez une analyse d'abord — les rapports seront "
            "ensuite générables depuis cette section."
        )
        return

    # ── Toggle : inclure la fiscalité dans la présentation ──
    _include_fiscal = st.checkbox(
        "Inclure les slides fiscalité (retrait, transmission, projections)",
        value=True,
        key="reports_include_fiscal",
        help="Décochez pour générer une présentation sans les slides fiscales (slides 3 et 6).",
    )
    report_data = {**report_data, "include_fiscal_slides": _include_fiscal}

    st.markdown("---")

    _nom_cli = (report_data.get("nom_client", "") or "client").replace(" ", "_")
    _c1, _c2, _c3 = st.columns(3)

    with _c1:
        with st.container(border=True):
            st.markdown("#### 📄 Rapport PDF")
            st.caption("Rapport détaillé avec composition, performance, frais, diversification.")
            try:
                _pdf = generate_pdf_report(report_data)
                st.download_button(
                    "⬇️ Télécharger le PDF",
                    data=_pdf,
                    file_name=f"rapport_{_nom_cli}.pdf",
                    mime="application/pdf",
                    key="dl_pdf_reports",
                    use_container_width=True,
                )
            except Exception as e:
                st.error(f"PDF indisponible : {e}")

    with _c2:
        with st.container(border=True):
            st.markdown("#### 📊 Présentation client")
            st.caption("PowerPoint 9 slides : Bilan & Recommandation.")
            try:
                _pptx = generate_pptx_report(report_data)
                st.download_button(
                    "⬇️ Télécharger la présentation",
                    data=_pptx,
                    file_name=f"presentation_{_nom_cli}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="dl_pptx_reports",
                    use_container_width=True,
                )
            except Exception as e:
                st.error(f"PPTX indisponible : {e}")

    with _c3:
        with st.container(border=True):
            st.markdown("#### 📝 Proposition d'arbitrage")
            st.caption("One-pager PDF résumant la proposition.")
            try:
                _onepager = _build_onepager_pdf(report_data)
                if _onepager:
                    st.download_button(
                        "⬇️ Télécharger le one-pager",
                        data=_onepager,
                        file_name=f"arbitrage_{_nom_cli}.pdf",
                        mime="application/pdf",
                        key="dl_onepager_reports",
                        use_container_width=True,
                    )
                else:
                    st.info("One-pager non disponible (Reportlab requis).")
            except Exception as e:
                st.error(f"One-pager indisponible : {e}")


def render_mode_router():
    st.set_page_config(page_title=APP_TITLE, layout="wide")
    with st.sidebar:
        _sidebar_title = st.session_state.get("NOM_CABINET", "").strip()
        st.markdown(f"## {_sidebar_title if _sidebar_title else APP_TITLE}")
        st.caption(APP_SUBTITLE)
        st.session_state.setdefault("NOM_CLIENT", "")
        st.session_state.setdefault("NOM_CABINET", "")
        _nom_client = st.text_input(
            "Nom du client",
            value=st.session_state.get("NOM_CLIENT", ""),
            key="_global_nom_client",
            placeholder="Ex : M. et Mme Dupont",
        )
        st.session_state["NOM_CLIENT"] = _nom_client
        _nom_cabinet = st.text_input(
            "Nom du cabinet",
            value=st.session_state.get("NOM_CABINET", ""),
            key="_global_nom_cabinet",
            placeholder="Ex : Cabinet Dupont Patrimoine",
        )
        st.session_state["NOM_CABINET"] = _nom_cabinet
        st.divider()
        _nav_options = [
            "📊 Comparateur",
            "🏗️ Construction optimisée",
            "🧾 Fiscalité & Avantages AV",
            "📥 Rapports & Présentations",
        ]
        # Permettre la navigation programmatique via APP_MODE
        st.session_state.pop("APP_MODE", None)  # Nettoyage — plus utilisé
        mode = st.radio(
            "Navigation",
            _nav_options,
            label_visibility="collapsed",
            key="_nav_mode_radio",
        )
    # Reset du mode amélioration quand on quitte le builder
    if mode != "🏗️ Construction optimisée":
        st.session_state.pop("PP_IMPROVE_MODE", None)
        st.session_state.pop("PP_IMPROVE_LINES", None)
        st.session_state.pop("PP_IMPROVE_SOURCE", None)
        st.session_state.pop("_PP_IMPROVE_APPLIED", None)

    if mode == "📊 Comparateur":
        run_comparator()
    elif mode == "🏗️ Construction optimisée":
        run_perfect_portfolio()
    elif mode == "🧾 Fiscalité & Avantages AV":
        render_tax_module()
    else:
        render_reports_section()  # "📥 Rapports & Présentations"


def _render_with_crash_shield():
    try:
        render_mode_router()
        st.session_state["APP_STATUS"] = "OK"
    except Exception as e:
        st.session_state["APP_STATUS"] = "KO"
        st.session_state["LAST_EXCEPTION"] = str(e)
        st.title(APP_TITLE)
        st.info("App chargée, statut KO")
        st.error("Une erreur est survenue pendant le rendu.")
        st.exception(e)
        st.markdown("""
Conseils :
- Vérifiez vos dépendances (reportlab/matplotlib).
- Vérifiez la clé EODHD dans les secrets.
- Réessayez après avoir vidé le cache Streamlit.
""")


_render_with_crash_shield()
