from __future__ import annotations

import json
import os
import sys
import textwrap
import importlib.util
import itertools
import uuid  # FIXED: stable UUID keys for fund lines (Bug 5)
from datetime import date
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

if importlib.util.find_spec("matplotlib") is not None:
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
    MATPLOTLIB_ERROR = ""
else:
    plt = None
    MATPLOTLIB_AVAILABLE = False
    MATPLOTLIB_ERROR = "matplotlib non installé"
import numpy as np
import pandas as pd
import requests
import streamlit as st
import altair as alt
if importlib.util.find_spec("pypfopt") is not None:
    from pypfopt import EfficientFrontier, risk_models, expected_returns
    PYPFOPT_AVAILABLE = True
    PYPFOPT_ERROR = ""
else:
    EfficientFrontier = None
    risk_models = None
    expected_returns = None
    PYPFOPT_AVAILABLE = False
    PYPFOPT_ERROR = "pyportfolioopt non installé"
if importlib.util.find_spec("reportlab") is not None:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
    REPORTLAB_ERROR = ""
else:
    SimpleDocTemplate = Paragraph = Spacer = Image = Table = TableStyle = PageBreak = None
    A4 = None
    getSampleStyleSheet = None
    colors = None
    ParagraphStyle = None
    canvas = None
    REPORTLAB_AVAILABLE = False
    REPORTLAB_ERROR = "reportlab non installé"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import pptx.util as pptx_util
    PPTX_AVAILABLE = True
    PPTX_ERROR = ""
except ImportError as e:
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    PPTX_AVAILABLE = False
    PPTX_ERROR = str(e)


# ------------------------------------------------------------
# Constantes & univers de fonds recommandés
# ------------------------------------------------------------
TODAY = pd.Timestamp.today().normalize()
APP_TITLE = "Analyse Patrimoniale"
APP_SUBTITLE = "Outil de conseil en gestion de patrimoine"
ANNUAL_FEE_EURO_PCT = 0.9
ANNUAL_FEE_UC_PCT = 1.2
RISK_FREE_RATE_FALLBACK = 0.026  # 2.6 % — valeur de repli si API Bund indisponible

# Registre des contrats disponibles
# Clé = label affiché à l'utilisateur
# Valeur = dict avec chemin des fichiers
CONTRACTS_REGISTRY: Dict[str, Any] = {
    "Linxea Avenir 2": {
        "path": "linxea/avenir2",
        "funds_filename": "Liste des fonds LINXEA AVENIR 2.csv",
        "euro_funds": {
            "Suravenir Opportunités 2": "suravenir_opportunites2_historique.csv",
            "Suravenir Rendement 2":    "suravenir_rendement2_historique.csv",
        },
        "assureur": "Suravenir",
        "entry_fee_max_pct": 0.0,
    },
    "Linxea Spirit 2": {
        "path": "linxea/spirit2",
        "funds_filename": "SPIRIT2_UC_frais_avec_categories.csv",
        "euro_funds": {
            "Euro Nouvelle Génération": "spirit2_euro_nouvelle_generation.csv",
            "Euro Objectif Climat":     "spirit2_euro_objectif_climat.csv",
        },
        "assureur": "Spirica",
        "entry_fee_max_pct": 0.0,
    },
}

RECO_FUNDS_CORE = [
    ("R-co Valor C EUR", "FR0011253624"),
    ("Vivalor International", "FR0014001LS1"),
    ("CARMIGNAC Investissement A EUR", "FR0010148981"),
    ("FIDELITY FUNDS - WORLD FUND", "LU0069449576"),
    ("CLARTAN VALEURS", "LU1100076550"),
    ("CARMIGNAC PATRIMOINE", "FR0010135103"),
]

RECO_FUNDS_DEF = [
    ("Fonds en euros (EUROFUND)", "EUROFUND"),
    ("SYCOYIELD 2030 RC", "FR001400MCQ6"),
    ("R-Co Target 2029 HY", "FR0014002XJ3"),
    ("Euro Bond 1-3 Years", "LU0321462953"),
]

FUND_NAME_MAP = {isin: name for name, isin in RECO_FUNDS_CORE + RECO_FUNDS_DEF}

# Univers UC explicites (hors fonds en euros)
EQUITY_FUNDS = [isin for _, isin in RECO_FUNDS_CORE]
BOND_FUNDS = [isin for _, isin in RECO_FUNDS_DEF if isin != "EUROFUND"]

# Libellés FR -> codes internes pour l'affectation des versements
ALLOC_LABELS = {
    "Répartition égale": "equal",
    "Personnalisé": "custom",
    "Tout sur une ligne": "single",
}


# ------------------------------------------------------------
# Utils format
# ------------------------------------------------------------

def params_hash(values: Tuple[Any, ...]) -> str:
    try:
        payload = json.dumps(values, default=str, sort_keys=True)
    except Exception:
        payload = repr(values)
    return str(hash(payload))

def to_eur(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return s + " €"


def fmt_date(x: Any) -> str:
    try:
        return pd.Timestamp(x).strftime("%d/%m/%Y")
    except Exception:
        return "—"


def fmt_eur_fr(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return f"{s} €"


def fmt_pct_fr(x: Any) -> str:
    try:
        v = float(x)
    except Exception:
        return "—"
    s = f"{v:,.2f}".replace(",", "X").replace(".", ",").replace("X", " ")
    return f"{s} %"


# ------------------------------------------------------------
# XIRR
# ------------------------------------------------------------

def _npv(rate: float, cfs: List[Tuple[pd.Timestamp, float]]) -> float:
    t0 = cfs[0][0]
    return sum(cf / ((1 + rate) ** ((t - t0).days / 365.25)) for t, cf in cfs)


def xirr(cash_flows: List[Tuple[pd.Timestamp, float]], guess: float = 0.1) -> Optional[float]:
    if not cash_flows:
        return None
    cfs = sorted(cash_flows, key=lambda x: x[0])
    try:
        r = guess
        for _ in range(100):
            f = _npv(r, cfs)
            h = 1e-6
            f1 = _npv(r + h, cfs)
            d = (f1 - f) / h
            if abs(d) < 1e-12:
                break
            r2 = r - f / d
            if abs(r2 - r) < 1e-9:
                r = r2
                break
            r = r2
        if abs(_npv(r, cfs)) > 1.0:  # NPV non nulle = pas convergé
            return None
        return r
    except Exception:
        return None


# ------------------------------------------------------------
# API EODHD
# ------------------------------------------------------------

def _get_api_key() -> str:
    return st.secrets.get("EODHD_API_KEY", "")


@st.cache_data(show_spinner=False, ttl=3600)
def eodhd_get(path: str, params: Dict[str, Any] | None = None) -> Any:
    # IMPORTANT : ne jamais appeler st.* (spinner, warning, etc.)
    # dans une fonction @st.cache_data — Streamlit intercepte ces appels
    # lors du replay du cache au premier rendu et lève une exception
    # silencieuse qui vide le résultat.
    base = "https://eodhd.com/api"
    token = _get_api_key()
    p = {"api_token": token, "fmt": "json"}
    if params:
        p.update(params)
    r = requests.get(f"{base}{path}", params=p, timeout=20)
    r.raise_for_status()
    try:
        return r.json()
    except Exception:
        return None


@st.cache_data(show_spinner=False, ttl=3600)
def eodhd_search(q: str) -> List[Dict[str, Any]]:
    try:
        js = eodhd_get(f"/search/{q}")
        if isinstance(js, list):
            return js
    except Exception:
        pass
    return []


@st.cache_data(show_spinner=False, ttl=86400)
def _fetch_bund_rate_from_api() -> Optional[float]:
    """
    Tente de récupérer le rendement du Bund allemand 10 ans via EODHD.
    Retourne le taux en décimal (ex: 0.026) ou None si indisponible.
    Essaie plusieurs tickers dans l'ordre en cas d'échec.
    TTL : 24 h (le taux change peu en intra-journalier).
    """
    tickers = ["DE10Y.GBOND", "BUND10Y.INDX", "IRLTLT01DEM156N.FRED"]
    for ticker in tickers:
        try:
            data = eodhd_get(
                f"/eod/{ticker}",
                {"fmt": "json", "order": "d", "limit": 1},
            )
            if isinstance(data, list) and len(data) > 0:
                rate_pct = float(data[0].get("close", 0.0))
                if 0.0 < rate_pct < 15.0:
                    return rate_pct / 100.0
        except Exception:
            continue
    return None


def get_risk_free_rate() -> float:
    """
    Retourne le taux sans risque actif en décimal.
    Priorité :
    1. API EODHD (Bund 10 ans, si disponible)
    2. Saisie manuelle dans la sidebar (RISK_FREE_RATE_MANUAL, stocké en décimal)
    3. Fallback RISK_FREE_RATE_FALLBACK
    """
    api_rate = _fetch_bund_rate_from_api()
    if api_rate is not None:
        return api_rate
    return float(
        st.session_state.get("RISK_FREE_RATE_MANUAL", RISK_FREE_RATE_FALLBACK)
    )


@st.cache_data(show_spinner=False, ttl=86400)
def eodhd_prices_daily(symbol: str) -> pd.DataFrame:
    try:
        js = eodhd_get(f"/eod/{symbol}", params={"period": "d"})
        if not isinstance(js, list) or len(js) == 0:
            return pd.DataFrame()
        df = pd.DataFrame(js)
        df["date"] = pd.to_datetime(df["date"])
        df.set_index("date", inplace=True)
        if "adjusted_close" in df.columns and pd.notnull(df["adjusted_close"]).any():
            df["Close"] = df["adjusted_close"].astype(float)
        elif "close" in df.columns:
            df["Close"] = df["close"].astype(float)
        else:
            return pd.DataFrame()
        return df[["Close"]].sort_index()
    except Exception:
        return pd.DataFrame()


def _symbol_candidates(isin_or_name: str) -> List[str]:
    val = str(isin_or_name).strip()
    if not val:
        return []
    if val.upper() == "EUROFUND":
        return ["EUROFUND"]
    candidates = [f"{val}.EUFUND", f"{val}.FUND", val]
    try:
        res = eodhd_search(val)
        for it in res:
            code = it.get("Code")
            exch = it.get("Exchange")
            if code and exch:
                candidates.append(f"{code}.{exch}")
    except Exception:
        pass
    seen = set()
    uniq = []
    for c in candidates:
        if c not in seen:
            seen.add(c)
            uniq.append(c)
    return uniq


def _get_close_on(df: pd.DataFrame, d: pd.Timestamp) -> float:
    if df.empty:
        return np.nan
    if d in df.index:
        return float(df.loc[d, "Close"])
    after = df.loc[df.index >= d]
    if not after.empty:
        return float(after.iloc[0]["Close"])
    return float(df.iloc[-1]["Close"])


def apply_annual_fee(
    df: pd.DataFrame,
    annual_fee_pct: float,
    buy_date: Optional[pd.Timestamp] = None,  # FIXED: anchor fee deduction to buy date, not fund inception (Bug 1)
) -> pd.DataFrame:
    # ⚠️ USAGE CORRECT de cette fonction :
    # - Toujours fournir buy_date (date d'achat du client), jamais None
    # - Réservée aux calculs de simulation et de performance client
    #   (XIRR, val_now = qty × last_px_net)
    # - NE PAS utiliser pour afficher une VL de marché absolue
    # - Pour la VL de marché brute : utiliser get_current_nav(isin)
    #   ou get_price_series(isin, None, 0.0)
    # - Pour les calculs de corrélation/vol/Sharpe : utiliser
    #   get_series_for_line(..., apply_fees=False)
    if df.empty or annual_fee_pct == 0:
        return df
    df = df.copy()
    fee_rate = float(annual_fee_pct) / 100.0
    # FIXED: use buy_date as base to avoid applying 20+ years of fictitious fees (Bug 1)
    base_date = pd.Timestamp(buy_date).normalize() if buy_date is not None else df.index[0]
    # FIXED: clamp to zero so pre-buy-date rows (if any) are never penalised (Bug 1)
    day_offsets = np.maximum(0.0, (df.index - base_date).days.astype(float))
    fee_factors = (1.0 - fee_rate) ** (day_offsets / 365.25)
    df["Close"] = df["Close"].astype(float).to_numpy() * fee_factors
    return df


# ------------------------------------------------------------
# Loaders données contrat
# ------------------------------------------------------------

# Chemin racine des données contrats
CONTRACT_DATA_ROOT = "data"


@st.cache_data(show_spinner=False, ttl=3600)
def load_contract_funds(contract_path: str, funds_filename: str) -> pd.DataFrame:
    """
    Charge le fichier Liste des fonds d'un contrat.
    Nettoie les frais (virgule→point, supprime %) et les convertit en float.
    Retourne un DataFrame avec colonnes normalisées :
      isin, name, manager, category, sri, fee_uc_pct, fee_contract_pct, fee_total_pct
    """
    path = os.path.join(CONTRACT_DATA_ROOT, contract_path, funds_filename)
    try:
        df = pd.read_csv(path, sep=";", encoding="utf-8-sig", dtype=str)
    except FileNotFoundError:
        st.error(
            f"Fichier référentiel introuvable : {path}\n"
            f"Vérifiez que le fichier existe bien dans le dépôt GitHub "
            f"avec ce nom exact (attention aux espaces et majuscules)."
        )
        return pd.DataFrame()
    except Exception as e:
        st.error(f"Erreur chargement référentiel : {e}")
        return pd.DataFrame()
    df.columns = df.columns.str.strip()

    def _parse_pct(val: str) -> float:
        try:
            return float(str(val).replace("%", "").replace(",", ".").strip())
        except Exception:
            return 0.0

    result = pd.DataFrame({
        "isin":             df["Code ISIN"].str.strip(),
        "name":             df["Libellé du fonds"].str.strip(),
        "manager":          df["Société de gestion"].str.strip(),
        "category":         df["Catégorie Morningstar"].str.strip(),
        "sri":              pd.to_numeric(df["Risque (SRI)"], errors="coerce").fillna(0).astype(int),
        "fee_uc_pct":       df["Frais UC (B) %"].apply(_parse_pct),
        "fee_contract_pct": df["Frais contrat (C) %"].apply(_parse_pct),
        "fee_total_pct":    df["Frais totaux (B+C) %"].apply(_parse_pct),
    })
    return result.dropna(subset=["isin"]).reset_index(drop=True)


@st.cache_data(show_spinner=False, ttl=3600)
def load_euro_fund_history(contract_path: str, fund_filename: str) -> pd.DataFrame:
    """
    Charge l'historique d'un fonds euros.
    Retourne un DataFrame avec colonnes : annee (int), taux_net_publie_pct (float).
    La colonne utilisée pour la simulation est taux_net_publie_% (net de frais,
    avant prélèvements sociaux — c'est le rendement servi sur le contrat).
    """
    path = os.path.join(CONTRACT_DATA_ROOT, contract_path, fund_filename)
    df = pd.read_csv(path, sep=";", encoding="utf-8-sig")
    df.columns = df.columns.str.strip()
    df["annee"] = pd.to_numeric(df["annee"], errors="coerce").astype("Int64")
    df["taux_net_publie_pct"] = pd.to_numeric(
        df["taux_net_publie_%"], errors="coerce"
    )
    return df[["annee", "taux_net_publie_pct"]].dropna().reset_index(drop=True)


def get_euro_fund_avg_rate(history_df: pd.DataFrame, years: int = 5) -> float:
    """
    Calcule la moyenne du taux net publié sur les N dernières années disponibles.
    Utilisé comme taux par défaut du fonds euros dans la simulation.
    """
    if history_df.empty:
        return 2.0
    recent = history_df.nlargest(years, "annee")
    return round(float(recent["taux_net_publie_pct"].mean()), 2)


def _compute_auto_euro_rate(
    history_df: pd.DataFrame,
    start: pd.Timestamp,
    end: pd.Timestamp,
) -> Optional[float]:
    """
    Calcule le taux moyen du fonds en euros sur les années couvertes par la fenêtre
    [start, end]. Retourne None si l'historique est insuffisant.
    """
    if history_df.empty:
        return None
    start_year = int(start.year)
    end_year = int(end.year)
    mask = (history_df["annee"] >= start_year) & (history_df["annee"] <= end_year)
    filtered = history_df[mask]
    if filtered.empty:
        return None
    return round(float(filtered["taux_net_publie_pct"].mean()), 2)


@st.cache_data(show_spinner=False, ttl=86400)
def _mstarpy_nav_series(isin: str) -> pd.DataFrame:
    """Fallback mstarpy désactivé — retourne toujours un DataFrame vide."""
    return pd.DataFrame()


def get_price_series(
    isin_or_name: str, start: Optional[pd.Timestamp], euro_rate: float
) -> Tuple[pd.DataFrame, str, str]:
    """
    Retourne la série de prix BRUTE (sans frais appliqués).
    • EUROFUND : start et euro_rate comptent dans la clé de cache (ils définissent la série).
    • Instruments EODHD : toujours appeler avec start=None pour maximiser le taux de cache ;
      le filtrage et l'application des frais sont délégués à get_price_series_with_fees().
    """
    # FIXED: fees and start-filtering for EODHD removed from cached fn (Correction 1 — cache explosion)
    debug = {"cands": []}
    val = str(isin_or_name).strip()
    if not val:
        return pd.DataFrame(), "", json.dumps(debug)

    # ✅ Fonds en euros — capitalisation annualisée (jours calendaires)
    if val.upper() == "EUROFUND":
        start_dt = (
            pd.Timestamp(start).normalize()
            if start is not None
            else pd.Timestamp("2000-01-03")
        )
        start_dt = max(start_dt, pd.Timestamp("2000-01-03"))

        idx = pd.bdate_range(start=start_dt, end=TODAY, freq="B")
        if len(idx) == 0:
            return pd.DataFrame(), "", "{}"

        # FIXED: vectorise with NumPy (Bug 8) — fees applied in get_price_series_with_fees()
        days = (idx - idx[0]).days.values.astype(float)
        r = float(euro_rate) / 100.0
        close_values = (1.0 + r) ** (days / 365.25)
        df = pd.DataFrame({"Close": close_values}, index=idx)
        return df, "EUROFUND", "{}"

    # Source prioritaire : VL importée manuellement (prend le dessus sur EODHD et mstarpy)
    # get_price_series() n'est pas décorée @st.cache_data — accès direct à st.session_state autorisé
    _manual_store = st.session_state.get("MANUAL_NAV_STORE", {})
    if val in _manual_store:
        return _manual_store[val], val, json.dumps({"source": "manual_upload"})

    # ✅ Instruments EODHD — historique complet, sans filtrage ni frais
    # Passer start=None depuis les appelants pour qu'un seul enregistrement de cache
    # serve toutes les lignes du même fonds quelle que soit leur date d'achat.
    cands = _symbol_candidates(val)
    debug["cands"] = cands

    for sym in cands:
        df = eodhd_prices_daily(sym)
        if not df.empty:
            return df, sym, json.dumps(debug)

    # Fallback mstarpy si EODHD n'a rien retourné
    # Uniquement pour les ISIN valides (12 caractères, 2 premières lettres alpha)
    _is_isin = len(val) == 12 and val[:2].isalpha()
    if _is_isin:
        _mstar_df = _mstarpy_nav_series(val)
        if not _mstar_df.empty:
            return _mstar_df, val, json.dumps({"source": "mstarpy_nav"})

    return pd.DataFrame(), "", json.dumps(debug)


def get_price_series_with_fees(
    isin_or_name: str,
    start: Optional[pd.Timestamp],
    euro_rate: float,
    annual_uc_fee_pct: float = ANNUAL_FEE_UC_PCT,  # frais UC de la ligne (contract-aware)
) -> Tuple[pd.DataFrame, str, str]:
    """
    Wrapper non-caché : appelle get_price_series() (start=None pour les instruments EODHD
    afin de maximiser le taux de cache) puis applique les frais annuels ancrés sur la
    date d'achat.  La fonction cachée est ainsi partagée entre toutes les lignes d'un
    même fonds indépendamment de leur date d'achat individuelle.

    ⚠️ USAGE CORRECT de cette fonction :
    - Toujours fournir start = date d'achat du client (jamais None)
    - Réservée aux calculs de simulation et de performance client
    - NE PAS utiliser pour afficher une VL de marché absolue
    - Pour la VL de marché brute : utiliser get_current_nav(isin)
    - Pour les calculs de corrélation/vol/Sharpe : utiliser
      get_series_for_line(..., apply_fees=False)
    """
    # FIXED: fee application and start-trimming moved outside the cache (Correction 1)
    val = str(isin_or_name).strip()

    if val.upper() == "EUROFUND":
        # EUROFUND : euro_rate est déjà net de tous frais de gestion assureur
        # → pas d'apply_annual_fee pour ne pas doubler-compter les frais
        df, sym, dbg = get_price_series(isin_or_name, start, euro_rate)
        return df, sym, dbg

    # Instrument EODHD : start=None → un seul enregistrement de cache par ISIN
    df, sym, dbg = get_price_series(isin_or_name, None, euro_rate)
    if df.empty:
        return df, sym, dbg
    if start is not None:
        df = df.loc[df.index >= pd.Timestamp(start)]
    if not df.empty:
        df = apply_annual_fee(df, annual_uc_fee_pct, buy_date=start)
    return df, sym, dbg


def get_current_nav(
    isin_or_name: str,
) -> Tuple[float, Optional[pd.Timestamp]]:
    """
    Retourne (dernière_VL_brute, date_de_cette_VL).
    Retourne (np.nan, None) si indisponible.
    Ne jamais appeler apply_annual_fee() ici.
    """
    val = str(isin_or_name).strip().upper()
    if val in ("EUROFUND", "STRUCTURED", ""):
        return np.nan, None
    df, _, _ = get_price_series(isin_or_name, None, 0.0)
    if df.empty:
        return np.nan, None
    last_date = df.index[-1]
    return float(df["Close"].iloc[-1]), pd.Timestamp(last_date)


@st.cache_data(show_spinner=False, ttl=3600)
def structured_series(
    start: pd.Timestamp,
    end: pd.Timestamp,
    annual_rate_pct: float,
    redemption_years: int,
) -> pd.DataFrame:
    """
    Série synthétique autocall (simplifiée) :
    - Prix d'achat = 1.0
    - Plat jusqu'à la date de remboursement estimée
    - Saut à 1 + (rate * years) le jour de remboursement, puis figé
    """
    start_dt = pd.Timestamp(start).normalize()
    end_dt = pd.Timestamp(end).normalize()
    idx = pd.bdate_range(start=start_dt, end=end_dt, freq="B")
    if len(idx) == 0:
        return pd.DataFrame()

    df = pd.DataFrame(index=idx, columns=["Close"], dtype=float)
    df.iloc[0, 0] = 1.0

    r = float(annual_rate_pct) / 100.0
    yrs = int(redemption_years)

    redemption_dt = start_dt + pd.DateOffset(years=yrs)

    # série plate + saut à partir du 1er jour ouvré >= redemption_dt
    redeemed = False
    for i in range(1, len(df)):
        d = df.index[i]
        df.iloc[i, 0] = df.iloc[i - 1, 0]

        if (not redeemed) and (d >= redemption_dt):
            df.iloc[i, 0] = 1.0 + r * yrs
            df.iloc[i:, 0] = df.iloc[i, 0]
            redeemed = True
            break

    # sécurité : propagation si besoin
    for i in range(1, len(df)):
        if pd.isna(df.iloc[i, 0]):
            df.iloc[i, 0] = df.iloc[i - 1, 0]

    return df


def _warn_once(key: str, msg: str) -> None:
    seen = st.session_state.get("WARN_ONCE")
    if not isinstance(seen, set):
        if isinstance(seen, list):
            seen = set(seen)
        else:
            seen = set()
    if key in seen:
        return
    seen.add(key)
    st.session_state["WARN_ONCE"] = seen
    st.warning(msg)


def _safe_struct_params(line: Dict[str, Any]) -> Tuple[float, int]:
    try:
        years = int(line.get("struct_years", 6))
    except Exception:
        years = 6
    years = max(1, years)
    try:
        rate = float(line.get("struct_rate", 8.0))
    except Exception:
        rate = 0.0
    rate = max(0.0, min(rate, 25.0))
    return rate, years


def get_series_for_line(
    line: Dict[str, Any],
    start: Optional[pd.Timestamp],
    euro_rate: float,
    apply_fees: bool = True,  # FIXED: route to fee-aware wrapper or raw cache fn (Correction 1)
) -> Tuple[pd.DataFrame, str]:
    isin_or_name = str(line.get("isin") or line.get("name") or "").strip()
    sym_upper = isin_or_name.upper()

    if sym_upper == "EUROFUND":
        if apply_fees:
            df, _, _ = get_price_series_with_fees("EUROFUND", start, euro_rate)
        else:
            df, _, _ = get_price_series("EUROFUND", start, euro_rate)
        if df.empty:
            _warn_once("series_empty:EUROFUND", "Série indisponible pour le fonds en euros (EUROFUND).")
        return df, "EUROFUND"

    if sym_upper == "STRUCTURED":
        buy_ts = pd.Timestamp(line.get("buy_date") or start or TODAY).normalize()
        rate, years = _safe_struct_params(line)
        df = structured_series(
            start=buy_ts,
            end=TODAY,
            annual_rate_pct=rate,
            redemption_years=years,
        )
        if start is not None:
            df = df.loc[df.index >= pd.Timestamp(start)]
        if df.empty:
            label = line.get("name") or "Produit structuré"
            _warn_once(
                f"series_empty:STRUCTURED:{label}",
                f"Série indisponible pour {label} (STRUCTURED) : ligne ignorée.",
            )
        return df, "STRUCTURED"

    # Appliquer uniquement fee_contract_pct (C) — la VL EODHD est déjà
    # nette de fee_uc_pct (B/TER). Utiliser fee_total_pct uniquement en
    # fallback si fee_contract_pct est absent (ligne sans référentiel contrat).
    _raw_fee_gsl = line.get("fee_contract_pct")
    if _raw_fee_gsl is None or _raw_fee_gsl == "":
        _raw_fee_gsl = line.get("fee_total_pct")
        if _raw_fee_gsl is None or _raw_fee_gsl == "":
            fee_uc = ANNUAL_FEE_UC_PCT
        else:
            try:
                fee_uc = float(_raw_fee_gsl)
            except (TypeError, ValueError):
                fee_uc = ANNUAL_FEE_UC_PCT
    else:
        try:
            fee_uc = float(_raw_fee_gsl)
        except (TypeError, ValueError):
            fee_uc = ANNUAL_FEE_UC_PCT
    # fee_uc == 0.0 explicite → pas de frais appliqués (ETF sans frais contrat)

    if apply_fees:
        df, sym, _ = get_price_series_with_fees(isin_or_name, start, euro_rate,
                                                 annual_uc_fee_pct=fee_uc)
    else:
        df, sym, _ = get_price_series(isin_or_name, None, euro_rate)
        if start is not None and not df.empty:
            df = df.loc[df.index >= pd.Timestamp(start)]
    if df.empty:
        label = line.get("name") or isin_or_name or "Ligne"
        _warn_once(
            f"series_empty:{isin_or_name}",
            f"Série indisponible pour {label} ({isin_or_name}) : ligne ignorée.",
        )
    return df, sym or isin_or_name

# ------------------------------------------------------------
# Alternatives si date < 1ère VL
# ------------------------------------------------------------

def suggest_alternative_funds(buy_date: pd.Timestamp, euro_rate: float) -> List[Tuple[str, str, pd.Timestamp]]:
    """
    Propose des fonds recommandés (core + défensifs) dont la première VL
    est antérieure ou égale à la date d'achat donnée.
    Retourne (nom, isin, date_inception).
    """
    alternatives: List[Tuple[str, str, pd.Timestamp]] = []
    universe = RECO_FUNDS_CORE + RECO_FUNDS_DEF

    for name, isin in universe:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty:
            continue
        inception = df.index.min()
        if inception <= buy_date:
            alternatives.append((name, isin, inception))

    return alternatives


# ------------------------------------------------------------
# Calendrier versements & poids
# ------------------------------------------------------------

def _month_schedule(d0: pd.Timestamp, d1: pd.Timestamp) -> List[pd.Timestamp]:
    if d0 > d1:
        return []
    out = []
    cur = pd.Timestamp(d0.year, d0.month, 1)
    stop = pd.Timestamp(d1.year, d1.month, 1)
    while cur <= stop:
        bdays = pd.bdate_range(start=cur, end=cur + pd.offsets.MonthEnd(0))
        if len(bdays) > 0:
            out.append(bdays[0])
        cur = cur + pd.offsets.MonthBegin(1)
    return out


def _weights_for(
    lines: List[Dict[str, Any]],
    alloc_mode: str,
    custom_weights: Dict[Any, float],
    single_target: Optional[Any],
) -> Dict[Any, float]:
    # FIXED: use stable UUID key (ln.get("id")) with id(ln) fallback for old lines (Bug 5)
    keys = [ln.get("id") or id(ln) for ln in lines]
    if len(keys) == 0:
        return {}
    if alloc_mode == "equal":
        w = 1.0 / len(keys)
        return {k: w for k in keys}
    if alloc_mode == "custom":
        tot = sum(max(0.0, float(custom_weights.get(ln.get("id") or id(ln), 0.0))) for ln in lines)
        if tot <= 0:
            w = 1.0 / len(keys)
            return {k: w for k in keys}
        return {ln.get("id") or id(ln): max(0.0, float(custom_weights.get(ln.get("id") or id(ln), 0.0))) / tot for ln in lines}
    if alloc_mode == "single":
        target = single_target
        return {ln.get("id") or id(ln): (1.0 if (ln.get("id") or id(ln)) == target else 0.0) for ln in lines}
    w = 1.0 / len(keys)
    return {k: w for k in keys}


# ------------------------------------------------------------
# Calcul par ligne (avec frais)
# ------------------------------------------------------------

def compute_line_metrics(
    line: Dict[str, Any], fee_pct: float, euro_rate: float
) -> Tuple[float, float, float]:
    amt_brut = float(line.get("amount_gross", 0.0))
    net_amt = amt_brut * (1.0 - fee_pct / 100.0)
    buy_ts = pd.Timestamp(line.get("buy_date"))
    px_manual = line.get("buy_px", None)

    dfp, sym_used = get_series_for_line(line, buy_ts, euro_rate)
    if dfp.empty:
        return float(net_amt), np.nan, 0.0

    sym_upper = str(sym_used or line.get("isin") or "").upper()
    if sym_upper == "EUROFUND":
        px = _get_close_on(dfp, buy_ts)
    else:
        if px_manual not in (None, "", 0, "0"):
            try:
                px = float(px_manual)
            except Exception:
                px = _get_close_on(dfp, buy_ts)
        else:
            px = _get_close_on(dfp, buy_ts)

    qty = (net_amt / px) if px and px > 0 else 0.0
    return float(net_amt), float(px), float(qty)


# ------------------------------------------------------------
# Simulation d'un portefeuille (avec contrôle 1ère VL)
# + distinction poids mensuels / ponctuels
# ------------------------------------------------------------

def simulate_portfolio(
    lines: List[Dict[str, Any]],
    monthly_amt_gross: float,
    one_amt_gross: float,
    one_date: date,
    alloc_mode: str,
    custom_weights_monthly: Optional[Dict[int, float]],
    custom_weights_oneoff: Optional[Dict[int, float]],
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    portfolio_label: str = "",
) -> Tuple[pd.DataFrame, float, float, float, Optional[float], pd.Timestamp, pd.Timestamp]:
    if not lines:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    price_map: Dict[Any, pd.Series] = {}
    eff_buy_date: Dict[Any, pd.Timestamp] = {}
    buy_price_used: Dict[Any, float] = {}

    invalid_found = False
    date_warnings = st.session_state.setdefault("DATE_WARNINGS", [])

    for ln in lines:
        key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key, avoids rerun ID changes (Bug 5)
        # FIXED (Bug B): lire la série BRUTE (sans frais) pour déterminer l'inception
        # et vérifier la date d'achat ; les frais seront appliqués ci-dessous
        # ancrés sur d_buy et non sur l'inception du fonds.
        df_full, sym = get_series_for_line(ln, None, euro_rate, apply_fees=False)

        # Sécurité
        if df_full.empty:
            continue

        inception = df_full.index.min()
        d_buy = pd.Timestamp(ln["buy_date"])

        if d_buy < inception:
            invalid_found = True
            ln["invalid_date"] = True
            ln["inception_date"] = inception

            alts = suggest_alternative_funds(d_buy, euro_rate)
            if alts:
                alt_lines = [
                    f"- {name} ({isin}), historique depuis le {fmt_date(incep)}"
                    for name, isin, incep in alts
                ]
                alt_msg = "\n".join(alt_lines)
            else:
                alt_msg = "Aucun fonds recommandé ne dispose d'un historique suffisant pour cette date."

            date_warnings.append(
                f"[{portfolio_label}] {ln.get('name','(sans nom)')} "
                f"({ln.get('isin','—')}) :\n"
                f"- Date d'achat saisie : {fmt_date(d_buy)}\n"
                f"- 1ère VL disponible : {fmt_date(inception)}\n\n"
                f"Impossible de simuler ce fonds sur toute la période demandée.\n"
                f"Propositions d'alternatives pour l'analyse historique :\n{alt_msg}"
            )
            continue

        # FIXED (Bug B): appliquer les frais ancrés sur d_buy (date client réelle)
        # FIXED (Bug E): distinguer fee_total_pct absent (fallback) de 0.0 explicite
        # Appliquer uniquement fee_contract_pct (C) — VL EODHD déjà nette du TER.
        # Fallback fee_total_pct si fee_contract_pct absent, puis ANNUAL_FEE_UC_PCT.
        _isin_sim = str(ln.get("isin", "")).upper()
        _raw_fee_sim = ln.get("fee_contract_pct")
        if _raw_fee_sim is None or _raw_fee_sim == "":
            _raw_fee_sim = ln.get("fee_total_pct")
            if _raw_fee_sim is None or _raw_fee_sim == "":
                _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
            else:
                try:
                    _fee_sim = float(_raw_fee_sim)
                except (TypeError, ValueError):
                    _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
        else:
            try:
                _fee_sim = float(_raw_fee_sim)
            except (TypeError, ValueError):
                _fee_sim = ANNUAL_FEE_UC_PCT if _isin_sim not in ("EUROFUND", "STRUCTURED") else 0.0
        # EUROFUND et STRUCTURED restent à 0.0 (inchangé)
        if _fee_sim > 0 and _isin_sim not in ("EUROFUND", "STRUCTURED"):
            df = apply_annual_fee(df_full.copy(), _fee_sim, buy_date=d_buy)
        else:
            df = df_full

        ln["sym_used"] = sym

        if d_buy in df.index:
            px_buy = float(df.loc[d_buy, "Close"])
            eff_dt = d_buy
        else:
            after = df.loc[df.index >= d_buy]
            if after.empty:
                px_buy = float(df.iloc[-1]["Close"])
                eff_dt = df.index[-1]
            else:
                px_buy = float(after.iloc[0]["Close"])
                eff_dt = after.index[0]

        px_manual = ln.get("buy_px", None)
        px_for_qty = float(px_manual) if (px_manual not in (None, "", 0, "0")) else px_buy

        price_map[key_id] = df["Close"].astype(float)
        eff_buy_date[key_id] = eff_dt
        buy_price_used[key_id] = px_for_qty

    if invalid_found and not price_map:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY
    if not price_map:
        return pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    start_min = min(eff_buy_date.values())
    start_full = max(eff_buy_date.values())

    bidx = pd.bdate_range(start=start_min, end=TODAY, freq="B")
    prices = pd.DataFrame(index=bidx)
    for key_id, s in price_map.items():
        prices[key_id] = s.reindex(bidx).ffill()

    qty_events = pd.DataFrame(0.0, index=bidx, columns=prices.columns)
    total_brut = 0.0
    total_net = 0.0
    cash_flows: List[Tuple[pd.Timestamp, float]] = []

    # Achats initiaux
    for ln in lines:
        key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
        if key_id not in prices.columns:
            continue
        brut = float(ln.get("amount_gross", 0.0))
        net = brut * (1.0 - fee_pct / 100.0)
        px = float(buy_price_used[key_id])
        dt = eff_buy_date[key_id]
        if brut > 0 and px > 0:
            q = net / px
            tgt = dt if dt in qty_events.index else qty_events.index[qty_events.index >= dt][0]
            qty_events.loc[tgt, key_id] += q
            total_brut += brut
            total_net += net
            cash_flows.append((tgt, -brut))

    # Poids pour versements mensuels / ponctuels
    weights_monthly = _weights_for(
        lines,
        alloc_mode,
        custom_weights_monthly or {},
        single_target,
    )
    weights_oneoff = _weights_for(
        lines,
        alloc_mode,
        custom_weights_oneoff or {},
        single_target,
    )

    # Versement ponctuel
    if one_amt_gross > 0:
        dt = pd.Timestamp(one_date)
        if dt not in qty_events.index:
            after = qty_events.index[qty_events.index >= dt]
            if len(after) > 0:
                dt = after[0]
            else:
                dt = None
        if dt is not None:
            net_amt = one_amt_gross * (1.0 - fee_pct / 100.0)
            for ln in lines:
                key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
                w = weights_oneoff.get(key_id, 0.0)
                if w <= 0 or key_id not in prices.columns:
                    continue
                px = float(prices.loc[dt, key_id])
                if px > 0:
                    qty_events.loc[dt, key_id] += (net_amt * w) / px
            total_brut += float(one_amt_gross)
            total_net += float(net_amt)
            cash_flows.append((dt, -float(one_amt_gross)))

    # Mensuels
    if monthly_amt_gross > 0:
        sched = _month_schedule(start_min, TODAY)
        for dt in sched:
            if dt not in qty_events.index:
                after = qty_events.index[qty_events.index >= dt]
                if len(after) == 0:
                    continue
                dt = after[0]
            net_m = monthly_amt_gross * (1.0 - fee_pct / 100.0)
            for ln in lines:
                key_id = ln.get("id") or id(ln)  # FIXED: stable UUID key (Bug 5)
                w = weights_monthly.get(key_id, 0.0)
                if w <= 0 or key_id not in prices.columns:
                    continue
                px = float(prices.loc[dt, key_id])
                if px > 0:
                    qty_events.loc[dt, key_id] += (net_m * w) / px
            total_brut += float(monthly_amt_gross)
            total_net += float(net_m)
            cash_flows.append((dt, -float(monthly_amt_gross)))

    qty_cum = qty_events.cumsum()
    values = (qty_cum * prices).sum(axis=1)
    df_val = pd.DataFrame({"Valeur": values})
    final_val = float(df_val["Valeur"].iloc[-1]) if not df_val.empty else 0.0

    cash_flows.append((TODAY, final_val))
    irr = xirr(cash_flows)

    return df_val, total_brut, total_net, final_val, (irr * 100.0 if irr is not None else None), start_min, start_full


@st.cache_data(show_spinner=False, ttl=3600)
def _simulate_portfolio_cached(
    lines_json: str,
    monthly_amt_gross: float,
    one_amt_gross: float,
    one_date_str: str,
    alloc_mode: str,
    custom_weights_monthly_json: str,
    custom_weights_oneoff_json: str,
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    portfolio_label: str = "",
) -> Tuple[pd.DataFrame, float, float, float, Optional[float], pd.Timestamp, pd.Timestamp]:
    """Wrapper caché de simulate_portfolio. Tous les paramètres sont hashables.
    La clé de cache est le hash de la sérialisation JSON des lignes + scalaires.
    Si le portefeuille n'a pas changé entre deux reruns, résultat instantané."""
    lines = json.loads(lines_json)
    custom_weights_monthly = json.loads(custom_weights_monthly_json)
    custom_weights_oneoff = json.loads(custom_weights_oneoff_json)
    one_date = pd.Timestamp(one_date_str).date()
    return simulate_portfolio(
        lines=lines,
        monthly_amt_gross=monthly_amt_gross,
        one_amt_gross=one_amt_gross,
        one_date=one_date,
        alloc_mode=alloc_mode,
        custom_weights_monthly=custom_weights_monthly or None,
        custom_weights_oneoff=custom_weights_oneoff or None,
        single_target=single_target,
        euro_rate=euro_rate,
        fee_pct=fee_pct,
        portfolio_label=portfolio_label,
    )


def _make_sim_args(
    lines: List[Dict[str, Any]],
    monthly_amt: float,
    one_amt: float,
    one_date,
    alloc_mode: str,
    custom_weights_monthly: Optional[Dict],
    custom_weights_oneoff: Optional[Dict],
    single_target: Optional[int],
    euro_rate: float,
    fee_pct: float,
    label: str,
) -> dict:
    """Sérialise les arguments pour _simulate_portfolio_cached.
    Convertit les clés en str pour la sérialisation JSON."""
    def _str_keys(d: Optional[Dict]) -> Dict[str, float]:
        if not d:
            return {}
        return {str(k): v for k, v in d.items()}
    return dict(
        lines_json=json.dumps(lines, default=str, sort_keys=True),
        monthly_amt_gross=float(monthly_amt),
        one_amt_gross=float(one_amt),
        one_date_str=str(pd.Timestamp(one_date).date()),
        alloc_mode=str(alloc_mode),
        custom_weights_monthly_json=json.dumps(_str_keys(custom_weights_monthly)),
        custom_weights_oneoff_json=json.dumps(_str_keys(custom_weights_oneoff)),
        single_target=single_target,
        euro_rate=float(euro_rate),
        fee_pct=float(fee_pct),
        portfolio_label=str(label),
    )


# ------------------------------------------------------------
# Cartes lignes (édition / suppression)
# ------------------------------------------------------------

def _line_card(line: Dict[str, Any], idx: int, port_key: str):
    # FIXED: use stable UUID-based key so widget state survives st.rerun() (Bug 5)
    _card_id = line.get("id", str(idx))
    state_key = f"edit_mode_{port_key}_{_card_id}"
    if state_key not in st.session_state:
        st.session_state[state_key] = False

    fee_pct = st.session_state.get("FEE_A", 0.0) if port_key == "A_lines" else st.session_state.get("FEE_B", 0.0)
    # FIXED: use per-portfolio euro rate instead of shared EURO_RATE_PREVIEW (Résidu Bug 4)
    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )
    net_amt, buy_px, qty_disp = compute_line_metrics(line, fee_pct, euro_rate)

    with st.container(border=True):
        cols = st.columns([3, 2, 2, 2, 1])
        with cols[0]:
            st.markdown(f"**{line.get('name','—')}**")
            st.caption(f"ISIN / Code : `{line.get('isin','—')}`")
            st.caption(f"Symbole EODHD : `{line.get('sym_used','—')}`")
            if line.get("invalid_date"):
                st.markdown(
                    f"⚠️ Date d'achat antérieure à la 1ère VL ({fmt_date(line.get('inception_date'))}).",
                )
        with cols[1]:
            st.markdown(f"Investi (brut)\n\n**{to_eur(line.get('amount_gross', 0.0))}**")
            st.caption(f"Net après frais {fee_pct:.1f}% : **{to_eur(net_amt)}**")
            st.caption(f"Date d'achat : {fmt_date(line.get('buy_date'))}")
            if line.get("date_overridden"):
                st.caption("📌 Date individuelle (non liée à la date globale)")
            _buy_ts_check = pd.Timestamp(line.get("buy_date", TODAY))
            if _buy_ts_check > pd.Timestamp(TODAY):
                st.caption("⚠️ Date d'achat dans le futur — simulation non disponible.")
        with cols[2]:
            st.markdown(f"VL d'achat\n\n**{to_eur(buy_px)}**")
            st.caption(f"Quantité : {qty_disp:.6f}")
            if line.get("note"):
                st.caption(line["note"])
        with cols[3]:
            try:
                # FIXED (Bug A + P2): VL brute sans frais + date de cotation
                last, nav_date = get_current_nav(
                    line.get("isin") or line.get("name") or ""
                )
                if last == last:  # not nan
                    date_str = fmt_date(nav_date) if nav_date is not None else "—"
                    st.markdown(f"Dernière VL connue ({date_str}) : **{to_eur(last)}**")
                    st.caption(
                        "VL publiée par EODHD (clôture J-1 à J-3 selon le fonds)"
                    )
                else:
                    st.markdown("Dernière VL connue : —")
            except Exception:
                st.markdown("Dernière VL connue : —")
        with cols[4]:
            if not st.session_state[state_key]:
                # FIXED: use stable UUID-based widget key (Bug 5)
                if st.button("✏️", key=f"edit_{port_key}_{_card_id}", help="Modifier"):
                    st.session_state[state_key] = True
                    st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)
            # FIXED: use stable UUID-based widget key (Bug 5)
            if st.button("🗑️", key=f"del_{port_key}_{_card_id}", help="Supprimer"):
                st.session_state[port_key].pop(idx)
                st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)

        if st.session_state[state_key]:
            # FIXED: use stable UUID-based form key (Bug 5)
            with st.form(key=f"form_edit_{port_key}_{_card_id}", clear_on_submit=False):
                c1, c2, c3, c4 = st.columns([2, 2, 2, 1])
                with c1:
                    new_amount = st.text_input("Montant investi (brut) €", value=str(line.get("amount_gross", "")))
                with c2:
                    new_date = st.date_input("Date d'achat", value=pd.Timestamp(line.get("buy_date")).date())
                with c3:
                    new_px = st.text_input("Prix d'achat (optionnel)", value=str(line.get("buy_px", "")))
                with c4:
                    st.caption(" ")
                    submitted = st.form_submit_button("💾 Enregistrer")
                if submitted:
                    try:
                        amt_gross = float(str(new_amount).replace(" ", "").replace(",", "."))
                        assert amt_gross > 0
                    except Exception:
                        st.warning("Montant brut invalide.")
                        st.stop()
                    buy_ts = pd.Timestamp(new_date)
                    line["amount_gross"] = float(amt_gross)
                    line["buy_date"] = buy_ts
                    line["date_overridden"] = True
                    if new_px.strip():
                        try:
                            line["buy_px"] = float(str(new_px).replace(",", "."))
                        except Exception:
                            line["buy_px"] = ""
                    else:
                        line["buy_px"] = ""
                    line.pop("invalid_date", None)
                    line.pop("inception_date", None)
                    st.session_state[state_key] = False
                    st.success("Ligne mise à jour.")
                    st.rerun()  # FIXED: st.experimental_rerun() deprecated since 1.27 (Bug 6)

    # Import VL manuelle si la série est vide (EODHD + mstarpy ont échoué)
    _isin_card = str(line.get("isin") or "").strip()
    _has_data = not get_series_for_line(
        line, pd.Timestamp(line.get("buy_date", TODAY)), euro_rate, apply_fees=False
    )[0].empty
    _has_manual = _isin_card in st.session_state.get("MANUAL_NAV_STORE", {})
    if not _has_data and not _has_manual and _isin_card not in ("EUROFUND", "STRUCTURED", ""):
        with st.expander(
            f"📥 Importer l'historique VL — {line.get('name', _isin_card)}", expanded=False
        ):
            st.caption(
                "Ce fonds n'est pas disponible via EODHD ni Morningstar. "
                "Importez son historique de VL au format CSV (colonnes : date ; vl). "
                "Sources recommandées : OPCVM360, Fundkis, Morningstar Direct."
            )
            _uploaded_nav = st.file_uploader(
                "CSV historique VL (séparateur ; décimale ,)",
                type=["csv"],
                key=f"nav_upload_{_isin_card}_{_card_id}",
            )
            if _uploaded_nav is not None:
                try:
                    _nav_raw = pd.read_csv(
                        _uploaded_nav, sep=";", decimal=",",
                        encoding="utf-8-sig", dtype=str,
                    )
                    _nav_raw.columns = _nav_raw.columns.str.strip().str.lower()
                    _vl_col = next(
                        (c for c in _nav_raw.columns
                         if c in ("vl", "nav", "close", "valeur liquidative", "price")),
                        None,
                    )
                    if _vl_col and "date" in _nav_raw.columns:
                        _nav_raw["date"] = pd.to_datetime(
                            _nav_raw["date"], dayfirst=True, errors="coerce"
                        )
                        _nav_raw[_vl_col] = pd.to_numeric(
                            _nav_raw[_vl_col].str.replace(",", "."), errors="coerce"
                        )
                        _nav_raw = _nav_raw.dropna(subset=["date", _vl_col])
                        _nav_df = (
                            _nav_raw.set_index("date")
                            .rename(columns={_vl_col: "Close"})[["Close"]]
                            .sort_index()
                        )
                        if not _nav_df.empty:
                            _store = st.session_state.get("MANUAL_NAV_STORE", {})
                            _store[_isin_card] = _nav_df
                            st.session_state["MANUAL_NAV_STORE"] = _store
                            st.success(
                                f"✅ {len(_nav_df)} VL importées pour "
                                f"{line.get('name', _isin_card)}. Relancez le calcul."
                            )
                        else:
                            st.error("Fichier valide mais aucune donnée exploitable.")
                    else:
                        st.error(
                            "Colonnes attendues : 'date' et 'vl' (ou 'nav' / 'close'). "
                            f"Colonnes détectées : {list(_nav_raw.columns)}"
                        )
                except Exception as _e:
                    st.error(f"Erreur lecture CSV : {_e}")


def build_positions_dataframe(port_key: str) -> pd.DataFrame:
    """
    Construit un DataFrame par ligne :
    Nom, ISIN, Date d'achat, Net investi, Valeur actuelle, Perf € et Perf %.
    """
    fee_pct = (
        st.session_state.get("FEE_A", 0.0)
        if port_key == "A_lines"
        else st.session_state.get("FEE_B", 0.0)
    )

    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )

    lines = st.session_state.get(port_key, [])
    rows: List[Dict[str, Any]] = []

    for ln in lines:
        buy_ts = pd.Timestamp(ln.get("buy_date"))
        net_amt, buy_px, qty = compute_line_metrics(ln, fee_pct, euro_rate)
        dfl, _ = get_series_for_line(ln, buy_ts, euro_rate)
        if dfl.empty:
            continue

        # FIXED: removed dead `if not dfl.empty` block — unreachable after `continue` above (Bug 9)
        last_px = float(dfl["Close"].iloc[-1])

        val_now = qty * last_px if last_px == last_px else 0.0
        perf_abs = val_now - net_amt
        perf_pct = (val_now / net_amt - 1.0) * 100.0 if net_amt > 0 else np.nan

        # FIXED (Bug D + P2): VL marché brute séparée de val_now ; date ignorée ici
        nav_brute, _ = get_current_nav(ln.get("isin") or "")

        rows.append(
            {
                "Nom": ln.get("name", ""),
                "ISIN / Code": ln.get("isin", ""),
                "Date d'achat": fmt_date(ln.get("buy_date")),
                "Date individuelle": "📌" if ln.get("date_overridden") else "",
                "Net investi €": net_amt,
                "Valeur actuelle €": val_now,
                "VL marché": round(nav_brute, 4) if nav_brute == nav_brute else None,
                "Perf €": perf_abs,
                "Perf %": perf_pct,
            }
        )

    return pd.DataFrame(rows)

# ------------------------------------------------------------
# Tableau synthétique par ligne (un seul tableau par portefeuille)
# ------------------------------------------------------------

def positions_table(title: str, port_key: str):
    """
    Affiche un seul tableau synthétique par portefeuille :
    Nom, ISIN, Date d'achat, Net investi, Valeur actuelle, Perf € et Perf %.
    """
    fee_pct = (
        st.session_state.get("FEE_A", 0.0)
        if port_key == "A_lines"
        else st.session_state.get("FEE_B", 0.0)
    )

    # ✅ Taux fonds euros par portefeuille (au lieu de EURO_RATE_PREVIEW)
    euro_rate = (
        st.session_state.get("EURO_RATE_A", 2.0)
        if port_key == "A_lines"
        else st.session_state.get("EURO_RATE_B", 2.5)
    )

    lines = st.session_state.get(port_key, [])
    rows: List[Dict[str, Any]] = []

    for ln in lines:
        buy_ts = pd.Timestamp(ln.get("buy_date"))

        # Montant net investi, VL d'achat et quantité
        net_amt, buy_px, qty = compute_line_metrics(ln, fee_pct, euro_rate)

        # ✅ IMPORTANT : on récupère la série "depuis buy_ts" pour éviter le mismatch EUROFUND
        dfl, _ = get_series_for_line(ln, buy_ts, euro_rate)
        if dfl.empty:
            continue

        # FIXED: removed dead `if not dfl.empty` block — unreachable after `continue` above (Bug 9)
        last_px = float(dfl["Close"].iloc[-1])

        # Valeur actuelle et performance
        val_now = qty * last_px if last_px == last_px else 0.0
        perf_abs = val_now - net_amt
        perf_pct = (val_now / net_amt - 1.0) * 100.0 if net_amt > 0 else np.nan

        # FIXED (Bug D + P2): VL marché brute séparée de val_now ; date ignorée ici
        nav_brute, _ = get_current_nav(ln.get("isin") or "")

        rows.append(
            {
                "Nom": ln.get("name", ""),
                "ISIN / Code": ln.get("isin", ""),
                "Date d'achat": fmt_date(ln.get("buy_date")),
                "Date individuelle": "📌" if ln.get("date_overridden") else "",
                "Net investi €": net_amt,
                "Valeur actuelle €": val_now,
                "VL marché": round(nav_brute, 4) if nav_brute == nav_brute else None,
                "Perf €": perf_abs,
                "Perf %": perf_pct,
            }
        )

    st.markdown(f"### {title}")
    df = pd.DataFrame(rows)
    if df.empty:
        st.info("Aucune ligne.")
    else:
        st.dataframe(
            df.style.format(
                {
                    "Net investi €": to_eur,
                    "Valeur actuelle €": to_eur,
                    "VL marché": "{:.4f}".format,
                    "Perf €": to_eur,
                    "Perf %": "{:,.2f}%".format,
                }
            ),
            hide_index=True,
            use_container_width=True,
        )
        st.caption("VL marché : dernière cotation brute EODHD (hors frais de gestion)")


def _prepare_pie_df(df_positions: pd.DataFrame, max_items: int = 8, min_pct: float = 0.03) -> pd.DataFrame:
    if df_positions.empty:
        return df_positions
    df = df_positions.copy()
    df = df[df["Valeur actuelle €"] > 0]
    if df.empty:
        return df
    total = df["Valeur actuelle €"].sum()
    df["Part %"] = df["Valeur actuelle €"] / total
    df = df.sort_values("Valeur actuelle €", ascending=False)
    if len(df) > max_items:
        df_main = df.iloc[:max_items].copy()
        df_other = df.iloc[max_items:]
        df_main = pd.concat(
            [
                df_main,
                pd.DataFrame(
                    {
                        "Nom": ["Autres"],
                        "Valeur actuelle €": [df_other["Valeur actuelle €"].sum()],
                        "Part %": [df_other["Valeur actuelle €"].sum() / total],
                    }
                ),
            ],
            ignore_index=True,
        )
        df = df_main
    else:
        small = df[df["Part %"] < min_pct]
        if not small.empty and len(df) > 1:
            df_main = df[df["Part %"] >= min_pct]
            df_other = pd.DataFrame(
                {
                    "Nom": ["Autres"],
                    "Valeur actuelle €": [small["Valeur actuelle €"].sum()],
                    "Part %": [small["Valeur actuelle €"].sum() / total],
                }
            )
            df = pd.concat([df_main, df_other], ignore_index=True)
    df["Part %"] = df["Part %"] * 100.0
    return df


# ------------------------------------------------------------
# Analytics internes : retours, corrélation, volatilité
# ------------------------------------------------------------


def _build_returns_df(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    years: int = 3,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Construit un DataFrame de rendements journaliers (pct_change)
    pour toutes les lignes du portefeuille avec un historique suffisant.
    Index = dates, colonnes = "Nom (ISIN)".
    """
    cutoff = TODAY - pd.Timedelta(days=365 * years)
    series_map: Dict[str, pd.Series] = {}

    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()
        key = f"{label} ({isin})" if isin else label

        # FIXED (Bug F): apply_fees=False — les frais biaisent les rendements si
        # ancrés sur l'inception du fonds ; pct_change() ne nécessite pas de frais nets
        df, _ = get_series_for_line(ln, None, euro_rate, apply_fees=False)
        if df.empty:
            continue

        s = df["Close"].astype(float)
        s = s[s.index >= cutoff]
        if s.size < min_points:
            continue

        series_map[key] = s

    if not series_map:
        return pd.DataFrame()

    df_prices = pd.DataFrame(series_map).dropna(how="any")
    if df_prices.shape[0] < min_points:
        return pd.DataFrame()

    returns = df_prices.pct_change().dropna(how="any")
    return returns



def correlation_matrix_from_lines(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    years: int = 3,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Matrice de corrélation entre les lignes du portefeuille,
    basée sur les rendements journaliers.
    """
    rets = _build_returns_df(lines, euro_rate, years=years, min_points=min_points)
    if rets.empty:
        return pd.DataFrame()
    return rets.corr()


def volatility_table_from_lines(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    years: int = 3,
    min_points: int = 60,
) -> pd.DataFrame:
    """
    Volatilité annuelle par ligne (et écart-type quotidien).
    """
    rets = _build_returns_df(lines, euro_rate, years=years, min_points=min_points)
    if rets.empty:
        return pd.DataFrame()

    rows = []
    for col in rets.columns:
        r = rets[col]
        daily_std = float(r.std())
        ann_std = daily_std * np.sqrt(252.0)
        rows.append(
            {
                "Ligne": col,
                "Écart-type quotidien %": daily_std * 100.0,
                "Volatilité annuelle %": ann_std * 100.0,
                "Nombre de points": int(r.count()),
            }
        )
    df = pd.DataFrame(rows)
    return df.sort_values("Volatilité annuelle %", ascending=False)


def portfolio_risk_stats(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    years: int = 3,
    min_points: int = 60,
    fee_pct: float = 0.0,  # FIXED: explicit fee rate eliminates max(net_A, net_B) proxy (Bug 3)
) -> Optional[Dict[str, float]]:
    """
    Calcule quelques stats globales de risque pour le portefeuille :
    - volatilité annuelle
    - max drawdown sur la période.
    Pondération par montant net investi.
    """
    rets = _build_returns_df(lines, euro_rate, years=years, min_points=min_points)
    if rets.empty:
        return None

    # FIXED: use caller-supplied fee_pct instead of guessing A vs B from session state (Bug 3)
    net_by_col: Dict[str, float] = {}
    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()
        key = f"{label} ({isin})" if isin else label

        net, _, _ = compute_line_metrics(ln, fee_pct, euro_rate)
        if net > 0:
            net_by_col[key] = net

    # normalisation des poids
    weights = {}
    tot = sum(net_by_col.get(c, 0.0) for c in rets.columns)
    if tot <= 0:
        return None
    for c in rets.columns:
        w = net_by_col.get(c, 0.0) / tot
        weights[c] = w

    # série de rendement du portefeuille
    w_vec = np.array([weights[c] for c in rets.columns])
    rp = rets.to_numpy().dot(w_vec)
    rp = pd.Series(rp, index=rets.index)

    daily_std = float(rp.std())
    vol_ann = daily_std * np.sqrt(252.0)

    # max drawdown
    cum = (1.0 + rp).cumprod()
    running_max = cum.cummax()
    dd = cum / running_max - 1.0
    max_dd = float(dd.min())  # négatif

    return {
        "vol_ann_pct": vol_ann * 100.0,
        "max_dd_pct": max_dd * 100.0,
    }


def compute_diversification_score(
    lines: List[Dict[str, Any]],
    euro_rate: float,
) -> Optional[Dict[str, Any]]:
    """
    Score de diversification 0–100 basé sur la corrélation moyenne hors-diagonale.
    0 = tout corrélé (fausse diversification), 100 = parfaitement décorrélé.
    Retourne aussi les paires fortement corrélées (doublons) et en vigilance.
    """
    corr = correlation_matrix_from_lines(lines, euro_rate)
    if corr.empty or corr.shape[0] < 2:
        return None

    avg_corr = _avg_offdiag_corr(corr)
    score = max(0.0, min(100.0, (1.0 - avg_corr) * 100.0))

    doublons: List[tuple] = []   # > 0.90 : quasi-identiques
    vigilance: List[tuple] = []  # > 0.80 : diversification limitée
    cols = list(corr.columns)
    for i in range(len(cols)):
        for j in range(i + 1, len(cols)):
            c = float(corr.iloc[i, j])
            name_i = cols[i].split(" (")[0][:30]
            name_j = cols[j].split(" (")[0][:30]
            if c > 0.90:
                doublons.append((name_i, name_j, c))
            elif c > 0.80:
                vigilance.append((name_i, name_j, c))

    # Nombre de fonds réellement utiles (diversification effective)
    # Un fonds est "utile" s'il n'est pas corrélé >0.80 avec un fonds déjà retenu
    effective_funds: List[str] = []
    for col_i in cols:
        is_redundant = False
        for col_e in effective_funds:
            if float(corr.loc[col_i, col_e]) > 0.80:
                is_redundant = True
                break
        if not is_redundant:
            effective_funds.append(col_i)
    n_effective = len(effective_funds)

    return {
        "score": round(score, 0),
        "avg_corr": avg_corr,
        "doublons": doublons,
        "vigilance": vigilance,
        "n_lines": corr.shape[0],
        "n_effective": n_effective,
    }


def _portfolio_weighted_returns(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    min_points: int = 60,
) -> Optional[pd.Series]:
    """Série de rendements quotidiens pondérés du portefeuille. Réutilisé par Sharpe/Sortino/Beta."""
    rets = _build_returns_df(lines, euro_rate, min_points=min_points)
    if rets.empty:
        return None
    net_by_col: Dict[str, float] = {}
    for ln in lines:
        label = (ln.get("name") or ln.get("isin") or "Ligne").strip()
        isin = (ln.get("isin") or "").strip()
        key = f"{label} ({isin})" if isin else label
        net, _, _ = compute_line_metrics(ln, fee_pct, euro_rate)
        if net > 0:
            net_by_col[key] = net
    tot = sum(net_by_col.get(c, 0.0) for c in rets.columns)
    if tot <= 0:
        return None
    w_vec = np.array([net_by_col.get(c, 0.0) / tot for c in rets.columns])
    rp = pd.Series(rets.to_numpy().dot(w_vec), index=rets.index)
    return rp if len(rp) >= min_points else None


def compute_sharpe_ratio(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
) -> Optional[float]:
    """
    Ratio de Sharpe annualisé.
    Taux sans risque = euro_rate / 100 si non fourni (alternative naturelle en AV).
    Retourne None si données insuffisantes (< 60 jours) ou volatilité nulle.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points)
    if rp is None:
        return None
    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann = rp.mean() * 252.0
    vol_ann = rp.std() * np.sqrt(252.0)
    if vol_ann == 0:
        return None
    return float((ret_ann - rf) / vol_ann)


def compute_sortino_ratio(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
) -> Optional[float]:
    """
    Ratio de Sortino annualisé — pénalise uniquement la volatilité à la baisse.
    Retourne None si données insuffisantes ou aucun rendement négatif.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points)
    if rp is None:
        return None
    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann = rp.mean() * 252.0
    downside = rp[rp < 0]
    if len(downside) == 0:
        return None
    downside_dev = downside.std() * np.sqrt(252.0)
    if downside_dev == 0:
        return None
    return float((ret_ann - rf) / downside_dev)


def compute_beta_alpha(
    lines: List[Dict[str, Any]],
    euro_rate: float,
    fee_pct: float = 0.0,
    benchmark_symbol: str = "CW8.PA",
    risk_free_rate: Optional[float] = None,
    min_points: int = 60,
) -> Optional[Dict[str, Any]]:
    """
    Bêta et Alpha de Jensen par rapport à l'indice de référence.
    Retourne {"beta": float, "alpha_pct": float, "benchmark_name": str} ou None.
    """
    rp = _portfolio_weighted_returns(lines, euro_rate, fee_pct, min_points)
    if rp is None:
        return None

    # Récupérer la série benchmark avec fallbacks
    df_bench: pd.DataFrame = pd.DataFrame()
    bench_name = benchmark_symbol
    for sym in [benchmark_symbol, "CW8.PA", "IWDA.AS"]:
        df_b, _, _ = get_price_series(sym, None, 0.0)
        if not df_b.empty:
            df_bench = df_b
            bench_name = sym
            break
    if df_bench.empty:
        return None

    rb = df_bench["Close"].astype(float).pct_change().dropna()
    rb.index = pd.DatetimeIndex([pd.Timestamp(x).normalize() for x in rb.index])

    # Aligner les deux séries sur les dates communes
    common = rp.index.intersection(rb.index)
    if len(common) < min_points:
        return None
    rp_c = rp.loc[common]
    rb_c = rb.loc[common]

    var_b = float(rb_c.var())
    if var_b == 0:
        return None
    beta = float(np.cov(rp_c.values, rb_c.values)[0, 1] / var_b)

    rf = risk_free_rate if risk_free_rate is not None else euro_rate / 100.0
    ret_ann_p = rp_c.mean() * 252.0
    ret_ann_b = rb_c.mean() * 252.0
    alpha_pct = float((ret_ann_p - (rf + beta * (ret_ann_b - rf))) * 100.0)

    return {"beta": beta, "alpha_pct": alpha_pct, "benchmark_name": bench_name}


def _corr_heatmap_chart(corr: pd.DataFrame, title: str) -> Optional[alt.Chart]:
    """
    Heatmap de corrélation — triangle inférieur uniquement,
    palette divergente RdYlGn inversée (rouge=corrélé, vert=décorrélé).
    """
    if corr.empty or corr.shape[0] < 2:
        return None
    labels = list(corr.index)
    rows = []
    for i, r in enumerate(labels):
        for j, c in enumerate(labels):
            if j <= i:  # triangle inférieur + diagonale
                rows.append({
                    "Ligne1": r,
                    "Ligne2": c,
                    "corr": float(corr.loc[r, c]),
                    "sort_i": i,
                    "sort_j": j,
                })
    df_melt = pd.DataFrame(rows)
    base = alt.Chart(df_melt).encode(
        x=alt.X("Ligne1:O", sort=labels, title="", axis=alt.Axis(labelAngle=-35)),
        y=alt.Y("Ligne2:O", sort=labels[::-1], title=""),
    )
    heat = base.mark_rect().encode(
        color=alt.Color(
            "corr:Q",
            scale=alt.Scale(
                domain=[-1, 0, 1],
                range=["#1a9641", "#ffffbf", "#d7191c"],
                type="linear",
            ),
            legend=alt.Legend(title="Corrélation", format=".2f"),
        ),
        tooltip=[
            alt.Tooltip("Ligne1:N", title="Ligne 1"),
            alt.Tooltip("Ligne2:N", title="Ligne 2"),
            alt.Tooltip("corr:Q", title="Corrélation", format=".2f"),
        ],
    )
    text = base.mark_text(
        baseline="middle",
        fontSize=11,
        fontWeight="bold",
    ).encode(
        text=alt.Text("corr:Q", format=".2f"),
        color=alt.condition(
            "datum.corr > 0.5 || datum.corr < -0.5",
            alt.value("white"),
            alt.value("#333333"),
        ),
    )
    return (heat + text).properties(title=title, height=max(220, len(labels) * 48))

# ------------------------------------------------------------
# Blocs de saisie : soit fonds recommandés, soit saisie libre
# FIXED: removed 2 duplicate comment blocks (Bug 9)
# ------------------------------------------------------------

def _add_from_reco_block(port_key: str, label: str):
    st.subheader(label)

    cat = st.selectbox(
        "Catégorie",
        ["Core (référence)", "Défensif", "Produits structurés"],
        key=f"reco_cat_{port_key}",
    )

    # ✅ Date d'achat centralisée (versement initial uniquement)
    buy_date = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if port_key == "A_lines"
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    # ============================
    # CAS 1 — PRODUIT STRUCTURÉ
    # ============================
    if cat == "Produits structurés":
        st.markdown("### Produit structuré (Autocall)")

        c1, c2 = st.columns(2)
        with c1:
            amount = st.text_input(
                "Montant investi (brut) €",
                value="",
                key=f"struct_amt_{port_key}",
            )
        with c2:
            struct_years = st.number_input(
                "Durée estimée avant remboursement (années)",
                min_value=1,
                max_value=12,
                value=6,
                step=1,
                key=f"struct_years_{port_key}",
            )

        struct_rate = st.number_input(
            "Rendement annuel estimé (%)",
            min_value=0.0,
            max_value=25.0,
            value=8.0,
            step=0.10,
            key=f"struct_rate_{port_key}",
        )

        st.caption(
            f"Date d'investissement initiale : {pd.Timestamp(buy_date).strftime('%d/%m/%Y')}"
        )

        if st.button("➕ Ajouter le produit structuré", key=f"struct_add_{port_key}"):
            try:
                amt = float(str(amount).replace(" ", "").replace(",", "."))
                assert amt > 0
            except Exception:
                st.warning("Montant invalide.")
                return

            ln = {
                "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
                "name": f"Produit structuré ({struct_rate:.2f}% / {int(struct_years)} ans)",
                "isin": "STRUCTURED",
                "amount_gross": float(amt),
                "buy_date": pd.Timestamp(buy_date),
                "buy_px": 1.0,
                "struct_rate": float(struct_rate),
                "struct_years": int(struct_years),
                "note": "",
                "sym_used": "STRUCTURED",
            }
            st.session_state[port_key].append(ln)
            st.toast("✅ Produit structuré ajouté.", icon="✅")
            st.rerun()
        return  # ✅ IMPORTANT : on sort de la fonction pour ne pas afficher la partie fonds

    # ============================
    # CAS 2 — FONDS CLASSIQUES
    # ============================
    if cat == "Core (référence)":
        fonds_list = RECO_FUNDS_CORE
    else:
        fonds_list = RECO_FUNDS_DEF

    options = [f"{nm} ({isin})" for nm, isin in fonds_list]
    choice = st.selectbox("Fonds recommandé", options, key=f"reco_choice_{port_key}")
    idx = options.index(choice) if choice in options else 0
    name, isin = fonds_list[idx]

    c1, c2 = st.columns([2, 2])
    with c1:
        amount = st.text_input("Montant investi (brut) €", value="", key=f"reco_amt_{port_key}")
    with c2:
        st.caption(f"Date d'achat (versement initial) : {pd.Timestamp(buy_date).strftime('%d/%m/%Y')}")

    px = st.text_input("Prix d'achat (optionnel)", value="", key=f"reco_px_{port_key}")

    if st.button("➕ Ajouter ce fonds recommandé", key=f"reco_add_{port_key}"):
        try:
            amt = float(str(amount).replace(" ", "").replace(",", "."))
            assert amt > 0
        except Exception:
            st.warning("Montant invalide.")
            return

        ln = {
            "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
            "name": name,
            "isin": isin,
            "amount_gross": float(amt),
            "buy_date": pd.Timestamp(buy_date),
            "buy_px": float(str(px).replace(",", ".")) if px.strip() else "",
            "note": "",
            "sym_used": "",
        }
        st.session_state[port_key].append(ln)
        st.toast("✅ Fonds recommandé ajouté.", icon="✅")
        st.rerun()


def _add_line_form_free(port_key: str, label: str):
    st.subheader(label)

    # ✅ Date d'achat centralisée (versement initial)
    buy_date_central = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if port_key == "A_lines"
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    with st.form(key=f"form_add_free_{port_key}", clear_on_submit=False):
        c1, c2 = st.columns([3, 2])

        with c1:
            name = st.text_input("Nom du fonds (libre)", value="")
            isin = st.text_input("ISIN ou code (peut être 'EUROFUND')", value="")

        with c2:
            amount = st.text_input("Montant investi (brut) €", value="")
            st.caption(
                f"Date d'achat (versement initial) : "
                f"{pd.Timestamp(buy_date_central).strftime('%d/%m/%Y')}"
            )

        px = st.text_input("Prix d'achat (optionnel)", value="")
        note = st.text_input("Note (optionnel)", value="")
        add_btn = st.form_submit_button("➕ Ajouter cette ligne")

    if not add_btn:
        return

    isin_final = isin.strip()
    name_final = name.strip()

    # Si nom vide mais ISIN renseigné : tentative de récupération du nom
    if not name_final and isin_final:
        res = eodhd_search(isin_final)
        match = None
        for it in res:
            if it.get("ISIN") == isin_final:
                match = it
                break
        if match is None and res:
            match = res[0]
        if match:
            name_final = match.get("Name", isin_final)

    if not name_final and isin_final.upper() == "EUROFUND":
        name_final = "Fonds en euros (EUROFUND)"

    if not name_final:
        name_final = isin_final or "—"

    try:
        amt = float(str(amount).replace(" ", "").replace(",", "."))
        assert amt > 0
    except Exception:
        st.warning("Montant invalide.")
        return

    ln = {
        "id": str(uuid.uuid4()),  # FIXED: stable UUID key for line identity (Bug 5)
        "name": name_final,
        "isin": isin_final or name_final,
        "amount_gross": float(amt),
        "buy_date": pd.Timestamp(buy_date_central),  # ✅ applique la date centrale
        "buy_px": float(str(px).replace(",", ".")) if px.strip() else "",
        "note": note.strip(),
        "sym_used": "",
    }

    st.session_state[port_key].append(ln)
    st.success("Ligne ajoutée.")


@st.cache_data(show_spinner=False, ttl=3600)
def _build_returns_matrix(
    isins: Tuple[str, ...],
    euro_rate: float,
    start: pd.Timestamp,
    end: pd.Timestamp,
    min_points: int = 60,
) -> Tuple[pd.DataFrame, List[str], Dict[str, str]]:
    series_map: Dict[str, pd.Series] = {}
    warnings: List[str] = []
    status: Dict[str, str] = {}

    for isin in isins:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty or df["Close"].dropna().shape[0] < min_points:
            warnings.append(f"{isin} : historique insuffisant")
            status[isin] = "insufficient"
            continue
        s = df["Close"].astype(float)
        s = s[(s.index >= start) & (s.index <= end)]
        if s.dropna().shape[0] < min_points:
            warnings.append(f"{isin} : historique insuffisant sur la fenêtre")
            status[isin] = "insufficient"
            continue
        series_map[isin] = s
        status[isin] = "ok"

    if not series_map:
        return pd.DataFrame(), warnings, status

    prices = pd.DataFrame(series_map).ffill().dropna(how="any")
    if prices.shape[0] < min_points:
        return pd.DataFrame(), warnings, status
    returns = prices.pct_change().dropna(how="any")
    return returns, warnings, status


def _suggest_weights(
    returns: pd.DataFrame,
    max_weight: float,
    min_funds: int,
) -> Dict[str, float]:
    if returns.empty:
        return {}
    corr = returns.corr()
    vols = returns.std() * np.sqrt(252.0)
    avg_corr = corr.mean()
    ranked = avg_corr.sort_values().index.tolist()
    min_funds = max(1, min(min_funds, len(ranked)))
    selected = ranked[:min_funds]
    inv_vol = 1 / vols[selected]
    weights = (inv_vol / inv_vol.sum()).to_dict()

    # cap weights if needed
    if max_weight > 0 and max_weight < 1:
        for _ in range(10):
            over = {k: v for k, v in weights.items() if v > max_weight}
            if not over:
                break
            excess = sum(v - max_weight for v in over.values())
            for k in over:
                weights[k] = max_weight
            remaining = {k: v for k, v in weights.items() if k not in over}
            if not remaining:
                break
            total_remaining = sum(remaining.values())
            for k in remaining:
                weights[k] += excess * (remaining[k] / total_remaining)

    total = sum(weights.values())
    if total > 0:
        weights = {k: v / total for k, v in weights.items()}
    return weights


def _round_allocations(amounts: Dict[str, float]) -> Dict[str, int]:
    floors = {k: int(np.floor(v)) for k, v in amounts.items()}
    remainder = int(round(sum(amounts.values()) - sum(floors.values())))
    if remainder <= 0:
        return floors
    fractions = sorted(
        ((k, amounts[k] - floors[k]) for k in amounts),
        key=lambda x: x[1],
        reverse=True,
    )
    for i in range(min(remainder, len(fractions))):
        k, _ = fractions[i]
        floors[k] += 1
    return floors


def _apply_weight_caps(weights: Dict[str, float], max_weight: float) -> Dict[str, float]:
    if not weights or max_weight <= 0:
        return weights
    weights = weights.copy()
    for _ in range(10):
        over = {k: v for k, v in weights.items() if v > max_weight}
        if not over:
            break
        excess = sum(v - max_weight for v in over.values())
        for k in over:
            weights[k] = max_weight
        remaining = {k: v for k, v in weights.items() if k not in over}
        if not remaining or excess <= 0:
            break
        total_remaining = sum(remaining.values())
        for k in remaining:
            weights[k] += excess * (remaining[k] / total_remaining)
    total = sum(weights.values())
    if total > 0:
        weights = {k: v / total for k, v in weights.items()}
    return weights


def _apply_min_floor_preserve_count(
    weights: Dict[str, float],
    floor: float,
) -> Dict[str, float]:
    if not weights:
        return {}
    keys = list(weights.keys())
    n = len(keys)
    if n == 0:
        return {}

    cleaned: Dict[str, float] = {}
    for k, v in weights.items():
        try:
            w = float(v)
        except Exception:
            w = 0.0
        if not np.isfinite(w) or w < 0:
            w = 0.0
        cleaned[k] = w

    total = sum(cleaned.values())
    if total <= 0:
        base = {k: 1.0 / n for k in keys}
    else:
        base = {k: cleaned[k] / total for k in keys}

    max_floor = (1.0 / n) - 1e-9
    eff_floor = min(max(floor, 0.0), max_floor) if max_floor > 0 else 0.0
    residual = 1.0 - n * eff_floor
    if residual < 0:
        residual = 0.0

    out = {k: eff_floor + residual * base[k] for k in keys}
    return out


def _round_weights_to_step_preserve_count(
    weights: Dict[str, float],
    step: float,
    min_w: float,
    max_w: Optional[float] = None,
) -> Dict[str, float]:
    if not weights:
        return {}
    if step <= 0:
        return weights

    keys = list(weights.keys())
    n = len(keys)
    if n == 0:
        return {}

    def _round_half_up(v: float) -> int:
        return int(np.floor((v / step) + 0.5 + 1e-12))

    eff_min = min_w
    if n > 0:
        max_floor = (1.0 / n) - 1e-9
        if max_floor > 0:
            eff_min = min(min_w, max_floor)
        else:
            eff_min = 0.0
    min_units = int(np.ceil(eff_min / step)) if eff_min > 0 else 0
    max_units = int(np.floor(max_w / step)) if (max_w is not None and max_w > 0) else None
    if max_units is not None and max_units < min_units:
        max_units = min_units
    target_units = int(round(1.0 / step))

    units: Dict[str, int] = {}
    for k in keys:
        u = _round_half_up(weights.get(k, 0.0))
        if u < min_units:
            u = min_units
        if max_units is not None and u > max_units:
            u = max_units
        units[k] = u

    delta = target_units - sum(units.values())
    if delta > 0:
        while delta > 0:
            candidates = [k for k in keys if max_units is None or units[k] < max_units]
            if not candidates:
                break
            k = max(candidates, key=lambda x: units[x])
            units[k] += 1
            delta -= 1
    elif delta < 0:
        while delta < 0:
            candidates = [k for k in keys if units[k] > min_units]
            if not candidates:
                break
            k = max(candidates, key=lambda x: units[x])
            units[k] -= 1
            delta += 1

    if delta != 0:
        k = max(keys, key=lambda x: units[x])
        units[k] = max(units[k] + delta, min_units)

    result = {k: units[k] * step for k in keys}
    total = sum(result.values())
    if total > 0 and abs(total - 1.0) > 1e-9:
        result = {k: v / total for k, v in result.items()}
    return result


def _apply_practical_constraints(
    weights: Dict[str, float],
    min_w: float = 0.10,
    step: float = 0.05,
    max_w: Optional[float] = None,
) -> Dict[str, float]:
    if not weights:
        return {}

    floored = _apply_min_floor_preserve_count(weights, min_w)

    if max_w is not None and max_w > 0:
        floored = _apply_weight_caps(floored, max_w)
        total = sum(floored.values())
        if total > 0:
            floored = {k: v / total for k, v in floored.items()}

    rounded = _round_weights_to_step_preserve_count(floored, step, min_w, max_w=max_w)
    return rounded


def _returns_for_isins(
    isins: List[str],
    start: pd.Timestamp,
    end: pd.Timestamp,
    euro_rate: float,
    min_points: int = 60,
) -> Tuple[pd.DataFrame, Dict[str, str]]:
    series_map: Dict[str, pd.Series] = {}
    status: Dict[str, str] = {}
    for isin in isins:
        df, _, _ = get_price_series(isin, None, euro_rate)
        if df.empty:
            status[isin] = "insufficient"
            continue
        s = df["Close"].astype(float)
        s = s[(s.index >= start) & (s.index <= end)]
        if s.dropna().shape[0] < min_points:
            status[isin] = "insufficient"
            continue
        series_map[isin] = s
        status[isin] = "ok"
    if not series_map:
        return pd.DataFrame(), status
    prices = pd.DataFrame(series_map).ffill().dropna(how="any")
    if prices.empty or prices.shape[0] < min_points:
        return pd.DataFrame(), status
    returns = prices.pct_change().dropna(how="any")
    return returns, status


def _annualized_stats(returns: pd.DataFrame) -> Tuple[pd.Series, pd.Series]:
    if returns.empty:
        return pd.Series(dtype=float), pd.Series(dtype=float)
    # FIXED: arithmetic compound formula — was mean()*252 (linear, not compounded) (Résidu Bug 2)
    ann_return = (1 + returns.mean()) ** 252 - 1
    ann_vol = returns.std() * np.sqrt(252.0)
    return ann_return, ann_vol


def _avg_correlation(corr: pd.DataFrame) -> float:
    if corr.empty or corr.shape[0] < 2:
        return np.nan
    vals = corr.values[np.triu_indices_from(corr, 1)]
    return float(np.nanmean(vals)) if vals.size else np.nan


def _avg_offdiag_corr(corr: pd.DataFrame) -> float:
    if corr.empty or corr.shape[0] < 2:
        return np.nan
    vals = corr.values[np.triu_indices_from(corr, 1)]
    return float(np.nanmean(vals)) if vals.size else np.nan


def _select_min_corr_subset(
    candidates: List[str],
    returns: pd.DataFrame,
    k: int,
    anchor: Optional[str] = None,
) -> List[str]:
    if k <= 0 or not candidates:
        return []
    if anchor and anchor not in candidates:
        return []
    if k >= len(candidates):
        return candidates
    corr = returns.corr()
    if corr.empty:
        return candidates[:k]
    if anchor:
        pool = [c for c in candidates if c != anchor]
        best_combo = [anchor]
        best_score = None
        max_checks = 2000
        combos = itertools.combinations(pool, k - 1)
        for idx, combo in enumerate(combos):
            if idx >= max_checks:
                break
            subset = [anchor, *combo]
            score = _avg_offdiag_corr(corr.loc[subset, subset])
            if best_score is None or score < best_score:
                best_score = score
                best_combo = subset
        if best_score is None:
            best_combo = [anchor] + pool[: k - 1]
        return best_combo
    max_checks = 2000
    best_combo: List[str] = []
    best_score = None
    for idx, combo in enumerate(itertools.combinations(candidates, k)):
        if idx >= max_checks:
            break
        subset = list(combo)
        score = _avg_offdiag_corr(corr.loc[subset, subset])
        if best_score is None or score < best_score:
            best_score = score
            best_combo = subset
    if not best_combo:
        best_combo = candidates[:k]
    return best_combo


def _greedy_select(
    candidates: List[str],
    returns: pd.DataFrame,
    target_count: int,
    forced: Optional[str] = None,
    corr_penalty: float = 0.35,
) -> List[str]:
    selected: List[str] = []
    if forced and forced in candidates:
        selected.append(forced)
    remaining = [c for c in candidates if c not in selected]
    if returns.empty or target_count <= 0:
        return selected
    ann_return, ann_vol = _annualized_stats(returns[remaining + selected] if remaining else returns)
    rfr = get_risk_free_rate()  # taux sans risque dynamique (Bund 10 ans ou saisie manuelle)
    sharpe = (ann_return - rfr) / ann_vol.replace(0, np.nan)
    corr = returns.corr() if not returns.empty else pd.DataFrame()

    while len(selected) < target_count and remaining:
        best = None
        best_score = None
        for cand in remaining:
            base_sharpe = float(sharpe.get(cand, np.nan))
            if np.isnan(base_sharpe):
                base_sharpe = -1e9
            if selected and not corr.empty:
                avg_corr = float(corr.loc[cand, selected].mean())
            else:
                avg_corr = 0.0
            score = base_sharpe - corr_penalty * avg_corr
            if best_score is None or score > best_score:
                best_score = score
                best = cand
        if best is None:
            break
        selected.append(best)
        remaining = [c for c in remaining if c != best]
    return selected


def _optimize_uc_weights(
    returns: pd.DataFrame,
    objective: str,
    min_weight: float,
    max_weight: float,
    target_vol: Optional[float],
    target_return: Optional[float],
) -> Dict[str, float]:
    if returns.empty:
        return {}
    n_assets = returns.shape[1]
    if n_assets == 0:
        return {}
    bounds = (min_weight, max_weight)
    try:
        if not PYPFOPT_AVAILABLE:
            raise RuntimeError(PYPFOPT_ERROR)
        mu = expected_returns.mean_historical_return(returns, frequency=252)
        cov = risk_models.sample_cov(returns, frequency=252)
        ef = EfficientFrontier(mu, cov, weight_bounds=bounds)
        if objective == "max_sharpe":
            weights = ef.max_sharpe(risk_free_rate=get_risk_free_rate())
        elif objective == "target_vol" and target_vol is not None:
            weights = ef.efficient_risk(target_vol)
        elif objective == "target_return" and target_return is not None:
            weights = ef.efficient_return(target_return)
        else:
            weights = ef.max_sharpe(risk_free_rate=get_risk_free_rate())
        cleaned = ef.clean_weights()
        total = sum(cleaned.values())
        if total > 0:
            return {k: v / total for k, v in cleaned.items()}
    except Exception:
        pass
    equal_weight = 1.0 / n_assets
    return {col: equal_weight for col in returns.columns}


def _select_min_corr_combo(
    returns: pd.DataFrame,
    k: int,
    anchor: Optional[str] = None,
) -> List[str]:
    if returns.empty:
        return []
    cols = list(returns.columns)
    if anchor and anchor not in cols:
        return []
    if k <= 0:
        return []
    if k > len(cols):
        return cols
    corr = returns.corr()
    best_combo: List[str] = []
    best_score = None
    if anchor:
        others = [c for c in cols if c != anchor]
        for combo in itertools.combinations(others, k - 1):
            candidate = [anchor, *combo]
            sub = corr.loc[candidate, candidate].to_numpy()
            triu = sub[np.triu_indices_from(sub, 1)]
            score = float(triu.mean()) if triu.size else 1.0
            if best_score is None or score < best_score:
                best_score = score
                best_combo = candidate
    else:
        for combo in itertools.combinations(cols, k):
            sub = corr.loc[list(combo), list(combo)].to_numpy()
            triu = sub[np.triu_indices_from(sub, 1)]
            score = float(triu.mean()) if triu.size else 1.0
            if best_score is None or score < best_score:
                best_score = score
                best_combo = list(combo)
    return best_combo


def _compute_drawdown(returns: pd.Series) -> Optional[float]:
    if returns.empty:
        return None
    cum = (1.0 + returns).cumprod()
    running_max = cum.cummax()
    dd = cum / running_max - 1.0
    return float(dd.min())


def _fund_name(isin: str) -> str:
    return FUND_NAME_MAP.get(isin, isin)


def _safe_fund_label(name: str, isin: str) -> str:
    cleaned = str(name or "").strip()
    if cleaned:
        return f"{cleaned} ({isin})"
    if isin:
        res = eodhd_search(isin)
        if res:
            maybe = res[0].get("Name") or res[0].get("name")
            if maybe:
                return f"{maybe} ({isin})"
    return f"{isin} ({isin})"


def render_portfolio_builder():
    st.title("Construction de portefeuille optimisé")

    st.session_state.setdefault("PP_RUN", False)
    st.session_state.setdefault("PP_PARAMS_HASH", "")
    st.session_state.setdefault("PP_OPT_START_DATE", (TODAY - pd.DateOffset(years=3)).date())
    st.session_state.setdefault("PP_OPT_END_DATE", TODAY.date())
    st.session_state.setdefault("PP_CONTRACT_LABEL", list(CONTRACTS_REGISTRY.keys())[0])
    st.session_state.setdefault("PP_TOTAL_UC_COUNT", 4)

    # Charger le référentiel du contrat sélectionné (indépendant de render_app)
    pp_contract_label = st.session_state.get("PP_CONTRACT_LABEL", list(CONTRACTS_REGISTRY.keys())[0])
    pp_contract_cfg = CONTRACTS_REGISTRY.get(pp_contract_label, list(CONTRACTS_REGISTRY.values())[0])
    funds_df = load_contract_funds(
        pp_contract_cfg["path"],
        pp_contract_cfg["funds_filename"],
    )
    CONTRACT_FUND_NAMES = dict(zip(funds_df["isin"], funds_df["name"])) if not funds_df.empty else FUND_NAME_MAP.copy()
    CONTRACT_FUND_FEES  = dict(zip(funds_df["isin"], funds_df["fee_total_pct"])) if not funds_df.empty else {}

    profile_map = {
        "Prudent": 50,
        "Equilibre": 30,
        "Dynamique": 20,
        "Agressif": 10,
    }

    with st.sidebar:
        st.header("Paramètres clés")

        # ── Contrat & fonds en euros ───────────────────────────
        pp_contract_label = st.selectbox(
            "Contrat",
            list(CONTRACTS_REGISTRY.keys()),
            key="PP_CONTRACT_LABEL",
        )
        pp_contract_cfg = CONTRACTS_REGISTRY[pp_contract_label]
        pp_euro_options = list(pp_contract_cfg.get("euro_funds", {}).keys())
        pp_euro_fund_label = st.selectbox(
            "Fonds en euros",
            pp_euro_options if pp_euro_options else ["—"],
            key="PP_EURO_FUND_LABEL",
        )
        st.markdown("---")

        # ── Profil & budget ────────────────────────────────────
        profile = st.selectbox(
            "Profil de risque client",
            list(profile_map.keys()),
            index=1,
            key="PP_PROFILE",
        )
        _recommended_euro_pct = profile_map[profile]
        st.info(
            f"📋 Allocation fonds en euros recommandée "
            f"pour un profil **{profile}** : **{_recommended_euro_pct}%**"
        )
        euro_pct = int(st.number_input(
            "Part fonds en euros (%)",
            min_value=0,
            max_value=100,
            value=int(st.session_state.get("PP_EURO_PCT", _recommended_euro_pct)),
            step=5,
            key="PP_EURO_PCT",
            help="Valeur recommandée selon le profil. Modifiable selon la situation du client.",
        ))

        total_budget = st.number_input(
            "Budget total (EUR)",
            min_value=0,
            max_value=10_000_000,
            value=100_000,
            step=10,
            key="PP_TOTAL_BUDGET",
        )

        # ── Fenêtre d'analyse (détermine opt_start / opt_end) ──
        opt_window_mode = st.radio(
            "Fenêtre d'analyse",
            ["1 an", "3 ans", "5 ans", "10 ans", "Dates personnalisées"],
            horizontal=False,
            key="PP_WINDOW_MODE",
        )

        if opt_window_mode == "Dates personnalisées":
            opt_start_date = st.date_input(
                "Date de début",
                value=st.session_state.get("PP_OPT_START_DATE", (TODAY - pd.DateOffset(years=3)).date()),
                key="PP_OPT_START_DATE",
            )
            opt_end_date = st.date_input(
                "Date de fin",
                value=st.session_state.get("PP_OPT_END_DATE", TODAY.date()),
                key="PP_OPT_END_DATE",
            )
            opt_start = pd.Timestamp(opt_start_date)
            opt_end = pd.Timestamp(opt_end_date)
        else:
            years_map = {"1 an": 1, "3 ans": 3, "5 ans": 5, "10 ans": 10}
            opt_years = years_map[opt_window_mode]
            opt_start = TODAY - pd.DateOffset(years=opt_years)
            opt_end = TODAY

        # ── Taux fonds en euros (auto depuis historique) ───────
        st.markdown("---")
        st.markdown("**Rendement fonds en euros**")
        try:
            _pp_ef_filename = pp_contract_cfg["euro_funds"].get(pp_euro_fund_label, "")
            _pp_ef_history = (
                load_euro_fund_history(pp_contract_cfg["path"], _pp_ef_filename)
                if _pp_ef_filename else pd.DataFrame()
            )
        except Exception:
            _pp_ef_history = pd.DataFrame()

        _auto_euro_rate = _compute_auto_euro_rate(_pp_ef_history, opt_start, opt_end)
        if _auto_euro_rate is not None:
            st.success(
                f"Taux net moyen **{pp_euro_fund_label}** "
                f"sur la fenêtre : **{_auto_euro_rate:.2f}%**"
            )
        else:
            st.info("Historique insuffisant sur la fenêtre — saisissez le taux manuellement.")

        _override_euro = st.checkbox(
            "Saisir un taux personnalisé",
            value=False,
            key="PP_OVERRIDE_EURO_RATE",
        )
        if _override_euro:
            euro_rate = float(st.number_input(
                "Rendement annuel du fonds en euros (%)",
                min_value=0.0, max_value=10.0,
                value=float(_auto_euro_rate if _auto_euro_rate is not None else 2.0),
                step=0.10,
                key="PP_EURO_RATE",
            ))
        else:
            euro_rate = float(_auto_euro_rate if _auto_euro_rate is not None else 2.0)
            st.session_state["PP_EURO_RATE"] = euro_rate
            st.caption(f"Taux utilisé : **{euro_rate:.2f}%**")

        # ── Taux sans risque (Sharpe) ──────────────────────────
        st.markdown("---")
        st.markdown("**Taux sans risque (Sharpe)**")
        rfr_api = _fetch_bund_rate_from_api()
        if rfr_api is not None:
            st.success(
                f"📡 Bund 10 ans : **{rfr_api * 100:.2f}%** — récupéré automatiquement"
            )
            st.session_state["RISK_FREE_RATE_SOURCE"] = "api"
            st.session_state["RISK_FREE_RATE_VALUE"] = rfr_api
        else:
            st.warning("⚠️ Bund 10 ans indisponible via API")
        default_pct = (
            rfr_api * 100.0
            if rfr_api is not None
            else float(st.session_state.get("RISK_FREE_RATE_MANUAL_PCT",
                                            RISK_FREE_RATE_FALLBACK * 100.0))
        )
        manual_rfr_pct = st.number_input(
            "Taux sans risque (%)",
            min_value=0.0,
            max_value=15.0,
            value=default_pct,
            step=0.05,
            format="%.2f",
            help=(
                "Taux de référence : rendement du Bund allemand 10 ans. "
                "Valeur actuelle : https://www.investing.com/rates-bonds/"
                "germany-10-year-bond-yield"
            ),
            key="RFR_MANUAL_INPUT",
        )
        st.session_state["RISK_FREE_RATE_MANUAL_PCT"] = manual_rfr_pct
        st.session_state["RISK_FREE_RATE_MANUAL"] = manual_rfr_pct / 100.0
        if rfr_api is None:
            st.session_state["RISK_FREE_RATE_SOURCE"] = "manual"

    # ── Fonds UC configurables ─────────────────────────────────
    st.markdown("### Fonds UC")

    # Utiliser _is_bond_category (définie au niveau module) pour une
    # classification cohérente avec le comparateur.
    _all_cats = sorted(funds_df["category"].dropna().unique().tolist()) if not funds_df.empty else []
    _bond_cats = [c for c in _all_cats if _is_bond_category(c)]
    _other_cats = [c for c in _all_cats if not _is_bond_category(c)]
    _bond_cat_options  = ["(Toutes obligataires)"] + _bond_cats
    _other_cat_options = ["(Toutes)"] + _other_cats

    # ── Fonds obligataires ─────────────────────────────────────
    st.markdown("#### Fonds obligataires")
    n_bond = int(st.number_input(
        "Nombre de fonds obligataires",
        min_value=0, max_value=8,
        value=int(st.session_state.get("PP_N_BOND", 1)),
        step=1,
        key="PP_N_BOND",
        help="Nombre de fonds UC à dominante obligataire à inclure dans le portefeuille.",
    ))
    bond_slot_configs: List[Dict[str, Any]] = []
    for _i in range(n_bond):
        with st.container(border=True):
            st.caption(f"Fonds obligataire {_i + 1}")
            _sc1, _sc2 = st.columns([3, 1])
            with _sc1:
                _cat = st.selectbox(
                    "Catégorie",
                    _bond_cat_options,
                    key=f"PP_BOND_CAT_{_i}",
                )
                _cat_normalized = "(Toutes)" if _cat == "(Toutes obligataires)" else _cat
            with _sc2:
                _sri = st.selectbox(
                    "SRI max",
                    ["Tous"] + [str(j) for j in range(1, 8)],
                    key=f"PP_BOND_SRI_{_i}",
                )
        bond_slot_configs.append({"cat": _cat_normalized, "sri": _sri})

    st.markdown("---")

    # ── Autres fonds UC ────────────────────────────────────────
    st.markdown("#### Autres fonds UC")
    n_other = int(st.number_input(
        "Nombre d'autres fonds UC",
        min_value=0, max_value=10,
        value=int(st.session_state.get("PP_N_OTHER", 3)),
        step=1,
        key="PP_N_OTHER",
        help="Fonds actions, diversifiés, thématiques, etc.",
    ))
    other_slot_configs: List[Dict[str, Any]] = []
    for _i in range(n_other):
        with st.container(border=True):
            st.caption(f"Fonds UC {_i + 1}")
            _sc1, _sc2 = st.columns([3, 1])
            with _sc1:
                _cat = st.selectbox(
                    "Catégorie",
                    _other_cat_options,
                    key=f"PP_OTHER_CAT_{_i}",
                )
            with _sc2:
                _sri = st.selectbox(
                    "SRI max",
                    ["Tous"] + [str(j) for j in range(1, 8)],
                    key=f"PP_OTHER_SRI_{_i}",
                )
        other_slot_configs.append({"cat": _cat, "sri": _sri})

    # ── Fusion pour compatibilité avec le reste du code ────────
    uc_slot_configs: List[Dict[str, Any]] = bond_slot_configs + other_slot_configs
    total_uc_count = len(uc_slot_configs)
    st.session_state["PP_TOTAL_UC_COUNT"] = total_uc_count

    objective_options = [
        "Maximiser Sharpe",
        "Minimiser volatilite",
        "Maximiser rendement annualise",
        "Meilleur compromis (Sharpe + diversification)",
        "Diversification maximale (decorrelation)",
    ]
    if st.session_state.get("PP_OBJECTIVE") == "Risk parity":
        st.session_state["PP_OBJECTIVE"] = objective_options[0]
    objective_choice = st.selectbox(
        "Objectif",
        objective_options,
        key="PP_OBJECTIVE",
    )

    practical_mode = st.checkbox(
        "Optimisation terrain (min 10% + arrondi 5%)",
        value=True,
        key="PP_PRACTICAL_MODE",
    )

    force_fund = st.checkbox("Forcer un fonds (ancre)", value=False, key="PP_FORCE_ANCHOR")
    forced_isin = None
    if force_fund:
        if not funds_df.empty:
            _anchor_options = [
                f"{row['name']} ({row['isin']})"
                for _, row in funds_df.iterrows()
            ]
            _anchor_lookup = {
                f"{row['name']} ({row['isin']})": row["isin"]
                for _, row in funds_df.iterrows()
            }
            _forced_choice = st.selectbox("Fonds ancre", _anchor_options, key="PP_ANCHOR_ISIN")
            forced_isin = _anchor_lookup.get(_forced_choice)

    params = (
        profile,
        int(euro_pct),
        float(euro_rate),
        int(total_budget),
        opt_window_mode,
        str(opt_start.date()),
        str(opt_end.date()),
        tuple((s["cat"], s["sri"]) for s in uc_slot_configs),
        objective_choice,
        bool(force_fund),
        forced_isin,
    )
    current_hash = params_hash(params)

    if st.session_state.get("PP_RUN") and st.session_state.get("PP_PARAMS_HASH") != current_hash:
        st.session_state["PP_RUN"] = False
        st.info("Paramètres modifiés. Cliquez de nouveau sur '✅ Lancer l'optimisation'.")

    run_clicked = st.button("✅ Lancer l'optimisation", type="primary")
    if run_clicked:
        st.session_state["PP_RUN"] = True
        st.session_state["PP_PARAMS_HASH"] = current_hash

    if not st.session_state.get("PP_RUN"):
        st.info("Renseignez tous les paramètres puis cliquez sur '✅ Lancer l'optimisation'.")
        return

    if opt_start > opt_end:
        st.warning("La date de debut doit etre anterieure a la date de fin.")
        return

    try:
        if funds_df.empty:
            st.warning(
                "Référentiel de fonds non chargé. "
                "Sélectionnez un contrat dans la barre latérale."
            )
            return

        # Construire l'univers de candidats depuis les slots UC
        def _slot_candidates(slot: Dict[str, Any]) -> List[str]:
            """Retourne la liste des ISIN du contrat correspondant au filtre slot."""
            df = funds_df.copy()
            if slot["cat"] != "(Toutes)":
                df = df[df["category"] == slot["cat"]]
            if slot["sri"] != "Tous":
                df = df[df["sri"] <= int(slot["sri"])]
            return df["isin"].tolist()

        all_slot_isins: List[str] = []
        for slot in uc_slot_configs:
            all_slot_isins.extend(_slot_candidates(slot))

        all_candidates = sorted(set(all_slot_isins))
        if not all_candidates:
            st.info("Aucun fonds UC disponible avec les critères de slots choisis.")
            return

        returns_all, status_all = _returns_for_isins(all_candidates, opt_start, opt_end, euro_rate=float(euro_rate))

        insufficient = [isin for isin, status in status_all.items() if status != "ok"]
        valid_all = [isin for isin in all_candidates if status_all.get(isin) == "ok"]

        if insufficient:
            st.warning("Certains fonds ont été exclus (historique insuffisant sur la fenêtre).")

        if forced_isin and forced_isin not in valid_all:
            st.warning("Le fonds ancre est indisponible sur la période et a été ignoré.")
            forced_isin = None

        if returns_all.empty:
            st.info("Historique insuffisant pour calculer les rendements sur la fenêtre choisie.")
            return

        def _zscore(s: pd.Series) -> pd.Series:
            if s.empty:
                return s
            std = float(s.std())
            if std == 0 or np.isnan(std):
                return pd.Series(0.0, index=s.index)
            return (s - s.mean()) / std

        def _stats_for(candidates: List[str], ret_df: pd.DataFrame) -> Tuple[pd.Series, pd.Series, pd.Series, pd.DataFrame, pd.Series]:
            if not candidates or ret_df.empty:
                return pd.Series(dtype=float), pd.Series(dtype=float), pd.Series(dtype=float), pd.DataFrame(), pd.Series(dtype=float)
            # FIXED: arithmetic annualisation formula — was incorrectly treating pct_change() returns
            # as log-returns then applying exp(), double-compounding the result (Bug 2)
            rfr = get_risk_free_rate()  # taux sans risque dynamique (Bund 10 ans ou saisie manuelle)
            ann_return = (1 + ret_df[candidates].mean()) ** 252 - 1
            ann_vol = ret_df[candidates].std() * np.sqrt(252.0)
            sharpe = ((ann_return - rfr) / ann_vol.replace(0, np.nan)).replace([np.inf, -np.inf], np.nan)
            corr = ret_df[candidates].corr()
            avg_corr = corr.apply(lambda row: row.drop(labels=row.name, errors="ignore").mean(), axis=1).fillna(0.0) if not corr.empty else pd.Series(0.0, index=candidates)
            return ann_return, ann_vol, sharpe, corr, avg_corr

        def _greedy_min_corr(candidates: List[str], corr: pd.DataFrame, k: int, seed: Optional[List[str]] = None) -> List[str]:
            if k <= 0 or not candidates:
                return []
            selected = [x for x in (seed or []) if x in candidates]
            remaining = [c for c in candidates if c not in selected]
            while len(selected) < k and remaining:
                best, best_score = None, None
                for cand in remaining:
                    score = float(corr.loc[cand, selected].mean()) if selected and not corr.empty else 0.0
                    if best_score is None or score < best_score:
                        best, best_score = cand, score
                if best is None:
                    break
                selected.append(best)
                remaining = [c for c in remaining if c != best]
            return selected

        def _select_by_objective(candidates: List[str], ret_df: pd.DataFrame, k: int, objective: str, forced: Optional[str] = None) -> List[str]:
            if k <= 0 or not candidates:
                return []
            if ret_df.empty:
                base = candidates[:k]
                if forced and forced in candidates and forced not in base:
                    base = [forced] + [x for x in base if x != forced]
                    base = base[:k]
                return base

            ann_ret, ann_vol, sharpe, corr, avg_corr = _stats_for(candidates, ret_df)
            sharpe_rank = sharpe.fillna(-1e9).sort_values(ascending=False)

            if objective == "Maximiser rendement annualise":
                selected = ann_ret.sort_values(ascending=False).index.tolist()[:k]
            elif objective == "Minimiser volatilite":
                selected = ann_vol.sort_values(ascending=True).index.tolist()[:k]
            elif objective == "Maximiser Sharpe":
                selected = sharpe_rank.index.tolist()[:k]
            elif objective in ("Diversification maximale (decorrelation)", "Risk parity"):
                seed = sharpe_rank.index.tolist()[:1]
                selected = _greedy_min_corr(candidates, corr, k, seed=seed)
            elif objective == "Meilleur compromis (Sharpe + diversification)":
                score = _zscore(sharpe.fillna(0.0)) + _zscore(1.0 - avg_corr)
                selected = score.sort_values(ascending=False).index.tolist()[:k]
            else:
                selected = sharpe_rank.index.tolist()[:k]

            if forced and forced in candidates:
                if forced not in selected:
                    if objective in ("Diversification maximale (decorrelation)", "Risk parity"):
                        selected = _greedy_min_corr(candidates, corr, k, seed=[forced])
                    else:
                        rest = [x for x in selected if x != forced]
                        selected = [forced] + rest
                        selected = selected[:k]
                else:
                    selected = [forced] + [x for x in selected if x != forced]
            return selected[:k]

        # Sélection par slot : 1 fonds par slot, dédupliqué
        _seen: set = set()
        selected_isins: List[str] = []
        _isin_slot_label: Dict[str, str] = {}  # isin → label catégorie pour affichage
        _anchor_placed = False
        for _si, _slot in enumerate(uc_slot_configs):
            _slot_valid = [
                isin for isin in _slot_candidates(_slot)
                if isin in valid_all and isin not in _seen
            ]
            if not _slot_valid:
                continue
            _forced_for_slot = forced_isin if (not _anchor_placed and forced_isin in _slot_valid) else None
            _slot_returns = returns_all[_slot_valid] if _slot_valid else pd.DataFrame()
            _picked = _select_by_objective(_slot_valid, _slot_returns, 1, objective_choice, forced=_forced_for_slot)
            for _isin in _picked:
                if _isin not in _seen:
                    selected_isins.append(_isin)
                    _seen.add(_isin)
                    _cat_label = _slot["cat"] if _slot["cat"] != "(Toutes)" else (
                        funds_df.loc[funds_df["isin"] == _isin, "category"].values[0]
                        if not funds_df.loc[funds_df["isin"] == _isin].empty else "UC"
                    )
                    _isin_slot_label[_isin] = _cat_label
                    if _forced_for_slot and _isin == forced_isin:
                        _anchor_placed = True
        if not selected_isins:
            st.info("Aucun fonds selectionne apres filtrage.")
            return

        returns_selected = returns_all[selected_isins].dropna(how="any")
        if returns_selected.empty:
            st.info("Historique insuffisant apres selection des UC.")
            return

        uc_total = max(0.0, 1.0 - (float(euro_pct) / 100.0))
        if uc_total <= 0.0:
            st.info("Part UC nulle avec le profil choisi.")
            return

        cap_uc_final = 0.20  # max 20% du portefeuille total par fonds
        uc_max_bound = min(cap_uc_final / uc_total, 1.0)

        def _risk_parity_weights(ret_df: pd.DataFrame) -> Dict[str, float]:
            cov = ret_df.cov()
            if cov.empty:
                return {}
            vols = np.sqrt(np.diag(cov.values))
            inv = np.where(vols > 0, 1.0 / vols, 0.0)
            if float(inv.sum()) <= 0:
                return {c: 1.0 / len(ret_df.columns) for c in ret_df.columns}
            w = inv / inv.sum()
            return {c: float(w[i]) for i, c in enumerate(ret_df.columns)}

        weights_uc_raw: Dict[str, float] = {}
        if PYPFOPT_AVAILABLE:
            try:
                mu = expected_returns.mean_historical_return(returns_selected, returns_data=True, frequency=252)
                cov = risk_models.sample_cov(returns_selected, returns_data=True, frequency=252)
                ef = EfficientFrontier(mu, cov, weight_bounds=(0.0, uc_max_bound))

                if objective_choice == "Maximiser Sharpe":
                    ef.max_sharpe(risk_free_rate=get_risk_free_rate())
                    weights_uc_raw = ef.clean_weights()
                elif objective_choice == "Minimiser volatilite":
                    ef.min_volatility()
                    weights_uc_raw = ef.clean_weights()
                elif objective_choice == "Maximiser rendement annualise":
                    if not mu.empty:
                        target = float(mu.max()) - 1e-6
                        try:
                            ef.efficient_return(target_return=target)
                        except Exception:
                            ef.max_sharpe(risk_free_rate=get_risk_free_rate())
                    else:
                        ef.max_sharpe(risk_free_rate=get_risk_free_rate())
                    weights_uc_raw = ef.clean_weights()
                elif objective_choice == "Risk parity":
                    weights_uc_raw = _risk_parity_weights(returns_selected)
                else:
                    ef.min_volatility() if objective_choice == "Diversification maximale (decorrelation)" else ef.max_sharpe(risk_free_rate=get_risk_free_rate())
                    weights_uc_raw = ef.clean_weights()
            except Exception:
                weights_uc_raw = {}

        if not weights_uc_raw:
            if objective_choice == "Risk parity":
                weights_uc_raw = _risk_parity_weights(returns_selected)
            elif objective_choice == "Minimiser volatilite":
                vol = returns_selected.std() * np.sqrt(252.0)
                score = (1.0 / vol.replace(0, np.nan)).fillna(0.0)
                weights_uc_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}
            elif objective_choice == "Maximiser rendement annualise":
                # FIXED: arithmetic compound annualisation (Résidu Bug 2)
                ann_ret = (1 + returns_selected.mean()) ** 252 - 1
                score = ann_ret.clip(lower=0.0)
                weights_uc_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}
            elif objective_choice == "Meilleur compromis (Sharpe + diversification)":
                # FIXED: arithmetic compound annualisation for Sharpe (Résidu Bug 2)
                rfr = get_risk_free_rate()
                ann_ret = (1 + returns_selected.mean()) ** 252 - 1
                vol = returns_selected.std() * np.sqrt(252.0)
                sharpe = ((ann_ret - rfr) / vol.replace(0, np.nan)).fillna(0.0)
                corr = returns_selected.corr()
                avg_corr = corr.apply(lambda row: row.drop(labels=row.name, errors="ignore").mean(), axis=1).fillna(0.0) if not corr.empty else pd.Series(0.0, index=returns_selected.columns)
                score = _zscore(sharpe) + _zscore(1.0 - avg_corr)
                score = score.clip(lower=0.0)
                weights_uc_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}
            elif objective_choice == "Diversification maximale (decorrelation)":
                corr = returns_selected.corr()
                avg_corr = corr.apply(lambda row: row.drop(labels=row.name, errors="ignore").mean(), axis=1).fillna(0.0) if not corr.empty else pd.Series(0.0, index=returns_selected.columns)
                score = (1.0 - avg_corr).clip(lower=0.0)
                weights_uc_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}
            else:
                # FIXED: arithmetic compound annualisation for default Sharpe (Résidu Bug 2)
                rfr = get_risk_free_rate()
                ann_ret = (1 + returns_selected.mean()) ** 252 - 1
                vol = returns_selected.std() * np.sqrt(252.0)
                score = ((ann_ret - rfr) / vol.replace(0, np.nan)).fillna(0.0).clip(lower=0.0)
                weights_uc_raw = (score / score.sum()).to_dict() if float(score.sum()) > 0 else {}

        if not weights_uc_raw:
            weights_uc_raw = {isin: 1.0 / len(selected_isins) for isin in selected_isins}

        weights_uc_raw = {k: float(v) for k, v in weights_uc_raw.items() if k in selected_isins}
        if not weights_uc_raw:
            weights_uc_raw = {isin: 1.0 / len(selected_isins) for isin in selected_isins}

        weights_uc_raw = _apply_weight_caps(weights_uc_raw, uc_max_bound)
        total_uc_raw = float(sum(weights_uc_raw.values()))
        if total_uc_raw <= 0:
            weights_uc_raw = {isin: 1.0 / len(selected_isins) for isin in selected_isins}
            weights_uc_raw = _apply_weight_caps(weights_uc_raw, uc_max_bound)
            total_uc_raw = float(sum(weights_uc_raw.values()))
        if total_uc_raw > 0:
            weights_uc_raw = {k: v / total_uc_raw for k, v in weights_uc_raw.items()}

        if practical_mode:
            weights_uc_raw = _apply_practical_constraints(
                weights_uc_raw,
                min_w=0.10,
                step=0.05,
                max_w=uc_max_bound,
            )

        weights_final = {"EUROFUND": float(euro_pct) / 100.0}
        for isin in selected_isins:
            weights_final[isin] = float(weights_uc_raw.get(isin, 0.0)) * uc_total

        # Security clamp + exact renormalization
        weights_final = {k: max(0.0, float(v)) for k, v in weights_final.items()}
        total_weight = float(sum(weights_final.values()))
        if total_weight > 0:
            weights_final = {k: v / total_weight for k, v in weights_final.items()}

        def _round_allocations_to_step(amounts: Dict[str, float], step: int, total: int) -> Dict[str, int]:
            if step <= 0:
                step = 1
            rounded = {k: int(np.floor(max(0.0, v) / step)) * step for k, v in amounts.items()}
            remaining = int(total - sum(rounded.values()))
            if remaining > 0 and rounded:
                frac_rank = sorted(
                    ((k, (max(0.0, amounts[k]) / step) - np.floor(max(0.0, amounts[k]) / step)) for k in rounded.keys()),
                    key=lambda x: x[1],
                    reverse=True,
                )
                idx = 0
                while remaining >= step and frac_rank:
                    k = frac_rank[idx % len(frac_rank)][0]
                    rounded[k] += step
                    remaining -= step
                    idx += 1
                if remaining != 0:
                    last_key = list(rounded.keys())[-1]
                    rounded[last_key] = max(0, rounded[last_key] + remaining)
            elif remaining < 0 and rounded:
                last_key = list(rounded.keys())[-1]
                rounded[last_key] = max(0, rounded[last_key] + remaining)

            delta = int(total - sum(rounded.values()))
            if rounded and delta != 0:
                last_key = list(rounded.keys())[-1]
                rounded[last_key] = max(0, rounded[last_key] + delta)
            return rounded

        amounts_raw = {k: float(total_budget) * float(w) for k, w in weights_final.items()}
        amounts = _round_allocations_to_step(amounts_raw, step=10, total=int(total_budget))

        rows = [
            {
                "Support": "Fonds en euros (EUROFUND)",
                "Categorie": "EUROFUND",
                "%": weights_final.get("EUROFUND", 0.0) * 100.0,
                "Montant EUR": amounts.get("EUROFUND", 0),
            }
        ]
        for isin in selected_isins:
            rows.append(
                {
                    "Support": CONTRACT_FUND_NAMES.get(isin, FUND_NAME_MAP.get(isin, isin)),
                    "Categorie": _isin_slot_label.get(isin, "UC"),
                    "%": weights_final.get(isin, 0.0) * 100.0,
                    "Montant EUR": amounts.get(isin, 0),
                }
            )
        df_alloc = pd.DataFrame(rows)

        st.markdown("**Allocation finale**")
        st.dataframe(
            df_alloc.style.format({"%": "{:,.2f}%".format, "Montant EUR": to_eur}),
            use_container_width=True,
            hide_index=True,
        )

        if MATPLOTLIB_AVAILABLE:
            fig, ax = plt.subplots(figsize=(5.0, 3.2))
            ax.pie(df_alloc["Montant EUR"], labels=df_alloc["Support"], autopct="%1.1f%%")
            ax.set_title("Repartition finale")
            st.pyplot(fig)
            plt.close(fig)
        else:
            st.warning(f"Graphique indisponible ({MATPLOTLIB_ERROR}).")

        if len(selected_isins) >= 2:
            st.markdown("**Heatmap correlation UC**")
            corr_uc = returns_selected.corr().copy()
            corr_uc["Ligne1"] = corr_uc.index
            heat_df = corr_uc.melt(id_vars="Ligne1", var_name="Ligne2", value_name="corr")
            heat = (
                alt.Chart(heat_df)
                .mark_rect()
                .encode(
                    x=alt.X("Ligne1:O", sort=None, title=""),
                    y=alt.Y("Ligne2:O", sort=None, title=""),
                    color=alt.Color("corr:Q", scale=alt.Scale(domain=[-1, 1])),
                    tooltip=[
                        alt.Tooltip("Ligne1:N", title="Ligne 1"),
                        alt.Tooltip("Ligne2:N", title="Ligne 2"),
                        alt.Tooltip("corr:Q", title="Correlation", format=".2f"),
                    ],
                )
                .properties(height=260)
            )
            st.altair_chart(heat, use_container_width=True)

        uc_weights_norm = pd.Series({k: weights_uc_raw.get(k, 0.0) for k in selected_isins}, index=selected_isins)
        uc_weights_norm = uc_weights_norm / uc_weights_norm.sum() if float(uc_weights_norm.sum()) > 0 else pd.Series(1.0 / len(selected_isins), index=selected_isins)

        port_log_ret = returns_selected[selected_isins].dot(uc_weights_norm.values)
        # FIXED: arithmetic compound annualisation — was exp(mean*252)-1 on arithmetic returns (Résidu Bug 2)
        rfr = get_risk_free_rate()  # taux sans risque dynamique (Bund 10 ans ou saisie manuelle)
        ann_ret_uc = float((1 + port_log_ret.mean()) ** 252 - 1)
        ann_vol_uc = float(port_log_ret.std() * np.sqrt(252.0))
        sharpe_uc = (ann_ret_uc - rfr) / ann_vol_uc if ann_vol_uc > 0 else np.nan

        # FIXED: use fee-aware wrapper so management fees are included in the return metric (Correction 1)
        euro_df, _, _ = get_price_series_with_fees("EUROFUND", None, float(euro_rate))
        euro_df = euro_df.loc[(euro_df.index >= opt_start) & (euro_df.index <= opt_end)]
        if euro_df.empty:
            euro_total_ret = 0.0
        else:
            euro_total_ret = float(euro_df["Close"].iloc[-1] / euro_df["Close"].iloc[0] - 1.0)

        uc_path = (1.0 + port_log_ret).cumprod()
        uc_total_ret = float(uc_path.iloc[-1] - 1.0) if not uc_path.empty else 0.0 if not uc_path.empty else 0.0

        total_ret = (float(euro_pct) / 100.0) * euro_total_ret + uc_total * uc_total_ret

        k1, k2, k3, k4 = st.columns(4)
        with k1:
            st.metric("Rendement annualise UC", fmt_pct_fr(ann_ret_uc * 100.0))
        with k2:
            st.metric("Volatilite annualisee UC", fmt_pct_fr(ann_vol_uc * 100.0))
        with k3:
            st.metric("Sharpe UC", f"{sharpe_uc:.2f}" if sharpe_uc == sharpe_uc else "-")
        with k4:
            st.metric("Rendement total (EUROFUND+UC)", fmt_pct_fr(total_ret * 100.0))

        rfr_display = get_risk_free_rate()
        rfr_src = st.session_state.get("RISK_FREE_RATE_SOURCE", "manual")
        rfr_label = (
            f"📡 Bund 10 ans en temps réel ({rfr_display * 100:.2f}%)"
            if rfr_src == "api"
            else f"✏️ Taux saisi manuellement ({rfr_display * 100:.2f}%)"
        )
        st.caption(
            f"Ratio de Sharpe calculé avec un taux sans risque de "
            f"**{rfr_display * 100:.2f}%** — {rfr_label}"
        )

        st.markdown(
            """
- Allocation calibree sur la fenetre d'analyse et l'objectif choisi.
- EUROFUND est traite uniquement via son taux annuel parametre.
- Les UC respectent un cap strict de 25% du portefeuille final par fonds.
- En mode ancre, le fonds impose est conserve dans la poche actions.
- En cas de donnees insuffisantes, la selection est reduite automatiquement.
"""
        )

        if practical_mode:
            st.caption(
                "Les ponderations UC sont ajustees pour rester coherentes en pratique "
                "(min 10 % et arrondi par paliers de 5 %)."
            )

        if insufficient:
            st.caption("Fonds exclus : " + ", ".join(insufficient))
        st.caption(f"Fenetre utilisee : {fmt_date(opt_start)} a {fmt_date(opt_end)}")

    except Exception as e:
        st.error("Une erreur est survenue dans le builder.")
        st.exception(e)


# ------------------------------------------------------------
# Sélection fonds depuis le référentiel contrat
# ------------------------------------------------------------

_BOND_CATEGORY_PREFIXES = ("OBLIGATIONS",)
_BOND_CATEGORY_KEYWORDS = (
    "CONVERTIBLES",
    "SUBORDINATED BOND",
    "REVENUS FIXES",
)


def _is_bond_category(cat: str) -> bool:
    """
    Retourne True si la catégorie Morningstar est de nature obligataire.
    Couvre :
    - Toutes catégories commençant par "OBLIGATIONS" (ex : OBLIGATIONS EUR...)
    - Convertibles (CONVERTIBLES EUR, CONVERTIBLES EUROPE...)
    - Dettes subordonnées (EUR SUBORDINATED BOND)
    - Revenus fixes (REVENUS FIXES EUR — fonds à échéance)
    - Produits leveragés/inversés sur obligations
      (TRADING - LEVERAGED/INVERSE OBLIGATIONS)
    """
    c = str(cat).upper().strip()
    return (
        c.startswith(_BOND_CATEGORY_PREFIXES)
        or any(k in c for k in _BOND_CATEGORY_KEYWORDS)
        or "OBLIGATIONS" in c
    )


def _get_assureur(contract_label: str) -> str:
    cfg = CONTRACTS_REGISTRY.get(contract_label, {})
    return cfg.get("assureur", "—")


def _add_fund_from_contract(port_key: str, label: str):
    """
    Interface d'ajout de fonds structurée en deux blocs :
      Bloc A — Actifs défensifs (fonds euros + fonds obligataires)
      Bloc B — Autres UC (actions, diversifiés, etc.)
    Utilise le contrat spécifique au portefeuille (A ou B).
    """
    st.subheader(label)

    is_a = port_key == "A_lines"
    contract_label = st.session_state.get(
        "CONTRACT_LABEL_A" if is_a else "CONTRACT_LABEL_B", ""
    )
    contract_cfg = CONTRACTS_REGISTRY.get(contract_label, {})
    funds_df = st.session_state.get(
        "CONTRACT_FUNDS_DF_A" if is_a else "CONTRACT_FUNDS_DF_B", pd.DataFrame()
    )
    buy_date_central = (
        st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
        if is_a
        else st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    )

    if not contract_label or not contract_cfg:
        st.warning("Contrat non sélectionné — configurez-le dans les paramètres.")
        return

    assureur = contract_cfg.get("assureur", "—")
    st.caption(
        f"Contrat : **{contract_label}** ({assureur}) "
        f"— {len(funds_df)} UC disponibles "
        f"| Date d'achat : {pd.Timestamp(buy_date_central).strftime('%d/%m/%Y')}"
    )

    # ═══════════════════════════════════════════════════════════
    # BLOC A — Actifs défensifs
    # ═══════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 🛡️ Actifs défensifs")

    # ── Sous-section 1 : Fonds en euros ──────────────────────
    st.markdown("#### Fonds en euros — capital garanti à 98%")

    ef_options = list(contract_cfg.get("euro_funds", {}).keys())
    if ef_options:
        ef_selected = st.selectbox("Fonds en euros", ef_options, key=f"ef_sel_{port_key}")
        ef_filename = contract_cfg["euro_funds"][ef_selected]
        try:
            ef_history = load_euro_fund_history(contract_cfg["path"], ef_filename)
        except Exception:
            ef_history = pd.DataFrame()
        ef_avg = get_euro_fund_avg_rate(ef_history, years=5)
        ef_last = (
            float(ef_history.iloc[-1]["taux_net_publie_pct"])
            if not ef_history.empty else ef_avg
        )
        with st.container(border=True):
            hc1, hc2 = st.columns([3, 1])
            with hc1:
                st.markdown(f"**{ef_selected}**")
                st.caption(
                    f"Assureur : {assureur} | SRI 1 | Capital garanti à 98%"
                    f" | Taux net {int(ef_history['annee'].max()) if not ef_history.empty else '—'}"
                    f" : **{ef_last:.2f}%**"
                )
            with hc2:
                st.metric("Moy. 5 ans", f"{ef_avg:.2f}%")
        euro_amt = st.number_input(
            "Montant (€)",
            min_value=0.0, max_value=10_000_000.0,
            value=0.0, step=1000.0,
            key=f"euro_amt_{port_key}",
        )
        if st.button("➕ Ajouter fonds en euros", key=f"add_euro_{port_key}"):
            if euro_amt <= 0:
                st.warning("Montant invalide.")
            else:
                st.session_state[port_key].append({
                    "id":            str(uuid.uuid4()),
                    "name":          ef_selected,
                    "isin":          "EUROFUND",
                    "amount_gross":  float(euro_amt),
                    "buy_date":      pd.Timestamp(buy_date_central),
                    "buy_px":        "",
                    "note":          f"Contrat : {contract_label}",
                    "sym_used":      "EUROFUND",
                    "fee_total_pct": 0.0,
                })
                st.toast(f"✅ {ef_selected} ajouté.", icon="✅")
                st.rerun()

    # ── Sous-section 2 : Fonds obligataires ──────────────────
    st.markdown("#### Fonds obligataires")

    if funds_df.empty:
        st.info("Référentiel de fonds non chargé.")
    else:
        bond_df = funds_df[funds_df["category"].apply(_is_bond_category)].copy()

        if bond_df.empty:
            st.info("Aucun fonds obligataire dans le référentiel.")
        else:
            bf1, bf2, bf3 = st.columns([3, 2, 1])
            with bf1:
                bond_search = st.text_input(
                    "Rechercher (nom, société de gestion)",
                    value="", key=f"bond_search_{port_key}",
                    placeholder="Ex : AXA, Amundi...",
                )
            with bf2:
                bond_cats = ["Toutes"] + sorted(bond_df["category"].dropna().unique().tolist())
                bond_cat = st.selectbox("Catégorie", bond_cats, key=f"bond_cat_{port_key}")
            with bf3:
                bond_sri = st.selectbox("SRI", ["Tous"] + [str(i) for i in range(1, 8)], key=f"bond_sri_{port_key}")

            if bond_search.strip():
                q = bond_search.strip().upper()
                bond_df = bond_df[
                    bond_df["name"].str.upper().str.contains(q, na=False) |
                    bond_df["manager"].str.upper().str.contains(q, na=False)
                ]
            if bond_cat != "Toutes":
                bond_df = bond_df[bond_df["category"] == bond_cat]
            if bond_sri != "Tous":
                bond_df = bond_df[bond_df["sri"] == int(bond_sri)]

            if bond_df.empty:
                st.info("Aucun résultat.")
            else:
                # Téléchargement liste filtrée
                _bond_xlsx = BytesIO()
                bond_df[["isin", "name", "manager", "category", "sri", "fee_total_pct"]].rename(
                    columns={
                        "isin": "ISIN", "name": "Libellé", "manager": "Société",
                        "category": "Catégorie", "sri": "SRI", "fee_total_pct": "Frais totaux %",
                    }
                ).to_excel(_bond_xlsx, index=False, engine="openpyxl")
                st.download_button(
                    f"⬇️ Télécharger la liste filtrée ({len(bond_df)} fonds)",
                    data=_bond_xlsx.getvalue(),
                    file_name=f"obligataires_{contract_label.replace(' ', '_')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"bond_dl_{port_key}",
                )
                # Sélection directe par selectbox
                bond_options = [
                    f"{row['name']}  |  {row['manager']}  |  SRI {row['sri']}  |  {row['fee_total_pct']:.2f}%"
                    for _, row in bond_df.iterrows()
                ]
                bond_sel_idx = st.selectbox(
                    f"{len(bond_df)} fonds obligataires disponibles — sélectionner",
                    options=range(len(bond_options)),
                    format_func=lambda i: bond_options[i],
                    key=f"bond_sel_{port_key}",
                )
                bond_row = bond_df.iloc[bond_sel_idx]
                with st.container(border=True):
                    bc1, bc2 = st.columns([3, 1])
                    with bc1:
                        st.markdown(f"**{bond_row['name']}**")
                        st.caption(
                            f"{bond_row['manager']}  ·  {bond_row['category']}"
                            f"  ·  ISIN : `{bond_row['isin']}`"
                        )
                    with bc2:
                        st.metric("Frais totaux", f"{bond_row['fee_total_pct']:.2f}%")
                bb1, bb2, bb3 = st.columns([2, 2, 1])
                with bb1:
                    bond_amt = st.text_input(
                        "Montant (€)", value="", key=f"bond_amt_{port_key}",
                        placeholder="Ex : 10000",
                    )
                with bb2:
                    bond_px = st.text_input(
                        "Prix d'achat (optionnel)", value="", key=f"bond_px_{port_key}",
                    )
                with bb3:
                    if st.button("➕ Ajouter", key=f"bond_add_{port_key}"):
                        try:
                            amt = float(str(bond_amt).replace(" ", "").replace(",", "."))
                            assert amt > 0
                        except Exception:
                            st.warning("Montant invalide.")
                        else:
                            st.session_state[port_key].append({
                                "id":               str(uuid.uuid4()),
                                "name":             bond_row["name"],
                                "isin":             bond_row["isin"],
                                "amount_gross":     float(amt),
                                "buy_date":         pd.Timestamp(buy_date_central),
                                "buy_px":           float(str(bond_px).replace(",", ".")) if bond_px.strip() else "",
                                "note":             f"Contrat : {contract_label} | SRI {bond_row['sri']}",
                                "sym_used":         "",
                                "fee_uc_pct":       bond_row["fee_uc_pct"],
                                "fee_contract_pct": bond_row["fee_contract_pct"],
                                "fee_total_pct":    bond_row["fee_total_pct"],
                            })
                            st.toast(f"✅ {bond_row['name']} ajouté.", icon="✅")
                            st.rerun()

    # ═══════════════════════════════════════════════════════════
    # BLOC B — Autres UC
    # ═══════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("### 📈 Autres UC — Actions, diversifiés, etc.")

    if funds_df.empty:
        st.info("Référentiel de fonds non chargé.")
        return

    other_df = funds_df[~funds_df["category"].apply(_is_bond_category)].copy()

    if other_df.empty:
        st.info("Aucun fonds UC dans le référentiel.")
        return

    oc1, oc2, oc3 = st.columns([3, 2, 1])
    with oc1:
        other_search = st.text_input(
            "Rechercher (nom, ISIN)",
            value="", key=f"other_search_{port_key}",
            placeholder="Ex : Carmignac, FR0010135103...",
        )
    with oc2:
        other_cats = ["Toutes"] + sorted(other_df["category"].dropna().unique().tolist())
        other_cat = st.selectbox("Catégorie", other_cats, key=f"other_cat_{port_key}")
    with oc3:
        other_sri = st.selectbox("SRI", ["Tous"] + [str(i) for i in range(1, 8)], key=f"other_sri_{port_key}")

    if other_search.strip():
        q = other_search.strip().upper()
        other_df = other_df[
            other_df["isin"].str.upper().str.contains(q, na=False) |
            other_df["name"].str.upper().str.contains(q, na=False)
        ]
    if other_cat != "Toutes":
        other_df = other_df[other_df["category"] == other_cat]
    if other_sri != "Tous":
        other_df = other_df[other_df["sri"] == int(other_sri)]

    if other_df.empty:
        st.info("Aucun fonds ne correspond aux critères.")
        return

    # Téléchargement liste filtrée
    _other_xlsx = BytesIO()
    other_df[["isin", "name", "manager", "category", "sri", "fee_total_pct"]].rename(
        columns={
            "isin": "ISIN", "name": "Libellé", "manager": "Société",
            "category": "Catégorie", "sri": "SRI", "fee_total_pct": "Frais totaux %",
        }
    ).to_excel(_other_xlsx, index=False, engine="openpyxl")
    st.download_button(
        f"⬇️ Télécharger la liste filtrée ({len(other_df)} UC)",
        data=_other_xlsx.getvalue(),
        file_name=f"uc_{contract_label.replace(' ', '_')}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key=f"other_dl_{port_key}",
    )
    # Sélection directe par selectbox
    other_options = [
        f"{row['name']}  |  {row['manager']}  |  SRI {row['sri']}  |  {row['fee_total_pct']:.2f}%"
        for _, row in other_df.iterrows()
    ]
    other_sel_idx = st.selectbox(
        f"{len(other_df)} UC disponibles — sélectionner",
        options=range(len(other_options)),
        format_func=lambda i: other_options[i],
        key=f"other_sel_{port_key}",
    )
    other_row = other_df.iloc[other_sel_idx]
    with st.container(border=True):
        oc1, oc2 = st.columns([3, 1])
        with oc1:
            st.markdown(f"**{other_row['name']}**")
            st.caption(
                f"{other_row['manager']}  ·  {other_row['category']}"
                f"  ·  ISIN : `{other_row['isin']}`"
            )
        with oc2:
            st.metric("Frais totaux", f"{other_row['fee_total_pct']:.2f}%")
    ob1, ob2, ob3 = st.columns([2, 2, 1])
    with ob1:
        other_amt = st.text_input(
            "Montant (€)", value="", key=f"other_amt_{port_key}",
            placeholder="Ex : 10000",
        )
    with ob2:
        other_px = st.text_input(
            "Prix d'achat (optionnel)", value="", key=f"other_px_{port_key}",
        )
    with ob3:
        if st.button("➕ Ajouter", key=f"other_add_{port_key}"):
            try:
                amt = float(str(other_amt).replace(" ", "").replace(",", "."))
                assert amt > 0
            except Exception:
                st.warning("Montant invalide.")
            else:
                st.session_state[port_key].append({
                    "id":               str(uuid.uuid4()),
                    "name":             other_row["name"],
                    "isin":             other_row["isin"],
                    "amount_gross":     float(amt),
                    "buy_date":         pd.Timestamp(buy_date_central),
                    "buy_px":           float(str(other_px).replace(",", ".")) if other_px.strip() else "",
                    "note":             f"Contrat : {contract_label} | SRI {other_row['sri']}",
                    "sym_used":         "",
                    "fee_uc_pct":       other_row["fee_uc_pct"],
                    "fee_contract_pct": other_row["fee_contract_pct"],
                    "fee_total_pct":    other_row["fee_total_pct"],
                })
                st.toast(f"✅ {other_row['name']} ajouté.", icon="✅")
                st.rerun()


def _build_onepager_pdf(report: Dict[str, Any]) -> bytes:
    """
    Génère un PDF one-pager 'Proposition d'Arbitrage' prêt à présenter.
    Inclut : état actuel vs cible, gain estimé, économie frais/fiscalité.
    """
    if not REPORTLAB_AVAILABLE:
        return b""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor, white, black
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    )
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from io import BytesIO
    NAVY   = HexColor("#1B2A4A")
    GOLD   = HexColor("#C9A84C")
    STEEL  = HexColor("#4A6FA5")
    LIGHT  = HexColor("#F5F7FA")
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=1.8*cm, rightMargin=1.8*cm,
        topMargin=1.5*cm, bottomMargin=1.5*cm,
    )
    styles = getSampleStyleSheet()
    story = []
    def _h(txt, color=NAVY, size=16):
        return Paragraph(
            f'<font color="#{color.hexval()[2:]}" size="{size}"><b>{txt}</b></font>',
            styles["Normal"]
        )
    def _p(txt, size=9):
        return Paragraph(f'<font size="{size}">{txt}</font>', styles["Normal"])
    def _kpi(label, value, delta=None):
        val_str = f'<font size="14" color="#1B2A4A"><b>{value}</b></font>'
        lbl_str = f'<font size="8" color="#666666">{label}</font>'
        d_str = f'<br/><font size="9" color="#2E86AB">{delta}</font>' if delta else ""
        return Paragraph(f'{lbl_str}<br/>{val_str}{d_str}', styles["Normal"])
    # ── En-tête ──────────────────────────────────────────────────────
    nom_client = report.get("nom_client", "Client") or "Client"
    nom_cabinet = report.get("nom_cabinet", "Cabinet") or "Cabinet"
    as_of = report.get("as_of", "")
    header_data = [[
        Paragraph(f'<font color="white" size="18"><b>PROPOSITION D\'ARBITRAGE</b></font>'
                  f'<br/><font color="white" size="9">{nom_cabinet}</font>',
                  styles["Normal"]),
        Paragraph(f'<font color="white" size="10"><b>{nom_client}</b></font>'
                  f'<br/><font color="white" size="8">Au {as_of}</font>',
                  styles["Normal"]),
    ]]
    header_tbl = Table(header_data, colWidths=[11*cm, 6*cm])
    header_tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,-1), NAVY),
        ("TEXTCOLOR",   (0,0), (-1,-1), white),
        ("ALIGN",       (1,0), (1,0),   "RIGHT"),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",  (0,0), (-1,-1), 12),
        ("BOTTOMPADDING",(0,0),(-1,-1), 12),
        ("LEFTPADDING", (0,0), (-1,-1), 12),
        ("ROUNDEDCORNERS", [4]),
    ]))
    story.append(header_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 1 : État actuel vs Cible ────────────────────────────────
    story.append(_h("① État actuel vs Portefeuille cible", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    synth_a = report.get("client_summary", {})
    synth_b = report.get("valority_summary", {})
    comp    = report.get("comparison", {})
    def _eur(v):
        try:
            return f"{float(v):,.0f} €".replace(",", " ")
        except Exception:
            return "—"
    def _pct(v):
        try:
            return f"{float(v):.2f} %"
        except Exception:
            return "—"
    compare_data = [
        ["", "Portefeuille actuel (Client)", "Portefeuille cible (Cabinet)"],
        ["Valeur actuelle",
         _eur(synth_a.get("val", 0)), _eur(synth_b.get("val", 0))],
        ["Versements nets",
         _eur(synth_a.get("net", 0)), _eur(synth_b.get("net", 0))],
        ["Gain net",
         _eur(synth_a.get("val", 0) - synth_a.get("net", 0)),
         _eur(synth_b.get("val", 0) - synth_b.get("net", 0))],
        ["XIRR (annualisé)",
         _pct(synth_a.get("irr_pct", 0)), _pct(synth_b.get("irr_pct", 0))],
        ["Performance totale",
         _pct(synth_a.get("perf_tot_pct", 0)), _pct(synth_b.get("perf_tot_pct", 0))],
    ]
    cmp_tbl = Table(compare_data, colWidths=[5*cm, 5.5*cm, 5.5*cm])
    cmp_tbl.setStyle(TableStyle([
        ("BACKGROUND",   (0,0), (-1,0),  NAVY),
        ("TEXTCOLOR",    (0,0), (-1,0),  white),
        ("BACKGROUND",   (0,1), (-1,-1), LIGHT),
        ("BACKGROUND",   (2,1), (2,-1),  HexColor("#EAF4FB")),
        ("FONTSIZE",     (0,0), (-1,-1), 8),
        ("FONTNAME",     (0,0), (-1,0),  "Helvetica-Bold"),
        ("ALIGN",        (1,0), (-1,-1), "CENTER"),
        ("ROWBACKGROUNDS",(0,1),(-1,-1), [LIGHT, white]),
        ("GRID",         (0,0), (-1,-1), 0.3, HexColor("#CCCCCC")),
        ("TOPPADDING",   (0,0), (-1,-1), 5),
        ("BOTTOMPADDING",(0,0),(-1,-1),  5),
    ]))
    story.append(cmp_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 2 : Gain de performance estimé ──────────────────────────
    story.append(_h("② Gain de performance net estimé", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    delta_val  = comp.get("delta_val", 0.0)
    delta_xirr = comp.get("delta_perf_pct", 0.0)
    kpi_data = [[
        _kpi("Gain de valeur estimé", _eur(delta_val),
             "Valeur Cabinet − Valeur Client"),
        _kpi("Gain XIRR annualisé", _pct(delta_xirr),
             "Performance supplémentaire annualisée"),
        _kpi("Horizon recommandé", "Long terme (≥8 ans)",
             "Pour bénéficier de l'abattement AV"),
    ]]
    kpi_tbl = Table(kpi_data, colWidths=[5.5*cm, 5.5*cm, 5.5*cm])
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0), (-1,-1), LIGHT),
        ("ALIGN",         (0,0), (-1,-1), "CENTER"),
        ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0), (-1,-1), 10),
        ("BOTTOMPADDING", (0,0), (-1,-1), 10),
        ("BOX",           (0,0), (-1,-1), 1, GOLD),
        ("INNERGRID",     (0,0), (-1,-1), 0.5, HexColor("#CCCCCC")),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 14))
    # ── Bloc 3 : Économie frais / fiscalité ──────────────────────────
    story.append(_h("③ Économie de frais & levier fiscal", NAVY, 12))
    story.append(HRFlowable(width="100%", thickness=1, color=GOLD, spaceAfter=6))
    fees = report.get("fees_analysis", {})
    val_a = float(synth_a.get("val", 0) or 0)
    val_b = float(synth_b.get("val", 0) or 0)
    fee_a_pct = float(report.get("fee_a_pct", 0.6))
    fee_b_pct = float(report.get("fee_b_pct", 0.6))
    econo_frais_an = val_a * max(0.0, fee_a_pct - fee_b_pct) / 100.0
    eco_data = [
        ["Levier", "Impact estimé", "Note"],
        ["Économie frais contrat/an",
         _eur(econo_frais_an),
         f"Si passage de {fee_a_pct:.2f}% → {fee_b_pct:.2f}%/an"],
        ["Abattement AV (célibataire, ≥8 ans)",
         "4 600 €/an d'IR évité",
         "Sur la quote-part de gains dans le rachat"],
        ["Abattement AV (couple, ≥8 ans)",
         "9 200 €/an d'IR évité",
         "Sur la quote-part de gains dans le rachat"],
        ["Transmission (Art. 990I)",
         "152 500 € exonérés/bénéficiaire",
         "Pour versements avant 70 ans"],
    ]
    eco_tbl = Table(eco_data, colWidths=[5.5*cm, 4.5*cm, 7*cm])
    eco_tbl.setStyle(TableStyle([
        ("BACKGROUND",   (0,0), (-1,0),  STEEL),
        ("TEXTCOLOR",    (0,0), (-1,0),  white),
        ("BACKGROUND",   (0,1), (-1,-1), LIGHT),
        ("ROWBACKGROUNDS",(0,1),(-1,-1), [LIGHT, white]),
        ("FONTSIZE",     (0,0), (-1,-1), 8),
        ("FONTNAME",     (0,0), (-1,0),  "Helvetica-Bold"),
        ("GRID",         (0,0), (-1,-1), 0.3, HexColor("#CCCCCC")),
        ("TOPPADDING",   (0,0), (-1,-1), 5),
        ("BOTTOMPADDING",(0,0),(-1,-1),  5),
    ]))
    story.append(eco_tbl)
    story.append(Spacer(1, 14))
    # ── Pied de page légal ───────────────────────────────────────────
    story.append(HRFlowable(width="100%", thickness=0.5, color=HexColor("#CCCCCC")))
    story.append(Spacer(1, 4))
    story.append(_p(
        "Document à titre indicatif — ne constitue pas un conseil en investissement. "
        "Performances passées ne préjugent pas des performances futures. "
        f"Généré par {nom_cabinet} • {as_of}" if nom_cabinet else f"Document confidentiel • {as_of}",
        size=7,
    ))
    doc.build(story)
    return buf.getvalue()


# ── Helpers d'interprétation des ratios (module level) ──────────────────────

def _interpret_sharpe(val: Optional[float]) -> Tuple[str, str, str]:
    """Retourne (emoji_couleur, verdict_court, phrase_client)."""
    if val is None:
        return "⚪", "Données insuffisantes", "Pas assez d'historique pour calculer ce ratio."
    if val >= 1.0:
        return "🟢", "Excellent", f"Pour chaque unité de risque prise, le portefeuille génère {val:.2f} unité de rendement excédentaire. C'est un très bon ratio."
    if val >= 0.5:
        return "🟢", "Bon", f"Le portefeuille offre un bon compromis rendement/risque (Sharpe {val:.2f})."
    if val >= 0.0:
        return "🟠", "Moyen", f"Le rendement compense à peine le risque pris (Sharpe {val:.2f}). Une optimisation est possible."
    return "🔴", "Négatif", f"Le portefeuille perd de l'argent par rapport au fonds euros après ajustement du risque (Sharpe {val:.2f})."


def _interpret_sharpe_short(val: Optional[float]) -> str:
    """Verdict court pour st.metric delta."""
    if val is None:
        return "—"
    if val >= 1.0:
        return "🟢 Excellent"
    if val >= 0.5:
        return "🟢 Bon"
    if val >= 0.0:
        return "🟠 Moyen"
    return "🔴 Négatif"


def _interpret_sortino_short(val: Optional[float]) -> str:
    if val is None:
        return "—"
    if val >= 1.5:
        return "🟢 Excellente protection"
    if val >= 0.8:
        return "🟢 Bonne protection"
    if val >= 0.0:
        return "🟠 Protection limitée"
    return "🔴 Vulnérable"


def _interpret_beta_short(val: Optional[float]) -> str:
    if val is None:
        return "—"
    if val < 0.5:
        return "🟢 Très défensif"
    if val < 0.8:
        return "🟢 Défensif"
    if val < 1.2:
        return "🟠 Neutre (marché)"
    return "🔴 Agressif"


def _interpret_sortino(val: Optional[float]) -> Tuple[str, str, str]:
    if val is None:
        return "⚪", "—", ""
    if val >= 1.5:
        return "🟢", "Excellente protection", f"Le Sortino de {val:.2f} montre que le portefeuille limite très bien les baisses."
    if val >= 0.8:
        return "🟢", "Bonne protection", f"Sortino de {val:.2f} : les pertes sont contenues par rapport aux gains."
    if val >= 0.0:
        return "🟠", "Protection limitée", f"Sortino de {val:.2f} : le portefeuille est vulnérable en cas de baisse."
    return "🔴", "Vulnérable", f"Sortino négatif ({val:.2f}) : le portefeuille souffre particulièrement dans les phases baissières."


def _interpret_beta(val: Optional[float]) -> Tuple[str, str, str]:
    if val is None:
        return "⚪", "—", ""
    if val < 0.5:
        return "🟢", "Très défensif", f"Bêta de {val:.2f} : le portefeuille ne bouge que de {val*100:.0f}% quand le marché bouge de 100%. Très protecteur."
    if val < 0.8:
        return "🟢", "Défensif", f"Bêta de {val:.2f} : le portefeuille amortit les mouvements du marché."
    if val < 1.2:
        return "🟠", "Neutre", f"Bêta de {val:.2f} : le portefeuille suit globalement le marché."
    return "🔴", "Agressif", f"Bêta de {val:.2f} : le portefeuille amplifie les mouvements du marché. Plus de potentiel mais plus de risque."


def _render_ratios_card(
    sharpe: Optional[float],
    sortino: Optional[float],
    beta_alpha: Optional[Dict[str, Any]],
    euro_rate: float,
) -> None:
    """Affiche les 3 ratios dans un container bordé avec interprétation colorée."""
    with st.container(border=True):
        # ── Sharpe ──
        _sh_emoji, _sh_verdict, _sh_phrase = _interpret_sharpe(sharpe)
        st.metric(
            f"{_sh_emoji} Ratio de Sharpe",
            f"{sharpe:.2f}" if sharpe is not None else "—",
            delta=_sh_verdict,
        )
        st.caption(_sh_phrase)

        st.markdown("---")

        # ── Sortino ──
        _so_emoji, _so_verdict, _so_phrase = _interpret_sortino(sortino)
        st.metric(
            f"{_so_emoji} Ratio de Sortino",
            f"{sortino:.2f}" if sortino is not None else "—",
            delta=_so_verdict,
        )
        if _so_phrase:
            st.caption(_so_phrase)

        st.markdown("---")

        # ── Bêta ──
        if beta_alpha is not None:
            _b = beta_alpha["beta"]
            _a = beta_alpha["alpha_pct"]
            _bn = beta_alpha["benchmark_name"]
            _be_emoji, _be_verdict, _be_phrase = _interpret_beta(_b)
            st.metric(
                f"{_be_emoji} Bêta (vs {_bn})",
                f"{_b:.2f}",
                delta=_be_verdict,
            )
            st.caption(_be_phrase)
            if _a is not None:
                _alpha_color = "🟢" if _a > 0 else "🔴"
                st.caption(
                    f"{_alpha_color} Alpha de Jensen : {_a:+.2f}%/an — "
                    + ("le portefeuille crée de la valeur au-delà du marché."
                       if _a > 0
                       else "le portefeuille sous-performe par rapport à ce que son niveau de risque devrait générer.")
                )
        else:
            st.metric("⚪ Bêta", "—")
            st.caption("Indice de référence indisponible ou historique insuffisant.")

        # ── Note méthodologique ──
        st.caption(
            f"_Taux sans risque utilisé : {euro_rate:.2f}% (taux fonds euros du contrat). "
            f"Calculs sur l'historique complet des VL disponibles._"
        )


def render_app(run_page_config: bool = True):
    # ------------------------------------------------------------
    # Layout principal
    # ------------------------------------------------------------
    if run_page_config:
        st.set_page_config(page_title=APP_TITLE, layout="wide")
    _nom_cab_display = st.session_state.get("NOM_CABINET", "").strip()
    st.title(_nom_cab_display if _nom_cab_display else APP_TITLE)
    st.caption(APP_SUBTITLE)
    _nom = st.session_state.get("NOM_CLIENT", "").strip()
    if _nom:
        st.markdown(f"**Dossier client :** {_nom}")
    # Init state
    st.session_state.setdefault("NOM_CLIENT", "")
    st.session_state.setdefault("NOM_CABINET", "")
    st.session_state.setdefault("A_lines", [])
    st.session_state.setdefault("B_lines", [])
    st.session_state.setdefault("FEE_A", 3.0)
    st.session_state.setdefault("FEE_B", 2.0)
    st.session_state.setdefault("M_A", 0.0)
    st.session_state.setdefault("M_B", 0.0)
    st.session_state.setdefault("ONE_A", 0.0)
    st.session_state.setdefault("ONE_B", 0.0)
    st.session_state.setdefault("ONE_A_DATE", pd.Timestamp("2024-07-01").date())
    st.session_state.setdefault("ONE_B_DATE", pd.Timestamp("2024-07-01").date())
    st.session_state.setdefault("ALLOC_MODE", "equal")
    st.session_state.setdefault("MANUAL_NAV_STORE", {})
    st.session_state.setdefault("DATE_WARNINGS", [])
    st.session_state.setdefault("INIT_A_DATE", pd.Timestamp("2024-01-02").date())
    st.session_state.setdefault("INIT_B_DATE", pd.Timestamp("2024-01-02").date())
    st.session_state.setdefault("EURO_RATE_A", 2.0)
    st.session_state.setdefault("EURO_RATE_B", 2.5)

    # -------------------------------------------------------------------
    # Sidebar : paramètres globaux
    # -------------------------------------------------------------------
    with st.sidebar:
        st.markdown("---")
        # ── Nom du client et du cabinet ────────────────────────────
        st.session_state.setdefault("NOM_CLIENT", "")
        NOM_CLIENT = st.text_input(
            "Nom du client",
            value=st.session_state.get("NOM_CLIENT", ""),
            key="NOM_CLIENT",
            placeholder="Ex : M. et Mme Dupont",
        )
        st.session_state.setdefault("NOM_CABINET", "")
        NOM_CABINET = st.text_input(
            "Nom du cabinet",
            value=st.session_state.get("NOM_CABINET", ""),
            key="NOM_CABINET",
            placeholder="Ex : Cabinet Dupont Patrimoine",
        )
        st.divider()

        # ── Section 1 — Contrats ────────────────────────────────────
        with st.expander("① Contrats", expanded=True):
            st.markdown("**Portefeuille Client**")
            contract_label_A = st.selectbox(
                "Contrat client",
                list(CONTRACTS_REGISTRY.keys()),
                key="CONTRACT_LABEL_A",
            )
            contract_cfg_A = CONTRACTS_REGISTRY[contract_label_A]
            try:
                funds_df_A = load_contract_funds(
                    contract_cfg_A["path"], contract_cfg_A["funds_filename"]
                )
            except Exception:
                funds_df_A = pd.DataFrame()
            st.session_state["CONTRACT_FUNDS_DF_A"] = funds_df_A

            euro_fund_label_A = st.selectbox(
                "Fonds en euros client",
                list(contract_cfg_A["euro_funds"].keys()),
                key="EURO_FUND_LABEL_A",
            )
            try:
                euro_history_df_A = load_euro_fund_history(
                    contract_cfg_A["path"], contract_cfg_A["euro_funds"][euro_fund_label_A]
                )
            except Exception:
                euro_history_df_A = pd.DataFrame()
            avg_rate_A = get_euro_fund_avg_rate(euro_history_df_A, years=5)
            st.caption(
                f"Fonds euros client — moyenne 5 ans : **{avg_rate_A:.2f}%** "
                f"({euro_fund_label_A})"
            )

            st.markdown("---")
            st.markdown("**Portefeuille Cabinet**")
            contract_label_B = st.selectbox(
                "Contrat cabinet",
                list(CONTRACTS_REGISTRY.keys()),
                key="CONTRACT_LABEL_B",
            )
            contract_cfg_B = CONTRACTS_REGISTRY[contract_label_B]
            try:
                funds_df_B = load_contract_funds(
                    contract_cfg_B["path"], contract_cfg_B["funds_filename"]
                )
            except Exception:
                funds_df_B = pd.DataFrame()
            st.session_state["CONTRACT_FUNDS_DF_B"] = funds_df_B

            euro_fund_label_B = st.selectbox(
                "Fonds en euros cabinet",
                list(contract_cfg_B["euro_funds"].keys()),
                key="EURO_FUND_LABEL_B",
            )
            try:
                euro_history_df_B = load_euro_fund_history(
                    contract_cfg_B["path"], contract_cfg_B["euro_funds"][euro_fund_label_B]
                )
            except Exception:
                euro_history_df_B = pd.DataFrame()
            avg_rate_B = get_euro_fund_avg_rate(euro_history_df_B, years=5)
            st.caption(
                f"Fonds euros cabinet — moyenne 5 ans : **{avg_rate_B:.2f}%** "
                f"({euro_fund_label_B})"
            )

            # Rétrocompatibilité : render_portfolio_builder() lit CONTRACT_LABEL / CONTRACT_FUNDS_DF
            st.session_state["CONTRACT_LABEL"]    = contract_label_A
            st.session_state["CONTRACT_CFG"]      = contract_cfg_A
            st.session_state["CONTRACT_FUNDS_DF"] = funds_df_A
            st.session_state["EURO_FUND_LABEL"]   = euro_fund_label_A
            st.session_state["EURO_FUND_HISTORY"] = euro_history_df_A
            st.session_state["EURO_FUND_AVG_RATE"] = avg_rate_A

        # ── Section 2 — Paramètres de simulation ───────────────────
        with st.expander("② Paramètres de simulation", expanded=True):
            st.markdown("**Taux fonds en euros**")
            st.caption(
                f"Client : **{avg_rate_A:.2f}%** ({euro_fund_label_A}) | "
                f"Cabinet : **{avg_rate_B:.2f}%** ({euro_fund_label_B})"
            )
            override_euro = st.checkbox(
                "Utiliser un taux personnalisé (fonds euros boosté, etc.)",
                value=False,
                key="OVERRIDE_EURO_RATE",
            )
            if override_euro:
                euro_rate_value_A = st.number_input(
                    "Portefeuille Client — taux annuel (%)",
                    min_value=0.0, max_value=10.0,
                    value=float(st.session_state.get("EURO_RATE_A", avg_rate_A)),
                    step=0.10, key="EURO_RATE_A",
                )
                euro_rate_value_B = st.number_input(
                    "Portefeuille Cabinet — taux annuel (%)",
                    min_value=0.0, max_value=10.0,
                    value=float(st.session_state.get("EURO_RATE_B", avg_rate_B)),
                    step=0.10, key="EURO_RATE_B",
                )
            else:
                euro_rate_value_A = avg_rate_A
                euro_rate_value_B = avg_rate_B
                st.session_state["EURO_RATE_A"] = avg_rate_A
                st.session_state["EURO_RATE_B"] = avg_rate_B
                if avg_rate_A == avg_rate_B:
                    st.caption(f"Taux appliqué aux deux portefeuilles : **{avg_rate_A:.2f}%**")
                else:
                    st.caption(
                        f"Taux client : **{avg_rate_A:.2f}%** | Taux cabinet : **{avg_rate_B:.2f}%**"
                    )
            st.caption(
                "Frais de gestion contrat fonds euros inclus dans le taux net publié. "
                "Frais UC contrat : lus depuis le référentiel fonds."
            )

            st.markdown("---")
            st.markdown("**Frais d'entrée (%)**")
            FEE_A = st.number_input(
                "Frais d'entrée — Portefeuille 1 (Client)",
                0.0,
                10.0,
                st.session_state.get("FEE_A", 3.0),
                0.10,
                key="FEE_A",
            )
            FEE_B = st.number_input(
                "Frais d'entrée — Portefeuille 2 (Cabinet)",
                0.0,
                10.0,
                st.session_state.get("FEE_B", 2.0),
                0.10,
                key="FEE_B",
            )
            st.caption("Les frais s'appliquent sur chaque investissement (initial, mensuel, ponctuel).")

            st.markdown("---")
            st.markdown("**Date du versement initial**")
            st.date_input(
                "Portefeuille 1 (Client) — date d'investissement initiale",
                value=st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date()),
                key="INIT_A_DATE",
            )
            st.date_input(
                "Portefeuille 2 (Cabinet) — date d'investissement initiale",
                value=st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date()),
                key="INIT_B_DATE",
            )

        # ── Section 3 — Versements ──────────────────────────────────
        with st.expander("③ Versements", expanded=False):
            st.markdown("**Portefeuille 1 — Client**")
            M_A = st.number_input(
                "Mensuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("M_A", 0.0),
                100.0,
                key="M_A",
            )
            ONE_A = st.number_input(
                "Ponctuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("ONE_A", 0.0),
                100.0,
                key="ONE_A",
            )
            ONE_A_DATE = st.date_input(
                "Date du ponctuel",
                value=st.session_state.get("ONE_A_DATE", pd.Timestamp("2024-07-01").date()),
                key="ONE_A_DATE",
            )
            st.markdown("---")
            st.markdown("**Portefeuille 2 — Cabinet**")
            M_B = st.number_input(
                "Mensuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("M_B", 0.0),
                100.0,
                key="M_B",
            )
            ONE_B = st.number_input(
                "Ponctuel brut (€)",
                0.0,
                1_000_000.0,
                st.session_state.get("ONE_B", 0.0),
                100.0,
                key="ONE_B",
            )
            ONE_B_DATE = st.date_input(
                "Date du ponctuel",
                value=st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date()),
                key="ONE_B_DATE",
            )

        # ── Section 4 — Options d'analyse ──────────────────────────
        with st.expander("④ Options d'analyse", expanded=False):
            st.markdown("**Règle d'affectation des versements**")
            current_code = st.session_state.get("ALLOC_MODE", "equal")
            inv_labels = {v: k for k, v in ALLOC_LABELS.items()}
            current_label = inv_labels.get(current_code, "Répartition égale")

            mode_label = st.selectbox(
                "Mode",
                list(ALLOC_LABELS.keys()),
                index=list(ALLOC_LABELS.keys()).index(current_label),
                help="Répartition des versements entre les lignes.",
            )
            st.session_state["ALLOC_MODE"] = ALLOC_LABELS[mode_label]

            st.markdown("---")
            st.markdown("**Mode d'analyse**")
            mode_ui = st.radio(
                "Choix",
                ["Comparer Client vs Cabinet", "Analyser uniquement Cabinet", "Analyser uniquement Client"],
                index=0,
                key="MODE_ANALYSE_UI",
            )
            if "Comparer" in mode_ui:
                st.session_state["MODE_ANALYSE"] = "compare"
            elif "Cabinet" in mode_ui:
                st.session_state["MODE_ANALYSE"] = "valority"
            else:
                st.session_state["MODE_ANALYSE"] = "client"

            st.markdown("---")
            st.markdown("**Indice de référence (Bêta)**")
            _bench_options = {
                "MSCI World (CW8.PA)": "CW8.PA",
                "CAC 40 (^FCHI)": "^FCHI",
                "Euro Stoxx 50 (^STOXX50E)": "^STOXX50E",
                "S&P 500 (^GSPC)": "^GSPC",
            }
            _bench_label = st.selectbox(
                "Indice",
                list(_bench_options.keys()),
                index=0,
                key="BENCHMARK_LABEL",
            )
            st.session_state["BENCHMARK_SYMBOL"] = _bench_options[_bench_label]

        # ── Debug (masqué en production) ────────────────────────────
        _debug_enabled = st.secrets.get("debug_mode", False)
        if _debug_enabled:
            with st.expander("🔧 Debug", expanded=False):
                st.caption("Versions & état")
                st.code(
                    f"Python: {sys.version.split()[0]}\n"
                    f"Streamlit: {st.__version__}\n"
                    f"Pandas: {pd.__version__}"
                )
                st.caption("Modules")
                st.code(
                    f"Matplotlib: {MATPLOTLIB_AVAILABLE} ({MATPLOTLIB_ERROR})\n"
                    f"Reportlab: {REPORTLAB_AVAILABLE} ({REPORTLAB_ERROR})"
                )
                st.caption("Session state (clés)")
                st.write(sorted(list(st.session_state.keys())))
                st.caption("Dernière exception")
                st.write(st.session_state.get("LAST_EXCEPTION", "—"))
                st.caption("Test rapide EODHD")
                if _get_api_key():
                    try:
                        res = eodhd_get("/status")
                        st.write(res if res is not None else "Réponse vide")
                    except Exception as e:
                        st.write(f"Erreur EODHD: {e}")
                else:
                    st.write("Token EODHD absent")


    mode = st.session_state.get("MODE_ANALYSE", "compare")
    show_client = mode in ("compare", "client")
    show_valority = mode in ("compare", "valority")

    # ── Synchronisation date globale → lignes (pour cohérence carte + tableau + simulation) ──
    _global_date_A = pd.Timestamp(st.session_state.get("INIT_A_DATE", pd.Timestamp("2024-01-02").date()))
    _global_date_B = pd.Timestamp(st.session_state.get("INIT_B_DATE", pd.Timestamp("2024-01-02").date()))
    for _ln in st.session_state.get("A_lines", []):
        if not _ln.get("date_overridden"):
            _ln["buy_date"] = _global_date_A
    for _ln in st.session_state.get("B_lines", []):
        if not _ln.get("date_overridden"):
            _ln["buy_date"] = _global_date_B

    # Onglets principaux : Client / Cabinet (conditionnés au mode)
    tab_labels = []
    if show_client:
        tab_labels.append("Portefeuille Client")
    if show_valority:
        tab_labels.append("Portefeuille Cabinet")
    tabs = st.tabs(tab_labels) if tab_labels else []

    idx = 0
    if show_client:
        with tabs[idx]:
            _add_fund_from_contract("A_lines", "Ajouter un fonds — Portefeuille Client")
            st.markdown("#### Lignes actuelles — Portefeuille Client")
            for i, ln in enumerate(st.session_state.get("A_lines", [])):
                _line_card(ln, i, "A_lines")
        idx += 1

    if show_valority:
        with tabs[idx]:
            _add_fund_from_contract("B_lines", "Ajouter un fonds — Portefeuille Cabinet")
            st.markdown("#### Lignes actuelles — Portefeuille Cabinet")
            for i, ln in enumerate(st.session_state.get("B_lines", [])):
                _line_card(ln, i, "B_lines")

    # ------------------------------------------------------------
    # Simulation (selon mode)
    # ------------------------------------------------------------
    # FIXED: use stable UUID key instead of id() which changes on st.rerun() (Bug 5)
    _ln_a0 = st.session_state["A_lines"][0] if (show_client and st.session_state["A_lines"]) else None
    single_target_A = (_ln_a0.get("id") or id(_ln_a0)) if _ln_a0 is not None else None
    _ln_b0 = st.session_state["B_lines"][0] if (show_valority and st.session_state["B_lines"]) else None
    single_target_B = (_ln_b0.get("id") or id(_ln_b0)) if _ln_b0 is not None else None

    alloc_mode_code = st.session_state.get("ALLOC_MODE", "equal")

    custom_month_weights_A: Optional[Dict[int, float]] = None
    custom_oneoff_weights_A: Optional[Dict[int, float]] = None
    custom_month_weights_B: Optional[Dict[int, float]] = None
    custom_oneoff_weights_B: Optional[Dict[int, float]] = None

    if alloc_mode_code == "custom":
        if show_client:
            cmA = st.session_state.get("CUSTOM_M_A", {}) or {}
            coA = st.session_state.get("CUSTOM_O_A", {}) or {}
            tot_mA = sum(v for v in cmA.values() if v > 0)
            tot_oA = sum(v for v in coA.values() if v > 0)
            if tot_mA > 0:
                custom_month_weights_A = {k: v / tot_mA for k, v in cmA.items() if v > 0}
            if tot_oA > 0:
                custom_oneoff_weights_A = {k: v / tot_oA for k, v in coA.items() if v > 0}

        if show_valority:
            cmB = st.session_state.get("CUSTOM_M_B", {}) or {}
            coB = st.session_state.get("CUSTOM_O_B", {}) or {}
            tot_mB = sum(v for v in cmB.values() if v > 0)
            tot_oB = sum(v for v in coB.values() if v > 0)
            if tot_mB > 0:
                custom_month_weights_B = {k: v / tot_mB for k, v in cmB.items() if v > 0}
            if tot_oB > 0:
                custom_oneoff_weights_B = {k: v / tot_oB for k, v in coB.items() if v > 0}

    # Reset warnings avant chaque run
    st.session_state["DATE_WARNINGS"] = []

    # Valeurs par défaut (si on ne simule pas un des portefeuilles)
    dfA, brutA, netA, valA, xirrA, startA_min, fullA = pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY
    dfB, brutB, netB, valB, xirrB, startB_min, fullB = pd.DataFrame(), 0.0, 0.0, 0.0, None, TODAY, TODAY

    if show_client:
        _args_A = _make_sim_args(
            lines=st.session_state.get("A_lines", []),
            monthly_amt=st.session_state.get("M_A", 0.0),
            one_amt=st.session_state.get("ONE_A", 0.0),
            one_date=st.session_state.get("ONE_A_DATE", pd.Timestamp("2024-07-01").date()),
            alloc_mode=alloc_mode_code,
            custom_weights_monthly=custom_month_weights_A,
            custom_weights_oneoff=custom_oneoff_weights_A,
            single_target=single_target_A,
            euro_rate=st.session_state.get("EURO_RATE_A", 2.0),
            fee_pct=st.session_state.get("FEE_A", 0.0),
            label="Client",
        )
        dfA, brutA, netA, valA, xirrA, startA_min, fullA = _simulate_portfolio_cached(**_args_A)
        st.session_state["_LAST_VAL_A"] = valA
        st.session_state["_LAST_NET_A"] = netA
        st.session_state["_LAST_XIRR_A"] = xirrA

    if show_valority:
        _args_B = _make_sim_args(
            lines=st.session_state.get("B_lines", []),
            monthly_amt=st.session_state.get("M_B", 0.0),
            one_amt=st.session_state.get("ONE_B", 0.0),
            one_date=st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date()),
            alloc_mode=alloc_mode_code,
            custom_weights_monthly=custom_month_weights_B,
            custom_weights_oneoff=custom_oneoff_weights_B,
            single_target=single_target_B,
            euro_rate=st.session_state.get("EURO_RATE_B", 2.5),
            fee_pct=st.session_state.get("FEE_B", 0.0),
            label="Cabinet",
        )
        dfB, brutB, netB, valB, xirrB, startB_min, fullB = _simulate_portfolio_cached(**_args_B)
        st.session_state["_LAST_VAL_B"] = valB
        st.session_state["_LAST_NET_B"] = netB
        st.session_state["_LAST_XIRR_B"] = xirrB

    # ------------------------------------------------------------
    # Avertissements sur les dates / 1ère VL
    # ------------------------------------------------------------
    if st.session_state.get("DATE_WARNINGS"):
        with st.expander("⚠️ Problèmes d'historique / dates de VL"):
            for msg in st.session_state["DATE_WARNINGS"]:
                st.warning(msg)

    # ------------------------------------------------------------
    # Graphique (évolution des portefeuilles)
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📈 Comment s'est comporté votre portefeuille")

    # Déterminer le start_plot uniquement sur les portefeuilles affichés
    full_dates: List[pd.Timestamp] = []
    if show_client and isinstance(fullA, pd.Timestamp):
        full_dates.append(fullA)
    if show_valority and isinstance(fullB, pd.Timestamp):
        full_dates.append(fullB)

    start_plot = max(full_dates) if full_dates else TODAY

    idx = pd.bdate_range(start=start_plot, end=TODAY, freq="B")
    chart_df = pd.DataFrame(index=idx)

    if show_client and not dfA.empty:
        chart_df["Client"] = dfA.reindex(idx)["Valeur"].ffill()

    if show_valority and not dfB.empty:
        chart_df["Cabinet"] = dfB.reindex(idx)["Valeur"].ffill()

    # Passage en format long pour Altair
    chart_long = chart_df.reset_index().rename(columns={"index": "Date"})
    chart_long = chart_long.melt("Date", var_name="variable", value_name="Valeur (€)")

    if chart_long.dropna().empty:
        st.info("Ajoutez des lignes et/ou vérifiez vos paramètres pour afficher le graphique.")
    else:
        # Calcul du domaine Y dynamique avec marges
        valid_vals = chart_long["Valeur (€)"].dropna()
        if not valid_vals.empty:
            y_min = float(valid_vals.min())
            y_max = float(valid_vals.max())
            padding = (y_max - y_min) * 0.05 if y_max > y_min else y_max * 0.05
            y_domain = [max(0.0, y_min - padding), y_max + padding]
        else:
            y_domain = [0, 1]

        # ── Bandeau KPIs adapté au mode ──────────────────────────────
        with st.container(border=True):
            if mode == "compare":
                _has_A = show_client and not dfA.empty and netA > 0
                _has_B = show_valority and not dfB.empty and netB > 0
                if _has_A and _has_B:
                    _kpi_cols = st.columns(4)
                    with _kpi_cols[0]:
                        st.metric(
                            "Rendement client (%/an)",
                            f"{xirrA:.2f} %" if xirrA is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _kpi_cols[1]:
                        st.metric(
                            "Rendement simulation (%/an)",
                            f"{xirrB:.2f} %" if xirrB is not None else "—",
                            delta=f"{(xirrB - xirrA):.2f} pts"
                            if (xirrA is not None and xirrB is not None) else None,
                        )
                    with _kpi_cols[2]:
                        _gain_client = (valA - netA) if netA > 0 else 0.0
                        st.metric("Gain net Client", to_eur(_gain_client))
                    with _kpi_cols[3]:
                        _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                        st.metric(
                            "Gain net Cabinet", to_eur(_gain_cabinet),
                            delta=to_eur(_gain_cabinet - _gain_client) if netA > 0 else None,
                        )
                elif _has_A:
                    _pad_l, _k0, _k1, _k2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                    with _k0:
                        st.metric(
                            "Rendement client (%/an)",
                            f"{xirrA:.2f} %" if xirrA is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _k1:
                        _gain_client = (valA - netA) if netA > 0 else 0.0
                        st.metric("Gain net", to_eur(_gain_client))
                    with _k2:
                        _perf_cli = (valA / netA - 1.0) * 100.0 if netA > 0 else 0.0
                        st.metric("Performance totale", f"{_perf_cli:+.2f}%")
                elif _has_B:
                    _pad_l, _k0, _k1, _k2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                    with _k0:
                        st.metric(
                            "Rendement simulation (%/an)",
                            f"{xirrB:.2f} %" if xirrB is not None else "—",
                            help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                        )
                    with _k1:
                        _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                        st.metric("Gain net", to_eur(_gain_cabinet))
                    with _k2:
                        _perf_val = (valB / netB - 1.0) * 100.0 if netB > 0 else 0.0
                        st.metric("Performance totale", f"{_perf_val:+.2f}%")
            elif mode == "client":
                _pad_l, _kpi_cols0, _kpi_cols1, _kpi_cols2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                _kpi_cols = [_kpi_cols0, _kpi_cols1, _kpi_cols2]
                with _kpi_cols[0]:
                    st.metric(
                        "Rendement client (%/an)",
                        f"{xirrA:.2f} %" if xirrA is not None else "—",
                        help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                    )
                with _kpi_cols[1]:
                    _gain_client = (valA - netA) if netA > 0 else 0.0
                    st.metric("Gain net", to_eur(_gain_client))
                with _kpi_cols[2]:
                    _perf_cli = (valA / netA - 1.0) * 100.0 if netA > 0 else 0.0
                    st.metric("Performance totale", f"{_perf_cli:+.2f}%")
            else:  # valority
                _pad_l, _kpi_cols0, _kpi_cols1, _kpi_cols2, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                _kpi_cols = [_kpi_cols0, _kpi_cols1, _kpi_cols2]
                with _kpi_cols[0]:
                    st.metric(
                        "Rendement simulation (%/an)",
                        f"{xirrB:.2f} %" if xirrB is not None else "—",
                        help="Rendement annualisé (XIRR), net des frais de gestion du contrat",
                    )
                with _kpi_cols[1]:
                    _gain_cabinet = (valB - netB) if netB > 0 else 0.0
                    st.metric("Gain net", to_eur(_gain_cabinet))
                with _kpi_cols[2]:
                    _perf_val = (valB / netB - 1.0) * 100.0 if netB > 0 else 0.0
                    st.metric("Performance totale", f"{_perf_val:+.2f}%")

        # ── Graphique line chart (rendu direct, sans placeholder) ────
        base = (
            alt.Chart(chart_long)
            .mark_line()
            .encode(
                x=alt.X("Date:T", title="Date"),
                y=alt.Y(
                    "Valeur (€):Q",
                    title="Valeur (€)",
                    scale=alt.Scale(domain=y_domain),
                    axis=alt.Axis(format=",.0f"),
                ),
                color=alt.Color("variable:N", title="Portefeuille"),
                tooltip=[
                    alt.Tooltip("Date:T", title="Date"),
                    alt.Tooltip("variable:N", title="Portefeuille"),
                    alt.Tooltip("Valeur (€):Q", title="Valeur", format=",.2f"),
                ],
            )
            .properties(height=360, width="container")
        )
        st.altair_chart(base, use_container_width=True)

        # ── Bloc storytelling percutant ──────────────────────────────
        if mode == "compare" and netA > 0 and netB > 0:
            _perf_a = (valA / netA - 1.0) * 100.0
            _perf_b = (valB / netB - 1.0) * 100.0
            _gain_a = valA - netA
            _gain_b = valB - netB
            _manque = _gain_b - _gain_a
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric(
                        "Votre rendement sur la période",
                        f"{_perf_a:+.2f}%",
                        help="Performance totale nette des frais contrat",
                    )
                    st.metric(
                        "Vous avez gagné",
                        to_eur(_gain_a),
                        help="Valeur actuelle − versements nets investis",
                    )
                with _st_c2:
                    st.metric(
                        "Rendement de la proposition cabinet",
                        f"{_perf_b:+.2f}%",
                        delta=f"{(_perf_b - _perf_a):+.2f} pts vs votre situation actuelle",
                    )
                    st.metric(
                        "Avec le cabinet, vous auriez gagné",
                        to_eur(_gain_b),
                        delta=to_eur(_manque),
                        help="Manque à gagner par rapport à la proposition",
                    )
                if _manque > 0:
                    st.info(
                        f"💡 La proposition du cabinet aurait généré "
                        f"**{to_eur(_manque)} de plus** sur la même période "
                        f"et le même capital investi."
                    )
        elif mode == "client" and netA > 0:
            _perf_a = (valA / netA - 1.0) * 100.0
            _gain_a = valA - netA
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric("Votre rendement sur la période",
                               f"{_perf_a:+.2f}%")
                with _st_c2:
                    st.metric("Vous avez gagné", to_eur(_gain_a))
        elif mode == "valority" and netB > 0:
            _perf_b = (valB / netB - 1.0) * 100.0
            _gain_b = valB - netB
            with st.container(border=True):
                st.markdown("#### 📊 Ce que ces chiffres signifient concrètement")
                _st_c1, _st_c2 = st.columns(2)
                with _st_c1:
                    st.metric("Rendement de la simulation",
                               f"{_perf_b:+.2f}%")
                with _st_c2:
                    st.metric("Capital créé", to_eur(_gain_b))

    # ------------------------------------------------------------
    # Synthèse chiffrée : cartes Client / Cabinet
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("💡 Synthèse")

    mode = st.session_state.get("MODE_ANALYSE", "compare")

    # Période analysée (uniquement sur ce qui est affiché)
    period_dates: List[pd.Timestamp] = []
    if mode in ("compare", "client") and isinstance(startA_min, pd.Timestamp):
        period_dates.append(startA_min)
    if mode in ("compare", "valority") and isinstance(startB_min, pd.Timestamp):
        period_dates.append(startB_min)

    if period_dates:
        start_global = min(period_dates)
        st.caption(f"Période analysée : du **{fmt_date(start_global)}** au **{fmt_date(TODAY)}**")

    perf_tot_client = (valA / netA - 1.0) * 100.0 if (show_client and netA > 0) else None
    perf_tot_valority = (valB / netB - 1.0) * 100.0 if (show_valority and netB > 0) else None

    # ✅ 2 colonnes si compare, sinon 1 colonne (container)
    if mode == "compare":
        col_client, col_valority = st.columns(2)
    else:
        col_client = st.container()
        col_valority = st.container()

    # ----- Carte Client -----
    if mode in ("compare", "client"):
        with col_client:
            with st.container(border=True):
                st.markdown("#### 🧍 Situation actuelle — Client")
                st.metric("Valeur actuelle", to_eur(valA))
                if _nom := st.session_state.get("NOM_CLIENT", "").strip():
                    st.caption(f"Client : {_nom}")
                _s_c1, _s_c2 = st.columns(2)
                with _s_c1:
                    st.metric("Net investi", to_eur(netA),
                              help="Versements bruts − frais d'entrée")
                with _s_c2:
                    st.metric("Performance totale",
                              f"{perf_tot_client:+.2f}%" if perf_tot_client is not None else "—")
                _fee_detail_A = []
                for _ln in st.session_state.get("A_lines", []):
                    _isin = str(_ln.get("isin", "")).upper()
                    if _isin in ("EUROFUND", "STRUCTURED"):
                        continue
                    _name = _ln.get("name") or _isin
                    _fc = _ln.get("fee_contract_pct")
                    _fb = _ln.get("fee_uc_pct")
                    if _fc is not None and _fb is not None:
                        try:
                            _fee_detail_A.append(
                                f"{_name[:30]} : frais contrat {float(_fc):.2f}%/an "
                                f"(+ TER {float(_fb):.2f}%/an intégré dans la VL)"
                            )
                        except Exception:
                            pass
                _fee_note_A = (
                    "ℹ️ **Performance nette des frais de gestion du contrat assureur** "
                    f"({st.session_state.get('CONTRACT_LABEL_A', 'contrat')}). "
                    "Les frais de gestion internes des fonds (TER) sont déjà intégrés "
                    "dans les valeurs liquidatives publiées et ne sont pas déduits en plus. "
                    "À titre de comparaison, Quantalys et Morningstar affichent les "
                    "performances brutes de frais contrat."
                )
                with st.expander("ℹ️ Détail du calcul de performance", expanded=False):
                    st.markdown(_fee_note_A)
                    if _fee_detail_A:
                        st.markdown("**Frais appliqués par fonds :**")
                        for _d in _fee_detail_A:
                            st.markdown(f"- {_d}")
                    st.markdown(
                        "**Formule :** XIRR calculé sur les flux réels "
                        "(versements initiaux, mensuels, ponctuels) et la valeur "
                        "liquidative nette à ce jour."
                    )


    # ----- Carte Cabinet -----
    if mode in ("compare", "valority"):
        with col_valority:
            with st.container(border=True):
                st.markdown("#### 🏢 Simulation — Allocation Cabinet")
                st.metric("Valeur actuelle simulée", to_eur(valB))
                _s_c1b, _s_c2b = st.columns(2)
                with _s_c1b:
                    st.metric("Net investi", to_eur(netB),
                              help="Versements bruts − frais d'entrée")
                with _s_c2b:
                    st.metric("Performance totale",
                              f"{perf_tot_valority:+.2f}%" if perf_tot_valority is not None else "—")
                _fee_detail_B = []
                for _ln in st.session_state.get("B_lines", []):
                    _isin = str(_ln.get("isin", "")).upper()
                    if _isin in ("EUROFUND", "STRUCTURED"):
                        continue
                    _name = _ln.get("name") or _isin
                    _fc = _ln.get("fee_contract_pct")
                    _fb = _ln.get("fee_uc_pct")
                    if _fc is not None and _fb is not None:
                        try:
                            _fee_detail_B.append(
                                f"{_name[:30]} : frais contrat {float(_fc):.2f}%/an "
                                f"(+ TER {float(_fb):.2f}%/an intégré dans la VL)"
                            )
                        except Exception:
                            pass
                _fee_note_B = (
                    "ℹ️ **Performance nette des frais de gestion du contrat assureur** "
                    f"({st.session_state.get('CONTRACT_LABEL_B', 'contrat')}). "
                    "Les frais de gestion internes des fonds (TER) sont déjà intégrés "
                    "dans les valeurs liquidatives publiées et ne sont pas déduits en plus. "
                    "À titre de comparaison, Quantalys et Morningstar affichent les "
                    "performances brutes de frais contrat."
                )
                with st.expander("ℹ️ Détail du calcul de performance", expanded=False):
                    st.markdown(_fee_note_B)
                    if _fee_detail_B:
                        st.markdown("**Frais appliqués par fonds :**")
                        for _d in _fee_detail_B:
                            st.markdown(f"- {_d}")
                    st.markdown(
                        "**Formule :** XIRR calculé sur les flux réels "
                        "(versements initiaux, mensuels, ponctuels) et la valeur "
                        "liquidative nette à ce jour."
                    )


    def build_html_report(report: Dict[str, Any]) -> str:
        """
        Construit un rapport HTML exportable pour le client.
        Le contenu repose sur 'report', préparé plus bas dans le code.
        """
        as_of = report.get("as_of", "")
        mode = report.get("mode", "compare")
        synthA = report.get("client_summary", {})
        synthB = report.get("valority_summary", {})
        comp = report.get("comparison", {})

        dfA_lines = report.get("df_client_lines")
        dfB_lines = report.get("df_valority_lines")
        dfA_val = report.get("dfA_val")
        dfB_val = report.get("dfB_val")

        def _fmt_eur(x):
            try:
                return f"{x:,.2f} €".replace(",", " ").replace(".", ",")
            except Exception:
                return str(x)

        # Tables en HTML
        html_client_lines = dfA_lines.to_html(index=False, border=0, justify="left") if dfA_lines is not None else ""
        html_valority_lines = dfB_lines.to_html(index=False, border=0, justify="left") if dfB_lines is not None else ""

        if dfA_val is not None:
            html_A_val = dfA_val.to_html(index=False, border=0, justify="left")
        else:
            html_A_val = ""

        if dfB_val is not None:
            html_B_val = dfB_val.to_html(index=False, border=0, justify="left")
        else:
            html_B_val = ""

        if mode == "compare":
            synth_html = f"""
<div class="block">
  <h3>Situation actuelle — Client</h3>
  <ul>
    <li>Valeur actuelle : <b>{_fmt_eur(synthA.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthA.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthA.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthA.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthA.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>

<div class="block">
  <h3>Simulation — Allocation Cabinet</h3>
  <ul>
    <li>Valeur actuelle simulée : <b>{_fmt_eur(synthB.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthB.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthB.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthB.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthB.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>

<div class="block">
  <h3>Comparaison Client vs Cabinet</h3>
  <ul>
    <li>Différence de valeur finale : <b>{_fmt_eur(comp.get("delta_val", 0))}</b></li>
    <li>Écart de performance totale (Cabinet – Client) :
        <b>{comp.get("delta_perf_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Client</h3>
{html_client_lines}

<h3>Portefeuille Cabinet</h3>
{html_valority_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur des portefeuilles</h2>

<h3>Client – Valeur du portefeuille par date</h3>
{html_A_val}

<h3>Cabinet – Valeur du portefeuille par date</h3>
{html_B_val}
"""
        elif mode == "valority":
            synth_html = f"""
<div class="block">
  <h3>Simulation — Allocation Cabinet</h3>
  <ul>
    <li>Valeur actuelle simulée : <b>{_fmt_eur(synthB.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthB.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthB.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthB.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthB.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Cabinet</h3>
{html_valority_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur du portefeuille</h2>

<h3>Cabinet – Valeur du portefeuille par date</h3>
{html_B_val}
"""
        else:
            synth_html = f"""
<div class="block">
  <h3>Situation actuelle — Client</h3>
  <ul>
    <li>Valeur actuelle : <b>{_fmt_eur(synthA.get("val", 0))}</b></li>
    <li>Montants réellement investis (net) : {_fmt_eur(synthA.get("net", 0))}</li>
    <li>Montants versés (brut) : {_fmt_eur(synthA.get("brut", 0))}</li>
    <li>Rendement total depuis le début : <b>{synthA.get("perf_tot_pct", 0):.2f} %</b></li>
    <li>Rendement annualisé (XIRR) : <b>{synthA.get("irr_pct", 0):.2f} %</b></li>
  </ul>
</div>
"""
            detail_html = f"""
<h2>2. Détail des lignes</h2>

<h3>Portefeuille Client</h3>
{html_client_lines}
"""
            hist_html = f"""
<h2>3. Historique de la valeur du portefeuille</h2>

<h3>Client – Valeur du portefeuille par date</h3>
{html_A_val}
"""

        html = f"""
<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="utf-8" />
<title>Rapport de portefeuille</title>
<style>
body {{
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
  margin: 24px;
  color: #222;
}}
h1, h2, h3 {{
  margin-top: 24px;
}}
table {{
  border-collapse: collapse;
  width: 100%;
  margin: 8px 0 16px 0;
  font-size: 14px;
}}
th, td {{
  border: 1px solid #ddd;
  padding: 6px 8px;
}}
th {{
  background-color: #f5f5f5;
  text-align: left;
}}
.small {{
  font-size: 12px;
  color: #666;
}}
.block {{
  border: 1px solid #eee;
  border-radius: 8px;
  padding: 12px 16px;
  margin-bottom: 16px;
  background-color: #fafafa;
}}
</style>
</head>
<body>

<h1>Rapport de portefeuille</h1>
<p class="small">Date de génération : {as_of}</p>

<h2>1. Synthèse chiffrée</h2>

{synth_html}
{detail_html}
{hist_html}

<p class="small">
Ce document est fourni à titre informatif uniquement et ne constitue pas un conseil en investissement
personnalisé.
</p>

</body>
    </html>
"""
        return html


    def _add_table_to_story(
        story: List[Any],
        df: pd.DataFrame,
        col_widths: Optional[List[float]] = None,
        font_size: int = 9,
    ):
        if df.empty:
            story.append(Paragraph("Données indisponibles.", getSampleStyleSheet()["Normal"]))
            return
        headers = list(df.columns)
        fmt_rows = []
        for _, row in df.iterrows():
            formatted = []
            for col, val in row.items():
                if "€" in col:
                    formatted.append(val if isinstance(val, str) else fmt_eur_fr(val))
                elif "%" in col:
                    formatted.append(val if isinstance(val, str) else fmt_pct_fr(val))
                else:
                    formatted.append(str(val))
            fmt_rows.append(formatted)
        data = [headers] + fmt_rows
        table = Table(data, repeatRows=1, colWidths=col_widths)
        style = [
            ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), font_size),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]
        for i in range(1, len(data)):
            if i % 2 == 0:
                style.append(("BACKGROUND", (0, i), (-1, i), colors.whitesmoke))
        table.setStyle(TableStyle(style))
        story.append(table)


    def _fig_to_rl_image(fig: plt.Figure, width: float = 480, height: float = 270) -> Image:
        if not MATPLOTLIB_AVAILABLE:
            return None
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        buf.seek(0)
        plt.close(fig)
        return Image(buf, width=width, height=height)


    def _build_value_chart(df_map: Dict[str, pd.DataFrame]) -> Optional[Image]:
        if not MATPLOTLIB_AVAILABLE:
            return None
        if not df_map:
            return None
        plt.rcParams["font.family"] = "DejaVu Sans"
        _CHART_COLORS = ["#1F3B6D", "#C8963E"]
        fig, ax = plt.subplots(figsize=(6, 3))
        fig.patch.set_facecolor("white")
        has_data = False
        for idx, (label, df) in enumerate(df_map.items()):
            if df is None or df.empty or "Valeur" not in df.columns:
                continue
            ax.plot(
                df.index, df["Valeur"],
                label=label,
                color=_CHART_COLORS[idx % len(_CHART_COLORS)],
                linewidth=1.8,
            )
            has_data = True
        if not has_data:
            plt.close(fig)
            return None
        from matplotlib.ticker import FuncFormatter
        ax.yaxis.set_major_formatter(FuncFormatter(lambda x, _: f"{x:,.0f}".replace(",", " ")))
        ax.set_title("Évolution de la valeur du portefeuille")
        ax.set_xlabel("Date")
        ax.set_ylabel("Valeur (€)")
        ax.legend(loc="best")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.grid(axis="y", color="#E2E8F0", linewidth=0.5)
        fig.autofmt_xdate()
        return _fig_to_rl_image(fig)


    def _wrap_label(label: str, width: int = 28) -> str:
        if not label:
            return "—"
        return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)


    def _allocation_from_positions(df_positions: pd.DataFrame) -> Tuple[pd.DataFrame, str]:
        df = df_positions.copy()
        df = df[df["Valeur actuelle €"] >= 0]
        df["Nom"] = df["Nom"].fillna("—")
        df["ISIN / Code"] = df["ISIN / Code"].fillna("—")

        total_value = df["Valeur actuelle €"].sum()
        if total_value > 0:
            df["Poids"] = df["Valeur actuelle €"] / total_value
            basis_label = "Valeur actuelle"
        else:
            total_net = df["Net investi €"].sum()
            if total_net > 0:
                df["Poids"] = df["Net investi €"] / total_net
            else:
                df["Poids"] = 0.0
                if len(df) > 0:
                    df.loc[df.index[0], "Poids"] = 1.0
            basis_label = "Net investi"

        df = df.sort_values("Poids", ascending=False)
        if len(df) > 8:
            df_main = df.iloc[:8].copy()
            df_other = df.iloc[8:]
            other_row = pd.DataFrame(
                {
                    "Nom": ["Autres"],
                    "ISIN / Code": ["—"],
                    "Net investi €": [df_other["Net investi €"].sum()],
                    "Valeur actuelle €": [df_other["Valeur actuelle €"].sum()],
                    "Poids": [df_other["Poids"].sum()],
                }
            )
            df = pd.concat([df_main, other_row], ignore_index=True)

        df["Part %"] = df["Poids"] * 100.0
        return df, basis_label


    def _build_allocation_donut(
        df_alloc: pd.DataFrame,
        title: str,
        figsize: Tuple[float, float] = (6.0, 3.4),
    ) -> Optional[Image]:
        if not MATPLOTLIB_AVAILABLE or df_alloc.empty:
            return None
        plt.rcParams["font.family"] = "DejaVu Sans"
        _DONUT_COLORS = ["#1F3B6D", "#C8963E", "#1A7A4A", "#F1F4F9", "#2E5FA3",
                         "#4B5563", "#6B7280", "#9CA3AF"]
        fig, ax = plt.subplots(figsize=figsize)
        fig.patch.set_facecolor("white")
        wedges, _ = ax.pie(
            df_alloc["Poids"],
            startangle=90,
            labels=None,
            colors=_DONUT_COLORS[:len(df_alloc)],
            wedgeprops=dict(width=0.35, edgecolor="white"),
        )
        labels = [
            f"{_wrap_label(nm)} ({pct:.1f}%)"
            for nm, pct in zip(df_alloc["Nom"], df_alloc["Part %"])
        ]
        ax.legend(
            wedges,
            labels,
            loc="center left",
            bbox_to_anchor=(1.02, 0.5),
            frameon=False,
            fontsize=8,
        )
        ax.set_title(title)
        ax.set_aspect("equal")
        fig.tight_layout(rect=[0, 0, 0.78, 1])
        return _fig_to_rl_image(fig, width=460, height=280)


    def _build_envelope_breakdown(
        lines: List[Dict[str, Any]],
        title: str,
    ) -> Tuple[Optional[Image], Optional[str]]:
        if not lines:
            return None, "Répartition par enveloppe : —"
        categories = {"Fonds euros": 0.0, "UC": 0.0, "Structurés": 0.0}
        for ln in lines:
            isin = str(ln.get("isin", "")).upper()
            val = float(ln.get("value", 0.0))
            if isin == "EUROFUND":
                categories["Fonds euros"] += val
            elif isin == "STRUCTURED":
                categories["Structurés"] += val
            else:
                categories["UC"] += val
        total = sum(categories.values())
        if total <= 0:
            return None, "Répartition par enveloppe : —"

        shares = {k: v / total for k, v in categories.items() if v > 0}
        major = max(shares.items(), key=lambda x: x[1])
        if sum(1 for v in shares.values() if v >= 0.01) < 2:
            return None, f"Répartition par enveloppe : {major[1] * 100:.1f}% {major[0]}"

        if not MATPLOTLIB_AVAILABLE:
            return None, None
        labels = list(shares.keys())
        values = [shares[k] * 100 for k in labels]
        fig, ax = plt.subplots(figsize=(6.0, 1.8))
        ax.barh(labels, values, color="#4C78A8")
        ax.set_xlim(0, 100)
        ax.set_xlabel("%")
        ax.set_title(title)
        for i, v in enumerate(values):
            ax.text(min(v + 1, 98), i, f"{v:.1f}%", va="center", fontsize=8)
        fig.tight_layout()
        return _fig_to_rl_image(fig, width=460, height=140), None


    def _build_contribution_bar(df_positions: pd.DataFrame) -> Optional[Image]:
        if not MATPLOTLIB_AVAILABLE or df_positions.empty:
            return None
        df = df_positions.copy()
        if not {"Nom", "Valeur actuelle €", "Net investi €"}.issubset(df.columns):
            return None
        plt.rcParams["font.family"] = "DejaVu Sans"
        df["Contribution €"] = df["Valeur actuelle €"] - df["Net investi €"]
        df = df.sort_values("Contribution €", ascending=False)
        bar_colors = ["#1A7A4A" if v >= 0 else "#CC2200" for v in df["Contribution €"]]
        fig_height = max(2.0, min(4.2, 0.35 * len(df) + 1.2))
        fig, ax = plt.subplots(figsize=(6.2, fig_height))
        fig.patch.set_facecolor("white")
        ax.barh(df["Nom"], df["Contribution €"], color=bar_colors)
        ax.invert_yaxis()
        ax.set_title("Contribution à la performance (€)")
        ax.axvline(0, color="black", linewidth=0.5)
        ax.tick_params(axis="y", labelsize=8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        for i, v in enumerate(df["Contribution €"]):
            offset = 0.01 * abs(v) if v != 0 else 0.5
            x_pos = v + offset if v >= 0 else v - offset
            ax.text(x_pos, i, fmt_eur_fr(v), va="center", fontsize=8)
        fig.tight_layout()
        return _fig_to_rl_image(fig, width=460, height=200)


    def _build_single_line_bar(label: str, value: float, title: str) -> Optional[Image]:
        if not MATPLOTLIB_AVAILABLE:
            return None
        fig, ax = plt.subplots(figsize=(6.0, 1.3))
        ax.barh([label], [100], color="#4C78A8")
        ax.set_xlim(0, 100)
        ax.set_title(title)
        ax.set_xticks([])
        ax.tick_params(axis="y", labelsize=8)
        fig.tight_layout()
        return _fig_to_rl_image(fig, width=460, height=90)


    def fmt_eur_pdf(x: Any) -> str:
        return fmt_eur_fr(x)


    def generate_pdf_report(report: Dict[str, Any]) -> bytes:
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError(f"PDF indisponible: {REPORTLAB_ERROR}")
        if not MATPLOTLIB_AVAILABLE:
            raise RuntimeError(f"PDF indisponible: {MATPLOTLIB_ERROR}")

        class NumberedCanvas(canvas.Canvas):
            def __init__(self, *args, **kwargs):
                super().__init__(*args, **kwargs)
                self._saved_page_states = []

            def showPage(self):
                self._saved_page_states.append(dict(self.__dict__))
                self._startPage()

            def save(self):
                page_count = len(self._saved_page_states)
                for state in self._saved_page_states:
                    self.__dict__.update(state)
                    self._draw_header_footer(page_count)
                    super().showPage()
                super().save()

            def _draw_header_footer(self, page_count: int):
                width, height = A4
                self.setFillColor(colors.HexColor("#1F3B6D"))
                self.setFont("Helvetica-Bold", 10)
                header_text = report.get("nom_cabinet", "") or "Rapport de portefeuille"
                self.drawString(36, height - 30, header_text)
                self.setFillColor(colors.grey)
                self.setFont("Helvetica", 8)
                as_of_str = report.get("as_of", "")
                nom_cli_str = report.get("nom_client", "").strip()
                right_str = f"{nom_cli_str}  ·  {as_of_str}" if nom_cli_str else as_of_str
                self.drawRightString(width - 36, height - 30, right_str)
                self.setFillColor(colors.grey)
                self.setFont("Helvetica", 7)
                self.drawString(
                    36,
                    24,
                    "Document informatif, ne constitue pas un conseil en investissement.",
                )
                self.drawRightString(width - 36, 24, f"Page {self.getPageNumber()} / {page_count}")

        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=36,
            rightMargin=36,
            topMargin=54,
            bottomMargin=48,
        )
        base_styles = getSampleStyleSheet()
        styles = {
            "title": ParagraphStyle(
                "Title",
                parent=base_styles["Title"],
                textColor=colors.HexColor("#1F3B6D"),
                fontSize=20,
                spaceAfter=12,
            ),
            "h1": ParagraphStyle(
                "H1",
                parent=base_styles["Heading1"],
                textColor=colors.HexColor("#1F3B6D"),
                fontSize=14,
                spaceAfter=8,
            ),
            "h2": ParagraphStyle(
                "H2",
                parent=base_styles["Heading2"],
                textColor=colors.HexColor("#4B5563"),
                fontSize=12,
                spaceAfter=6,
            ),
            "small": ParagraphStyle(
                "Small",
                parent=base_styles["Normal"],
                fontSize=8,
                textColor=colors.grey,
            ),
            "kpi": ParagraphStyle(
                "KPI",
                parent=base_styles["Normal"],
                fontSize=10,
                textColor=colors.HexColor("#111827"),
            ),
        }
        story: List[Any] = []

        def _kpi_table(title: str, rows: List[Tuple[str, str]]) -> Table:
            data = [[Paragraph(f"<b>{title}</b>", styles["h2"]), ""]]
            for label, value in rows:
                data.append([Paragraph(label, styles["small"]), Paragraph(value, styles["kpi"])])
            table = Table(data, colWidths=[160, 120])
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF2F7")),
                        ("BACKGROUND", (0, 1), (-1, -1), colors.whitesmoke),
                        ("BOX", (0, 0), (-1, -1), 0.5, colors.lightgrey),
                        ("INNERGRID", (0, 1), (-1, -1), 0.25, colors.lightgrey),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ]
                )
            )
            return table

        story.append(Paragraph("Rapport client", styles["title"]))
        _cover_nom_cli = report.get("nom_client", "").strip()
        _cover_nom_cab = report.get("nom_cabinet", "").strip()
        if _cover_nom_cli or _cover_nom_cab:
            _cover_parts = []
            if _cover_nom_cli:
                _cover_parts.append(f"Client : {_cover_nom_cli}")
            if _cover_nom_cab:
                _cover_parts.append(f"Cabinet : {_cover_nom_cab}")
            story.append(Paragraph("  ·  ".join(_cover_parts), styles["small"]))
        story.append(Paragraph(f"Date de génération : {report.get('as_of', '')}", styles["small"]))
        story.append(Spacer(1, 12))

        mode = report.get("mode", "compare")
        synthA = report.get("client_summary", {})
        synthB = report.get("valority_summary", {})
        comp = report.get("comparison", {})

        story.append(Paragraph("Synthèse", styles["h1"]))
        if mode == "compare":
            client_rows = [
                ("Valeur actuelle", fmt_eur_fr(synthA.get("val", 0))),
                ("Net investi", fmt_eur_fr(synthA.get("net", 0))),
                ("Brut versé", fmt_eur_fr(synthA.get("brut", 0))),
                ("Perf totale", fmt_pct_fr(synthA.get("perf_tot_pct", 0))),
                ("XIRR", fmt_pct_fr(synthA.get("irr_pct", 0))),
            ]
            valority_rows = [
                ("Valeur actuelle", fmt_eur_fr(synthB.get("val", 0))),
                ("Net investi", fmt_eur_fr(synthB.get("net", 0))),
                ("Brut versé", fmt_eur_fr(synthB.get("brut", 0))),
                ("Perf totale", fmt_pct_fr(synthB.get("perf_tot_pct", 0))),
                ("XIRR", fmt_pct_fr(synthB.get("irr_pct", 0))),
            ]
            table = Table(
                [[_kpi_table("Client", client_rows), _kpi_table("Cabinet", valority_rows)]],
                colWidths=[240, 240],
            )
            story.append(table)
            story.append(Spacer(1, 4))
            story.append(Paragraph(
                "Performance nette des frais de gestion du contrat. "
                "Les frais internes des fonds (TER) sont intégrés dans les valeurs liquidatives.",
                styles["small"],
            ))
            story.append(Spacer(1, 8))
            comp_rows = [
                ("Différence de valeur", fmt_eur_fr(comp.get("delta_val", 0))),
                ("Écart de performance", fmt_pct_fr(comp.get("delta_perf_pct", 0))),
            ]
            story.append(_kpi_table("Comparaison", comp_rows))
        else:
            title = "Cabinet" if mode == "valority" else "Client"
            synth = synthB if mode == "valority" else synthA
            rows = [
                ("Valeur actuelle", fmt_eur_fr(synth.get("val", 0))),
                ("Net investi", fmt_eur_fr(synth.get("net", 0))),
                ("Brut versé", fmt_eur_fr(synth.get("brut", 0))),
                ("Perf totale", fmt_pct_fr(synth.get("perf_tot_pct", 0))),
                ("XIRR", fmt_pct_fr(synth.get("irr_pct", 0))),
            ]
            story.append(_kpi_table(title, rows))
            story.append(Spacer(1, 4))
            story.append(Paragraph(
                "Performance nette des frais de gestion du contrat. "
                "Les frais internes des fonds (TER) sont intégrés dans les valeurs liquidatives.",
                styles["small"],
            ))
            fees = report.get("fees_analysis", {})
            if fees:
                story.append(Spacer(1, 8))
                fees_rows = [
                    ("Frais d'entrée payés", fmt_eur_fr(fees.get("fees_paid", 0))),
                    ("Valeur créée", fmt_eur_fr(fees.get("value_created", 0))),
                    ("Valeur/an", fmt_eur_fr(fees.get("value_per_year", 0))),
                ]
                story.append(_kpi_table("Frais & valeur créée", fees_rows))

        story.append(Spacer(1, 12))

        # ---- Diversification & Indicateurs de risque ----
        def _fmt_ratio(v, decimals=2):
            if v is None:
                return "—"
            return f"{v:.{decimals}f}"

        def _div_label_pdf(score):
            if score >= 70:
                return "Bonne diversification"
            elif score >= 40:
                return "Diversification moyenne"
            return "Fausse diversification"

        div_A = report.get("diversification_A")
        div_B = report.get("diversification_B")
        risk_pdf_A = report.get("risk_A")
        risk_pdf_B = report.get("risk_B")
        has_div = div_A is not None or div_B is not None

        if has_div:
            story.append(Paragraph("Diversification du portefeuille", styles["h1"]))

            def _div_rows(div_res, risk_res, port_val):
                rows = []
                if div_res:
                    rows.append(("Score de diversification", f"{div_res['score']:.0f}/100 — {_div_label_pdf(div_res['score'])}"))
                    rows.append(("Corrélation moyenne", f"{div_res['avg_corr']:.0%}"))
                    rows.append(("Nb lignes analysées", str(div_res["n_lines"])))
                    if div_res["doublons"]:
                        pairs = ", ".join(f"{a}/{b}" for a, b, _ in div_res["doublons"])
                        rows.append(("Doublons détectés", pairs))
                if risk_res:
                    rows.append(("Volatilité annuelle", f"{risk_res['vol_ann_pct']:.1f}%"))
                    rows.append(("Max drawdown", f"{risk_res['max_dd_pct']:.1f}%"))
                return rows

            if mode == "compare":
                _val_A = report.get("client_summary", {}).get("val", 0)
                _val_B = report.get("valority_summary", {}).get("val", 0)
                _div_tbl = Table(
                    [[
                        _kpi_table("Client", _div_rows(div_A, risk_pdf_A, _val_A)),
                        _kpi_table("Cabinet", _div_rows(div_B, risk_pdf_B, _val_B)),
                    ]],
                    colWidths=[240, 240],
                )
                story.append(_div_tbl)
            else:
                _d = div_B if mode == "valority" else div_A
                _r = risk_pdf_B if mode == "valority" else risk_pdf_A
                _lbl = "Cabinet" if mode == "valority" else "Client"
                if _d or _r:
                    story.append(_kpi_table(_lbl, _div_rows(_d, _r, None)))

            story.append(Spacer(1, 8))

        # ---- Ratios techniques ----
        ratios_A = report.get("ratios_A", {})
        ratios_B = report.get("ratios_B", {})
        has_ratios = any(ratios_A.get(k) is not None for k in ("sharpe", "sortino", "beta_alpha")) or \
                     any(ratios_B.get(k) is not None for k in ("sharpe", "sortino", "beta_alpha"))
        if has_ratios:
            story.append(Paragraph("Ratios techniques", styles["h1"] if not has_div else styles["h2"]))
            if mode == "compare":
                def _ratio_rows(rat):
                    ba = rat.get("beta_alpha") or {}
                    return [
                        ("Ratio de Sharpe", _fmt_ratio(rat.get("sharpe"))),
                        ("Ratio de Sortino", _fmt_ratio(rat.get("sortino"))),
                        ("Bêta", _fmt_ratio(ba.get("beta"))),
                        ("Alpha (%/an)", _fmt_ratio(ba.get("alpha_pct"), 2)),
                    ]
                _rt = Table(
                    [[_kpi_table("Client", _ratio_rows(ratios_A)), _kpi_table("Cabinet", _ratio_rows(ratios_B))]],
                    colWidths=[240, 240],
                )
                story.append(_rt)
            else:
                rat = ratios_B if mode == "valority" else ratios_A
                ba = rat.get("beta_alpha") or {}
                bench_name = ba.get("benchmark_name", "")
                ratio_rows_single = [
                    ("Ratio de Sharpe", _fmt_ratio(rat.get("sharpe"))),
                    ("Ratio de Sortino", _fmt_ratio(rat.get("sortino"))),
                    (f"Beta{' vs ' + bench_name if bench_name else ''}", _fmt_ratio(ba.get("beta"))),
                    ("Alpha (%/an)", _fmt_ratio(ba.get("alpha_pct"), 2)),
                ]
                story.append(_kpi_table("Ratios de risque", ratio_rows_single))
            story.append(Paragraph(
                "Sharpe = (rendement annualisé – taux sans risque) / volatilité. "
                "Sortino = idem avec volatilité à la baisse uniquement. "
                "Beta = sensibilité relative au marché de référence.",
                styles["small"],
            ))
            story.append(Spacer(1, 12))

        story.append(Paragraph("Graphiques", styles["h1"]))

        value_chart = _build_value_chart(report.get("df_map", {}))
        if value_chart is not None:
            story.append(value_chart)
            story.append(Spacer(1, 8))
        else:
            story.append(Paragraph("Données indisponibles pour la courbe de valeur.", styles["small"]))

        def add_portfolio_details_section(
            label: str,
            positions_df: pd.DataFrame,
            lines_values: List[Dict[str, Any]],
        ):
            # PAGE 2/4 — Détail du portefeuille
            story.append(PageBreak())
            story.append(Paragraph(f"Détail du portefeuille — {label}", styles["h1"]))

            if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
                df_alloc, basis_label = _allocation_from_positions(positions_df)

                if len(df_alloc) >= 2:
                    story.append(Paragraph(f"Allocation par ligne ({basis_label})", styles["h2"]))
                    donut = _build_allocation_donut(df_alloc, "Allocation par ligne")
                    if donut is not None:
                        story.append(donut)
                        story.append(Spacer(1, 6))
                else:
                    line = df_alloc.iloc[0] if not df_alloc.empty else None
                    name = line["Nom"] if line is not None else "—"
                    story.append(
                        Paragraph(
                            f"Portefeuille concentré : 100% sur <b>{name}</b>.",
                            styles["kpi"],
                        )
                    )
                    bar = _build_single_line_bar(_wrap_label(name), 100.0, "Répartition 100%")
                    if bar is not None:
                        story.append(bar)
                        story.append(Spacer(1, 6))

                alloc_table = df_alloc[
                    ["Nom", "ISIN / Code", "Part %", "Net investi €", "Valeur actuelle €"]
                ].copy()
                _add_table_to_story(
                    story,
                    alloc_table,
                    col_widths=[170, 80, 55, 95, 95],
                    font_size=9,
                )
                story.append(Spacer(1, 6))

                envelope_chart, envelope_text = _build_envelope_breakdown(
                    lines_values,
                    "Répartition par enveloppe",
                )
                if envelope_chart is not None:
                    story.append(envelope_chart)
                elif envelope_text:
                    story.append(Paragraph(envelope_text, styles["small"]))
            else:
                story.append(Paragraph("Données indisponibles pour la composition.", styles["small"]))

            # PAGE 3/5 — Contribution & Positions
            story.append(PageBreak())
            story.append(Paragraph(f"Contribution & positions — {label}", styles["h1"]))

            if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
                if len(positions_df) == 1:
                    ln = positions_df.iloc[0]
                    story.append(
                        Paragraph(
                            f"Contribution : <b>{ln['Nom']}</b> = {fmt_eur_pdf(ln['Valeur actuelle €'] - ln['Net investi €'])}",
                            styles["kpi"],
                        )
                    )
                    bar = _build_single_line_bar(_wrap_label(ln["Nom"]), 100.0, "Contribution (ligne unique)")
                    if bar is not None:
                        story.append(bar)
                        story.append(Spacer(1, 6))
                else:
                    contrib_chart = _build_contribution_bar(positions_df)
                    if contrib_chart is not None:
                        story.append(contrib_chart)
                        story.append(Spacer(1, 6))
            else:
                story.append(Paragraph("Données indisponibles pour la contribution.", styles["small"]))

            story.append(Paragraph("Positions", styles["h2"]))
            if isinstance(positions_df, pd.DataFrame) and not positions_df.empty:
                positions_table = positions_df[
                    [
                        "Nom",
                        "ISIN / Code",
                        "Date d'achat",
                        "Net investi €",
                        "Valeur actuelle €",
                        "Perf €",
                        "Perf %",
                    ]
                ].copy()
                _add_table_to_story(
                    story,
                    positions_table,
                    col_widths=[150, 70, 65, 80, 80, 50, 45],
                    font_size=9,
                )
            else:
                story.append(Paragraph("Données indisponibles.", styles["small"]))

        if mode == "compare":
            add_portfolio_details_section(
                "Client",
                report.get("positions_df_client", pd.DataFrame()),
                report.get("lines_client", []),
            )
            add_portfolio_details_section(
                "Cabinet",
                report.get("positions_df_valority", pd.DataFrame()),
                report.get("lines_valority", []),
            )
        else:
            add_portfolio_details_section(
                "Cabinet" if mode == "valority" else "Client",
                report.get("positions_df", pd.DataFrame()),
                report.get("lines", []),
            )

        doc.build(story, canvasmaker=NumberedCanvas)
        buffer.seek(0)
        return buffer.read()


    def _chart_to_pptx_image(
        df: pd.DataFrame,
        col_y: str,
        label: str,
        color: str = "#1B2A4A",
    ) -> Optional[bytes]:
        """
        Génère un graphique matplotlib en mémoire (PNG bytes)
        pour insertion dans le PPTX.
        df doit avoir un DatetimeIndex ou colonne Date.
        """
        if not MATPLOTLIB_AVAILABLE or df is None or df.empty:
            return None
        try:
            import io as _io
            _fig, _ax = plt.subplots(figsize=(7, 2.8))
            _ax.fill_between(
                df.index if hasattr(df.index, "dtype") and "datetime" in str(df.index.dtype)
                else range(len(df)),
                df[col_y].ffill(),
                alpha=0.18,
                color=color,
            )
            _ax.plot(
                df.index if hasattr(df.index, "dtype") and "datetime" in str(df.index.dtype)
                else range(len(df)),
                df[col_y].ffill(),
                color=color,
                linewidth=1.8,
            )
            _ax.set_ylabel("€", fontsize=8)
            _ax.set_title(label, fontsize=9, color="#1B2A4A", fontweight="bold")
            _ax.spines["top"].set_visible(False)
            _ax.spines["right"].set_visible(False)
            _ax.yaxis.set_major_formatter(
                plt.FuncFormatter(lambda x, _: f"{x:,.0f}")
            )
            _fig.tight_layout()
            _buf = _io.BytesIO()
            _fig.savefig(_buf, format="png", dpi=130, bbox_inches="tight")
            plt.close(_fig)
            _buf.seek(0)
            return _buf.read()
        except Exception:
            return None

    def generate_pptx_report(report: Dict[str, Any]) -> bytes:
        """Génère une présentation 16:9 de 5 slides pour le client."""
        if not PPTX_AVAILABLE:
            raise RuntimeError(f"PPTX indisponible: {PPTX_ERROR}")

        # ── Palette ──────────────────────────────────────────────────────────
        NAVY   = RGBColor(0x1F, 0x3B, 0x6D)
        GOLD   = RGBColor(0xC8, 0x96, 0x3E)
        ICE    = RGBColor(0xE8, 0xEE, 0xF7)
        GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
        GREY   = RGBColor(0x64, 0x74, 0x8B)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        LGREY  = RGBColor(0xF1, 0xF4, 0xF9)
        BORDER = RGBColor(0xD1, 0xD9, 0xE6)
        BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
        BLACK  = RGBColor(0x11, 0x18, 0x27)

        # ── Données du rapport ────────────────────────────────────────────────
        nom_cab  = report.get("nom_cabinet", "") or "Votre cabinet"
        nom_cli  = report.get("nom_client", "") or ""
        as_of    = report.get("as_of", "")
        contrat  = report.get("contrat_label", "Assurance-vie") or "Assurance-vie"
        date_ouv = report.get("date_ouverture", "")
        mode     = report.get("mode", "compare")

        synthA  = report.get("client_summary", {})
        synthB  = report.get("valority_summary", {})
        comp    = report.get("comparison", {})

        # Synthèse principale selon le mode
        if mode == "valority":
            synth_main = synthB
            _raw_df = report.get("positions_df_valority")
            positions_df_main = _raw_df if isinstance(_raw_df, pd.DataFrame) else pd.DataFrame()
        else:
            synth_main = synthA
            _raw_df = report.get("positions_df_client")
            positions_df_main = _raw_df if isinstance(_raw_df, pd.DataFrame) else pd.DataFrame()

        val  = synth_main.get("val", 0.0)
        net  = synth_main.get("net", 0.0)
        brut = synth_main.get("brut", 0.0)
        perf = synth_main.get("perf_tot_pct", 0.0)
        xirr_val = synth_main.get("irr_pct", 0.0)
        # Guard : si val = 0, utiliser valB ou valA selon disponibilité
        if val == 0.0 and mode == "compare":
            val  = report.get("valority_summary", {}).get("val", 0.0) or \
                   report.get("client_summary", {}).get("val", 0.0)
            net  = report.get("valority_summary", {}).get("net", 0.0) or \
                   report.get("client_summary", {}).get("net", 0.0)
            brut = report.get("valority_summary", {}).get("brut", 0.0) or \
                   report.get("client_summary", {}).get("brut", 0.0)

        def _fe(x: Any) -> str:
            try:
                return f"{float(x):,.0f} €".replace(",", " ")
            except Exception:
                return "— €"

        def _fp(x: Any) -> str:
            try:
                return f"{float(x):+.2f}%"
            except Exception:
                return "—%"

        # ── Helpers ──────────────────────────────────────────────────────────
        def _add_rect(slide, l: float, t: float, w: float, h: float,
                      fill_rgb: RGBColor, line_rgb: Optional[RGBColor] = None):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE RECTANGLE
                Inches(l), Inches(t), Inches(w), Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            if line_rgb is None:
                shape.line.fill.background()
            else:
                shape.line.color.rgb = line_rgb
            return shape

        def _add_text(slide, text: str, l: float, t: float, w: float, h: float,
                      size: float, color: RGBColor, bold: bool = False,
                      align: str = "left", italic: bool = False) -> None:
            txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.italic = italic
            p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.RIGHT if align == "right" else PP_ALIGN.LEFT

        def _footer(slide):
            parts = [nom_cab]
            if nom_cli:
                parts.append(f"Document pour {nom_cli}")
            if as_of:
                parts.append(as_of)
            _add_text(slide, "  ·  ".join(parts),
                      0.28, 5.38, 9.44, 0.20, 7, GREY, align="center")

        def _left_stripe(slide, color: RGBColor = NAVY):
            _add_rect(slide, 0, 0, 0.18, 5.625, color)

        # ── Présentation ──────────────────────────────────────────────────────
        from pptx import Presentation as _Prs
        from pptx.util import Inches as _In, Pt as _Pt, Emu as _Emu

        prs = _Prs()
        prs.slide_width  = _In(10)
        prs.slide_height = _In(5.625)

        blank_layout = prs.slide_layouts[6]  # blank layout

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 1 — Page de garde
        # ══════════════════════════════════════════════════════════════════════
        s1 = prs.slides.add_slide(blank_layout)
        _add_rect(s1, 0, 0, 10, 5.625, NAVY)          # fond navy
        _add_rect(s1, 0, 0, 0.18, 5.625, GOLD)         # bande gauche or
        _add_rect(s1, 0.55, 0.70, 8.80, 2.00, LGREY)   # bandeau clair

        _add_text(s1, nom_cab.upper(),
                  0.55, 0.80, 8.80, 0.32, 9, GOLD, bold=False, align="center")
        if mode == "client":
            _titre_s1 = f"Audit Patrimonial"
            _sous_titre_s1 = f"Situation actuelle de {nom_cli}" if nom_cli else "Situation actuelle"
        elif mode == "valority":
            _titre_s1 = f"Proposition d'Investissement"
            _sous_titre_s1 = f"Une stratégie élaborée par {nom_cab}" if nom_cab else "Stratégie conseiller"
        else:  # compare
            _titre_s1 = f"Synthèse d'Arbitrage"
            _sous_titre_s1 = f"Audit & Recommandation {nom_cab}" if nom_cab else "Comparatif complet"
        _add_text(s1, _titre_s1,
                  0.55, 1.00, 8.80, 0.55, 28, NAVY, bold=True, align="center")
        _add_text(s1, _sous_titre_s1,
                  0.55, 1.58, 8.80, 0.32, 13, GREY, align="center")

        # séparateur or
        _add_rect(s1, 3.20, 2.80, 3.60, 0.04, GOLD)

        # Nom client
        _add_text(s1, nom_cli if nom_cli else "—",
                  0.55, 2.95, 8.80, 0.42, 20, WHITE, bold=True, align="center")
        # Contrat
        _add_text(s1, contrat,
                  0.55, 3.50, 8.80, 0.28, 9, ICE, align="center")
        # Dates
        _dates_str = ""
        if date_ouv:
            _dates_str = f"Ouverture : {date_ouv}"
        if as_of:
            _dates_str += (f"  ·  Arrêté au : {as_of}" if _dates_str else f"Arrêté au : {as_of}")
        _add_text(s1, _dates_str,
                  0.55, 3.80, 8.80, 0.28, 8, GREY, align="center")
        # Disclaimer
        _add_text(s1,
                  "Document à usage interne — Simulation historique. "
                  "Les performances passées ne préjugent pas des performances futures.",
                  0.55, 5.20, 8.80, 0.28, 6, GREY, align="center")

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 2 — Votre portefeuille aujourd'hui
        # ══════════════════════════════════════════════════════════════════════
        s2 = prs.slides.add_slide(blank_layout)
        _add_rect(s2, 0, 0, 10, 5.625, LGREY)
        _left_stripe(s2)
        _footer(s2)

        _add_text(s2, "Votre portefeuille aujourd'hui",
                  0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
        _add_text(s2, f"Situation au {as_of}",
                  0.30, 0.68, 9.30, 0.25, 9, GREY)

        # 4 cartes KPI
        _kpi_data = [
            ("Valeur actuelle", _fe(val), f"Versé : {_fe(brut)}", WHITE),
            ("Capital investi", _fe(net), "Après frais d'entrée", WHITE),
            ("Performance totale", f"{perf:+.2f}%", "Depuis l'ouverture", GREEN),
            ("Rendement annualisé", f"{xirr_val:+.2f}%", "XIRR", BLUE),
        ]
        _kpi_x = [0.28, 2.58, 4.88, 7.18]
        for idx, (lbl, val_str, sub_lbl, val_color) in enumerate(_kpi_data):
            _xk = _kpi_x[idx]
            _add_rect(s2, _xk, 1.05, 2.20, 1.10, WHITE, BORDER)
            _add_text(s2, lbl,  _xk + 0.08, 1.10, 2.04, 0.22, 8, GREY)
            _add_text(s2, val_str, _xk + 0.08, 1.33, 2.04, 0.30, 18, val_color, bold=True)
            _add_text(s2, sub_lbl, _xk + 0.08, 1.65, 2.04, 0.22, 7, GREY)

        # Bande "Valeur créée"
        _vc = val - net
        _fees_paid = max(0.0, brut - net)
        _add_rect(s2, 0.28, 2.28, 9.10, 0.82, NAVY)
        _add_text(s2, "Valeur créée pour vous",
                  0.36, 2.30, 2.50, 0.28, 10, ICE, bold=True)
        _add_text(s2, _fe(_vc),
                  0.36, 2.55, 2.50, 0.30, 18, GOLD, bold=True)
        _fees_note = f"Frais d'entrée payés : {_fe(_fees_paid)}  ·  Performance nette des frais contrat"
        _add_text(s2, _fees_note, 0.36, 2.90, 9.00, 0.18, 7, ICE)

        # Tableau positions
        _add_text(s2, "Composition du portefeuille",
                  0.28, 3.22, 9.10, 0.28, 11, NAVY, bold=True)

        if not positions_df_main.empty:
            _cols_needed = {"Nom", "ISIN / Code", "Net investi €", "Valeur actuelle €", "Perf %"}
            _df_tbl = positions_df_main.copy()
            _row_y = 3.52
            _col_w = [3.5, 0.85, 1.5, 1.65, 1.0]
            _col_x = [0.28]
            for _cw in _col_w[:-1]:
                _col_x.append(_col_x[-1] + _cw)

            _headers = ["Fonds", "Part %", "Net investi", "Valeur actuelle", "Perf %"]
            for ci, (_cx, _ch) in enumerate(zip(_col_x, _headers)):
                _add_rect(s2, _cx, _row_y, _col_w[ci], 0.25, NAVY)
                _add_text(s2, _ch, _cx + 0.04, _row_y + 0.02, _col_w[ci] - 0.08, 0.22, 9, WHITE, bold=True)

            _total_val_tbl = _df_tbl["Valeur actuelle €"].sum() if "Valeur actuelle €" in _df_tbl.columns else 0.0
            for ri, row in enumerate(_df_tbl.head(6).itertuples(index=False)):
                _ry = _row_y + 0.25 + ri * 0.26
                _row_bg = WHITE if ri % 2 == 0 else LGREY
                _add_rect(s2, 0.28, _ry, 9.10, 0.26, _row_bg)
                _nom_tbl = str(getattr(row, "Nom", "—"))[:38]
                _net_tbl  = float(getattr(row, "Net_investi_€",  getattr(row, "Net investi €", 0)) or 0)
                _val_tbl  = float(getattr(row, "Valeur_actuelle_€", getattr(row, "Valeur actuelle €", 0)) or 0)
                _perf_tbl = ((_val_tbl / _net_tbl - 1) * 100) if _net_tbl > 0 else 0.0
                _part_tbl = (_val_tbl / _total_val_tbl * 100) if _total_val_tbl > 0 else 0.0
                _perf_color = GREEN if _perf_tbl >= 0 else RGBColor(0xCC, 0x22, 0x00)
                _row_vals = [_nom_tbl, f"{_part_tbl:.1f}%", _fe(_net_tbl), _fe(_val_tbl), f"{_perf_tbl:+.2f}%"]
                for ci2, (_cx2, _rv) in enumerate(zip(_col_x, _row_vals)):
                    _tc = _perf_color if ci2 == 4 else BLACK
                    _add_text(s2, _rv, _cx2 + 0.04, _ry + 0.03, _col_w[ci2] - 0.08, 0.20, 8, _tc)

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 3 — Allocation recommandée (ou analyse simple)
        # ══════════════════════════════════════════════════════════════════════
        s3 = prs.slides.add_slide(blank_layout)
        _add_rect(s3, 0, 0, 10, 5.625, WHITE)
        _left_stripe(s3)
        _footer(s3)

        if mode == "compare":
            val_B  = synthB.get("val", 0.0)
            net_B  = synthB.get("net", 0.0)
            perf_B = synthB.get("perf_tot_pct", 0.0)
            xirr_B = synthB.get("irr_pct", 0.0)
            delta_val  = comp.get("delta_val", 0.0)
            delta_perf = comp.get("delta_perf_pct", 0.0)
            delta_xirr = (xirr_B - synth_main.get("irr_pct", 0.0))

            _add_text(s3, "L'allocation recommandée",
                      0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
            _add_text(s3, "Simulation sur la même période avec l'allocation de votre conseiller",
                      0.30, 0.68, 9.30, 0.25, 9, GREY)

            # Col gauche — 3 cartes KPI verticales
            _s3_kpis = [
                ("Valeur simulée", _fe(val_B), "Allocation conseiller", GREEN),
                ("Performance totale", f"{perf_B:+.2f}%", "Depuis l'ouverture", GREEN),
                ("Rendement annualisé", f"{xirr_B:+.2f}%", "XIRR", BLUE),
            ]
            for ki, (klbl, kval, ksub, kcolor) in enumerate(_s3_kpis):
                _ky = 1.08 + ki * 1.12
                _add_rect(s3, 0.28, _ky, 2.85, 1.00, LGREY, BORDER)
                _add_text(s3, klbl,  0.36, _ky + 0.06, 2.69, 0.22, 8, GREY)
                _add_text(s3, kval,  0.36, _ky + 0.28, 2.69, 0.32, 18, kcolor, bold=True)
                _add_text(s3, ksub,  0.36, _ky + 0.65, 2.69, 0.22, 7, GREY)

            # Graphique évolution si disponible
            _dfA_val = report.get("dfA_val")
            _dfB_val = report.get("dfB_val")
            if isinstance(_dfA_val, pd.DataFrame) and not _dfA_val.empty \
                    and "Valeur" in _dfA_val.columns:
                _img_bytes = _chart_to_pptx_image(
                    _dfA_val.set_index("Date") if "Date" in _dfA_val.columns else _dfA_val,
                    "Valeur", "Évolution — Portefeuille actuel", "#1B2A4A"
                )
            elif isinstance(_dfB_val, pd.DataFrame) and not _dfB_val.empty \
                    and "Valeur" in _dfB_val.columns:
                _img_bytes = _chart_to_pptx_image(
                    _dfB_val.set_index("Date") if "Date" in _dfB_val.columns else _dfB_val,
                    "Valeur", "Évolution — Portefeuille proposé", "#C9A84C"
                )
            else:
                _img_bytes = None
            if _img_bytes is not None:
                _img_stream = BytesIO(_img_bytes)
                s3.shapes.add_picture(
                    _img_stream,
                    left=Inches(0.28), top=Inches(1.05),
                    width=Inches(2.85), height=Inches(2.00),
                )
            # Col droite — bloc delta
            _add_rect(s3, 3.32, 1.00, 6.16, 3.55, ICE, BORDER)
            _add_text(s3, "Ce que l'allocation recommandée\naurait généré de plus",
                      3.50, 1.08, 5.80, 0.60, 13, NAVY, bold=True)
            _delta_sign = "+" if delta_val >= 0 else ""
            _add_text(s3, f"{_delta_sign}{_fe(delta_val)}",
                      3.32, 1.72, 6.16, 0.58, 36, GOLD, bold=True, align="center")
            _add_text(s3, "de valeur supplémentaire",
                      3.32, 2.28, 6.16, 0.28, 10, GREY, align="center")
            _add_rect(s3, 3.52, 2.64, 5.76, 0.02, BORDER)
            # Deux colonnes delta
            _d_kpis = [
                (_fp(delta_perf), "de performance en plus"),
                (_fp(delta_xirr), "de rendement annualisé"),
            ]
            for di, (dv, dl) in enumerate(_d_kpis):
                _dx = 3.52 + di * 3.00
                _add_text(s3, dv, _dx, 2.75, 2.80, 0.38, 20, GREEN, bold=True, align="center")
                _add_text(s3, dl, _dx, 3.14, 2.80, 0.25, 9, GREY, align="center")

            _add_text(s3,
                      "⚠️ Simulation historique. Les performances passées ne préjugent pas des performances futures.",
                      0.28, 4.72, 9.44, 0.22, 8, GREY)
        elif mode == "client":
            # Slide 3 Audit — Points forts / Points de vigilance
            _add_text(s3, "Audit & Diagnostic de votre portefeuille",
                      0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
            _add_text(s3, f"Analyse au {as_of}", 0.30, 0.68, 9.30, 0.25, 9, GREY)
            # Logique Python d'audit automatique
            _points_forts = []
            _points_vigilance = []
            _xirr_a = synth_main.get("irr_pct", 0.0) or 0.0
            _val_a  = synth_main.get("val", 0.0) or 0.0
            _net_a  = synth_main.get("net", 0.0) or 0.0
            _perf_a = synth_main.get("perf_tot_pct", 0.0) or 0.0
            _fee_c  = float(report.get("fee_contract_pct", 0.6) or 0.6)
            if _xirr_a > 5.0:
                _points_forts.append(f"Rendement solide : XIRR de {_xirr_a:.1f}%/an")
            elif _xirr_a > 2.5:
                _points_forts.append(f"Rendement positif : XIRR de {_xirr_a:.1f}%/an")
            else:
                _points_vigilance.append(f"Rendement limité : XIRR de {_xirr_a:.1f}%/an — sous l'inflation")
            if _fee_c > 1.5:
                _points_vigilance.append(f"Empilement de frais : {_fee_c:.2f}%/an limite la performance nette")
            elif _fee_c < 0.8:
                _points_forts.append(f"Enveloppe compétitive : frais contrat de {_fee_c:.2f}%/an")
            if _perf_a > 10.0:
                _points_forts.append(f"Performance historique de {_perf_a:.1f}% depuis l'ouverture")
            elif _perf_a < 0:
                _points_vigilance.append(f"Performance négative : {_perf_a:.1f}% depuis l'ouverture")
            if not _points_forts:
                _points_forts.append("Portefeuille en cours de construction")
            if not _points_vigilance:
                _points_vigilance.append("Aucun point de vigilance identifié à ce stade")
            # Col gauche — Points forts
            _add_rect(s3, 0.28, 1.05, 4.50, 0.32, GREEN)
            _add_text(s3, "✓ Points forts", 0.36, 1.09, 4.34, 0.24, 10, WHITE, bold=True)
            for pfi, pf in enumerate(_points_forts[:3]):
                _pfy = 1.42 + pfi * 0.52
                _add_rect(s3, 0.28, _pfy, 4.50, 0.45, WHITE, BORDER)
                _add_text(s3, f"• {pf}", 0.36, _pfy + 0.05, 4.34, 0.36, 9, BLACK)
            # Col droite — Points de vigilance
            _add_rect(s3, 5.00, 1.05, 4.72, 0.32, RGBColor(0xCC, 0x44, 0x00))
            _add_text(s3, "⚠ Points de vigilance", 5.08, 1.09, 4.56, 0.24, 10, WHITE, bold=True)
            for pvi, pv in enumerate(_points_vigilance[:3]):
                _pvy = 1.42 + pvi * 0.52
                _add_rect(s3, 5.00, _pvy, 4.72, 0.45, WHITE, BORDER)
                _add_text(s3, f"• {pv}", 5.08, _pvy + 0.05, 4.56, 0.36, 9, BLACK)
            _add_text(s3,
                      "⚠️ Simulation indicative basée sur l'historique de VL. "
                      "Les performances passées ne préjugent pas des performances futures.",
                      0.28, 5.00, 9.44, 0.22, 7, GREY)
        else:  # valority seul
            _add_text(s3, f"La Stratégie {nom_cab}",
                      0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
            _add_text(s3,
                      f"Une allocation construite pour {nom_cli}" if nom_cli
                      else "Allocation conseiller",
                      0.30, 0.68, 9.30, 0.25, 9, GREY)
            _val_b  = synthB.get("val", 0.0) or 0.0
            _net_b  = synthB.get("net", 0.0) or 0.0
            _xirr_b = synthB.get("irr_pct", 0.0) or 0.0
            _perf_b = synthB.get("perf_tot_pct", 0.0) or 0.0
            _s3b_kpis = [
                ("Valeur simulée", _fe(_val_b), "Notre allocation", GREEN),
                ("Rendement annualisé", f"{_xirr_b:+.2f}%", "XIRR net de frais", BLUE),
                ("Performance totale", f"{_perf_b:+.2f}%", "Depuis l'ouverture", GREEN),
            ]
            for ki3b, (klbl3b, kval3b, ksub3b, kcolor3b) in enumerate(_s3b_kpis):
                _ky3b = 1.10 + ki3b * 1.10
                _add_rect(s3, 0.28, _ky3b, 4.20, 0.95, LGREY, BORDER)
                _add_text(s3, klbl3b, 0.36, _ky3b + 0.06, 4.04, 0.22, 8, GREY)
                _add_text(s3, kval3b, 0.36, _ky3b + 0.28, 4.04, 0.32, 18, kcolor3b, bold=True)
                _add_text(s3, ksub3b, 0.36, _ky3b + 0.68, 4.04, 0.20, 7, GREY)
            _add_rect(s3, 4.72, 1.05, 4.90, 3.55, ICE, BORDER)
            _add_text(s3, "Notre approche",
                      4.86, 1.12, 4.62, 0.28, 11, NAVY, bold=True)
            _engagements_s3 = [
                ("Maîtrise de la volatilité",
                 "Allocation diversifiée limitant le drawdown"),
                ("Performance nette optimisée",
                 "Sélection de fonds à rapport qualité/frais élevé"),
                ("Suivi actif",
                 "Rééquilibrage selon les conditions de marché"),
            ]
            for ei3b, (et3b, ex3b) in enumerate(_engagements_s3):
                _ey3b = 1.50 + ei3b * 0.90
                _add_rect(s3, 4.80, _ey3b, 0.06, 0.55, GOLD)
                _add_text(s3, et3b, 4.96, _ey3b + 0.04, 4.40, 0.24, 9, NAVY, bold=True)
                _add_text(s3, ex3b, 4.96, _ey3b + 0.28, 4.40, 0.24, 8, GREY)
            _add_text(s3,
                      "⚠️ Simulation historique. Les performances passées ne préjugent pas des performances futures.",
                      0.28, 5.00, 9.44, 0.22, 7, GREY)

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 4 — Transparence des frais
        # ══════════════════════════════════════════════════════════════════════
        s4 = prs.slides.add_slide(blank_layout)
        _add_rect(s4, 0, 0, 10, 5.625, LGREY)
        _left_stripe(s4, GOLD)
        _footer(s4)
        _add_text(s4, "Transparence des frais",
                  0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
        _add_text(s4, "Décomposition du coût total de détention de votre portefeuille",
                  0.30, 0.68, 9.30, 0.25, 9, GREY)
        # Tableau frais par couche
        _frais_headers = ["Support", "TER fonds (%/an)", "Frais contrat (%/an)", "Coût total"]
        _col_w4 = [4.2, 1.8, 1.8, 1.8]
        _col_x4 = [0.28]
        for _cw4 in _col_w4[:-1]:
            _col_x4.append(_col_x4[-1] + _cw4)
        _row_y4 = 1.10
        for ci4, (_cx4, _ch4) in enumerate(zip(_col_x4, _frais_headers)):
            _add_rect(s4, _cx4, _row_y4, _col_w4[ci4], 0.28, NAVY)
            _add_text(s4, _ch4, _cx4 + 0.05, _row_y4 + 0.04, _col_w4[ci4] - 0.10, 0.22, 8, WHITE, bold=True)
        _positions_frais = positions_df_main if not positions_df_main.empty else pd.DataFrame()
        _fee_contract = report.get("fee_contract_pct", 0.6)
        if not _positions_frais.empty:
            for ri4, row4 in enumerate(_positions_frais.head(7).itertuples(index=False)):
                _ry4 = _row_y4 + 0.28 + ri4 * 0.30
                _bg4 = WHITE if ri4 % 2 == 0 else LGREY
                _add_rect(s4, 0.28, _ry4, 9.60, 0.30, _bg4)
                _nom4 = str(getattr(row4, "Nom", "—"))[:42]
                _ter4 = float(getattr(row4, "TER_%", getattr(row4, "fee_uc_pct", 0)) or 0)
                _tot4 = _ter4 + float(_fee_contract or 0)
                _vals4 = [_nom4, f"{_ter4:.2f}%", f"{_fee_contract:.2f}%", f"{_tot4:.2f}%"]
                for ci4b, (_cx4b, _rv4) in enumerate(zip(_col_x4, _vals4)):
                    _tc4 = RGBColor(0xCC, 0x22, 0x00) if (ci4b == 3 and _tot4 > 2.0) else BLACK
                    _add_text(s4, _rv4, _cx4b + 0.05, _ry4 + 0.05, _col_w4[ci4b] - 0.10, 0.22, 8, _tc4)
        else:
            _add_text(s4, "Données de frais non disponibles — ajoutez des fonds pour les voir ici.",
                      0.28, 1.50, 9.44, 0.40, 10, GREY)
        # Alerte Clean Shares si assureur Spirica
        _assureur_s4 = CONTRACTS_REGISTRY.get(
            report.get("contrat_label", ""), {}
        ).get("assureur", "")
        if _assureur_s4 == "Spirica":
            _add_rect(s4, 0.28, 3.10, 9.44, 0.44, RGBColor(0xFF, 0xF8, 0xE1), GOLD)
            _add_text(s4,
                      "✨ Optimisation possible — Spirica : accès aux parts institutionnelles (I) "
                      "sur certains fonds éligibles → économie potentielle de −0,5% à −1%/an.",
                      0.38, 3.15, 9.00, 0.34, 8, RGBColor(0x7B, 0x5E, 0x00))
        # Encart "Impact des frais sur 10 ans"
        _add_rect(s4, 0.28, 3.62, 9.44, 1.30, NAVY)
        _add_text(s4, "Impact des frais sur 10 ans (simulation)",
                  0.38, 3.68, 9.00, 0.28, 10, GOLD, bold=True)
        _cap = float(val or 10000.0)
        _fee_total_est = float(_fee_contract or 0.6)
        _brut_10 = _cap * (1.07 ** 10)
        _net_10  = _cap * ((1.07 - _fee_total_est / 100) ** 10)
        _drag_10 = _brut_10 - _net_10
        _add_text(s4,
                  f"Capital initial {_fe(_cap)} · Rendement brut hypothétique 7%/an",
                  0.38, 3.94, 9.00, 0.22, 8, ICE)
        _add_text(s4,
                  f"Valeur brute : {_fe(_brut_10)}   →   Valeur nette de frais : {_fe(_net_10)}",
                  0.38, 4.14, 9.00, 0.22, 8, ICE)
        _add_text(s4,
                  f"Coût cumulé des frais sur 10 ans : {_fe(_drag_10)}",
                  0.38, 4.42, 5.00, 0.30, 14, GOLD, bold=True)

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 4bis — Fiscalité & Levier fiscal
        # ══════════════════════════════════════════════════════════════════════
        s4b = prs.slides.add_slide(blank_layout)
        _add_rect(s4b, 0, 0, 10, 5.625, WHITE)
        _left_stripe(s4b, BLUE)
        _footer(s4b)
        _add_text(s4b, "Levier fiscal — Assurance-vie",
                  0.30, 0.15, 9.30, 0.55, 20, NAVY, bold=True)
        _cli_label = f"pour {nom_cli}" if nom_cli else ""
        _add_text(s4b,
                  f"Optimisation de l'abattement annuel {_cli_label} — Art. 158-6 CGI",
                  0.30, 0.68, 9.30, 0.25, 9, GREY)
        # Bloc abattement annuel
        _sit_fam = report.get("situation_familiale", "Célibataire / veuf / divorcé")
        _is_couple = "Couple" in str(_sit_fam)
        _abat_annuel = 9200.0 if _is_couple else 4600.0
        _abat_label  = "9 200 € (couple)" if _is_couple else "4 600 € (célibataire)"
        _add_rect(s4b, 0.28, 1.10, 4.50, 2.50, LGREY, BORDER)
        _add_text(s4b, "Abattement annuel disponible",
                  0.36, 1.16, 4.34, 0.28, 10, NAVY, bold=True)
        _add_text(s4b, _abat_label,
                  0.36, 1.48, 4.34, 0.42, 24, GREEN, bold=True)
        _add_text(s4b,
                  "Sur la quote-part de gains dans chaque rachat\n"
                  "Après 8 ans de détention du contrat\n"
                  "Aucun IR sur les gains dans cette limite",
                  0.36, 1.96, 4.34, 0.70, 8, GREY)
        _add_rect(s4b, 0.36, 2.75, 0.06, 0.60, GOLD)
        _add_text(s4b, "Stratégie recommandée",
                  0.52, 2.78, 4.00, 0.24, 9, NAVY, bold=True)
        _add_text(s4b,
                  f"Rachat annuel optimisé pour utiliser l'abattement\n"
                  f"de {_abat_label} sans générer d'IR",
                  0.52, 3.02, 4.00, 0.40, 8, GREY)
        # Bloc projection économie fiscale
        _add_rect(s4b, 5.00, 1.10, 4.72, 2.50, ICE, BORDER)
        _add_text(s4b, "Économie IR estimée sur 10 ans",
                  5.08, 1.16, 4.56, 0.28, 10, NAVY, bold=True)
        _tmi_ref = 0.30
        _eco_ir_annual = _abat_annuel * _tmi_ref
        _eco_ir_10ans  = _eco_ir_annual * 10
        _add_text(s4b, f"{_fe(_eco_ir_10ans)}",
                  5.08, 1.48, 4.56, 0.48, 28, GOLD, bold=True)
        _add_text(s4b,
                  f"Base : TMI 30% · {_abat_label}/an · 10 ans",
                  5.08, 1.96, 4.56, 0.22, 8, GREY)
        _eco_rows = [
            (f"Économie IR/an",     f"{_fe(_eco_ir_annual)}"),
            (f"Économie IR/10 ans", f"{_fe(_eco_ir_10ans)}"),
            ("PS inévitables",     "17,2% sur la fraction de gains"),
        ]
        for eri, (el, ev) in enumerate(_eco_rows):
            _ery = 2.25 + eri * 0.42
            _add_text(s4b, el, 5.08, _ery, 2.50, 0.30, 8, GREY)
            _add_text(s4b, ev, 7.60, _ery, 2.00, 0.30, 9, NAVY, bold=True, align="right")
        # Bloc transmission
        _add_rect(s4b, 0.28, 3.75, 9.44, 1.15, NAVY)
        _add_text(s4b, "Transmission — Art. 990I",
                  0.38, 3.80, 9.00, 0.28, 10, GOLD, bold=True)
        _add_text(s4b,
                  "Versements avant 70 ans : 152 500 € exonérés par bénéficiaire · "
                  "Hors succession · Transmission facilitée entre générations",
                  0.38, 4.08, 9.00, 0.30, 8, ICE)
        _add_text(s4b,
                  "⚠️ Simulation indicative. Consultez un fiscaliste pour votre situation personnelle.",
                  0.38, 4.50, 9.00, 0.22, 7, GREY)

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 4b — Diversification (Niveau 1 uniquement — pas de ratios)
        # ══════════════════════════════════════════════════════════════════════
        _pptx_div_A = report.get("diversification_A")
        _pptx_div_B = report.get("diversification_B")
        _pptx_risk_A = report.get("risk_A")
        _pptx_risk_B = report.get("risk_B")
        _has_pptx_div = _pptx_div_A is not None or _pptx_div_B is not None

        if _has_pptx_div:
            s_div = prs.slides.add_slide(blank_layout)
            _add_rect(s_div, 0, 0, 10, 5.625, WHITE)
            _add_rect(s_div, 0, 0, 0.18, 5.625, NAVY)
            _add_rect(s_div, 0.18, 0, 9.82, 0.72, NAVY)
            _add_text(s_div, "Diversification de votre portefeuille",
                      0.42, 0.12, 9.20, 0.52, 18, WHITE, bold=True)
            _add_text(s_div, "  ".join(filter(None, [nom_cab, as_of])),
                      0.42, 5.38, 9.00, 0.20, 7, GREY, align="center")

            def _pptx_div_col(slide, div_res, risk_res, x, w, col_label):
                _add_text(slide, col_label, x, 0.82, w, 0.28, 11, NAVY, bold=True)
                _y = 1.15
                if div_res is None:
                    _add_text(slide, "Données insuffisantes (min. 2 lignes).", x, _y, w, 0.28, 9, GREY)
                    return
                _n_tot = div_res["n_lines"]
                _n_eff = div_res.get("n_effective", _n_tot)
                _avg = div_res["avg_corr"]
                if _n_eff == _n_tot:
                    _dlbl = f"{_n_tot} lignes → {_n_eff} sources réelles"
                    _dcol = GREEN
                elif _n_eff >= _n_tot - 1:
                    _dlbl = f"{_n_tot} lignes → {_n_eff} sources réelles"
                    _dcol = RGBColor(0xF0, 0x8C, 0x00)
                else:
                    _dlbl = f"{_n_tot} lignes → seulement {_n_eff} sources réelles"
                    _dcol = RGBColor(0xCC, 0x00, 0x00)
                _add_text(slide, _dlbl, x, _y, w, 0.40, 18, _dcol, bold=True)
                _y += 0.50
                _add_text(slide,
                          f"Corrélation moyenne : {_avg:.0%}  |  {_n_tot} lignes analysées",
                          x, _y, w, 0.24, 8, GREY)
                _y += 0.30
                for _ni, _nj, _c in div_res["doublons"]:
                    _add_text(slide, f"⚠ Doublon : {_ni} / {_nj} ({_c:.0%})", x, _y, w, 0.22, 8, RGBColor(0xCC, 0x00, 0x00))
                    _y += 0.25
                for _ni, _nj, _c in div_res["vigilance"][:2]:
                    _add_text(slide, f"~ Similaire : {_ni} / {_nj} ({_c:.0%})", x, _y, w, 0.22, 8, RGBColor(0xF0, 0x8C, 0x00))
                    _y += 0.25
                if risk_res:
                    _y = max(_y, 2.80)
                    _add_text(slide, f"Volatilité annuelle : {risk_res['vol_ann_pct']:.1f}%", x, _y, w, 0.22, 9, NAVY)
                    _y += 0.26
                    _add_text(slide, f"Pire baisse historique : {risk_res['max_dd_pct']:.1f}%", x, _y, w, 0.22, 9, NAVY)

            _pptx_mode = report.get("mode", "compare")
            if _pptx_mode == "compare":
                _add_rect(s_div, 5.28, 0.78, 0.02, 4.50, LGREY)  # séparateur
                _pptx_div_col(s_div, _pptx_div_A, _pptx_risk_A, 0.42, 4.60, "Portefeuille Client")
                _pptx_div_col(s_div, _pptx_div_B, _pptx_risk_B, 5.40, 4.30, "Proposition Cabinet")
            elif _pptx_mode == "valority":
                _pptx_div_col(s_div, _pptx_div_B, _pptx_risk_B, 0.42, 9.20, "Portefeuille Cabinet")
            else:
                _pptx_div_col(s_div, _pptx_div_A, _pptx_risk_A, 0.42, 9.20, "Portefeuille Client")

        # ══════════════════════════════════════════════════════════════════════
        # SLIDE 5 — Conclusion / Engagements
        # ══════════════════════════════════════════════════════════════════════
        s5 = prs.slides.add_slide(blank_layout)
        _add_rect(s5, 0, 0, 10, 5.625, NAVY)
        _add_rect(s5, 0, 0, 0.18, 5.625, GOLD)

        _add_text(s5, "Votre conseiller\nà vos côtés",
                  0.75, 0.45, 5.50, 1.10, 28, WHITE, bold=True)
        _add_rect(s5, 0.75, 1.70, 2.00, 0.04, GOLD)

        _engagements = [
            ("Transparence",
             "Simulations basées sur les frais réels de votre contrat."),
            ("Indépendance",
             "Sélection de fonds fondée sur la performance."),
            ("Suivi",
             "Point annuel sur l'évolution de votre allocation."),
        ]
        for ei, (etitle, etxt) in enumerate(_engagements):
            _ey = 1.90 + ei * 0.95
            _add_rect(s5, 0.75, _ey, 0.05, 0.60, GOLD)
            _add_text(s5, etitle, 0.90, _ey + 0.04, 4.50, 0.28, 12, WHITE, bold=True)
            _add_text(s5, etxt,   0.90, _ey + 0.32, 4.50, 0.28, 9, ICE)

        _footer_parts = [nom_cab, "Document confidentiel"]
        if nom_cli:
            _footer_parts.append(nom_cli)
        if as_of:
            _footer_parts.append(as_of)
        _add_text(s5, "  ·  ".join(_footer_parts),
                  0.28, 5.38, 9.44, 0.20, 7, GREY, align="center")

        # ── Sérialisation ─────────────────────────────────────────────────────
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()


    def _years_between(d0: pd.Timestamp, d1: pd.Timestamp) -> float:
        return max(0.0, (d1 - d0).days / 365.25)


    report_data = {
        "as_of": fmt_date(TODAY),
        "mode": st.session_state.get("MODE_ANALYSE", "compare"),
    }
    report_data["nom_client"] = st.session_state.get("NOM_CLIENT", "").strip()
    report_data["nom_cabinet"] = st.session_state.get("NOM_CABINET", "").strip()
    report_data["fee_contract_pct"] = st.session_state.get("FEE_A", 0.6)
    report_data["situation_familiale"] = st.session_state.get(
        "tax_situation_familiale", "Célibataire / veuf / divorcé"
    )
    report_data["contrat_label"] = st.session_state.get("CONTRACT_LABEL_A", "")
    report_data["date_ouverture"] = ""
    _all_lines = st.session_state.get("A_lines", [])
    if _all_lines:
        _dates = [ln.get("buy_date") for ln in _all_lines if ln.get("buy_date")]
        if _dates:
            _min_date = min(pd.Timestamp(d) for d in _dates)
            report_data["date_ouverture"] = _min_date.strftime("%d/%m/%Y")

    mode_report = report_data["mode"]
    df_client_lines = build_positions_dataframe("A_lines") if show_client else pd.DataFrame()
    df_valority_lines = build_positions_dataframe("B_lines") if show_valority else pd.DataFrame()

    report_data["df_client_lines"] = df_client_lines
    report_data["df_valority_lines"] = df_valority_lines
    report_data["positions_df_client"] = df_client_lines
    report_data["positions_df_valority"] = df_valority_lines
    report_data["dfA_val"] = (
        dfA.reset_index().rename(columns={"index": "Date"}) if (show_client and not dfA.empty) else pd.DataFrame()
    )
    report_data["dfB_val"] = (
        dfB.reset_index().rename(columns={"index": "Date"}) if (show_valority and not dfB.empty) else pd.DataFrame()
    )
    report_data["client_summary"] = {
        "val": valA,
        "net": netA,
        "brut": brutA,
        "perf_tot_pct": perf_tot_client or 0.0,
        "irr_pct": xirrA or 0.0,
    }
    report_data["valority_summary"] = {
        "val": valB,
        "net": netB,
        "brut": brutB,
        "perf_tot_pct": perf_tot_valority or 0.0,
        "irr_pct": xirrB or 0.0,
    }
    report_data["comparison"] = {
        "delta_val": (valB - valA) if (valA is not None and valB is not None) else 0.0,
        "delta_perf_pct": (perf_tot_valority - perf_tot_client)
        if (perf_tot_client is not None and perf_tot_valority is not None)
        else 0.0,
    }

    df_map: Dict[str, pd.DataFrame] = {}
    if mode_report in ("compare", "client") and not dfA.empty:
        df_map["Client"] = dfA
    if mode_report in ("compare", "valority") and not dfB.empty:
        df_map["Cabinet"] = dfB
    report_data["df_map"] = df_map

    if mode_report == "compare":
        positions_df = df_valority_lines if not df_valority_lines.empty else df_client_lines
        lines = st.session_state.get("B_lines", []) or st.session_state.get("A_lines", [])
    else:
        if mode_report == "valority":
            positions_df = df_valority_lines
            lines = st.session_state.get("B_lines", [])
            start_min = startB_min
            brut = brutB
            net = netB
            val = valB
        else:
            positions_df = df_client_lines
            lines = st.session_state.get("A_lines", [])
            start_min = startA_min
            brut = brutA
            net = netA
            val = valA

        years = _years_between(start_min, TODAY) if isinstance(start_min, pd.Timestamp) else 0.0
        fees_paid = max(0.0, brut - net) if brut is not None and net is not None else 0.0
        value_created = (val - net) if val is not None and net is not None else 0.0
        value_per_year = (value_created / years) if years > 0 else 0.0
        report_data["fees_analysis"] = {
            "fees_paid": fees_paid,
            "value_created": value_created,
            "value_per_year": value_per_year,
        }

    def _build_lines_with_values(
        lines_src: List[Dict[str, Any]],
        positions_src: pd.DataFrame,
    ) -> List[Dict[str, Any]]:
        items: List[Dict[str, Any]] = []
        if isinstance(positions_src, pd.DataFrame) and not positions_src.empty:
            for ln in lines_src:
                isin = ln.get("isin", "")
                name = ln.get("name", "")
                match = positions_src[
                    (positions_src["Nom"] == name) & (positions_src["ISIN / Code"] == isin)
                ]
                val = float(match["Valeur actuelle €"].iloc[0]) if not match.empty else 0.0
                items.append({"isin": isin, "value": val})
        return items

    report_data["positions_df"] = positions_df
    report_data["lines"] = _build_lines_with_values(lines, positions_df)
    report_data["lines_client"] = _build_lines_with_values(
        st.session_state.get("A_lines", []),
        df_client_lines,
    )
    report_data["lines_valority"] = _build_lines_with_values(
        st.session_state.get("B_lines", []),
        df_valority_lines,
    )

    st.session_state["REPORT_DATA"] = report_data
    # Pré-générer le PDF une seule fois (évite double-rendu au 1er chargement)
    try:
        _pdf_bytes_cache = generate_pdf_report(report_data)
    except Exception:
        _pdf_bytes_cache = b""
    st.session_state["PDF_BYTES_CACHE"] = _pdf_bytes_cache

    # ------------------------------------------------------------
    # Bloc final : Comparaison OU "Frais & valeur créée"
    # ------------------------------------------------------------
    mode = st.session_state.get("MODE_ANALYSE", "compare")

    # ============================
    # CAS 1 — MODE COMPARAISON
    # ============================
    if mode == "compare":
        st.markdown("---")
        st.subheader("⚖️ Comparaison")

        gain_vs_client = (valB - valA) if (valA is not None and valB is not None) else 0.0
        delta_xirr = (xirrB - xirrA) if (xirrA is not None and xirrB is not None) else None
        perf_diff_tot = (
            (perf_tot_valority - perf_tot_client)
            if (perf_tot_client is not None and perf_tot_valority is not None)
            else None
        )

        with st.container(border=True):
            c1, c2, c3 = st.columns(3)
            with c1:
                st.metric("Gain en valeur", to_eur(gain_vs_client))
            with c2:
                st.metric(
                    "Surperformance totale",
                    f"{perf_diff_tot:+.2f}%" if perf_diff_tot is not None else "—",
                )
            with c3:
                st.metric(
                    "Surperformance annualisée (Î” XIRR)",
                    f"{delta_xirr:+.2f}%" if delta_xirr is not None else "—",
                )



    # ============================
    # CAS 2 — MODE ANALYSE SIMPLE
    # ============================
    else:
        # Sélection des variables selon le mode
        if mode == "valority":
            brut = brutB
            net = netB
            val = valB
            start_min = startB_min
            irr = xirrB
            fee_pct = st.session_state.get("FEE_B", 0.0)
            title = "🏢 Allocation Cabinet — Frais & valeur créée"
        else:  # mode == "client"
            brut = brutA
            net = netA
            val = valA
            start_min = startA_min
            irr = xirrA
            fee_pct = st.session_state.get("FEE_A", 0.0)
            title = "🧍 Portefeuille — Frais & valeur créée"

        st.markdown("---")
        st.subheader("📌 Analyse : frais & valeur créée")

        if brut > 0 and net >= 0 and val >= 0 and isinstance(start_min, pd.Timestamp):
            fees_paid = max(0.0, brut - net)     # frais d'entrée réellement payés
            value_created = val - net            # valeur créée vs capital réellement investi
            years = _years_between(start_min, TODAY)
            value_per_year = (value_created / years) if years > 0 else None

            with st.container(border=True):
                st.markdown(f"#### {title}")
                st.caption(
                    f"Période : **{fmt_date(start_min)} → {fmt_date(TODAY)}** "
                    f"• Frais d'entrée : **{fee_pct:.2f}%**"
                )

                _pad_l, c1, c2, c3, _pad_r = st.columns([0.5, 1, 1, 1, 0.5])
                with c1:
                    st.metric("Frais d'entrée payés", to_eur(fees_paid))
                with c2:
                    st.metric("Valeur créée (net)", to_eur(value_created))
                with c3:
                    st.metric(
                        "Valeur créée / an (moyenne)",
                        to_eur(value_per_year) if value_per_year is not None else "—",
                    )

                st.markdown(
                    f"""
- Montants versés (brut) : **{to_eur(brut)}**
- Montants réellement investis (après frais) : **{to_eur(net)}**
- Valeur actuelle : **{to_eur(val)}**
"""
                )

                if irr is not None:
                    st.markdown(f"- Rendement annualisé (XIRR) : **{irr:.2f}%**")
                else:
                    st.markdown("- Rendement annualisé (XIRR) : **—**")

                # Message factuel — différencié selon que la valeur créée est positive ou négative
                if fees_paid > 0 and value_created > 0:
                    ratio = value_created / fees_paid
                    st.markdown(
                        f"**Lecture :** {to_eur(fees_paid)} de frais d'entrée ont généré "
                        f"**{to_eur(value_created)}** de valeur nette créée à date "
                        f"(**×{ratio:.1f}**)."
                    )
                elif fees_paid > 0 and value_created <= 0:
                    st.markdown(
                        f"**Lecture :** {to_eur(fees_paid)} de frais d'entrée payés. "
                        f"Le portefeuille affiche une moins-value nette de **{to_eur(abs(value_created))}** à date."
                    )
        else:
            st.info("Ajoutez des lignes (et/ou des versements) pour afficher l'analyse frais & valeur créée.")

    # ------------------------------------------------------------
    # Indicateurs de risque & performance (2 niveaux)
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📊 Indicateurs de risque & performance")

    _bench_sym = st.session_state.get("BENCHMARK_SYMBOL", "CW8.PA")
    _euro_rate_A = st.session_state.get("EURO_RATE_A", 2.0)
    _euro_rate_B = st.session_state.get("EURO_RATE_B", 2.5)
    _linesA = st.session_state.get("A_lines", [])
    _linesB = st.session_state.get("B_lines", [])
    _fee_A = float(st.session_state.get("FEE_A", 0.0))
    _fee_B = float(st.session_state.get("FEE_B", 0.0))

    # ── Calculs diversification ────────────────────────────────────────────
    _div_A = compute_diversification_score(_linesA, _euro_rate_A) if (show_client and _linesA) else None
    _div_B = compute_diversification_score(_linesB, _euro_rate_B) if (show_valority and _linesB) else None
    _risk_A = portfolio_risk_stats(_linesA, _euro_rate_A, fee_pct=_fee_A) if (show_client and _linesA) else None
    _risk_B = portfolio_risk_stats(_linesB, _euro_rate_B, fee_pct=_fee_B) if (show_valority and _linesB) else None

    # ── Calculs ratios techniques ──────────────────────────────────────────
    _sharpe_A = compute_sharpe_ratio(_linesA, _euro_rate_A, _fee_A) if (show_client and _linesA) else None
    _sharpe_B = compute_sharpe_ratio(_linesB, _euro_rate_B, _fee_B) if (show_valority and _linesB) else None
    _sortino_A = compute_sortino_ratio(_linesA, _euro_rate_A, _fee_A) if (show_client and _linesA) else None
    _sortino_B = compute_sortino_ratio(_linesB, _euro_rate_B, _fee_B) if (show_valority and _linesB) else None
    _ba_A = compute_beta_alpha(_linesA, _euro_rate_A, _fee_A, _bench_sym) if (show_client and _linesA) else None
    _ba_B = compute_beta_alpha(_linesB, _euro_rate_B, _fee_B, _bench_sym) if (show_valority and _linesB) else None

    # Stocker le bêta pour le stress-test
    st.session_state["BETA_AUTO_A"] = _ba_A["beta"] if _ba_A else None
    st.session_state["BETA_AUTO_B"] = _ba_B["beta"] if _ba_B else None

    # Indicateurs de risque pour le PDF — inséré ICI après calcul des variables
    report_data["diversification_A"] = _div_A
    report_data["diversification_B"] = _div_B
    report_data["risk_A"] = _risk_A
    report_data["risk_B"] = _risk_B
    report_data["sharpe_A"] = _sharpe_A
    report_data["sharpe_B"] = _sharpe_B
    report_data["sortino_A"] = _sortino_A
    report_data["sortino_B"] = _sortino_B
    report_data["beta_alpha_A"] = _ba_A
    report_data["beta_alpha_B"] = _ba_B
    report_data["ratios_A"] = {
        "sharpe": _sharpe_A,
        "sortino": _sortino_A,
        "beta_alpha": _ba_A,
    }
    report_data["ratios_B"] = {
        "sharpe": _sharpe_B,
        "sortino": _sortino_B,
        "beta_alpha": _ba_B,
    }

    _mode_risk = st.session_state.get("MODE_ANALYSE", "compare")

    # ══════════════════════════════════════════════════════════════════════
    # ① DIVERSIFICATION
    # ══════════════════════════════════════════════════════════════════════
    st.markdown("##### 🧩 Diversification du portefeuille")

    def _render_diversification(div_res, label):
        """Affiche la diversification d'un portefeuille en langage clair."""
        if div_res is None:
            st.info("Ajoutez au moins 2 lignes pour analyser la diversification.")
            return
        _n = div_res["n_lines"]
        _n_eff = div_res.get("n_effective", _n)
        _avg = div_res["avg_corr"]

        # ── Fonds réellement utiles ──
        if _n_eff == _n:
            st.markdown(f"✅ **{_n} lignes → {_n_eff} sources de diversification réelles**")
            st.caption("Chaque fonds apporte une exposition distincte au portefeuille.")
        elif _n_eff >= _n - 1:
            st.markdown(f"🟠 **{_n} lignes → {_n_eff} sources de diversification réelles**")
            st.caption(f"{_n - _n_eff} fonds est redondant avec un autre — il double les frais sans améliorer la diversification.")
        else:
            st.markdown(f"🔴 **{_n} lignes → seulement {_n_eff} sources de diversification réelles**")
            st.caption(f"{_n - _n_eff} fonds sont redondants — vous payez des frais sur {_n} lignes pour la diversification de {_n_eff}.")

        # ── Diversification globale ──
        if _avg < 0.30:
            st.caption(f"🟢 Bonne couverture multi-classes (corrélation moyenne : {_avg:.0%}) — les actifs se complètent bien.")
        elif _avg < 0.60:
            st.caption(f"🟠 Couverture correcte (corrélation moyenne : {_avg:.0%}) — des améliorations sont possibles.")
        else:
            st.caption(f"🔴 Couverture insuffisante (corrélation moyenne : {_avg:.0%}) — les actifs évoluent trop dans le même sens.")

        # ── Doublons nommés ──
        for _ni, _nj, _c in div_res["doublons"]:
            st.warning(
                f"⚠️ **{_ni}** et **{_nj}** sont quasi-identiques (corrélation {_c:.0%}). "
                f"Garder les deux double les frais sans apporter de diversification."
            )
        for _ni, _nj, _c in div_res["vigilance"]:
            st.caption(f"🟠 {_ni} et {_nj} ont un comportement très similaire (corrélation {_c:.0%}).")

    if _mode_risk == "compare":
        _col_div_a, _col_div_b = st.columns(2)
        with _col_div_a:
            st.markdown("**🧍 Client**")
            _render_diversification(_div_A, "Client")
        with _col_div_b:
            st.markdown("**🏢 Cabinet**")
            _render_diversification(_div_B, "Cabinet")
        # Narratif comparatif
        if _div_A and _div_B:
            with st.container(border=True):
                _eff_A = _div_A["n_effective"]
                _eff_B = _div_B["n_effective"]

                if _eff_B > _eff_A and _eff_A > 0:
                    _ratio = _eff_B / _eff_A
                    if _ratio >= 2.0:
                        _ratio_txt = f"{_ratio:.0f}× plus de sources"
                    else:
                        _pct = (_ratio - 1) * 100
                        _ratio_txt = f"{_pct:.0f}% de sources en plus"
                    st.success(
                        f"💬 **Concrètement, la proposition du cabinet offre {_ratio_txt} "
                        f"de diversification réelles** ({_eff_B} contre {_eff_A} actuellement). "
                        f"Quand un secteur ou une classe d'actifs subit une baisse, "
                        f"les autres lignes du portefeuille ne sont pas entraînées dans la chute "
                        f"— le capital est mieux protégé et la reprise est plus rapide."
                    )
                elif _eff_A > _eff_B and _eff_B > 0:
                    st.warning(
                        f"💬 Votre portefeuille actuel est mieux diversifié "
                        f"({_eff_A} sources réelles contre {_eff_B} pour la proposition). "
                        f"La proposition compense peut-être par une meilleure sélection de fonds "
                        f"— vérifiez les indicateurs de rendement ci-dessous."
                    )
                else:
                    st.info(
                        f"💬 Les deux portefeuilles ont une diversification comparable "
                        f"({_eff_A} sources réelles). L'avantage de la proposition se joue "
                        f"sur d'autres critères (rendement, frais, protection baissière)."
                    )
    elif _mode_risk == "client":
        _render_diversification(_div_A, "Client")
    else:
        _render_diversification(_div_B, "Cabinet")

    # ══════════════════════════════════════════════════════════════════════
    # ② RENDEMENT, RISQUE & VALEUR AJOUTÉE (langage client)
    # ══════════════════════════════════════════════════════════════════════
    st.markdown("---")
    st.markdown("##### 💡 Rendement, risque & valeur ajoutée")

    _VS_DIV = "<div style='text-align:center; padding-top:1.4rem; font-size:1.1rem; color:#aaa;'>vs</div>"

    def _render_single_ratios_narrative(sharpe, sortino, ba, euro_rate, val, risk, label):
        """Affiche les ratios en langage client pour un portefeuille seul."""
        # ── Rendement/risque vs fonds euros ──
        if sharpe is not None:
            with st.container(border=True):
                st.markdown("**Le rendement justifie-t-il le risque pris ?**")
                _gain_sh = sharpe * 1000
                st.metric("Rendement excédentaire", f"{_gain_sh:,.0f} €", delta="pour 1 000 € de risque pris")
                if sharpe >= 1.0:
                    _sh_comment = "C'est un bon ratio — au-dessus de 1 000 € c'est excellent."
                elif sharpe >= 0.5:
                    _sh_comment = "C'est correct — au-dessus de 1 000 € c'est considéré excellent, en-dessous de 500 € le risque n'est pas assez rémunéré."
                else:
                    _sh_comment = "Le risque pris n'est pas suffisamment rémunéré. En-dessous de 500 €, il est légitime de se demander si le jeu en vaut la chandelle."
                st.caption(
                    f"💬 Par rapport au fonds euros ({euro_rate:.2f}%/an, sans risque), {label} "
                    f"génère **{_gain_sh:,.0f} €** de rendement supplémentaire pour chaque 1 000 € de risque accepté. "
                    f"{_sh_comment}"
                )
        # ── Protection baissière ──
        if sortino and sortino > 0:
            with st.container(border=True):
                st.markdown("**Protection en cas de baisse**")
                _perte = 1.0 / sortino
                st.metric("Perte par euro gagné", f"{_perte:.2f} €", delta="en phase de baisse")
                if _perte < 0.7:
                    _so_comment = "C'est une bonne protection."
                elif _perte < 1.0:
                    _so_comment = "C'est correct."
                else:
                    _so_comment = "Le portefeuille est vulnérable aux baisses — les pertes dépassent les gains."
                st.caption(
                    f"💬 Pour chaque 1 € gagné dans les bonnes phases, {label} "
                    f"perd {_perte:.2f} € dans les mauvaises. {_so_comment}"
                )
        # ── Alpha ──
        _al = ba["alpha_pct"] if ba else None
        if _al is not None and abs(_al) > 0.5:
            with st.container(border=True):
                _bench_name = st.session_state.get("BENCHMARK_LABEL", "MSCI World").split(" (")[0]
                st.markdown("**Valeur ajoutée de l'allocation**")
                st.metric("Performance vs marché", f"{_al:+.2f} %/an",
                          delta="surperformance" if _al > 0 else "sous-performance")
                _al_comment = "l'allocation ajoute de la valeur." if _al > 0 else "une optimisation est possible."
                st.caption(
                    f"💬 Comparé à un simple ETF {_bench_name} au même niveau de risque, "
                    f"{label} fait **{_al:+.1f}%/an** "
                    f"{'de plus — ' + _al_comment if _al > 0 else 'de moins — ' + _al_comment}"
                )

    if _mode_risk == "compare" and _sharpe_A is not None and _sharpe_B is not None:
        # ── Rendement vs risque ──
        with st.container(border=True):
            st.markdown("**Le rendement justifie-t-il le risque pris ?**")
            _col_rr_a, _col_rr_vs, _col_rr_b = st.columns([5, 1, 5])
            with _col_rr_a:
                _gain_A_sh = _sharpe_A * 1000
                st.metric("🧍 Client", f"{_gain_A_sh:,.0f} €", delta="pour 1 000 € de risque")
            with _col_rr_vs:
                st.markdown(_VS_DIV, unsafe_allow_html=True)
            with _col_rr_b:
                _gain_B_sh = _sharpe_B * 1000
                st.metric("🏢 Cabinet", f"{_gain_B_sh:,.0f} €", delta="pour 1 000 € de risque")
            _avantage_pct = ((_sharpe_B / _sharpe_A) - 1) * 100 if _sharpe_A > 0 else 0
            if _sharpe_B > _sharpe_A:
                st.caption(
                    f"💬 À risque comparable, la proposition cabinet génère **{_avantage_pct:+.0f}%** de rendement "
                    f"excédentaire en plus par rapport au fonds euros. "
                    f"Concrètement, pour le même niveau de risque, chaque euro investi travaille {abs(_avantage_pct):.0f}% plus efficacement."
                )
            else:
                st.caption(
                    f"💬 Le portefeuille actuel génère un meilleur rendement par unité de risque ({_gain_A_sh:,.0f} € vs {_gain_B_sh:,.0f} €)."
                )

        # ── Protection à la baisse ──
        if _sortino_A is not None and _sortino_B is not None and _sortino_A > 0 and _sortino_B > 0:
            with st.container(border=True):
                st.markdown("**Comment le portefeuille résiste-t-il aux baisses ?**")
                _perte_A = 1.0 / _sortino_A
                _perte_B = 1.0 / _sortino_B
                _col_pb_a, _col_pb_vs, _col_pb_b = st.columns([5, 1, 5])
                with _col_pb_a:
                    st.metric("🧍 Client", f"{_perte_A:.2f} €", delta="perdu pour 1 € gagné en hausse")
                with _col_pb_vs:
                    st.markdown(_VS_DIV, unsafe_allow_html=True)
                with _col_pb_b:
                    st.metric("🏢 Cabinet", f"{_perte_B:.2f} €", delta="perdu pour 1 € gagné en hausse")
                _reduction = (1 - _perte_B / _perte_A) * 100 if _perte_A > 0 else 0
                if _reduction > 0:
                    st.caption(
                        f"💬 Dans les phases de baisse, la proposition limite les pertes : "
                        f"**{_reduction:.0f}% de pertes en moins** par rapport au portefeuille actuel."
                    )
                elif _reduction < -10:
                    st.caption(
                        f"💬 Le portefeuille actuel résiste mieux aux baisses "
                        f"({_perte_A:.2f} € vs {_perte_B:.2f} € perdu par euro gagné)."
                    )

        # ── Valeur ajoutée (Alpha) ──
        _alpha_A = _ba_A["alpha_pct"] if _ba_A else None
        _alpha_B = _ba_B["alpha_pct"] if _ba_B else None
        if _alpha_A is not None and _alpha_B is not None:
            with st.container(border=True):
                st.markdown("**L'allocation crée-t-elle de la valeur ?**")
                _bench_label_disp = st.session_state.get("BENCHMARK_LABEL", "MSCI World (CW8.PA)").split(" (")[0]
                _col_al_a, _col_al_vs, _col_al_b = st.columns([5, 1, 5])
                with _col_al_a:
                    st.metric("🧍 Client", f"{_alpha_A:+.2f} %/an",
                              delta="surperformance" if _alpha_A > 0 else "sous-performance")
                with _col_al_vs:
                    st.markdown(_VS_DIV, unsafe_allow_html=True)
                with _col_al_b:
                    st.metric("🏢 Cabinet", f"{_alpha_B:+.2f} %/an",
                              delta="surperformance" if _alpha_B > 0 else "sous-performance")
                _alpha_conclu = "La proposition crée davantage de valeur ajoutée." if _alpha_B > _alpha_A + 0.5 else ""
                st.caption(
                    f"💬 Par rapport à un simple ETF {_bench_label_disp} avec le même niveau de risque : "
                    f"le portefeuille client fait **{_alpha_A:+.1f}%/an** "
                    f"{'de plus' if _alpha_A > 0 else 'de moins'}, "
                    f"la proposition fait **{_alpha_B:+.1f}%/an** "
                    f"{'de plus' if _alpha_B > 0 else 'de moins'}. "
                    f"{_alpha_conclu}"
                )

        # ── Verdict global ──
        with st.container(border=True):
            _wins = []
            _losses = []
            if _sharpe_B > _sharpe_A + 0.05:
                _r_sh = abs(((_sharpe_B / _sharpe_A) - 1) * 100) if _sharpe_A > 0 else 0
                _wins.append(f"rendement/risque {_r_sh:.0f}% plus efficace")
            elif _sharpe_A > _sharpe_B + 0.05:
                _losses.append("rendement/risque")
            if _sortino_A and _sortino_B and _sortino_B > _sortino_A + 0.05:
                _r_so = (1 - (1 / _sortino_B) / (1 / _sortino_A)) * 100 if _sortino_A > 0 else 0
                _wins.append(f"{abs(_r_so):.0f}% de pertes en moins dans les baisses")
            elif _sortino_A and _sortino_B and _sortino_A > _sortino_B + 0.05:
                _losses.append("protection baissière")
            _alpha_A_v = _ba_A["alpha_pct"] if _ba_A else None
            _alpha_B_v = _ba_B["alpha_pct"] if _ba_B else None
            if _alpha_B_v is not None and _alpha_A_v is not None and _alpha_B_v > _alpha_A_v + 0.5:
                _wins.append(f"alpha supérieur de {_alpha_B_v - _alpha_A_v:.1f} pts/an")
            if len(_wins) >= 2:
                st.success(f"✅ **La proposition cabinet est plus efficace** : {', '.join(_wins)}.")
            elif len(_wins) == 1:
                st.info(f"Amélioration côté Cabinet : {_wins[0]}."
                        + (f" Point d'attention : {', '.join(_losses)}." if _losses else ""))
            elif _losses:
                st.warning(f"⚠️ Le portefeuille actuel fait mieux sur : {', '.join(_losses)}.")
            else:
                st.info("Profils de risque comparables. L'intérêt du changement repose sur d'autres critères (diversification, frais).")

    elif _mode_risk == "client" and _sharpe_A is not None:
        _render_single_ratios_narrative(_sharpe_A, _sortino_A, _ba_A, _euro_rate_A, valA, _risk_A, "votre portefeuille")

    elif _mode_risk == "valority" and _sharpe_B is not None:
        _render_single_ratios_narrative(_sharpe_B, _sortino_B, _ba_B, _euro_rate_B, valB, _risk_B, "ce portefeuille")

    # ── Détails techniques (expander fermé, pour le CGP) ──────────────────
    with st.expander("📐 Détails techniques — Ratios", expanded=False):
        st.caption("Données brutes à destination du conseiller.")
        _rows_tech = []
        if _sharpe_A is not None or _sharpe_B is not None:
            _rows_tech.append({
                "Indicateur": "Sharpe",
                "Client": f"{_sharpe_A:.2f}" if _sharpe_A is not None else "—",
                "Cabinet": f"{_sharpe_B:.2f}" if _sharpe_B is not None else "—",
            })
        if _sortino_A is not None or _sortino_B is not None:
            _rows_tech.append({
                "Indicateur": "Sortino",
                "Client": f"{_sortino_A:.2f}" if _sortino_A is not None else "—",
                "Cabinet": f"{_sortino_B:.2f}" if _sortino_B is not None else "—",
            })
        if _ba_A or _ba_B:
            _bench_lbl_tech = st.session_state.get("BENCHMARK_LABEL", "MSCI World")
            _rows_tech.append({
                "Indicateur": f"Bêta (vs {_bench_lbl_tech})",
                "Client": f"{_ba_A['beta']:.2f}" if _ba_A else "—",
                "Cabinet": f"{_ba_B['beta']:.2f}" if _ba_B else "—",
            })
            _rows_tech.append({
                "Indicateur": "Alpha",
                "Client": f"{_ba_A['alpha_pct']:+.2f}%" if _ba_A else "—",
                "Cabinet": f"{_ba_B['alpha_pct']:+.2f}%" if _ba_B else "—",
            })
        if _risk_A or _risk_B:
            _rows_tech.append({
                "Indicateur": "Volatilité ann.",
                "Client": f"{_risk_A['vol_ann_pct']:.2f}%" if _risk_A else "—",
                "Cabinet": f"{_risk_B['vol_ann_pct']:.2f}%" if _risk_B else "—",
            })
            _rows_tech.append({
                "Indicateur": "Max drawdown",
                "Client": f"{_risk_A['max_dd_pct']:.2f}%" if _risk_A else "—",
                "Cabinet": f"{_risk_B['max_dd_pct']:.2f}%" if _risk_B else "—",
            })
        if _rows_tech:
            st.dataframe(pd.DataFrame(_rows_tech), hide_index=True, use_container_width=True)
        st.caption(f"_Taux sans risque : fonds euros du contrat. Indice : {st.session_state.get('BENCHMARK_LABEL', 'MSCI World')}._")

    with st.expander("ℹ️ Comprendre ces indicateurs", expanded=False):
        st.markdown("""
**Diversification effective** — Nombre de fonds qui apportent réellement une exposition nouvelle. Un fonds corrélé à plus de 80% avec un autre ne réduit pas le risque global — il double les frais sans améliorer la protection.

**Rendement par unité de risque** — Mesure combien de rendement supplémentaire votre portefeuille génère au-delà du fonds euros (sans risque), pour chaque unité de volatilité subie. Plus c'est haut, mieux c'est — 1 000 € pour 1 000 € de risque est un bon seuil.

**Protection en phase de baisse** — Compare les gains en période haussière aux pertes en période baissière. Un portefeuille bien construit perd moins qu'il ne gagne.

**Valeur ajoutée de l'allocation (Alpha)** — Mesure si l'allocation fait mieux qu'un simple ETF indiciel avec le même niveau de risque. Un alpha positif signifie que le choix des fonds crée de la valeur au-delà du marché.

**Ratio de Sharpe / Sortino / Bêta** — Détails dans l'expander "Détails techniques" ci-dessus.

_Taux sans risque utilisé : le taux du fonds euros du contrat, car c'est l'alternative naturelle du client en assurance-vie._
""")

    # ------------------------------------------------------------
    # Tables positions
    # ------------------------------------------------------------
    st.markdown("---")
    st.subheader("📋 Positions & composition")
    if show_client:
        positions_table("Portefeuille 1 — Client", "A_lines")
    if show_valority:
        positions_table("Portefeuille 2 — Cabinet", "B_lines")


    def _render_portfolio_pie(port_key: str, title: str):
        df_positions = build_positions_dataframe(port_key)
        if df_positions.empty:
            st.info(f"{title} : Données indisponibles.")
            return
        df_pie = _prepare_pie_df(df_positions)
        if df_pie.empty:
            st.info(f"{title} : Données indisponibles.")
            return
        # Palette premium Bleu Nuit / Acier / Or / complémentaires
        _WM_PALETTE = [
            "#1B2A4A", "#4A6FA5", "#C9A84C", "#2E86AB",
            "#5C4033", "#6B8F71", "#8B5E3C", "#A0A0A0",
        ]
        df_pie["color"] = [
            _WM_PALETTE[i % len(_WM_PALETTE)] for i in range(len(df_pie))
        ]
        df_pie["Part_frac"] = df_pie["Part %"] / 100.0
        donut = (
            alt.Chart(df_pie)
            .mark_arc(innerRadius=60, outerRadius=115)
            .encode(
                theta=alt.Theta("Valeur actuelle €:Q", stack=True),
                color=alt.Color(
                    "Nom:N",
                    scale=alt.Scale(
                        domain=df_pie["Nom"].tolist(),
                        range=_WM_PALETTE[: len(df_pie)],
                    ),
                    legend=alt.Legend(title="Support", orient="right"),
                ),
                tooltip=[
                    alt.Tooltip("Nom:N", title="Support"),
                    alt.Tooltip("Valeur actuelle €:Q", title="Valeur", format=",.0f"),
                    alt.Tooltip("Part %:Q", title="Part %", format=".1f"),
                ],
            )
            .properties(title="", height=260)
        )
        st.altair_chart(donut, use_container_width=True)
        st.dataframe(
            df_pie[["Nom", "Valeur actuelle €", "Part %"]].style.format(
                {
                    "Valeur actuelle €": to_eur,
                    "Part %": "{:,.2f}%".format,
                }
            ),
            hide_index=True,
            use_container_width=True,
        )

    if show_client and show_valority:
        col_a, col_b = st.columns(2)
        with col_a:
            st.markdown("#### Portefeuille Client")
            _render_portfolio_pie("A_lines", "Portefeuille Client")
        with col_b:
            st.markdown("#### Portefeuille Cabinet")
            _render_portfolio_pie("B_lines", "Portefeuille Cabinet")
    elif show_client:
        st.markdown("#### Portefeuille Client")
        _render_portfolio_pie("A_lines", "Portefeuille Client")
    elif show_valority:
        st.markdown("#### Portefeuille Cabinet")
        _render_portfolio_pie("B_lines", "Portefeuille Cabinet")

    # ── Transparence frais & Clean Shares ────────────────────────
    st.markdown("---")
    with st.expander("🔍 Transparence des frais & Clean Shares", expanded=False):
        st.caption(
            "Décomposition des frais par couche : frais du support (TER, intégré dans la VL) "
            "et frais de l'enveloppe contrat (déduits de la performance)."
        )
        # Identifier les fonds UC du portefeuille actif
        _cs_port_key = "A_lines" if show_client else "B_lines"
        _cs_contract = st.session_state.get(
            "CONTRACT_LABEL_A" if show_client else "CONTRACT_LABEL_B", ""
        )
        _cs_assureur = CONTRACTS_REGISTRY.get(_cs_contract, {}).get("assureur", "—")
        # Clean Shares par assureur (référentiel simplifié)
        _CLEAN_SHARE_ASSUREURS = {"Spirica"}  # Spirica permet les parts I sur certains fonds
        _cs_lines = [
            ln for ln in st.session_state.get(_cs_port_key, [])
            if str(ln.get("isin", "")).upper() not in ("EUROFUND", "STRUCTURED")
        ]
        if not _cs_lines:
            st.info("Aucun fonds UC dans le portefeuille sélectionné.")
        else:
            _frais_rows = []
            for _ln in _cs_lines:
                _name = (_ln.get("name") or _ln.get("isin") or "—")[:35]
                _ter = float(_ln.get("fee_uc_pct") or 0.0)
                _contrat_fee = float(_ln.get("fee_contract_pct") or
                                     st.session_state.get("FEE_A" if show_client else "FEE_B", 0.6))
                _total = _ter + _contrat_fee
                _frais_rows.append({
                    "Support": _name,
                    "TER fonds (%/an)": round(_ter, 2),
                    "Frais contrat (%/an)": round(_contrat_fee, 2),
                    "Coût total (%/an)": round(_total, 2),
                })
            _frais_df = pd.DataFrame(_frais_rows)
            st.dataframe(
                _frais_df.style.format({
                    "TER fonds (%/an)": "{:.2f}%",
                    "Frais contrat (%/an)": "{:.2f}%",
                    "Coût total (%/an)": "{:.2f}%",
                }).background_gradient(subset=["Coût total (%/an)"], cmap="YlOrRd"),
                hide_index=True, use_container_width=True,
            )
            st.caption(
                "Le TER est intégré dans la valeur liquidative publiée. "
                "Les frais contrat sont déduits en sus par l'assureur."
            )
            # Alerte Clean Shares si assureur éligible
            if _cs_assureur in _CLEAN_SHARE_ASSUREURS:
                st.info(
                    f"✨ **Optimisation possible — {_cs_assureur}** : Ce contrat donne "
                    "accès aux parts institutionnelles (parts I / Clean Shares) sur certains "
                    "fonds éligibles, potentiellement **−0.5% à −1%/an de frais de gestion**. "
                    "Vérifiez l'éligibilité de chaque fonds directement auprès de l'assureur."
                )

    # APP – Composition
    def _wrap_label_app(label: str, width: int = 28) -> str:
        if not label:
            return "—"
        return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)

    if show_valority and not dfB.empty:
        st.markdown("---")
        st.subheader("📈 Évolution du portefeuille Cabinet")
        st.caption(
            "Comparaison année par année entre le capital net investi "
            "et la valeur du portefeuille."
        )
        # ── Construire les données année par année ────────────────
        _blines_b = st.session_state.get("B_lines", [])
        _fee_b = float(st.session_state.get("FEE_B", 0.0) or 0.0)
        _M_B_evol = float(st.session_state.get("M_B", 0.0) or 0.0)
        _ONE_B = float(st.session_state.get("ONE_B", 0.0) or 0.0)
        _ONE_B_DATE = st.session_state.get("ONE_B_DATE", pd.Timestamp("2024-07-01").date())
        _start_b = startB_min if isinstance(startB_min, pd.Timestamp) else TODAY

        _years_b = sorted(set(dfB.index.year))
        _bar_rows: List[Dict[str, Any]] = []

        for _y in _years_b:
            _mask_y = dfB.index.year == _y
            if not _mask_y.any():
                continue
            _val_end = float(dfB.loc[_mask_y, "Valeur"].iloc[-1])
            _year_end = pd.Timestamp(f"{_y}-12-31")

            # 1) Versements initiaux nets
            _net_init = sum(
                float(ln.get("amount_gross", 0)) * (1.0 - _fee_b / 100.0)
                for ln in _blines_b
                if pd.Timestamp(ln.get("buy_date") or _start_b) <= _year_end
            )

            # 2) Versement ponctuel net
            _net_one = 0.0
            if _ONE_B > 0 and pd.Timestamp(_ONE_B_DATE) <= _year_end:
                _net_one = _ONE_B * (1.0 - _fee_b / 100.0)

            # 3) Versements mensuels nets cumulés
            _net_monthly_evol = 0.0
            if _M_B_evol > 0:
                _to_date = min(_year_end, TODAY)
                _n_months = max(
                    0,
                    (_to_date.year - _start_b.year) * 12
                    + (_to_date.month - _start_b.month),
                )
                _net_monthly_evol = _n_months * _M_B_evol * (1.0 - _fee_b / 100.0)

            _net_cumul = _net_init + _net_one + _net_monthly_evol

            _bar_rows.append({
                "Année": str(_y),
                "Capital net investi": round(_net_cumul, 2),
                "Valeur du portefeuille": round(_val_end, 2),
            })

        _bar_df_evol = pd.DataFrame(_bar_rows)

        if not _bar_df_evol.empty:
            _bar_long_evol = _bar_df_evol.melt(
                "Année", var_name="Type", value_name="Montant (€)"
            )
            _evol_chart = (
                alt.Chart(_bar_long_evol)
                .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                .encode(
                    x=alt.X("Année:N", title="Année", axis=alt.Axis(labelAngle=0)),
                    y=alt.Y(
                        "Montant (€):Q",
                        title="Montant (€)",
                        axis=alt.Axis(format=",.0f"),
                        scale=alt.Scale(zero=True),
                    ),
                    color=alt.Color(
                        "Type:N",
                        scale=alt.Scale(
                            domain=["Capital net investi", "Valeur du portefeuille"],
                            range=["#1B2A4A", "#C9A84C"],
                        ),
                        legend=alt.Legend(title=None, orient="bottom"),
                    ),
                    xOffset="Type:N",
                    tooltip=[
                        alt.Tooltip("Année:N"),
                        alt.Tooltip("Type:N"),
                        alt.Tooltip("Montant (€):Q", format=",.2f"),
                    ],
                )
                .properties(height=400)
            )
            st.altair_chart(_evol_chart, use_container_width=True)

            # Récapitulatif chiffré de la dernière année disponible
            _last_row = _bar_df_evol.iloc[-1]
            _pv_total = _last_row["Valeur du portefeuille"] - _last_row["Capital net investi"]
            _pv_pct = (
                _pv_total / _last_row["Capital net investi"] * 100.0
            ) if _last_row["Capital net investi"] > 0 else 0.0
            _cols_legend = st.columns(3)
            with _cols_legend[0]:
                st.metric("Capital net investi", to_eur(_last_row["Capital net investi"]))
            with _cols_legend[1]:
                st.metric("Valeur actuelle", to_eur(_last_row["Valeur du portefeuille"]))
            with _cols_legend[2]:
                st.metric("Plus-value", to_eur(_pv_total), delta=f"{_pv_pct:+.1f}%")
        else:
            st.info("Données insuffisantes pour afficher l'évolution du capital.")

    # ----------------------------------------------------------------
    # Section téléchargement des rapports
    # ----------------------------------------------------------------
    def _generate_pdf_safe(rd: Dict[str, Any]) -> bytes:
        try:
            return generate_pdf_report(rd)
        except Exception as e:
            st.warning(f"PDF indisponible : {e}")
            return b""

    st.markdown("---")
    st.subheader("📥 Télécharger")
    _report_data = st.session_state.get("REPORT_DATA")
    if _report_data is not None:
        _report_data_standard = {**_report_data}
        _html_standard = build_html_report(_report_data_standard)
        _col1, _col2 = st.columns(2)
        with _col1:
            st.download_button(
                "📄 Rapport standard (PDF)",
                data=st.session_state.get("PDF_BYTES_CACHE", b""),
                file_name="rapport_portefeuille.pdf",
                mime="application/pdf",
                help="Rapport de performance et composition, sans analyse Morningstar",
            )
        with _col2:
            st.download_button(
                "📄 Rapport standard (HTML)",
                data=_html_standard.encode("utf-8"),
                file_name="rapport_portefeuille.html",
                mime="text/html",
                help="Version HTML du rapport standard, consultable dans un navigateur",
            )

        # ── One-Pager Proposition d'Arbitrage ─────────────────────
        if _report_data is not None:
            st.markdown("---")
            if REPORTLAB_AVAILABLE:
                _report_onepager = {
                    **_report_data_standard,
                    "fee_a_pct": st.session_state.get("FEE_A", 0.6),
                    "fee_b_pct": st.session_state.get("FEE_B", 0.6),
                }
                st.download_button(
                    "📋 Proposition d'arbitrage (PDF one-pager)",
                    data=_build_onepager_pdf(_report_onepager),
                    file_name="proposition_arbitrage.pdf",
                    mime="application/pdf",
                    help="PDF synthétique prêt à présenter : état actuel vs cible, "
                         "gain estimé, levier fiscal.",
                )

        # ── Présentation PPTX ─────────────────────────────────────────────
        st.markdown("---")
        if PPTX_AVAILABLE:
            def _generate_pptx_safe(rd: Dict[str, Any]) -> bytes:
                try:
                    return generate_pptx_report(rd)
                except Exception as e:
                    st.session_state["_PPTX_LAST_ERROR"] = str(e)
                    return b""

            _nom_client_safe = (
                st.session_state.get("NOM_CLIENT", "") or "client"
            ).strip().replace(" ", "_")

            _col_p1, _col_p2 = st.columns(2)
            _report_pptx_standard = {**_report_data_standard}
            with _col_p1:
                st.download_button(
                    "📊 Présentation client (PPTX)",
                    data=_generate_pptx_safe(_report_pptx_standard),
                    file_name=f"presentation_{_nom_client_safe}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    help="Présentation PowerPoint 5 slides pour le client",
                )
            _pptx_err = st.session_state.pop("_PPTX_LAST_ERROR", None)
            if _pptx_err:
                st.warning(f"Erreur génération PPTX : {_pptx_err}")
        else:
            st.caption(f"PPTX indisponible : {PPTX_ERROR}")
    else:
        st.info("Les rapports seront disponibles après le calcul du portefeuille.")

    with st.expander("Aide rapide"):
        st.markdown(
            """
- Dans chaque portefeuille, vous pouvez **soit** ajouter des *fonds recommandés* (onglet dédié),
  **soit** utiliser la *saisie libre* avec ISIN / code.
- Pour le **fonds en euros**, utilisez le symbole **EUROFUND** (taux paramétrable dans la barre de gauche).
- Les frais d'entrée s'appliquent à chaque investissement.
- Le **rendement total** est la performance globale depuis l'origine (valeur actuelle / net investi).
- Le **rendement annualisé** utilise le XIRR (prise en compte des dates et montants).
- En mode **Personnalisé**, vous pouvez affecter précisément les versements mensuels et ponctuels à chaque ligne,
  avec un contrôle automatique de cohérence par rapport aux montants bruts saisis.
            """
        )

    # ------------------------------------------------------------
    # Analyse interne — Corrélation & volatilité (réservé conseiller)
    # ------------------------------------------------------------
    st.markdown("---")
    with st.expander("🔒 Analyse interne — Corrélation, volatilité et profil de risque", expanded=False):
        st.caption(
            "Section réservée au conseiller : analyse technique basée sur les valeurs liquidatives "
            "(corrélations, volatilités, drawdown)."
        )

        # FIXED: use per-portfolio euro rates instead of shared EURO_RATE_PREVIEW (Résidu Bug 4)
        euro_rate_A = st.session_state.get("EURO_RATE_A", 2.0)
        euro_rate_B = st.session_state.get("EURO_RATE_B", 2.5)
        linesA = st.session_state.get("A_lines", [])
        linesB = st.session_state.get("B_lines", [])

        # Portefeuille Client
        if show_client:
            st.markdown("### Portefeuille 1 — Client")
            corrA = correlation_matrix_from_lines(linesA, euro_rate_A)
            volA = volatility_table_from_lines(linesA, euro_rate_A)
            riskA = portfolio_risk_stats(linesA, euro_rate_A, fee_pct=st.session_state.get("FEE_A", 0.0))

            if corrA.empty and volA.empty:
                st.info("Pas assez d'historique ou de lignes pour analyser ce portefeuille.")
            else:
                if riskA is not None:
                    c1, c2 = st.columns(2)
                    with c1:
                        st.metric(
                            "Volatilité annuelle estimée",
                            f"{riskA['vol_ann_pct']:.2f} %",
                        )
                    with c2:
                        st.metric(
                            "Max drawdown (historique sur la période)",
                            f"{riskA['max_dd_pct']:.2f} %",
                        )

                if not volA.empty:
                    st.markdown("**Volatilité par ligne**")
                    st.dataframe(
                        volA.style.format(
                            {
                                "Écart-type quotidien %": "{:,.2f}%".format,
                                "Volatilité annuelle %": "{:,.2f}%".format,
                            }
                        ),
                        use_container_width=True,
                    )

                if not corrA.empty:
                    chartA = _corr_heatmap_chart(corrA, "Corrélation des lignes — Portefeuille Client")
                    if chartA is not None:
                        st.altair_chart(chartA, use_container_width=True)
                    # Alerte si corrélation excessive
                    _corrA_vals = corrA.where(
                        ~pd.DataFrame(
                            [[i == j for j in range(len(corrA.columns))]
                             for i in range(len(corrA.index))],
                            index=corrA.index, columns=corrA.columns
                        )
                    )
                    if not _corrA_vals.isna().all().all() and float(_corrA_vals.max().max()) > 0.90:
                        st.caption(
                            "⚠️ Information : Certains fonds sont très fortement corrélés (>0.90), "
                            "limitant l'impact réel de la diversification."
                        )

        if show_client and show_valority:
            st.markdown("---")

        if show_valority:
            st.markdown("### Portefeuille 2 — Cabinet")
            corrB = correlation_matrix_from_lines(linesB, euro_rate_B)
            volB = volatility_table_from_lines(linesB, euro_rate_B)
            riskB = portfolio_risk_stats(linesB, euro_rate_B, fee_pct=st.session_state.get("FEE_B", 0.0))

            if corrB.empty and volB.empty:
                st.info("Pas assez d'historique ou de lignes pour analyser ce portefeuille.")
            else:
                if riskB is not None:
                    c1, c2 = st.columns(2)
                    with c1:
                        st.metric(
                            "Volatilité annuelle estimée",
                            f"{riskB['vol_ann_pct']:.2f} %",
                        )
                    with c2:
                        st.metric(
                            "Max drawdown (historique sur la période)",
                            f"{riskB['max_dd_pct']:.2f} %",
                        )

                if not volB.empty:
                    st.markdown("**Volatilité par ligne**")
                    st.dataframe(
                        volB.style.format(
                            {
                                "Écart-type quotidien %": "{:,.2f}%".format,
                                "Volatilité annuelle %": "{:,.2f}%".format,
                            }
                        ),
                        use_container_width=True,
                    )

                if not corrB.empty:
                    chartB = _corr_heatmap_chart(corrB, "Corrélation des lignes — Portefeuille Cabinet")
                    if chartB is not None:
                        st.altair_chart(chartB, use_container_width=True)
                    # Alerte si corrélation excessive
                    _corrB_vals = corrB.where(
                        ~pd.DataFrame(
                            [[i == j for j in range(len(corrB.columns))]
                             for i in range(len(corrB.index))],
                            index=corrB.index, columns=corrB.columns
                        )
                    )
                    if not _corrB_vals.isna().all().all() and float(_corrB_vals.max().max()) > 0.90:
                        st.caption(
                            "⚠️ Information : Certains fonds sont très fortement corrélés (>0.90), "
                            "limitant l'impact réel de la diversification."
                        )

    # ----------------------------------------------------------------
    # Module "Drag" — Traînée de frais contrat
    # ----------------------------------------------------------------
    st.markdown("---")
    with st.expander("⚡ Analyse du Drag — Traînée de frais contrat", expanded=False):
        st.caption(
            "Visualise l'impact cumulé des frais de gestion du contrat sur la performance "
            "d'un fonds. Plus le 'Drag' est élevé, plus l'enveloppe érode la performance du support."
        )
        # Time to Market par assureur
        _TTM = {
            "Suravenir": {"délai": "J+2 (arbitrage)", "note": "Ordres transmis avant 12h exécutés J+1 VL"},
            "Spirica":   {"délai": "J+3 (arbitrage)", "note": "Délai réglementaire Spirica standard"},
        }
        _assureur_A = CONTRACTS_REGISTRY.get(
            st.session_state.get("CONTRACT_LABEL_A", ""), {}
        ).get("assureur", "—")
        _assureur_B = CONTRACTS_REGISTRY.get(
            st.session_state.get("CONTRACT_LABEL_B", ""), {}
        ).get("assureur", "—")
        _tc1, _tc2 = st.columns(2)
        with _tc1:
            _ttm_a = _TTM.get(_assureur_A, {"délai": "—", "note": "—"})
            st.metric("Time to Market — Contrat Client", _ttm_a["délai"])
            st.caption(f"({_assureur_A}) {_ttm_a['note']}")
        with _tc2:
            _ttm_b = _TTM.get(_assureur_B, {"délai": "—", "note": "—"})
            st.metric("Time to Market — Contrat Cabinet", _ttm_b["délai"])
            st.caption(f"({_assureur_B}) {_ttm_b['note']}")
        st.markdown("---")
        st.markdown("**Simulation du Drag sur un fonds UC**")
        _drag_c1, _drag_c2 = st.columns(2)
        with _drag_c1:
            _drag_perf = st.number_input(
                "Performance brute annuelle du fonds (%)",
                min_value=-20.0, max_value=40.0,
                value=7.0, step=0.5, key="drag_perf_brute",
            )
            _drag_fee = st.number_input(
                "Frais de gestion contrat (%/an)",
                min_value=0.0, max_value=3.0,
                value=0.6, step=0.1, key="drag_fee_contrat",
            )
        with _drag_c2:
            _drag_ter = st.number_input(
                "TER du fonds (%/an, déjà dans la VL)",
                min_value=0.0, max_value=3.0,
                value=1.5, step=0.1, key="drag_ter",
            )
            _drag_years = st.slider(
                "Horizon (années)", min_value=1, max_value=30,
                value=10, key="drag_years",
            )
        _drag_net = _drag_perf - _drag_fee
        _cap_init = 10_000.0
        _years_range = list(range(_drag_years + 1))
        _brut_series = [_cap_init * (1 + _drag_perf / 100) ** y for y in _years_range]
        _net_series  = [_cap_init * (1 + _drag_net  / 100) ** y for y in _years_range]
        _drag_df = pd.DataFrame({
            "Année": _years_range * 2,
            "Valeur (€)": _brut_series + _net_series,
            "Série": ["Brut (hors frais contrat)"] * len(_years_range)
                   + ["Net (après frais contrat)"]  * len(_years_range),
        })
        _drag_loss = _brut_series[-1] - _net_series[-1]
        _drag_area = (
            alt.Chart(_drag_df)
            .mark_area(opacity=0.18, interpolate="monotone")
            .encode(
                x=alt.X("Année:Q", title="Années"),
                y=alt.Y("Valeur (€):Q", axis=alt.Axis(format=",.0f")),
                color=alt.Color(
                    "Série:N",
                    scale=alt.Scale(
                        domain=["Brut (hors frais contrat)", "Net (après frais contrat)"],
                        range=["#1B2A4A", "#C9A84C"],
                    ),
                ),
                tooltip=[
                    alt.Tooltip("Année:Q"),
                    alt.Tooltip("Série:N"),
                    alt.Tooltip("Valeur (€):Q", format=",.0f"),
                ],
            )
        )
        _drag_line = (
            alt.Chart(_drag_df)
            .mark_line(strokeWidth=2.5, interpolate="monotone")
            .encode(
                x=alt.X("Année:Q"),
                y=alt.Y("Valeur (€):Q"),
                color=alt.Color(
                    "Série:N",
                    scale=alt.Scale(
                        domain=["Brut (hors frais contrat)", "Net (après frais contrat)"],
                        range=["#1B2A4A", "#C9A84C"],
                    ),
                ),
            )
        )
        st.altair_chart((_drag_area + _drag_line).properties(height=300), use_container_width=True)
        # ── Graphique comparatif barres : Brut vs Net en % et en € ──
        st.markdown("**Comparatif annualisé : Performance brute vs nette de frais**")
        _bar_data = pd.DataFrame({
            "Métrique": [
                "Rendement brut (%/an)",
                "Rendement net (%/an)",
                "Écart (Drag)",
            ],
            "Valeur": [
                _drag_perf,
                _drag_net,
                _drag_perf - _drag_net,
            ],
            "Couleur": ["Brut", "Net", "Drag"],
        })
        _bar_chart = (
            alt.Chart(_bar_data)
            .mark_bar(cornerRadiusTopLeft=4, cornerRadiusTopRight=4)
            .encode(
                x=alt.X(
                    "Métrique:N",
                    sort=None,
                    axis=alt.Axis(labelAngle=0, labelFontSize=10),
                    title="",
                ),
                y=alt.Y(
                    "Valeur:Q",
                    title="%/an",
                    axis=alt.Axis(format=".1f"),
                ),
                color=alt.Color(
                    "Couleur:N",
                    scale=alt.Scale(
                        domain=["Brut", "Net", "Drag"],
                        range=["#1B2A4A", "#C9A84C", "#CC2200"],
                    ),
                    legend=None,
                ),
                tooltip=[
                    alt.Tooltip("Métrique:N", title="Métrique"),
                    alt.Tooltip("Valeur:Q", title="%/an", format=".2f"),
                ],
            )
            .properties(height=220)
        )
        _text_bar = (
            alt.Chart(_bar_data)
            .mark_text(dy=-8, fontSize=11, fontWeight="bold")
            .encode(
                x=alt.X("Métrique:N", sort=None),
                y=alt.Y("Valeur:Q"),
                text=alt.Text("Valeur:Q", format=".2f"),
                color=alt.value("#333333"),
            )
        )
        st.altair_chart((_bar_chart + _text_bar).properties(height=220),
                        use_container_width=True)
        # Tableau récapitulatif € pour l'horizon choisi
        _drag_table_df = pd.DataFrame({
            "": ["Performance brute", "Performance nette", "Drag (coût frais)"],
            "Taux (%/an)": [
                f"{_drag_perf:.2f}%",
                f"{_drag_net:.2f}%",
                f"−{_drag_fee:.2f}%",
            ],
            f"Impact sur {_drag_years} ans (base {to_eur(_cap_init)})": [
                f"+{to_eur(_brut_series[-1] - _cap_init)}",
                f"+{to_eur(_net_series[-1] - _cap_init)}",
                f"−{to_eur(_drag_loss)}",
            ],
        })
        st.dataframe(_drag_table_df, hide_index=True, use_container_width=True)
        st.warning(
            f"💸 **Drag total sur {_drag_years} ans** : **{to_eur(_drag_loss)}** de valeur "
            f"perdue uniquement à cause des frais de gestion contrat ({_drag_fee:.1f}%/an). "
            f"Le TER du fonds ({_drag_ter:.1f}%/an) est déjà intégré dans la VL publiée."
        )

    # ----------------------------------------------------------------
    # Module Stress-Test
    # ----------------------------------------------------------------
    st.markdown("---")
    with st.expander("🧨 Stress-Test — Scénarios de crise", expanded=False):
        st.caption(
            "Estimation de l'impact d'un choc de marché ou de taux sur le portefeuille. "
            "Les calculs de bêta sont approximatifs (base : corrélation historique avec l'indice)."
        )
        # Paramètres portefeuille actif
        _st_val = valA if show_client else valB
        _st_net = netA if show_client else netB
        _st_lines = (
            st.session_state.get("A_lines", []) if show_client
            else st.session_state.get("B_lines", [])
        )
        _st_euro_rate = (
            st.session_state.get("EURO_RATE_A", 2.0) if show_client
            else st.session_state.get("EURO_RATE_B", 2.5)
        )
        # Calculer part euro vs UC
        _uc_val = sum(
            float(ln.get("montant_net", 0) or 0)
            for ln in _st_lines
            if str(ln.get("isin", "")).upper() not in ("EUROFUND", "STRUCTURED")
        )
        _euro_val = _st_val - _uc_val if _st_val > 0 else 0.0
        _uc_pct = _uc_val / _st_val if _st_val > 0 else 0.5
        _euro_pct = 1.0 - _uc_pct
        st.info(
            f"Portefeuille analysé : valeur **{to_eur(_st_val)}** — "
            f"UC estimées **{_uc_pct*100:.0f}%** / Fonds euros **{_euro_pct*100:.0f}%**"
        )
        st.markdown("#### Scénario 1 — Choc de marché : CAC 40 −20%")
        _auto_beta = (
            st.session_state.get("BETA_AUTO_A")
            if show_client
            else st.session_state.get("BETA_AUTO_B")
        )
        if _auto_beta is not None:
            _bench_lbl = st.session_state.get("BENCHMARK_LABEL", "MSCI World (CW8.PA)")
            st.info(
                f"Bêta calculé automatiquement : **{_auto_beta:.2f}** vs {_bench_lbl} "
                f"(voir section « Indicateurs de risque & performance » ci-dessus)"
            )
            _use_auto_beta = st.checkbox(
                "Utiliser le bêta calculé automatiquement",
                value=True,
                key="stress_use_auto_beta",
                help="Décochez pour saisir un bêta manuel.",
            )
            if _use_auto_beta:
                _beta = float(_auto_beta)
            else:
                _beta = st.slider(
                    "Bêta manuel du portefeuille UC",
                    min_value=0.0, max_value=2.0,
                    value=round(float(_auto_beta) * 20) / 20,
                    step=0.05,
                    key="stress_beta_manual",
                    help="Bêta > 1 : plus sensible que le marché. < 1 : plus défensif.",
                )
        else:
            _beta = st.slider(
                "Bêta estimé du portefeuille UC (1 = suit le marché)",
                min_value=0.0, max_value=2.0, value=0.80, step=0.05,
                key="stress_beta",
                help="Bêta > 1 : plus sensible que le marché. < 1 : plus défensif.",
            )
        _choc_marche = -0.20
        _perte_uc = _st_val * _uc_pct * _beta * _choc_marche
        _val_post_marche = _st_val + _perte_uc
        _sc1_c1, _sc1_c2, _sc1_c3 = st.columns(3)
        _sc1_c1.metric("Perte estimée UC", to_eur(_perte_uc), delta=f"{_beta * _choc_marche * _uc_pct * 100:.1f}% du total")
        _sc1_c2.metric("Valeur post-choc", to_eur(_val_post_marche))
        _sc1_c3.metric(
            "Perte vs investi",
            f"{(_val_post_marche / _st_net - 1) * 100:.1f}%" if _st_net > 0 else "—",
        )
        st.markdown("---")
        st.markdown("#### Scénario 2 — Choc de taux : +1% (impact fonds euros)")
        _duration_euro = st.slider(
            "Duration implicite du fonds euros (années)",
            min_value=0.0, max_value=10.0, value=3.0, step=0.5,
            key="stress_duration",
            help="Approximation : un fonds euros majoritairement obligataire a une duration de 3-5 ans.",
        )
        _choc_taux = 0.01  # +1%
        _perte_euro_pct = -_duration_euro * _choc_taux
        _perte_euro = _st_val * _euro_pct * _perte_euro_pct
        _val_post_taux = _st_val + _perte_euro
        _sc2_c1, _sc2_c2, _sc2_c3 = st.columns(3)
        _sc2_c1.metric("Perte estimée fonds €", to_eur(_perte_euro), delta=f"{_perte_euro_pct * _euro_pct * 100:.1f}% du total")
        _sc2_c2.metric("Valeur post-choc", to_eur(_val_post_taux))
        _sc2_c3.metric("Impact duration", f"−{_duration_euro * _choc_taux * 100:.1f}%")
        st.markdown("---")
        st.markdown("#### Scénario combiné — Krach global")
        _crash_result = simulate_market_crash(
            portfolio_value=_st_val,
            uc_pct=_uc_pct,
            euro_pct=_euro_pct,
            beta=_beta,
            crash_magnitude=0.20,
            duration_euro=_duration_euro,
            rate_shock=0.01,
        )
        _sc3_c1, _sc3_c2, _sc3_c3 = st.columns(3)
        _sc3_c1.metric(
            "Perte totale estimée",
            to_eur(_crash_result["perte_totale"]),
            delta=f"{_crash_result['pct_perte_total']:.1f}% du portefeuille",
        )
        _sc3_c2.metric(
            "Valeur post-krach",
            to_eur(_crash_result["valeur_post_choc"]),
        )
        _sc3_c3.metric(
            "Retour à l'équilibre",
            f"+{abs(_crash_result['pct_perte_total']) / (1 - abs(_crash_result['pct_perte_total']) / 100):.1f}% requis"
            if _crash_result["pct_perte_total"] < 0 else "—",
            help="Performance nécessaire pour récupérer la perte",
        )
        # Graphique waterfall simplifié (barres empilées)
        _crash_df = pd.DataFrame({
            "Composante": [
                "Portefeuille initial",
                "Perte UC (krach −20%)",
                "Perte fonds € (taux +1%)",
                "Valeur post-krach",
            ],
            "Valeur": [
                _st_val,
                _crash_result["perte_uc"],
                _crash_result["perte_euro"],
                _crash_result["valeur_post_choc"],
            ],
            "Type": ["Base", "Perte", "Perte", "Résultat"],
        })
        _crash_bar = (
            alt.Chart(_crash_df)
            .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
            .encode(
                x=alt.X("Composante:N", sort=None,
                        axis=alt.Axis(labelAngle=-15, labelFontSize=9), title=""),
                y=alt.Y("Valeur:Q", axis=alt.Axis(format=",.0f"), title="€"),
                color=alt.Color(
                    "Type:N",
                    scale=alt.Scale(
                        domain=["Base", "Perte", "Résultat"],
                        range=["#1B2A4A", "#CC2200", "#C9A84C"],
                    ),
                    legend=alt.Legend(title=""),
                ),
                tooltip=[
                    alt.Tooltip("Composante:N"),
                    alt.Tooltip("Valeur:Q", format=",.0f", title="€"),
                ],
            )
            .properties(height=260)
        )
        st.altair_chart(_crash_bar, use_container_width=True)
        st.caption(
            "⚠️ Simulation indicative uniquement — sensibilités linéaires. "
            "Ne constitue pas une prévision ni un conseil en investissement."
        )


def simulate_market_crash(
    portfolio_value: float,
    uc_pct: float,
    euro_pct: float,
    beta: float = 0.85,
    crash_magnitude: float = 0.20,
    duration_euro: float = 3.0,
    rate_shock: float = 0.01,
) -> Dict[str, float]:
    """
    Simule l'impact d'un krach combiné (marché + taux) sur un portefeuille.
    Retourne les pertes estimées par composante et la valeur post-choc.
    """
    perte_uc     = portfolio_value * uc_pct   * beta   * (-crash_magnitude)
    perte_euro   = portfolio_value * euro_pct * (-duration_euro * rate_shock)
    perte_totale = perte_uc + perte_euro
    return {
        "perte_uc":          perte_uc,
        "perte_euro":        perte_euro,
        "perte_totale":      perte_totale,
        "valeur_post_choc":  portfolio_value + perte_totale,
        "pct_perte_total":   (perte_totale / portfolio_value * 100) if portfolio_value > 0 else 0.0,
    }


def run_comparator():
    render_app(run_page_config=False)


# ============================================================
# MODULE FISCALITÉ ASSURANCE-VIE — fonctions de calcul pures
# ============================================================

def calc_quote_part_gains(
    valeur_contrat: float,
    versements_nets: float,
    montant_rachat: float,
) -> float:
    """Quote-part de gains dans le rachat."""
    if valeur_contrat <= 0 or versements_nets >= valeur_contrat:
        return 0.0
    return max(0.0, montant_rachat * (1.0 - versements_nets / valeur_contrat))


def calc_imposition_rachat(
    gains: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
    montant_rachat: float,
    option_ir: bool = False,
) -> Dict[str, Any]:
    PS_RATE = 0.172
    PFU_RATE = 0.128
    PFL_LOW_RATE = 0.075
    PFL_HIGH_RATE = 0.128
    ABATTEMENT_SEUL = 4_600.0
    ABATTEMENT_COUPLE = 9_200.0
    SEUIL_150K = 150_000.0

    abattement = 0.0
    base_ir = gains
    taux_ir = 0.0
    regime = ""
    tmi_applicable = ""

    if anciennete_annees < 8:
        taux_ir = PFU_RATE
        regime = "PFU 30% (contrat < 8 ans)"
        tmi_applicable = "Flat Tax 12,8% IR + 17,2% PS"
        if option_ir:
            tmi_applicable = "Option barème IR — à déclarer case 2CH"
    else:
        abattement = ABATTEMENT_COUPLE if "Couple" in situation_familiale else ABATTEMENT_SEUL
        base_ir = max(0.0, gains - abattement)
        if versements_nets_total <= SEUIL_150K:
            taux_ir = PFL_LOW_RATE
            regime = "PFL 7,5% + PS (contrat ≥8 ans, versements ≤150k€)"
            tmi_applicable = "Prélèvement forfaitaire libératoire 7,5%"
        else:
            taux_ir = PFL_HIGH_RATE
            regime = "PFU 12,8% + PS (contrat ≥8 ans, versements >150k€)"
            tmi_applicable = "Flat Tax 12,8% (versements post-27/09/2017 >150k€)"

    montant_ir = base_ir * taux_ir
    montant_ps = gains * PS_RATE

    return {
        "gains": gains,
        "regime": regime,
        "base_ir": base_ir,
        "abattement_applique": abattement,
        "taux_ir": taux_ir,
        "montant_ir": montant_ir,
        "montant_ps": montant_ps,
        "total_impots": montant_ir + montant_ps,
        "net_percu": montant_rachat - (montant_ir + montant_ps),
        "tmi_applicable": tmi_applicable,
    }


def calc_rachat_depuis_net(
    montant_net_souhaite: float,
    valeur_contrat: float,
    versements_nets: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
) -> Tuple[float, Dict[str, Any]]:
    """Bisection : trouve le brut tel que brut - impôts == net_souhaité."""
    if montant_net_souhaite <= 0:
        return 0.0, calc_imposition_rachat(0.0, anciennete_annees, situation_familiale, versements_nets_total, 0.0)

    lo, hi = montant_net_souhaite, min(valeur_contrat, montant_net_souhaite * 2.5)
    # s'assurer que hi couvre bien le cas
    for _ in range(50):
        g = calc_quote_part_gains(valeur_contrat, versements_nets, hi)
        d = calc_imposition_rachat(g, anciennete_annees, situation_familiale, versements_nets_total, hi)
        if hi - d["total_impots"] >= montant_net_souhaite:
            break
        hi = min(valeur_contrat, hi * 1.5)

    for _ in range(1000):
        mid = (lo + hi) / 2.0
        g = calc_quote_part_gains(valeur_contrat, versements_nets, mid)
        d = calc_imposition_rachat(g, anciennete_annees, situation_familiale, versements_nets_total, mid)
        net = mid - d["total_impots"]
        if abs(net - montant_net_souhaite) < 0.01:
            return mid, d
        if net < montant_net_souhaite:
            lo = mid
        else:
            hi = mid
    return mid, d


def calc_optimisation_abattement(
    valeur_contrat: float,
    versements_nets: float,
    anciennete_annees: float,
    situation_familiale: str,
    versements_nets_total: float,
    abattement_deja_utilise: float = 0.0,
) -> Dict[str, Any]:
    PS_RATE = 0.172
    ABATTEMENT_SEUL = 4_600.0
    ABATTEMENT_COUPLE = 9_200.0

    abattement = ABATTEMENT_COUPLE if "Couple" in situation_familiale else ABATTEMENT_SEUL
    abattement_restant = max(0.0, abattement - abattement_deja_utilise)

    if valeur_contrat <= 0 or versements_nets >= valeur_contrat:
        return {
            "rachat_optimal": 0.0, "gains_realises": 0.0,
            "ir_du": 0.0, "ps_du": 0.0, "economies_ir": 0.0,
            "abattement_utilise": 0.0, "abattement_restant_apres": abattement_restant,
        }

    taux_pv = 1.0 - versements_nets / valeur_contrat
    if taux_pv <= 0:
        return {
            "rachat_optimal": 0.0, "gains_realises": 0.0,
            "ir_du": 0.0, "ps_du": 0.0, "economies_ir": 0.0,
            "abattement_utilise": 0.0, "abattement_restant_apres": abattement_restant,
        }

    rachat_optimal = abattement_restant / taux_pv
    rachat_optimal = min(rachat_optimal, valeur_contrat)
    gains_realises = calc_quote_part_gains(valeur_contrat, versements_nets, rachat_optimal)

    # IR si on avait dépassé l'abattement (scénario sans optimisation)
    SEUIL_150K = 150_000.0
    taux_ir = 0.075 if versements_nets_total <= SEUIL_150K else 0.128
    economies_ir = gains_realises * taux_ir

    return {
        "rachat_optimal": rachat_optimal,
        "gains_realises": gains_realises,
        "ir_du": 0.0,
        "ps_du": gains_realises * PS_RATE,
        "economies_ir": economies_ir,
        "abattement_utilise": gains_realises,
        "abattement_restant_apres": 0.0,
    }


def calc_transmission_990I(
    capital_deces: float,
    nb_beneficiaires: int,
    parts_pct: List[float],
    types_beneficiaires: List[str],
) -> List[Dict[str, Any]]:
    ABATTEMENT_990I = 152_500.0
    TAUX_990I_T1 = 0.20
    TAUX_990I_T2 = 0.3125
    SEUIL_990I_T2 = 700_000.0

    results = []
    for i in range(nb_beneficiaires):
        t = types_beneficiaires[i] if i < len(types_beneficiaires) else "Enfant"
        p = parts_pct[i] if i < len(parts_pct) else 0.0
        part_eur = capital_deces * p / 100.0
        note = ""

        if "Conjoint" in t or "PACS" in t:
            taxe = 0.0
            taxable = 0.0
            abat = part_eur
        else:
            if "Démembré" in t:
                note = (
                    "Le démembrement de la clause bénéficiaire nécessite une évaluation "
                    "actuarielle des droits d'usufruit et nue-propriété selon le barème "
                    "fiscal de l'article 669 CGI (fonction de l'âge de l'usufruitier). "
                    "Cette simulation affiche la taxation globale sur la part démembrée "
                    "comme si elle revenait à un bénéficiaire standard, à affiner avec "
                    "un notaire."
                )
            abat = min(part_eur, ABATTEMENT_990I)
            taxable = max(0.0, part_eur - ABATTEMENT_990I)
            if taxable <= SEUIL_990I_T2:
                taxe = taxable * TAUX_990I_T1
            else:
                taxe = SEUIL_990I_T2 * TAUX_990I_T1 + (taxable - SEUIL_990I_T2) * TAUX_990I_T2

        results.append({
            "index": i + 1,
            "type": t,
            "part_pct": p,
            "part_eur": part_eur,
            "abattement": abat,
            "taxable": taxable if "Conjoint" not in t and "PACS" not in t else 0.0,
            "taxe": taxe,
            "net_recu": part_eur - taxe,
            "note": note,
        })
    return results


def calc_transmission_757B(
    primes_versees_apres_70: float,
    nb_beneficiaires: int,
    parts_pct: List[float],
    types_beneficiaires: List[str],
    capital_deces: float,
    liens_succession: List[str],
) -> List[Dict[str, Any]]:
    ABATTEMENT_757B = 30_500.0
    TAUX_LIEN = {
        "Enfant": 0.20,
        "Frère/Sœur": 0.40,
        "Neveu/Nièce": 0.55,
        "Tiers": 0.60,
        "Conjoint/PACS": 0.0,
    }

    # Parts des non-conjoints pour répartir l'abattement global
    non_conjoint_total_pct = sum(
        parts_pct[i] for i in range(nb_beneficiaires)
        if i < len(types_beneficiaires)
        and "Conjoint" not in types_beneficiaires[i]
        and "PACS" not in types_beneficiaires[i]
    )

    results = []
    for i in range(nb_beneficiaires):
        t = types_beneficiaires[i] if i < len(types_beneficiaires) else "Enfant"
        p = parts_pct[i] if i < len(parts_pct) else 0.0
        lien = liens_succession[i] if i < len(liens_succession) else "Enfant"
        part_eur = capital_deces * p / 100.0

        if "Conjoint" in t or "PACS" in t or lien == "Conjoint/PACS":
            results.append({
                "index": i + 1, "type": t, "part_pct": p, "part_eur": part_eur,
                "primes_part": 0.0, "abattement_part": 0.0,
                "taxable": 0.0, "taxe": 0.0, "net_recu": part_eur,
                "taux_applique": 0.0,
            })
            continue

        # Répartition des primes post-70 ans selon la quote-part
        primes_part = primes_versees_apres_70 * p / max(non_conjoint_total_pct, 1.0)
        # Abattement global 30 500 € réparti proportionnellement
        abat_part = ABATTEMENT_757B * p / max(non_conjoint_total_pct, 1.0)
        taxable = max(0.0, primes_part - abat_part)
        taux = TAUX_LIEN.get(lien, 0.60)
        taxe = taxable * taux

        results.append({
            "index": i + 1, "type": t, "part_pct": p, "part_eur": part_eur,
            "primes_part": primes_part, "abattement_part": abat_part,
            "taxable": taxable, "taxe": taxe, "net_recu": part_eur - taxe,
            "taux_applique": taux,
        })
    return results


# ============================================================
# MODULE FISCALITÉ ASSURANCE-VIE — onglets UI
# ============================================================

def _fmt_eur(v: float) -> str:
    """Formatage monétaire uniforme pour le module fiscal."""
    return f"{v:,.0f} €".replace(",", "\u202f")


def _tab_rachat() -> None:
    # ── Données du contrat — lues depuis le Tableau de bord ──────────
    _date_ouv_r = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
    anciennete_annees = (date.today() - _date_ouv_r).days / 365.25
    valeur_contrat = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
    versements_nets = float(st.session_state.get("tax_versements_nets", 80_000.0))
    versements_nets_total = float(st.session_state.get("tax_versements_nets_total", 80_000.0))
    situation_familiale = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

    with st.container(border=True):
        _ri1, _ri2, _ri3 = st.columns(3)
        with _ri1:
            st.metric("Valeur du contrat", to_eur(valeur_contrat))
        with _ri2:
            st.metric("Versements nets", to_eur(versements_nets))
        with _ri3:
            _anc_disp = anciennete_annees
            st.metric("Ancienneté", f"{_anc_disp:.1f} ans")
        st.caption(f"📅 Ouverture : **{_date_ouv_r.strftime('%d/%m/%Y')}** | Situation : **{situation_familiale}** — _modifiez dans le 📋 Tableau de bord_")

    ps_deja_preleves = st.number_input(
        "PS déjà prélevés sur fonds euros (€)",
        min_value=0.0, max_value=100_000.0,
        value=0.0, step=100.0,
        key="tax_ps_deja_preleves",
        help="Sur les fonds en euros, les PS sont prélevés chaque année par l'assureur.",
    )
    st.markdown("#### Montant du rachat")

    # Synchronisation brut ↔ net via session_state
    last_edited = st.session_state.get("tax_last_edited", "brut")

    col_brut, col_net = st.columns(2)
    with col_brut:
        brut_default = float(st.session_state.get("tax_montant_brut", 10_000.0))
        montant_brut = st.number_input(
            "Montant brut du rachat (€)",
            min_value=0.0,
            max_value=float(max(valeur_contrat, 0.01)),
            value=brut_default,
            step=500.0,
            key="_tax_brut_input",
        )
        if montant_brut != brut_default:
            st.session_state["tax_last_edited"] = "brut"
            st.session_state["tax_montant_brut"] = montant_brut

    with col_net:
        net_default = float(st.session_state.get("tax_montant_net", 9_000.0))
        montant_net_input = st.number_input(
            "Montant net souhaité (€)",
            min_value=0.0,
            max_value=float(max(valeur_contrat, 0.01)),
            value=net_default,
            step=500.0,
            key="_tax_net_input",
        )
        if montant_net_input != net_default:
            st.session_state["tax_last_edited"] = "net"
            st.session_state["tax_montant_net"] = montant_net_input

    # Résoudre le champ non-prioritaire
    if st.session_state.get("tax_last_edited", "brut") == "net":
        montant_brut, _ = calc_rachat_depuis_net(
            montant_net_input, valeur_contrat, versements_nets,
            anciennete_annees, situation_familiale, versements_nets_total,
        )
    else:
        montant_brut = st.session_state.get("tax_montant_brut", montant_brut)

    # Calcul et affichage
    if valeur_contrat > 0 and versements_nets > 0 and montant_brut > 0:
        gains = calc_quote_part_gains(valeur_contrat, versements_nets, montant_brut)
        result = calc_imposition_rachat(
            gains, anciennete_annees, situation_familiale,
            versements_nets_total, montant_brut,
        )
        st.session_state["tax_rachat_result"] = result

        # Bandeau régime
        if anciennete_annees < 8:
            st.error(f"⏱️ Contrat de {anciennete_annees:.1f} an(s) — Régime PFU 30%")
        else:
            st.success(f"✅ Contrat de {anciennete_annees:.1f} ans — Régime favorable ≥8 ans")

        # Tableau récapitulatif
        capital_pur = montant_brut - gains
        abat = result["abattement_applique"]
        ps_net = max(0.0, result["montant_ps"] - ps_deja_preleves)
        total_impots_net = result["montant_ir"] + ps_net
        net_percu_net = montant_brut - total_impots_net

        rows_display = [
            ("Montant brut du rachat", _fmt_eur(montant_brut)),
            ("  dont quote-part gains", _fmt_eur(gains)),
            ("  dont quote-part capital", _fmt_eur(capital_pur)),
            ("─" * 35, ""),
            (f"Abattement IR appliqué", _fmt_eur(abat)),
            ("Base imposable IR", _fmt_eur(result["base_ir"])),
            (f"IR dû ({result['taux_ir']*100:.1f}%)", _fmt_eur(result["montant_ir"])),
            ("Prélèvements sociaux (17,2%)", _fmt_eur(result["montant_ps"])),
        ]
        if ps_deja_preleves > 0:
            rows_display.append((f"PS déjà prélevés (fonds €)", f"− {_fmt_eur(ps_deja_preleves)}"))
        rows_display += [
            ("─" * 35, ""),
            ("TOTAL prélèvements", _fmt_eur(total_impots_net)),
            ("MONTANT NET PERÇU", _fmt_eur(net_percu_net)),
        ]
        df_recap = pd.DataFrame(rows_display, columns=["", "Montant"])
        st.dataframe(df_recap, hide_index=True, use_container_width=True)

        st.info(
            f"**Régime appliqué :** {result['regime']}\n\n"
            f"**Taux :** {result['tmi_applicable']}"
        )

        # Conseil stratégique pour contrats ≥8 ans
        if anciennete_annees >= 8 and gains > 0:
            abat_total = result["abattement_applique"]
            rachat_capital_pur = versements_nets / (valeur_contrat / montant_brut) if valeur_contrat > 0 else 0
            brut_sans_ir, _ = calc_rachat_depuis_net(
                capital_pur, valeur_contrat, versements_nets,
                anciennete_annees, situation_familiale, versements_nets_total,
            ) if abat_total > 0 else (capital_pur, {})
            g_limite = calc_quote_part_gains(valeur_contrat, versements_nets, brut_sans_ir)
            ps_limite = g_limite * 0.172
            st.caption(
                f"💡 **Stratégie :** un rachat limité à {_fmt_eur(brut_sans_ir)} "
                f"(capital pur + abattement) ne génèrerait aucun IR, "
                f"uniquement {_fmt_eur(ps_limite)} de PS."
            )

        # ── Graphiques : Donut décomposition + Waterfall ──────────
        _capital_restitue = montant_brut - gains
        _gains_exoneres = min(gains, result["abattement_applique"])
        _gains_taxes_net = max(0.0, gains - result["abattement_applique"])
        _ir_chart = result["montant_ir"]
        _ps_net_chart = max(0.0, result["montant_ps"] - ps_deja_preleves)
        _net_chart = montant_brut - _ir_chart - _ps_net_chart

        _donut_rows = [
            {"Composante": "Capital restitué (non taxé)", "Montant": _capital_restitue},
            {"Composante": "Gains exonérés (abattement)", "Montant": _gains_exoneres},
            {"Composante": "Impôt sur le revenu", "Montant": _ir_chart},
            {"Composante": "Prélèvements sociaux", "Montant": _ps_net_chart},
        ]
        _donut_df = pd.DataFrame([r for r in _donut_rows if r["Montant"] > 0.5])

        if not _donut_df.empty:
            _dc1, _dc2 = st.columns(2)
            with _dc1:
                _donut_chart = (
                    alt.Chart(_donut_df)
                    .mark_arc(innerRadius=55, outerRadius=105)
                    .encode(
                        theta=alt.Theta("Montant:Q"),
                        color=alt.Color(
                            "Composante:N",
                            scale=alt.Scale(
                                domain=["Capital restitué (non taxé)", "Gains exonérés (abattement)",
                                        "Impôt sur le revenu", "Prélèvements sociaux"],
                                range=["#2E7D32", "#81C784", "#E53935", "#EF9A9A"],
                            ),
                            legend=alt.Legend(title="", orient="bottom"),
                        ),
                        tooltip=[
                            alt.Tooltip("Composante:N"),
                            alt.Tooltip("Montant:Q", format=",.0f", title="€"),
                        ],
                    )
                    .properties(height=280, title="Décomposition du rachat")
                )
                st.altair_chart(_donut_chart, use_container_width=True)

            with _dc2:
                _wf_data = pd.DataFrame([
                    {"Étape": "Rachat brut", "Montant": montant_brut, "Type": "Base"},
                    {"Étape": "− IR", "Montant": -_ir_chart, "Type": "Impôt"},
                    {"Étape": "− PS", "Montant": -_ps_net_chart, "Type": "Impôt"},
                    {"Étape": "Net perçu", "Montant": _net_chart, "Type": "Net"},
                ])
                _wf_chart = (
                    alt.Chart(_wf_data)
                    .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                    .encode(
                        x=alt.X("Étape:N", sort=None, axis=alt.Axis(labelAngle=0), title=""),
                        y=alt.Y("Montant:Q", axis=alt.Axis(format=",.0f"), title="€"),
                        color=alt.Color(
                            "Type:N",
                            scale=alt.Scale(
                                domain=["Base", "Impôt", "Net"],
                                range=["#1B2A4A", "#E53935", "#2E7D32"],
                            ),
                            legend=None,
                        ),
                        tooltip=[
                            alt.Tooltip("Étape:N"),
                            alt.Tooltip("Montant:Q", format=",.0f", title="€"),
                        ],
                    )
                    .properties(height=280, title="Brut → Net")
                )
                st.altair_chart(_wf_chart, use_container_width=True)


def _tab_optimisation_abattement() -> None:
    st.markdown("#### Optimisation de l'abattement annuel")
    st.info(
        "L'abattement de 4 600 € (célibataire) ou 9 200 € (couple) est global "
        "tous contrats d'assurance-vie. Cette simulation calcule le rachat optimal "
        "pour 'purger' les plus-values progressivement sans payer d'IR."
    )

    # ── Données lues depuis le Tableau de bord ──────────────────
    _date_ouv_o = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
    anc2 = (date.today() - _date_ouv_o).days / 365.25
    valeur2 = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
    versements2 = float(st.session_state.get("tax_versements_nets", 80_000.0))
    versements_total2 = float(st.session_state.get("tax_versements_nets_total", 80_000.0))
    sit2 = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

    with st.container(border=True):
        _oi1, _oi2, _oi3 = st.columns(3)
        _oi1.metric("Valeur du contrat", to_eur(valeur2))
        _oi2.metric("Versements nets", to_eur(versements2))
        _oi3.metric("Ancienneté", f"{anc2:.1f} ans")
        st.caption(f"Situation : **{sit2}** — _modifiez dans le 📋 Tableau de bord_")

    abat_deja = st.number_input(
        "Abattement déjà utilisé cette année sur d'autres contrats (€)",
        min_value=0.0, max_value=9_200.0,
        value=float(st.session_state.get("tax_abattement_deja_utilise", 0.0)),
        step=100.0, key="tax_abattement_deja_utilise",
        help="L'abattement de 4 600€/9 200€ est global tous contrats AV.",
    )

    if anc2 < 8:
        st.warning(
            f"Ce contrat n'a pas encore atteint 8 ans ({anc2:.1f} ans). "
            "L'optimisation par abattement ne s'applique pas."
        )
        return

    if valeur2 <= 0 or versements2 >= valeur2:
        st.info("Pas de plus-value latente — aucun calcul nécessaire.")
        return

    res = calc_optimisation_abattement(
        valeur2, versements2, anc2, sit2, versements_total2, abat_deja,
    )

    abat_total = 9_200.0 if "Couple" in sit2 else 4_600.0
    abat_restant = max(0.0, abat_total - abat_deja)

    mc1, mc2, mc3 = st.columns(3)
    mc1.metric("Abattement disponible restant", _fmt_eur(abat_restant))
    mc2.metric("Rachat recommandé", _fmt_eur(res["rachat_optimal"]))
    mc3.metric("IR dû", _fmt_eur(res["ir_du"]))

    mc4, mc5 = st.columns(2)
    mc4.metric("PS dus (inévitables)", _fmt_eur(res["ps_du"]))
    mc5.metric("Économie IR vs rachat non optimisé", _fmt_eur(res["economies_ir"]))

    # Graphique
    chart_data = pd.DataFrame({
        "Catégorie": ["Capital récupéré", "PS dus", "IR économisé"],
        "Montant (€)": [
            res["rachat_optimal"] - res["ps_du"],
            res["ps_du"],
            res["economies_ir"],
        ],
    })
    if MATPLOTLIB_AVAILABLE:
        fig, ax = plt.subplots(figsize=(5, 2.8))
        colors_bar = ["#2ecc71", "#e67e22", "#3498db"]
        bars = ax.bar(chart_data["Catégorie"], chart_data["Montant (€)"], color=colors_bar)
        ax.bar_label(bars, fmt="%.0f €", padding=3, fontsize=9)
        ax.set_ylabel("€")
        ax.set_title("Optimisation abattement — répartition")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        st.pyplot(fig)
        plt.close(fig)
    else:
        st.bar_chart(chart_data.set_index("Catégorie"))

    # Projection pluriannuelle
    nb_annees = st.slider(
        "Nombre d'années de la stratégie",
        min_value=1, max_value=30,
        value=int(st.session_state.get("tax_nb_annees", 10)),
        key="tax_nb_annees",
    )
    economie_cumulee = res["economies_ir"] * nb_annees
    st.info(
        f"📅 Cette optimisation se renouvelle chaque année civile. "
        f"Sur **{nb_annees} années**, l'économie d'IR cumulée estimée serait de "
        f"**{_fmt_eur(economie_cumulee)}** (hypothèse : situation stable)."
    )


def _tab_transmission() -> None:
    st.markdown("#### Fiscalité au décès — Clause bénéficiaire")

    regime_versements = st.radio(
        "Les versements ont été effectués :",
        [
            "Principalement avant 70 ans (Art. 990I)",
            "Principalement après 70 ans (Art. 757B)",
            "Mix avant et après 70 ans",
        ],
        key="tax_regime_versements",
    )

    c1, c2 = st.columns(2)
    with c1:
        capital_deces = st.number_input(
            "Valeur du contrat au décès (€)",
            min_value=0.0, max_value=10_000_000.0,
            value=float(st.session_state.get("tax_capital_deces", 200_000.0)),
            step=5_000.0, key="tax_capital_deces",
        )
    with c2:
        primes_apres_70 = 0.0
        if "757B" in regime_versements or "Mix" in regime_versements:
            primes_apres_70 = st.number_input(
                "Primes versées après 70 ans (€)",
                min_value=0.0, max_value=capital_deces,
                value=0.0, step=5_000.0,
                help="Seules les primes (pas les gains) au-delà de 30 500€ sont soumises aux droits de succession.",
            )

    st.markdown("---")

    # ── Lire les bénéficiaires depuis le Tableau de bord ──────────
    _benef_data_t = st.session_state.get("tax_dash_beneficiaires", [])
    if not _benef_data_t:
        st.warning("⚠️ Renseignez les bénéficiaires dans l'onglet 📋 Tableau de bord.")
        return

    nb_benef = len(_benef_data_t)
    noms = [b.get("nom", "") for b in _benef_data_t]
    types_benef = [b.get("type", "Enfant") for b in _benef_data_t]
    parts_pct = [float(b.get("part", 100.0 / nb_benef)) for b in _benef_data_t]
    liens_succ = [b.get("type", "Enfant") for b in _benef_data_t]
    # Mapper les types AV vers liens succession
    _lien_map = {
        "Conjoint/PACS": "Conjoint/PACS",
        "Enfant": "Enfant",
        "Frère/Sœur": "Frère/Sœur",
        "Neveu/Nièce": "Neveu/Nièce",
        "Tiers": "Tiers",
    }
    liens_succ = [_lien_map.get(t, "Enfant") for t in types_benef]

    total_parts = sum(parts_pct)
    if abs(total_parts - 100.0) > 0.5:
        st.error(f"⚠️ La somme des quotes-parts est {total_parts:.1f}% — elle doit être égale à 100%.")
        return

    st.markdown("---")

    if capital_deces <= 0:
        st.info("Saisissez la valeur du contrat au décès.")
        return

    # Calcul selon régime
    show_990i = "990I" in regime_versements or "Mix" in regime_versements
    show_757b = "757B" in regime_versements or "Mix" in regime_versements

    if show_990i:
        st.markdown("##### Art. 990I — Versements avant 70 ans")
        types_990i = [
            "Conjoint/PACS" if ("Conjoint" in t or "PACS" in t) else t
            for t in types_benef
        ]
        results_990i = calc_transmission_990I(capital_deces, nb_benef, parts_pct, types_990i)
        cap_990i = capital_deces
        if "Mix" in regime_versements:
            cap_990i_pct = 1.0 - (primes_apres_70 / capital_deces) if capital_deces > 0 else 1.0
            results_990i = calc_transmission_990I(
                capital_deces * cap_990i_pct, nb_benef, parts_pct, types_990i
            )
        _render_table_transmission(results_990i, noms, "990I")
        # Note démembrement
        for r in results_990i:
            if r.get("note"):
                st.info(r["note"])

    if show_757b:
        st.markdown("##### Art. 757B — Versements après 70 ans")
        st.info(
            "Les **gains** générés par les primes versées après 70 ans sont exonérés. "
            "Seules les **primes** (hors gains) au-delà de 30 500 € (abattement global) "
            "sont soumises aux droits de succession selon le lien de parenté."
        )
        st.caption(
            "Les droits de succession réels dépendent du patrimoine global, des donations "
            "antérieures et des abattements de droit commun (ex. 100 000€ par enfant). "
            "Cette simulation utilise des taux simplifiés."
        )
        results_757b = calc_transmission_757B(
            primes_apres_70, nb_benef, parts_pct, types_benef, capital_deces, liens_succ,
        )
        _render_table_transmission(results_757b, noms, "757B")


def _render_table_transmission(results: List[Dict[str, Any]], noms: List[str], regime: str) -> None:
    """Affiche le tableau de résultats par bénéficiaire."""
    rows = []
    for r in results:
        nom = noms[r["index"] - 1] if r["index"] - 1 < len(noms) and noms[r["index"] - 1] else f"Bénéficiaire {r['index']}"
        if regime == "990I":
            rows.append({
                "Bénéficiaire": nom,
                "Lien": r["type"],
                "Part (€)": _fmt_eur(r["part_eur"]),
                "Abattement": _fmt_eur(r.get("abattement", 0.0)),
                "Taxable": _fmt_eur(r["taxable"]),
                "Taxe due": _fmt_eur(r["taxe"]),
                "Net reçu": _fmt_eur(r["net_recu"]),
            })
        else:
            rows.append({
                "Bénéficiaire": nom,
                "Lien": r["type"],
                "Part (€)": _fmt_eur(r["part_eur"]),
                "Primes part": _fmt_eur(r.get("primes_part", 0.0)),
                "Abattement": _fmt_eur(r.get("abattement_part", 0.0)),
                "Taxable": _fmt_eur(r["taxable"]),
                "Taxe due": _fmt_eur(r["taxe"]),
                "Net reçu": _fmt_eur(r["net_recu"]),
            })
    total_taxe = sum(r["taxe"] for r in results)
    total_net = sum(r["net_recu"] for r in results)
    total_part = sum(r["part_eur"] for r in results)

    if rows:
        st.dataframe(pd.DataFrame(rows), hide_index=True, use_container_width=True)

    tc1, tc2, tc3 = st.columns(3)
    tc1.metric("Capital total transmis", _fmt_eur(total_part))
    tc2.metric("Total taxes", _fmt_eur(total_taxe))
    taux_eff = (total_taxe / total_part * 100) if total_part > 0 else 0.0
    tc3.metric("Taux effectif global", f"{taux_eff:.1f}%")


def _tab_exoneration() -> None:
    st.info(
        "Dans certaines situations, les gains issus d'un rachat sont totalement exonérés "
        "d'IR **ET** de prélèvements sociaux. Ces exonérations s'appliquent sur justificatif."
    )

    st.markdown("#### Cochez votre situation")
    cas_licenciement = st.checkbox(
        "Licenciement (inscription à Pôle Emploi / France Travail)"
    )
    cas_retraite_anticipee = st.checkbox(
        "Mise à la retraite anticipée par l'employeur"
    )
    cas_invalidite = st.checkbox(
        "Invalidité 2e ou 3e catégorie (assuré, conjoint ou enfant à charge)"
    )
    cas_liquidation = st.checkbox(
        "Liquidation judiciaire (non-salariés : cessation d'activité suite à jugement de liquidation)"
    )
    cas_cessation_ns = st.checkbox(
        "Cessation d'activité non salariée suite à jugement de liquidation"
    )

    cas_coches = [
        (cas_licenciement, "Licenciement",
         "Justificatif : lettre de licenciement + attestation Pôle Emploi / France Travail"),
        (cas_retraite_anticipee, "Retraite anticipée",
         "Justificatif : notification de mise à la retraite par l'employeur"),
        (cas_invalidite, "Invalidité",
         "Justificatif : notification CPAM de classement en invalidité 2e ou 3e catégorie"),
        (cas_liquidation, "Liquidation judiciaire",
         "Justificatif : jugement du tribunal de commerce"),
        (cas_cessation_ns, "Cessation NS",
         "Justificatif : jugement du tribunal de commerce"),
    ]

    au_moins_un = any(c[0] for c in cas_coches)

    if au_moins_un:
        st.success(
            "✅ **Exonération totale applicable** : ni IR ni prélèvements sociaux "
            "ne sont dus sur les gains du rachat."
        )
        # Tableau comparatif
        rachat_result: Optional[Dict[str, Any]] = st.session_state.get("tax_rachat_result")
        if rachat_result:
            ir_normal = rachat_result.get("montant_ir", 0.0)
            ps_normal = rachat_result.get("montant_ps", 0.0)
            total_normal = rachat_result.get("total_impots", 0.0)
            df_comp = pd.DataFrame([
                {"": "IR dû",
                 "Sans exonération": _fmt_eur(ir_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(ir_normal)},
                {"": "PS dûs",
                 "Sans exonération": _fmt_eur(ps_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(ps_normal)},
                {"": "TOTAL",
                 "Sans exonération": _fmt_eur(total_normal),
                 "Avec exonération": _fmt_eur(0.0),
                 "Économie": _fmt_eur(total_normal)},
            ])
            st.dataframe(df_comp, hide_index=True, use_container_width=True)
        else:
            st.caption("Renseignez d'abord l'onglet **Rachat** pour voir l'économie chiffrée.")

        st.markdown("#### Justificatifs requis")
        for coched, label, justif in cas_coches:
            if coched:
                st.caption(f"**{label}** — {justif}")

    st.markdown("---")
    st.warning(
        "⚠️ **Procédure :** l'exonération n'est pas appliquée automatiquement par l'assureur. "
        "Le bénéficiaire doit fournir les justificatifs à l'assureur au moment du rachat "
        "**ET** mentionner l'exonération dans sa déclaration de revenus (case spécifique)."
    )


def _tab_rente_viagere() -> None:
    """Simulation de transformation du capital en rente viagère."""
    st.markdown("#### Simulation — Sortie en rente viagère")
    st.info(
        "La rente viagère issue d'un contrat d'assurance-vie bénéficie d'une "
        "fiscalité allégée. Seule une fraction de la rente est imposable à l'IR, "
        "selon l'âge au moment de la mise en rente (Art. 158-6 CGI)."
    )
    # Fractions imposables selon l'âge (Art. 158-6 CGI)
    _FRACTIONS_IMPOSABLES = {
        "Moins de 50 ans":     0.70,
        "50 à 59 ans":         0.50,
        "60 à 69 ans":         0.40,
        "70 ans et plus":      0.30,
    }
    _rc1, _rc2 = st.columns(2)
    with _rc1:
        _capital_rente = st.number_input(
            "Capital à convertir en rente (€)",
            min_value=0.0, max_value=5_000_000.0,
            value=200_000.0, step=5_000.0, key="rente_capital",
        )
        _taux_rente_annuel = st.number_input(
            "Taux de conversion annuel estimé (%)",
            min_value=1.0, max_value=10.0,
            value=4.5, step=0.1, key="rente_taux",
            help="Dépend de l'âge, du sexe et de la table de mortalité de l'assureur. "
                 "Typiquement 4-6% pour un rentier de 65 ans.",
        )
    with _rc2:
        _tranche_age = st.selectbox(
            "Tranche d'âge au moment de la mise en rente",
            list(_FRACTIONS_IMPOSABLES.keys()),
            index=2,  # 60-69 ans par défaut
            key="rente_age",
        )
        _tmi = st.selectbox(
            "Taux marginal d'imposition (TMI)",
            [0.11, 0.30, 0.41, 0.45],
            index=1,
            format_func=lambda x: f"{x*100:.0f}%",
            key="rente_tmi",
        )
    _rente_brute_annuelle = _capital_rente * (_taux_rente_annuel / 100.0)
    _rente_brute_mensuelle = _rente_brute_annuelle / 12.0
    _fraction_imposable = _FRACTIONS_IMPOSABLES[_tranche_age]
    _base_ir = _rente_brute_annuelle * _fraction_imposable
    _ir_annuel = _base_ir * _tmi
    _ps_annuel = _base_ir * 0.172  # PS à 17,2%
    _rente_nette_annuelle = _rente_brute_annuelle - _ir_annuel - _ps_annuel
    _rente_nette_mensuelle = _rente_nette_annuelle / 12.0
    _taux_prelevement_effectif = (_ir_annuel + _ps_annuel) / _rente_brute_annuelle if _rente_brute_annuelle > 0 else 0.0
    st.markdown("---")
    _rm1, _rm2, _rm3, _rm4 = st.columns(4)
    _rm1.metric("Rente brute mensuelle", f"{_rente_brute_mensuelle:,.0f} €")
    _rm2.metric("Rente nette mensuelle", f"{_rente_nette_mensuelle:,.0f} €")
    _rm3.metric("Fraction imposable", f"{_fraction_imposable*100:.0f}%")
    _rm4.metric("Taux réel de prélèvement", f"{_taux_prelevement_effectif*100:.1f}%")
    st.markdown("---")
    _detail_rows = [
        ("Rente brute annuelle", f"{_rente_brute_annuelle:,.0f} €"),
        (f"Fraction imposable ({_fraction_imposable*100:.0f}%)", f"{_base_ir:,.0f} €"),
        (f"IR ({_tmi*100:.0f}% × fraction imposable)", f"−{_ir_annuel:,.0f} €"),
        ("Prélèvements sociaux (17,2% × fraction imposable)", f"−{_ps_annuel:,.0f} €"),
        ("Rente nette annuelle", f"{_rente_nette_annuelle:,.0f} €"),
        ("Rente nette mensuelle", f"{_rente_nette_mensuelle:,.0f} €"),
    ]
    st.dataframe(
        pd.DataFrame(_detail_rows, columns=["", "Montant"]),
        hide_index=True, use_container_width=True,
    )
    st.caption(
        "⚠️ Simulation indicative. Le taux de conversion réel dépend de la table de mortalité "
        "TGH05/TGF05, du type de rente (simple, réversible, avec annuités garanties) et de "
        "l'assureur. Consultez un conseiller pour un chiffrage personnalisé."
    )


def render_tax_module() -> None:
    """Module Fiscalité & Avantages AV — entièrement autonome."""
    # Constantes fiscales (locales, non globales)
    PS_RATE            = 0.172
    PFU_RATE           = 0.128
    PFL_LOW_RATE       = 0.075
    PFL_HIGH_RATE      = 0.128
    ABATTEMENT_SEUL    = 4_600.0
    ABATTEMENT_COUPLE  = 9_200.0
    SEUIL_150K         = 150_000.0
    ABATTEMENT_990I    = 152_500.0
    TAUX_990I_TRANCHE1 = 0.20
    TAUX_990I_TRANCHE2 = 0.3125
    SEUIL_990I_T2      = 700_000.0
    ABATTEMENT_757B    = 30_500.0
    AGE_SEUIL          = 70

    # Initialisation session_state
    st.session_state.setdefault("tax_date_ouverture", date(2016, 1, 2))
    st.session_state.setdefault("tax_valeur_contrat", 100_000.0)
    st.session_state.setdefault("tax_versements_nets", 80_000.0)
    st.session_state.setdefault("tax_versements_nets_total", 80_000.0)
    st.session_state.setdefault("tax_situation_familiale", "Célibataire / veuf / divorcé")
    st.session_state.setdefault("tax_montant_brut", 10_000.0)
    st.session_state.setdefault("tax_montant_net", 9_000.0)
    st.session_state.setdefault("tax_last_edited", "brut")
    st.session_state.setdefault("tax_rachat_result", None)
    st.session_state.setdefault("tax_nb_beneficiaires", 2)
    st.session_state.setdefault("tax_regime_versements", "Principalement avant 70 ans (Art. 990I)")
    st.session_state.setdefault("tax_capital_deces", 200_000.0)
    st.session_state.setdefault("tax_abattement_deja_utilise", 0.0)
    st.session_state.setdefault("tax_nb_annees", 10)

    st.title("Fiscalité & Avantages de l'assurance-vie")
    st.warning(
        "⚠️ Simulation indicative à titre pédagogique. Ne constitue pas un conseil "
        "fiscal ou juridique. Consultez un professionnel pour toute décision."
    )

    # ── Sélecteur de portefeuille ──────────────────────────────────
    _has_A = bool(st.session_state.get("A_lines", []))
    _has_B = bool(st.session_state.get("B_lines", []))
    _options_pf = []
    if _has_A:
        _options_pf.append("🧍 Portefeuille Client")
    if _has_B:
        _options_pf.append("🏢 Portefeuille Cabinet")
    if not _options_pf:
        _options_pf = ["🧍 Portefeuille Client"]

    _selected_pf = st.radio(
        "Portefeuille analysé",
        _options_pf,
        horizontal=True,
        key="tax_portfolio_selector",
    )
    _is_client = "Client" in _selected_pf
    st.session_state["tax_is_client"] = _is_client

    # ── Synchronisation depuis le comparateur ──────────────────────
    if _is_client:
        _sync_val = float(st.session_state.get("_LAST_VAL_A", 0.0) or 0.0)
        _sync_net = float(st.session_state.get("_LAST_NET_A", 0.0) or 0.0)
        _sync_xirr = st.session_state.get("_LAST_XIRR_A")
        _sync_contract = st.session_state.get("CONTRACT_LABEL_A", "")
        _sync_euro_rate = float(st.session_state.get("EURO_RATE_A", 2.0))
    else:
        _sync_val = float(st.session_state.get("_LAST_VAL_B", 0.0) or 0.0)
        _sync_net = float(st.session_state.get("_LAST_NET_B", 0.0) or 0.0)
        _sync_xirr = st.session_state.get("_LAST_XIRR_B")
        _sync_contract = st.session_state.get("CONTRACT_LABEL_B", "")
        _sync_euro_rate = float(st.session_state.get("EURO_RATE_B", 2.5))

    if _sync_val > 0:
        st.session_state["tax_valeur_contrat"] = _sync_val
        st.session_state["tax_capital_deces"] = _sync_val
    if _sync_net > 0:
        st.session_state["tax_versements_nets"] = _sync_net
        st.session_state["tax_versements_nets_total"] = _sync_net

    st.session_state["_tax_sync_contract"] = _sync_contract
    st.session_state["_tax_sync_xirr"] = _sync_xirr
    st.session_state["_tax_sync_euro_rate"] = _sync_euro_rate

    tab_dashboard, tab_retrait, tab_transmission_v2, tab_avantages = st.tabs([
        "📋 Tableau de bord",
        "💶 Je veux retirer de l'argent",
        "🏠 Je veux transmettre",
        "⚖️ Pourquoi l'assurance-vie ?",
    ])

    with tab_dashboard:
        _tab_dashboard_fiscal()
    with tab_retrait:
        _tab_retrait_wrapper()
    with tab_transmission_v2:
        _tab_transmission_v2()
    with tab_avantages:
        _tab_avantages_av()

    st.markdown("---")
    with st.expander("📖 Cas d'exonération & Sortie en rente", expanded=False):
        _ex1, _ex2 = st.tabs(["Cas d'exonération", "Sortie en rente"])
        with _ex1:
            _tab_exoneration()
        with _ex2:
            _tab_rente_viagere()


def _tab_dashboard_fiscal():
    """Hub central : infos contrat + situation perso + bénéficiaires + frise."""
    st.markdown("#### 📋 Votre contrat en un coup d'œil")

    _c1, _c2 = st.columns(2)
    with _c1:
        _contract = st.session_state.get("_tax_sync_contract", "") or st.session_state.get("CONTRACT_LABEL_A", "")
        if _contract:
            st.markdown(f"**Contrat** : {_contract}")
        _date_ouv = st.date_input(
            "Date d'ouverture du contrat",
            value=st.session_state.get("tax_date_ouverture", date(2016, 1, 2)),
            key="tax_dash_date_ouverture",
            help="Date réelle d'ouverture du contrat.",
        )
        st.session_state["tax_date_ouverture"] = _date_ouv

        _sit_fam = st.selectbox(
            "Situation familiale",
            ["Célibataire / veuf / divorcé", "Marié / pacsé"],
            key="tax_dash_situation_familiale",
        )
        st.session_state["tax_situation_familiale"] = _sit_fam

        _age = st.number_input(
            "Âge du souscripteur",
            min_value=18, max_value=100, value=55, step=1,
            key="tax_age_souscripteur",
        )

        _tmi = st.selectbox(
            "Taux marginal d'imposition (TMI)",
            [0.0, 0.11, 0.30, 0.41, 0.45],
            index=2,
            format_func=lambda x: f"{int(x*100)}%",
            key="tax_tmi",
            help="Taux marginal de l'IR. Utilisé pour estimer les économies fiscales.",
        )

    with _c2:
        _val = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
        _net = float(st.session_state.get("tax_versements_nets", 80_000.0))
        _pv = _val - _net
        _pv_pct = (_pv / _net * 100) if _net > 0 else 0
        _xirr = st.session_state.get("_tax_sync_xirr")

        st.metric("Valeur actuelle du contrat", to_eur(_val))
        st.metric("Versements nets", to_eur(_net))
        st.metric("Plus-values latentes", to_eur(_pv), delta=f"{_pv_pct:+.1f}%")
        if _xirr is not None:
            st.metric("Rendement annualisé (XIRR)", f"{float(_xirr):+.2f}%")

    # ── Frise chronologique maturité fiscale ──
    st.markdown("---")
    _anciennete = (_date_ouv.toordinal() - date.today().toordinal()) * -1 / 365.25
    _anciennete = max(0.0, _anciennete)
    _ans = int(_anciennete)
    _mois = int((_anciennete - _ans) * 12)

    # Frise simple via Altair
    _ouv_ts = pd.Timestamp(_date_ouv)
    _today_ts = pd.Timestamp(date.today())
    _4ans_ts = _ouv_ts + pd.DateOffset(years=4)
    _8ans_ts = _ouv_ts + pd.DateOffset(years=8)
    _12ans_ts = _ouv_ts + pd.DateOffset(years=12)

    _frise_pts = pd.DataFrame([
        {"label": "Ouverture", "date": _ouv_ts, "type": "jalon"},
        {"label": "4 ans", "date": _4ans_ts, "type": "jalon"},
        {"label": "8 ans (maturité)", "date": _8ans_ts, "type": "jalon_majeur"},
    ])
    _aujourd_hui_df = pd.DataFrame([
        {"label": "Aujourd'hui", "date": _today_ts, "type": "position"},
    ])

    _pos_color = "#2E7D32" if _anciennete >= 8 else ("#FF8F00" if _anciennete >= 4 else "#C62828")

    _rule_df = pd.DataFrame({"x": [_ouv_ts], "x2": [_12ans_ts], "y": [0.5]})
    _rule = (
        alt.Chart(_rule_df)
        .mark_rule(strokeWidth=3, color="#CCCCCC")
        .encode(x=alt.X("x:T"), x2="x2:T", y=alt.value(40))
    )

    _pts = (
        alt.Chart(_frise_pts)
        .mark_point(size=120, filled=True, color="#1B2A4A")
        .encode(
            x=alt.X("date:T", axis=alt.Axis(format="%Y", title="")),
            y=alt.value(40),
            tooltip=["label:N", alt.Tooltip("date:T", format="%d/%m/%Y")],
        )
    )
    _lbl = (
        alt.Chart(_frise_pts)
        .mark_text(dy=-18, fontSize=11, fontWeight="bold", color="#1B2A4A")
        .encode(x="date:T", y=alt.value(40), text="label:N")
    )
    _pos_pt = (
        alt.Chart(_aujourd_hui_df)
        .mark_point(size=200, filled=True, shape="triangle-down")
        .encode(
            x="date:T",
            y=alt.value(40),
            color=alt.value(_pos_color),
            tooltip=["label:N"],
        )
    )
    _pos_lbl = (
        alt.Chart(_aujourd_hui_df)
        .mark_text(dy=20, fontSize=11, fontWeight="bold")
        .encode(x="date:T", y=alt.value(40), text="label:N", color=alt.value(_pos_color))
    )

    _frise_chart = (_rule + _pts + _lbl + _pos_pt + _pos_lbl).properties(height=80, width="container").configure_view(strokeOpacity=0)
    st.altair_chart(_frise_chart, use_container_width=True)

    # ── Statut de maturité ──
    if _anciennete >= 8:
        _emoji_mat, _label_mat = "🟢", "Maturité fiscale atteinte"
        _abat = "9 200 €" if ("Marié" in _sit_fam or "pacsé" in _sit_fam) else "4 600 €"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Régime fiscal optimal. "
            f"Abattement annuel de **{_abat}** sur les gains en cas de rachat. "
            f"Taux réduit de 7,5% (PFL) pour les versements ≤ 150 000 €."
        )
    elif _anciennete >= 4:
        _restant = 8 - _anciennete
        _r_ans, _r_mois = int(_restant), int((_restant - int(_restant)) * 12) + 1
        _emoji_mat, _label_mat = "🟠", "Maturité en cours"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Maturité fiscale dans "
            f"**{_r_ans} an(s) {_r_mois} mois**. Rachats soumis au PFU 30% en attendant."
        )
    else:
        _restant = 8 - _anciennete
        _r_ans, _r_mois = int(_restant), int((_restant - int(_restant)) * 12) + 1
        _emoji_mat, _label_mat = "🔴", "Contrat récent"
        _detail = (
            f"Contrat de **{_ans} ans {_mois} mois** — Flat tax 30% sur les rachats. "
            f"Maturité fiscale dans **{_r_ans} ans {_r_mois} mois**."
        )
    st.markdown(f"### {_emoji_mat} {_label_mat}")
    st.markdown(_detail)

    # ── Situation familiale & Clause bénéficiaire ──
    st.markdown("---")
    st.markdown("#### 👨‍👩‍👧‍👦 Situation familiale & Clause bénéficiaire")

    _fc1, _fc2 = st.columns(2)
    with _fc1:
        _nb_enfants = st.number_input(
            "Nombre d'enfants",
            min_value=0, max_value=10,
            value=int(st.session_state.get("tax_nb_enfants", 2)),
            step=1, key="tax_nb_enfants",
            help="Pour le calcul des droits de succession classique (abattement 100 000€ par enfant).",
        )
    with _fc2:
        _nb_benef = st.number_input(
            "Nombre de bénéficiaires (clause AV)",
            min_value=1, max_value=5,
            value=int(st.session_state.get("tax_nb_beneficiaires", 2)),
            step=1, key="tax_nb_beneficiaires",
        )

    _TYPES_BENEF = ["Conjoint/PACS", "Enfant", "Frère/Sœur", "Neveu/Nièce", "Tiers"]
    _benef_data = []
    for i in range(int(_nb_benef)):
        with st.expander(f"Bénéficiaire {i + 1}", expanded=(i == 0)):
            _bc1, _bc2, _bc3 = st.columns([2, 2, 1])
            with _bc1:
                _nom = st.text_input("Nom (facultatif)", key=f"tax_dash_benef_nom_{i}", value=st.session_state.get(f"tax_dash_benef_nom_{i}", ""))
            with _bc2:
                _type = st.selectbox("Lien", _TYPES_BENEF, key=f"tax_dash_benef_type_{i}")
            with _bc3:
                _part = st.number_input(
                    "Quote-part %", min_value=0.0, max_value=100.0,
                    value=round(100.0 / int(_nb_benef), 0), step=5.0,
                    key=f"tax_dash_benef_part_{i}",
                )
            _benef_data.append({"nom": _nom, "type": _type, "part": float(_part)})

    st.session_state["tax_dash_beneficiaires"] = _benef_data

    _total_parts = sum(b["part"] for b in _benef_data)
    if abs(_total_parts - 100.0) > 0.5:
        st.error(f"⚠️ La somme des quotes-parts est {_total_parts:.1f}% — elle doit être égale à 100%.")

    # ── Transmission en bref ──
    st.markdown("---")
    st.markdown("#### 🏠 Transmission — En bref")
    if _age < 70:
        _abat_trans = 152_500 * int(_nb_benef)
        st.success(
            f"✅ Vous avez moins de 70 ans — **Art. 990I** : **{to_eur(152_500)}** "
            f"exonérés par bénéficiaire. Avec **{int(_nb_benef)} bénéficiaire(s)**, "
            f"jusqu'à **{to_eur(_abat_trans)}** transmis hors succession."
        )
        _marge = _abat_trans - _val
        if _marge > 0:
            st.info(f"💡 Marge de versement avant plafond d'exonération : **{to_eur(_marge)}**")
    else:
        st.warning(
            "⚠️ Après 70 ans — **Art. 757B** : abattement global de **30 500 €**. "
            "Les **intérêts et plus-values** restent exonérés de droits de succession."
        )


def _tab_retrait_wrapper():
    """Regroupe : combien je paie / combien je peux retirer sans impôt / dois-je attendre."""
    st.markdown("#### 💶 Retirer de l'argent — Simulateur de rachat")

    _q1, _q2, _q3 = st.tabs([
        "Combien je paie si je retire X€ ?",
        "Combien puis-je retirer sans impôt ?",
        "Ai-je intérêt à attendre ?",
    ])

    with _q1:
        _tab_rachat()

    with _q2:
        _tab_optimisation_abattement()

    with _q3:
        _date_ouv_att = st.session_state.get("tax_date_ouverture", date(2016, 1, 2))
        _anciennete_att = (date.today() - _date_ouv_att).days / 365.25

        if _anciennete_att >= 8:
            st.success(
                "✅ Votre contrat a déjà atteint sa maturité fiscale (>8 ans). "
                "Aucun intérêt fiscal à attendre davantage pour effectuer un rachat."
            )
        else:
            _restant_att = 8 - _anciennete_att
            _r_ans_att = int(_restant_att)
            _r_mois_att = int((_restant_att - _r_ans_att) * 12) + 1
            st.markdown(f"**Maturité fiscale dans : {_r_ans_att} an(s) et {_r_mois_att} mois**")

            _montant_test = st.number_input(
                "Montant du rachat envisagé (€)",
                min_value=0.0, max_value=5_000_000.0,
                value=30_000.0, step=1_000.0, key="attente_montant",
            )

            _val_att = float(st.session_state.get("tax_valeur_contrat", 100_000.0))
            _net_att = float(st.session_state.get("tax_versements_nets", 80_000.0))
            _net_total_att = float(st.session_state.get("tax_versements_nets_total", _net_att))
            _sit_att = st.session_state.get("tax_situation_familiale", "Célibataire / veuf / divorcé")

            if _montant_test > 0 and _val_att > 0:
                _gains_now = calc_quote_part_gains(_val_att, _net_att, _montant_test)
                _gains_after = calc_quote_part_gains(_val_att, _net_att, _montant_test)
                _res_now = calc_imposition_rachat(
                    _gains_now, _anciennete_att, _sit_att, _net_total_att, _montant_test
                )
                _res_after = calc_imposition_rachat(
                    _gains_after, 8.1, _sit_att, _net_total_att, _montant_test
                )
                _eco = _res_now["total_impots"] - _res_after["total_impots"]

                _ca1, _ca2 = st.columns(2)
                with _ca1:
                    st.markdown("**Rachat maintenant**")
                    st.metric("Impôt + PS", to_eur(_res_now["total_impots"]))
                    st.metric("Net perçu", to_eur(_res_now["net_percu"]))
                with _ca2:
                    st.markdown(f"**Rachat après maturité** ({_r_ans_att}a {_r_mois_att}m)")
                    st.metric("Impôt + PS", to_eur(_res_after["total_impots"]))
                    st.metric("Net perçu", to_eur(_res_after["net_percu"]))

                if _eco > 0:
                    st.success(
                        f"💰 **En attendant {_r_ans_att} an(s) et {_r_mois_att} mois**, vous économisez "
                        f"**{to_eur(_eco)}** d'impôt sur ce rachat de {to_eur(_montant_test)}."
                    )
                    # Bar chart comparatif
                    _bar_att_df = pd.DataFrame([
                        {"Scénario": "Rachat\nmaintenant", "Impôt (€)": _res_now["total_impots"], "Type": "Maintenant"},
                        {"Scénario": f"Rachat après\nmaturité ({_r_ans_att}a)", "Impôt (€)": _res_after["total_impots"], "Type": "Après"},
                    ])
                    _bar_att = (
                        alt.Chart(_bar_att_df)
                        .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
                        .encode(
                            x=alt.X("Scénario:N", sort=None, axis=alt.Axis(labelAngle=0), title=""),
                            y=alt.Y("Impôt (€):Q", axis=alt.Axis(format=",.0f"), title="Impôt total (€)"),
                            color=alt.Color("Type:N", scale=alt.Scale(
                                domain=["Maintenant", "Après"], range=["#E53935", "#2E7D32"]
                            ), legend=None),
                            tooltip=[alt.Tooltip("Scénario:N"), alt.Tooltip("Impôt (€):Q", format=",.0f")],
                        )
                        .properties(height=250, title=f"Économie : {to_eur(_eco)}")
                    )
                    st.altair_chart(_bar_att, use_container_width=True)
                else:
                    st.info("L'écart fiscal est négligeable dans votre situation.")


def _tab_transmission_v2():
    """Transmission enrichie : calcul existant + comparaison AV vs succession + projection."""

    _q1, _q2, _q3 = st.tabs([
        "Combien mes proches recevront-ils ?",
        "AV vs succession classique",
        "Ai-je intérêt à verser avant 70 ans ?",
    ])

    with _q1:
        _tab_transmission()

    with _q2:
        st.markdown("#### AV vs succession classique — Comparaison en euros")
        st.caption(
            "Combien vos bénéficiaires recevraient-ils via l'assurance-vie "
            "vs via une succession classique, pour le même capital ?"
        )

        _capital = float(st.session_state.get("tax_capital_deces", 200_000.0))
        _nb = int(st.session_state.get("tax_nb_beneficiaires", 2))
        _age = st.session_state.get("tax_age_souscripteur", 55)

        if _capital <= 0 or _nb <= 0:
            st.info("Renseignez la valeur du contrat et le nombre de bénéficiaires dans l'onglet 'Tableau de bord'.")
            return

        _part_par_benef = _capital / _nb

        # ── Via l'AV (art. 990I si <70 ans) ──
        _abat_av = 152_500.0 if _age < 70 else 30_500.0 / _nb
        _taxable_av = max(0.0, _part_par_benef - _abat_av)
        if _taxable_av <= 700_000:
            _droits_av = _taxable_av * 0.20
        else:
            _droits_av = 700_000 * 0.20 + (_taxable_av - 700_000) * 0.3125
        _droits_av_total = _droits_av * _nb
        _net_av = _capital - _droits_av_total

        # ── Via succession classique (barème ligne directe simplifié) ──
        def _droits_succession_ligne_directe(taxable: float) -> float:
            if taxable <= 0:
                return 0.0
            _tranches = [
                (8_072, 0.05), (12_109, 0.10), (15_932, 0.15),
                (552_324, 0.20), (902_838, 0.30), (1_805_677, 0.40),
                (float("inf"), 0.45),
            ]
            _total = 0.0
            _prev = 0.0
            for _limit, _rate in _tranches:
                _tranche = min(taxable, _limit) - _prev
                if _tranche <= 0:
                    break
                _total += _tranche * _rate
                _prev = _limit
            return _total

        _nb_enfants_succ = max(1, int(st.session_state.get("tax_nb_enfants", _nb)))
        _abat_succ = 100_000.0
        _part_par_enfant_succ = _capital / _nb_enfants_succ
        _taxable_succ = max(0.0, _part_par_enfant_succ - _abat_succ)
        _droits_succ = _droits_succession_ligne_directe(_taxable_succ)
        _droits_succ_total = _droits_succ * _nb_enfants_succ
        _net_succ = _capital - _droits_succ_total

        _eco_av = _droits_succ_total - _droits_av_total

        _col_av, _col_vs2, _col_succ = st.columns([5, 1, 5])
        with _col_av:
            st.markdown("**🛡️ Via l'assurance-vie**")
            st.metric("Capital transmis", to_eur(_capital))
            st.metric("Droits dus", to_eur(_droits_av_total))
            st.metric("Net reçu par les bénéficiaires", to_eur(_net_av))
            if _age < 70:
                st.caption(f"Art. 990I — Abattement {to_eur(152_500)} par bénéficiaire")
            else:
                st.caption(f"Art. 757B — Abattement global {to_eur(30_500)}")
        with _col_vs2:
            st.markdown(
                "<div style='text-align:center; padding-top:4rem; font-size:1.2rem; color:#aaa;'>vs</div>",
                unsafe_allow_html=True,
            )
        with _col_succ:
            st.markdown("**⚖️ Succession classique**")
            st.metric("Capital transmis", to_eur(_capital))
            st.metric("Droits dus", to_eur(_droits_succ_total))
            st.metric("Net reçu par les bénéficiaires", to_eur(_net_succ))
            st.caption(f"Abattement {to_eur(100_000)} par enfant en ligne directe")

        if _eco_av > 0:
            st.success(
                f"✅ **L'assurance-vie économise {to_eur(_eco_av)} de droits** "
                f"par rapport à une succession classique pour {_nb} bénéficiaire(s). "
                f"Vos proches reçoivent {to_eur(_net_av)} au lieu de {to_eur(_net_succ)}."
            )
        else:
            st.info("Dans cette configuration, la succession classique est aussi avantageuse que l'AV.")

        st.caption(
            "⚠️ Calcul simplifié en ligne directe (enfants). Les droits réels dépendent "
            "du patrimoine global, des donations antérieures et du lien de parenté."
        )

        # Stacked bar AV vs Succession
        _bar_trans = pd.DataFrame([
            {"Voie": "Assurance-vie", "Composante": "Net reçu", "Montant": _net_av},
            {"Voie": "Assurance-vie", "Composante": "Droits dus", "Montant": _droits_av_total},
            {"Voie": "Succession classique", "Composante": "Net reçu", "Montant": _net_succ},
            {"Voie": "Succession classique", "Composante": "Droits dus", "Montant": _droits_succ_total},
        ])
        _stacked_trans = (
            alt.Chart(_bar_trans)
            .mark_bar(cornerRadiusTopLeft=3, cornerRadiusTopRight=3)
            .encode(
                x=alt.X("Voie:N", sort=None, axis=alt.Axis(labelAngle=0), title=""),
                y=alt.Y("Montant:Q", stack="zero", axis=alt.Axis(format=",.0f"), title="€"),
                color=alt.Color(
                    "Composante:N",
                    scale=alt.Scale(
                        domain=["Net reçu", "Droits dus"],
                        range=["#2E7D32", "#E53935"],
                    ),
                    legend=alt.Legend(title="", orient="bottom"),
                ),
                order=alt.Order("Composante:N", sort="descending"),
                tooltip=[
                    alt.Tooltip("Voie:N"),
                    alt.Tooltip("Composante:N"),
                    alt.Tooltip("Montant:Q", format=",.0f", title="€"),
                ],
            )
            .properties(height=350)
        )
        st.altair_chart(_stacked_trans, use_container_width=True)

    with _q3:
        st.markdown("#### Ai-je intérêt à verser avant 70 ans ?")

        _age = st.session_state.get("tax_age_souscripteur", 55)
        _nb = int(st.session_state.get("tax_nb_beneficiaires", 2))
        _val_actuelle = float(st.session_state.get("tax_valeur_contrat", 100_000.0))

        if _age >= 70:
            st.warning(
                "Vous avez déjà dépassé 70 ans. Les nouveaux versements relèvent de "
                "l'article 757B (abattement global de 30 500€). "
                "Les **intérêts et plus-values** sur ces versements restent toutefois exonérés."
            )
            return

        _ans_avant_70 = 70 - _age
        _plafond_exo = 152_500 * _nb
        _marge = max(0, _plafond_exo - _val_actuelle)

        st.markdown(f"**Vous avez {_age} ans — encore {_ans_avant_70} ans pour verser avant 70 ans.**")
        st.metric("Plafond d'exonération total (990I)", to_eur(_plafond_exo),
                  delta=f"{_nb} bénéficiaire(s) × {to_eur(152_500)}")
        st.metric("Valeur actuelle du contrat", to_eur(_val_actuelle))
        st.metric("Marge de versement avant plafond", to_eur(_marge))

        st.markdown("---")
        st.markdown("**Projection : que devient votre versement ?**")

        _xirr = st.session_state.get("_LAST_XIRR_A") or st.session_state.get("_LAST_XIRR_B")
        _default_rdt = float(_xirr) if (_xirr and _xirr > 0) else 5.0

        _p1, _p2 = st.columns(2)
        with _p1:
            _versement_proj = st.number_input(
                "Versement envisagé (€)",
                min_value=0.0, max_value=5_000_000.0,
                value=min(_marge, 100_000.0) if _marge > 0 else 50_000.0,
                step=5_000.0, key="proj_versement",
            )
        with _p2:
            _rdt_proj = st.number_input(
                "Rendement annuel estimé (%)",
                min_value=-5.0, max_value=20.0,
                value=round(_default_rdt, 1), step=0.5,
                key="proj_rendement",
                help="Pré-rempli avec le rendement réel du portefeuille (XIRR) si disponible.",
            )

        _horizon = _ans_avant_70
        if _horizon > 0 and _versement_proj > 0:
            _capital_futur = _versement_proj * (1 + _rdt_proj / 100) ** _horizon
            _gains = _capital_futur - _versement_proj

            st.markdown("---")
            _pc1, _pc2, _pc3 = st.columns(3)
            with _pc1:
                st.metric("Versement aujourd'hui", to_eur(_versement_proj))
            with _pc2:
                st.metric(f"Valeur dans {_horizon} ans", to_eur(_capital_futur))
            with _pc3:
                st.metric("Plus-values générées", to_eur(_gains),
                          delta=f"+{(_capital_futur / _versement_proj - 1) * 100:.0f}%")

            if _capital_futur <= _plafond_exo:
                st.success(
                    f"✅ Si vous versez **{to_eur(_versement_proj)}** aujourd'hui et que le portefeuille "
                    f"fait **{_rdt_proj:.1f}%/an**, dans **{_horizon} ans** ce sera "
                    f"**{to_eur(_capital_futur)}** transmis à vos {_nb} bénéficiaire(s) — "
                    f"**entièrement exonéré de droits de succession** (art. 990I)."
                )
            else:
                st.info(
                    f"Le capital projeté ({to_eur(_capital_futur)}) dépasse le plafond d'exonération "
                    f"({to_eur(_plafond_exo)}). La part au-delà sera taxée à 20% puis 31.25%."
                )

            # Area chart projection + seuil exonération
            _proj_years = list(range(_horizon + 1))
            _proj_vals = [_versement_proj * (1 + _rdt_proj / 100) ** y for y in _proj_years]
            _proj_df = pd.DataFrame({"Année": _proj_years, "Valeur (€)": _proj_vals})

            _area_proj = (
                alt.Chart(_proj_df)
                .mark_area(opacity=0.25, color="#1B2A4A", interpolate="monotone")
                .encode(
                    x=alt.X("Année:Q", title="Années après versement"),
                    y=alt.Y("Valeur (€):Q", axis=alt.Axis(format=",.0f"), title="Valeur (€)"),
                )
            )
            _line_proj = (
                alt.Chart(_proj_df)
                .mark_line(strokeWidth=2.5, color="#1B2A4A", interpolate="monotone")
                .encode(
                    x="Année:Q",
                    y="Valeur (€):Q",
                    tooltip=[alt.Tooltip("Année:Q"), alt.Tooltip("Valeur (€):Q", format=",.0f")],
                )
            )
            _seuil_df = pd.DataFrame([{"Année": 0, "Seuil": _plafond_exo}, {"Année": _horizon, "Seuil": _plafond_exo}])
            _seuil_line = (
                alt.Chart(_seuil_df)
                .mark_line(strokeDash=[5, 5], strokeWidth=2, color="#E53935")
                .encode(x="Année:Q", y=alt.Y("Seuil:Q"))
            )
            _seuil_lbl_df = pd.DataFrame([{
                "Année": _horizon * 0.5, "Seuil": _plafond_exo,
                "text": f"Seuil exonération : {to_eur(_plafond_exo)}",
            }])
            _seuil_label = (
                alt.Chart(_seuil_lbl_df)
                .mark_text(dy=-12, fontSize=11, color="#E53935", fontWeight="bold")
                .encode(x="Année:Q", y="Seuil:Q", text="text:N")
            )
            st.altair_chart(
                (_area_proj + _line_proj + _seuil_line + _seuil_label).properties(height=320),
                use_container_width=True,
            )


def _tab_avantages_av():
    """Comparaison AV vs CTO + tableau AV vs PER."""

    _q1, _q2 = st.tabs([
        "AV vs Compte-titres (CTO)",
        "AV vs PER",
    ])

    with _q1:
        st.markdown("#### L'effet capitalisation : AV vs CTO")
        st.caption(
            "Même capital, même performance. Mais en AV la capitalisation se fait "
            "sans frottement fiscal (taxé uniquement à la sortie). "
            "En CTO, chaque arbitrage est taxé à 30%."
        )

        _xirr = st.session_state.get("_LAST_XIRR_A") or st.session_state.get("_LAST_XIRR_B")
        _default_rdt = float(_xirr) if (_xirr and _xirr > 0) else 6.0

        _ac1, _ac2 = st.columns(2)
        with _ac1:
            _capital_init = st.number_input(
                "Capital initial (€)", min_value=1_000.0, max_value=5_000_000.0,
                value=100_000.0, step=5_000.0, key="avcto_capital",
            )
            _perf_brute = st.number_input(
                "Performance brute annuelle (%)",
                min_value=0.0, max_value=20.0,
                value=round(_default_rdt, 1), step=0.5, key="avcto_perf",
                help="Pré-rempli avec le rendement réel du portefeuille si disponible.",
            )
        with _ac2:
            _horizon_av = st.slider(
                "Horizon (années)", min_value=5, max_value=40,
                value=20, key="avcto_horizon",
            )
            _freq_arb = st.selectbox(
                "Fréquence d'arbitrage (CTO)",
                ["Annuel", "Tous les 2 ans", "Tous les 5 ans"],
                index=0, key="avcto_freq",
                help="En CTO, chaque arbitrage déclenche la flat tax sur les PV réalisées.",
            )
        _freq_years = {"Annuel": 1, "Tous les 2 ans": 2, "Tous les 5 ans": 5}[_freq_arb]

        # ── Calcul AV : capitalisation sans frottement ──
        _val_av = _capital_init * (1 + _perf_brute / 100) ** _horizon_av
        _gains_av = _val_av - _capital_init
        _abat_av_sortie = 4_600.0
        _gains_taxables_av = max(0.0, _gains_av - _abat_av_sortie)
        _tax_av = _gains_taxables_av * (0.075 + 0.172)
        _net_av = _val_av - _tax_av

        # ── Calcul CTO : frottement à chaque arbitrage ──
        _series_cto = [_capital_init]
        _cto_val = _capital_init
        _cto_base = _capital_init
        for _y in range(1, _horizon_av + 1):
            _cto_val *= (1 + _perf_brute / 100)
            if _y % _freq_years == 0 and _y < _horizon_av:
                _pv = _cto_val - _cto_base
                if _pv > 0:
                    _cto_val -= _pv * 0.30
                    _cto_base = _cto_val
            _series_cto.append(_cto_val)
        _gains_cto = _cto_val - _capital_init
        _tax_cto_final = max(0.0, _gains_cto) * 0.30
        _net_cto = _cto_val - _tax_cto_final

        # ── Graphique ──
        _years = list(range(_horizon_av + 1))
        _series_av_list = [_capital_init * (1 + _perf_brute / 100) ** y for y in _years]
        _chart_df = pd.DataFrame({
            "Année": _years * 2,
            "Valeur (€)": _series_av_list + _series_cto,
            "Enveloppe": (["Assurance-vie (sans frottement)"] * len(_years)
                          + ["Compte-titres (flat tax 30%)"] * len(_years)),
        })

        # Zone d'écart colorée entre AV et CTO
        _ecart_df = pd.DataFrame({
            "Année": _years,
            "AV": _series_av_list,
            "CTO": _series_cto,
        })
        _area_ecart = (
            alt.Chart(_ecart_df.melt("Année", value_vars=["AV", "CTO"], var_name="X", value_name="Y")
               .pivot_table(index="Année", columns="X", values="Y").reset_index())
            .mark_area(opacity=0.12, color="#2E7D32")
            .encode(
                x=alt.X("Année:Q"),
                y=alt.Y("CTO:Q"),
                y2="AV:Q",
            )
        )

        _chart = (
            alt.Chart(_chart_df)
            .mark_line(strokeWidth=2.5)
            .encode(
                x=alt.X("Année:Q", title="Années"),
                y=alt.Y("Valeur (€):Q", axis=alt.Axis(format=",.0f"), title="Valeur (€)"),
                color=alt.Color(
                    "Enveloppe:N",
                    scale=alt.Scale(
                        domain=["Assurance-vie (sans frottement)", "Compte-titres (flat tax 30%)"],
                        range=["#1B2A4A", "#CC2200"],
                    ),
                    legend=alt.Legend(title="", orient="bottom"),
                ),
                tooltip=[
                    alt.Tooltip("Année:Q"),
                    alt.Tooltip("Enveloppe:N"),
                    alt.Tooltip("Valeur (€):Q", format=",.0f"),
                ],
            )
            .properties(height=400)
        )
        st.altair_chart((_area_ecart + _chart).properties(height=400), use_container_width=True)

        _ecart = _net_av - _net_cto
        _ecart_pct = (_ecart / _net_cto * 100) if _net_cto > 0 else 0
        _mc1, _mc2, _mc3 = st.columns(3)
        with _mc1:
            st.metric("Net après impôt — AV", to_eur(_net_av))
        with _mc2:
            st.metric("Net après impôt — CTO", to_eur(_net_cto))
        with _mc3:
            st.metric("Avantage AV", to_eur(_ecart), delta=f"+{_ecart_pct:.0f}%")

        if _ecart > 0:
            st.success(
                f"✅ Sur **{_horizon_av} ans**, l'assurance-vie génère **{to_eur(_ecart)} de plus** "
                f"que le compte-titres, à performance identique. C'est l'effet de la capitalisation "
                f"sans frottement fiscal."
            )

    with _q2:
        st.markdown("#### AV vs PER — Tableau comparatif")
        _comp_data = pd.DataFrame({
            "Critère": [
                "Disponibilité du capital",
                "Fiscalité à l'entrée",
                "Fiscalité à la sortie (capital)",
                "Fiscalité à la sortie (rente)",
                "Transmission",
                "Plafond de versement",
                "Cas de déblocage anticipé",
            ],
            "Assurance-vie": [
                "✅ Totale, à tout moment",
                "Aucune déduction",
                "PFL 7,5% + PS après 8 ans (avec abattement)",
                "Fraction imposable selon l'âge (30-70%)",
                "✅ Hors succession (152 500€/bénéficiaire)",
                "Aucun plafond",
                "Non applicable (capital toujours disponible)",
            ],
            "PER": [
                "❌ Bloqué jusqu'à la retraite",
                "✅ Déductible du revenu imposable",
                "IR + PS sur le capital (si déduit à l'entrée)",
                "IR + PS sur la totalité",
                "❌ Succession classique (sauf exceptions)",
                "Plafonné (10% revenus N-1)",
                "Achat résidence principale, invalidité, décès conjoint...",
            ],
        })
        st.dataframe(_comp_data, hide_index=True, use_container_width=True)
        st.caption(
            "L'AV privilégie la **disponibilité et la transmission**. "
            "Le PER privilégie la **déduction fiscale à l'entrée** (intéressant pour les TMI élevées)."
        )


def run_perfect_portfolio():
    render_portfolio_builder()


def render_mode_router():
    st.set_page_config(page_title=APP_TITLE, layout="wide")
    with st.sidebar:
        _sidebar_title = st.session_state.get("NOM_CABINET", "").strip()
        st.markdown(f"## {_sidebar_title if _sidebar_title else APP_TITLE}")
        st.caption(APP_SUBTITLE)
        st.divider()
        mode = st.radio(
            "Navigation",
            [
                "📊 Comparateur",
                "🏗️ Construction optimisée",
                "🧾 Fiscalité & Avantages AV",
            ],
            label_visibility="collapsed",
        )
    if mode == "📊 Comparateur":
        run_comparator()
    elif mode == "🏗️ Construction optimisée":
        run_perfect_portfolio()
    else:
        render_tax_module()  # "🧾 Fiscalité & Avantages AV"


def _render_with_crash_shield():
    try:
        render_mode_router()
        st.session_state["APP_STATUS"] = "OK"
    except Exception as e:
        st.session_state["APP_STATUS"] = "KO"
        st.session_state["LAST_EXCEPTION"] = str(e)
        st.title(APP_TITLE)
        st.info("App chargée, statut KO")
        st.error("Une erreur est survenue pendant le rendu.")
        st.exception(e)
        st.markdown("""
Conseils :
- Vérifiez vos dépendances (reportlab/matplotlib).
- Vérifiez la clé EODHD dans les secrets.
- Réessayez après avoir vidé le cache Streamlit.
""")


_render_with_crash_shield()
